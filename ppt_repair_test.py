import os
import sys
import re
import traceback
import importlib
import io
from typing import Dict, Any, Tuple

# Force UTF-8 for Windows output if it fails to detect
if sys.stdout.encoding != 'utf-8':
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

# Import project-specific modules
from src.ppt_generator import (
    generate_slide_code, save_slide_code, load_slide_code, 
    _patch_common_errors, _SCRIPT_HEADER, _rename_func,
    SafeProxy, SlideProxy, SlidesListProxy # Import the hardened proxies from production code
)
from src import ppt_generator
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION

# ── Dynamic Patching Sanitizer (Mock production logic) ─────────────

def _sanitize_text(text: Any) -> Any:
    if not isinstance(text, str):
        return text
    # Remove characters outside the Basic Multilingual Plane (0x0000 - 0xFFFF)
    return "".join(c for c in text if ord(c) <= 0xFFFF)

# ── Hardened Execution Engine ───────────────────────────────────────

def execute_with_fallback(code: str, slide_obj: Any, globals_dict: Dict[str, Any]) -> Tuple[bool, str]:
    """Execute code with a proxy and multiple retry/patch attempts."""
    
    # Wrap the slide object in a proxy (ensuring it's the latest SlideProxy)
    proxy_slide = SlideProxy(slide_obj)
    globals_dict['slide'] = proxy_slide
    
    # 1. Initial attempt with current patch logic
    code = _patch_common_errors(code)
    
    attempts = 0
    max_attempts = 3
    current_code = code

    # Local storage for build_slide function
    local_vars = {}
    
    while attempts < max_attempts:
        try:
            exec(current_code, globals_dict, local_vars)
            
            # Use the correct function from locals
            func = None
            if 'build_slide' in local_vars:
                func = local_vars['build_slide']
            else:
                for k in local_vars:
                    if k.startswith('build_slide_'):
                        func = local_vars[k]
                        break
            
            if func:
                func(proxy_slide)
            
            return True, ""
        except Exception as e:
            attempts += 1
            err_msg = traceback.format_exc()
            print(f"Attempt {attempts} failed: {e}")
            
            # ── Intelligent Patching ──
            if "AttributeError" in err_msg:
                # Fix: 'LineFormat' object has no attribute 'background'
                if "'LineFormat' object has no attribute 'background'" in err_msg:
                    current_code = current_code.replace(".line.background()", ".line.fill.background()")
                    continue
                if "'FillFormat' object has no attribute 'background'" in err_msg:
                    current_code = current_code.replace(".fill.background()", ".fill.solid(); .fill.fore_color.rgb = WHITE")
                    continue
            
            if "NameError" in err_msg:
                match = re.search(r"name '(\w+)' is not defined", str(e))
                if match:
                    missing_name = match.group(1)
                    found = False
                    for modname in ["pptx.enum.shapes", "pptx.enum.text", "pptx.enum.chart", "pptx.util"]:
                        try:
                            mod = importlib.import_module(modname)
                            if hasattr(mod, missing_name):
                                globals_dict[missing_name] = getattr(mod, missing_name)
                                found = True
                                break
                        except ImportError: continue
                    if found: continue
            
            if attempts >= max_attempts:
                return False, err_msg
                
    return False, "Unknown failure"

# ── Test Runner ─────────────────────────────────────────────────────

def test_project_ppt_generation(project_name: str):
    print(f"Testing project: {project_name}")
    proj_dir = os.path.join("projects", project_name)
    img_dir = os.path.join(proj_dir, "生成的图片")
    output_pptx = os.path.join(proj_dir, "最终文档", f"{project_name}_test_hardened.pptx")
    
    if not os.path.exists(img_dir):
        print(f"Error: Image directory not found: {img_dir}")
        return

    img_files = sorted([f for f in os.listdir(img_dir) if f.lower().endswith(('.jpg', '.png'))])
    if not img_files:
        print("No images found in project.")
        return

    print(f"Found {len(img_files)} images. Starting per-page generation...")

    # Import the proxy presentation factory from production
    from src.ppt_generator import _exec_script, ProxyPresentation

    execution_globals = {}
    from src.ppt_generator import _SCRIPT_HEADER
    exec(_SCRIPT_HEADER, execution_globals)
    execution_globals.update({
        'Inches': Inches, 'Pt': Pt, 'Emu': Emu,
        'RGBColor': RGBColor, 'PP_ALIGN': PP_ALIGN, 
        'MSO_SHAPE': MSO_SHAPE, 'MSO_CONNECTOR': MSO_CONNECTOR,
        'XL_CHART_TYPE': XL_CHART_TYPE, 'XL_LABEL_POSITION': XL_LABEL_POSITION,
        'CategoryChartData': ppt_generator.CategoryChartData
    })

    results = []
    
    # Create the main master presentation using Proxy
    # Note: We won't set a global background for the whole PRS, but per-slide
    main_prs = Presentation() 
    main_prs.slide_width = Inches(13.333)
    main_prs.slide_height = Inches(7.5)

    for i, img_file in enumerate(img_files):
        page_num = i + 1
        img_path = os.path.abspath(os.path.join(img_dir, img_file))
        print(f"Processing Page {page_num}...")
        
        code = load_slide_code(proj_dir, page_num)
        if not code:
            print(f"  Requesting AI for Page {page_num}...")
            # Use real generation logic
            from src.ppt_generator import generate_slide_code
            code = generate_slide_code(img_path, page_num, len(img_files))
            save_slide_code(proj_dir, page_num, code)
        
        # 1. Main Slide (for the final merged file)
        # We use a Proxy wrapper for this specific slide to add the background
        from src.ppt_generator import SlidesListProxy, _exec_script
        proxy_slides = SlidesListProxy(main_prs.slides, bg_image_path=img_path)
        main_slide = proxy_slides.add_slide(main_prs.slide_layouts[6])
        
        # USE THE PRODUCTION HARDENED ENGINE DIRECTLY
        # This now handles redirection, patching, and proxies internally
        success, err = _exec_script(code, bg_image_path=img_path, redirect_slide=main_slide)
        
        if success:
            print(f"  [OK] Page {page_num} Success")
            results.append((page_num, True, ""))
        else:
            print(f"  [FAIL] Page {page_num} Failed")
            results.append((page_num, False, err))

        # 2. Individual Debug Slide (High Fidelity)
        page_pptx = output_pptx.replace(".pptx", f"_page_{page_num:02d}.pptx")
        temp_prs = Presentation()
        temp_prs.slide_width = main_prs.slide_width
        temp_prs.slide_height = main_prs.slide_height
        
        temp_proxy_slides = SlidesListProxy(temp_prs.slides, bg_image_path=img_path)
        dst_slide = temp_proxy_slides.add_slide(temp_prs.slide_layouts[6])
        
        # Use production engine for individual slide too
        _exec_script(code, bg_image_path=img_path, redirect_slide=dst_slide)
        temp_prs.save(page_pptx)
        print(f"    - Saved individual slide to {os.path.basename(page_pptx)}")

    # 3. Final save
    try:
        main_prs.save(output_pptx)
        print(f"\nFinal PPT saved to: {output_pptx}")
    except Exception as e:
        print(f"Failed to save final PPT: {e}")

    # Summary
    print("\n" + "="*30)
    print("TEST SUMMARY")
    print("="*30)
    for page, success, error in results:
        status = "PASS" if success else "FAIL"
        print(f"Page {page}: {status}")
        if not success:
            print(f"  Error snippet: {error.splitlines()[-1]}")
    print("="*30)

if __name__ == "__main__":
    projects_dir = "projects"
    if os.path.exists(projects_dir):
        projects = [d for d in os.listdir(projects_dir) if os.path.isdir(os.path.join(projects_dir, d))]
        if projects:
            projects.sort(key=lambda x: os.path.getmtime(os.path.join(projects_dir, x)), reverse=True)
            test_project_ppt_generation(projects[0])
        else:
            print("No projects found in projects/ directory.")
    else:
        print("projects/ directory not found.")
