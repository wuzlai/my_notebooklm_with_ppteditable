"""Generate PPTX from infographic images using Gemini to write python-pptx code."""

import os
import re
import traceback
from typing import Any
from pptx.util import Pt

from src.gemini_client import generate_text_with_images
from src.prompts import PPT_CODE_GEN_SYSTEM_PROMPT, PPT_CODE_GEN_USER_PROMPT

DEFAULT_MODEL = "gemini-3-pro-preview"

# ── Shared boilerplate embedded in every generated script ──
_SCRIPT_HEADER = '''\
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# ── Color Palette ──
BLUE_HEADER = RGBColor(0x5B, 0x9B, 0xD5)
BLUE_DARK   = RGBColor(0x4A, 0x86, 0xC8)
CYAN        = RGBColor(0x00, 0xBC, 0xD4)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x33, 0x33, 0x33)
GRAY_TEXT   = RGBColor(0x55, 0x55, 0x55)
GRAY_BAR    = RGBColor(0xB0, 0xBE, 0xC5)
RED         = RGBColor(0xE5, 0x39, 0x35)
GREEN       = RGBColor(0x43, 0xA0, 0x47)
ORANGE      = RGBColor(0xFF, 0x98, 0x00)

ICON_BG  = RGBColor(0xE3, 0xE8, 0xED)
ICON_FG  = RGBColor(0x54, 0x6E, 0x7A)
FONT_NAME = "Microsoft YaHei"
SLIDE_WIDTH = Inches(13.333)
HEADER_H    = Inches(0.75)
SUBTITLE_Y  = Inches(0.95)


def add_header_banner(slide, title_text, bg_color=None):
    if bg_color is None:
        bg_color = BLUE_HEADER
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, HEADER_H
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = bg_color
    banner.line.fill.background()
    tf = banner.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(26)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME


def add_subtitle(slide, text, left, top, width=Inches(12), font_size=Pt(18)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.font.name = FONT_NAME
    return txBox


def add_icon_box(slide, left, top, symbol, size=Inches(0.48)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ICON_BG
    shape.line.fill.background()
    shape.adjustments[0] = 0.25
    tf = shape.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = symbol
    p.font.size = Pt(18)
    p.font.color.rgb = ICON_FG
    p.font.bold = False
    return shape


def add_bullet_item(slide, left, top, symbol, label, description,
                    width=Inches(5.5), desc_size=Pt(13)):
    add_icon_box(slide, left, top, symbol)
    text_left = left + Inches(0.65)
    txBox = slide.shapes.add_textbox(text_left, top - Inches(0.02), width, Inches(0.65))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run_label = p.add_run()
    run_label.text = label + "\\uff1a"
    run_label.font.size = Pt(14)
    run_label.font.color.rgb = BLACK
    run_label.font.bold = True
    run_label.font.name = FONT_NAME
    run_desc = p.add_run()
    run_desc.text = description
    run_desc.font.size = desc_size
    run_desc.font.color.rgb = GRAY_TEXT
    run_desc.font.bold = False
    run_desc.font.name = FONT_NAME
    return txBox


def add_conclusion_box(slide, left, top, width, text, font_size=Pt(13)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = BLACK
    run.font.bold = True
    run.font.name = FONT_NAME
    return txBox


def add_table(slide, left, top, width, height, rows, cols, data,
              header_color=None, col_widths=None):
    if header_color is None:
        header_color = BLUE_HEADER
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if data[r][c] is not None else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(11)
                    run.font.name = FONT_NAME
                    if r == 0:
                        run.font.color.rgb = WHITE
                        run.font.bold = True
                    else:
                        run.font.color.rgb = BLACK
                        run.font.bold = False
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else RGBColor(0xF5, 0xF5, 0xF5)
            cell.margin_left = Pt(5)
            cell.margin_right = Pt(5)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
    return table_shape


def add_bar_chart(slide, left, top, width, height,
                  categories, values, title="", bar_colors=None):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', values)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    chart.chart_style = 2
    plot = chart.plots[0]
    plot.gap_width = 100
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CYAN
    series.has_data_labels = True
    dl = series.data_labels
    dl.font.size = Pt(13)
    dl.font.bold = True
    dl.font.color.rgb = BLACK
    dl.number_format = '0.#'
    dl.show_value = True
    dl.label_position = XL_LABEL_POSITION.OUTSIDE_END
    if bar_colors:
        for i, color in enumerate(bar_colors):
            pt = series.points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = color
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(12)
    cat_axis.tick_labels.font.name = FONT_NAME
    cat_axis.major_tick_mark = 2
    cat_axis.format.line.fill.background()
    val_axis = chart.value_axis
    val_axis.visible = False
    val_axis.major_tick_mark = 2
    val_axis.format.line.fill.background()
    val_axis.major_gridlines.format.line.fill.background()
    if title:
        chart.has_title = True
        ct = chart.chart_title.text_frame.paragraphs[0]
        ct.text = title
        ct.font.size = Pt(14)
        ct.font.bold = True
        ct.font.name = FONT_NAME
    else:
        chart.has_title = False
    return chart_frame


def add_callout_label(slide, left, top, text, bg_color=None, font_size=Pt(11)):
    if bg_color is None:
        bg_color = CYAN
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.3), Inches(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME
    return shape


def add_data_card(slide, left, top, width, height, value, label,
                  value_color=None, bg_color=None):
    if value_color is None:
        value_color = CYAN
    if bg_color is None:
        bg_color = WHITE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(3)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run1 = p1.add_run()
    run1.text = str(value)
    run1.font.size = Pt(24)
    run1.font.color.rgb = value_color
    run1.font.bold = True
    run1.font.name = FONT_NAME
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(10)
    run2.font.color.rgb = GRAY_TEXT
    run2.font.bold = False
    run2.font.name = FONT_NAME
    return shape
'''


def _extract_code(response_text: str) -> str:
    """Extract ONLY the build_slide function from the AI response."""
    # Step 1: extract from markdown code block if present
    match = re.search(r'```python\s*\n(.*?)```', response_text, re.DOTALL)
    code = match.group(1).strip() if match else response_text.strip()

    # Step 2: find the build_slide function and extract only it
    match = re.search(r'(def build_slide\(slide\):.*)', code, re.DOTALL)
    if not match:
        return code

    func_text = match.group(1)
    lines = func_text.split('\n')
    result = [lines[0]]  # "def build_slide(slide):"
    for line in lines[1:]:
        # Stop at the next top-level definition or non-indented code
        # (but allow blank lines and comments inside the function)
        if line and not line[0].isspace() and not line.startswith('#') and line.strip():
            break
        result.append(line)

    # Remove trailing blank lines
    while result and not result[-1].strip():
        result.pop()

    return '\n'.join(result)


def generate_slide_code(
    image_path: str,
    page_num: int,
    total_pages: int,
    model: str = DEFAULT_MODEL,
) -> str:
    """Send an infographic image to Gemini and get back python-pptx code."""
    user_prompt = PPT_CODE_GEN_USER_PROMPT.format(
        page_num=page_num,
        total_pages=total_pages,
    )
    response = generate_text_with_images(
        model=model,
        system_prompt=PPT_CODE_GEN_SYSTEM_PROMPT,
        user_prompt=user_prompt,
        image_paths=[image_path],
    )
    return _extract_code(response)


def _make_pptx_script(build_func_codes: list[tuple[str, str]], output_path: str) -> str:
    """Assemble a full runnable script.

    build_func_codes: list of (func_name, func_body) pairs.
    """
    parts = [_SCRIPT_HEADER]
    parts.append(f'\nOUTPUT_PATH = r"{output_path}"\n')

    slide_calls = []
    for i, (func_name, code) in enumerate(build_func_codes):
        parts.append(f"\n# ── Slide {i + 1} ──\n")
        parts.append(code)
        parts.append("\n")
        slide_calls.append(
            f"s{i} = prs.slides.add_slide(prs.slide_layouts[6])\n"
            f"{func_name}(s{i})"
        )

    parts.append(f"""
# ── Main ──
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

{chr(10).join(slide_calls)}
prs.save(OUTPUT_PATH)
""")
    return "\n".join(parts)


def _rename_func(code: str, new_name: str) -> str:
    """Rename build_slide -> new_name in the code."""
    return code.replace("def build_slide(slide)", f"def {new_name}(slide)")


def build_single_slide_pptx(slide_code: str, output_path: str) -> tuple[bool, str]:
    """Generate a single-slide PPTX from one slide's code. Returns (success, error)."""
    func_name = "build_slide_1"
    renamed = _rename_func(slide_code, func_name)
    script = _make_pptx_script([(func_name, renamed)], output_path)

    # Save assembled script for debugging (use _full suffix to avoid overwriting code file)
    script_path = output_path.replace(".pptx", "_full.py")
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    with open(script_path, "w", encoding="utf-8") as f:
        f.write(script)

    return _exec_script(script)


def build_full_pptx(slide_codes: dict[int, str], output_path: str) -> tuple[bool, str]:
    """Generate a full PPTX from multiple slide codes.

    slide_codes: {page_num: code_string} (1-based page numbers).
    Returns (success, error).
    """
    func_pairs = []
    for page_num in sorted(slide_codes.keys()):
        func_name = f"build_slide_{page_num}"
        renamed = _rename_func(slide_codes[page_num], func_name)
        func_pairs.append((func_name, renamed))

    script = _make_pptx_script(func_pairs, output_path)

    script_path = output_path.replace(".pptx", "_full.py")
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    with open(script_path, "w", encoding="utf-8") as f:
        f.write(script)

    return _exec_script(script)


# ── Safe Proxy Enhancement ──────────────────────────────────────────

def sanitize_text(text: Any) -> Any:
    if not isinstance(text, str):
        return text
    
    def is_safe(c):
        cp = ord(c)
        # 1. Basic ASCII (0x20-0x7E) plus common whitespace
        if 0x20 <= cp <= 0x7E or c in "\n\r\t":
            return True
        # 2. CJK Unified Ideographs (Common Chinese characters)
        if 0x4E00 <= cp <= 0x9FFF:
            return True
        # 3. CJK Symbols and Punctuation (Chinese full-width punctuation)
        if 0x3000 <= cp <= 0x303F:
            return True
        # 4. Hiragana/Katakana (Japanese)
        if 0x3040 <= cp <= 0x30FF:
            return True
        # 5. Full-width alphanumeric (0xFF01-0xFFEE)
        if 0xFF00 <= cp <= 0xFFEF:
            return True
        return False

    # Filter text to keep only safe characters
    return "".join(c for c in text if is_safe(c))

class NullProxy:
    """A 'Black Hole' object that swallows all calls, attributes, and items to prevent crashes."""
    def __getattr__(self, name): return self
    def __setattr__(self, name, value): pass
    def __call__(self, *args, **kwargs): return self
    def __getitem__(self, key): return self
    def __setitem__(self, key, value): pass
    def __len__(self): return 0
    def __iter__(self): return iter([])
    def __bool__(self): return False
    def __repr__(self): return "<NullProxy>"

class SafeProxy:
    """A proxy object to intercept and fix common python-pptx AI errors."""
    def __init__(self, target):
        object.__setattr__(self, "_target", target)

    def _unwrap(self, val):
        if hasattr(val, "_target") and "SafeProxy" in str(type(val)):
            return val._target
        # Only recurse on standard JSON-like structures that AI commonly outputs
        if type(val) in (list, tuple):
            return type(val)(self._unwrap(x) for x in val)
        if type(val) is dict:
            return {k: self._unwrap(v) for k, v in val.items()}
        # Treat other subclasses (like RGBColor) as leaf values
        return val

    def __getattr__(self, name):
        if self._target is None:
            return NullProxy()
        
        # ── 1. Specific Smart Path Mappings ──
        # Fix: text_frame.alignment -> paragraphs[0].alignment
        if name == "alignment" and "TextFrame" in str(type(self._target)):
            try: return SafeProxy(self._target.paragraphs[0].alignment)
            except: pass
        if name == "font" and "TextFrame" in str(type(self._target)):
            try: return SafeProxy(self._target.paragraphs[0].font)
            except: pass
        
        # ── 2. Standard Native Access with Auto-Healing ──
        try:
            # Special case for 'background' to 'fill.background' (common in AI)
            if name == "background" and hasattr(self._target, "fill"):
                return SafeProxy(self._target.fill.background)
            
            # ATTEMPT ACCESS
            try:
                attr = getattr(self._target, name)
            except (AttributeError, TypeError, ValueError):
                # AUTO-HEALING: fill.fore_color fails if .solid() wasn't called
                if name == "fore_color" and "Fill" in str(type(self._target)):
                    try:
                        self._target.solid()
                        attr = getattr(self._target, name)
                    except: return NullProxy()
                elif name == "rgb" and "Color" in str(type(self._target)):
                    # Fix: If accessing .rgb directly on something that has fore_color
                    try: return SafeProxy(self._target.fore_color.rgb)
                    except: return NullProxy()
                else: return NullProxy()

            if hasattr(attr, "__dict__") or "pptx." in str(type(attr)):
                return SafeProxy(attr)
            return attr
        except (AttributeError, TypeError, ValueError, KeyError):
            # ── 3. Silent Fallback (Null Object Pattern) ──
            # Returns a NullProxy instead of crashing, allowing script to finish
            print(f"WARNING: Hallucination or sequence error! Skipping '{name}' on {type(self._target)}")
            return NullProxy()

    def __setattr__(self, name, value):
        if self._target is None:
            return
        # Sanitize text assignments to prevent XML corruption
        if name == "text" and isinstance(value, str):
            value = sanitize_text(value)
        # Unwrap if the value is another proxy (e.g. chart_data = other_proxy)
        try:
            setattr(self._target, name, self._unwrap(value))
        except (AttributeError, TypeError):
            # If the attribute doesn't exist for assignment, just skip it (Resilience)
            pass

    def __setitem__(self, key, value):
        if self._target is None: return
        try:
            self._target[key] = self._unwrap(value)
        except: pass

    def __getitem__(self, key):
        if self._target is None: return NullProxy()
        try:
            res = self._target[key]
            if hasattr(res, "__dict__") or "pptx." in str(type(res)):
                return SafeProxy(res)
            return res
        except: return NullProxy()

    def __call__(self, *args, **kwargs):
        if self._target is None: return NullProxy()
        try:
            # Unwrap all incoming arguments so the native pptx library can handle them
            u_args = [self._unwrap(a) for a in args]
            u_kwargs = {k: self._unwrap(v) for k, v in kwargs.items()}
            res = self._target(*u_args, **u_kwargs)
            if hasattr(res, "__dict__") or "pptx." in str(type(res)):
                return SafeProxy(res)
            return res
        except: return NullProxy()

    def __bool__(self):
        if self._target is None: return False
        return bool(self._target)

    def __contains__(self, item):
        if self._target is None: return False
        try: return self._unwrap(item) in self._target
        except: return False

    def __len__(self):
        return len(self._target)

    def __iter__(self):
        if self._target is None: return iter([])
        try:
            # Try to start iteration
            it = iter(self._target)
            for item in it:
                if hasattr(item, "__dict__") or "pptx." in str(type(item)):
                    yield SafeProxy(item)
                else:
                    yield item
        except TypeError:
            # ── Smart Fallback: Non-iterable object (like DataLabels) ──
            # If the AI thinks it's a sequence but it's not, just skip the loop
            print(f"WARNING: Tried to iterate over non-iterable {type(self._target)}. Skipping loop.")
            return # Yields nothing (empty generator)

class ShapesProxy(SafeProxy):
    def _to_emu(self, val):
        if isinstance(val, float) and 0 < val < 50:
            from pptx.util import Inches
            return Inches(val)
        return val

    def add_shape(self, shape_type, left, top, width, height):
        from pptx.enum.shapes import MSO_SHAPE
        left, top, width, height = map(self._to_emu, [left, top, width, height])
        if shape_type == MSO_SHAPE.ROUNDED_RECTANGULAR_CALLOUT:
            shape_type = MSO_SHAPE.ROUNDED_RECTANGLE
        try:
            return SafeProxy(self._target.add_shape(shape_type, left, top, width, height))
        except Exception:
            return SafeProxy(self._target.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height))

    def add_textbox(self, left, top, width, height):
        left, top, width, height = map(self._to_emu, [left, top, width, height])
        return SafeProxy(self._target.add_textbox(left, top, width, height))

    def add_table(self, rows, cols, left, top, width, height):
        left, top, width, height = map(self._to_emu, [left, top, width, height])
        return SafeProxy(self._target.add_table(rows, cols, left, top, width, height))

    def add_chart(self, chart_type, left, top, width, height, chart_data):
        left, top, width, height = map(self._to_emu, [left, top, width, height])
        try:
            return SafeProxy(self._target.add_chart(chart_type, left, top, width, height, chart_data))
        except Exception:
            # Fallback: Instead of crashing, add a placeholder rectangle
            print("ERROR: add_chart failed, using placeholder rectangle")
            from pptx.dml.color import RGBColor
            dummy = self._target.add_shape(1, left, top, width, height)
            dummy.fill.solid()
            dummy.fill.fore_color.rgb = RGBColor(200, 200, 200)
            return SafeProxy(dummy)

    def add_connector(self, connector_type, x1, y1, x2, y2):
        x1, y1, x2, y2 = map(self._to_emu, [x1, y1, x2, y2])
        from pptx.enum.shapes import MSO_SHAPE
        left = min(x1, x2)
        top = min(y1, y2)
        width = max(Pt(1), abs(x2 - x1))
        height = max(Pt(1), abs(y2 - y1))
        # Use RECTANGLE as a safe fallback for lines to avoid connector relationship issues
        return SafeProxy(self._target.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height))

class SlideProxy(SafeProxy):
    def __getattr__(self, name):
        if name == "shapes":
            return ShapesProxy(self._target.shapes)
        return super().__getattr__(name)

    # ... other methods if needed ...

class PptxProxy:
    def __init__(self, real_prs):
        object.__setattr__(self, "_real", real_prs)
        object.__setattr__(self, "slides", SlidesListProxy(real_prs.slides))
    def __getattr__(self, name): return getattr(self._real, name)
    def __setattr__(self, name, value): setattr(self._real, name, value)
    def save(self, *args, **kwargs): return self._real.save(*args, **kwargs)

class SlidesListProxy:
    def __init__(self, real_slides): self._real = real_slides
    def add_slide(self, layout): return SlideProxy(self._real.add_slide(layout))
    def __len__(self): return len(self._real)
    def __getitem__(self, idx): return SlideProxy(self._real[idx])
    def __getattr__(self, name): return getattr(self._real, name)


def _patch_common_errors(code: str) -> str:
    """Auto-fix common mistakes and syntax errors the AI makes in generated code."""
    # ── 1. Local Syntax Fixes: 'Peeled Patching' for string literals ──
    new_lines = []
    for line in code.split('\n'):
        if not line.strip() or '"""' in line or "'''" in line:
            new_lines.append(line)
            continue
        
        # Heuristic: Remove all valid, paired strings first to find the culprit
        # We use non-greedy matching to leave lone quotes alone
        stripped = re.sub(r'"[^"]*"', '', line)
        stripped = re.sub(r"'[^']*'", '', stripped)
        
        # If we have a lone quote leftover, append it at the end
        if '"' in stripped and not line.rstrip().endswith('\\'):
            line = line + '"'
        elif "'" in stripped and not line.rstrip().endswith('\\'):
            line = line + "'"
            
        new_lines.append(line)
        
    code = '\n'.join(new_lines)

    # ── 2. Global Syntax Fix: Balanced Parentheses (Script-wide) ──
    # If the script was truncated, close ALL remaining parentheses at the very end
    open_p = code.count('(')
    close_p = code.count(')')
    if open_p > close_p:
        code += ')' * (open_p - close_p)

    # ── 3. Structural Fixes ──
    # .line.background() -> .line.fill.background()
    code = re.sub(r'\.line\.background\(\)', '.line.fill.background()', code)
    # .line.no_fill() -> .line.fill.background()
    code = re.sub(r'\.line\.no_fill\(\)', '.line.fill.background()', code)
    
    # MSO_CONNECTOR not imported -> add import if used
    if 'MSO_CONNECTOR' in code and 'from pptx.enum.shapes import MSO_CONNECTOR' not in code:
        code = 'from pptx.enum.shapes import MSO_CONNECTOR\n' + code
    
    # Replace invalid MSO_SHAPE members with ROUNDED_RECTANGLE
    from pptx.enum.shapes import MSO_SHAPE
    _valid_shapes = set(MSO_SHAPE.__members__.keys())
    def _fix_shape(m):
        name = m.group(1)
        if name in _valid_shapes:
            return m.group(0)
        return f'MSO_SHAPE.ROUNDED_RECTANGLE'
    code = re.sub(r'MSO_SHAPE\.([A-Z_0-9]+)', _fix_shape, code)
    return code


def _exec_script(script: str) -> tuple[bool, str]:
    """Execute a generated pptx script in-process using proxy and fallback. Returns (success, error)."""
    import importlib
    script = _patch_common_errors(script)
    
    attempts = 0
    max_attempts = 3
    current_script = script
    
    while attempts < max_attempts:
        try:
            # 1. Strip standard Presentation import to allow Proxy injection
            script_to_run = re.sub(r'from pptx import Presentation', '# removed to allow proxy', current_script)
            
            # Prepare execution environment
            from pptx import Presentation as RealPresentation
            from pptx.chart.data import CategoryChartData as RealCCD

            def ProxyPresentation(*args, **kwargs):
                return PptxProxy(RealPresentation(*args, **kwargs))

            def ProxyCategoryChartData(*args, **kwargs):
                return SafeProxy(RealCCD(*args, **kwargs))

            exec_globals = {
                "__builtins__": __builtins__, 
                "Presentation": ProxyPresentation,
                "CategoryChartData": ProxyCategoryChartData,
                "re": re, 
                "traceback": traceback
            }
            exec(compile(script_to_run, "<pptx_gen>", "exec"), exec_globals)
            return True, ""
            
        except Exception as e:
            attempts += 1
            err_msg = traceback.format_exc()
            
            # ── Dynamic Fallback Patches ──
            if "AttributeError" in err_msg:
                if "'LineFormat' object has no attribute 'background'" in err_msg:
                    current_script = current_script.replace(".line.background()", ".line.fill.background()")
                    continue
            
            if "NameError" in err_msg:
                match = re.search(r"name '(\w+)' is not defined", str(e))
                if match:
                    missing_name = match.group(1)
                    found = False
                    for modname in ["pptx.enum.shapes", "pptx.enum.text", "pptx.enum.chart", "pptx.util"]:
                        try:
                            mod = importlib.import_module(modname)
                            if hasattr(mod, missing_name):
                                current_script = f"from {modname} import {missing_name}\n" + current_script
                                found = True
                                break
                        except ImportError: continue
                    if found: continue
            
            if attempts >= max_attempts:
                return False, traceback.format_exc()
    
    return False, "Failed after max retries"


# ── Slide code persistence ──

def get_slides_dir(proj_dir: str) -> str:
    d = os.path.join(proj_dir, "最终文档", "ppt_slides")
    os.makedirs(d, exist_ok=True)
    return d


def save_slide_code(proj_dir: str, page_num: int, code: str):
    path = os.path.join(get_slides_dir(proj_dir), f"slide_{page_num:02d}.py")
    with open(path, "w", encoding="utf-8") as f:
        f.write(code)


def load_slide_code(proj_dir: str, page_num: int) -> str | None:
    path = os.path.join(get_slides_dir(proj_dir), f"slide_{page_num:02d}.py")
    if os.path.exists(path):
        with open(path, "r", encoding="utf-8") as f:
            return f.read()
    return None


def load_all_slide_codes(proj_dir: str) -> dict[int, str]:
    """Return {page_num: code} for all saved slide codes."""
    slides_dir = get_slides_dir(proj_dir)
    result = {}
    for fname in sorted(os.listdir(slides_dir)):
        m = re.match(r"slide_(\d+)\.py$", fname)
        if m:
            page_num = int(m.group(1))
            with open(os.path.join(slides_dir, fname), "r", encoding="utf-8") as f:
                result[page_num] = f.read()
    return result


def get_single_pptx_path(proj_dir: str, page_num: int) -> str:
    return os.path.join(get_slides_dir(proj_dir), f"slide_{page_num:02d}.pptx")
