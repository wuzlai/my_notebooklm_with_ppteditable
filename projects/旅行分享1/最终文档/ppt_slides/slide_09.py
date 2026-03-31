def build_slide(slide):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

    # Colors
    BG_COLOR = RGBColor(245, 242, 235)
    ORANGE_DECOR = RGBColor(235, 104, 65)
    PURPLE_DECOR = RGBColor(103, 58, 183)
    YELLOW_HL = RGBColor(255, 235, 59)
    BLACK = RGBColor(0, 0, 0)
    WHITE = RGBColor(255, 255, 255)
    BOX_ORANGE = RGBColor(255, 112, 67)
    BOX_GREEN = RGBColor(129, 199, 132)
    BOX_YELLOW = RGBColor(255, 235, 59)

    # 1. Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_COLOR
    bg.line.fill.background()

    # Corner decorations (simulating torn paper)
    tr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(11), Inches(0), Inches(2.333), Inches(1.5))
    tr.rotation = 180
    tr.fill.solid()
    tr.fill.fore_color.rgb = ORANGE_DECOR
    tr.line.fill.background()

    bl = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(0), Inches(5.5), Inches(2.5), Inches(2))
    bl.fill.solid()
    bl.fill.fore_color.rgb = ORANGE_DECOR
    bl.line.fill.background()

    br = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(12), Inches(6.5), Inches(1.333), Inches(1))
    br.rotation = 270
    br.fill.solid()
    br.fill.fore_color.rgb = PURPLE_DECOR
    br.line.fill.background()

    # 2. Title & Subtitle
    # Title Highlight
    hl = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.9), Inches(8.0), Inches(0.6))
    hl.fill.solid()
    hl.fill.fore_color.rgb = YELLOW_HL
    hl.line.fill.background()

    # Title Text
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.7), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "桂林生存法则：如何优雅地避坑"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = BLACK

    # Subtitle Text
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(5), Inches(0.5))
    tf = sub_box.text_frame
    p = tf.add_paragraph()
    p.text = "猎奇博主的血泪教训"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = BLACK

    # 3. Left Visuals
    # Landscape Photo Base (White border)
    photo_base = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.0), Inches(5.0), Inches(3.5))
    photo_base.fill.solid()
    photo_base.fill.fore_color.rgb = WHITE
    photo_base.line.fill.background()
    photo_base.rotation = -2

    # Landscape Inner
    land = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(3.2), Inches(4.6), Inches(3.1))
    land.fill.solid()
    land.fill.fore_color.rgb = RGBColor(46, 139, 87)
    land.line.fill.background()
    land.rotation = -2
    tf = land.text_frame
    tf.text = "\n\n桂林山水"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE

    # 20 RMB Note
    note = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.5), Inches(2.0), Inches(2.8), Inches(1.4))
    note.fill.solid()
    note.fill.fore_color.rgb = RGBColor(245, 235, 235)
    note.line.color.rgb = RGBColor(200, 200, 200)
    note.rotation = -10
    tf = note.text_frame
    tf.text = "20\n中国人民银行"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(150, 50, 50)

    # Blogger (Simulated)
    blogger_body = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.5), Inches(1.8), Inches(2.5))
    blogger_body.fill.solid()
    blogger_body.fill.fore_color.rgb = RGBColor(255, 193, 7)
    blogger_body.line.color.rgb = BLACK
    blogger_body.line.width = Pt(2)
    tf = blogger_body.text_frame
    tf.text = "\n\n博主"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = BLACK

    # Boom!
    boom = slide.shapes.add_shape(MSO_SHAPE.EXPLOSION1, Inches(0.5), Inches(2.2), Inches(1.8), Inches(1.2))
    boom.fill.solid()
    boom.fill.fore_color.rgb = RGBColor(255, 87, 34)
    boom.line.color.rgb = BLACK
    boom.line.width = Pt(2)
    boom.rotation = -15
    tf = boom.text_frame
    p = tf.paragraphs[0]
    p.text = "Boom!"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = BLACK

    # OMG!
    omg = slide.shapes.add_shape(MSO_SHAPE.EXPLOSION2, Inches(4.2), Inches(5.5), Inches(1.5), Inches(1.2))
    omg.fill.solid()
    omg.fill.fore_color.rgb = RGBColor(255, 87, 34)
    omg.line.color.rgb = BLACK
    omg.line.width = Pt(2)
    omg.rotation = 10
    tf = omg.text_frame
    p = tf.paragraphs[0]
    p.text = "OMG!"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = BLACK

    # 4. Right Pitfall Boxes
    def add_pitfall_box(slide, left, top, img_bg_color, text_bg_color, text_color, img_text, desc_text):
        box_w = Inches(2.8)
        img_h = Inches(1.5)
        txt_h = Inches(0.9)

        # Image part
        img_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, box_w, img_h)
        img_box.fill.solid()
        img_box.fill.fore_color.rgb = img_bg_color
        img_box.line.color.rgb = BLACK
        img_box.line.width = Pt(2)
        tf = img_box.text_frame
        tf.text = img_text
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = BLACK

        # Text part
        txt_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + img_h, box_w, txt_h)
        txt_box.fill.solid()
        txt_box.fill.fore_color.rgb = text_bg_color
        txt_box.line.color.rgb = BLACK
        txt_box.line.width = Pt(2)
        tf = txt_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc_text
        p.alignment = PP_ALIGN.LEFT
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = text_color

        # Warning Triangle
        warn = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, left - Inches(0.25), top - Inches(0.25), Inches(0.6), Inches(0.6))
        warn.fill.solid()
        warn.fill.fore_color.rgb = YELLOW_HL
        warn.line.color.rgb = BLACK
        warn.line.width = Pt(2)
        tf = warn.text_frame
        p = tf.paragraphs[0]
        p.text = "!"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = BLACK

    # Box 1: Auntie
    add_pitfall_box(slide, Inches(6.5), Inches(1.5),
                    BOX_ORANGE, BOX_YELLOW, BLACK,
                    "帅哥/美女，来吃鱼！\n\nFAKE NEWS!  ❌ 坑！",
                    "别相信路边喊你“帅哥/美女”去吃鱼的阿姨")

    # Box 2: Cormorant 1
    add_pitfall_box(slide, Inches(10.0), Inches(1.2),
                    BOX_GREEN, BLACK, WHITE,
                    "¥999+\n\n(鸬鹚拍照)",
                    "拍照前先问价格，否则鸬鹚的身价比你还高")

    # Box 3: Cormorant 2
    add_pitfall_box(slide, Inches(6.5), Inches(4.5),
                    BOX_GREEN, BLACK, WHITE,
                    "¥999+\n\n(鸬鹚拍照)",
                    "拍照前先问价格，否则鸬鹚的身价比你还高")

    # Box 4: Luosifen
    add_pitfall_box(slide, Inches(10.0), Inches(4.5),
                    BOX_ORANGE, BLACK, WHITE,
                    "螺蛳粉 + 肠胃药\n\n🔥 危险！",
                    "备好肠胃药，螺蛳粉的后劲比酒还大")

    # 5. Arrows (Using connectors)
    conn1 = slide.shapes.add_connector(MSO_CONNECTOR.CURVE, Inches(9.3), Inches(2.2), Inches(10.0), Inches(2.2))
    conn1.line.color.rgb = BLACK
    conn1.line.width = Pt(2)

    conn2 = slide.shapes.add_connector(MSO_CONNECTOR.CURVE, Inches(6.3), Inches(3.0), Inches(6.3), Inches(4.5))
    conn2.line.color.rgb = BLACK
    conn2.line.width = Pt(2)

    conn3 = slide.shapes.add_connector(MSO_CONNECTOR.CURVE, Inches(9.3), Inches(5.2), Inches(10.0), Inches(5.2))
    conn3.line.color.rgb = BLACK
    conn3.line.width = Pt(2)

    conn4 = slide.shapes.add_connector(MSO_CONNECTOR.CURVE, Inches(11.4), Inches(3.6), Inches(11.4), Inches(4.5))
    conn4.line.color.rgb = BLACK
    conn4.line.width = Pt(2)

    # 6. Page Number & Icon
    page_oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.5), Inches(6.8), Inches(0.6), Inches(0.4))
    page_oval.fill.solid()
    page_oval.fill.fore_color.rgb = WHITE
    page_oval.line.color.rgb = BLACK
    page_oval.line.width = Pt(1.5)
    tf = page_oval.text_frame
    p = tf.paragraphs[0]
    p.text = "09"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = BLACK

    face = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(12.2), Inches(6.6), Inches(0.8), Inches(0.8))
    face.fill.solid()
    face.fill.fore_color.rgb = RGBColor(255, 224, 178)
    face.line.color.rgb = BLACK
    face.line.width = Pt(1.5)
    tf = face.text_frame
    p = tf.paragraphs[0]
    p.text = "博主"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = BLACK