def build_slide(slide):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
    
    # 定义颜色
    C_BG = RGBColor(253, 249, 238)
    C_RED = RGBColor(228, 57, 50)
    C_YELLOW = RGBColor(255, 214, 89)
    C_GREEN = RGBColor(0, 150, 100)
    C_ORANGE = RGBColor(244, 121, 32)
    C_BLACK = RGBColor(0, 0, 0)
    C_WHITE = RGBColor(255, 255, 255)

    # 设置背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_BG
    bg.line.fill.background()
    
    # 添加背景装饰色块
    blob1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(-2), Inches(5), Inches(5))
    blob1.fill.solid(); blob1.fill.fore_color.rgb = C_RED; blob1.line.fill.background()
    
    blob2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3), Inches(-2), Inches(6), Inches(4))
    blob2.fill.solid(); blob2.fill.fore_color.rgb = C_YELLOW; blob2.line.fill.background()
    
    blob3 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(-2), Inches(6), Inches(5))
    blob3.fill.solid(); blob3.fill.fore_color.rgb = C_GREEN; blob3.line.fill.background()
    
    blob4 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1), Inches(6), Inches(5), Inches(4))
    blob4.fill.solid(); blob4.fill.fore_color.rgb = C_GREEN; blob4.line.fill.background()
    
    blob5 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8), Inches(6), Inches(7), Inches(4))
    blob5.fill.solid(); blob5.fill.fore_color.rgb = C_ORANGE; blob5.line.fill.background()

    # 标题 (使用多重偏移模拟描边效果)
    title_text = "螺蛳粉：桂林的“生化武器”诱惑"
    offsets = [(-0.03, -0.03), (0.03, -0.03), (-0.03, 0.03), (0.03, 0.03), (-0.04, 0), (0.04, 0), (0, -0.04), (0, 0.04)]
    for ox, oy in offsets:
        tb = slide.shapes.add_textbox(Inches(0.5 + ox), Inches(0.3 + oy), Inches(10), Inches(1))
        p = tb.text_frame.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = C_BLACK
        p.font.name = "Microsoft YaHei"
        
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(10), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p.font.name = "Microsoft YaHei"

    # 副标题
    sub_text = "闻着臭，吃着爽，回味想撞墙"
    for ox, oy in offsets:
        tb = slide.shapes.add_textbox(Inches(0.5 + ox), Inches(1.1 + oy), Inches(8), Inches(0.8))
        p = tb.text_frame.paragraphs[0]
        p.text = sub_text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = C_BLACK
        p.font.name = "Microsoft YaHei"
        
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(8), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.text = sub_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p.font.name = "Microsoft YaHei"

    # 页码
    page_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.8), Inches(0.4), Inches(1.0), Inches(1.0))
    page_circle.fill.solid(); page_circle.fill.fore_color.rgb = C_BG
    page_circle.line.color.rgb = C_BLACK; page_circle.line.width = Pt(2)
    
    page_inner = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.85), Inches(0.45), Inches(0.9), Inches(0.9))
    page_inner.fill.background()
    page_inner.line.color.rgb = C_BLACK; page_inner.line.width = Pt(1)
    
    tb = slide.shapes.add_textbox(Inches(11.8), Inches(0.5), Inches(1.0), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.text = "PAGE\n5"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = C_BLACK

    # --- 左侧：内容要点 ---
    shadow_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(2.6), Inches(2.0), Inches(0.6))
    shadow_box.fill.solid(); shadow_box.fill.fore_color.rgb = C_GREEN
    shadow_box.line.color.rgb = C_BLACK; shadow_box.line.width = Pt(1.5)
    
    main_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.5), Inches(2.0), Inches(0.6))
    main_box.fill.solid(); main_box.fill.fore_color.rgb = C_YELLOW
    main_box.line.color.rgb = C_BLACK; main_box.line.width = Pt(1.5)
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.55), Inches(2.0), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = "内容要点"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = C_BLACK

    # 要点 1
    bullet1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(3.6), Inches(0.15), Inches(0.15))
    bullet1.fill.solid(); bullet1.fill.fore_color.rgb = C_RED; bullet1.line.color.rgb = C_BLACK
    
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(3.4), Inches(3.0), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = "气味等级："
    p.font.size = Pt(20)
    p.font.bold = True
    
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(3.9), Inches(3.0), Inches(1.0))
    p = tb.text_frame.paragraphs[0]
    p.text = "路过店门口，衣服\n自动“腌制”三天"
    p.font.size = Pt(18)
    
    # 腌制 圈注
    oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.5), Inches(4.2), Inches(0.9), Inches(0.4))
    oval.fill.background()
    oval.line.color.rgb = C_ORANGE; oval.line.width = Pt(2)

    # 要点 2
    bullet2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(5.4), Inches(0.15), Inches(0.15))
    bullet2.fill.solid(); bullet2.fill.fore_color.rgb = C_RED; bullet2.line.color.rgb = C_BLACK
    
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(5.2), Inches(3.0), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = "猎奇吃法："
    p.font.size = Pt(20)
    p.font.bold = True
    
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(5.7), Inches(3.0), Inches(1.0))
    p = tb.text_frame.paragraphs[0]
    run1 = p.add_run()
    run1.text = "加辣加酸笋"
    run1.font.size = Pt(18)
    run1.font.underline = True
    run2 = p.add_run()
    run2.text = "，\n挑战味蕾极限"
    run2.font.size = Pt(18)

    # --- 中间：爆炸图与杀伤力 ---
    exp_red = slide.shapes.add_shape(MSO_SHAPE.EXPLOSION2, Inches(3.8), Inches(2.2), Inches(5.5), Inches(5.0))
    exp_red.fill.solid(); exp_red.fill.fore_color.rgb = C_RED
    exp_red.line.color.rgb = C_BLACK; exp_red.line.width = Pt(2)
    
    exp_yellow = slide.shapes.add_shape(MSO_SHAPE.EXPLOSION1, Inches(4.3), Inches(2.7), Inches(4.5), Inches(4.0))
    exp_yellow.fill.solid(); exp_yellow.fill.fore_color.rgb = C_YELLOW
    exp_yellow.line.color.rgb = C_BLACK; exp_yellow.line.width = Pt(2)
    
    # 中心碗
    bowl_outer = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.4), Inches(3.8), Inches(2.2), Inches(1.6))
    bowl_outer.fill.solid(); bowl_outer.fill.fore_color.rgb = C_WHITE
    bowl_outer.line.color.rgb = C_BLACK; bowl_outer.line.width = Pt(2)
    
    bowl_inner = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.6), Inches(3.9), Inches(1.8), Inches(0.8))
    bowl_inner.fill.solid(); bowl_inner.fill.fore_color.rgb = C_ORANGE
    bowl_inner.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(5.4), Inches(4.7), Inches(2.2), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = "杀伤力"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)
    p.font.bold = True

    # 放射箭头
    def add_custom_arrow(slide, x, y, width, height, rotation):
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(width), Inches(height))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = C_YELLOW
        arrow.line.color.rgb = C_RED; arrow.line.width = Pt(1.5)
        arrow.rotation = rotation

    add_custom_arrow(slide, 5.4, 3.4, 0.8, 0.4, 225)
    add_custom_arrow(slide, 6.1, 3.0, 0.8, 0.4, 270)
    add_custom_arrow(slide, 7.1, 3.4, 0.8, 0.4, 315)
    add_custom_arrow(slide, 7.6, 4.5, 0.8, 0.4, 0)
    add_custom_arrow(slide, 7.1, 5.5, 0.8, 0.4, 45)
    add_custom_arrow(slide, 6.1, 5.9, 0.8, 0.4, 90)
    add_custom_arrow(slide, 5.4, 5.5, 0.8, 0.4, 135)
    add_custom_arrow(slide, 4.8, 4.5, 0.8, 0.4, 180)

    # 标签文字
    labels = [
        ("辣油", 4.2, 2.8), ("臭！", 6.0, 2.0), ("酸笋", 7.8, 2.6),
        ("辣！", 8.6, 4.2), ("腐竹", 8.0, 5.6), ("螺蛳汤", 6.0, 6.5),
        ("臭味", 4.2, 5.8), ("爽！", 3.8, 4.2)
    ]
    for text, lx, ly in labels:
        tb = slide.shapes.add_textbox(Inches(lx), Inches(ly), Inches(1.2), Inches(0.6))
        p = tb.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = C_BLACK

    # BOOM 装饰字
    tb = slide.shapes.add_textbox(Inches(7.8), Inches(1.8), Inches(1.5), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = "BOOM!"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = C_RED
    tb.rotation = -15

    # --- 右侧：表情包九宫格 ---
    bullet3 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(2.2), Inches(0.15), Inches(0.15))
    bullet3.fill.solid(); bullet3.fill.fore_color.rgb = C_RED; bullet3.line.color.rgb = C_BLACK
    
    tb = slide.shapes.add_textbox(Inches(9.7), Inches(2.0), Inches(3.0), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = "搞笑记录："
    p.font.size = Pt(22)
    p.font.bold = True
    
    tb = slide.shapes.add_textbox(Inches(9.7), Inches(2.5), Inches(3.0), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = "表情包九宫格"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = C_ORANGE

    # 第一次吃 小插图
    tb = slide.shapes.add_textbox(Inches(11.5), Inches(1.8), Inches(1.5), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    p.text = "第一次吃！"
    p.font.size = Pt(12)
    p.font.bold = True
    
    bowl_small = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.8), Inches(2.2), Inches(0.8), Inches(0.4))
    bowl_small.fill.solid(); bowl_small.fill.fore_color.rgb = C_WHITE
    bowl_small.line.color.rgb = C_BLACK

    # 九宫格
    captions = ["OMG!", "上头了", "救命！", "辣哭", "真香！", "——", "爱了爱了", "？？", "想撞墙"]
    start_x, start_y = 9.5, 3.2
    cell_w, cell_h = 1.1, 1.3
    spacing = 0.15
    
    photo_colors = [
        RGBColor(255, 200, 200), RGBColor(200, 255, 200), RGBColor(200, 200, 255),
        RGBColor(255, 255, 200), RGBColor(255, 200, 255), RGBColor(200, 255, 255),
        RGBColor(240, 240, 240), RGBColor(255, 220, 180), RGBColor(180, 220, 255)
    ]
    
    for i in range(9):
        row, col = i // 3, i % 3
        x = start_x + col * (cell_w + spacing)
        y = start_y + row * (cell_h + spacing)
        
        # 外框
        frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(cell_w), Inches(cell_h))
        frame.fill.solid(); frame.fill.fore_color.rgb = C_WHITE
        frame.line.color.rgb = C_BLACK; frame.line.width = Pt(1)
        
        # 照片区
        photo = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 0.05), Inches(y + 0.05), Inches(cell_w - 0.1), Inches(cell_h - 0.4))
        photo.fill.solid(); photo.fill.fore_color.rgb = photo_colors[i]
        photo.line.color.rgb = C_BLACK; photo.line.width = Pt(1)
        
        # 底部文字
        tb = slide.shapes.add_textbox(Inches(x), Inches(y + cell_h - 0.35), Inches(cell_w), Inches(0.3))
        p = tb.text_frame.paragraphs[0]
        p.text = captions[i]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(12)
        p.font.bold = True

    # 胶带装饰
    def add_tape(slide, x, y, rotation):
        tape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.8), Inches(0.25))
        tape.fill.solid(); tape.fill.fore_color.rgb = RGBColor(240, 230, 210)
        tape.line.fill.background()
        tape.rotation = rotation
        
    add_tape(slide, 0.2, 2.3, -45)
    add_tape(slide, 2.5, 2.3, 30)
    add_tape(slide, 9.2, 3.0, -20)
    add_tape(slide, 12.8, 3.0, 45)

    return slide