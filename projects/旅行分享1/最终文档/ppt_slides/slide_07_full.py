from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# ── Color Palette ──
BLUE_HEADER = RGBColor(0x5B, 0x9B, 0xD5)
BLUE_DARK   = RGBColor(0x4A, 0x86, 0xC8)
CYAN        = RGBColor(0x00, 0xBC, 0xD4)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x33, 0x33, 0x33)
GRAY_TEXT   = RGBColor(0x55, 0x55, 0x55)
GRAY_BAR    = RGBColor(0xB0, 0xBE, 0xC5)
RED         = RGBColor(0xE5, 0x39, 0x35)
GREEN       = RGBColor(0x43, 0xA0, 0x47)
ORANGE      = RGBColor(0xFF, 0x98, 0x00)

ICON_BG  = RGBColor(0xE3, 0xE8, 0xED)
ICON_FG  = RGBColor(0x54, 0x6E, 0x7A)
FONT_NAME = "Microsoft YaHei"
SLIDE_WIDTH = Inches(13.333)
HEADER_H    = Inches(0.75)
SUBTITLE_Y  = Inches(0.95)


def add_header_banner(slide, title_text, bg_color=None):
    if bg_color is None:
        bg_color = BLUE_HEADER
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, HEADER_H
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = bg_color
    banner.line.fill.background()
    tf = banner.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(26)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME


def add_subtitle(slide, text, left, top, width=Inches(12), font_size=Pt(18)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.font.name = FONT_NAME
    return txBox


def add_icon_box(slide, left, top, symbol, size=Inches(0.48)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ICON_BG
    shape.line.fill.background()
    shape.adjustments[0] = 0.25
    tf = shape.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = symbol
    p.font.size = Pt(18)
    p.font.color.rgb = ICON_FG
    p.font.bold = False
    return shape


def add_bullet_item(slide, left, top, symbol, label, description,
                    width=Inches(5.5), desc_size=Pt(13)):
    add_icon_box(slide, left, top, symbol)
    text_left = left + Inches(0.65)
    txBox = slide.shapes.add_textbox(text_left, top - Inches(0.02), width, Inches(0.65))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run_label = p.add_run()
    run_label.text = label + "\uff1a"
    run_label.font.size = Pt(14)
    run_label.font.color.rgb = BLACK
    run_label.font.bold = True
    run_label.font.name = FONT_NAME
    run_desc = p.add_run()
    run_desc.text = description
    run_desc.font.size = desc_size
    run_desc.font.color.rgb = GRAY_TEXT
    run_desc.font.bold = False
    run_desc.font.name = FONT_NAME
    return txBox


def add_conclusion_box(slide, left, top, width, text, font_size=Pt(13)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = BLACK
    run.font.bold = True
    run.font.name = FONT_NAME
    return txBox


def add_table(slide, left, top, width, height, rows, cols, data,
              header_color=None, col_widths=None):
    if header_color is None:
        header_color = BLUE_HEADER
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if data[r][c] is not None else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(11)
                    run.font.name = FONT_NAME
                    if r == 0:
                        run.font.color.rgb = WHITE
                        run.font.bold = True
                    else:
                        run.font.color.rgb = BLACK
                        run.font.bold = False
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else RGBColor(0xF5, 0xF5, 0xF5)
            cell.margin_left = Pt(5)
            cell.margin_right = Pt(5)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
    return table_shape


def add_bar_chart(slide, left, top, width, height,
                  categories, values, title="", bar_colors=None):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', values)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    chart.chart_style = 2
    plot = chart.plots[0]
    plot.gap_width = 100
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CYAN
    series.has_data_labels = True
    dl = series.data_labels
    dl.font.size = Pt(13)
    dl.font.bold = True
    dl.font.color.rgb = BLACK
    dl.number_format = '0.#'
    dl.show_value = True
    dl.label_position = XL_LABEL_POSITION.OUTSIDE_END
    if bar_colors:
        for i, color in enumerate(bar_colors):
            pt = series.points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = color
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(12)
    cat_axis.tick_labels.font.name = FONT_NAME
    cat_axis.major_tick_mark = 2
    cat_axis.format.line.fill.background()
    val_axis = chart.value_axis
    val_axis.visible = False
    val_axis.major_tick_mark = 2
    val_axis.format.line.fill.background()
    val_axis.major_gridlines.format.line.fill.background()
    if title:
        chart.has_title = True
        ct = chart.chart_title.text_frame.paragraphs[0]
        ct.text = title
        ct.font.size = Pt(14)
        ct.font.bold = True
        ct.font.name = FONT_NAME
    else:
        chart.has_title = False
    return chart_frame


def add_callout_label(slide, left, top, text, bg_color=None, font_size=Pt(11)):
    if bg_color is None:
        bg_color = CYAN
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.3), Inches(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME
    return shape


def add_data_card(slide, left, top, width, height, value, label,
                  value_color=None, bg_color=None):
    if value_color is None:
        value_color = CYAN
    if bg_color is None:
        bg_color = WHITE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(3)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run1 = p1.add_run()
    run1.text = str(value)
    run1.font.size = Pt(24)
    run1.font.color.rgb = value_color
    run1.font.bold = True
    run1.font.name = FONT_NAME
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(10)
    run2.font.color.rgb = GRAY_TEXT
    run2.font.bold = False
    run2.font.name = FONT_NAME
    return shape


OUTPUT_PATH = r"C:\Users\Administrator\Desktop\notebook\my_notebooklm_with_ppteditable\projects\旅行分享1\最终文档\ppt_slides\slide_07.pptx"


# ── Slide 1 ──

def build_slide_1(slide):
    # --- 颜色定义 ---
    ORANGE = RGBColor(0xF2, 0x71, 0x27)
    DARK_TEXT = RGBColor(0x20, 0x20, 0x20)
    GRAY_TEXT = RGBColor(0x60, 0x60, 0x60)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    BLUE_LIGHT = RGBColor(0x00, 0xBC, 0xD4)
    RED = RGBColor(0xE5, 0x39, 0x35)
    IMG_BG = RGBColor(0xE8, 0xF0, 0xF4)

    # --- 1. 左上角页码 ---
    page_num = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.2), Inches(0.2), Inches(1.0), Inches(1.0))
    page_num.fill.solid()
    page_num.fill.fore_color.rgb = ORANGE
    page_num.line.fill.background()
    tf = page_num.text_frame
    p = tf.paragraphs[0]
    p.text = "第7页"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DARK_TEXT
    p.alignment = PP_ALIGN.CENTER

    # --- 2. 主标题与副标题 ---
    title_box = slide.shapes.add_textbox(Inches(1.5), Inches(0.3), Inches(10), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "漓江竹筏：塑料椅子上的“速度与激情”"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.name = "Microsoft YaHei"
    p.font.color.rgb = DARK_TEXT

    # 标题下划线
    underline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.6), Inches(0.9), Inches(8.5), Inches(0.1))
    underline.fill.solid()
    underline.fill.fore_color.rgb = ORANGE
    underline.line.fill.background()

    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(2.5), Inches(1.1), Inches(8), Inches(0.5))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "别被照片骗了，这其实是水上拖拉机"
    p.font.size = Pt(20)
    p.font.name = "Microsoft YaHei"
    p.font.color.rgb = GRAY_TEXT
    p.alignment = PP_ALIGN.CENTER

    # --- 辅助函数：添加带下划线的区块标题 ---
    def add_section_title(left, top, text):
        box = slide.shapes.add_textbox(left, top, Inches(2), Inches(0.5))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.name = "Microsoft YaHei"
        
        ul = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + Inches(0.4), Inches(1.5), Inches(0.08))
        ul.fill.solid()
        ul.fill.fore_color.rgb = ORANGE
        ul.line.fill.background()

    # ==========================================
    # --- Section 1: 现实反差 (左侧) ---
    # ==========================================
    add_section_title(Inches(0.8), Inches(1.8), "现实反差")
    
    t1 = slide.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(2), Inches(0.4))
    t1.text_frame.text = "理想中：优雅的竹筏"
    
    # 图片占位 1 (传统竹筏)
    img1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.8), Inches(2.2), Inches(1.5))
    img1.fill.solid()
    img1.fill.fore_color.rgb = IMG_BG
    img1.line.color.rgb = WHITE
    img1.line.width = Pt(3)
    
    # 红色大叉
    cross = slide.shapes.add_shape(MSO_SHAPE.MATH_MULTIPLY, Inches(1.5), Inches(3.1), Inches(0.8), Inches(0.8))
    cross.fill.solid()
    cross.fill.fore_color.rgb = RED
    cross.line.fill.background()

    # 气泡提示
    bubble = slide.shapes.add_shape(MSO_SHAPE.CLOUD_CALLOUT, Inches(2.5), Inches(2.5), Inches(1.5), Inches(1.0))
    bubble.fill.solid()
    bubble.fill.fore_color.rgb = BLUE_LIGHT
    bubble.line.fill.background()
    tf = bubble.text_frame
    p = tf.paragraphs[0]
    p.text = "其实是\nPVC管+马达"
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 图片占位 2 (PVC竹筏)
    img2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(3.8), Inches(2.8), Inches(2.0))
    img2.fill.solid()
    img2.fill.fore_color.rgb = IMG_BG
    img2.line.color.rgb = WHITE
    img2.line.width = Pt(3)

    # 底部说明文字
    cap1 = slide.shapes.add_textbox(Inches(0.8), Inches(5.9), Inches(3.5), Inches(0.5))
    tf = cap1.text_frame
    p = tf.paragraphs[0]
    p.text = "优雅的竹筏其实是PVC管+马达"
    p.font.size = Pt(14)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # ==========================================
    # --- Section 2: 搞笑瞬间 (中间) ---
    # ==========================================
    add_section_title(Inches(5.2), Inches(1.8), "搞笑瞬间")

    # 图片占位
    img3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.0), Inches(2.5), Inches(3.5), Inches(2.8))
    img3.fill.solid()
    img3.fill.fore_color.rgb = IMG_BG
    img3.line.color.rgb = WHITE
    img3.line.width = Pt(3)

    # 拟声词
    t_boom = slide.shapes.add_textbox(Inches(6.5), Inches(2.8), Inches(1.5), Inches(0.5))
    t_boom.text_frame.text = "轰隆隆!\nBOOM!"
    t_boom.text_frame.paragraphs[0].font.bold = True

    # 吐槽气泡
    t_omg = slide.shapes.add_shape(MSO_SHAPE.OVAL_CALLOUT, Inches(7.0), Inches(4.0), Inches(1.2), Inches(0.8))
    t_omg.fill.solid()
    t_omg.fill.fore_color.rgb = RGBColor(0xA0, 0xD0, 0xA0)
    t_omg.line.fill.background()
    tf = t_omg.text_frame
    p = tf.paragraphs[0]
    p.text = "OMG!\n诗意呢?"
    p.font.size = Pt(10)
    p.alignment = PP_ALIGN.CENTER

    # 底部说明文字
    cap2 = slide.shapes.add_textbox(Inches(4.8), Inches(5.4), Inches(4.0), Inches(0.5))
    tf = cap2.text_frame
    p = tf.paragraphs[0]
    p.text = "马达启动时的黑烟与诗意山水的完美融合"
    p.font.size = Pt(14)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # ==========================================
    # --- Section 3: 猎奇体验 (右侧) ---
    # ==========================================
    add_section_title(Inches(9.5), Inches(1.8), "猎奇体验")

    # 图片占位
    img4 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.2), Inches(2.5), Inches(3.5), Inches(2.8))
    img4.fill.solid()
    img4.fill.fore_color.rgb = IMG_BG
    img4.line.color.rgb = WHITE
    img4.line.width = Pt(3)

    # 标签：移动超市
    t_shop = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10.0), Inches(2.8), Inches(1.2), Inches(0.4))
    t_shop.fill.solid()
    t_shop.fill.fore_color.rgb = RGBColor(0xFF, 0xDD, 0xAA)
    t_shop.line.fill.background()
    tf = t_shop.text_frame
    p = tf.paragraphs[0]
    p.text = "移动超市"
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_TEXT
    p.alignment = PP_ALIGN.CENTER

    # 标签：现烤鱼
    t_fish = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.5), Inches(4.5), Inches(1.0), Inches(0.4))
    t_fish.fill.solid()
    t_fish.fill.fore_color.rgb = RGBColor(0xFF, 0xDD, 0xAA)
    t_fish.line.fill.background()
    tf = t_fish.text_frame
    p = tf.paragraphs[0]
    p.text = "现烤鱼"
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_TEXT
    p.alignment = PP_ALIGN.CENTER

    # 底部说明文字
    cap3 = slide.shapes.add_textbox(Inches(9.0), Inches(5.4), Inches(4.0), Inches(0.5))
    tf = cap3.text_frame
    p = tf.paragraphs[0]
    p.text = "江上的“移动超市”，划着竹筏卖烤鱼"
    p.font.size = Pt(14)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # --- 区块之间的引导箭头 ---
    arrow1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.5), Inches(1.9), Inches(0.8), Inches(0.3))
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = GRAY_TEXT
    arrow1.line.fill.background()

    arrow2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(7.8), Inches(1.9), Inches(0.8), Inches(0.3))
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = GRAY_TEXT
    arrow2.line.fill.background()

    # ==========================================
    # --- 底部图表区域: 马达轰鸣声 vs 优雅度 ---
    # ==========================================
    bot_title = slide.shapes.add_textbox(Inches(6.5), Inches(5.8), Inches(3.5), Inches(0.5))
    tf = bot_title.text_frame
    p = tf.paragraphs[0]
    p.text = "马达轰鸣声 vs 优雅度"
    p.font.size = Pt(18)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    bot_ul = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(6.2), Inches(2.9), Inches(0.08))
    bot_ul.fill.solid()
    bot_ul.fill.fore_color.rgb = ORANGE
    bot_ul.line.fill.background()

    # 渐变进度条 (使用纯色箭头代替)
    bar = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.0), Inches(6.4), Inches(5.0), Inches(0.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0xFF, 0x98, 0x00) # 橙色
    bar.line.color.rgb = DARK_TEXT
    bar.line.width = Pt(1)

    # 左侧文本 (优雅度)
    left_txt = slide.shapes.add_textbox(Inches(5.0), Inches(6.1), Inches(1.0), Inches(0.8))
    tf = left_txt.text_frame
    tf.text = "🕊️\n优雅度\n(最高)"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    if len(tf.paragraphs) > 1:
        tf.paragraphs[1].font.size = Pt(12)
        tf.paragraphs[1].font.bold = True
        tf.paragraphs[1].alignment = PP_ALIGN.CENTER
    if len(tf.paragraphs) > 2:
        tf.paragraphs[2].font.size = Pt(10)
        tf.paragraphs[2].alignment = PP_ALIGN.CENTER

    # 右侧文本 (马达轰鸣声)
    right_txt = slide.shapes.add_textbox(Inches(11.2), Inches(6.1), Inches(1.5), Inches(0.8))
    tf = right_txt.text_frame
    tf.text = "🚜\n马达轰鸣声\n(最大)"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    if len(tf.paragraphs) > 1:
        tf.paragraphs[1].font.size = Pt(12)
        tf.paragraphs[1].font.bold = True
        tf.paragraphs[1].alignment = PP_ALIGN.CENTER
    if len(tf.paragraphs) > 2:
        tf.paragraphs[2].font.size = Pt(10)
        tf.paragraphs[2].alignment = PP_ALIGN.CENTER

    # 进度条上的标签
    labels = ["静音", "嗡嗡...", "轰隆隆!!!", "VROOM!!!"]
    x_positions = [6.3, 7.3, 8.5, 9.8]
    for i, text in enumerate(labels):
        lbl = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_positions[i]), Inches(6.8), Inches(1.0), Inches(0.3))
        lbl.fill.solid()
        lbl.fill.fore_color.rgb = WHITE
        lbl.line.color.rgb = DARK_TEXT
        lbl.line.width = Pt(1)
        tf = lbl.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = DARK_TEXT
        p.alignment = PP_ALIGN.CENTER

    # 底部总结文字
    bot_cap = slide.shapes.add_textbox(Inches(6.0), Inches(7.2), Inches(5.0), Inches(0.4))
    tf = bot_cap.text_frame
    p = tf.paragraphs[0]
    p.text = "马达一响，优雅全无 (进度条图表)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER



# ── Main ──
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

s0 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_1(s0)
prs.save(OUTPUT_PATH)
