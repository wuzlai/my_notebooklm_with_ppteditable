from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# ── Color Palette ──
BLUE_HEADER = RGBColor(0x5B, 0x9B, 0xD5)
BLUE_DARK   = RGBColor(0x4A, 0x86, 0xC8)
CYAN        = RGBColor(0x00, 0xBC, 0xD4)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x33, 0x33, 0x33)
GRAY_TEXT   = RGBColor(0x55, 0x55, 0x55)
GRAY_BAR    = RGBColor(0xB0, 0xBE, 0xC5)
RED         = RGBColor(0xE5, 0x39, 0x35)
GREEN       = RGBColor(0x43, 0xA0, 0x47)
ORANGE      = RGBColor(0xFF, 0x98, 0x00)

ICON_BG  = RGBColor(0xE3, 0xE8, 0xED)
ICON_FG  = RGBColor(0x54, 0x6E, 0x7A)
FONT_NAME = "Microsoft YaHei"
SLIDE_WIDTH = Inches(13.333)
HEADER_H    = Inches(0.75)
SUBTITLE_Y  = Inches(0.95)


def add_header_banner(slide, title_text, bg_color=None):
    if bg_color is None:
        bg_color = BLUE_HEADER
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, HEADER_H
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = bg_color
    banner.line.fill.background()
    tf = banner.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(26)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME


def add_subtitle(slide, text, left, top, width=Inches(12), font_size=Pt(18)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.font.name = FONT_NAME
    return txBox


def add_icon_box(slide, left, top, symbol, size=Inches(0.48)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ICON_BG
    shape.line.fill.background()
    shape.adjustments[0] = 0.25
    tf = shape.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = symbol
    p.font.size = Pt(18)
    p.font.color.rgb = ICON_FG
    p.font.bold = False
    return shape


def add_bullet_item(slide, left, top, symbol, label, description,
                    width=Inches(5.5), desc_size=Pt(13)):
    add_icon_box(slide, left, top, symbol)
    text_left = left + Inches(0.65)
    txBox = slide.shapes.add_textbox(text_left, top - Inches(0.02), width, Inches(0.65))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run_label = p.add_run()
    run_label.text = label + "\uff1a"
    run_label.font.size = Pt(14)
    run_label.font.color.rgb = BLACK
    run_label.font.bold = True
    run_label.font.name = FONT_NAME
    run_desc = p.add_run()
    run_desc.text = description
    run_desc.font.size = desc_size
    run_desc.font.color.rgb = GRAY_TEXT
    run_desc.font.bold = False
    run_desc.font.name = FONT_NAME
    return txBox


def add_conclusion_box(slide, left, top, width, text, font_size=Pt(13)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = BLACK
    run.font.bold = True
    run.font.name = FONT_NAME
    return txBox


def add_table(slide, left, top, width, height, rows, cols, data,
              header_color=None, col_widths=None):
    if header_color is None:
        header_color = BLUE_HEADER
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if data[r][c] is not None else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(11)
                    run.font.name = FONT_NAME
                    if r == 0:
                        run.font.color.rgb = WHITE
                        run.font.bold = True
                    else:
                        run.font.color.rgb = BLACK
                        run.font.bold = False
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else RGBColor(0xF5, 0xF5, 0xF5)
            cell.margin_left = Pt(5)
            cell.margin_right = Pt(5)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
    return table_shape


def add_bar_chart(slide, left, top, width, height,
                  categories, values, title="", bar_colors=None):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', values)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    chart.chart_style = 2
    plot = chart.plots[0]
    plot.gap_width = 100
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CYAN
    series.has_data_labels = True
    dl = series.data_labels
    dl.font.size = Pt(13)
    dl.font.bold = True
    dl.font.color.rgb = BLACK
    dl.number_format = '0.#'
    dl.show_value = True
    dl.label_position = XL_LABEL_POSITION.OUTSIDE_END
    if bar_colors:
        for i, color in enumerate(bar_colors):
            pt = series.points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = color
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(12)
    cat_axis.tick_labels.font.name = FONT_NAME
    cat_axis.major_tick_mark = 2
    cat_axis.format.line.fill.background()
    val_axis = chart.value_axis
    val_axis.visible = False
    val_axis.major_tick_mark = 2
    val_axis.format.line.fill.background()
    val_axis.major_gridlines.format.line.fill.background()
    if title:
        chart.has_title = True
        ct = chart.chart_title.text_frame.paragraphs[0]
        ct.text = title
        ct.font.size = Pt(14)
        ct.font.bold = True
        ct.font.name = FONT_NAME
    else:
        chart.has_title = False
    return chart_frame


def add_callout_label(slide, left, top, text, bg_color=None, font_size=Pt(11)):
    if bg_color is None:
        bg_color = CYAN
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.3), Inches(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME
    return shape


def add_data_card(slide, left, top, width, height, value, label,
                  value_color=None, bg_color=None):
    if value_color is None:
        value_color = CYAN
    if bg_color is None:
        bg_color = WHITE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(3)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run1 = p1.add_run()
    run1.text = str(value)
    run1.font.size = Pt(24)
    run1.font.color.rgb = value_color
    run1.font.bold = True
    run1.font.name = FONT_NAME
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(10)
    run2.font.color.rgb = GRAY_TEXT
    run2.font.bold = False
    run2.font.name = FONT_NAME
    return shape


OUTPUT_PATH = r"C:\Users\Administrator\Desktop\notebook\my_notebooklm_with_ppteditable\projects\旅行分享1\最终文档\ppt_slides\slide_10.pptx"


# ── Slide 1 ──

def build_slide_1(slide):
    # 1. Background (Green irregular shape approximated by a rectangle)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.6), Inches(11.7), Inches(6.3))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x2E, 0x9E, 0x7B)
    bg.line.fill.background()

    # 2. Title
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11.333), Inches(1.0))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "总结：桂林，一个越怪越美的地方"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.font.name = "Microsoft YaHei"
    p.font.shadow = True

    # 3. Subtitle
    sub_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(11.333), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "这里的山水有灵，这里的人们有戏"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.font.name = "Microsoft YaHei"
    p.font.shadow = True

    # 4. Item 1: Core Viewpoint
    # Camera Icon (Constructed with shapes)
    cam_base = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.6), Inches(3.4), Inches(0.8), Inches(0.6))
    cam_base.fill.solid()
    cam_base.fill.fore_color.rgb = RGBColor(0xE0, 0xF7, 0xFA)
    cam_base.line.color.rgb = RGBColor(0x00, 0x00, 0x00)
    cam_base.line.width = Pt(2)
    
    cam_lens = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.75), Inches(3.45), Inches(0.5), Inches(0.5))
    cam_lens.fill.solid()
    cam_lens.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    cam_lens.line.color.rgb = RGBColor(0x00, 0x00, 0x00)
    cam_lens.line.width = Pt(2)

    cam_heart = slide.shapes.add_shape(MSO_SHAPE.HEART, Inches(1.85), Inches(3.55), Inches(0.3), Inches(0.3))
    cam_heart.fill.solid()
    cam_heart.fill.fore_color.rgb = RGBColor(0xFF, 0x66, 0x00)
    cam_heart.line.fill.background()

    cam_flash = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.15), Inches(3.48), Inches(0.12), Inches(0.12))
    cam_flash.fill.solid()
    cam_flash.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    cam_flash.line.color.rgb = RGBColor(0x00, 0x00, 0x00)
    cam_flash.line.width = Pt(1)

    # Text
    tb1 = slide.shapes.add_textbox(Inches(2.6), Inches(3.4), Inches(9.0), Inches(0.8))
    p1 = tb1.text_frame.paragraphs[0]
    p1.text = "核心观点：风景是背景，有趣才是旅行的灵魂"
    p1.font.size = Pt(28)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p1.font.name = "Microsoft YaHei"
    p1.font.shadow = True

    # 5. Item 2: Novelty Index
    # Star Icon
    star_icon1 = slide.shapes.add_shape(MSO_SHAPE.5_POINT_STAR, Inches(1.7), Inches(4.5), Inches(0.6), Inches(0.6))
    star_icon1.fill.solid()
    star_icon1.fill.fore_color.rgb = RGBColor(0xFF, 0x66, 0x00)
    star_icon1.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    star_icon1.line.width = Pt(1.5)

    # Text
    tb2 = slide.shapes.add_textbox(Inches(2.6), Inches(4.4), Inches(2.5), Inches(0.8))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = "猎奇指数："
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p2.font.name = "Microsoft YaHei"
    p2.font.shadow = True

    # Rating Stars
    for i in range(5):
        star = slide.shapes.add_shape(MSO_SHAPE.5_POINT_STAR, Inches(4.8 + i*0.7), Inches(4.5), Inches(0.6), Inches(0.6))
        star.fill.solid()
        star.fill.fore_color.rgb = RGBColor(0xFF, 0x66, 0x00)
        star.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        star.line.width = Pt(1.5)

    # 6. Item 3: Recommendation Index
    # Star Icon
    star_icon2 = slide.shapes.add_shape(MSO_SHAPE.5_POINT_STAR, Inches(1.7), Inches(5.5), Inches(0.6), Inches(0.6))
    star_icon2.fill.solid()
    star_icon2.fill.fore_color.rgb = RGBColor(0xFF, 0x66, 0x00)
    star_icon2.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    star_icon2.line.width = Pt(1.5)

    # Text
    tb3 = slide.shapes.add_textbox(Inches(2.6), Inches(5.4), Inches(2.5), Inches(0.8))
    p3 = tb3.text_frame.paragraphs[0]
    p3.text = "推荐指数："
    p3.font.size = Pt(28)
    p3.font.bold = True
    p3.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p3.font.name = "Microsoft YaHei"
    p3.font.shadow = True

    # Rating Stars
    for i in range(5):
        star = slide.shapes.add_shape(MSO_SHAPE.5_POINT_STAR, Inches(4.8 + i*0.7), Inches(5.5), Inches(0.6), Inches(0.6))
        if i < 4:
            star.fill.solid()
            star.fill.fore_color.rgb = RGBColor(0xFF, 0x66, 0x00)
            star.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        else:
            star.fill.solid()
            star.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            star.line.color.rgb = RGBColor(0xFF, 0x66, 0x00)
        star.line.width = Pt(1.5)

    # Comment Text
    tb_comment = slide.shapes.add_textbox(Inches(8.2), Inches(5.45), Inches(4.0), Inches(0.8))
    p_comment = tb_comment.text_frame.paragraphs[0]
    p_comment.text = "(扣一星怕它太骄傲)"
    p_comment.font.size = Pt(24)
    p_comment.font.bold = True
    p_comment.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p_comment.font.name = "Microsoft YaHei"
    p_comment.font.shadow = True

    # 7. Decorations
    # OMG!
    omg_box = slide.shapes.add_textbox(Inches(10.2), Inches(3.2), Inches(2.0), Inches(1.0))
    omg_box.rotation = 15
    p_omg = omg_box.text_frame.paragraphs[0]
    p_omg.text = "OMG!"
    p_omg.font.size = Pt(32)
    p_omg.font.bold = True
    p_omg.font.color.rgb = RGBColor(0xFF, 0x66, 0x00)
    p_omg.font.name = "Arial"

    # Boom!
    boom_box = slide.shapes.add_textbox(Inches(10.5), Inches(4.8), Inches(2.0), Inches(1.0))
    boom_box.rotation = -15
    p_boom = boom_box.text_frame.paragraphs[0]
    p_boom.text = "Boom!"
    p_boom.font.size = Pt(32)
    p_boom.font.bold = True
    p_boom.font.color.rgb = RGBColor(0x80, 0x00, 0x80)
    p_boom.font.name = "Arial"

    # Arrow
    arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(7.6), Inches(4.2), Inches(0.3), Inches(0.5))
    arrow.rotation = 225
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RGBColor(0xFF, 0x33, 0x33)
    arrow.line.fill.background()

    # Decorative circles
    circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.5), Inches(3.1), Inches(0.15), Inches(0.15))
    circle1.fill.background()
    circle1.line.color.rgb = RGBColor(0xFF, 0x66, 0x00)
    circle1.line.width = Pt(2)

    circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.8), Inches(4.8), Inches(0.15), Inches(0.15))
    circle2.fill.background()
    circle2.line.color.rgb = RGBColor(0xFF, 0x33, 0x33)
    circle2.line.width = Pt(2)
    
    # Decorative arc
    arc = slide.shapes.add_shape(MSO_SHAPE.ARC, Inches(10.5), Inches(6.0), Inches(0.8), Inches(0.4))
    arc.rotation = 180
    arc.fill.background()
    arc.line.color.rgb = RGBColor(0x80, 0x00, 0x80)
    arc.line.width = Pt(2)

    # 8. Page Number
    page_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(11.8), Inches(6.5), Inches(1.0), Inches(0.4))
    page_box.fill.solid()
    page_box.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    page_box.line.color.rgb = RGBColor(0x2E, 0x9E, 0x7B)
    page_box.line.width = Pt(2)
    
    tf_page = page_box.text_frame
    tf_page.margin_left = 0
    tf_page.margin_right = 0
    tf_page.margin_top = 0
    tf_page.margin_bottom = 0
    p_page = tf_page.paragraphs[0]
    p_page.text = "第10页"
    p_page.alignment = PP_ALIGN.CENTER
    p_page.font.size = Pt(14)
    p_page.font.bold = True
    p_page.font.color.rgb = RGBColor(0x2E, 0x9E, 0x7B)
    p_page.font.name = "Microsoft YaHei"



# ── Main ──
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

s0 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_1(s0)
prs.save(OUTPUT_PATH)
