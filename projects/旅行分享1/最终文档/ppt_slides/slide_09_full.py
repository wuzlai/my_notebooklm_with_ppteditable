from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# ── Color Palette ──
BLUE_HEADER = RGBColor(0x5B, 0x9B, 0xD5)
BLUE_DARK   = RGBColor(0x4A, 0x86, 0xC8)
CYAN        = RGBColor(0x00, 0xBC, 0xD4)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x33, 0x33, 0x33)
GRAY_TEXT   = RGBColor(0x55, 0x55, 0x55)
GRAY_BAR    = RGBColor(0xB0, 0xBE, 0xC5)
RED         = RGBColor(0xE5, 0x39, 0x35)
GREEN       = RGBColor(0x43, 0xA0, 0x47)
ORANGE      = RGBColor(0xFF, 0x98, 0x00)

ICON_BG  = RGBColor(0xE3, 0xE8, 0xED)
ICON_FG  = RGBColor(0x54, 0x6E, 0x7A)
FONT_NAME = "Microsoft YaHei"
SLIDE_WIDTH = Inches(13.333)
HEADER_H    = Inches(0.75)
SUBTITLE_Y  = Inches(0.95)


def add_header_banner(slide, title_text, bg_color=None):
    if bg_color is None:
        bg_color = BLUE_HEADER
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, HEADER_H
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = bg_color
    banner.line.fill.background()
    tf = banner.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(26)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME


def add_subtitle(slide, text, left, top, width=Inches(12), font_size=Pt(18)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.font.name = FONT_NAME
    return txBox


def add_icon_box(slide, left, top, symbol, size=Inches(0.48)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ICON_BG
    shape.line.fill.background()
    shape.adjustments[0] = 0.25
    tf = shape.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = symbol
    p.font.size = Pt(18)
    p.font.color.rgb = ICON_FG
    p.font.bold = False
    return shape


def add_bullet_item(slide, left, top, symbol, label, description,
                    width=Inches(5.5), desc_size=Pt(13)):
    add_icon_box(slide, left, top, symbol)
    text_left = left + Inches(0.65)
    txBox = slide.shapes.add_textbox(text_left, top - Inches(0.02), width, Inches(0.65))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run_label = p.add_run()
    run_label.text = label + "\uff1a"
    run_label.font.size = Pt(14)
    run_label.font.color.rgb = BLACK
    run_label.font.bold = True
    run_label.font.name = FONT_NAME
    run_desc = p.add_run()
    run_desc.text = description
    run_desc.font.size = desc_size
    run_desc.font.color.rgb = GRAY_TEXT
    run_desc.font.bold = False
    run_desc.font.name = FONT_NAME
    return txBox


def add_conclusion_box(slide, left, top, width, text, font_size=Pt(13)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = BLACK
    run.font.bold = True
    run.font.name = FONT_NAME
    return txBox


def add_table(slide, left, top, width, height, rows, cols, data,
              header_color=None, col_widths=None):
    if header_color is None:
        header_color = BLUE_HEADER
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if data[r][c] is not None else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(11)
                    run.font.name = FONT_NAME
                    if r == 0:
                        run.font.color.rgb = WHITE
                        run.font.bold = True
                    else:
                        run.font.color.rgb = BLACK
                        run.font.bold = False
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else RGBColor(0xF5, 0xF5, 0xF5)
            cell.margin_left = Pt(5)
            cell.margin_right = Pt(5)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
    return table_shape


def add_bar_chart(slide, left, top, width, height,
                  categories, values, title="", bar_colors=None):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', values)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    chart.chart_style = 2
    plot = chart.plots[0]
    plot.gap_width = 100
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CYAN
    series.has_data_labels = True
    dl = series.data_labels
    dl.font.size = Pt(13)
    dl.font.bold = True
    dl.font.color.rgb = BLACK
    dl.number_format = '0.#'
    dl.show_value = True
    dl.label_position = XL_LABEL_POSITION.OUTSIDE_END
    if bar_colors:
        for i, color in enumerate(bar_colors):
            pt = series.points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = color
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(12)
    cat_axis.tick_labels.font.name = FONT_NAME
    cat_axis.major_tick_mark = 2
    cat_axis.format.line.fill.background()
    val_axis = chart.value_axis
    val_axis.visible = False
    val_axis.major_tick_mark = 2
    val_axis.format.line.fill.background()
    val_axis.major_gridlines.format.line.fill.background()
    if title:
        chart.has_title = True
        ct = chart.chart_title.text_frame.paragraphs[0]
        ct.text = title
        ct.font.size = Pt(14)
        ct.font.bold = True
        ct.font.name = FONT_NAME
    else:
        chart.has_title = False
    return chart_frame


def add_callout_label(slide, left, top, text, bg_color=None, font_size=Pt(11)):
    if bg_color is None:
        bg_color = CYAN
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.3), Inches(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME
    return shape


def add_data_card(slide, left, top, width, height, value, label,
                  value_color=None, bg_color=None):
    if value_color is None:
        value_color = CYAN
    if bg_color is None:
        bg_color = WHITE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(3)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run1 = p1.add_run()
    run1.text = str(value)
    run1.font.size = Pt(24)
    run1.font.color.rgb = value_color
    run1.font.bold = True
    run1.font.name = FONT_NAME
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(10)
    run2.font.color.rgb = GRAY_TEXT
    run2.font.bold = False
    run2.font.name = FONT_NAME
    return shape


OUTPUT_PATH = r"C:\Users\Administrator\Desktop\notebook\my_notebooklm_with_ppteditable\projects\旅行分享1\最终文档\ppt_slides\slide_09.pptx"


# ── Slide 1 ──

def build_slide_1(slide):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

    # Colors
    BG_COLOR = RGBColor(245, 242, 235)
    ORANGE_DECOR = RGBColor(235, 104, 65)
    PURPLE_DECOR = RGBColor(103, 58, 183)
    YELLOW_HL = RGBColor(255, 235, 59)
    BLACK = RGBColor(0, 0, 0)
    WHITE = RGBColor(255, 255, 255)
    BOX_ORANGE = RGBColor(255, 112, 67)
    BOX_GREEN = RGBColor(129, 199, 132)
    BOX_YELLOW = RGBColor(255, 235, 59)

    # 1. Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_COLOR
    bg.line.fill.background()

    # Corner decorations (simulating torn paper)
    tr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(11), Inches(0), Inches(2.333), Inches(1.5))
    tr.rotation = 180
    tr.fill.solid()
    tr.fill.fore_color.rgb = ORANGE_DECOR
    tr.line.fill.background()

    bl = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(0), Inches(5.5), Inches(2.5), Inches(2))
    bl.fill.solid()
    bl.fill.fore_color.rgb = ORANGE_DECOR
    bl.line.fill.background()

    br = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(12), Inches(6.5), Inches(1.333), Inches(1))
    br.rotation = 270
    br.fill.solid()
    br.fill.fore_color.rgb = PURPLE_DECOR
    br.line.fill.background()

    # 2. Title & Subtitle
    # Title Highlight
    hl = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.9), Inches(8.0), Inches(0.6))
    hl.fill.solid()
    hl.fill.fore_color.rgb = YELLOW_HL
    hl.line.fill.background()

    # Title Text
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.7), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "桂林生存法则：如何优雅地避坑"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = BLACK

    # Subtitle Text
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(5), Inches(0.5))
    tf = sub_box.text_frame
    p = tf.add_paragraph()
    p.text = "猎奇博主的血泪教训"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = BLACK

    # 3. Left Visuals
    # Landscape Photo Base (White border)
    photo_base = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.0), Inches(5.0), Inches(3.5))
    photo_base.fill.solid()
    photo_base.fill.fore_color.rgb = WHITE
    photo_base.line.fill.background()
    photo_base.rotation = -2

    # Landscape Inner
    land = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(3.2), Inches(4.6), Inches(3.1))
    land.fill.solid()
    land.fill.fore_color.rgb = RGBColor(46, 139, 87)
    land.line.fill.background()
    land.rotation = -2
    tf = land.text_frame
    tf.text = "\n\n桂林山水"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE

    # 20 RMB Note
    note = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.5), Inches(2.0), Inches(2.8), Inches(1.4))
    note.fill.solid()
    note.fill.fore_color.rgb = RGBColor(245, 235, 235)
    note.line.color.rgb = RGBColor(200, 200, 200)
    note.rotation = -10
    tf = note.text_frame
    tf.text = "20\n中国人民银行"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(150, 50, 50)

    # Blogger (Simulated)
    blogger_body = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.5), Inches(1.8), Inches(2.5))
    blogger_body.fill.solid()
    blogger_body.fill.fore_color.rgb = RGBColor(255, 193, 7)
    blogger_body.line.color.rgb = BLACK
    blogger_body.line.width = Pt(2)
    tf = blogger_body.text_frame
    tf.text = "\n\n博主"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = BLACK

    # Boom!
    boom = slide.shapes.add_shape(MSO_SHAPE.EXPLOSION1, Inches(0.5), Inches(2.2), Inches(1.8), Inches(1.2))
    boom.fill.solid()
    boom.fill.fore_color.rgb = RGBColor(255, 87, 34)
    boom.line.color.rgb = BLACK
    boom.line.width = Pt(2)
    boom.rotation = -15
    tf = boom.text_frame
    p = tf.paragraphs[0]
    p.text = "Boom!"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = BLACK

    # OMG!
    omg = slide.shapes.add_shape(MSO_SHAPE.EXPLOSION2, Inches(4.2), Inches(5.5), Inches(1.5), Inches(1.2))
    omg.fill.solid()
    omg.fill.fore_color.rgb = RGBColor(255, 87, 34)
    omg.line.color.rgb = BLACK
    omg.line.width = Pt(2)
    omg.rotation = 10
    tf = omg.text_frame
    p = tf.paragraphs[0]
    p.text = "OMG!"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = BLACK

    # 4. Right Pitfall Boxes
    def add_pitfall_box(slide, left, top, img_bg_color, text_bg_color, text_color, img_text, desc_text):
        box_w = Inches(2.8)
        img_h = Inches(1.5)
        txt_h = Inches(0.9)

        # Image part
        img_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, box_w, img_h)
        img_box.fill.solid()
        img_box.fill.fore_color.rgb = img_bg_color
        img_box.line.color.rgb = BLACK
        img_box.line.width = Pt(2)
        tf = img_box.text_frame
        tf.text = img_text
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = BLACK

        # Text part
        txt_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + img_h, box_w, txt_h)
        txt_box.fill.solid()
        txt_box.fill.fore_color.rgb = text_bg_color
        txt_box.line.color.rgb = BLACK
        txt_box.line.width = Pt(2)
        tf = txt_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc_text
        p.alignment = PP_ALIGN.LEFT
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = text_color

        # Warning Triangle
        warn = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, left - Inches(0.25), top - Inches(0.25), Inches(0.6), Inches(0.6))
        warn.fill.solid()
        warn.fill.fore_color.rgb = YELLOW_HL
        warn.line.color.rgb = BLACK
        warn.line.width = Pt(2)
        tf = warn.text_frame
        p = tf.paragraphs[0]
        p.text = "!"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = BLACK

    # Box 1: Auntie
    add_pitfall_box(slide, Inches(6.5), Inches(1.5),
                    BOX_ORANGE, BOX_YELLOW, BLACK,
                    "帅哥/美女，来吃鱼！\n\nFAKE NEWS!  ❌ 坑！",
                    "别相信路边喊你“帅哥/美女”去吃鱼的阿姨")

    # Box 2: Cormorant 1
    add_pitfall_box(slide, Inches(10.0), Inches(1.2),
                    BOX_GREEN, BLACK, WHITE,
                    "¥999+\n\n(鸬鹚拍照)",
                    "拍照前先问价格，否则鸬鹚的身价比你还高")

    # Box 3: Cormorant 2
    add_pitfall_box(slide, Inches(6.5), Inches(4.5),
                    BOX_GREEN, BLACK, WHITE,
                    "¥999+\n\n(鸬鹚拍照)",
                    "拍照前先问价格，否则鸬鹚的身价比你还高")

    # Box 4: Luosifen
    add_pitfall_box(slide, Inches(10.0), Inches(4.5),
                    BOX_ORANGE, BLACK, WHITE,
                    "螺蛳粉 + 肠胃药\n\n🔥 危险！",
                    "备好肠胃药，螺蛳粉的后劲比酒还大")

    # 5. Arrows (Using connectors)
    conn1 = slide.shapes.add_connector(MSO_CONNECTOR.CURVE, Inches(9.3), Inches(2.2), Inches(10.0), Inches(2.2))
    conn1.line.color.rgb = BLACK
    conn1.line.width = Pt(2)

    conn2 = slide.shapes.add_connector(MSO_CONNECTOR.CURVE, Inches(6.3), Inches(3.0), Inches(6.3), Inches(4.5))
    conn2.line.color.rgb = BLACK
    conn2.line.width = Pt(2)

    conn3 = slide.shapes.add_connector(MSO_CONNECTOR.CURVE, Inches(9.3), Inches(5.2), Inches(10.0), Inches(5.2))
    conn3.line.color.rgb = BLACK
    conn3.line.width = Pt(2)

    conn4 = slide.shapes.add_connector(MSO_CONNECTOR.CURVE, Inches(11.4), Inches(3.6), Inches(11.4), Inches(4.5))
    conn4.line.color.rgb = BLACK
    conn4.line.width = Pt(2)

    # 6. Page Number & Icon
    page_oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.5), Inches(6.8), Inches(0.6), Inches(0.4))
    page_oval.fill.solid()
    page_oval.fill.fore_color.rgb = WHITE
    page_oval.line.color.rgb = BLACK
    page_oval.line.width = Pt(1.5)
    tf = page_oval.text_frame
    p = tf.paragraphs[0]
    p.text = "09"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = BLACK

    face = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(12.2), Inches(6.6), Inches(0.8), Inches(0.8))
    face.fill.solid()
    face.fill.fore_color.rgb = RGBColor(255, 224, 178)
    face.line.color.rgb = BLACK
    face.line.width = Pt(1.5)
    tf = face.text_frame
    p = tf.paragraphs[0]
    p.text = "博主"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = BLACK



# ── Main ──
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

s0 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_1(s0)
prs.save(OUTPUT_PATH)
