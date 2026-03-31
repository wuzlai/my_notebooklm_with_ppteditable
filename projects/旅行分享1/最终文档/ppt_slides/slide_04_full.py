from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# ── Color Palette ──
BLUE_HEADER = RGBColor(0x5B, 0x9B, 0xD5)
BLUE_DARK   = RGBColor(0x4A, 0x86, 0xC8)
CYAN        = RGBColor(0x00, 0xBC, 0xD4)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x33, 0x33, 0x33)
GRAY_TEXT   = RGBColor(0x55, 0x55, 0x55)
GRAY_BAR    = RGBColor(0xB0, 0xBE, 0xC5)
RED         = RGBColor(0xE5, 0x39, 0x35)
GREEN       = RGBColor(0x43, 0xA0, 0x47)
ORANGE      = RGBColor(0xFF, 0x98, 0x00)

ICON_BG  = RGBColor(0xE3, 0xE8, 0xED)
ICON_FG  = RGBColor(0x54, 0x6E, 0x7A)
FONT_NAME = "Microsoft YaHei"
SLIDE_WIDTH = Inches(13.333)
HEADER_H    = Inches(0.75)
SUBTITLE_Y  = Inches(0.95)


def add_header_banner(slide, title_text, bg_color=None):
    if bg_color is None:
        bg_color = BLUE_HEADER
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, HEADER_H
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = bg_color
    banner.line.fill.background()
    tf = banner.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(26)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME


def add_subtitle(slide, text, left, top, width=Inches(12), font_size=Pt(18)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.font.name = FONT_NAME
    return txBox


def add_icon_box(slide, left, top, symbol, size=Inches(0.48)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ICON_BG
    shape.line.fill.background()
    shape.adjustments[0] = 0.25
    tf = shape.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = symbol
    p.font.size = Pt(18)
    p.font.color.rgb = ICON_FG
    p.font.bold = False
    return shape


def add_bullet_item(slide, left, top, symbol, label, description,
                    width=Inches(5.5), desc_size=Pt(13)):
    add_icon_box(slide, left, top, symbol)
    text_left = left + Inches(0.65)
    txBox = slide.shapes.add_textbox(text_left, top - Inches(0.02), width, Inches(0.65))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run_label = p.add_run()
    run_label.text = label + "\uff1a"
    run_label.font.size = Pt(14)
    run_label.font.color.rgb = BLACK
    run_label.font.bold = True
    run_label.font.name = FONT_NAME
    run_desc = p.add_run()
    run_desc.text = description
    run_desc.font.size = desc_size
    run_desc.font.color.rgb = GRAY_TEXT
    run_desc.font.bold = False
    run_desc.font.name = FONT_NAME
    return txBox


def add_conclusion_box(slide, left, top, width, text, font_size=Pt(13)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = BLACK
    run.font.bold = True
    run.font.name = FONT_NAME
    return txBox


def add_table(slide, left, top, width, height, rows, cols, data,
              header_color=None, col_widths=None):
    if header_color is None:
        header_color = BLUE_HEADER
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if data[r][c] is not None else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(11)
                    run.font.name = FONT_NAME
                    if r == 0:
                        run.font.color.rgb = WHITE
                        run.font.bold = True
                    else:
                        run.font.color.rgb = BLACK
                        run.font.bold = False
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else RGBColor(0xF5, 0xF5, 0xF5)
            cell.margin_left = Pt(5)
            cell.margin_right = Pt(5)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
    return table_shape


def add_bar_chart(slide, left, top, width, height,
                  categories, values, title="", bar_colors=None):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', values)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    chart.chart_style = 2
    plot = chart.plots[0]
    plot.gap_width = 100
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CYAN
    series.has_data_labels = True
    dl = series.data_labels
    dl.font.size = Pt(13)
    dl.font.bold = True
    dl.font.color.rgb = BLACK
    dl.number_format = '0.#'
    dl.show_value = True
    dl.label_position = XL_LABEL_POSITION.OUTSIDE_END
    if bar_colors:
        for i, color in enumerate(bar_colors):
            pt = series.points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = color
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(12)
    cat_axis.tick_labels.font.name = FONT_NAME
    cat_axis.major_tick_mark = 2
    cat_axis.format.line.fill.background()
    val_axis = chart.value_axis
    val_axis.visible = False
    val_axis.major_tick_mark = 2
    val_axis.format.line.fill.background()
    val_axis.major_gridlines.format.line.fill.background()
    if title:
        chart.has_title = True
        ct = chart.chart_title.text_frame.paragraphs[0]
        ct.text = title
        ct.font.size = Pt(14)
        ct.font.bold = True
        ct.font.name = FONT_NAME
    else:
        chart.has_title = False
    return chart_frame


def add_callout_label(slide, left, top, text, bg_color=None, font_size=Pt(11)):
    if bg_color is None:
        bg_color = CYAN
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.3), Inches(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME
    return shape


def add_data_card(slide, left, top, width, height, value, label,
                  value_color=None, bg_color=None):
    if value_color is None:
        value_color = CYAN
    if bg_color is None:
        bg_color = WHITE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(3)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run1 = p1.add_run()
    run1.text = str(value)
    run1.font.size = Pt(24)
    run1.font.color.rgb = value_color
    run1.font.bold = True
    run1.font.name = FONT_NAME
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(10)
    run2.font.color.rgb = GRAY_TEXT
    run2.font.bold = False
    run2.font.name = FONT_NAME
    return shape


OUTPUT_PATH = r"C:\Users\Administrator\Desktop\notebook\my_notebooklm_with_ppteditable\projects\旅行分享1\最终文档\ppt_slides\slide_04.pptx"


# ── Slide 1 ──

def build_slide_1(slide):
    # 1. 添加左上角页码
    tx_box_page = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(2.0), Inches(0.8))
    tf_page = tx_box_page.text_frame
    p_page = tf_page.paragraphs[0]
    p_page.text = "第4页"
    p_page.font.name = "Microsoft YaHei"
    p_page.font.size = Pt(32)
    p_page.font.bold = True
    p_page.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # 2. 添加主标题
    tx_box_title = slide.shapes.add_textbox(Inches(3.5), Inches(0.6), Inches(9.0), Inches(0.8))
    tf_title = tx_box_title.text_frame
    p_title = tf_title.paragraphs[0]
    p_title.text = "象鼻山：这头大象到底在喝什么？"
    p_title.font.name = "Microsoft YaHei"
    p_title.font.size = Pt(36)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # 3. 添加副标题
    tx_box_subtitle = slide.shapes.add_textbox(Inches(6.5), Inches(1.3), Inches(6.0), Inches(0.6))
    tf_subtitle = tx_box_subtitle.text_frame
    p_subtitle = tf_subtitle.paragraphs[0]
    p_subtitle.text = "关于桂林城徽的终极猜想"
    p_subtitle.font.name = "Microsoft YaHei"
    p_subtitle.font.size = Pt(24)
    p_subtitle.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
    p_subtitle.alignment = PP_ALIGN.RIGHT

    # 4. 左侧大象喝奶茶插图占位符及“奶茶？”文字
    # 插图占位符
    pic_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.5))
    pic_shape.fill.solid()
    pic_shape.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    pic_shape.line.fill.background()
    tf_pic = pic_shape.text_frame
    p_pic = tf_pic.paragraphs[0]
    p_pic.text = "[大象喝奶茶插图区域]"
    p_pic.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
    p_pic.alignment = PP_ALIGN.CENTER

    # “奶茶？”文字
    tx_box_tea = slide.shapes.add_textbox(Inches(5.2), Inches(4.8), Inches(1.5), Inches(0.8))
    tx_box_tea.rotation = -15
    tf_tea = tx_box_tea.text_frame
    p_tea = tf_tea.paragraphs[0]
    p_tea.text = "奶茶？"
    p_tea.font.name = "Microsoft YaHei"
    p_tea.font.size = Pt(28)
    p_tea.font.bold = True
    p_tea.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # 5. 右侧图文列表项
    # 列表项 1
    icon1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.8), Inches(2.6), Inches(0.8), Inches(0.8))
    icon1.fill.solid()
    icon1.fill.fore_color.rgb = RGBColor(0xE1, 0xF5, 0xFE)
    icon1.line.fill.background()
    icon1.text_frame.text = "📷"
    
    tx1 = slide.shapes.add_textbox(Inches(7.8), Inches(2.5), Inches(5.0), Inches(1.0))
    tf1 = tx1.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "视觉错位："
    p1.font.name = "Microsoft YaHei"
    p1.font.bold = True
    p1.font.size = Pt(18)
    p1.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    run1 = p1.add_run()
    run1.text = "从哪个角度看最像一只喝醉的大象"
    run1.font.bold = False
    run1.font.size = Pt(18)

    # 列表项 2
    icon2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.8), Inches(4.1), Inches(0.8), Inches(0.8))
    icon2.fill.solid()
    icon2.fill.fore_color.rgb = RGBColor(0xF3, 0xE5, 0xF5)
    icon2.line.fill.background()
    icon2.text_frame.text = "🍷"
    
    tx2 = slide.shapes.add_textbox(Inches(7.8), Inches(4.0), Inches(5.0), Inches(1.0))
    tf2 = tx2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "猎奇冷知识："
    p2.font.name = "Microsoft YaHei"
    p2.font.bold = True
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    run2 = p2.add_run()
    run2.text = "象鼻山内部其实是空的（藏酒洞）"
    run2.font.bold = False
    run2.font.size = Pt(18)

    # 列表项 3
    icon3 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.8), Inches(5.6), Inches(0.8), Inches(0.8))
    icon3.fill.solid()
    icon3.fill.fore_color.rgb = RGBColor(0xE0, 0xF2, 0xF1)
    icon3.line.fill.background()
    icon3.text_frame.text = "💦"
    
    tx3 = slide.shapes.add_textbox(Inches(7.8), Inches(5.5), Inches(5.0), Inches(1.0))
    tf3 = tx3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.text = "吐槽点："
    p3.font.name = "Microsoft YaHei"
    p3.font.bold = True
    p3.font.size = Pt(18)
    p3.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    run3 = p3.add_run()
    run3.text = "为了拍一张“喂象照”，我差点掉进江里"
    run3.font.bold = False
    run3.font.size = Pt(18)

    # 6. 底部右侧装饰元素 (BooM!, OMG!, 箭头)
    # BooM!
    boom_box = slide.shapes.add_textbox(Inches(8.5), Inches(6.8), Inches(1.2), Inches(0.6))
    boom_box.rotation = -10
    tf_boom = boom_box.text_frame
    p_boom = tf_boom.paragraphs[0]
    p_boom.text = "BooM!"
    p_boom.font.name = "Microsoft YaHei"
    p_boom.font.size = Pt(16)
    p_boom.font.bold = True
    p_boom.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # 青色箭头
    arrow1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(9.8), Inches(6.9), Inches(0.4), Inches(0.25))
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = RGBColor(0x00, 0xBC, 0xD4)
    arrow1.line.fill.background()

    # OMG!
    omg_box = slide.shapes.add_textbox(Inches(10.3), Inches(6.8), Inches(1.0), Inches(0.6))
    tf_omg = omg_box.text_frame
    p_omg = tf_omg.paragraphs[0]
    p_omg.text = "OMG!"
    p_omg.font.name = "Microsoft YaHei"
    p_omg.font.size = Pt(16)
    p_omg.font.bold = True
    p_omg.font.color.rgb = RGBColor(0x00, 0xBC, 0xD4)

    # 紫色箭头
    arrow2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(11.5), Inches(6.9), Inches(0.4), Inches(0.25))
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = RGBColor(0x9C, 0x27, 0xB0)
    arrow2.line.fill.background()

    # 绿色箭头
    arrow3 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(12.1), Inches(6.9), Inches(0.4), Inches(0.25))
    arrow3.fill.solid()
    arrow3.fill.fore_color.rgb = RGBColor(0x4C, 0xAF, 0x50)
    arrow3.line.fill.background()



# ── Main ──
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

s0 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_1(s0)
prs.save(OUTPUT_PATH)
