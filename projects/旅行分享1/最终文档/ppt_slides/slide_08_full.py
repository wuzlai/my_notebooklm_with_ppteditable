from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# ── Color Palette ──
BLUE_HEADER = RGBColor(0x5B, 0x9B, 0xD5)
BLUE_DARK   = RGBColor(0x4A, 0x86, 0xC8)
CYAN        = RGBColor(0x00, 0xBC, 0xD4)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x33, 0x33, 0x33)
GRAY_TEXT   = RGBColor(0x55, 0x55, 0x55)
GRAY_BAR    = RGBColor(0xB0, 0xBE, 0xC5)
RED         = RGBColor(0xE5, 0x39, 0x35)
GREEN       = RGBColor(0x43, 0xA0, 0x47)
ORANGE      = RGBColor(0xFF, 0x98, 0x00)

ICON_BG  = RGBColor(0xE3, 0xE8, 0xED)
ICON_FG  = RGBColor(0x54, 0x6E, 0x7A)
FONT_NAME = "Microsoft YaHei"
SLIDE_WIDTH = Inches(13.333)
HEADER_H    = Inches(0.75)
SUBTITLE_Y  = Inches(0.95)


def add_header_banner(slide, title_text, bg_color=None):
    if bg_color is None:
        bg_color = BLUE_HEADER
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, HEADER_H
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = bg_color
    banner.line.fill.background()
    tf = banner.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(26)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME


def add_subtitle(slide, text, left, top, width=Inches(12), font_size=Pt(18)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.font.name = FONT_NAME
    return txBox


def add_icon_box(slide, left, top, symbol, size=Inches(0.48)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ICON_BG
    shape.line.fill.background()
    shape.adjustments[0] = 0.25
    tf = shape.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = symbol
    p.font.size = Pt(18)
    p.font.color.rgb = ICON_FG
    p.font.bold = False
    return shape


def add_bullet_item(slide, left, top, symbol, label, description,
                    width=Inches(5.5), desc_size=Pt(13)):
    add_icon_box(slide, left, top, symbol)
    text_left = left + Inches(0.65)
    txBox = slide.shapes.add_textbox(text_left, top - Inches(0.02), width, Inches(0.65))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run_label = p.add_run()
    run_label.text = label + "\uff1a"
    run_label.font.size = Pt(14)
    run_label.font.color.rgb = BLACK
    run_label.font.bold = True
    run_label.font.name = FONT_NAME
    run_desc = p.add_run()
    run_desc.text = description
    run_desc.font.size = desc_size
    run_desc.font.color.rgb = GRAY_TEXT
    run_desc.font.bold = False
    run_desc.font.name = FONT_NAME
    return txBox


def add_conclusion_box(slide, left, top, width, text, font_size=Pt(13)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = BLACK
    run.font.bold = True
    run.font.name = FONT_NAME
    return txBox


def add_table(slide, left, top, width, height, rows, cols, data,
              header_color=None, col_widths=None):
    if header_color is None:
        header_color = BLUE_HEADER
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if data[r][c] is not None else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(11)
                    run.font.name = FONT_NAME
                    if r == 0:
                        run.font.color.rgb = WHITE
                        run.font.bold = True
                    else:
                        run.font.color.rgb = BLACK
                        run.font.bold = False
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else RGBColor(0xF5, 0xF5, 0xF5)
            cell.margin_left = Pt(5)
            cell.margin_right = Pt(5)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
    return table_shape


def add_bar_chart(slide, left, top, width, height,
                  categories, values, title="", bar_colors=None):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', values)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    chart.chart_style = 2
    plot = chart.plots[0]
    plot.gap_width = 100
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CYAN
    series.has_data_labels = True
    dl = series.data_labels
    dl.font.size = Pt(13)
    dl.font.bold = True
    dl.font.color.rgb = BLACK
    dl.number_format = '0.#'
    dl.show_value = True
    dl.label_position = XL_LABEL_POSITION.OUTSIDE_END
    if bar_colors:
        for i, color in enumerate(bar_colors):
            pt = series.points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = color
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(12)
    cat_axis.tick_labels.font.name = FONT_NAME
    cat_axis.major_tick_mark = 2
    cat_axis.format.line.fill.background()
    val_axis = chart.value_axis
    val_axis.visible = False
    val_axis.major_tick_mark = 2
    val_axis.format.line.fill.background()
    val_axis.major_gridlines.format.line.fill.background()
    if title:
        chart.has_title = True
        ct = chart.chart_title.text_frame.paragraphs[0]
        ct.text = title
        ct.font.size = Pt(14)
        ct.font.bold = True
        ct.font.name = FONT_NAME
    else:
        chart.has_title = False
    return chart_frame


def add_callout_label(slide, left, top, text, bg_color=None, font_size=Pt(11)):
    if bg_color is None:
        bg_color = CYAN
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.3), Inches(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME
    return shape


def add_data_card(slide, left, top, width, height, value, label,
                  value_color=None, bg_color=None):
    if value_color is None:
        value_color = CYAN
    if bg_color is None:
        bg_color = WHITE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(3)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run1 = p1.add_run()
    run1.text = str(value)
    run1.font.size = Pt(24)
    run1.font.color.rgb = value_color
    run1.font.bold = True
    run1.font.name = FONT_NAME
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(10)
    run2.font.color.rgb = GRAY_TEXT
    run2.font.bold = False
    run2.font.name = FONT_NAME
    return shape


OUTPUT_PATH = r"C:\Users\Administrator\Desktop\notebook\my_notebooklm_with_ppteditable\projects\旅行分享1\最终文档\ppt_slides\slide_08.pptx"


# ── Slide 1 ──

def build_slide_1(slide):
    # --- Background & Decorative Elements ---
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(245, 245, 240)
    bg.line.fill.background()

    # Colorful blobs in corners
    blob1 = slide.shapes.add_shape(MSO_SHAPE.CLOUD, Inches(-1.0), Inches(-1.0), Inches(3.0), Inches(3.0))
    blob1.fill.solid(); blob1.fill.fore_color.rgb = RGBColor(200, 255, 50)
    blob1.line.fill.background()

    blob2 = slide.shapes.add_shape(MSO_SHAPE.CLOUD, Inches(11.5), Inches(-1.0), Inches(3.0), Inches(3.0))
    blob2.fill.solid(); blob2.fill.fore_color.rgb = RGBColor(255, 100, 200)
    blob2.line.fill.background()

    blob3 = slide.shapes.add_shape(MSO_SHAPE.CLOUD, Inches(-1.0), Inches(5.5), Inches(3.0), Inches(3.0))
    blob3.fill.solid(); blob3.fill.fore_color.rgb = RGBColor(255, 150, 0)
    blob3.line.fill.background()

    # --- Titles ---
    # Main Title Shadow
    title1_shadow = slide.shapes.add_textbox(Inches(0.53), Inches(0.33), Inches(5), Inches(0.8))
    p_shadow = title1_shadow.text_frame.paragraphs[0]
    p_shadow.text = "阳朔西街："
    p_shadow.font.name = "Microsoft YaHei"
    p_shadow.font.size = Pt(48)
    p_shadow.font.bold = True
    p_shadow.font.color.rgb = RGBColor(0, 0, 0)

    # Main Title
    title1 = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(5), Inches(0.8))
    p = title1.text_frame.paragraphs[0]
    p.text = "阳朔西街："
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(50, 20, 80)

    # Subtitle Shadow
    title2_shadow = slide.shapes.add_textbox(Inches(0.53), Inches(1.03), Inches(8), Inches(0.8))
    p_shadow2 = title2_shadow.text_frame.paragraphs[0]
    p_shadow2.text = "中西合璧的“迷惑行为”大赏"
    p_shadow2.font.name = "Microsoft YaHei"
    p_shadow2.font.size = Pt(36)
    p_shadow2.font.bold = True
    p_shadow2.font.color.rgb = RGBColor(0, 0, 0)

    # Subtitle
    title2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(8), Inches(0.8))
    p2 = title2.text_frame.paragraphs[0]
    p2.text = "中西合璧的“迷惑行为”大赏"
    p2.font.name = "Microsoft YaHei"
    p2.font.size = Pt(36)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(230, 80, 0)

    # Banner
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.9), Inches(6.0), Inches(0.6))
    banner.fill.solid(); banner.fill.fore_color.rgb = RGBColor(160, 80, 230)
    banner.line.fill.background()
    p = banner.text_frame.paragraphs[0]
    p.text = "这里的外国人比外地人还多"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # --- Map Section (Top Right) ---
    ff = slide.shapes.build_freeform(Inches(8.0), Inches(0.5))
    ff.add_line_segments([
        (Inches(10.0), Inches(0.2)), (Inches(12.5), Inches(1.0)),
        (Inches(12.0), Inches(2.5)), (Inches(9.5), Inches(3.5)),
        (Inches(8.0), Inches(3.0)), (Inches(4.5), Inches(3.8)),
        (Inches(4.0), Inches(3.0)), (Inches(7.5), Inches(2.0)),
        (Inches(8.0), Inches(0.5))
    ])
    map_shape = ff.convert_to_shape()
    map_shape.fill.solid(); map_shape.fill.fore_color.rgb = RGBColor(40, 40, 40)
    map_shape.line.color.rgb = RGBColor(150, 80, 220); map_shape.line.width = Pt(3)

    # Neon streets
    def add_neon_line(x1, y1, x2, y2, color):
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        line.line.color.rgb = color; line.line.width = Pt(3)

    add_neon_line(4.5, 3.2, 8.0, 2.5, RGBColor(255, 150, 0))
    add_neon_line(8.0, 2.5, 12.0, 1.8, RGBColor(255, 150, 0))
    add_neon_line(6.0, 2.9, 6.2, 2.0, RGBColor(255, 50, 150))
    add_neon_line(9.0, 2.3, 9.5, 1.0, RGBColor(150, 255, 50))
    add_neon_line(10.5, 2.0, 11.0, 3.0, RGBColor(255, 50, 150))

    # Map Text
    def add_map_text(x, y, text, rot, color):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(1.5), Inches(0.4))
        tb.rotation = rot
        p = tb.text_frame.paragraphs[0]
        p.text = text; p.font.size = Pt(10); p.font.color.rgb = color; p.font.bold = True

    add_map_text(6.5, 2.6, "阳朔西街", -10, RGBColor(255,255,255))
    add_map_text(10.0, 1.8, "阳朔西街", -5, RGBColor(255,255,255))
    add_map_text(9.2, 1.2, "阳朔西街", 70, RGBColor(255,255,255))

    # Comic Texts
    def add_comic_text(x, y, text, color, rot):
        tb_shadow = slide.shapes.add_textbox(Inches(x+0.03), Inches(y+0.03), Inches(2.0), Inches(1.0))
        tb_shadow.rotation = rot
        p_shadow = tb_shadow.text_frame.paragraphs[0]
        p_shadow.text = text; p_shadow.font.name = "Arial Black"; p_shadow.font.size = Pt(28); p_shadow.font.color.rgb = RGBColor(0,0,0)

        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(2.0), Inches(1.0))
        tb.rotation = rot
        p = tb.text_frame.paragraphs[0]
        p.text = text; p.font.name = "Arial Black"; p.font.size = Pt(28); p.font.color.rgb = color

    add_comic_text(7.0, 0.5, "Boom!", RGBColor(255, 50, 150), -15)
    add_comic_text(4.5, 4.0, "OMG!", RGBColor(150, 255, 50), 10)
    add_comic_text(11.0, 2.5, "WoW!", RGBColor(255, 50, 150), -20)

    # Decor
    bolt1 = slide.shapes.add_shape(MSO_SHAPE.LIGHTNING_BOLT, Inches(6.5), Inches(0.5), Inches(0.4), Inches(0.8))
    bolt1.fill.solid(); bolt1.fill.fore_color.rgb = RGBColor(150, 80, 220); bolt1.line.color.rgb = RGBColor(0,0,0)
    bolt2 = slide.shapes.add_shape(MSO_SHAPE.LIGHTNING_BOLT, Inches(12.0), Inches(1.0), Inches(0.4), Inches(0.8))
    bolt2.fill.solid(); bolt2.fill.fore_color.rgb = RGBColor(255, 150, 0); bolt2.line.color.rgb = RGBColor(0,0,0)
    star1 = slide.shapes.add_shape(MSO_SHAPE.5_POINT_STAR, Inches(10.0), Inches(2.8), Inches(0.4), Inches(0.4))
    star1.fill.solid(); star1.fill.fore_color.rgb = RGBColor(255, 150, 0); star1.line.color.rgb = RGBColor(0,0,0)

    # --- Section Titles ---
    def add_section_title(x, y, text):
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(1.6), Inches(0.4))
        box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0,0,0); box.line.fill.background()
        p = box.text_frame.paragraphs[0]
        p.text = text; p.font.name = "Microsoft YaHei"; p.font.size = Pt(16); p.font.color.rgb = RGBColor(255,255,255); p.font.bold = True; p.alignment = PP_ALIGN.CENTER

    add_section_title(0.5, 2.8, "猎奇景观：")
    add_section_title(6.0, 4.0, "搞笑互动：")
    add_section_title(9.5, 3.5, "避坑指南：")

    # --- Section 1: Beer Fish ---
    desc1 = slide.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(3.0), Inches(0.8))
    p1 = desc1.text_frame.paragraphs[0]
    p1.text = "卖啤酒鱼的店比"
    p1.font.name = "Microsoft YaHei"; p1.font.size = Pt(16); p1.font.bold = True
    p2 = desc1.text_frame.add_paragraph()
    p2.text = "整条街的树还多"
    p2.font.name = "Microsoft YaHei"; p2.font.size = Pt(16); p2.font.bold = True

    # Buildings
    b1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(4.5), Inches(1.5), Inches(2.5))
    b1.fill.solid(); b1.fill.fore_color.rgb = RGBColor(50,50,50)
    b2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.2), Inches(4.0), Inches(1.2), Inches(3.0))
    b2.fill.solid(); b2.fill.fore_color.rgb = RGBColor(40,40,40)
    b3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.5), Inches(4.5), Inches(1.2), Inches(2.5))
    b3.fill.solid(); b3.fill.fore_color.rgb = RGBColor(60,60,60)

    # Trees
    tree_trunk1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(6.5), Inches(0.2), Inches(0.8))
    tree_trunk1.fill.solid(); tree_trunk1.fill.fore_color.rgb = RGBColor(100, 50, 0)
    tree_leaves1 = slide.shapes.add_shape(MSO_SHAPE.CLOUD, Inches(1.0), Inches(6.0), Inches(1.2), Inches(0.8))
    tree_leaves1.fill.solid(); tree_leaves1.fill.fore_color.rgb = RGBColor(50, 150, 50)
    tree_trunk2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.0), Inches(6.8), Inches(0.2), Inches(0.6))
    tree_trunk2.fill.solid(); tree_trunk2.fill.fore_color.rgb = RGBColor(100, 50, 0)
    tree_leaves2 = slide.shapes.add_shape(MSO_SHAPE.CLOUD, Inches(2.6), Inches(6.4), Inches(1.0), Inches(0.7))
    tree_leaves2.fill.solid(); tree_leaves2.fill.fore_color.rgb = RGBColor(50, 150, 50)

    # Signs
    def add_beer_fish_sign(x, y, w, h, bg_rgb, text_rgb, rot=0):
        sign = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        sign.fill.solid(); sign.fill.fore_color.rgb = bg_rgb
        sign.line.color.rgb = RGBColor(0,0,0); sign.line.width = Pt(2); sign.rotation = rot
        tf = sign.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = "\n".join(list("啤酒鱼")) if h > w else "啤酒鱼"
        p.font.name = "Microsoft YaHei"; p.font.bold = True; p.font.size = Pt(14); p.font.color.rgb = text_rgb; p.alignment = PP_ALIGN.CENTER

    add_beer_fish_sign(0.5, 4.5, 0.8, 1.5, RGBColor(150,255,50), RGBColor(0,0,0))
    add_beer_fish_sign(1.2, 4.2, 0.6, 1.8, RGBColor(255,50,150), RGBColor(255,255,255))
    add_beer_fish_sign(0.6, 5.5, 1.2, 0.5, RGBColor(255,150,0), RGBColor(0,0,0))
    add_beer_fish_sign(0.4, 6.2, 1.5, 0.6, RGBColor(150,255,50), RGBColor(0,0,0))
    add_beer_fish_sign(2.0, 4.8, 0.7, 1.6, RGBColor(255,50,150), RGBColor(255,255,255))
    add_beer_fish_sign(2.5, 4.0, 0.8, 2.0, RGBColor(150,255,50), RGBColor(0,0,0))
    add_beer_fish_sign(2.2, 5.8, 1.0, 0.5, RGBColor(255,150,0), RGBColor(0,0,0))
    add_beer_fish_sign(2.8, 5.2, 0.6, 1.4, RGBColor(255,50,150), RGBColor(255,255,255))
    add_beer_fish_sign(3.5, 4.5, 1.2, 0.6, RGBColor(150,255,50), RGBColor(0,0,0))
    add_beer_fish_sign(3.4, 5.2, 1.4, 0.7, RGBColor(150,255,50), RGBColor(0,0,0))
    add_beer_fish_sign(3.6, 6.0, 1.2, 0.6, RGBColor(255,150,0), RGBColor(0,0,0))
    add_beer_fish_sign(1.0, 5.0, 1.0, 0.5, RGBColor(255,255,50), RGBColor(0,0,0), rot=-10)
    add_beer_fish_sign(2.8, 4.5, 0.5, 1.2, RGBColor(150,80,220), RGBColor(255,255,255), rot=5)

    # --- Section 2: Interaction ---
    # Snail -> Person icon
    snail_body = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.7), Inches(4.0), Inches(0.4), Inches(0.3))
    snail_body.fill.solid(); snail_body.fill.fore_color.rgb = RGBColor(150, 100, 50)
    arrow1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.2), Inches(4.1), Inches(0.3), Inches(0.15))
    arrow1.fill.solid(); arrow1.fill.fore_color.rgb = RGBColor(255, 150, 0)
    head_icon = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.6), Inches(4.0), Inches(0.3), Inches(0.3))
    head_icon.fill.solid(); head_icon.fill.fore_color.rgb = RGBColor(255, 200, 150)

    # People
    head1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.5), Inches(5.5), Inches(0.6), Inches(0.6))
    head1.fill.solid(); head1.fill.fore_color.rgb = RGBColor(255, 220, 200)
    body1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.3), Inches(6.1), Inches(1.0), Inches(0.8))
    body1.fill.solid(); body1.fill.fore_color.rgb = RGBColor(150, 80, 220)

    head2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.0), Inches(5.4), Inches(0.6), Inches(0.6))
    head2.fill.solid(); head2.fill.fore_color.rgb = RGBColor(255, 220, 200)
    body2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.8), Inches(6.0), Inches(1.0), Inches(0.9))
    body2.fill.solid(); body2.fill.fore_color.rgb = RGBColor(100, 220, 50)

    # Speech Bubbles
    bubble1 = slide.shapes.add_shape(MSO_SHAPE.OVAL_CALLOUT, Inches(5.3), Inches(4.5), Inches(2.2), Inches(0.9))
    bubble1.fill.solid(); bubble1.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bubble1.line.color.rgb = RGBColor(0,0,0); bubble1.line.width = Pt(2)
    p1 = bubble1.text_frame.paragraphs[0]
    p1.text = "How to eat... um...\nTianluo (snails)?"
    p1.font.size = Pt(11); p1.font.color.rgb = RGBColor(0,0,0)

    bubble2 = slide.shapes.add_shape(MSO_SHAPE.OVAL_CALLOUT, Inches(7.7), Inches(4.2), Inches(2.0), Inches(0.8))
    bubble2.fill.solid(); bubble2.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bubble2.line.color.rgb = RGBColor(0,0,0); bubble2.line.width = Pt(2)
    p2 = bubble2.text_frame.paragraphs[0]
    p2.text = "Haha, use a\ntoothpick, mate!"
    p2.font.size = Pt(11); p2.font.color.rgb = RGBColor(0,0,0)

    # Bottom Text Box
    interact_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.0), Inches(6.8), Inches(3.2), Inches(0.6))
    interact_box.fill.solid(); interact_box.fill.fore_color.rgb = RGBColor(0,0,0); interact_box.line.fill.background()
    p = interact_box.text_frame.paragraphs[0]
    p.text = "尝试用蹩脚英语和老外聊怎么吃田螺"
    p.font.name = "Microsoft YaHei"; p.font.size = Pt(12); p.font.color.rgb = RGBColor(255,255,255); p.font.bold = True; p.alignment = PP_ALIGN.CENTER

    # --- Section 3: Tourist Trap ---
    # Shop -> Shop icon
    shop1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(11.2), Inches(3.5), Inches(0.4), Inches(0.4))
    shop1.fill.solid(); shop1.fill.fore_color.rgb = RGBColor(200, 100, 200)
    arrow2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(11.7), Inches(3.6), Inches(0.3), Inches(0.15))
    arrow2.fill.solid(); arrow2.fill.fore_color.rgb = RGBColor(150, 255, 50)
    shop2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(12.1), Inches(3.5), Inches(0.4), Inches(0.4))
    shop2.fill.solid(); shop2.fill.fore_color.rgb = RGBColor(200, 100, 200)

    # Shops
    frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.5), Inches(4.2), Inches(3.5), Inches(1.8))
    frame.fill.solid(); frame.fill.fore_color.rgb = RGBColor(120, 80, 180)
    frame.line.color.rgb = RGBColor(0,0,0); frame.line.width = Pt(2)

    stall_w = 3.5 / 4
    for i in range(4):
        stall = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.5 + i*stall_w + 0.1), Inches(4.8), Inches(stall_w - 0.2), Inches(1.2))
        stall.fill.solid(); stall.fill.fore_color.rgb = RGBColor(60, 40, 90); stall.line.color.rgb = RGBColor(0,0,0)
        
        counter = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.5 + i*stall_w + 0.1), Inches(5.5), Inches(stall_w - 0.2), Inches(0.5))
        counter.fill.solid(); counter.fill.fore_color.rgb = RGBColor(150, 100, 50); counter.line.color.rgb = RGBColor(0,0,0)

        sign = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.5 + i*stall_w + 0.05), Inches(4.3), Inches(stall_w - 0.1), Inches(0.4))
        sign.fill.solid(); sign.fill.fore_color.rgb = RGBColor(255, 100, 100) if i%2==0 else RGBColor(100, 255, 100)
        sign.line.color.rgb = RGBColor(0,0,0)
        tf = sign.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = "网红店"; p.font.size = Pt(10); p.font.bold = True; p.alignment = PP_ALIGN.CENTER

    # Red Cross
    cross1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10.0), Inches(4.5), Inches(12.5), Inches(5.8))
    cross1.line.color.rgb = RGBColor(220, 40, 40); cross1.line.width = Pt(12)
    cross2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(12.5), Inches(4.5), Inches(10.0), Inches(5.8))
    cross2.line.color.rgb = RGBColor(220, 40, 40); cross2.line.width = Pt(12)

    # Bottom Text
    trap_text = slide.shapes.add_textbox(Inches(9.5), Inches(6.2), Inches(3.5), Inches(0.8))
    p = trap_text.text_frame.paragraphs[0]
    p.text = "西街的“网红”其实都是同一个模版"
    p.font.name = "Microsoft YaHei"; p.font.size = Pt(14); p.font.bold = True

    # --- Footer ---
    mt1 = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(11.0), Inches(6.8), Inches(0.5), Inches(0.5))
    mt1.fill.solid(); mt1.fill.fore_color.rgb = RGBColor(200, 200, 200); mt1.line.color.rgb = RGBColor(0,0,0)
    mt2 = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(11.3), Inches(6.9), Inches(0.4), Inches(0.4))
    mt2.fill.solid(); mt2.fill.fore_color.rgb = RGBColor(180, 180, 180); mt2.line.color.rgb = RGBColor(0,0,0)

    footer = slide.shapes.add_textbox(Inches(12.0), Inches(6.8), Inches(1.0), Inches(0.5))
    p = footer.text_frame.paragraphs[0]
    p.text = "08/11"; p.font.name = "Arial"; p.font.size = Pt(20); p.font.bold = True



# ── Main ──
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

s0 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_1(s0)
prs.save(OUTPUT_PATH)
