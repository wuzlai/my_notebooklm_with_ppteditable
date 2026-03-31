from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# ── Color Palette ──
BLUE_HEADER = RGBColor(0x5B, 0x9B, 0xD5)
BLUE_DARK   = RGBColor(0x4A, 0x86, 0xC8)
CYAN        = RGBColor(0x00, 0xBC, 0xD4)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x33, 0x33, 0x33)
GRAY_TEXT   = RGBColor(0x55, 0x55, 0x55)
GRAY_BAR    = RGBColor(0xB0, 0xBE, 0xC5)
RED         = RGBColor(0xE5, 0x39, 0x35)
GREEN       = RGBColor(0x43, 0xA0, 0x47)
ORANGE      = RGBColor(0xFF, 0x98, 0x00)

ICON_BG  = RGBColor(0xE3, 0xE8, 0xED)
ICON_FG  = RGBColor(0x54, 0x6E, 0x7A)
FONT_NAME = "Microsoft YaHei"
SLIDE_WIDTH = Inches(13.333)
HEADER_H    = Inches(0.75)
SUBTITLE_Y  = Inches(0.95)


def add_header_banner(slide, title_text, bg_color=None):
    if bg_color is None:
        bg_color = BLUE_HEADER
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, HEADER_H
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = bg_color
    banner.line.fill.background()
    tf = banner.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(26)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME


def add_subtitle(slide, text, left, top, width=Inches(12), font_size=Pt(18)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.font.name = FONT_NAME
    return txBox


def add_icon_box(slide, left, top, symbol, size=Inches(0.48)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ICON_BG
    shape.line.fill.background()
    shape.adjustments[0] = 0.25
    tf = shape.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = symbol
    p.font.size = Pt(18)
    p.font.color.rgb = ICON_FG
    p.font.bold = False
    return shape


def add_bullet_item(slide, left, top, symbol, label, description,
                    width=Inches(5.5), desc_size=Pt(13)):
    add_icon_box(slide, left, top, symbol)
    text_left = left + Inches(0.65)
    txBox = slide.shapes.add_textbox(text_left, top - Inches(0.02), width, Inches(0.65))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run_label = p.add_run()
    run_label.text = label + "\uff1a"
    run_label.font.size = Pt(14)
    run_label.font.color.rgb = BLACK
    run_label.font.bold = True
    run_label.font.name = FONT_NAME
    run_desc = p.add_run()
    run_desc.text = description
    run_desc.font.size = desc_size
    run_desc.font.color.rgb = GRAY_TEXT
    run_desc.font.bold = False
    run_desc.font.name = FONT_NAME
    return txBox


def add_conclusion_box(slide, left, top, width, text, font_size=Pt(13)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = BLACK
    run.font.bold = True
    run.font.name = FONT_NAME
    return txBox


def add_table(slide, left, top, width, height, rows, cols, data,
              header_color=None, col_widths=None):
    if header_color is None:
        header_color = BLUE_HEADER
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if data[r][c] is not None else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(11)
                    run.font.name = FONT_NAME
                    if r == 0:
                        run.font.color.rgb = WHITE
                        run.font.bold = True
                    else:
                        run.font.color.rgb = BLACK
                        run.font.bold = False
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else RGBColor(0xF5, 0xF5, 0xF5)
            cell.margin_left = Pt(5)
            cell.margin_right = Pt(5)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
    return table_shape


def add_bar_chart(slide, left, top, width, height,
                  categories, values, title="", bar_colors=None):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', values)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    chart.chart_style = 2
    plot = chart.plots[0]
    plot.gap_width = 100
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CYAN
    series.has_data_labels = True
    dl = series.data_labels
    dl.font.size = Pt(13)
    dl.font.bold = True
    dl.font.color.rgb = BLACK
    dl.number_format = '0.#'
    dl.show_value = True
    dl.label_position = XL_LABEL_POSITION.OUTSIDE_END
    if bar_colors:
        for i, color in enumerate(bar_colors):
            pt = series.points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = color
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(12)
    cat_axis.tick_labels.font.name = FONT_NAME
    cat_axis.major_tick_mark = 2
    cat_axis.format.line.fill.background()
    val_axis = chart.value_axis
    val_axis.visible = False
    val_axis.major_tick_mark = 2
    val_axis.format.line.fill.background()
    val_axis.major_gridlines.format.line.fill.background()
    if title:
        chart.has_title = True
        ct = chart.chart_title.text_frame.paragraphs[0]
        ct.text = title
        ct.font.size = Pt(14)
        ct.font.bold = True
        ct.font.name = FONT_NAME
    else:
        chart.has_title = False
    return chart_frame


def add_callout_label(slide, left, top, text, bg_color=None, font_size=Pt(11)):
    if bg_color is None:
        bg_color = CYAN
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.3), Inches(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME
    return shape


def add_data_card(slide, left, top, width, height, value, label,
                  value_color=None, bg_color=None):
    if value_color is None:
        value_color = CYAN
    if bg_color is None:
        bg_color = WHITE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(3)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run1 = p1.add_run()
    run1.text = str(value)
    run1.font.size = Pt(24)
    run1.font.color.rgb = value_color
    run1.font.bold = True
    run1.font.name = FONT_NAME
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(10)
    run2.font.color.rgb = GRAY_TEXT
    run2.font.bold = False
    run2.font.name = FONT_NAME
    return shape


OUTPUT_PATH = r"C:\Users\Administrator\Desktop\notebook\my_notebooklm_with_ppteditable\projects\旅行分享1\最终文档\ppt_slides\slide_01.pptx"


# ── Slide 1 ──

def build_slide_1(slide):
    # --- Colors ---
    BG_COLOR = RGBColor(245, 242, 235)
    TEXT_BLACK = RGBColor(20, 20, 20)
    ORANGE = RGBColor(244, 122, 32)
    GREEN = RGBColor(0, 168, 112)
    DARK_BLUE = RGBColor(20, 40, 80)
    LIGHT_TEAL = RGBColor(160, 220, 200)
    WHITE = RGBColor(255, 255, 255)

    # Set background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

    # --- Edge Decorations (Torn Paper Effect) ---
    # Top Left
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(-1.5), Inches(-1.5), Inches(3), Inches(3))
    shape.rotation = 45
    shape.fill.solid()
    shape.fill.fore_color.rgb = GREEN
    shape.line.fill.background()

    # Top Right
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(11.5), Inches(-1), Inches(3), Inches(2))
    shape.rotation = -20
    shape.fill.solid()
    shape.fill.fore_color.rgb = ORANGE
    shape.line.fill.background()
    
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(12), Inches(0), Inches(2), Inches(3))
    shape.rotation = 10
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    # Bottom Left
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(-1), Inches(6.5), Inches(3), Inches(2))
    shape.rotation = 15
    shape.fill.solid()
    shape.fill.fore_color.rgb = ORANGE
    shape.line.fill.background()

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(-0.5), Inches(7), Inches(3), Inches(2))
    shape.rotation = -10
    shape.fill.solid()
    shape.fill.fore_color.rgb = TEXT_BLACK
    shape.line.fill.background()

    # --- Main Title ---
    # Line 1
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(8), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "桂林山水“甲”天下，"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.name = "Microsoft YaHei"
    p.font.color.rgb = TEXT_BLACK

    # Highlight under Line 2
    highlight = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.6), Inches(6.2), Inches(0.25))
    highlight.fill.solid()
    highlight.fill.fore_color.rgb = ORANGE
    highlight.line.fill.background()

    # Line 2
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(8), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "我的脑洞“假”不了"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.name = "Microsoft YaHei"
    p.font.color.rgb = TEXT_BLACK

    # --- Subtitle ---
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(8), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "一个猎奇博主的桂林“真香”探险报告"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.name = "Microsoft YaHei"
    p.font.color.rgb = TEXT_BLACK

    # --- 20 RMB Graphic ---
    rmb_left = Inches(8.0)
    rmb_top = Inches(0.6)
    
    # Base
    rmb_base = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, rmb_left, rmb_top, Inches(4.8), Inches(2.5))
    rmb_base.fill.solid()
    rmb_base.fill.fore_color.rgb = WHITE
    rmb_base.line.color.rgb = DARK_BLUE
    rmb_base.line.width = Pt(2)

    # Inner fill
    rmb_inner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, rmb_left + Inches(0.1), rmb_top + Inches(0.1), Inches(4.6), Inches(2.3))
    rmb_inner.fill.solid()
    rmb_inner.fill.fore_color.rgb = LIGHT_TEAL
    rmb_inner.line.fill.background()

    # Sun
    sun = slide.shapes.add_shape(MSO_SHAPE.OVAL, rmb_left + Inches(2.0), rmb_top + Inches(0.3), Inches(0.6), Inches(0.6))
    sun.fill.solid()
    sun.fill.fore_color.rgb = WHITE
    sun.line.fill.background()

    # Mountains (Triangles)
    m1 = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, rmb_left + Inches(0.5), rmb_top + Inches(0.8), Inches(1.0), Inches(1.2))
    m1.fill.solid()
    m1.fill.fore_color.rgb = DARK_BLUE
    m1.line.fill.background()

    m2 = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, rmb_left + Inches(1.2), rmb_top + Inches(0.6), Inches(1.2), Inches(1.4))
    m2.fill.solid()
    m2.fill.fore_color.rgb = GREEN
    m2.line.fill.background()

    m3 = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, rmb_left + Inches(2.5), rmb_top + Inches(0.9), Inches(1.0), Inches(1.1))
    m3.fill.solid()
    m3.fill.fore_color.rgb = GREEN
    m3.line.fill.background()

    # River (Rectangle at bottom)
    river = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, rmb_left + Inches(0.1), rmb_top + Inches(1.8), Inches(4.6), Inches(0.6))
    river.fill.solid()
    river.fill.fore_color.rgb = DARK_BLUE
    river.line.fill.background()

    # Text "20" Top Left
    txBox = slide.shapes.add_textbox(rmb_left + Inches(0.2), rmb_top + Inches(0.1), Inches(1), Inches(0.5))
    p = txBox.text_frame.paragraphs[0]
    p.text = "20"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Text "20 RMB" Bottom Left
    txBox = slide.shapes.add_textbox(rmb_left + Inches(0.2), rmb_top + Inches(1.8), Inches(1.5), Inches(0.5))
    p = txBox.text_frame.paragraphs[0]
    p.text = "20 RMB"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Text "中国人民银行" Top Right
    txBox = slide.shapes.add_textbox(rmb_left + Inches(2.5), rmb_top + Inches(0.1), Inches(2), Inches(0.4))
    p = txBox.text_frame.paragraphs[0]
    p.text = "中国人民银行"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # --- Labels & Character (Bottom Left) ---
    # Character Placeholder
    char_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.0), Inches(4.5), Inches(1.5), Inches(2.5))
    char_bg.fill.solid()
    char_bg.fill.fore_color.rgb = ORANGE
    char_bg.line.color.rgb = TEXT_BLACK
    char_bg.line.width = Pt(2)
    p = char_bg.text_frame.paragraphs[0]
    p.text = "猎奇\n博主"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    def add_label(text, left, top, width, height, rotation, border_color):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.rotation = rotation
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = border_color
        shape.line.width = Pt(2.5)
        
        tf = shape.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.name = "Microsoft YaHei"
        p.font.color.rgb = TEXT_BLACK
        p.alignment = PP_ALIGN.CENTER
        tf.margin_top = Pt(3)
        tf.margin_bottom = Pt(3)
        tf.margin_left = Pt(5)
        tf.margin_right = Pt(5)

    add_label("奇葩角落", Inches(1.2), Inches(4.0), Inches(1.4), Inches(0.4), -15, GREEN)
    add_label("真香！！", Inches(1.2), Inches(5.0), Inches(1.4), Inches(0.4), -5, ORANGE)
    add_label("不错的青", Inches(1.5), Inches(6.0), Inches(1.4), Inches(0.4), 10, DARK_BLUE)
    
    add_label("拒绝中老年团", Inches(4.8), Inches(4.2), Inches(2.0), Inches(0.4), -15, ORANGE)
    add_label("冷知识", Inches(5.2), Inches(5.2), Inches(1.2), Inches(0.4), 5, GREEN)
    add_label("美景背后", Inches(5.0), Inches(6.0), Inches(1.4), Inches(0.4), 15, DARK_BLUE)

    # Add some arrows/lines pointing to center
    line1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(2.5), Inches(4.2), Inches(2.9), Inches(4.8))
    line1.line.color.rgb = TEXT_BLACK
    line1.line.width = Pt(2)
    
    line2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(4.8), Inches(4.6), Inches(4.5), Inches(5.0))
    line2.line.color.rgb = TEXT_BLACK
    line2.line.width = Pt(2)

    # --- Right List ---
    def add_list_item(text, left, top, bg_color, icon_border_color, icon_type):
        # Background Parallelogram
        bg = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, left + Inches(0.4), top, Inches(4.5), Inches(0.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = bg_color
        bg.line.fill.background()

        # Text
        txBox = slide.shapes.add_textbox(left + Inches(0.8), top - Inches(0.05), Inches(4.0), Inches(0.6))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.name = "Microsoft YaHei"
        p.font.color.rgb = TEXT_BLACK

        # Icon Circle
        icon_bg = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top - Inches(0.1), Inches(0.7), Inches(0.7))
        icon_bg.fill.solid()
        icon_bg.fill.fore_color.rgb = WHITE
        icon_bg.line.color.rgb = icon_border_color
        icon_bg.line.width = Pt(2.5)

        # Simple icon drawing inside the circle
        if icon_type == 'bus':
            bus = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.15), top + Inches(0.1), Inches(0.4), Inches(0.3))
            bus.fill.solid()
            bus.fill.fore_color.rgb = DARK_BLUE
            bus.line.fill.background()
            cross = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left + Inches(0.1), top, left + Inches(0.6), top + Inches(0.5))
            cross.line.color.rgb = ORANGE
            cross.line.width = Pt(3)
        elif icon_type == 'mag':
            mag_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.15), top + Inches(0.05), Inches(0.3), Inches(0.3))
            mag_circle.fill.background()
            mag_circle.line.color.rgb = GREEN
            mag_circle.line.width = Pt(2)
            handle = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left + Inches(0.4), top + Inches(0.3), left + Inches(0.55), top + Inches(0.45))
            handle.line.color.rgb = ORANGE
            handle.line.width = Pt(3)
        elif icon_type == 'person':
            head = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.25), top + Inches(0.05), Inches(0.2), Inches(0.2))
            head.fill.solid()
            head.fill.fore_color.rgb = ORANGE
            head.line.fill.background()
            body = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.15), top + Inches(0.3), Inches(0.4), Inches(0.2))
            body.fill.solid()
            body.fill.fore_color.rgb = GREEN
            body.line.fill.background()

    add_list_item("拒绝传统中老年旅行团画风", Inches(7.8), Inches(4.0), GREEN, ORANGE, 'bus')
    add_list_item("深度挖掘桂林不为人知的“奇葩”角落", Inches(7.8), Inches(5.0), GREEN, ORANGE, 'mag')
    add_list_item("搞笑博主的生存视角：\n美景背后的冷知识", Inches(7.8), Inches(6.0), ORANGE, GREEN, 'person')

    # --- Page Number ---
    # Green background circle
    circle_bg = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(12.1), Inches(6.8), Inches(0.4), Inches(0.4))
    circle_bg.fill.solid()
    circle_bg.fill.fore_color.rgb = GREEN
    circle_bg.line.fill.background()

    # Orange Circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(12.2), Inches(6.8), Inches(0.4), Inches(0.4))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ORANGE
    circle.line.fill.background()

    # Text
    txBox = slide.shapes.add_textbox(Inches(12.4), Inches(6.7), Inches(0.6), Inches(0.6))
    p = txBox.text_frame.paragraphs[0]
    p.text = "01"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.name = "Microsoft YaHei"
    p.font.color.rgb = TEXT_BLACK



# ── Main ──
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

s0 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_1(s0)
prs.save(OUTPUT_PATH)
