from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# ── Color Palette ──
BLUE_HEADER = RGBColor(0x5B, 0x9B, 0xD5)
BLUE_DARK   = RGBColor(0x4A, 0x86, 0xC8)
CYAN        = RGBColor(0x00, 0xBC, 0xD4)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x33, 0x33, 0x33)
GRAY_TEXT   = RGBColor(0x55, 0x55, 0x55)
GRAY_BAR    = RGBColor(0xB0, 0xBE, 0xC5)
RED         = RGBColor(0xE5, 0x39, 0x35)
GREEN       = RGBColor(0x43, 0xA0, 0x47)
ORANGE      = RGBColor(0xFF, 0x98, 0x00)

ICON_BG  = RGBColor(0xE3, 0xE8, 0xED)
ICON_FG  = RGBColor(0x54, 0x6E, 0x7A)
FONT_NAME = "Microsoft YaHei"
SLIDE_WIDTH = Inches(13.333)
HEADER_H    = Inches(0.75)
SUBTITLE_Y  = Inches(0.95)


def add_header_banner(slide, title_text, bg_color=None):
    if bg_color is None:
        bg_color = BLUE_HEADER
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, HEADER_H
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = bg_color
    banner.line.fill.background()
    tf = banner.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(26)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME


def add_subtitle(slide, text, left, top, width=Inches(12), font_size=Pt(18)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.font.name = FONT_NAME
    return txBox


def add_icon_box(slide, left, top, symbol, size=Inches(0.48)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ICON_BG
    shape.line.fill.background()
    shape.adjustments[0] = 0.25
    tf = shape.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = symbol
    p.font.size = Pt(18)
    p.font.color.rgb = ICON_FG
    p.font.bold = False
    return shape


def add_bullet_item(slide, left, top, symbol, label, description,
                    width=Inches(5.5), desc_size=Pt(13)):
    add_icon_box(slide, left, top, symbol)
    text_left = left + Inches(0.65)
    txBox = slide.shapes.add_textbox(text_left, top - Inches(0.02), width, Inches(0.65))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run_label = p.add_run()
    run_label.text = label + "\uff1a"
    run_label.font.size = Pt(14)
    run_label.font.color.rgb = BLACK
    run_label.font.bold = True
    run_label.font.name = FONT_NAME
    run_desc = p.add_run()
    run_desc.text = description
    run_desc.font.size = desc_size
    run_desc.font.color.rgb = GRAY_TEXT
    run_desc.font.bold = False
    run_desc.font.name = FONT_NAME
    return txBox


def add_conclusion_box(slide, left, top, width, text, font_size=Pt(13)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = BLACK
    run.font.bold = True
    run.font.name = FONT_NAME
    return txBox


def add_table(slide, left, top, width, height, rows, cols, data,
              header_color=None, col_widths=None):
    if header_color is None:
        header_color = BLUE_HEADER
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if data[r][c] is not None else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(11)
                    run.font.name = FONT_NAME
                    if r == 0:
                        run.font.color.rgb = WHITE
                        run.font.bold = True
                    else:
                        run.font.color.rgb = BLACK
                        run.font.bold = False
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else RGBColor(0xF5, 0xF5, 0xF5)
            cell.margin_left = Pt(5)
            cell.margin_right = Pt(5)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
    return table_shape


def add_bar_chart(slide, left, top, width, height,
                  categories, values, title="", bar_colors=None):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', values)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    chart.chart_style = 2
    plot = chart.plots[0]
    plot.gap_width = 100
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CYAN
    series.has_data_labels = True
    dl = series.data_labels
    dl.font.size = Pt(13)
    dl.font.bold = True
    dl.font.color.rgb = BLACK
    dl.number_format = '0.#'
    dl.show_value = True
    dl.label_position = XL_LABEL_POSITION.OUTSIDE_END
    if bar_colors:
        for i, color in enumerate(bar_colors):
            pt = series.points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = color
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(12)
    cat_axis.tick_labels.font.name = FONT_NAME
    cat_axis.major_tick_mark = 2
    cat_axis.format.line.fill.background()
    val_axis = chart.value_axis
    val_axis.visible = False
    val_axis.major_tick_mark = 2
    val_axis.format.line.fill.background()
    val_axis.major_gridlines.format.line.fill.background()
    if title:
        chart.has_title = True
        ct = chart.chart_title.text_frame.paragraphs[0]
        ct.text = title
        ct.font.size = Pt(14)
        ct.font.bold = True
        ct.font.name = FONT_NAME
    else:
        chart.has_title = False
    return chart_frame


def add_callout_label(slide, left, top, text, bg_color=None, font_size=Pt(11)):
    if bg_color is None:
        bg_color = CYAN
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.3), Inches(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME
    return shape


def add_data_card(slide, left, top, width, height, value, label,
                  value_color=None, bg_color=None):
    if value_color is None:
        value_color = CYAN
    if bg_color is None:
        bg_color = WHITE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(3)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run1 = p1.add_run()
    run1.text = str(value)
    run1.font.size = Pt(24)
    run1.font.color.rgb = value_color
    run1.font.bold = True
    run1.font.name = FONT_NAME
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(10)
    run2.font.color.rgb = GRAY_TEXT
    run2.font.bold = False
    run2.font.name = FONT_NAME
    return shape


OUTPUT_PATH = r"C:\Users\Administrator\Desktop\notebook\my_notebooklm_with_ppteditable\projects\旅行分享1\最终文档\ppt_slides\slide_03.pptx"


# ── Slide 1 ──

def build_slide_1(slide):
    # Colors
    BLACK = RGBColor(0x00, 0x00, 0x00)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    GREEN_LBL = RGBColor(0x8A, 0xE0, 0xA1)
    ORANGE_LBL = RGBColor(0xFF, 0xB3, 0x47)
    GREEN_BORDER = RGBColor(0x4C, 0xAF, 0x50)
    ORANGE_BORDER = RGBColor(0xFF, 0x57, 0x22)
    PURPLE = RGBColor(0x9C, 0x27, 0xB0)
    GRAY_TEXT = RGBColor(0x75, 0x75, 0x75)

    # 1. Title
    title_box = slide.shapes.add_textbox(Inches(1.5), Inches(0.4), Inches(10.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    
    run1 = p.add_run()
    run1.text = "20元人民币打卡： "
    run1.font.size = Pt(36)
    run1.font.bold = True
    run1.font.name = "Microsoft YaHei"
    
    run2 = p.add_run()
    run2.text = "买家秀"
    run2.font.size = Pt(36)
    run2.font.bold = True
    run2.font.name = "Microsoft YaHei"
    
    run3 = p.add_run()
    run3.text = " vs "
    run3.font.size = Pt(36)
    run3.font.bold = True
    run3.font.name = "Microsoft YaHei"
    
    run4 = p.add_run()
    run4.text = "卖家秀"
    run4.font.size = Pt(36)
    run4.font.bold = True
    run4.font.name = "Microsoft YaHei"

    # Purple underline for "卖家秀" (approximate position)
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(9.3), Inches(1.15), Inches(10.8), Inches(1.15))
    line.line.color.rgb = PURPLE
    line.line.width = Pt(4)

    # 2. Subtitle
    sub_box = slide.shapes.add_textbox(Inches(3), Inches(1.1), Inches(7.333), Inches(0.6))
    tf_sub = sub_box.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.alignment = PP_ALIGN.CENTER
    run_sub = p_sub.add_run()
    run_sub.text = "理想很丰满，现实很骨感"
    run_sub.font.size = Pt(22)
    run_sub.font.bold = True
    run_sub.font.name = "Microsoft YaHei"

    # 3. Center Divider (Lightning Bolt)
    lightning = slide.shapes.add_shape(MSO_SHAPE.LIGHTNING_BOLT, Inches(6.3), Inches(1.8), Inches(0.8), Inches(5.5))
    lightning.fill.solid()
    lightning.fill.fore_color.rgb = ORANGE_BORDER
    lightning.line.fill.background()

    # --- LEFT SIDE (IDEAL) ---

    # Label: 理想 (Ideal)
    lbl_ideal = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.3), Inches(1.8), Inches(2.0), Inches(0.5))
    lbl_ideal.fill.solid()
    lbl_ideal.fill.fore_color.rgb = GREEN_LBL
    lbl_ideal.line.fill.background()
    tf_ideal = lbl_ideal.text_frame
    tf_ideal.text = "理想 (Ideal)"
    tf_ideal.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf_ideal.paragraphs[0].font.size = Pt(18)
    tf_ideal.paragraphs[0].font.bold = True
    tf_ideal.paragraphs[0].font.color.rgb = BLACK
    tf_ideal.paragraphs[0].font.name = "Microsoft YaHei"

    # Image Placeholder: 20 RMB
    img_ideal = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.4), Inches(5.0), Inches(2.6))
    img_ideal.fill.solid()
    img_ideal.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
    img_ideal.line.fill.solid()
    img_ideal.line.fore_color.rgb = GREEN_BORDER
    img_ideal.line.width = Pt(3)
    tf_img1 = img_ideal.text_frame
    tf_img1.text = "[20元人民币风景图]"
    tf_img1.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf_img1.paragraphs[0].font.color.rgb = RGBColor(0x9E, 0x9E, 0x9E)

    # Photographer Placeholder
    photo_guy = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.8), Inches(5.2), Inches(1.5), Inches(1.8))
    photo_guy.fill.solid()
    photo_guy.fill.fore_color.rgb = RGBColor(0xBB, 0xDE, 0xFB)
    photo_guy.line.fill.solid()
    photo_guy.line.fore_color.rgb = BLACK
    photo_guy.line.width = Pt(2)
    tf_guy = photo_guy.text_frame
    tf_guy.text = "摄影师\n(插画)"
    tf_guy.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf_guy.paragraphs[0].font.color.rgb = BLACK

    # Speech Bubble: PERFECT SHOT!
    bubble1 = slide.shapes.add_shape(MSO_SHAPE.OVAL_CALLOUT, Inches(3.1), Inches(4.8), Inches(1.5), Inches(0.8))
    bubble1.fill.solid()
    bubble1.fill.fore_color.rgb = WHITE
    bubble1.line.fill.solid()
    bubble1.line.fore_color.rgb = BLACK
    bubble1.line.width = Pt(2)
    tf_b1 = bubble1.text_frame
    tf_b1.text = "PERFECT\nSHOT!"
    tf_b1.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf_b1.paragraphs[0].font.size = Pt(11)
    tf_b1.paragraphs[0].font.bold = True
    tf_b1.paragraphs[0].font.color.rgb = BLACK

    # Conclusion Box
    box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(6.4), Inches(2.8), Inches(0.8))
    box1.fill.solid()
    box1.fill.fore_color.rgb = WHITE
    box1.line.fill.solid()
    box1.line.fore_color.rgb = BLACK
    box1.line.width = Pt(2)
    tf_box1 = box1.text_frame
    tf_box1.text = "结论：找准角度，你就是\n人民币上的那个男人/女人"
    tf_box1.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf_box1.paragraphs[0].font.size = Pt(11)
    tf_box1.paragraphs[0].font.bold = True
    tf_box1.paragraphs[0].font.color.rgb = BLACK
    tf_box1.paragraphs[0].font.name = "Microsoft YaHei"

    # Arrow to Conclusion Box
    arrow1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.1), Inches(6.7), Inches(0.4), Inches(0.2))
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = ORANGE_BORDER
    arrow1.line.fill.background()

    # --- RIGHT SIDE (REALITY) ---

    # Label: 现实 (Reality)
    lbl_reality = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.9), Inches(1.8), Inches(2.2), Inches(0.5))
    lbl_reality.fill.solid()
    lbl_reality.fill.fore_color.rgb = ORANGE_LBL
    lbl_reality.line.fill.background()
    tf_reality = lbl_reality.text_frame
    tf_reality.text = "现实 (Reality)"
    tf_reality.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf_reality.paragraphs[0].font.size = Pt(18)
    tf_reality.paragraphs[0].font.bold = True
    tf_reality.paragraphs[0].font.color.rgb = BLACK
    tf_reality.paragraphs[0].font.name = "Microsoft YaHei"

    # Image Placeholder: Crowd
    img_reality = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.5), Inches(2.4), Inches(5.0), Inches(3.0))
    img_reality.fill.solid()
    img_reality.fill.fore_color.rgb = RGBColor(0xFF, 0xE0, 0xB2)
    img_reality.line.fill.solid()
    img_reality.line.fore_color.rgb = ORANGE_BORDER
    img_reality.line.width = Pt(3)
    tf_img2 = img_reality.text_frame
    tf_img2.text = "[拥挤的人群举着20元拍照]"
    tf_img2.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf_img2.paragraphs[0].font.color.rgb = RGBColor(0x9E, 0x9E, 0x9E)

    # Stressed Guy Placeholder
    stress_guy = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.5), Inches(3.5), Inches(1.5), Inches(1.8))
    stress_guy.fill.solid()
    stress_guy.fill.fore_color.rgb = RGBColor(0xFF, 0xCD, 0xD2)
    stress_guy.line.fill.solid()
    stress_guy.line.fore_color.rgb = BLACK
    stress_guy.line.width = Pt(2)
    tf_guy2 = stress_guy.text_frame
    tf_guy2.text = "崩溃游客\n(插画)"
    tf_guy2.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf_guy2.paragraphs[0].font.color.rgb = BLACK

    # Speech Bubble: OMG! TOO MANY PEOPLE!
    bubble2 = slide.shapes.add_shape(MSO_SHAPE.CLOUD_CALLOUT, Inches(11.0), Inches(2.2), Inches(1.8), Inches(1.2))
    bubble2.fill.solid()
    bubble2.fill.fore_color.rgb = WHITE
    bubble2.line.fill.solid()
    bubble2.line.fore_color.rgb = BLACK
    bubble2.line.width = Pt(2)
    tf_b2 = bubble2.text_frame
    tf_b2.text = "OMG!\nTOO MANY\nPEOPLE!"
    tf_b2.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf_b2.paragraphs[0].font.size = Pt(11)
    tf_b2.paragraphs[0].font.bold = True
    tf_b2.paragraphs[0].font.color.rgb = BLACK

    # Reality Text Box
    box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.2), Inches(4.8), Inches(2.5), Inches(1.0))
    box2.fill.solid()
    box2.fill.fore_color.rgb = WHITE
    box2.line.fill.solid()
    box2.line.fore_color.rgb = BLACK
    box2.line.width = Pt(2)
    tf_box2 = box2.text_frame
    tf_box2.text = "现实：岸边挤满了100个\n同样拿着20块钱的人"
    tf_box2.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf_box2.paragraphs[0].font.size = Pt(11)
    tf_box2.paragraphs[0].font.bold = True
    tf_box2.paragraphs[0].font.color.rgb = BLACK
    tf_box2.paragraphs[0].font.name = "Microsoft YaHei"

    # Cormorant Placeholder
    bird = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.2), Inches(5.5), Inches(1.2), Inches(1.8))
    bird.fill.solid()
    bird.fill.fore_color.rgb = RGBColor(0xCF, 0xD8, 0xDC)
    bird.line.fill.solid()
    bird.line.fore_color.rgb = BLACK
    bird.line.width = Pt(2)
    tf_bird = bird.text_frame
    tf_bird.text = "鸬鹚\n(插画)"
    tf_bird.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf_bird.paragraphs[0].font.color.rgb = BLACK

    # RETIRED sign
    sign = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.3), Inches(6.5), Inches(1.0), Inches(0.4))
    sign.fill.solid()
    sign.fill.fore_color.rgb = RGBColor(0xFF, 0xF5, 0x9D)
    sign.line.fill.solid()
    sign.line.fore_color.rgb = BLACK
    sign.line.width = Pt(1)
    tf_sign = sign.text_frame
    tf_sign.text = "RETIRED"
    tf_sign.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf_sign.paragraphs[0].font.size = Pt(10)
    tf_sign.paragraphs[0].font.bold = True
    tf_sign.paragraphs[0].font.color.rgb = BLACK

    # Fact Text Box
    box3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.8), Inches(6.4), Inches(2.8), Inches(0.8))
    box3.fill.solid()
    box3.fill.fore_color.rgb = WHITE
    box3.line.fill.solid()
    box3.line.fore_color.rgb = BLACK
    box3.line.width = Pt(2)
    tf_box3 = box3.text_frame
    tf_box3.text = "猎奇点：那只配合拍照的\n鸬鹚其实已经“退休”了"
    tf_box3.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf_box3.paragraphs[0].font.size = Pt(11)
    tf_box3.paragraphs[0].font.bold = True
    tf_box3.paragraphs[0].font.color.rgb = BLACK
    tf_box3.paragraphs[0].font.name = "Microsoft YaHei"

    # Arrow to Fact Box
    arrow2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(9.3), Inches(6.7), Inches(0.4), Inches(0.2))
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = ORANGE_BORDER
    arrow2.line.fill.background()

    # 4. Page Number
    page_num = slide.shapes.add_textbox(Inches(12.5), Inches(7.0), Inches(0.8), Inches(0.4))
    tf_page = page_num.text_frame
    p_page = tf_page.paragraphs[0]
    p_page.text = "3 / 11"
    p_page.font.size = Pt(14)
    p_page.font.color.rgb = GRAY_TEXT
    p_page.alignment = PP_ALIGN.RIGHT



# ── Main ──
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

s0 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_1(s0)
prs.save(OUTPUT_PATH)
