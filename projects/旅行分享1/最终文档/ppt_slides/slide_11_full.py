from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# ── Color Palette ──
BLUE_HEADER = RGBColor(0x5B, 0x9B, 0xD5)
BLUE_DARK   = RGBColor(0x4A, 0x86, 0xC8)
CYAN        = RGBColor(0x00, 0xBC, 0xD4)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x33, 0x33, 0x33)
GRAY_TEXT   = RGBColor(0x55, 0x55, 0x55)
GRAY_BAR    = RGBColor(0xB0, 0xBE, 0xC5)
RED         = RGBColor(0xE5, 0x39, 0x35)
GREEN       = RGBColor(0x43, 0xA0, 0x47)
ORANGE      = RGBColor(0xFF, 0x98, 0x00)

ICON_BG  = RGBColor(0xE3, 0xE8, 0xED)
ICON_FG  = RGBColor(0x54, 0x6E, 0x7A)
FONT_NAME = "Microsoft YaHei"
SLIDE_WIDTH = Inches(13.333)
HEADER_H    = Inches(0.75)
SUBTITLE_Y  = Inches(0.95)


def add_header_banner(slide, title_text, bg_color=None):
    if bg_color is None:
        bg_color = BLUE_HEADER
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, HEADER_H
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = bg_color
    banner.line.fill.background()
    tf = banner.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(26)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME


def add_subtitle(slide, text, left, top, width=Inches(12), font_size=Pt(18)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.font.name = FONT_NAME
    return txBox


def add_icon_box(slide, left, top, symbol, size=Inches(0.48)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ICON_BG
    shape.line.fill.background()
    shape.adjustments[0] = 0.25
    tf = shape.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = symbol
    p.font.size = Pt(18)
    p.font.color.rgb = ICON_FG
    p.font.bold = False
    return shape


def add_bullet_item(slide, left, top, symbol, label, description,
                    width=Inches(5.5), desc_size=Pt(13)):
    add_icon_box(slide, left, top, symbol)
    text_left = left + Inches(0.65)
    txBox = slide.shapes.add_textbox(text_left, top - Inches(0.02), width, Inches(0.65))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run_label = p.add_run()
    run_label.text = label + "\uff1a"
    run_label.font.size = Pt(14)
    run_label.font.color.rgb = BLACK
    run_label.font.bold = True
    run_label.font.name = FONT_NAME
    run_desc = p.add_run()
    run_desc.text = description
    run_desc.font.size = desc_size
    run_desc.font.color.rgb = GRAY_TEXT
    run_desc.font.bold = False
    run_desc.font.name = FONT_NAME
    return txBox


def add_conclusion_box(slide, left, top, width, text, font_size=Pt(13)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = BLACK
    run.font.bold = True
    run.font.name = FONT_NAME
    return txBox


def add_table(slide, left, top, width, height, rows, cols, data,
              header_color=None, col_widths=None):
    if header_color is None:
        header_color = BLUE_HEADER
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if data[r][c] is not None else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(11)
                    run.font.name = FONT_NAME
                    if r == 0:
                        run.font.color.rgb = WHITE
                        run.font.bold = True
                    else:
                        run.font.color.rgb = BLACK
                        run.font.bold = False
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else RGBColor(0xF5, 0xF5, 0xF5)
            cell.margin_left = Pt(5)
            cell.margin_right = Pt(5)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
    return table_shape


def add_bar_chart(slide, left, top, width, height,
                  categories, values, title="", bar_colors=None):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', values)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    chart.chart_style = 2
    plot = chart.plots[0]
    plot.gap_width = 100
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CYAN
    series.has_data_labels = True
    dl = series.data_labels
    dl.font.size = Pt(13)
    dl.font.bold = True
    dl.font.color.rgb = BLACK
    dl.number_format = '0.#'
    dl.show_value = True
    dl.label_position = XL_LABEL_POSITION.OUTSIDE_END
    if bar_colors:
        for i, color in enumerate(bar_colors):
            pt = series.points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = color
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(12)
    cat_axis.tick_labels.font.name = FONT_NAME
    cat_axis.major_tick_mark = 2
    cat_axis.format.line.fill.background()
    val_axis = chart.value_axis
    val_axis.visible = False
    val_axis.major_tick_mark = 2
    val_axis.format.line.fill.background()
    val_axis.major_gridlines.format.line.fill.background()
    if title:
        chart.has_title = True
        ct = chart.chart_title.text_frame.paragraphs[0]
        ct.text = title
        ct.font.size = Pt(14)
        ct.font.bold = True
        ct.font.name = FONT_NAME
    else:
        chart.has_title = False
    return chart_frame


def add_callout_label(slide, left, top, text, bg_color=None, font_size=Pt(11)):
    if bg_color is None:
        bg_color = CYAN
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.3), Inches(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME
    return shape


def add_data_card(slide, left, top, width, height, value, label,
                  value_color=None, bg_color=None):
    if value_color is None:
        value_color = CYAN
    if bg_color is None:
        bg_color = WHITE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(3)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run1 = p1.add_run()
    run1.text = str(value)
    run1.font.size = Pt(24)
    run1.font.color.rgb = value_color
    run1.font.bold = True
    run1.font.name = FONT_NAME
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(10)
    run2.font.color.rgb = GRAY_TEXT
    run2.font.bold = False
    run2.font.name = FONT_NAME
    return shape


OUTPUT_PATH = r"C:\Users\Administrator\Desktop\notebook\my_notebooklm_with_ppteditable\projects\旅行分享1\最终文档\ppt_slides\slide_11.pptx"


# ── Slide 1 ──

def build_slide_1(slide):
    # 1. Background
    # Left background (Light beige)
    bg_left = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg_left.fill.solid()
    bg_left.fill.fore_color.rgb = RGBColor(0xF4, 0xF1, 0xEA)
    bg_left.line.fill.background()

    # Right background (Orange, angled)
    bg_right = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(-3), Inches(10), Inches(14))
    bg_right.rotation = 12
    bg_right.fill.solid()
    bg_right.fill.fore_color.rgb = RGBColor(0xFF, 0x7A, 0x00)
    bg_right.line.fill.background()

    # 2. Top Left Tag ("第11页")
    tag_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0.4), Inches(1.5), Inches(0.6))
    tag_box.fill.solid()
    tag_box.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)
    tag_box.line.fill.background()
    
    tf_tag = tag_box.text_frame
    tf_tag.clear()
    p_tag = tf_tag.paragraphs[0]
    p_tag.text = "第11页"
    p_tag.font.name = "Microsoft YaHei"
    p_tag.font.size = Pt(20)
    p_tag.font.bold = True
    p_tag.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p_tag.alignment = PP_ALIGN.CENTER

    # 3. Titles
    # Main Title Background
    title1_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.8), Inches(0.5), Inches(9.5), Inches(1.0))
    title1_bg.rotation = -2
    title1_bg.fill.solid()
    title1_bg.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    title1_bg.line.fill.background()

    # Main Title Text
    title1_tx = slide.shapes.add_textbox(Inches(2.9), Inches(0.5), Inches(9.3), Inches(1.0))
    title1_tx.rotation = -2
    tf1 = title1_tx.text_frame
    tf1.clear()
    p1 = tf1.paragraphs[0]
    p1.text = "感谢观看，记得给个“五星好评”"
    p1.font.name = "Microsoft YaHei"
    p1.font.size = Pt(40)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(0xE6, 0x4A, 0x19) # Deep Orange

    # Subtitle Background
    title2_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(1.5), Inches(7.5), Inches(0.7))
    title2_bg.fill.solid()
    title2_bg.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    title2_bg.line.fill.background()

    # Subtitle Text
    title2_tx = slide.shapes.add_textbox(Inches(5.8), Inches(1.55), Inches(7.0), Inches(0.6))
    tf2 = title2_tx.text_frame
    tf2.clear()
    p2 = tf2.paragraphs[0]
    p2.text = "关注我，带你解锁更多奇葩目的地"
    p2.font.name = "Microsoft YaHei"
    p2.font.size = Pt(26)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    # 4. Left Visuals (Mountains & River Placeholder)
    # Mountains (Overlapping green triangles)
    colors_mtn = [RGBColor(0x81, 0xC7, 0x84), RGBColor(0xA5, 0xD6, 0xA7), RGBColor(0x66, 0xBB, 0x6A)]
    mtn_coords = [(0.2, 3.5, 2.5, 4.0), (1.5, 2.5, 2.5, 5.0), (3.0, 3.5, 2.5, 4.0)]
    
    for i, (x, y, w, h) in enumerate(mtn_coords):
        mtn = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        mtn.fill.solid()
        mtn.fill.fore_color.rgb = colors_mtn[i]
        mtn.line.color.rgb = RGBColor(0x1B, 0x5E, 0x20)
        mtn.line.width = Pt(2)

    # River (Curved shape)
    river = slide.shapes.add_shape(MSO_SHAPE.MOON, Inches(1.0), Inches(5.0), Inches(4.5), Inches(2.5))
    river.rotation = 50
    river.fill.solid()
    river.fill.fore_color.rgb = RGBColor(0xE0, 0xF2, 0xF1)
    river.line.color.rgb = RGBColor(0x00, 0x4D, 0x40)
    river.line.width = Pt(2)

    # Decorations (Stars)
    star_coords = [(0.5, 2.0), (1.0, 2.8), (4.8, 4.2), (1.2, 6.5)]
    for x, y in star_coords:
        star = slide.shapes.add_shape(MSO_SHAPE.5_POINT_STAR, Inches(x), Inches(y), Inches(0.4), Inches(0.4))
        star.rotation = 15
        star.fill.background()
        star.line.color.rgb = RGBColor(0x9C, 0x27, 0xB0)
        star.line.width = Pt(2)

    # Text "BOOM!"
    boom_tx = slide.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(1.2), Inches(0.5))
    boom_tx.rotation = -15
    p_boom = boom_tx.text_frame.paragraphs[0]
    p_boom.text = "BOOM!"
    p_boom.font.name = "Arial"
    p_boom.font.size = Pt(18)
    p_boom.font.bold = True
    p_boom.font.color.rgb = RGBColor(0x9C, 0x27, 0xB0)

    # Text "OMG!"
    omg_coords = [(4.0, 2.2), (4.8, 3.2), (4.5, 5.5)]
    for x, y in omg_coords:
        omg_tx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(1.0), Inches(0.5))
        omg_tx.rotation = -10
        p_omg = omg_tx.text_frame.paragraphs[0]
        p_omg.text = "OMG!"
        p_omg.font.name = "Arial"
        p_omg.font.size = Pt(16)
        p_omg.font.bold = True
        p_omg.font.color.rgb = RGBColor(0xF4, 0x43, 0x36)

    # 5. Center Visual (Silhouette Placeholder)
    # Head
    head = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.5), Inches(2.6), Inches(0.8), Inches(1.0))
    head.fill.solid()
    head.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)
    head.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    head.line.width = Pt(4)

    # Body
    body = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(3.5), Inches(2.0), Inches(3.5))
    body.fill.solid()
    body.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)
    body.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    body.line.width = Pt(4)

    # Arm (waving)
    arm = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(2.8), Inches(0.6), Inches(1.5))
    arm.rotation = -30
    arm.fill.solid()
    arm.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)
    arm.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    arm.line.width = Pt(4)

    # 6. Right Visual (QR Code Area)
    # Sticky Note (Yellow)
    note = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.0), Inches(1.8), Inches(3.8), Inches(3.8))
    note.rotation = 3
    note.fill.solid()
    note.fill.fore_color.rgb = RGBColor(0xFF, 0xEB, 0x3B)
    note.line.fill.background()

    # QR Code Base (Orange)
    qr_base = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.4), Inches(2.2), Inches(3.0), Inches(3.0))
    qr_base.rotation = 3
    qr_base.fill.solid()
    qr_base.fill.fore_color.rgb = RGBColor(0xFF, 0x98, 0x00)
    qr_base.line.fill.background()

    # QR Code Inner details (Yellow squares)
    qr_inners = [(9.6, 2.4), (11.6, 2.5), (9.7, 4.4)]
    for x, y in qr_inners:
        inner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.6), Inches(0.6))
        inner.rotation = 3
        inner.fill.solid()
        inner.fill.fore_color.rgb = RGBColor(0xFF, 0xEB, 0x3B)
        inner.line.fill.background()

    # Hand-drawn circles around QR code
    circ1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.1), Inches(2.1), Inches(3.6), Inches(1.2))
    circ1.rotation = 3
    circ1.fill.background()
    circ1.line.color.rgb = RGBColor(0x00, 0x00, 0x00)
    circ1.line.width = Pt(1.5)

    circ2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.0), Inches(4.3), Inches(1.5), Inches(0.8))
    circ2.rotation = -5
    circ2.fill.background()
    circ2.line.color.rgb = RGBColor(0x00, 0x00, 0x00)
    circ2.line.width = Pt(1.5)

    # 7. Bottom Right Text Block ("内容要点")
    # Title
    tx_points_title = slide.shapes.add_textbox(Inches(8.8), Inches(5.3), Inches(3.0), Inches(0.5))
    p_title = tx_points_title.text_frame.paragraphs[0]
    p_title.text = "内容要点"
    p_title.font.name = "Microsoft YaHei"
    p_title.font.size = Pt(22)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    # Points Data
    points_data = [
        ("?", "互动区：你见过最奇葩的景点在哪里？"),
        ("@", "联系方式：微博/小红书/抖音同名"),
        ("!", "结束语：山水不改，脑洞常在！")
    ]

    start_y = 5.9
    for icon_char, text_content in points_data:
        # Icon circle
        icon_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.8), Inches(start_y), Inches(0.3), Inches(0.3))
        icon_shape.fill.solid()
        icon_shape.fill.fore_color.rgb = RGBColor(0xFF, 0xEB, 0x3B)
        icon_shape.line.color.rgb = RGBColor(0x00, 0x00, 0x00)
        icon_shape.line.width = Pt(1)
        
        icon_p = icon_shape.text_frame.paragraphs[0]
        icon_p.text = icon_char
        icon_p.font.name = "Arial"
        icon_p.font.size = Pt(12)
        icon_p.font.bold = True
        icon_p.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
        icon_p.alignment = PP_ALIGN.CENTER

        # Text
        tx_point = slide.shapes.add_textbox(Inches(9.2), Inches(start_y - 0.05), Inches(4.0), Inches(0.4))
        p_point = tx_point.text_frame.paragraphs[0]
        p_point.text = text_content
        p_point.font.name = "Microsoft YaHei"
        p_point.font.size = Pt(14)
        p_point.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
        
        start_y += 0.45



# ── Main ──
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

s0 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_1(s0)
prs.save(OUTPUT_PATH)
