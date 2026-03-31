from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# ── Color Palette ──
BLUE_HEADER = RGBColor(0x5B, 0x9B, 0xD5)
BLUE_DARK   = RGBColor(0x4A, 0x86, 0xC8)
CYAN        = RGBColor(0x00, 0xBC, 0xD4)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x33, 0x33, 0x33)
GRAY_TEXT   = RGBColor(0x55, 0x55, 0x55)
GRAY_BAR    = RGBColor(0xB0, 0xBE, 0xC5)
RED         = RGBColor(0xE5, 0x39, 0x35)
GREEN       = RGBColor(0x43, 0xA0, 0x47)
ORANGE      = RGBColor(0xFF, 0x98, 0x00)

ICON_BG  = RGBColor(0xE3, 0xE8, 0xED)
ICON_FG  = RGBColor(0x54, 0x6E, 0x7A)
FONT_NAME = "Microsoft YaHei"
SLIDE_WIDTH = Inches(13.333)
HEADER_H    = Inches(0.75)
SUBTITLE_Y  = Inches(0.95)


def add_header_banner(slide, title_text, bg_color=None):
    if bg_color is None:
        bg_color = BLUE_HEADER
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, HEADER_H
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = bg_color
    banner.line.fill.background()
    tf = banner.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(26)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME


def add_subtitle(slide, text, left, top, width=Inches(12), font_size=Pt(18)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.font.name = FONT_NAME
    return txBox


def add_icon_box(slide, left, top, symbol, size=Inches(0.48)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ICON_BG
    shape.line.fill.background()
    shape.adjustments[0] = 0.25
    tf = shape.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = symbol
    p.font.size = Pt(18)
    p.font.color.rgb = ICON_FG
    p.font.bold = False
    return shape


def add_bullet_item(slide, left, top, symbol, label, description,
                    width=Inches(5.5), desc_size=Pt(13)):
    add_icon_box(slide, left, top, symbol)
    text_left = left + Inches(0.65)
    txBox = slide.shapes.add_textbox(text_left, top - Inches(0.02), width, Inches(0.65))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run_label = p.add_run()
    run_label.text = label + "\uff1a"
    run_label.font.size = Pt(14)
    run_label.font.color.rgb = BLACK
    run_label.font.bold = True
    run_label.font.name = FONT_NAME
    run_desc = p.add_run()
    run_desc.text = description
    run_desc.font.size = desc_size
    run_desc.font.color.rgb = GRAY_TEXT
    run_desc.font.bold = False
    run_desc.font.name = FONT_NAME
    return txBox


def add_conclusion_box(slide, left, top, width, text, font_size=Pt(13)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = BLACK
    run.font.bold = True
    run.font.name = FONT_NAME
    return txBox


def add_table(slide, left, top, width, height, rows, cols, data,
              header_color=None, col_widths=None):
    if header_color is None:
        header_color = BLUE_HEADER
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if data[r][c] is not None else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(11)
                    run.font.name = FONT_NAME
                    if r == 0:
                        run.font.color.rgb = WHITE
                        run.font.bold = True
                    else:
                        run.font.color.rgb = BLACK
                        run.font.bold = False
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else RGBColor(0xF5, 0xF5, 0xF5)
            cell.margin_left = Pt(5)
            cell.margin_right = Pt(5)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
    return table_shape


def add_bar_chart(slide, left, top, width, height,
                  categories, values, title="", bar_colors=None):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', values)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    chart.chart_style = 2
    plot = chart.plots[0]
    plot.gap_width = 100
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CYAN
    series.has_data_labels = True
    dl = series.data_labels
    dl.font.size = Pt(13)
    dl.font.bold = True
    dl.font.color.rgb = BLACK
    dl.number_format = '0.#'
    dl.show_value = True
    dl.label_position = XL_LABEL_POSITION.OUTSIDE_END
    if bar_colors:
        for i, color in enumerate(bar_colors):
            pt = series.points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = color
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(12)
    cat_axis.tick_labels.font.name = FONT_NAME
    cat_axis.major_tick_mark = 2
    cat_axis.format.line.fill.background()
    val_axis = chart.value_axis
    val_axis.visible = False
    val_axis.major_tick_mark = 2
    val_axis.format.line.fill.background()
    val_axis.major_gridlines.format.line.fill.background()
    if title:
        chart.has_title = True
        ct = chart.chart_title.text_frame.paragraphs[0]
        ct.text = title
        ct.font.size = Pt(14)
        ct.font.bold = True
        ct.font.name = FONT_NAME
    else:
        chart.has_title = False
    return chart_frame


def add_callout_label(slide, left, top, text, bg_color=None, font_size=Pt(11)):
    if bg_color is None:
        bg_color = CYAN
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.3), Inches(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME
    return shape


def add_data_card(slide, left, top, width, height, value, label,
                  value_color=None, bg_color=None):
    if value_color is None:
        value_color = CYAN
    if bg_color is None:
        bg_color = WHITE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(3)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run1 = p1.add_run()
    run1.text = str(value)
    run1.font.size = Pt(24)
    run1.font.color.rgb = value_color
    run1.font.bold = True
    run1.font.name = FONT_NAME
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(10)
    run2.font.color.rgb = GRAY_TEXT
    run2.font.bold = False
    run2.font.name = FONT_NAME
    return shape


OUTPUT_PATH = r"C:\Users\Administrator\Desktop\notebook\my_notebooklm_with_ppteditable\projects\旅行分享1\最终文档\ppt_slides\slide_06.pptx"


# ── Slide 1 ──

def build_slide_1(slide):
    # Colors
    BG_COLOR = RGBColor(248, 245, 238)
    TITLE_COLOR = RGBColor(142, 36, 170)
    TEXT_COLOR = RGBColor(30, 30, 30)
    
    ORANGE = RGBColor(255, 152, 0)
    GREEN = RGBColor(76, 175, 80)
    CYAN = RGBColor(0, 188, 212)
    PURPLE = RGBColor(156, 39, 176)
    LIGHT_BLUE = RGBColor(225, 245, 254)
    YELLOW = RGBColor(255, 235, 59)

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_COLOR
    bg.line.fill.background()

    # --- Header ---
    # Main Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.2), Inches(11.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "芦笛岩：地底下的“80年代迪厅”"
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = TITLE_COLOR

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(2), Inches(1.2), Inches(9.333), Inches(0.6))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "这种审美真的很“硬核”"
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = TEXT_COLOR

    # --- Dividers ---
    line1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(4.5), Inches(2.0), Inches(4.5), Inches(6.5))
    line1.line.color.rgb = TEXT_COLOR
    line1.line.width = Pt(2)
    line1.line.dash_style = 7 

    line2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(8.8), Inches(2.0), Inches(8.8), Inches(6.5))
    line2.line.color.rgb = TEXT_COLOR
    line2.line.width = Pt(2)
    line2.line.dash_style = 7 

    # --- Left Column (Disco Cave) ---
    # Image Placeholder
    img_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(2.2), Inches(3.6), Inches(3.2))
    img_rect.fill.solid()
    img_rect.fill.fore_color.rgb = RGBColor(40, 20, 60)
    img_rect.line.color.rgb = YELLOW
    img_rect.line.width = Pt(3)
    
    # Laser lines
    laser1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(2.4), Inches(3.8), Inches(0.6), Inches(2.5))
    laser1.line.color.rgb = CYAN
    laser1.line.width = Pt(4)
    laser2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(2.4), Inches(3.8), Inches(4.2), Inches(2.8))
    laser2.line.color.rgb = PURPLE
    laser2.line.width = Pt(4)
    laser3 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(2.4), Inches(3.8), Inches(1.0), Inches(5.0))
    laser3.line.color.rgb = ORANGE
    laser3.line.width = Pt(4)

    # Inset Photo Placeholder
    inset_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(1.8), Inches(2.2), Inches(1.6))
    inset_rect.rotation = -10
    inset_rect.fill.solid()
    inset_rect.fill.fore_color.rgb = RGBColor(200, 200, 200)
    inset_rect.line.color.rgb = RGBColor(255, 255, 255)
    inset_rect.line.width = Pt(4)
    
    tape1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(1.7), Inches(0.8), Inches(0.2))
    tape1.rotation = -30
    tape1.fill.solid()
    tape1.fill.fore_color.rgb = GREEN
    tape1.line.fill.background()

    # Text
    left_text = slide.shapes.add_textbox(Inches(0.8), Inches(5.6), Inches(3.2), Inches(1.0))
    tf = left_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "五颜六色的LED灯把\n溶洞变成了夜店"
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = TEXT_COLOR

    # Boombox icon
    boombox = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.2), Inches(5.8), Inches(0.8), Inches(0.5))
    boombox.fill.solid()
    boombox.fill.fore_color.rgb = PURPLE
    boombox.line.color.rgb = TEXT_COLOR
    boombox.line.width = Pt(1.5)

    # --- Middle Column (Imagination) ---
    # Row 1
    rock1_bg = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4.8), Inches(2.2), Inches(1.6), Inches(1.6))
    rock1_bg.fill.solid()
    rock1_bg.fill.fore_color.rgb = ORANGE
    rock1_bg.line.color.rgb = TEXT_COLOR
    rock1_bg.line.width = Pt(2)
    rock1 = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(5.1), Inches(2.4), Inches(1.0), Inches(1.2))
    rock1.fill.solid()
    rock1.fill.fore_color.rgb = RGBColor(160, 140, 120)
    rock1.line.color.rgb = TEXT_COLOR
    
    arrow1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.5), Inches(2.8), Inches(0.6), Inches(0.4))
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = ORANGE
    arrow1.line.color.rgb = TEXT_COLOR

    pork_bg = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.2), Inches(2.2), Inches(1.6), Inches(1.6))
    pork_bg.fill.solid()
    pork_bg.fill.fore_color.rgb = GREEN
    pork_bg.line.color.rgb = TEXT_COLOR
    pork_bg.line.width = Pt(2)
    pork = slide.shapes.add_shape(MSO_SHAPE.CUBE, Inches(7.5), Inches(2.5), Inches(1.0), Inches(1.0))
    pork.fill.solid()
    pork.fill.fore_color.rgb = RGBColor(180, 80, 40)
    pork.line.color.rgb = TEXT_COLOR
    
    pork_text = slide.shapes.add_textbox(Inches(6.2), Inches(2.0), Inches(1.2), Inches(0.4))
    pork_text.text_frame.text = "红烧肉?"
    pork_text.text_frame.paragraphs[0].font.size = Pt(14)
    pork_text.text_frame.paragraphs[0].font.bold = True

    # Row 2
    rock2_bg = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4.8), Inches(4.0), Inches(1.6), Inches(1.6))
    rock2_bg.fill.solid()
    rock2_bg.fill.fore_color.rgb = CYAN
    rock2_bg.line.color.rgb = TEXT_COLOR
    rock2_bg.line.width = Pt(2)
    rock2 = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(5.1), Inches(4.2), Inches(1.0), Inches(1.2))
    rock2.fill.solid()
    rock2.fill.fore_color.rgb = RGBColor(160, 140, 120)
    rock2.line.color.rgb = TEXT_COLOR

    arrow2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.5), Inches(4.6), Inches(0.6), Inches(0.4))
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = GREEN
    arrow2.line.color.rgb = TEXT_COLOR

    alien_bg = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.2), Inches(4.0), Inches(1.6), Inches(1.6))
    alien_bg.fill.solid()
    alien_bg.fill.fore_color.rgb = PURPLE
    alien_bg.line.color.rgb = TEXT_COLOR
    alien_bg.line.width = Pt(2)
    alien = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.5), Inches(4.3), Inches(1.0), Inches(1.2))
    alien.fill.solid()
    alien.fill.fore_color.rgb = GREEN
    alien.line.color.rgb = TEXT_COLOR
    
    alien_text = slide.shapes.add_textbox(Inches(8.2), Inches(3.8), Inches(1.2), Inches(0.4))
    alien_text.text_frame.text = "外星人?"
    alien_text.text_frame.paragraphs[0].font.size = Pt(14)
    alien_text.text_frame.paragraphs[0].font.bold = True

    # Text
    mid_text = slide.shapes.add_textbox(Inches(5.2), Inches(5.6), Inches(3.5), Inches(1.0))
    tf = mid_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "猎奇想象：这块石头像红烧肉，\n那块石头像外星人"
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = TEXT_COLOR

    # Magnifying glass
    mag = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4.7), Inches(5.8), Inches(0.4), Inches(0.4))
    mag.fill.solid()
    mag.fill.fore_color.rgb = LIGHT_BLUE
    mag.line.color.rgb = TEXT_COLOR
    mag.line.width = Pt(1.5)
    mag_handle = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(4.8), Inches(6.1), Inches(4.6), Inches(6.3))
    mag_handle.line.color.rgb = ORANGE
    mag_handle.line.width = Pt(4)

    # --- Right Column (Comic Panels) ---
    # Panel 1
    panel1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.2), Inches(2.0), Inches(3.6), Inches(1.7))
    panel1.fill.solid()
    panel1.fill.fore_color.rgb = LIGHT_BLUE
    panel1.line.color.rgb = TEXT_COLOR
    panel1.line.width = Pt(2)
    
    face1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(2.4), Inches(0.8), Inches(0.8))
    face1.fill.solid()
    face1.fill.fore_color.rgb = RGBColor(255, 224, 189)
    face1.line.color.rgb = TEXT_COLOR
    
    bubble1 = slide.shapes.add_shape(MSO_SHAPE.CLOUD_CALLOUT, Inches(10.5), Inches(2.2), Inches(2.0), Inches(1.0))
    bubble1.fill.solid()
    bubble1.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bubble1.line.color.rgb = TEXT_COLOR
    bubble1.line.width = Pt(1.5)
    tf = bubble1.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "灵魂被洗涤"
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = CYAN

    # Panel 2
    panel2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.2), Inches(3.9), Inches(3.6), Inches(1.7))
    panel2.fill.solid()
    panel2.fill.fore_color.rgb = LIGHT_BLUE
    panel2.line.color.rgb = TEXT_COLOR
    panel2.line.width = Pt(2)

    face2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(4.3), Inches(0.8), Inches(0.8))
    face2.fill.solid()
    face2.fill.fore_color.rgb = RGBColor(255, 224, 189)
    face2.line.color.rgb = TEXT_COLOR
    
    scarf = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.4), Inches(4.9), Inches(1.0), Inches(0.3))
    scarf.fill.solid()
    scarf.fill.fore_color.rgb = ORANGE
    scarf.line.color.rgb = TEXT_COLOR

    bubble2 = slide.shapes.add_shape(MSO_SHAPE.OVAL_CALLOUT, Inches(10.8), Inches(4.0), Inches(1.5), Inches(1.0))
    bubble2.fill.solid()
    bubble2.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bubble2.line.color.rgb = TEXT_COLOR
    bubble2.line.width = Pt(1.5)
    tf = bubble2.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "冷！"
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = CYAN

    cold_text = slide.shapes.add_textbox(Inches(10.8), Inches(5.0), Inches(1.5), Inches(0.5))
    tf = cold_text.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "冻感冒了"
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = TEXT_COLOR

    # Text
    right_text = slide.shapes.add_textbox(Inches(9.0), Inches(5.6), Inches(4.0), Inches(1.0))
    tf = right_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "吐槽：在洞里走了一圈，感觉\n灵魂被洗涤（其实是冻感冒了）"
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = TEXT_COLOR

    # Thermometer
    thermo = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(12.8), Inches(5.8), Inches(0.2), Inches(0.8))
    thermo.fill.solid()
    thermo.fill.fore_color.rgb = RGBColor(200, 200, 200)
    thermo.line.color.rgb = TEXT_COLOR
    thermo_bulb = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(12.75), Inches(6.4), Inches(0.3), Inches(0.3))
    thermo_bulb.fill.solid()
    thermo_bulb.fill.fore_color.rgb = RGBColor(255, 0, 0)
    thermo_bulb.line.color.rgb = TEXT_COLOR

    # --- Footer ---
    page_bg = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(12.2), Inches(7.0), Inches(1.0), Inches(0.4))
    page_bg.fill.solid()
    page_bg.fill.fore_color.rgb = GREEN
    page_bg.line.color.rgb = TEXT_COLOR
    page_bg.line.width = Pt(1.5)
    
    page_text = slide.shapes.add_textbox(Inches(12.3), Inches(6.95), Inches(0.8), Inches(0.4))
    tf = page_text.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "第6页"
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = TEXT_COLOR

    # Decorations
    excl1 = slide.shapes.add_textbox(Inches(4.5), Inches(1.8), Inches(0.5), Inches(0.5))
    excl1.text_frame.text = "!!"
    excl1.text_frame.paragraphs[0].font.size = Pt(24)
    excl1.text_frame.paragraphs[0].font.bold = True
    excl1.text_frame.paragraphs[0].font.color.rgb = ORANGE
    excl1.rotation = -15

    excl2 = slide.shapes.add_textbox(Inches(4.5), Inches(3.8), Inches(0.5), Inches(0.5))
    excl2.text_frame.text = "!!"
    excl2.text_frame.paragraphs[0].font.size = Pt(24)
    excl2.text_frame.paragraphs[0].font.bold = True
    excl2.text_frame.paragraphs[0].font.color.rgb = PURPLE
    excl2.rotation = -15

    wow = slide.shapes.add_textbox(Inches(2.8), Inches(1.5), Inches(1.0), Inches(0.5))
    wow.text_frame.text = "WOW!"
    wow.text_frame.paragraphs[0].font.size = Pt(20)
    wow.text_frame.paragraphs[0].font.bold = True
    wow.text_frame.paragraphs[0].font.color.rgb = PURPLE
    wow.rotation = -10

    omg = slide.shapes.add_textbox(Inches(8.0), Inches(1.8), Inches(1.0), Inches(0.5))
    omg.text_frame.text = "OMG!"
    omg.text_frame.paragraphs[0].font.size = Pt(20)
    omg.text_frame.paragraphs[0].font.bold = True
    omg.text_frame.paragraphs[0].font.color.rgb = GREEN
    omg.rotation = 15



# ── Main ──
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

s0 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_1(s0)
prs.save(OUTPUT_PATH)
