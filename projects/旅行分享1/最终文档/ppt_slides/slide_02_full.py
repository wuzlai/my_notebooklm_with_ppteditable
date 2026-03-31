from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# ── Color Palette ──
BLUE_HEADER = RGBColor(0x5B, 0x9B, 0xD5)
BLUE_DARK   = RGBColor(0x4A, 0x86, 0xC8)
CYAN        = RGBColor(0x00, 0xBC, 0xD4)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x33, 0x33, 0x33)
GRAY_TEXT   = RGBColor(0x55, 0x55, 0x55)
GRAY_BAR    = RGBColor(0xB0, 0xBE, 0xC5)
RED         = RGBColor(0xE5, 0x39, 0x35)
GREEN       = RGBColor(0x43, 0xA0, 0x47)
ORANGE      = RGBColor(0xFF, 0x98, 0x00)

ICON_BG  = RGBColor(0xE3, 0xE8, 0xED)
ICON_FG  = RGBColor(0x54, 0x6E, 0x7A)
FONT_NAME = "Microsoft YaHei"
SLIDE_WIDTH = Inches(13.333)
HEADER_H    = Inches(0.75)
SUBTITLE_Y  = Inches(0.95)


def add_header_banner(slide, title_text, bg_color=None):
    if bg_color is None:
        bg_color = BLUE_HEADER
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, HEADER_H
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = bg_color
    banner.line.fill.background()
    tf = banner.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(26)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME


def add_subtitle(slide, text, left, top, width=Inches(12), font_size=Pt(18)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.font.name = FONT_NAME
    return txBox


def add_icon_box(slide, left, top, symbol, size=Inches(0.48)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ICON_BG
    shape.line.fill.background()
    shape.adjustments[0] = 0.25
    tf = shape.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = symbol
    p.font.size = Pt(18)
    p.font.color.rgb = ICON_FG
    p.font.bold = False
    return shape


def add_bullet_item(slide, left, top, symbol, label, description,
                    width=Inches(5.5), desc_size=Pt(13)):
    add_icon_box(slide, left, top, symbol)
    text_left = left + Inches(0.65)
    txBox = slide.shapes.add_textbox(text_left, top - Inches(0.02), width, Inches(0.65))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run_label = p.add_run()
    run_label.text = label + "\uff1a"
    run_label.font.size = Pt(14)
    run_label.font.color.rgb = BLACK
    run_label.font.bold = True
    run_label.font.name = FONT_NAME
    run_desc = p.add_run()
    run_desc.text = description
    run_desc.font.size = desc_size
    run_desc.font.color.rgb = GRAY_TEXT
    run_desc.font.bold = False
    run_desc.font.name = FONT_NAME
    return txBox


def add_conclusion_box(slide, left, top, width, text, font_size=Pt(13)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = BLACK
    run.font.bold = True
    run.font.name = FONT_NAME
    return txBox


def add_table(slide, left, top, width, height, rows, cols, data,
              header_color=None, col_widths=None):
    if header_color is None:
        header_color = BLUE_HEADER
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if data[r][c] is not None else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(11)
                    run.font.name = FONT_NAME
                    if r == 0:
                        run.font.color.rgb = WHITE
                        run.font.bold = True
                    else:
                        run.font.color.rgb = BLACK
                        run.font.bold = False
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else RGBColor(0xF5, 0xF5, 0xF5)
            cell.margin_left = Pt(5)
            cell.margin_right = Pt(5)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
    return table_shape


def add_bar_chart(slide, left, top, width, height,
                  categories, values, title="", bar_colors=None):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', values)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    chart.chart_style = 2
    plot = chart.plots[0]
    plot.gap_width = 100
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CYAN
    series.has_data_labels = True
    dl = series.data_labels
    dl.font.size = Pt(13)
    dl.font.bold = True
    dl.font.color.rgb = BLACK
    dl.number_format = '0.#'
    dl.show_value = True
    dl.label_position = XL_LABEL_POSITION.OUTSIDE_END
    if bar_colors:
        for i, color in enumerate(bar_colors):
            pt = series.points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = color
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(12)
    cat_axis.tick_labels.font.name = FONT_NAME
    cat_axis.major_tick_mark = 2
    cat_axis.format.line.fill.background()
    val_axis = chart.value_axis
    val_axis.visible = False
    val_axis.major_tick_mark = 2
    val_axis.format.line.fill.background()
    val_axis.major_gridlines.format.line.fill.background()
    if title:
        chart.has_title = True
        ct = chart.chart_title.text_frame.paragraphs[0]
        ct.text = title
        ct.font.size = Pt(14)
        ct.font.bold = True
        ct.font.name = FONT_NAME
    else:
        chart.has_title = False
    return chart_frame


def add_callout_label(slide, left, top, text, bg_color=None, font_size=Pt(11)):
    if bg_color is None:
        bg_color = CYAN
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.3), Inches(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME
    return shape


def add_data_card(slide, left, top, width, height, value, label,
                  value_color=None, bg_color=None):
    if value_color is None:
        value_color = CYAN
    if bg_color is None:
        bg_color = WHITE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(3)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run1 = p1.add_run()
    run1.text = str(value)
    run1.font.size = Pt(24)
    run1.font.color.rgb = value_color
    run1.font.bold = True
    run1.font.name = FONT_NAME
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(10)
    run2.font.color.rgb = GRAY_TEXT
    run2.font.bold = False
    run2.font.name = FONT_NAME
    return shape


OUTPUT_PATH = r"C:\Users\Administrator\Desktop\notebook\my_notebooklm_with_ppteditable\projects\旅行分享1\最终文档\ppt_slides\slide_02.pptx"


# ── Slide 1 ──

def build_slide_1(slide):
    # 1. 设置背景颜色 (浅灰白)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(248, 248, 248)

    # 右侧边缘装饰条 (渐变橙/绿效果的简化)
    edge = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(13.1), Inches(0), Inches(0.233), Inches(7.5))
    edge.fill.solid()
    edge.fill.fore_color.rgb = RGBColor(230, 120, 50)
    edge.line.fill.background()

    # 2. 标题与副标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(10), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "角色设定：谁在带你逛桂林？"
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = RGBColor(10, 10, 10)

    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "颜值不够，脑洞来凑"
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(20, 20, 20)

    # 页码
    pg_box = slide.shapes.add_textbox(Inches(11.5), Inches(0.4), Inches(1.5), Inches(0.6))
    tf = pg_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "第2页"
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(22)
    run.font.bold = True

    # 3. 左侧视觉区域 (橙色泼墨背景 + 人物抠图占位)
    splash_color = RGBColor(255, 102, 0)
    
    # 使用多个椭圆组合模拟泼墨形状
    s1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), Inches(2.8), Inches(3.8), Inches(3.8))
    s1.fill.solid(); s1.fill.fore_color.rgb = splash_color; s1.line.fill.background()
    s2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), Inches(4.0), Inches(2.0), Inches(2.0))
    s2.fill.solid(); s2.fill.fore_color.rgb = splash_color; s2.line.fill.background()
    s3 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.5), Inches(2.5), Inches(2.0), Inches(2.0))
    s3.fill.solid(); s3.fill.fore_color.rgb = splash_color; s3.line.fill.background()
    s4 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.2), Inches(5.2), Inches(2.2), Inches(1.8))
    s4.fill.solid(); s4.fill.fore_color.rgb = splash_color; s4.line.fill.background()

    # 人物抠图占位符 (带白色粗边框的圆角矩形)
    cutout = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.8), Inches(2.8), Inches(2.4), Inches(3.8))
    cutout.fill.solid()
    cutout.fill.fore_color.rgb = RGBColor(30, 30, 30) # 深色衣服
    cutout.line.color.rgb = RGBColor(255, 255, 255)
    cutout.line.width = Pt(6)
    
    # 脸部占位
    face = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.4), Inches(3.0), Inches(1.2), Inches(1.5))
    face.fill.solid(); face.fill.fore_color.rgb = RGBColor(255, 218, 185); face.line.fill.background()
    
    # 相机占位
    cam = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.5), Inches(4.5), Inches(1.0), Inches(0.7))
    cam.fill.solid(); cam.fill.fore_color.rgb = RGBColor(50, 50, 50); cam.line.fill.background()
    lens = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.65), Inches(4.55), Inches(0.7), Inches(0.6))
    lens.fill.solid(); lens.fill.fore_color.rgb = RGBColor(10, 10, 10); lens.line.color.rgb = RGBColor(100, 100, 100)

    # "粗糙抠图" 文本
    cutout_txt = slide.shapes.add_textbox(Inches(2.2), Inches(6.7), Inches(2.0), Inches(0.5))
    tf = cutout_txt.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "粗糙抠图"
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(14)
    run.font.bold = True

    # 装饰性英文文本
    def add_deco_text(text, left, top, rotation, color, size=22):
        box = slide.shapes.add_textbox(left, top, Inches(1.5), Inches(0.8))
        box.rotation = rotation
        tf = box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.name = "Arial Black"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = True

    add_deco_text("OMG!", Inches(1.0), Inches(3.0), -15, RGBColor(0, 160, 140))
    add_deco_text("!?", Inches(3.5), Inches(2.8), 15, RGBColor(255, 255, 255), size=28)
    add_deco_text("BOOM!", Inches(3.8), Inches(3.8), 15, RGBColor(255, 255, 255))
    add_deco_text("LOOK!", Inches(4.2), Inches(5.5), 20, RGBColor(255, 255, 255))
    add_deco_text("NEW!", Inches(0.8), Inches(6.2), -10, RGBColor(0, 160, 140))
    add_deco_text("NEW!", Inches(11.5), Inches(2.0), 20, RGBColor(255, 112, 0))

    # 4. 右侧内容区域
    def add_content_block(left, top, width, title, desc):
        box = slide.shapes.add_textbox(left, top, width, Inches(1.5))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.line_spacing = 1.3
        
        run1 = p.add_run()
        run1.text = title
        run1.font.name = "Microsoft YaHei"
        run1.font.size = Pt(20)
        run1.font.bold = True
        run1.font.color.rgb = RGBColor(0, 0, 0)
        
        run2 = p.add_run()
        run2.text = desc
        run2.font.name = "Microsoft YaHei"
        run2.font.size = Pt(20)
        run2.font.bold = True
        run2.font.color.rgb = RGBColor(0, 0, 0)

    # --- 条目 1: 身份 ---
    i1_left, i1_top = Inches(5.5), Inches(2.5)
    # 图标占位 (人物 + 地图)
    p_body = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, i1_left+Inches(0.2), i1_top+Inches(0.6), Inches(0.8), Inches(0.6))
    p_body.fill.solid(); p_body.fill.fore_color.rgb = RGBColor(220, 50, 50); p_body.line.color.rgb = RGBColor(0, 0, 0)
    p_head = slide.shapes.add_shape(MSO_SHAPE.OVAL, i1_left+Inches(0.35), i1_top+Inches(0.1), Inches(0.5), Inches(0.5))
    p_head.fill.solid(); p_head.fill.fore_color.rgb = RGBColor(255, 200, 180); p_head.line.color.rgb = RGBColor(0, 0, 0)
    
    map_s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, i1_left+Inches(1.1), i1_top+Inches(0.3), Inches(0.8), Inches(0.7))
    map_s.rotation = 10
    map_s.fill.solid(); map_s.fill.fore_color.rgb = RGBColor(150, 200, 150); map_s.line.color.rgb = RGBColor(0, 0, 0)
    
    # 地图上的红叉
    x1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, i1_left+Inches(1.2), i1_top+Inches(0.4), i1_left+Inches(1.8), i1_top+Inches(0.9))
    x1.line.color.rgb = RGBColor(255, 0, 0); x1.line.width = Pt(3)
    x2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, i1_left+Inches(1.8), i1_top+Inches(0.4), i1_left+Inches(1.2), i1_top+Inches(0.9))
    x2.line.color.rgb = RGBColor(255, 0, 0); x2.line.width = Pt(3)

    # allergy 文本
    alg_box = slide.shapes.add_textbox(i1_left+Inches(1.0), i1_top+Inches(1.0), Inches(1.0), Inches(0.3))
    alg_box.rotation = -5
    tf = alg_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "allergy"
    run.font.name = "Arial"
    run.font.size = Pt(12)
    run.font.bold = True
    
    add_content_block(Inches(7.8), Inches(2.5), Inches(5.0), "身份：", "一个对“正常景点”\n过敏的猎奇博主")

    # --- 条目 2: 装备 ---
    add_content_block(Inches(5.8), Inches(4.2), Inches(4.5), "装备：", "自拍杆、扩音器、\n以及随时准备跑路的运动鞋")
    
    i2_left, i2_top = Inches(10.2), Inches(4.0)
    # 图标占位 (紫色背景 + 喇叭 + 鞋子)
    bg2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, i2_left, i2_top, Inches(1.2), Inches(1.2))
    bg2.fill.solid(); bg2.fill.fore_color.rgb = RGBColor(160, 32, 240); bg2.line.fill.background()
    
    mega = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, i2_left+Inches(0.6), i2_top-Inches(0.2), Inches(0.6), Inches(0.8))
    mega.rotation = 90
    mega.fill.solid(); mega.fill.fore_color.rgb = RGBColor(255, 255, 255); mega.line.color.rgb = RGBColor(160, 32, 240); mega.line.width = Pt(2)
    
    shoe = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, i2_left+Inches(1.0), i2_top+Inches(0.8), Inches(1.0), Inches(0.4))
    shoe.rotation = -15
    shoe.fill.solid(); shoe.fill.fore_color.rgb = RGBColor(255, 150, 0); shoe.line.color.rgb = RGBColor(0, 0, 0)
    
    add_deco_text("AHA!", i2_left+Inches(1.2), i2_top-Inches(0.5), -10, RGBColor(160, 32, 240))

    # --- 条目 3: 目标 ---
    add_content_block(Inches(5.8), Inches(6.0), Inches(3.5), "目标：\n", "寻找桂林最“野”的打开\n方式")
    
    i3_left, i3_top = Inches(9.5), Inches(5.8)
    # 图标占位 (青色背景 + 山水 + 指南针)
    bg3 = slide.shapes.add_shape(MSO_SHAPE.OVAL, i3_left+Inches(0.5), i3_top, Inches(1.4), Inches(1.4))
    bg3.fill.solid(); bg3.fill.fore_color.rgb = RGBColor(0, 160, 140); bg3.line.fill.background()
    
    m1 = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, i3_left+Inches(0.6), i3_top+Inches(0.4), Inches(0.6), Inches(0.8))
    m1.fill.solid(); m1.fill.fore_color.rgb = RGBColor(20, 20, 20); m1.line.fill.background()
    m2 = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, i3_left+Inches(1.0), i3_top+Inches(0.2), Inches(0.8), Inches(1.0))
    m2.fill.solid(); m2.fill.fore_color.rgb = RGBColor(10, 10, 10); m2.line.fill.background()
    
    comp = slide.shapes.add_shape(MSO_SHAPE.OVAL, i3_left, i3_top+Inches(0.5), Inches(0.8), Inches(0.8))
    comp.fill.solid(); comp.fill.fore_color.rgb = RGBColor(255, 220, 150); comp.line.color.rgb = RGBColor(0, 0, 0); comp.line.width = Pt(2)
    
    add_deco_text("WILD?", i3_left+Inches(1.5), i3_top-Inches(0.4), 15, RGBColor(0, 0, 0))

    # 5. 添加手绘风格的连接线/箭头指示
    a1 = slide.shapes.add_connector(MSO_CONNECTOR.CURVE, Inches(7.3), Inches(3.5), Inches(7.7), Inches(3.4))
    a1.line.color.rgb = RGBColor(150, 150, 150)
    
    a2 = slide.shapes.add_connector(MSO_CONNECTOR.CURVE, Inches(9.5), Inches(5.0), Inches(10.0), Inches(4.8))
    a2.line.color.rgb = RGBColor(150, 150, 150)
    
    a3 = slide.shapes.add_connector(MSO_CONNECTOR.CURVE, Inches(8.8), Inches(6.8), Inches(9.3), Inches(6.5))
    a3.line.color.rgb = RGBColor(150, 150, 150)



# ── Main ──
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

s0 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_1(s0)
prs.save(OUTPUT_PATH)
