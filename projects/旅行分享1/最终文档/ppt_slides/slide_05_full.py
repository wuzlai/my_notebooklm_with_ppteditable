from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# ── Color Palette ──
BLUE_HEADER = RGBColor(0x5B, 0x9B, 0xD5)
BLUE_DARK   = RGBColor(0x4A, 0x86, 0xC8)
CYAN        = RGBColor(0x00, 0xBC, 0xD4)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x33, 0x33, 0x33)
GRAY_TEXT   = RGBColor(0x55, 0x55, 0x55)
GRAY_BAR    = RGBColor(0xB0, 0xBE, 0xC5)
RED         = RGBColor(0xE5, 0x39, 0x35)
GREEN       = RGBColor(0x43, 0xA0, 0x47)
ORANGE      = RGBColor(0xFF, 0x98, 0x00)

ICON_BG  = RGBColor(0xE3, 0xE8, 0xED)
ICON_FG  = RGBColor(0x54, 0x6E, 0x7A)
FONT_NAME = "Microsoft YaHei"
SLIDE_WIDTH = Inches(13.333)
HEADER_H    = Inches(0.75)
SUBTITLE_Y  = Inches(0.95)


def add_header_banner(slide, title_text, bg_color=None):
    if bg_color is None:
        bg_color = BLUE_HEADER
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, HEADER_H
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = bg_color
    banner.line.fill.background()
    tf = banner.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(26)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME


def add_subtitle(slide, text, left, top, width=Inches(12), font_size=Pt(18)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.font.name = FONT_NAME
    return txBox


def add_icon_box(slide, left, top, symbol, size=Inches(0.48)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ICON_BG
    shape.line.fill.background()
    shape.adjustments[0] = 0.25
    tf = shape.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = symbol
    p.font.size = Pt(18)
    p.font.color.rgb = ICON_FG
    p.font.bold = False
    return shape


def add_bullet_item(slide, left, top, symbol, label, description,
                    width=Inches(5.5), desc_size=Pt(13)):
    add_icon_box(slide, left, top, symbol)
    text_left = left + Inches(0.65)
    txBox = slide.shapes.add_textbox(text_left, top - Inches(0.02), width, Inches(0.65))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run_label = p.add_run()
    run_label.text = label + "\uff1a"
    run_label.font.size = Pt(14)
    run_label.font.color.rgb = BLACK
    run_label.font.bold = True
    run_label.font.name = FONT_NAME
    run_desc = p.add_run()
    run_desc.text = description
    run_desc.font.size = desc_size
    run_desc.font.color.rgb = GRAY_TEXT
    run_desc.font.bold = False
    run_desc.font.name = FONT_NAME
    return txBox


def add_conclusion_box(slide, left, top, width, text, font_size=Pt(13)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = BLACK
    run.font.bold = True
    run.font.name = FONT_NAME
    return txBox


def add_table(slide, left, top, width, height, rows, cols, data,
              header_color=None, col_widths=None):
    if header_color is None:
        header_color = BLUE_HEADER
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if data[r][c] is not None else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(11)
                    run.font.name = FONT_NAME
                    if r == 0:
                        run.font.color.rgb = WHITE
                        run.font.bold = True
                    else:
                        run.font.color.rgb = BLACK
                        run.font.bold = False
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else RGBColor(0xF5, 0xF5, 0xF5)
            cell.margin_left = Pt(5)
            cell.margin_right = Pt(5)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
    return table_shape


def add_bar_chart(slide, left, top, width, height,
                  categories, values, title="", bar_colors=None):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', values)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    chart.chart_style = 2
    plot = chart.plots[0]
    plot.gap_width = 100
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CYAN
    series.has_data_labels = True
    dl = series.data_labels
    dl.font.size = Pt(13)
    dl.font.bold = True
    dl.font.color.rgb = BLACK
    dl.number_format = '0.#'
    dl.show_value = True
    dl.label_position = XL_LABEL_POSITION.OUTSIDE_END
    if bar_colors:
        for i, color in enumerate(bar_colors):
            pt = series.points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = color
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(12)
    cat_axis.tick_labels.font.name = FONT_NAME
    cat_axis.major_tick_mark = 2
    cat_axis.format.line.fill.background()
    val_axis = chart.value_axis
    val_axis.visible = False
    val_axis.major_tick_mark = 2
    val_axis.format.line.fill.background()
    val_axis.major_gridlines.format.line.fill.background()
    if title:
        chart.has_title = True
        ct = chart.chart_title.text_frame.paragraphs[0]
        ct.text = title
        ct.font.size = Pt(14)
        ct.font.bold = True
        ct.font.name = FONT_NAME
    else:
        chart.has_title = False
    return chart_frame


def add_callout_label(slide, left, top, text, bg_color=None, font_size=Pt(11)):
    if bg_color is None:
        bg_color = CYAN
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.3), Inches(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME
    return shape


def add_data_card(slide, left, top, width, height, value, label,
                  value_color=None, bg_color=None):
    if value_color is None:
        value_color = CYAN
    if bg_color is None:
        bg_color = WHITE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(3)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run1 = p1.add_run()
    run1.text = str(value)
    run1.font.size = Pt(24)
    run1.font.color.rgb = value_color
    run1.font.bold = True
    run1.font.name = FONT_NAME
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(10)
    run2.font.color.rgb = GRAY_TEXT
    run2.font.bold = False
    run2.font.name = FONT_NAME
    return shape


OUTPUT_PATH = r"C:\Users\Administrator\Desktop\notebook\my_notebooklm_with_ppteditable\projects\旅行分享1\最终文档\ppt_slides\slide_05.pptx"


# ── Slide 1 ──

def build_slide_1(slide):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
    
    # 定义颜色
    C_BG = RGBColor(253, 249, 238)
    C_RED = RGBColor(228, 57, 50)
    C_YELLOW = RGBColor(255, 214, 89)
    C_GREEN = RGBColor(0, 150, 100)
    C_ORANGE = RGBColor(244, 121, 32)
    C_BLACK = RGBColor(0, 0, 0)
    C_WHITE = RGBColor(255, 255, 255)

    # 设置背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_BG
    bg.line.fill.background()
    
    # 添加背景装饰色块
    blob1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(-2), Inches(5), Inches(5))
    blob1.fill.solid(); blob1.fill.fore_color.rgb = C_RED; blob1.line.fill.background()
    
    blob2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3), Inches(-2), Inches(6), Inches(4))
    blob2.fill.solid(); blob2.fill.fore_color.rgb = C_YELLOW; blob2.line.fill.background()
    
    blob3 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(-2), Inches(6), Inches(5))
    blob3.fill.solid(); blob3.fill.fore_color.rgb = C_GREEN; blob3.line.fill.background()
    
    blob4 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1), Inches(6), Inches(5), Inches(4))
    blob4.fill.solid(); blob4.fill.fore_color.rgb = C_GREEN; blob4.line.fill.background()
    
    blob5 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8), Inches(6), Inches(7), Inches(4))
    blob5.fill.solid(); blob5.fill.fore_color.rgb = C_ORANGE; blob5.line.fill.background()

    # 标题 (使用多重偏移模拟描边效果)
    title_text = "螺蛳粉：桂林的“生化武器”诱惑"
    offsets = [(-0.03, -0.03), (0.03, -0.03), (-0.03, 0.03), (0.03, 0.03), (-0.04, 0), (0.04, 0), (0, -0.04), (0, 0.04)]
    for ox, oy in offsets:
        tb = slide.shapes.add_textbox(Inches(0.5 + ox), Inches(0.3 + oy), Inches(10), Inches(1))
        p = tb.text_frame.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = C_BLACK
        p.font.name = "Microsoft YaHei"
        
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(10), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p.font.name = "Microsoft YaHei"

    # 副标题
    sub_text = "闻着臭，吃着爽，回味想撞墙"
    for ox, oy in offsets:
        tb = slide.shapes.add_textbox(Inches(0.5 + ox), Inches(1.1 + oy), Inches(8), Inches(0.8))
        p = tb.text_frame.paragraphs[0]
        p.text = sub_text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = C_BLACK
        p.font.name = "Microsoft YaHei"
        
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(8), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.text = sub_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p.font.name = "Microsoft YaHei"

    # 页码
    page_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.8), Inches(0.4), Inches(1.0), Inches(1.0))
    page_circle.fill.solid(); page_circle.fill.fore_color.rgb = C_BG
    page_circle.line.color.rgb = C_BLACK; page_circle.line.width = Pt(2)
    
    page_inner = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.85), Inches(0.45), Inches(0.9), Inches(0.9))
    page_inner.fill.background()
    page_inner.line.color.rgb = C_BLACK; page_inner.line.width = Pt(1)
    
    tb = slide.shapes.add_textbox(Inches(11.8), Inches(0.5), Inches(1.0), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.text = "PAGE\n5"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = C_BLACK

    # --- 左侧：内容要点 ---
    shadow_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(2.6), Inches(2.0), Inches(0.6))
    shadow_box.fill.solid(); shadow_box.fill.fore_color.rgb = C_GREEN
    shadow_box.line.color.rgb = C_BLACK; shadow_box.line.width = Pt(1.5)
    
    main_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.5), Inches(2.0), Inches(0.6))
    main_box.fill.solid(); main_box.fill.fore_color.rgb = C_YELLOW
    main_box.line.color.rgb = C_BLACK; main_box.line.width = Pt(1.5)
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.55), Inches(2.0), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = "内容要点"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = C_BLACK

    # 要点 1
    bullet1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(3.6), Inches(0.15), Inches(0.15))
    bullet1.fill.solid(); bullet1.fill.fore_color.rgb = C_RED; bullet1.line.color.rgb = C_BLACK
    
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(3.4), Inches(3.0), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = "气味等级："
    p.font.size = Pt(20)
    p.font.bold = True
    
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(3.9), Inches(3.0), Inches(1.0))
    p = tb.text_frame.paragraphs[0]
    p.text = "路过店门口，衣服\n自动“腌制”三天"
    p.font.size = Pt(18)
    
    # 腌制 圈注
    oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.5), Inches(4.2), Inches(0.9), Inches(0.4))
    oval.fill.background()
    oval.line.color.rgb = C_ORANGE; oval.line.width = Pt(2)

    # 要点 2
    bullet2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(5.4), Inches(0.15), Inches(0.15))
    bullet2.fill.solid(); bullet2.fill.fore_color.rgb = C_RED; bullet2.line.color.rgb = C_BLACK
    
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(5.2), Inches(3.0), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = "猎奇吃法："
    p.font.size = Pt(20)
    p.font.bold = True
    
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(5.7), Inches(3.0), Inches(1.0))
    p = tb.text_frame.paragraphs[0]
    run1 = p.add_run()
    run1.text = "加辣加酸笋"
    run1.font.size = Pt(18)
    run1.font.underline = True
    run2 = p.add_run()
    run2.text = "，\n挑战味蕾极限"
    run2.font.size = Pt(18)

    # --- 中间：爆炸图与杀伤力 ---
    exp_red = slide.shapes.add_shape(MSO_SHAPE.EXPLOSION2, Inches(3.8), Inches(2.2), Inches(5.5), Inches(5.0))
    exp_red.fill.solid(); exp_red.fill.fore_color.rgb = C_RED
    exp_red.line.color.rgb = C_BLACK; exp_red.line.width = Pt(2)
    
    exp_yellow = slide.shapes.add_shape(MSO_SHAPE.EXPLOSION1, Inches(4.3), Inches(2.7), Inches(4.5), Inches(4.0))
    exp_yellow.fill.solid(); exp_yellow.fill.fore_color.rgb = C_YELLOW
    exp_yellow.line.color.rgb = C_BLACK; exp_yellow.line.width = Pt(2)
    
    # 中心碗
    bowl_outer = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.4), Inches(3.8), Inches(2.2), Inches(1.6))
    bowl_outer.fill.solid(); bowl_outer.fill.fore_color.rgb = C_WHITE
    bowl_outer.line.color.rgb = C_BLACK; bowl_outer.line.width = Pt(2)
    
    bowl_inner = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.6), Inches(3.9), Inches(1.8), Inches(0.8))
    bowl_inner.fill.solid(); bowl_inner.fill.fore_color.rgb = C_ORANGE
    bowl_inner.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(5.4), Inches(4.7), Inches(2.2), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = "杀伤力"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)
    p.font.bold = True

    # 放射箭头
    def add_custom_arrow(slide, x, y, width, height, rotation):
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(width), Inches(height))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = C_YELLOW
        arrow.line.color.rgb = C_RED; arrow.line.width = Pt(1.5)
        arrow.rotation = rotation

    add_custom_arrow(slide, 5.4, 3.4, 0.8, 0.4, 225)
    add_custom_arrow(slide, 6.1, 3.0, 0.8, 0.4, 270)
    add_custom_arrow(slide, 7.1, 3.4, 0.8, 0.4, 315)
    add_custom_arrow(slide, 7.6, 4.5, 0.8, 0.4, 0)
    add_custom_arrow(slide, 7.1, 5.5, 0.8, 0.4, 45)
    add_custom_arrow(slide, 6.1, 5.9, 0.8, 0.4, 90)
    add_custom_arrow(slide, 5.4, 5.5, 0.8, 0.4, 135)
    add_custom_arrow(slide, 4.8, 4.5, 0.8, 0.4, 180)

    # 标签文字
    labels = [
        ("辣油", 4.2, 2.8), ("臭！", 6.0, 2.0), ("酸笋", 7.8, 2.6),
        ("辣！", 8.6, 4.2), ("腐竹", 8.0, 5.6), ("螺蛳汤", 6.0, 6.5),
        ("臭味", 4.2, 5.8), ("爽！", 3.8, 4.2)
    ]
    for text, lx, ly in labels:
        tb = slide.shapes.add_textbox(Inches(lx), Inches(ly), Inches(1.2), Inches(0.6))
        p = tb.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = C_BLACK

    # BOOM 装饰字
    tb = slide.shapes.add_textbox(Inches(7.8), Inches(1.8), Inches(1.5), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = "BOOM!"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = C_RED
    tb.rotation = -15

    # --- 右侧：表情包九宫格 ---
    bullet3 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(2.2), Inches(0.15), Inches(0.15))
    bullet3.fill.solid(); bullet3.fill.fore_color.rgb = C_RED; bullet3.line.color.rgb = C_BLACK
    
    tb = slide.shapes.add_textbox(Inches(9.7), Inches(2.0), Inches(3.0), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = "搞笑记录："
    p.font.size = Pt(22)
    p.font.bold = True
    
    tb = slide.shapes.add_textbox(Inches(9.7), Inches(2.5), Inches(3.0), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = "表情包九宫格"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = C_ORANGE

    # 第一次吃 小插图
    tb = slide.shapes.add_textbox(Inches(11.5), Inches(1.8), Inches(1.5), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    p.text = "第一次吃！"
    p.font.size = Pt(12)
    p.font.bold = True
    
    bowl_small = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.8), Inches(2.2), Inches(0.8), Inches(0.4))
    bowl_small.fill.solid(); bowl_small.fill.fore_color.rgb = C_WHITE
    bowl_small.line.color.rgb = C_BLACK

    # 九宫格
    captions = ["OMG!", "上头了", "救命！", "辣哭", "真香！", "——", "爱了爱了", "？？", "想撞墙"]
    start_x, start_y = 9.5, 3.2
    cell_w, cell_h = 1.1, 1.3
    spacing = 0.15
    
    photo_colors = [
        RGBColor(255, 200, 200), RGBColor(200, 255, 200), RGBColor(200, 200, 255),
        RGBColor(255, 255, 200), RGBColor(255, 200, 255), RGBColor(200, 255, 255),
        RGBColor(240, 240, 240), RGBColor(255, 220, 180), RGBColor(180, 220, 255)
    ]
    
    for i in range(9):
        row, col = i // 3, i % 3
        x = start_x + col * (cell_w + spacing)
        y = start_y + row * (cell_h + spacing)
        
        # 外框
        frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(cell_w), Inches(cell_h))
        frame.fill.solid(); frame.fill.fore_color.rgb = C_WHITE
        frame.line.color.rgb = C_BLACK; frame.line.width = Pt(1)
        
        # 照片区
        photo = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 0.05), Inches(y + 0.05), Inches(cell_w - 0.1), Inches(cell_h - 0.4))
        photo.fill.solid(); photo.fill.fore_color.rgb = photo_colors[i]
        photo.line.color.rgb = C_BLACK; photo.line.width = Pt(1)
        
        # 底部文字
        tb = slide.shapes.add_textbox(Inches(x), Inches(y + cell_h - 0.35), Inches(cell_w), Inches(0.3))
        p = tb.text_frame.paragraphs[0]
        p.text = captions[i]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(12)
        p.font.bold = True

    # 胶带装饰
    def add_tape(slide, x, y, rotation):
        tape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.8), Inches(0.25))
        tape.fill.solid(); tape.fill.fore_color.rgb = RGBColor(240, 230, 210)
        tape.line.fill.background()
        tape.rotation = rotation
        
    add_tape(slide, 0.2, 2.3, -45)
    add_tape(slide, 2.5, 2.3, 30)
    add_tape(slide, 9.2, 3.0, -20)
    add_tape(slide, 12.8, 3.0, 45)

    return slide



# ── Main ──
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

s0 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_1(s0)
prs.save(OUTPUT_PATH)
