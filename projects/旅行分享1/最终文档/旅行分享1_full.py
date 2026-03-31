from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# ── Color Palette ──
BLUE_HEADER = RGBColor(0x5B, 0x9B, 0xD5)
BLUE_DARK   = RGBColor(0x4A, 0x86, 0xC8)
CYAN        = RGBColor(0x00, 0xBC, 0xD4)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x33, 0x33, 0x33)
GRAY_TEXT   = RGBColor(0x55, 0x55, 0x55)
GRAY_BAR    = RGBColor(0xB0, 0xBE, 0xC5)
RED         = RGBColor(0xE5, 0x39, 0x35)
GREEN       = RGBColor(0x43, 0xA0, 0x47)
ORANGE      = RGBColor(0xFF, 0x98, 0x00)

ICON_BG  = RGBColor(0xE3, 0xE8, 0xED)
ICON_FG  = RGBColor(0x54, 0x6E, 0x7A)
FONT_NAME = "Microsoft YaHei"
SLIDE_WIDTH = Inches(13.333)
HEADER_H    = Inches(0.75)
SUBTITLE_Y  = Inches(0.95)


def add_header_banner(slide, title_text, bg_color=None):
    if bg_color is None:
        bg_color = BLUE_HEADER
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, HEADER_H
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = bg_color
    banner.line.fill.background()
    tf = banner.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(26)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME


def add_subtitle(slide, text, left, top, width=Inches(12), font_size=Pt(18)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.font.name = FONT_NAME
    return txBox


def add_icon_box(slide, left, top, symbol, size=Inches(0.48)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ICON_BG
    shape.line.fill.background()
    shape.adjustments[0] = 0.25
    tf = shape.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = symbol
    p.font.size = Pt(18)
    p.font.color.rgb = ICON_FG
    p.font.bold = False
    return shape


def add_bullet_item(slide, left, top, symbol, label, description,
                    width=Inches(5.5), desc_size=Pt(13)):
    add_icon_box(slide, left, top, symbol)
    text_left = left + Inches(0.65)
    txBox = slide.shapes.add_textbox(text_left, top - Inches(0.02), width, Inches(0.65))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run_label = p.add_run()
    run_label.text = label + "\uff1a"
    run_label.font.size = Pt(14)
    run_label.font.color.rgb = BLACK
    run_label.font.bold = True
    run_label.font.name = FONT_NAME
    run_desc = p.add_run()
    run_desc.text = description
    run_desc.font.size = desc_size
    run_desc.font.color.rgb = GRAY_TEXT
    run_desc.font.bold = False
    run_desc.font.name = FONT_NAME
    return txBox


def add_conclusion_box(slide, left, top, width, text, font_size=Pt(13)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = BLACK
    run.font.bold = True
    run.font.name = FONT_NAME
    return txBox


def add_table(slide, left, top, width, height, rows, cols, data,
              header_color=None, col_widths=None):
    if header_color is None:
        header_color = BLUE_HEADER
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if data[r][c] is not None else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(11)
                    run.font.name = FONT_NAME
                    if r == 0:
                        run.font.color.rgb = WHITE
                        run.font.bold = True
                    else:
                        run.font.color.rgb = BLACK
                        run.font.bold = False
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else RGBColor(0xF5, 0xF5, 0xF5)
            cell.margin_left = Pt(5)
            cell.margin_right = Pt(5)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
    return table_shape


def add_bar_chart(slide, left, top, width, height,
                  categories, values, title="", bar_colors=None):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', values)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    chart.chart_style = 2
    plot = chart.plots[0]
    plot.gap_width = 100
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CYAN
    series.has_data_labels = True
    dl = series.data_labels
    dl.font.size = Pt(13)
    dl.font.bold = True
    dl.font.color.rgb = BLACK
    dl.number_format = '0.#'
    dl.show_value = True
    dl.label_position = XL_LABEL_POSITION.OUTSIDE_END
    if bar_colors:
        for i, color in enumerate(bar_colors):
            pt = series.points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = color
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(12)
    cat_axis.tick_labels.font.name = FONT_NAME
    cat_axis.major_tick_mark = 2
    cat_axis.format.line.fill.background()
    val_axis = chart.value_axis
    val_axis.visible = False
    val_axis.major_tick_mark = 2
    val_axis.format.line.fill.background()
    val_axis.major_gridlines.format.line.fill.background()
    if title:
        chart.has_title = True
        ct = chart.chart_title.text_frame.paragraphs[0]
        ct.text = title
        ct.font.size = Pt(14)
        ct.font.bold = True
        ct.font.name = FONT_NAME
    else:
        chart.has_title = False
    return chart_frame


def add_callout_label(slide, left, top, text, bg_color=None, font_size=Pt(11)):
    if bg_color is None:
        bg_color = CYAN
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.3), Inches(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME
    return shape


def add_data_card(slide, left, top, width, height, value, label,
                  value_color=None, bg_color=None):
    if value_color is None:
        value_color = CYAN
    if bg_color is None:
        bg_color = WHITE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(3)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run1 = p1.add_run()
    run1.text = str(value)
    run1.font.size = Pt(24)
    run1.font.color.rgb = value_color
    run1.font.bold = True
    run1.font.name = FONT_NAME
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(10)
    run2.font.color.rgb = GRAY_TEXT
    run2.font.bold = False
    run2.font.name = FONT_NAME
    return shape


OUTPUT_PATH = r"C:\Users\Administrator\Desktop\notebook\my_notebooklm_with_ppteditable\projects\旅行分享1\最终文档\旅行分享1.pptx"


# ── Slide 1 ──

def build_slide_1(slide):
    # --- Colors ---
    BG_COLOR = RGBColor(245, 242, 235)
    TEXT_BLACK = RGBColor(20, 20, 20)
    ORANGE = RGBColor(244, 122, 32)
    GREEN = RGBColor(0, 168, 112)
    DARK_BLUE = RGBColor(20, 40, 80)
    LIGHT_TEAL = RGBColor(160, 220, 200)
    WHITE = RGBColor(255, 255, 255)

    # Set background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

    # --- Edge Decorations (Torn Paper Effect) ---
    # Top Left
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(-1.5), Inches(-1.5), Inches(3), Inches(3))
    shape.rotation = 45
    shape.fill.solid()
    shape.fill.fore_color.rgb = GREEN
    shape.line.fill.background()

    # Top Right
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(11.5), Inches(-1), Inches(3), Inches(2))
    shape.rotation = -20
    shape.fill.solid()
    shape.fill.fore_color.rgb = ORANGE
    shape.line.fill.background()
    
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(12), Inches(0), Inches(2), Inches(3))
    shape.rotation = 10
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    # Bottom Left
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(-1), Inches(6.5), Inches(3), Inches(2))
    shape.rotation = 15
    shape.fill.solid()
    shape.fill.fore_color.rgb = ORANGE
    shape.line.fill.background()

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(-0.5), Inches(7), Inches(3), Inches(2))
    shape.rotation = -10
    shape.fill.solid()
    shape.fill.fore_color.rgb = TEXT_BLACK
    shape.line.fill.background()

    # --- Main Title ---
    # Line 1
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(8), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "桂林山水“甲”天下，"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.name = "Microsoft YaHei"
    p.font.color.rgb = TEXT_BLACK

    # Highlight under Line 2
    highlight = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.6), Inches(6.2), Inches(0.25))
    highlight.fill.solid()
    highlight.fill.fore_color.rgb = ORANGE
    highlight.line.fill.background()

    # Line 2
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(8), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "我的脑洞“假”不了"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.name = "Microsoft YaHei"
    p.font.color.rgb = TEXT_BLACK

    # --- Subtitle ---
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(8), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "一个猎奇博主的桂林“真香”探险报告"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.name = "Microsoft YaHei"
    p.font.color.rgb = TEXT_BLACK

    # --- 20 RMB Graphic ---
    rmb_left = Inches(8.0)
    rmb_top = Inches(0.6)
    
    # Base
    rmb_base = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, rmb_left, rmb_top, Inches(4.8), Inches(2.5))
    rmb_base.fill.solid()
    rmb_base.fill.fore_color.rgb = WHITE
    rmb_base.line.color.rgb = DARK_BLUE
    rmb_base.line.width = Pt(2)

    # Inner fill
    rmb_inner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, rmb_left + Inches(0.1), rmb_top + Inches(0.1), Inches(4.6), Inches(2.3))
    rmb_inner.fill.solid()
    rmb_inner.fill.fore_color.rgb = LIGHT_TEAL
    rmb_inner.line.fill.background()

    # Sun
    sun = slide.shapes.add_shape(MSO_SHAPE.OVAL, rmb_left + Inches(2.0), rmb_top + Inches(0.3), Inches(0.6), Inches(0.6))
    sun.fill.solid()
    sun.fill.fore_color.rgb = WHITE
    sun.line.fill.background()

    # Mountains (Triangles)
    m1 = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, rmb_left + Inches(0.5), rmb_top + Inches(0.8), Inches(1.0), Inches(1.2))
    m1.fill.solid()
    m1.fill.fore_color.rgb = DARK_BLUE
    m1.line.fill.background()

    m2 = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, rmb_left + Inches(1.2), rmb_top + Inches(0.6), Inches(1.2), Inches(1.4))
    m2.fill.solid()
    m2.fill.fore_color.rgb = GREEN
    m2.line.fill.background()

    m3 = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, rmb_left + Inches(2.5), rmb_top + Inches(0.9), Inches(1.0), Inches(1.1))
    m3.fill.solid()
    m3.fill.fore_color.rgb = GREEN
    m3.line.fill.background()

    # River (Rectangle at bottom)
    river = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, rmb_left + Inches(0.1), rmb_top + Inches(1.8), Inches(4.6), Inches(0.6))
    river.fill.solid()
    river.fill.fore_color.rgb = DARK_BLUE
    river.line.fill.background()

    # Text "20" Top Left
    txBox = slide.shapes.add_textbox(rmb_left + Inches(0.2), rmb_top + Inches(0.1), Inches(1), Inches(0.5))
    p = txBox.text_frame.paragraphs[0]
    p.text = "20"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Text "20 RMB" Bottom Left
    txBox = slide.shapes.add_textbox(rmb_left + Inches(0.2), rmb_top + Inches(1.8), Inches(1.5), Inches(0.5))
    p = txBox.text_frame.paragraphs[0]
    p.text = "20 RMB"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Text "中国人民银行" Top Right
    txBox = slide.shapes.add_textbox(rmb_left + Inches(2.5), rmb_top + Inches(0.1), Inches(2), Inches(0.4))
    p = txBox.text_frame.paragraphs[0]
    p.text = "中国人民银行"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # --- Labels & Character (Bottom Left) ---
    # Character Placeholder
    char_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.0), Inches(4.5), Inches(1.5), Inches(2.5))
    char_bg.fill.solid()
    char_bg.fill.fore_color.rgb = ORANGE
    char_bg.line.color.rgb = TEXT_BLACK
    char_bg.line.width = Pt(2)
    p = char_bg.text_frame.paragraphs[0]
    p.text = "猎奇\n博主"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    def add_label(text, left, top, width, height, rotation, border_color):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.rotation = rotation
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = border_color
        shape.line.width = Pt(2.5)
        
        tf = shape.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.name = "Microsoft YaHei"
        p.font.color.rgb = TEXT_BLACK
        p.alignment = PP_ALIGN.CENTER
        tf.margin_top = Pt(3)
        tf.margin_bottom = Pt(3)
        tf.margin_left = Pt(5)
        tf.margin_right = Pt(5)

    add_label("奇葩角落", Inches(1.2), Inches(4.0), Inches(1.4), Inches(0.4), -15, GREEN)
    add_label("真香！！", Inches(1.2), Inches(5.0), Inches(1.4), Inches(0.4), -5, ORANGE)
    add_label("不错的青", Inches(1.5), Inches(6.0), Inches(1.4), Inches(0.4), 10, DARK_BLUE)
    
    add_label("拒绝中老年团", Inches(4.8), Inches(4.2), Inches(2.0), Inches(0.4), -15, ORANGE)
    add_label("冷知识", Inches(5.2), Inches(5.2), Inches(1.2), Inches(0.4), 5, GREEN)
    add_label("美景背后", Inches(5.0), Inches(6.0), Inches(1.4), Inches(0.4), 15, DARK_BLUE)

    # Add some arrows/lines pointing to center
    line1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(2.5), Inches(4.2), Inches(2.9), Inches(4.8))
    line1.line.color.rgb = TEXT_BLACK
    line1.line.width = Pt(2)
    
    line2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(4.8), Inches(4.6), Inches(4.5), Inches(5.0))
    line2.line.color.rgb = TEXT_BLACK
    line2.line.width = Pt(2)

    # --- Right List ---
    def add_list_item(text, left, top, bg_color, icon_border_color, icon_type):
        # Background Parallelogram
        bg = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, left + Inches(0.4), top, Inches(4.5), Inches(0.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = bg_color
        bg.line.fill.background()

        # Text
        txBox = slide.shapes.add_textbox(left + Inches(0.8), top - Inches(0.05), Inches(4.0), Inches(0.6))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.name = "Microsoft YaHei"
        p.font.color.rgb = TEXT_BLACK

        # Icon Circle
        icon_bg = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top - Inches(0.1), Inches(0.7), Inches(0.7))
        icon_bg.fill.solid()
        icon_bg.fill.fore_color.rgb = WHITE
        icon_bg.line.color.rgb = icon_border_color
        icon_bg.line.width = Pt(2.5)

        # Simple icon drawing inside the circle
        if icon_type == 'bus':
            bus = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.15), top + Inches(0.1), Inches(0.4), Inches(0.3))
            bus.fill.solid()
            bus.fill.fore_color.rgb = DARK_BLUE
            bus.line.fill.background()
            cross = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left + Inches(0.1), top, left + Inches(0.6), top + Inches(0.5))
            cross.line.color.rgb = ORANGE
            cross.line.width = Pt(3)
        elif icon_type == 'mag':
            mag_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.15), top + Inches(0.05), Inches(0.3), Inches(0.3))
            mag_circle.fill.background()
            mag_circle.line.color.rgb = GREEN
            mag_circle.line.width = Pt(2)
            handle = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left + Inches(0.4), top + Inches(0.3), left + Inches(0.55), top + Inches(0.45))
            handle.line.color.rgb = ORANGE
            handle.line.width = Pt(3)
        elif icon_type == 'person':
            head = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.25), top + Inches(0.05), Inches(0.2), Inches(0.2))
            head.fill.solid()
            head.fill.fore_color.rgb = ORANGE
            head.line.fill.background()
            body = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.15), top + Inches(0.3), Inches(0.4), Inches(0.2))
            body.fill.solid()
            body.fill.fore_color.rgb = GREEN
            body.line.fill.background()

    add_list_item("拒绝传统中老年旅行团画风", Inches(7.8), Inches(4.0), GREEN, ORANGE, 'bus')
    add_list_item("深度挖掘桂林不为人知的“奇葩”角落", Inches(7.8), Inches(5.0), GREEN, ORANGE, 'mag')
    add_list_item("搞笑博主的生存视角：\n美景背后的冷知识", Inches(7.8), Inches(6.0), ORANGE, GREEN, 'person')

    # --- Page Number ---
    # Green background circle
    circle_bg = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(12.1), Inches(6.8), Inches(0.4), Inches(0.4))
    circle_bg.fill.solid()
    circle_bg.fill.fore_color.rgb = GREEN
    circle_bg.line.fill.background()

    # Orange Circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(12.2), Inches(6.8), Inches(0.4), Inches(0.4))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ORANGE
    circle.line.fill.background()

    # Text
    txBox = slide.shapes.add_textbox(Inches(12.4), Inches(6.7), Inches(0.6), Inches(0.6))
    p = txBox.text_frame.paragraphs[0]
    p.text = "01"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.name = "Microsoft YaHei"
    p.font.color.rgb = TEXT_BLACK



# ── Slide 2 ──

def build_slide_2(slide):
    # 1. 设置背景颜色 (浅灰白)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(248, 248, 248)

    # 右侧边缘装饰条 (渐变橙/绿效果的简化)
    edge = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(13.1), Inches(0), Inches(0.233), Inches(7.5))
    edge.fill.solid()
    edge.fill.fore_color.rgb = RGBColor(230, 120, 50)
    edge.line.fill.background()

    # 2. 标题与副标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(10), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "角色设定：谁在带你逛桂林？"
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = RGBColor(10, 10, 10)

    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "颜值不够，脑洞来凑"
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(20, 20, 20)

    # 页码
    pg_box = slide.shapes.add_textbox(Inches(11.5), Inches(0.4), Inches(1.5), Inches(0.6))
    tf = pg_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "第2页"
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(22)
    run.font.bold = True

    # 3. 左侧视觉区域 (橙色泼墨背景 + 人物抠图占位)
    splash_color = RGBColor(255, 102, 0)
    
    # 使用多个椭圆组合模拟泼墨形状
    s1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), Inches(2.8), Inches(3.8), Inches(3.8))
    s1.fill.solid(); s1.fill.fore_color.rgb = splash_color; s1.line.fill.background()
    s2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), Inches(4.0), Inches(2.0), Inches(2.0))
    s2.fill.solid(); s2.fill.fore_color.rgb = splash_color; s2.line.fill.background()
    s3 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.5), Inches(2.5), Inches(2.0), Inches(2.0))
    s3.fill.solid(); s3.fill.fore_color.rgb = splash_color; s3.line.fill.background()
    s4 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.2), Inches(5.2), Inches(2.2), Inches(1.8))
    s4.fill.solid(); s4.fill.fore_color.rgb = splash_color; s4.line.fill.background()

    # 人物抠图占位符 (带白色粗边框的圆角矩形)
    cutout = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.8), Inches(2.8), Inches(2.4), Inches(3.8))
    cutout.fill.solid()
    cutout.fill.fore_color.rgb = RGBColor(30, 30, 30) # 深色衣服
    cutout.line.color.rgb = RGBColor(255, 255, 255)
    cutout.line.width = Pt(6)
    
    # 脸部占位
    face = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.4), Inches(3.0), Inches(1.2), Inches(1.5))
    face.fill.solid(); face.fill.fore_color.rgb = RGBColor(255, 218, 185); face.line.fill.background()
    
    # 相机占位
    cam = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.5), Inches(4.5), Inches(1.0), Inches(0.7))
    cam.fill.solid(); cam.fill.fore_color.rgb = RGBColor(50, 50, 50); cam.line.fill.background()
    lens = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.65), Inches(4.55), Inches(0.7), Inches(0.6))
    lens.fill.solid(); lens.fill.fore_color.rgb = RGBColor(10, 10, 10); lens.line.color.rgb = RGBColor(100, 100, 100)

    # "粗糙抠图" 文本
    cutout_txt = slide.shapes.add_textbox(Inches(2.2), Inches(6.7), Inches(2.0), Inches(0.5))
    tf = cutout_txt.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "粗糙抠图"
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(14)
    run.font.bold = True

    # 装饰性英文文本
    def add_deco_text(text, left, top, rotation, color, size=22):
        box = slide.shapes.add_textbox(left, top, Inches(1.5), Inches(0.8))
        box.rotation = rotation
        tf = box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.name = "Arial Black"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = True

    add_deco_text("OMG!", Inches(1.0), Inches(3.0), -15, RGBColor(0, 160, 140))
    add_deco_text("!?", Inches(3.5), Inches(2.8), 15, RGBColor(255, 255, 255), size=28)
    add_deco_text("BOOM!", Inches(3.8), Inches(3.8), 15, RGBColor(255, 255, 255))
    add_deco_text("LOOK!", Inches(4.2), Inches(5.5), 20, RGBColor(255, 255, 255))
    add_deco_text("NEW!", Inches(0.8), Inches(6.2), -10, RGBColor(0, 160, 140))
    add_deco_text("NEW!", Inches(11.5), Inches(2.0), 20, RGBColor(255, 112, 0))

    # 4. 右侧内容区域
    def add_content_block(left, top, width, title, desc):
        box = slide.shapes.add_textbox(left, top, width, Inches(1.5))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.line_spacing = 1.3
        
        run1 = p.add_run()
        run1.text = title
        run1.font.name = "Microsoft YaHei"
        run1.font.size = Pt(20)
        run1.font.bold = True
        run1.font.color.rgb = RGBColor(0, 0, 0)
        
        run2 = p.add_run()
        run2.text = desc
        run2.font.name = "Microsoft YaHei"
        run2.font.size = Pt(20)
        run2.font.bold = True
        run2.font.color.rgb = RGBColor(0, 0, 0)

    # --- 条目 1: 身份 ---
    i1_left, i1_top = Inches(5.5), Inches(2.5)
    # 图标占位 (人物 + 地图)
    p_body = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, i1_left+Inches(0.2), i1_top+Inches(0.6), Inches(0.8), Inches(0.6))
    p_body.fill.solid(); p_body.fill.fore_color.rgb = RGBColor(220, 50, 50); p_body.line.color.rgb = RGBColor(0, 0, 0)
    p_head = slide.shapes.add_shape(MSO_SHAPE.OVAL, i1_left+Inches(0.35), i1_top+Inches(0.1), Inches(0.5), Inches(0.5))
    p_head.fill.solid(); p_head.fill.fore_color.rgb = RGBColor(255, 200, 180); p_head.line.color.rgb = RGBColor(0, 0, 0)
    
    map_s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, i1_left+Inches(1.1), i1_top+Inches(0.3), Inches(0.8), Inches(0.7))
    map_s.rotation = 10
    map_s.fill.solid(); map_s.fill.fore_color.rgb = RGBColor(150, 200, 150); map_s.line.color.rgb = RGBColor(0, 0, 0)
    
    # 地图上的红叉
    x1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, i1_left+Inches(1.2), i1_top+Inches(0.4), i1_left+Inches(1.8), i1_top+Inches(0.9))
    x1.line.color.rgb = RGBColor(255, 0, 0); x1.line.width = Pt(3)
    x2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, i1_left+Inches(1.8), i1_top+Inches(0.4), i1_left+Inches(1.2), i1_top+Inches(0.9))
    x2.line.color.rgb = RGBColor(255, 0, 0); x2.line.width = Pt(3)

    # allergy 文本
    alg_box = slide.shapes.add_textbox(i1_left+Inches(1.0), i1_top+Inches(1.0), Inches(1.0), Inches(0.3))
    alg_box.rotation = -5
    tf = alg_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "allergy"
    run.font.name = "Arial"
    run.font.size = Pt(12)
    run.font.bold = True
    
    add_content_block(Inches(7.8), Inches(2.5), Inches(5.0), "身份：", "一个对“正常景点”\n过敏的猎奇博主")

    # --- 条目 2: 装备 ---
    add_content_block(Inches(5.8), Inches(4.2), Inches(4.5), "装备：", "自拍杆、扩音器、\n以及随时准备跑路的运动鞋")
    
    i2_left, i2_top = Inches(10.2), Inches(4.0)
    # 图标占位 (紫色背景 + 喇叭 + 鞋子)
    bg2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, i2_left, i2_top, Inches(1.2), Inches(1.2))
    bg2.fill.solid(); bg2.fill.fore_color.rgb = RGBColor(160, 32, 240); bg2.line.fill.background()
    
    mega = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, i2_left+Inches(0.6), i2_top-Inches(0.2), Inches(0.6), Inches(0.8))
    mega.rotation = 90
    mega.fill.solid(); mega.fill.fore_color.rgb = RGBColor(255, 255, 255); mega.line.color.rgb = RGBColor(160, 32, 240); mega.line.width = Pt(2)
    
    shoe = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, i2_left+Inches(1.0), i2_top+Inches(0.8), Inches(1.0), Inches(0.4))
    shoe.rotation = -15
    shoe.fill.solid(); shoe.fill.fore_color.rgb = RGBColor(255, 150, 0); shoe.line.color.rgb = RGBColor(0, 0, 0)
    
    add_deco_text("AHA!", i2_left+Inches(1.2), i2_top-Inches(0.5), -10, RGBColor(160, 32, 240))

    # --- 条目 3: 目标 ---
    add_content_block(Inches(5.8), Inches(6.0), Inches(3.5), "目标：\n", "寻找桂林最“野”的打开\n方式")
    
    i3_left, i3_top = Inches(9.5), Inches(5.8)
    # 图标占位 (青色背景 + 山水 + 指南针)
    bg3 = slide.shapes.add_shape(MSO_SHAPE.OVAL, i3_left+Inches(0.5), i3_top, Inches(1.4), Inches(1.4))
    bg3.fill.solid(); bg3.fill.fore_color.rgb = RGBColor(0, 160, 140); bg3.line.fill.background()
    
    m1 = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, i3_left+Inches(0.6), i3_top+Inches(0.4), Inches(0.6), Inches(0.8))
    m1.fill.solid(); m1.fill.fore_color.rgb = RGBColor(20, 20, 20); m1.line.fill.background()
    m2 = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, i3_left+Inches(1.0), i3_top+Inches(0.2), Inches(0.8), Inches(1.0))
    m2.fill.solid(); m2.fill.fore_color.rgb = RGBColor(10, 10, 10); m2.line.fill.background()
    
    comp = slide.shapes.add_shape(MSO_SHAPE.OVAL, i3_left, i3_top+Inches(0.5), Inches(0.8), Inches(0.8))
    comp.fill.solid(); comp.fill.fore_color.rgb = RGBColor(255, 220, 150); comp.line.color.rgb = RGBColor(0, 0, 0); comp.line.width = Pt(2)
    
    add_deco_text("WILD?", i3_left+Inches(1.5), i3_top-Inches(0.4), 15, RGBColor(0, 0, 0))

    # 5. 添加手绘风格的连接线/箭头指示
    a1 = slide.shapes.add_connector(MSO_CONNECTOR.CURVE, Inches(7.3), Inches(3.5), Inches(7.7), Inches(3.4))
    a1.line.color.rgb = RGBColor(150, 150, 150)
    
    a2 = slide.shapes.add_connector(MSO_CONNECTOR.CURVE, Inches(9.5), Inches(5.0), Inches(10.0), Inches(4.8))
    a2.line.color.rgb = RGBColor(150, 150, 150)
    
    a3 = slide.shapes.add_connector(MSO_CONNECTOR.CURVE, Inches(8.8), Inches(6.8), Inches(9.3), Inches(6.5))
    a3.line.color.rgb = RGBColor(150, 150, 150)



# ── Slide 3 ──

def build_slide_3(slide):
    # Colors
    BLACK = RGBColor(0x00, 0x00, 0x00)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    GREEN_LBL = RGBColor(0x8A, 0xE0, 0xA1)
    ORANGE_LBL = RGBColor(0xFF, 0xB3, 0x47)
    GREEN_BORDER = RGBColor(0x4C, 0xAF, 0x50)
    ORANGE_BORDER = RGBColor(0xFF, 0x57, 0x22)
    PURPLE = RGBColor(0x9C, 0x27, 0xB0)
    GRAY_TEXT = RGBColor(0x75, 0x75, 0x75)

    # 1. Title
    title_box = slide.shapes.add_textbox(Inches(1.5), Inches(0.4), Inches(10.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    
    run1 = p.add_run()
    run1.text = "20元人民币打卡： "
    run1.font.size = Pt(36)
    run1.font.bold = True
    run1.font.name = "Microsoft YaHei"
    
    run2 = p.add_run()
    run2.text = "买家秀"
    run2.font.size = Pt(36)
    run2.font.bold = True
    run2.font.name = "Microsoft YaHei"
    
    run3 = p.add_run()
    run3.text = " vs "
    run3.font.size = Pt(36)
    run3.font.bold = True
    run3.font.name = "Microsoft YaHei"
    
    run4 = p.add_run()
    run4.text = "卖家秀"
    run4.font.size = Pt(36)
    run4.font.bold = True
    run4.font.name = "Microsoft YaHei"

    # Purple underline for "卖家秀" (approximate position)
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(9.3), Inches(1.15), Inches(10.8), Inches(1.15))
    line.line.color.rgb = PURPLE
    line.line.width = Pt(4)

    # 2. Subtitle
    sub_box = slide.shapes.add_textbox(Inches(3), Inches(1.1), Inches(7.333), Inches(0.6))
    tf_sub = sub_box.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.alignment = PP_ALIGN.CENTER
    run_sub = p_sub.add_run()
    run_sub.text = "理想很丰满，现实很骨感"
    run_sub.font.size = Pt(22)
    run_sub.font.bold = True
    run_sub.font.name = "Microsoft YaHei"

    # 3. Center Divider (Lightning Bolt)
    lightning = slide.shapes.add_shape(MSO_SHAPE.LIGHTNING_BOLT, Inches(6.3), Inches(1.8), Inches(0.8), Inches(5.5))
    lightning.fill.solid()
    lightning.fill.fore_color.rgb = ORANGE_BORDER
    lightning.line.fill.background()

    # --- LEFT SIDE (IDEAL) ---

    # Label: 理想 (Ideal)
    lbl_ideal = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.3), Inches(1.8), Inches(2.0), Inches(0.5))
    lbl_ideal.fill.solid()
    lbl_ideal.fill.fore_color.rgb = GREEN_LBL
    lbl_ideal.line.fill.background()
    tf_ideal = lbl_ideal.text_frame
    tf_ideal.text = "理想 (Ideal)"
    tf_ideal.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf_ideal.paragraphs[0].font.size = Pt(18)
    tf_ideal.paragraphs[0].font.bold = True
    tf_ideal.paragraphs[0].font.color.rgb = BLACK
    tf_ideal.paragraphs[0].font.name = "Microsoft YaHei"

    # Image Placeholder: 20 RMB
    img_ideal = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.4), Inches(5.0), Inches(2.6))
    img_ideal.fill.solid()
    img_ideal.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
    img_ideal.line.fill.solid()
    img_ideal.line.fore_color.rgb = GREEN_BORDER
    img_ideal.line.width = Pt(3)
    tf_img1 = img_ideal.text_frame
    tf_img1.text = "[20元人民币风景图]"
    tf_img1.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf_img1.paragraphs[0].font.color.rgb = RGBColor(0x9E, 0x9E, 0x9E)

    # Photographer Placeholder
    photo_guy = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.8), Inches(5.2), Inches(1.5), Inches(1.8))
    photo_guy.fill.solid()
    photo_guy.fill.fore_color.rgb = RGBColor(0xBB, 0xDE, 0xFB)
    photo_guy.line.fill.solid()
    photo_guy.line.fore_color.rgb = BLACK
    photo_guy.line.width = Pt(2)
    tf_guy = photo_guy.text_frame
    tf_guy.text = "摄影师\n(插画)"
    tf_guy.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf_guy.paragraphs[0].font.color.rgb = BLACK

    # Speech Bubble: PERFECT SHOT!
    bubble1 = slide.shapes.add_shape(MSO_SHAPE.OVAL_CALLOUT, Inches(3.1), Inches(4.8), Inches(1.5), Inches(0.8))
    bubble1.fill.solid()
    bubble1.fill.fore_color.rgb = WHITE
    bubble1.line.fill.solid()
    bubble1.line.fore_color.rgb = BLACK
    bubble1.line.width = Pt(2)
    tf_b1 = bubble1.text_frame
    tf_b1.text = "PERFECT\nSHOT!"
    tf_b1.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf_b1.paragraphs[0].font.size = Pt(11)
    tf_b1.paragraphs[0].font.bold = True
    tf_b1.paragraphs[0].font.color.rgb = BLACK

    # Conclusion Box
    box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(6.4), Inches(2.8), Inches(0.8))
    box1.fill.solid()
    box1.fill.fore_color.rgb = WHITE
    box1.line.fill.solid()
    box1.line.fore_color.rgb = BLACK
    box1.line.width = Pt(2)
    tf_box1 = box1.text_frame
    tf_box1.text = "结论：找准角度，你就是\n人民币上的那个男人/女人"
    tf_box1.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf_box1.paragraphs[0].font.size = Pt(11)
    tf_box1.paragraphs[0].font.bold = True
    tf_box1.paragraphs[0].font.color.rgb = BLACK
    tf_box1.paragraphs[0].font.name = "Microsoft YaHei"

    # Arrow to Conclusion Box
    arrow1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.1), Inches(6.7), Inches(0.4), Inches(0.2))
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = ORANGE_BORDER
    arrow1.line.fill.background()

    # --- RIGHT SIDE (REALITY) ---

    # Label: 现实 (Reality)
    lbl_reality = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.9), Inches(1.8), Inches(2.2), Inches(0.5))
    lbl_reality.fill.solid()
    lbl_reality.fill.fore_color.rgb = ORANGE_LBL
    lbl_reality.line.fill.background()
    tf_reality = lbl_reality.text_frame
    tf_reality.text = "现实 (Reality)"
    tf_reality.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf_reality.paragraphs[0].font.size = Pt(18)
    tf_reality.paragraphs[0].font.bold = True
    tf_reality.paragraphs[0].font.color.rgb = BLACK
    tf_reality.paragraphs[0].font.name = "Microsoft YaHei"

    # Image Placeholder: Crowd
    img_reality = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.5), Inches(2.4), Inches(5.0), Inches(3.0))
    img_reality.fill.solid()
    img_reality.fill.fore_color.rgb = RGBColor(0xFF, 0xE0, 0xB2)
    img_reality.line.fill.solid()
    img_reality.line.fore_color.rgb = ORANGE_BORDER
    img_reality.line.width = Pt(3)
    tf_img2 = img_reality.text_frame
    tf_img2.text = "[拥挤的人群举着20元拍照]"
    tf_img2.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf_img2.paragraphs[0].font.color.rgb = RGBColor(0x9E, 0x9E, 0x9E)

    # Stressed Guy Placeholder
    stress_guy = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.5), Inches(3.5), Inches(1.5), Inches(1.8))
    stress_guy.fill.solid()
    stress_guy.fill.fore_color.rgb = RGBColor(0xFF, 0xCD, 0xD2)
    stress_guy.line.fill.solid()
    stress_guy.line.fore_color.rgb = BLACK
    stress_guy.line.width = Pt(2)
    tf_guy2 = stress_guy.text_frame
    tf_guy2.text = "崩溃游客\n(插画)"
    tf_guy2.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf_guy2.paragraphs[0].font.color.rgb = BLACK

    # Speech Bubble: OMG! TOO MANY PEOPLE!
    bubble2 = slide.shapes.add_shape(MSO_SHAPE.CLOUD_CALLOUT, Inches(11.0), Inches(2.2), Inches(1.8), Inches(1.2))
    bubble2.fill.solid()
    bubble2.fill.fore_color.rgb = WHITE
    bubble2.line.fill.solid()
    bubble2.line.fore_color.rgb = BLACK
    bubble2.line.width = Pt(2)
    tf_b2 = bubble2.text_frame
    tf_b2.text = "OMG!\nTOO MANY\nPEOPLE!"
    tf_b2.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf_b2.paragraphs[0].font.size = Pt(11)
    tf_b2.paragraphs[0].font.bold = True
    tf_b2.paragraphs[0].font.color.rgb = BLACK

    # Reality Text Box
    box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.2), Inches(4.8), Inches(2.5), Inches(1.0))
    box2.fill.solid()
    box2.fill.fore_color.rgb = WHITE
    box2.line.fill.solid()
    box2.line.fore_color.rgb = BLACK
    box2.line.width = Pt(2)
    tf_box2 = box2.text_frame
    tf_box2.text = "现实：岸边挤满了100个\n同样拿着20块钱的人"
    tf_box2.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf_box2.paragraphs[0].font.size = Pt(11)
    tf_box2.paragraphs[0].font.bold = True
    tf_box2.paragraphs[0].font.color.rgb = BLACK
    tf_box2.paragraphs[0].font.name = "Microsoft YaHei"

    # Cormorant Placeholder
    bird = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.2), Inches(5.5), Inches(1.2), Inches(1.8))
    bird.fill.solid()
    bird.fill.fore_color.rgb = RGBColor(0xCF, 0xD8, 0xDC)
    bird.line.fill.solid()
    bird.line.fore_color.rgb = BLACK
    bird.line.width = Pt(2)
    tf_bird = bird.text_frame
    tf_bird.text = "鸬鹚\n(插画)"
    tf_bird.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf_bird.paragraphs[0].font.color.rgb = BLACK

    # RETIRED sign
    sign = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.3), Inches(6.5), Inches(1.0), Inches(0.4))
    sign.fill.solid()
    sign.fill.fore_color.rgb = RGBColor(0xFF, 0xF5, 0x9D)
    sign.line.fill.solid()
    sign.line.fore_color.rgb = BLACK
    sign.line.width = Pt(1)
    tf_sign = sign.text_frame
    tf_sign.text = "RETIRED"
    tf_sign.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf_sign.paragraphs[0].font.size = Pt(10)
    tf_sign.paragraphs[0].font.bold = True
    tf_sign.paragraphs[0].font.color.rgb = BLACK

    # Fact Text Box
    box3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.8), Inches(6.4), Inches(2.8), Inches(0.8))
    box3.fill.solid()
    box3.fill.fore_color.rgb = WHITE
    box3.line.fill.solid()
    box3.line.fore_color.rgb = BLACK
    box3.line.width = Pt(2)
    tf_box3 = box3.text_frame
    tf_box3.text = "猎奇点：那只配合拍照的\n鸬鹚其实已经“退休”了"
    tf_box3.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf_box3.paragraphs[0].font.size = Pt(11)
    tf_box3.paragraphs[0].font.bold = True
    tf_box3.paragraphs[0].font.color.rgb = BLACK
    tf_box3.paragraphs[0].font.name = "Microsoft YaHei"

    # Arrow to Fact Box
    arrow2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(9.3), Inches(6.7), Inches(0.4), Inches(0.2))
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = ORANGE_BORDER
    arrow2.line.fill.background()

    # 4. Page Number
    page_num = slide.shapes.add_textbox(Inches(12.5), Inches(7.0), Inches(0.8), Inches(0.4))
    tf_page = page_num.text_frame
    p_page = tf_page.paragraphs[0]
    p_page.text = "3 / 11"
    p_page.font.size = Pt(14)
    p_page.font.color.rgb = GRAY_TEXT
    p_page.alignment = PP_ALIGN.RIGHT



# ── Slide 4 ──

def build_slide_4(slide):
    # 1. 添加左上角页码
    tx_box_page = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(2.0), Inches(0.8))
    tf_page = tx_box_page.text_frame
    p_page = tf_page.paragraphs[0]
    p_page.text = "第4页"
    p_page.font.name = "Microsoft YaHei"
    p_page.font.size = Pt(32)
    p_page.font.bold = True
    p_page.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # 2. 添加主标题
    tx_box_title = slide.shapes.add_textbox(Inches(3.5), Inches(0.6), Inches(9.0), Inches(0.8))
    tf_title = tx_box_title.text_frame
    p_title = tf_title.paragraphs[0]
    p_title.text = "象鼻山：这头大象到底在喝什么？"
    p_title.font.name = "Microsoft YaHei"
    p_title.font.size = Pt(36)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # 3. 添加副标题
    tx_box_subtitle = slide.shapes.add_textbox(Inches(6.5), Inches(1.3), Inches(6.0), Inches(0.6))
    tf_subtitle = tx_box_subtitle.text_frame
    p_subtitle = tf_subtitle.paragraphs[0]
    p_subtitle.text = "关于桂林城徽的终极猜想"
    p_subtitle.font.name = "Microsoft YaHei"
    p_subtitle.font.size = Pt(24)
    p_subtitle.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
    p_subtitle.alignment = PP_ALIGN.RIGHT

    # 4. 左侧大象喝奶茶插图占位符及“奶茶？”文字
    # 插图占位符
    pic_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.5))
    pic_shape.fill.solid()
    pic_shape.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    pic_shape.line.fill.background()
    tf_pic = pic_shape.text_frame
    p_pic = tf_pic.paragraphs[0]
    p_pic.text = "[大象喝奶茶插图区域]"
    p_pic.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
    p_pic.alignment = PP_ALIGN.CENTER

    # “奶茶？”文字
    tx_box_tea = slide.shapes.add_textbox(Inches(5.2), Inches(4.8), Inches(1.5), Inches(0.8))
    tx_box_tea.rotation = -15
    tf_tea = tx_box_tea.text_frame
    p_tea = tf_tea.paragraphs[0]
    p_tea.text = "奶茶？"
    p_tea.font.name = "Microsoft YaHei"
    p_tea.font.size = Pt(28)
    p_tea.font.bold = True
    p_tea.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # 5. 右侧图文列表项
    # 列表项 1
    icon1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.8), Inches(2.6), Inches(0.8), Inches(0.8))
    icon1.fill.solid()
    icon1.fill.fore_color.rgb = RGBColor(0xE1, 0xF5, 0xFE)
    icon1.line.fill.background()
    icon1.text_frame.text = "📷"
    
    tx1 = slide.shapes.add_textbox(Inches(7.8), Inches(2.5), Inches(5.0), Inches(1.0))
    tf1 = tx1.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "视觉错位："
    p1.font.name = "Microsoft YaHei"
    p1.font.bold = True
    p1.font.size = Pt(18)
    p1.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    run1 = p1.add_run()
    run1.text = "从哪个角度看最像一只喝醉的大象"
    run1.font.bold = False
    run1.font.size = Pt(18)

    # 列表项 2
    icon2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.8), Inches(4.1), Inches(0.8), Inches(0.8))
    icon2.fill.solid()
    icon2.fill.fore_color.rgb = RGBColor(0xF3, 0xE5, 0xF5)
    icon2.line.fill.background()
    icon2.text_frame.text = "🍷"
    
    tx2 = slide.shapes.add_textbox(Inches(7.8), Inches(4.0), Inches(5.0), Inches(1.0))
    tf2 = tx2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "猎奇冷知识："
    p2.font.name = "Microsoft YaHei"
    p2.font.bold = True
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    run2 = p2.add_run()
    run2.text = "象鼻山内部其实是空的（藏酒洞）"
    run2.font.bold = False
    run2.font.size = Pt(18)

    # 列表项 3
    icon3 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.8), Inches(5.6), Inches(0.8), Inches(0.8))
    icon3.fill.solid()
    icon3.fill.fore_color.rgb = RGBColor(0xE0, 0xF2, 0xF1)
    icon3.line.fill.background()
    icon3.text_frame.text = "💦"
    
    tx3 = slide.shapes.add_textbox(Inches(7.8), Inches(5.5), Inches(5.0), Inches(1.0))
    tf3 = tx3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.text = "吐槽点："
    p3.font.name = "Microsoft YaHei"
    p3.font.bold = True
    p3.font.size = Pt(18)
    p3.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    run3 = p3.add_run()
    run3.text = "为了拍一张“喂象照”，我差点掉进江里"
    run3.font.bold = False
    run3.font.size = Pt(18)

    # 6. 底部右侧装饰元素 (BooM!, OMG!, 箭头)
    # BooM!
    boom_box = slide.shapes.add_textbox(Inches(8.5), Inches(6.8), Inches(1.2), Inches(0.6))
    boom_box.rotation = -10
    tf_boom = boom_box.text_frame
    p_boom = tf_boom.paragraphs[0]
    p_boom.text = "BooM!"
    p_boom.font.name = "Microsoft YaHei"
    p_boom.font.size = Pt(16)
    p_boom.font.bold = True
    p_boom.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # 青色箭头
    arrow1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(9.8), Inches(6.9), Inches(0.4), Inches(0.25))
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = RGBColor(0x00, 0xBC, 0xD4)
    arrow1.line.fill.background()

    # OMG!
    omg_box = slide.shapes.add_textbox(Inches(10.3), Inches(6.8), Inches(1.0), Inches(0.6))
    tf_omg = omg_box.text_frame
    p_omg = tf_omg.paragraphs[0]
    p_omg.text = "OMG!"
    p_omg.font.name = "Microsoft YaHei"
    p_omg.font.size = Pt(16)
    p_omg.font.bold = True
    p_omg.font.color.rgb = RGBColor(0x00, 0xBC, 0xD4)

    # 紫色箭头
    arrow2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(11.5), Inches(6.9), Inches(0.4), Inches(0.25))
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = RGBColor(0x9C, 0x27, 0xB0)
    arrow2.line.fill.background()

    # 绿色箭头
    arrow3 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(12.1), Inches(6.9), Inches(0.4), Inches(0.25))
    arrow3.fill.solid()
    arrow3.fill.fore_color.rgb = RGBColor(0x4C, 0xAF, 0x50)
    arrow3.line.fill.background()



# ── Slide 5 ──

def build_slide_5(slide):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
    
    # 定义颜色
    C_BG = RGBColor(253, 249, 238)
    C_RED = RGBColor(228, 57, 50)
    C_YELLOW = RGBColor(255, 214, 89)
    C_GREEN = RGBColor(0, 150, 100)
    C_ORANGE = RGBColor(244, 121, 32)
    C_BLACK = RGBColor(0, 0, 0)
    C_WHITE = RGBColor(255, 255, 255)

    # 设置背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_BG
    bg.line.fill.background()
    
    # 添加背景装饰色块
    blob1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(-2), Inches(5), Inches(5))
    blob1.fill.solid(); blob1.fill.fore_color.rgb = C_RED; blob1.line.fill.background()
    
    blob2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3), Inches(-2), Inches(6), Inches(4))
    blob2.fill.solid(); blob2.fill.fore_color.rgb = C_YELLOW; blob2.line.fill.background()
    
    blob3 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(-2), Inches(6), Inches(5))
    blob3.fill.solid(); blob3.fill.fore_color.rgb = C_GREEN; blob3.line.fill.background()
    
    blob4 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1), Inches(6), Inches(5), Inches(4))
    blob4.fill.solid(); blob4.fill.fore_color.rgb = C_GREEN; blob4.line.fill.background()
    
    blob5 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8), Inches(6), Inches(7), Inches(4))
    blob5.fill.solid(); blob5.fill.fore_color.rgb = C_ORANGE; blob5.line.fill.background()

    # 标题 (使用多重偏移模拟描边效果)
    title_text = "螺蛳粉：桂林的“生化武器”诱惑"
    offsets = [(-0.03, -0.03), (0.03, -0.03), (-0.03, 0.03), (0.03, 0.03), (-0.04, 0), (0.04, 0), (0, -0.04), (0, 0.04)]
    for ox, oy in offsets:
        tb = slide.shapes.add_textbox(Inches(0.5 + ox), Inches(0.3 + oy), Inches(10), Inches(1))
        p = tb.text_frame.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = C_BLACK
        p.font.name = "Microsoft YaHei"
        
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(10), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p.font.name = "Microsoft YaHei"

    # 副标题
    sub_text = "闻着臭，吃着爽，回味想撞墙"
    for ox, oy in offsets:
        tb = slide.shapes.add_textbox(Inches(0.5 + ox), Inches(1.1 + oy), Inches(8), Inches(0.8))
        p = tb.text_frame.paragraphs[0]
        p.text = sub_text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = C_BLACK
        p.font.name = "Microsoft YaHei"
        
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(8), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.text = sub_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p.font.name = "Microsoft YaHei"

    # 页码
    page_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.8), Inches(0.4), Inches(1.0), Inches(1.0))
    page_circle.fill.solid(); page_circle.fill.fore_color.rgb = C_BG
    page_circle.line.color.rgb = C_BLACK; page_circle.line.width = Pt(2)
    
    page_inner = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.85), Inches(0.45), Inches(0.9), Inches(0.9))
    page_inner.fill.background()
    page_inner.line.color.rgb = C_BLACK; page_inner.line.width = Pt(1)
    
    tb = slide.shapes.add_textbox(Inches(11.8), Inches(0.5), Inches(1.0), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.text = "PAGE\n5"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = C_BLACK

    # --- 左侧：内容要点 ---
    shadow_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(2.6), Inches(2.0), Inches(0.6))
    shadow_box.fill.solid(); shadow_box.fill.fore_color.rgb = C_GREEN
    shadow_box.line.color.rgb = C_BLACK; shadow_box.line.width = Pt(1.5)
    
    main_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.5), Inches(2.0), Inches(0.6))
    main_box.fill.solid(); main_box.fill.fore_color.rgb = C_YELLOW
    main_box.line.color.rgb = C_BLACK; main_box.line.width = Pt(1.5)
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.55), Inches(2.0), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = "内容要点"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = C_BLACK

    # 要点 1
    bullet1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(3.6), Inches(0.15), Inches(0.15))
    bullet1.fill.solid(); bullet1.fill.fore_color.rgb = C_RED; bullet1.line.color.rgb = C_BLACK
    
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(3.4), Inches(3.0), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = "气味等级："
    p.font.size = Pt(20)
    p.font.bold = True
    
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(3.9), Inches(3.0), Inches(1.0))
    p = tb.text_frame.paragraphs[0]
    p.text = "路过店门口，衣服\n自动“腌制”三天"
    p.font.size = Pt(18)
    
    # 腌制 圈注
    oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.5), Inches(4.2), Inches(0.9), Inches(0.4))
    oval.fill.background()
    oval.line.color.rgb = C_ORANGE; oval.line.width = Pt(2)

    # 要点 2
    bullet2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(5.4), Inches(0.15), Inches(0.15))
    bullet2.fill.solid(); bullet2.fill.fore_color.rgb = C_RED; bullet2.line.color.rgb = C_BLACK
    
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(5.2), Inches(3.0), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = "猎奇吃法："
    p.font.size = Pt(20)
    p.font.bold = True
    
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(5.7), Inches(3.0), Inches(1.0))
    p = tb.text_frame.paragraphs[0]
    run1 = p.add_run()
    run1.text = "加辣加酸笋"
    run1.font.size = Pt(18)
    run1.font.underline = True
    run2 = p.add_run()
    run2.text = "，\n挑战味蕾极限"
    run2.font.size = Pt(18)

    # --- 中间：爆炸图与杀伤力 ---
    exp_red = slide.shapes.add_shape(MSO_SHAPE.EXPLOSION2, Inches(3.8), Inches(2.2), Inches(5.5), Inches(5.0))
    exp_red.fill.solid(); exp_red.fill.fore_color.rgb = C_RED
    exp_red.line.color.rgb = C_BLACK; exp_red.line.width = Pt(2)
    
    exp_yellow = slide.shapes.add_shape(MSO_SHAPE.EXPLOSION1, Inches(4.3), Inches(2.7), Inches(4.5), Inches(4.0))
    exp_yellow.fill.solid(); exp_yellow.fill.fore_color.rgb = C_YELLOW
    exp_yellow.line.color.rgb = C_BLACK; exp_yellow.line.width = Pt(2)
    
    # 中心碗
    bowl_outer = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.4), Inches(3.8), Inches(2.2), Inches(1.6))
    bowl_outer.fill.solid(); bowl_outer.fill.fore_color.rgb = C_WHITE
    bowl_outer.line.color.rgb = C_BLACK; bowl_outer.line.width = Pt(2)
    
    bowl_inner = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.6), Inches(3.9), Inches(1.8), Inches(0.8))
    bowl_inner.fill.solid(); bowl_inner.fill.fore_color.rgb = C_ORANGE
    bowl_inner.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(5.4), Inches(4.7), Inches(2.2), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = "杀伤力"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)
    p.font.bold = True

    # 放射箭头
    def add_custom_arrow(slide, x, y, width, height, rotation):
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(width), Inches(height))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = C_YELLOW
        arrow.line.color.rgb = C_RED; arrow.line.width = Pt(1.5)
        arrow.rotation = rotation

    add_custom_arrow(slide, 5.4, 3.4, 0.8, 0.4, 225)
    add_custom_arrow(slide, 6.1, 3.0, 0.8, 0.4, 270)
    add_custom_arrow(slide, 7.1, 3.4, 0.8, 0.4, 315)
    add_custom_arrow(slide, 7.6, 4.5, 0.8, 0.4, 0)
    add_custom_arrow(slide, 7.1, 5.5, 0.8, 0.4, 45)
    add_custom_arrow(slide, 6.1, 5.9, 0.8, 0.4, 90)
    add_custom_arrow(slide, 5.4, 5.5, 0.8, 0.4, 135)
    add_custom_arrow(slide, 4.8, 4.5, 0.8, 0.4, 180)

    # 标签文字
    labels = [
        ("辣油", 4.2, 2.8), ("臭！", 6.0, 2.0), ("酸笋", 7.8, 2.6),
        ("辣！", 8.6, 4.2), ("腐竹", 8.0, 5.6), ("螺蛳汤", 6.0, 6.5),
        ("臭味", 4.2, 5.8), ("爽！", 3.8, 4.2)
    ]
    for text, lx, ly in labels:
        tb = slide.shapes.add_textbox(Inches(lx), Inches(ly), Inches(1.2), Inches(0.6))
        p = tb.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = C_BLACK

    # BOOM 装饰字
    tb = slide.shapes.add_textbox(Inches(7.8), Inches(1.8), Inches(1.5), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = "BOOM!"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = C_RED
    tb.rotation = -15

    # --- 右侧：表情包九宫格 ---
    bullet3 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(2.2), Inches(0.15), Inches(0.15))
    bullet3.fill.solid(); bullet3.fill.fore_color.rgb = C_RED; bullet3.line.color.rgb = C_BLACK
    
    tb = slide.shapes.add_textbox(Inches(9.7), Inches(2.0), Inches(3.0), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = "搞笑记录："
    p.font.size = Pt(22)
    p.font.bold = True
    
    tb = slide.shapes.add_textbox(Inches(9.7), Inches(2.5), Inches(3.0), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = "表情包九宫格"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = C_ORANGE

    # 第一次吃 小插图
    tb = slide.shapes.add_textbox(Inches(11.5), Inches(1.8), Inches(1.5), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    p.text = "第一次吃！"
    p.font.size = Pt(12)
    p.font.bold = True
    
    bowl_small = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.8), Inches(2.2), Inches(0.8), Inches(0.4))
    bowl_small.fill.solid(); bowl_small.fill.fore_color.rgb = C_WHITE
    bowl_small.line.color.rgb = C_BLACK

    # 九宫格
    captions = ["OMG!", "上头了", "救命！", "辣哭", "真香！", "——", "爱了爱了", "？？", "想撞墙"]
    start_x, start_y = 9.5, 3.2
    cell_w, cell_h = 1.1, 1.3
    spacing = 0.15
    
    photo_colors = [
        RGBColor(255, 200, 200), RGBColor(200, 255, 200), RGBColor(200, 200, 255),
        RGBColor(255, 255, 200), RGBColor(255, 200, 255), RGBColor(200, 255, 255),
        RGBColor(240, 240, 240), RGBColor(255, 220, 180), RGBColor(180, 220, 255)
    ]
    
    for i in range(9):
        row, col = i // 3, i % 3
        x = start_x + col * (cell_w + spacing)
        y = start_y + row * (cell_h + spacing)
        
        # 外框
        frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(cell_w), Inches(cell_h))
        frame.fill.solid(); frame.fill.fore_color.rgb = C_WHITE
        frame.line.color.rgb = C_BLACK; frame.line.width = Pt(1)
        
        # 照片区
        photo = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 0.05), Inches(y + 0.05), Inches(cell_w - 0.1), Inches(cell_h - 0.4))
        photo.fill.solid(); photo.fill.fore_color.rgb = photo_colors[i]
        photo.line.color.rgb = C_BLACK; photo.line.width = Pt(1)
        
        # 底部文字
        tb = slide.shapes.add_textbox(Inches(x), Inches(y + cell_h - 0.35), Inches(cell_w), Inches(0.3))
        p = tb.text_frame.paragraphs[0]
        p.text = captions[i]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(12)
        p.font.bold = True

    # 胶带装饰
    def add_tape(slide, x, y, rotation):
        tape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.8), Inches(0.25))
        tape.fill.solid(); tape.fill.fore_color.rgb = RGBColor(240, 230, 210)
        tape.line.fill.background()
        tape.rotation = rotation
        
    add_tape(slide, 0.2, 2.3, -45)
    add_tape(slide, 2.5, 2.3, 30)
    add_tape(slide, 9.2, 3.0, -20)
    add_tape(slide, 12.8, 3.0, 45)

    return slide



# ── Slide 6 ──

def build_slide_6(slide):
    # Colors
    BG_COLOR = RGBColor(248, 245, 238)
    TITLE_COLOR = RGBColor(142, 36, 170)
    TEXT_COLOR = RGBColor(30, 30, 30)
    
    ORANGE = RGBColor(255, 152, 0)
    GREEN = RGBColor(76, 175, 80)
    CYAN = RGBColor(0, 188, 212)
    PURPLE = RGBColor(156, 39, 176)
    LIGHT_BLUE = RGBColor(225, 245, 254)
    YELLOW = RGBColor(255, 235, 59)

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_COLOR
    bg.line.fill.background()

    # --- Header ---
    # Main Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.2), Inches(11.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "芦笛岩：地底下的“80年代迪厅”"
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = TITLE_COLOR

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(2), Inches(1.2), Inches(9.333), Inches(0.6))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "这种审美真的很“硬核”"
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = TEXT_COLOR

    # --- Dividers ---
    line1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(4.5), Inches(2.0), Inches(4.5), Inches(6.5))
    line1.line.color.rgb = TEXT_COLOR
    line1.line.width = Pt(2)
    line1.line.dash_style = 7 

    line2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(8.8), Inches(2.0), Inches(8.8), Inches(6.5))
    line2.line.color.rgb = TEXT_COLOR
    line2.line.width = Pt(2)
    line2.line.dash_style = 7 

    # --- Left Column (Disco Cave) ---
    # Image Placeholder
    img_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(2.2), Inches(3.6), Inches(3.2))
    img_rect.fill.solid()
    img_rect.fill.fore_color.rgb = RGBColor(40, 20, 60)
    img_rect.line.color.rgb = YELLOW
    img_rect.line.width = Pt(3)
    
    # Laser lines
    laser1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(2.4), Inches(3.8), Inches(0.6), Inches(2.5))
    laser1.line.color.rgb = CYAN
    laser1.line.width = Pt(4)
    laser2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(2.4), Inches(3.8), Inches(4.2), Inches(2.8))
    laser2.line.color.rgb = PURPLE
    laser2.line.width = Pt(4)
    laser3 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(2.4), Inches(3.8), Inches(1.0), Inches(5.0))
    laser3.line.color.rgb = ORANGE
    laser3.line.width = Pt(4)

    # Inset Photo Placeholder
    inset_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(1.8), Inches(2.2), Inches(1.6))
    inset_rect.rotation = -10
    inset_rect.fill.solid()
    inset_rect.fill.fore_color.rgb = RGBColor(200, 200, 200)
    inset_rect.line.color.rgb = RGBColor(255, 255, 255)
    inset_rect.line.width = Pt(4)
    
    tape1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(1.7), Inches(0.8), Inches(0.2))
    tape1.rotation = -30
    tape1.fill.solid()
    tape1.fill.fore_color.rgb = GREEN
    tape1.line.fill.background()

    # Text
    left_text = slide.shapes.add_textbox(Inches(0.8), Inches(5.6), Inches(3.2), Inches(1.0))
    tf = left_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "五颜六色的LED灯把\n溶洞变成了夜店"
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = TEXT_COLOR

    # Boombox icon
    boombox = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.2), Inches(5.8), Inches(0.8), Inches(0.5))
    boombox.fill.solid()
    boombox.fill.fore_color.rgb = PURPLE
    boombox.line.color.rgb = TEXT_COLOR
    boombox.line.width = Pt(1.5)

    # --- Middle Column (Imagination) ---
    # Row 1
    rock1_bg = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4.8), Inches(2.2), Inches(1.6), Inches(1.6))
    rock1_bg.fill.solid()
    rock1_bg.fill.fore_color.rgb = ORANGE
    rock1_bg.line.color.rgb = TEXT_COLOR
    rock1_bg.line.width = Pt(2)
    rock1 = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(5.1), Inches(2.4), Inches(1.0), Inches(1.2))
    rock1.fill.solid()
    rock1.fill.fore_color.rgb = RGBColor(160, 140, 120)
    rock1.line.color.rgb = TEXT_COLOR
    
    arrow1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.5), Inches(2.8), Inches(0.6), Inches(0.4))
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = ORANGE
    arrow1.line.color.rgb = TEXT_COLOR

    pork_bg = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.2), Inches(2.2), Inches(1.6), Inches(1.6))
    pork_bg.fill.solid()
    pork_bg.fill.fore_color.rgb = GREEN
    pork_bg.line.color.rgb = TEXT_COLOR
    pork_bg.line.width = Pt(2)
    pork = slide.shapes.add_shape(MSO_SHAPE.CUBE, Inches(7.5), Inches(2.5), Inches(1.0), Inches(1.0))
    pork.fill.solid()
    pork.fill.fore_color.rgb = RGBColor(180, 80, 40)
    pork.line.color.rgb = TEXT_COLOR
    
    pork_text = slide.shapes.add_textbox(Inches(6.2), Inches(2.0), Inches(1.2), Inches(0.4))
    pork_text.text_frame.text = "红烧肉?"
    pork_text.text_frame.paragraphs[0].font.size = Pt(14)
    pork_text.text_frame.paragraphs[0].font.bold = True

    # Row 2
    rock2_bg = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4.8), Inches(4.0), Inches(1.6), Inches(1.6))
    rock2_bg.fill.solid()
    rock2_bg.fill.fore_color.rgb = CYAN
    rock2_bg.line.color.rgb = TEXT_COLOR
    rock2_bg.line.width = Pt(2)
    rock2 = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(5.1), Inches(4.2), Inches(1.0), Inches(1.2))
    rock2.fill.solid()
    rock2.fill.fore_color.rgb = RGBColor(160, 140, 120)
    rock2.line.color.rgb = TEXT_COLOR

    arrow2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.5), Inches(4.6), Inches(0.6), Inches(0.4))
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = GREEN
    arrow2.line.color.rgb = TEXT_COLOR

    alien_bg = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.2), Inches(4.0), Inches(1.6), Inches(1.6))
    alien_bg.fill.solid()
    alien_bg.fill.fore_color.rgb = PURPLE
    alien_bg.line.color.rgb = TEXT_COLOR
    alien_bg.line.width = Pt(2)
    alien = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.5), Inches(4.3), Inches(1.0), Inches(1.2))
    alien.fill.solid()
    alien.fill.fore_color.rgb = GREEN
    alien.line.color.rgb = TEXT_COLOR
    
    alien_text = slide.shapes.add_textbox(Inches(8.2), Inches(3.8), Inches(1.2), Inches(0.4))
    alien_text.text_frame.text = "外星人?"
    alien_text.text_frame.paragraphs[0].font.size = Pt(14)
    alien_text.text_frame.paragraphs[0].font.bold = True

    # Text
    mid_text = slide.shapes.add_textbox(Inches(5.2), Inches(5.6), Inches(3.5), Inches(1.0))
    tf = mid_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "猎奇想象：这块石头像红烧肉，\n那块石头像外星人"
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = TEXT_COLOR

    # Magnifying glass
    mag = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4.7), Inches(5.8), Inches(0.4), Inches(0.4))
    mag.fill.solid()
    mag.fill.fore_color.rgb = LIGHT_BLUE
    mag.line.color.rgb = TEXT_COLOR
    mag.line.width = Pt(1.5)
    mag_handle = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(4.8), Inches(6.1), Inches(4.6), Inches(6.3))
    mag_handle.line.color.rgb = ORANGE
    mag_handle.line.width = Pt(4)

    # --- Right Column (Comic Panels) ---
    # Panel 1
    panel1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.2), Inches(2.0), Inches(3.6), Inches(1.7))
    panel1.fill.solid()
    panel1.fill.fore_color.rgb = LIGHT_BLUE
    panel1.line.color.rgb = TEXT_COLOR
    panel1.line.width = Pt(2)
    
    face1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(2.4), Inches(0.8), Inches(0.8))
    face1.fill.solid()
    face1.fill.fore_color.rgb = RGBColor(255, 224, 189)
    face1.line.color.rgb = TEXT_COLOR
    
    bubble1 = slide.shapes.add_shape(MSO_SHAPE.CLOUD_CALLOUT, Inches(10.5), Inches(2.2), Inches(2.0), Inches(1.0))
    bubble1.fill.solid()
    bubble1.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bubble1.line.color.rgb = TEXT_COLOR
    bubble1.line.width = Pt(1.5)
    tf = bubble1.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "灵魂被洗涤"
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = CYAN

    # Panel 2
    panel2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.2), Inches(3.9), Inches(3.6), Inches(1.7))
    panel2.fill.solid()
    panel2.fill.fore_color.rgb = LIGHT_BLUE
    panel2.line.color.rgb = TEXT_COLOR
    panel2.line.width = Pt(2)

    face2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(4.3), Inches(0.8), Inches(0.8))
    face2.fill.solid()
    face2.fill.fore_color.rgb = RGBColor(255, 224, 189)
    face2.line.color.rgb = TEXT_COLOR
    
    scarf = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.4), Inches(4.9), Inches(1.0), Inches(0.3))
    scarf.fill.solid()
    scarf.fill.fore_color.rgb = ORANGE
    scarf.line.color.rgb = TEXT_COLOR

    bubble2 = slide.shapes.add_shape(MSO_SHAPE.OVAL_CALLOUT, Inches(10.8), Inches(4.0), Inches(1.5), Inches(1.0))
    bubble2.fill.solid()
    bubble2.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bubble2.line.color.rgb = TEXT_COLOR
    bubble2.line.width = Pt(1.5)
    tf = bubble2.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "冷！"
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = CYAN

    cold_text = slide.shapes.add_textbox(Inches(10.8), Inches(5.0), Inches(1.5), Inches(0.5))
    tf = cold_text.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "冻感冒了"
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = TEXT_COLOR

    # Text
    right_text = slide.shapes.add_textbox(Inches(9.0), Inches(5.6), Inches(4.0), Inches(1.0))
    tf = right_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "吐槽：在洞里走了一圈，感觉\n灵魂被洗涤（其实是冻感冒了）"
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = TEXT_COLOR

    # Thermometer
    thermo = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(12.8), Inches(5.8), Inches(0.2), Inches(0.8))
    thermo.fill.solid()
    thermo.fill.fore_color.rgb = RGBColor(200, 200, 200)
    thermo.line.color.rgb = TEXT_COLOR
    thermo_bulb = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(12.75), Inches(6.4), Inches(0.3), Inches(0.3))
    thermo_bulb.fill.solid()
    thermo_bulb.fill.fore_color.rgb = RGBColor(255, 0, 0)
    thermo_bulb.line.color.rgb = TEXT_COLOR

    # --- Footer ---
    page_bg = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(12.2), Inches(7.0), Inches(1.0), Inches(0.4))
    page_bg.fill.solid()
    page_bg.fill.fore_color.rgb = GREEN
    page_bg.line.color.rgb = TEXT_COLOR
    page_bg.line.width = Pt(1.5)
    
    page_text = slide.shapes.add_textbox(Inches(12.3), Inches(6.95), Inches(0.8), Inches(0.4))
    tf = page_text.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "第6页"
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = TEXT_COLOR

    # Decorations
    excl1 = slide.shapes.add_textbox(Inches(4.5), Inches(1.8), Inches(0.5), Inches(0.5))
    excl1.text_frame.text = "!!"
    excl1.text_frame.paragraphs[0].font.size = Pt(24)
    excl1.text_frame.paragraphs[0].font.bold = True
    excl1.text_frame.paragraphs[0].font.color.rgb = ORANGE
    excl1.rotation = -15

    excl2 = slide.shapes.add_textbox(Inches(4.5), Inches(3.8), Inches(0.5), Inches(0.5))
    excl2.text_frame.text = "!!"
    excl2.text_frame.paragraphs[0].font.size = Pt(24)
    excl2.text_frame.paragraphs[0].font.bold = True
    excl2.text_frame.paragraphs[0].font.color.rgb = PURPLE
    excl2.rotation = -15

    wow = slide.shapes.add_textbox(Inches(2.8), Inches(1.5), Inches(1.0), Inches(0.5))
    wow.text_frame.text = "WOW!"
    wow.text_frame.paragraphs[0].font.size = Pt(20)
    wow.text_frame.paragraphs[0].font.bold = True
    wow.text_frame.paragraphs[0].font.color.rgb = PURPLE
    wow.rotation = -10

    omg = slide.shapes.add_textbox(Inches(8.0), Inches(1.8), Inches(1.0), Inches(0.5))
    omg.text_frame.text = "OMG!"
    omg.text_frame.paragraphs[0].font.size = Pt(20)
    omg.text_frame.paragraphs[0].font.bold = True
    omg.text_frame.paragraphs[0].font.color.rgb = GREEN
    omg.rotation = 15



# ── Slide 7 ──

def build_slide_7(slide):
    # --- 颜色定义 ---
    ORANGE = RGBColor(0xF2, 0x71, 0x27)
    DARK_TEXT = RGBColor(0x20, 0x20, 0x20)
    GRAY_TEXT = RGBColor(0x60, 0x60, 0x60)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    BLUE_LIGHT = RGBColor(0x00, 0xBC, 0xD4)
    RED = RGBColor(0xE5, 0x39, 0x35)
    IMG_BG = RGBColor(0xE8, 0xF0, 0xF4)

    # --- 1. 左上角页码 ---
    page_num = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.2), Inches(0.2), Inches(1.0), Inches(1.0))
    page_num.fill.solid()
    page_num.fill.fore_color.rgb = ORANGE
    page_num.line.fill.background()
    tf = page_num.text_frame
    p = tf.paragraphs[0]
    p.text = "第7页"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DARK_TEXT
    p.alignment = PP_ALIGN.CENTER

    # --- 2. 主标题与副标题 ---
    title_box = slide.shapes.add_textbox(Inches(1.5), Inches(0.3), Inches(10), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "漓江竹筏：塑料椅子上的“速度与激情”"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.name = "Microsoft YaHei"
    p.font.color.rgb = DARK_TEXT

    # 标题下划线
    underline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.6), Inches(0.9), Inches(8.5), Inches(0.1))
    underline.fill.solid()
    underline.fill.fore_color.rgb = ORANGE
    underline.line.fill.background()

    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(2.5), Inches(1.1), Inches(8), Inches(0.5))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "别被照片骗了，这其实是水上拖拉机"
    p.font.size = Pt(20)
    p.font.name = "Microsoft YaHei"
    p.font.color.rgb = GRAY_TEXT
    p.alignment = PP_ALIGN.CENTER

    # --- 辅助函数：添加带下划线的区块标题 ---
    def add_section_title(left, top, text):
        box = slide.shapes.add_textbox(left, top, Inches(2), Inches(0.5))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.name = "Microsoft YaHei"
        
        ul = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + Inches(0.4), Inches(1.5), Inches(0.08))
        ul.fill.solid()
        ul.fill.fore_color.rgb = ORANGE
        ul.line.fill.background()

    # ==========================================
    # --- Section 1: 现实反差 (左侧) ---
    # ==========================================
    add_section_title(Inches(0.8), Inches(1.8), "现实反差")
    
    t1 = slide.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(2), Inches(0.4))
    t1.text_frame.text = "理想中：优雅的竹筏"
    
    # 图片占位 1 (传统竹筏)
    img1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.8), Inches(2.2), Inches(1.5))
    img1.fill.solid()
    img1.fill.fore_color.rgb = IMG_BG
    img1.line.color.rgb = WHITE
    img1.line.width = Pt(3)
    
    # 红色大叉
    cross = slide.shapes.add_shape(MSO_SHAPE.MATH_MULTIPLY, Inches(1.5), Inches(3.1), Inches(0.8), Inches(0.8))
    cross.fill.solid()
    cross.fill.fore_color.rgb = RED
    cross.line.fill.background()

    # 气泡提示
    bubble = slide.shapes.add_shape(MSO_SHAPE.CLOUD_CALLOUT, Inches(2.5), Inches(2.5), Inches(1.5), Inches(1.0))
    bubble.fill.solid()
    bubble.fill.fore_color.rgb = BLUE_LIGHT
    bubble.line.fill.background()
    tf = bubble.text_frame
    p = tf.paragraphs[0]
    p.text = "其实是\nPVC管+马达"
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 图片占位 2 (PVC竹筏)
    img2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(3.8), Inches(2.8), Inches(2.0))
    img2.fill.solid()
    img2.fill.fore_color.rgb = IMG_BG
    img2.line.color.rgb = WHITE
    img2.line.width = Pt(3)

    # 底部说明文字
    cap1 = slide.shapes.add_textbox(Inches(0.8), Inches(5.9), Inches(3.5), Inches(0.5))
    tf = cap1.text_frame
    p = tf.paragraphs[0]
    p.text = "优雅的竹筏其实是PVC管+马达"
    p.font.size = Pt(14)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # ==========================================
    # --- Section 2: 搞笑瞬间 (中间) ---
    # ==========================================
    add_section_title(Inches(5.2), Inches(1.8), "搞笑瞬间")

    # 图片占位
    img3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.0), Inches(2.5), Inches(3.5), Inches(2.8))
    img3.fill.solid()
    img3.fill.fore_color.rgb = IMG_BG
    img3.line.color.rgb = WHITE
    img3.line.width = Pt(3)

    # 拟声词
    t_boom = slide.shapes.add_textbox(Inches(6.5), Inches(2.8), Inches(1.5), Inches(0.5))
    t_boom.text_frame.text = "轰隆隆!\nBOOM!"
    t_boom.text_frame.paragraphs[0].font.bold = True

    # 吐槽气泡
    t_omg = slide.shapes.add_shape(MSO_SHAPE.OVAL_CALLOUT, Inches(7.0), Inches(4.0), Inches(1.2), Inches(0.8))
    t_omg.fill.solid()
    t_omg.fill.fore_color.rgb = RGBColor(0xA0, 0xD0, 0xA0)
    t_omg.line.fill.background()
    tf = t_omg.text_frame
    p = tf.paragraphs[0]
    p.text = "OMG!\n诗意呢?"
    p.font.size = Pt(10)
    p.alignment = PP_ALIGN.CENTER

    # 底部说明文字
    cap2 = slide.shapes.add_textbox(Inches(4.8), Inches(5.4), Inches(4.0), Inches(0.5))
    tf = cap2.text_frame
    p = tf.paragraphs[0]
    p.text = "马达启动时的黑烟与诗意山水的完美融合"
    p.font.size = Pt(14)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # ==========================================
    # --- Section 3: 猎奇体验 (右侧) ---
    # ==========================================
    add_section_title(Inches(9.5), Inches(1.8), "猎奇体验")

    # 图片占位
    img4 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.2), Inches(2.5), Inches(3.5), Inches(2.8))
    img4.fill.solid()
    img4.fill.fore_color.rgb = IMG_BG
    img4.line.color.rgb = WHITE
    img4.line.width = Pt(3)

    # 标签：移动超市
    t_shop = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10.0), Inches(2.8), Inches(1.2), Inches(0.4))
    t_shop.fill.solid()
    t_shop.fill.fore_color.rgb = RGBColor(0xFF, 0xDD, 0xAA)
    t_shop.line.fill.background()
    tf = t_shop.text_frame
    p = tf.paragraphs[0]
    p.text = "移动超市"
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_TEXT
    p.alignment = PP_ALIGN.CENTER

    # 标签：现烤鱼
    t_fish = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.5), Inches(4.5), Inches(1.0), Inches(0.4))
    t_fish.fill.solid()
    t_fish.fill.fore_color.rgb = RGBColor(0xFF, 0xDD, 0xAA)
    t_fish.line.fill.background()
    tf = t_fish.text_frame
    p = tf.paragraphs[0]
    p.text = "现烤鱼"
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_TEXT
    p.alignment = PP_ALIGN.CENTER

    # 底部说明文字
    cap3 = slide.shapes.add_textbox(Inches(9.0), Inches(5.4), Inches(4.0), Inches(0.5))
    tf = cap3.text_frame
    p = tf.paragraphs[0]
    p.text = "江上的“移动超市”，划着竹筏卖烤鱼"
    p.font.size = Pt(14)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # --- 区块之间的引导箭头 ---
    arrow1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.5), Inches(1.9), Inches(0.8), Inches(0.3))
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = GRAY_TEXT
    arrow1.line.fill.background()

    arrow2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(7.8), Inches(1.9), Inches(0.8), Inches(0.3))
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = GRAY_TEXT
    arrow2.line.fill.background()

    # ==========================================
    # --- 底部图表区域: 马达轰鸣声 vs 优雅度 ---
    # ==========================================
    bot_title = slide.shapes.add_textbox(Inches(6.5), Inches(5.8), Inches(3.5), Inches(0.5))
    tf = bot_title.text_frame
    p = tf.paragraphs[0]
    p.text = "马达轰鸣声 vs 优雅度"
    p.font.size = Pt(18)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    bot_ul = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(6.2), Inches(2.9), Inches(0.08))
    bot_ul.fill.solid()
    bot_ul.fill.fore_color.rgb = ORANGE
    bot_ul.line.fill.background()

    # 渐变进度条 (使用纯色箭头代替)
    bar = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.0), Inches(6.4), Inches(5.0), Inches(0.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0xFF, 0x98, 0x00) # 橙色
    bar.line.color.rgb = DARK_TEXT
    bar.line.width = Pt(1)

    # 左侧文本 (优雅度)
    left_txt = slide.shapes.add_textbox(Inches(5.0), Inches(6.1), Inches(1.0), Inches(0.8))
    tf = left_txt.text_frame
    tf.text = "🕊️\n优雅度\n(最高)"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    if len(tf.paragraphs) > 1:
        tf.paragraphs[1].font.size = Pt(12)
        tf.paragraphs[1].font.bold = True
        tf.paragraphs[1].alignment = PP_ALIGN.CENTER
    if len(tf.paragraphs) > 2:
        tf.paragraphs[2].font.size = Pt(10)
        tf.paragraphs[2].alignment = PP_ALIGN.CENTER

    # 右侧文本 (马达轰鸣声)
    right_txt = slide.shapes.add_textbox(Inches(11.2), Inches(6.1), Inches(1.5), Inches(0.8))
    tf = right_txt.text_frame
    tf.text = "🚜\n马达轰鸣声\n(最大)"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    if len(tf.paragraphs) > 1:
        tf.paragraphs[1].font.size = Pt(12)
        tf.paragraphs[1].font.bold = True
        tf.paragraphs[1].alignment = PP_ALIGN.CENTER
    if len(tf.paragraphs) > 2:
        tf.paragraphs[2].font.size = Pt(10)
        tf.paragraphs[2].alignment = PP_ALIGN.CENTER

    # 进度条上的标签
    labels = ["静音", "嗡嗡...", "轰隆隆!!!", "VROOM!!!"]
    x_positions = [6.3, 7.3, 8.5, 9.8]
    for i, text in enumerate(labels):
        lbl = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_positions[i]), Inches(6.8), Inches(1.0), Inches(0.3))
        lbl.fill.solid()
        lbl.fill.fore_color.rgb = WHITE
        lbl.line.color.rgb = DARK_TEXT
        lbl.line.width = Pt(1)
        tf = lbl.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = DARK_TEXT
        p.alignment = PP_ALIGN.CENTER

    # 底部总结文字
    bot_cap = slide.shapes.add_textbox(Inches(6.0), Inches(7.2), Inches(5.0), Inches(0.4))
    tf = bot_cap.text_frame
    p = tf.paragraphs[0]
    p.text = "马达一响，优雅全无 (进度条图表)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER



# ── Slide 8 ──

def build_slide_8(slide):
    # --- Background & Decorative Elements ---
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(245, 245, 240)
    bg.line.fill.background()

    # Colorful blobs in corners
    blob1 = slide.shapes.add_shape(MSO_SHAPE.CLOUD, Inches(-1.0), Inches(-1.0), Inches(3.0), Inches(3.0))
    blob1.fill.solid(); blob1.fill.fore_color.rgb = RGBColor(200, 255, 50)
    blob1.line.fill.background()

    blob2 = slide.shapes.add_shape(MSO_SHAPE.CLOUD, Inches(11.5), Inches(-1.0), Inches(3.0), Inches(3.0))
    blob2.fill.solid(); blob2.fill.fore_color.rgb = RGBColor(255, 100, 200)
    blob2.line.fill.background()

    blob3 = slide.shapes.add_shape(MSO_SHAPE.CLOUD, Inches(-1.0), Inches(5.5), Inches(3.0), Inches(3.0))
    blob3.fill.solid(); blob3.fill.fore_color.rgb = RGBColor(255, 150, 0)
    blob3.line.fill.background()

    # --- Titles ---
    # Main Title Shadow
    title1_shadow = slide.shapes.add_textbox(Inches(0.53), Inches(0.33), Inches(5), Inches(0.8))
    p_shadow = title1_shadow.text_frame.paragraphs[0]
    p_shadow.text = "阳朔西街："
    p_shadow.font.name = "Microsoft YaHei"
    p_shadow.font.size = Pt(48)
    p_shadow.font.bold = True
    p_shadow.font.color.rgb = RGBColor(0, 0, 0)

    # Main Title
    title1 = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(5), Inches(0.8))
    p = title1.text_frame.paragraphs[0]
    p.text = "阳朔西街："
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(50, 20, 80)

    # Subtitle Shadow
    title2_shadow = slide.shapes.add_textbox(Inches(0.53), Inches(1.03), Inches(8), Inches(0.8))
    p_shadow2 = title2_shadow.text_frame.paragraphs[0]
    p_shadow2.text = "中西合璧的“迷惑行为”大赏"
    p_shadow2.font.name = "Microsoft YaHei"
    p_shadow2.font.size = Pt(36)
    p_shadow2.font.bold = True
    p_shadow2.font.color.rgb = RGBColor(0, 0, 0)

    # Subtitle
    title2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(8), Inches(0.8))
    p2 = title2.text_frame.paragraphs[0]
    p2.text = "中西合璧的“迷惑行为”大赏"
    p2.font.name = "Microsoft YaHei"
    p2.font.size = Pt(36)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(230, 80, 0)

    # Banner
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.9), Inches(6.0), Inches(0.6))
    banner.fill.solid(); banner.fill.fore_color.rgb = RGBColor(160, 80, 230)
    banner.line.fill.background()
    p = banner.text_frame.paragraphs[0]
    p.text = "这里的外国人比外地人还多"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # --- Map Section (Top Right) ---
    ff = slide.shapes.build_freeform(Inches(8.0), Inches(0.5))
    ff.add_line_segments([
        (Inches(10.0), Inches(0.2)), (Inches(12.5), Inches(1.0)),
        (Inches(12.0), Inches(2.5)), (Inches(9.5), Inches(3.5)),
        (Inches(8.0), Inches(3.0)), (Inches(4.5), Inches(3.8)),
        (Inches(4.0), Inches(3.0)), (Inches(7.5), Inches(2.0)),
        (Inches(8.0), Inches(0.5))
    ])
    map_shape = ff.convert_to_shape()
    map_shape.fill.solid(); map_shape.fill.fore_color.rgb = RGBColor(40, 40, 40)
    map_shape.line.color.rgb = RGBColor(150, 80, 220); map_shape.line.width = Pt(3)

    # Neon streets
    def add_neon_line(x1, y1, x2, y2, color):
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        line.line.color.rgb = color; line.line.width = Pt(3)

    add_neon_line(4.5, 3.2, 8.0, 2.5, RGBColor(255, 150, 0))
    add_neon_line(8.0, 2.5, 12.0, 1.8, RGBColor(255, 150, 0))
    add_neon_line(6.0, 2.9, 6.2, 2.0, RGBColor(255, 50, 150))
    add_neon_line(9.0, 2.3, 9.5, 1.0, RGBColor(150, 255, 50))
    add_neon_line(10.5, 2.0, 11.0, 3.0, RGBColor(255, 50, 150))

    # Map Text
    def add_map_text(x, y, text, rot, color):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(1.5), Inches(0.4))
        tb.rotation = rot
        p = tb.text_frame.paragraphs[0]
        p.text = text; p.font.size = Pt(10); p.font.color.rgb = color; p.font.bold = True

    add_map_text(6.5, 2.6, "阳朔西街", -10, RGBColor(255,255,255))
    add_map_text(10.0, 1.8, "阳朔西街", -5, RGBColor(255,255,255))
    add_map_text(9.2, 1.2, "阳朔西街", 70, RGBColor(255,255,255))

    # Comic Texts
    def add_comic_text(x, y, text, color, rot):
        tb_shadow = slide.shapes.add_textbox(Inches(x+0.03), Inches(y+0.03), Inches(2.0), Inches(1.0))
        tb_shadow.rotation = rot
        p_shadow = tb_shadow.text_frame.paragraphs[0]
        p_shadow.text = text; p_shadow.font.name = "Arial Black"; p_shadow.font.size = Pt(28); p_shadow.font.color.rgb = RGBColor(0,0,0)

        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(2.0), Inches(1.0))
        tb.rotation = rot
        p = tb.text_frame.paragraphs[0]
        p.text = text; p.font.name = "Arial Black"; p.font.size = Pt(28); p.font.color.rgb = color

    add_comic_text(7.0, 0.5, "Boom!", RGBColor(255, 50, 150), -15)
    add_comic_text(4.5, 4.0, "OMG!", RGBColor(150, 255, 50), 10)
    add_comic_text(11.0, 2.5, "WoW!", RGBColor(255, 50, 150), -20)

    # Decor
    bolt1 = slide.shapes.add_shape(MSO_SHAPE.LIGHTNING_BOLT, Inches(6.5), Inches(0.5), Inches(0.4), Inches(0.8))
    bolt1.fill.solid(); bolt1.fill.fore_color.rgb = RGBColor(150, 80, 220); bolt1.line.color.rgb = RGBColor(0,0,0)
    bolt2 = slide.shapes.add_shape(MSO_SHAPE.LIGHTNING_BOLT, Inches(12.0), Inches(1.0), Inches(0.4), Inches(0.8))
    bolt2.fill.solid(); bolt2.fill.fore_color.rgb = RGBColor(255, 150, 0); bolt2.line.color.rgb = RGBColor(0,0,0)
    star1 = slide.shapes.add_shape(MSO_SHAPE.5_POINT_STAR, Inches(10.0), Inches(2.8), Inches(0.4), Inches(0.4))
    star1.fill.solid(); star1.fill.fore_color.rgb = RGBColor(255, 150, 0); star1.line.color.rgb = RGBColor(0,0,0)

    # --- Section Titles ---
    def add_section_title(x, y, text):
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(1.6), Inches(0.4))
        box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0,0,0); box.line.fill.background()
        p = box.text_frame.paragraphs[0]
        p.text = text; p.font.name = "Microsoft YaHei"; p.font.size = Pt(16); p.font.color.rgb = RGBColor(255,255,255); p.font.bold = True; p.alignment = PP_ALIGN.CENTER

    add_section_title(0.5, 2.8, "猎奇景观：")
    add_section_title(6.0, 4.0, "搞笑互动：")
    add_section_title(9.5, 3.5, "避坑指南：")

    # --- Section 1: Beer Fish ---
    desc1 = slide.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(3.0), Inches(0.8))
    p1 = desc1.text_frame.paragraphs[0]
    p1.text = "卖啤酒鱼的店比"
    p1.font.name = "Microsoft YaHei"; p1.font.size = Pt(16); p1.font.bold = True
    p2 = desc1.text_frame.add_paragraph()
    p2.text = "整条街的树还多"
    p2.font.name = "Microsoft YaHei"; p2.font.size = Pt(16); p2.font.bold = True

    # Buildings
    b1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(4.5), Inches(1.5), Inches(2.5))
    b1.fill.solid(); b1.fill.fore_color.rgb = RGBColor(50,50,50)
    b2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.2), Inches(4.0), Inches(1.2), Inches(3.0))
    b2.fill.solid(); b2.fill.fore_color.rgb = RGBColor(40,40,40)
    b3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.5), Inches(4.5), Inches(1.2), Inches(2.5))
    b3.fill.solid(); b3.fill.fore_color.rgb = RGBColor(60,60,60)

    # Trees
    tree_trunk1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(6.5), Inches(0.2), Inches(0.8))
    tree_trunk1.fill.solid(); tree_trunk1.fill.fore_color.rgb = RGBColor(100, 50, 0)
    tree_leaves1 = slide.shapes.add_shape(MSO_SHAPE.CLOUD, Inches(1.0), Inches(6.0), Inches(1.2), Inches(0.8))
    tree_leaves1.fill.solid(); tree_leaves1.fill.fore_color.rgb = RGBColor(50, 150, 50)
    tree_trunk2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.0), Inches(6.8), Inches(0.2), Inches(0.6))
    tree_trunk2.fill.solid(); tree_trunk2.fill.fore_color.rgb = RGBColor(100, 50, 0)
    tree_leaves2 = slide.shapes.add_shape(MSO_SHAPE.CLOUD, Inches(2.6), Inches(6.4), Inches(1.0), Inches(0.7))
    tree_leaves2.fill.solid(); tree_leaves2.fill.fore_color.rgb = RGBColor(50, 150, 50)

    # Signs
    def add_beer_fish_sign(x, y, w, h, bg_rgb, text_rgb, rot=0):
        sign = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        sign.fill.solid(); sign.fill.fore_color.rgb = bg_rgb
        sign.line.color.rgb = RGBColor(0,0,0); sign.line.width = Pt(2); sign.rotation = rot
        tf = sign.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = "\n".join(list("啤酒鱼")) if h > w else "啤酒鱼"
        p.font.name = "Microsoft YaHei"; p.font.bold = True; p.font.size = Pt(14); p.font.color.rgb = text_rgb; p.alignment = PP_ALIGN.CENTER

    add_beer_fish_sign(0.5, 4.5, 0.8, 1.5, RGBColor(150,255,50), RGBColor(0,0,0))
    add_beer_fish_sign(1.2, 4.2, 0.6, 1.8, RGBColor(255,50,150), RGBColor(255,255,255))
    add_beer_fish_sign(0.6, 5.5, 1.2, 0.5, RGBColor(255,150,0), RGBColor(0,0,0))
    add_beer_fish_sign(0.4, 6.2, 1.5, 0.6, RGBColor(150,255,50), RGBColor(0,0,0))
    add_beer_fish_sign(2.0, 4.8, 0.7, 1.6, RGBColor(255,50,150), RGBColor(255,255,255))
    add_beer_fish_sign(2.5, 4.0, 0.8, 2.0, RGBColor(150,255,50), RGBColor(0,0,0))
    add_beer_fish_sign(2.2, 5.8, 1.0, 0.5, RGBColor(255,150,0), RGBColor(0,0,0))
    add_beer_fish_sign(2.8, 5.2, 0.6, 1.4, RGBColor(255,50,150), RGBColor(255,255,255))
    add_beer_fish_sign(3.5, 4.5, 1.2, 0.6, RGBColor(150,255,50), RGBColor(0,0,0))
    add_beer_fish_sign(3.4, 5.2, 1.4, 0.7, RGBColor(150,255,50), RGBColor(0,0,0))
    add_beer_fish_sign(3.6, 6.0, 1.2, 0.6, RGBColor(255,150,0), RGBColor(0,0,0))
    add_beer_fish_sign(1.0, 5.0, 1.0, 0.5, RGBColor(255,255,50), RGBColor(0,0,0), rot=-10)
    add_beer_fish_sign(2.8, 4.5, 0.5, 1.2, RGBColor(150,80,220), RGBColor(255,255,255), rot=5)

    # --- Section 2: Interaction ---
    # Snail -> Person icon
    snail_body = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.7), Inches(4.0), Inches(0.4), Inches(0.3))
    snail_body.fill.solid(); snail_body.fill.fore_color.rgb = RGBColor(150, 100, 50)
    arrow1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.2), Inches(4.1), Inches(0.3), Inches(0.15))
    arrow1.fill.solid(); arrow1.fill.fore_color.rgb = RGBColor(255, 150, 0)
    head_icon = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.6), Inches(4.0), Inches(0.3), Inches(0.3))
    head_icon.fill.solid(); head_icon.fill.fore_color.rgb = RGBColor(255, 200, 150)

    # People
    head1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.5), Inches(5.5), Inches(0.6), Inches(0.6))
    head1.fill.solid(); head1.fill.fore_color.rgb = RGBColor(255, 220, 200)
    body1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.3), Inches(6.1), Inches(1.0), Inches(0.8))
    body1.fill.solid(); body1.fill.fore_color.rgb = RGBColor(150, 80, 220)

    head2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.0), Inches(5.4), Inches(0.6), Inches(0.6))
    head2.fill.solid(); head2.fill.fore_color.rgb = RGBColor(255, 220, 200)
    body2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.8), Inches(6.0), Inches(1.0), Inches(0.9))
    body2.fill.solid(); body2.fill.fore_color.rgb = RGBColor(100, 220, 50)

    # Speech Bubbles
    bubble1 = slide.shapes.add_shape(MSO_SHAPE.OVAL_CALLOUT, Inches(5.3), Inches(4.5), Inches(2.2), Inches(0.9))
    bubble1.fill.solid(); bubble1.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bubble1.line.color.rgb = RGBColor(0,0,0); bubble1.line.width = Pt(2)
    p1 = bubble1.text_frame.paragraphs[0]
    p1.text = "How to eat... um...\nTianluo (snails)?"
    p1.font.size = Pt(11); p1.font.color.rgb = RGBColor(0,0,0)

    bubble2 = slide.shapes.add_shape(MSO_SHAPE.OVAL_CALLOUT, Inches(7.7), Inches(4.2), Inches(2.0), Inches(0.8))
    bubble2.fill.solid(); bubble2.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bubble2.line.color.rgb = RGBColor(0,0,0); bubble2.line.width = Pt(2)
    p2 = bubble2.text_frame.paragraphs[0]
    p2.text = "Haha, use a\ntoothpick, mate!"
    p2.font.size = Pt(11); p2.font.color.rgb = RGBColor(0,0,0)

    # Bottom Text Box
    interact_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.0), Inches(6.8), Inches(3.2), Inches(0.6))
    interact_box.fill.solid(); interact_box.fill.fore_color.rgb = RGBColor(0,0,0); interact_box.line.fill.background()
    p = interact_box.text_frame.paragraphs[0]
    p.text = "尝试用蹩脚英语和老外聊怎么吃田螺"
    p.font.name = "Microsoft YaHei"; p.font.size = Pt(12); p.font.color.rgb = RGBColor(255,255,255); p.font.bold = True; p.alignment = PP_ALIGN.CENTER

    # --- Section 3: Tourist Trap ---
    # Shop -> Shop icon
    shop1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(11.2), Inches(3.5), Inches(0.4), Inches(0.4))
    shop1.fill.solid(); shop1.fill.fore_color.rgb = RGBColor(200, 100, 200)
    arrow2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(11.7), Inches(3.6), Inches(0.3), Inches(0.15))
    arrow2.fill.solid(); arrow2.fill.fore_color.rgb = RGBColor(150, 255, 50)
    shop2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(12.1), Inches(3.5), Inches(0.4), Inches(0.4))
    shop2.fill.solid(); shop2.fill.fore_color.rgb = RGBColor(200, 100, 200)

    # Shops
    frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.5), Inches(4.2), Inches(3.5), Inches(1.8))
    frame.fill.solid(); frame.fill.fore_color.rgb = RGBColor(120, 80, 180)
    frame.line.color.rgb = RGBColor(0,0,0); frame.line.width = Pt(2)

    stall_w = 3.5 / 4
    for i in range(4):
        stall = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.5 + i*stall_w + 0.1), Inches(4.8), Inches(stall_w - 0.2), Inches(1.2))
        stall.fill.solid(); stall.fill.fore_color.rgb = RGBColor(60, 40, 90); stall.line.color.rgb = RGBColor(0,0,0)
        
        counter = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.5 + i*stall_w + 0.1), Inches(5.5), Inches(stall_w - 0.2), Inches(0.5))
        counter.fill.solid(); counter.fill.fore_color.rgb = RGBColor(150, 100, 50); counter.line.color.rgb = RGBColor(0,0,0)

        sign = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.5 + i*stall_w + 0.05), Inches(4.3), Inches(stall_w - 0.1), Inches(0.4))
        sign.fill.solid(); sign.fill.fore_color.rgb = RGBColor(255, 100, 100) if i%2==0 else RGBColor(100, 255, 100)
        sign.line.color.rgb = RGBColor(0,0,0)
        tf = sign.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = "网红店"; p.font.size = Pt(10); p.font.bold = True; p.alignment = PP_ALIGN.CENTER

    # Red Cross
    cross1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10.0), Inches(4.5), Inches(12.5), Inches(5.8))
    cross1.line.color.rgb = RGBColor(220, 40, 40); cross1.line.width = Pt(12)
    cross2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(12.5), Inches(4.5), Inches(10.0), Inches(5.8))
    cross2.line.color.rgb = RGBColor(220, 40, 40); cross2.line.width = Pt(12)

    # Bottom Text
    trap_text = slide.shapes.add_textbox(Inches(9.5), Inches(6.2), Inches(3.5), Inches(0.8))
    p = trap_text.text_frame.paragraphs[0]
    p.text = "西街的“网红”其实都是同一个模版"
    p.font.name = "Microsoft YaHei"; p.font.size = Pt(14); p.font.bold = True

    # --- Footer ---
    mt1 = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(11.0), Inches(6.8), Inches(0.5), Inches(0.5))
    mt1.fill.solid(); mt1.fill.fore_color.rgb = RGBColor(200, 200, 200); mt1.line.color.rgb = RGBColor(0,0,0)
    mt2 = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(11.3), Inches(6.9), Inches(0.4), Inches(0.4))
    mt2.fill.solid(); mt2.fill.fore_color.rgb = RGBColor(180, 180, 180); mt2.line.color.rgb = RGBColor(0,0,0)

    footer = slide.shapes.add_textbox(Inches(12.0), Inches(6.8), Inches(1.0), Inches(0.5))
    p = footer.text_frame.paragraphs[0]
    p.text = "08/11"; p.font.name = "Arial"; p.font.size = Pt(20); p.font.bold = True



# ── Slide 9 ──

def build_slide_9(slide):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

    # Colors
    BG_COLOR = RGBColor(245, 242, 235)
    ORANGE_DECOR = RGBColor(235, 104, 65)
    PURPLE_DECOR = RGBColor(103, 58, 183)
    YELLOW_HL = RGBColor(255, 235, 59)
    BLACK = RGBColor(0, 0, 0)
    WHITE = RGBColor(255, 255, 255)
    BOX_ORANGE = RGBColor(255, 112, 67)
    BOX_GREEN = RGBColor(129, 199, 132)
    BOX_YELLOW = RGBColor(255, 235, 59)

    # 1. Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_COLOR
    bg.line.fill.background()

    # Corner decorations (simulating torn paper)
    tr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(11), Inches(0), Inches(2.333), Inches(1.5))
    tr.rotation = 180
    tr.fill.solid()
    tr.fill.fore_color.rgb = ORANGE_DECOR
    tr.line.fill.background()

    bl = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(0), Inches(5.5), Inches(2.5), Inches(2))
    bl.fill.solid()
    bl.fill.fore_color.rgb = ORANGE_DECOR
    bl.line.fill.background()

    br = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(12), Inches(6.5), Inches(1.333), Inches(1))
    br.rotation = 270
    br.fill.solid()
    br.fill.fore_color.rgb = PURPLE_DECOR
    br.line.fill.background()

    # 2. Title & Subtitle
    # Title Highlight
    hl = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.9), Inches(8.0), Inches(0.6))
    hl.fill.solid()
    hl.fill.fore_color.rgb = YELLOW_HL
    hl.line.fill.background()

    # Title Text
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.7), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "桂林生存法则：如何优雅地避坑"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = BLACK

    # Subtitle Text
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(5), Inches(0.5))
    tf = sub_box.text_frame
    p = tf.add_paragraph()
    p.text = "猎奇博主的血泪教训"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = BLACK

    # 3. Left Visuals
    # Landscape Photo Base (White border)
    photo_base = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.0), Inches(5.0), Inches(3.5))
    photo_base.fill.solid()
    photo_base.fill.fore_color.rgb = WHITE
    photo_base.line.fill.background()
    photo_base.rotation = -2

    # Landscape Inner
    land = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(3.2), Inches(4.6), Inches(3.1))
    land.fill.solid()
    land.fill.fore_color.rgb = RGBColor(46, 139, 87)
    land.line.fill.background()
    land.rotation = -2
    tf = land.text_frame
    tf.text = "\n\n桂林山水"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE

    # 20 RMB Note
    note = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.5), Inches(2.0), Inches(2.8), Inches(1.4))
    note.fill.solid()
    note.fill.fore_color.rgb = RGBColor(245, 235, 235)
    note.line.color.rgb = RGBColor(200, 200, 200)
    note.rotation = -10
    tf = note.text_frame
    tf.text = "20\n中国人民银行"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(150, 50, 50)

    # Blogger (Simulated)
    blogger_body = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.5), Inches(1.8), Inches(2.5))
    blogger_body.fill.solid()
    blogger_body.fill.fore_color.rgb = RGBColor(255, 193, 7)
    blogger_body.line.color.rgb = BLACK
    blogger_body.line.width = Pt(2)
    tf = blogger_body.text_frame
    tf.text = "\n\n博主"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = BLACK

    # Boom!
    boom = slide.shapes.add_shape(MSO_SHAPE.EXPLOSION1, Inches(0.5), Inches(2.2), Inches(1.8), Inches(1.2))
    boom.fill.solid()
    boom.fill.fore_color.rgb = RGBColor(255, 87, 34)
    boom.line.color.rgb = BLACK
    boom.line.width = Pt(2)
    boom.rotation = -15
    tf = boom.text_frame
    p = tf.paragraphs[0]
    p.text = "Boom!"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = BLACK

    # OMG!
    omg = slide.shapes.add_shape(MSO_SHAPE.EXPLOSION2, Inches(4.2), Inches(5.5), Inches(1.5), Inches(1.2))
    omg.fill.solid()
    omg.fill.fore_color.rgb = RGBColor(255, 87, 34)
    omg.line.color.rgb = BLACK
    omg.line.width = Pt(2)
    omg.rotation = 10
    tf = omg.text_frame
    p = tf.paragraphs[0]
    p.text = "OMG!"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = BLACK

    # 4. Right Pitfall Boxes
    def add_pitfall_box(slide, left, top, img_bg_color, text_bg_color, text_color, img_text, desc_text):
        box_w = Inches(2.8)
        img_h = Inches(1.5)
        txt_h = Inches(0.9)

        # Image part
        img_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, box_w, img_h)
        img_box.fill.solid()
        img_box.fill.fore_color.rgb = img_bg_color
        img_box.line.color.rgb = BLACK
        img_box.line.width = Pt(2)
        tf = img_box.text_frame
        tf.text = img_text
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = BLACK

        # Text part
        txt_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + img_h, box_w, txt_h)
        txt_box.fill.solid()
        txt_box.fill.fore_color.rgb = text_bg_color
        txt_box.line.color.rgb = BLACK
        txt_box.line.width = Pt(2)
        tf = txt_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc_text
        p.alignment = PP_ALIGN.LEFT
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = text_color

        # Warning Triangle
        warn = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, left - Inches(0.25), top - Inches(0.25), Inches(0.6), Inches(0.6))
        warn.fill.solid()
        warn.fill.fore_color.rgb = YELLOW_HL
        warn.line.color.rgb = BLACK
        warn.line.width = Pt(2)
        tf = warn.text_frame
        p = tf.paragraphs[0]
        p.text = "!"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = BLACK

    # Box 1: Auntie
    add_pitfall_box(slide, Inches(6.5), Inches(1.5),
                    BOX_ORANGE, BOX_YELLOW, BLACK,
                    "帅哥/美女，来吃鱼！\n\nFAKE NEWS!  ❌ 坑！",
                    "别相信路边喊你“帅哥/美女”去吃鱼的阿姨")

    # Box 2: Cormorant 1
    add_pitfall_box(slide, Inches(10.0), Inches(1.2),
                    BOX_GREEN, BLACK, WHITE,
                    "¥999+\n\n(鸬鹚拍照)",
                    "拍照前先问价格，否则鸬鹚的身价比你还高")

    # Box 3: Cormorant 2
    add_pitfall_box(slide, Inches(6.5), Inches(4.5),
                    BOX_GREEN, BLACK, WHITE,
                    "¥999+\n\n(鸬鹚拍照)",
                    "拍照前先问价格，否则鸬鹚的身价比你还高")

    # Box 4: Luosifen
    add_pitfall_box(slide, Inches(10.0), Inches(4.5),
                    BOX_ORANGE, BLACK, WHITE,
                    "螺蛳粉 + 肠胃药\n\n🔥 危险！",
                    "备好肠胃药，螺蛳粉的后劲比酒还大")

    # 5. Arrows (Using connectors)
    conn1 = slide.shapes.add_connector(MSO_CONNECTOR.CURVE, Inches(9.3), Inches(2.2), Inches(10.0), Inches(2.2))
    conn1.line.color.rgb = BLACK
    conn1.line.width = Pt(2)

    conn2 = slide.shapes.add_connector(MSO_CONNECTOR.CURVE, Inches(6.3), Inches(3.0), Inches(6.3), Inches(4.5))
    conn2.line.color.rgb = BLACK
    conn2.line.width = Pt(2)

    conn3 = slide.shapes.add_connector(MSO_CONNECTOR.CURVE, Inches(9.3), Inches(5.2), Inches(10.0), Inches(5.2))
    conn3.line.color.rgb = BLACK
    conn3.line.width = Pt(2)

    conn4 = slide.shapes.add_connector(MSO_CONNECTOR.CURVE, Inches(11.4), Inches(3.6), Inches(11.4), Inches(4.5))
    conn4.line.color.rgb = BLACK
    conn4.line.width = Pt(2)

    # 6. Page Number & Icon
    page_oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.5), Inches(6.8), Inches(0.6), Inches(0.4))
    page_oval.fill.solid()
    page_oval.fill.fore_color.rgb = WHITE
    page_oval.line.color.rgb = BLACK
    page_oval.line.width = Pt(1.5)
    tf = page_oval.text_frame
    p = tf.paragraphs[0]
    p.text = "09"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = BLACK

    face = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(12.2), Inches(6.6), Inches(0.8), Inches(0.8))
    face.fill.solid()
    face.fill.fore_color.rgb = RGBColor(255, 224, 178)
    face.line.color.rgb = BLACK
    face.line.width = Pt(1.5)
    tf = face.text_frame
    p = tf.paragraphs[0]
    p.text = "博主"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = BLACK



# ── Slide 10 ──

def build_slide_10(slide):
    # 1. Background (Green irregular shape approximated by a rectangle)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.6), Inches(11.7), Inches(6.3))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x2E, 0x9E, 0x7B)
    bg.line.fill.background()

    # 2. Title
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11.333), Inches(1.0))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "总结：桂林，一个越怪越美的地方"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.font.name = "Microsoft YaHei"
    p.font.shadow = True

    # 3. Subtitle
    sub_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(11.333), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "这里的山水有灵，这里的人们有戏"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.font.name = "Microsoft YaHei"
    p.font.shadow = True

    # 4. Item 1: Core Viewpoint
    # Camera Icon (Constructed with shapes)
    cam_base = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.6), Inches(3.4), Inches(0.8), Inches(0.6))
    cam_base.fill.solid()
    cam_base.fill.fore_color.rgb = RGBColor(0xE0, 0xF7, 0xFA)
    cam_base.line.color.rgb = RGBColor(0x00, 0x00, 0x00)
    cam_base.line.width = Pt(2)
    
    cam_lens = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.75), Inches(3.45), Inches(0.5), Inches(0.5))
    cam_lens.fill.solid()
    cam_lens.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    cam_lens.line.color.rgb = RGBColor(0x00, 0x00, 0x00)
    cam_lens.line.width = Pt(2)

    cam_heart = slide.shapes.add_shape(MSO_SHAPE.HEART, Inches(1.85), Inches(3.55), Inches(0.3), Inches(0.3))
    cam_heart.fill.solid()
    cam_heart.fill.fore_color.rgb = RGBColor(0xFF, 0x66, 0x00)
    cam_heart.line.fill.background()

    cam_flash = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.15), Inches(3.48), Inches(0.12), Inches(0.12))
    cam_flash.fill.solid()
    cam_flash.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    cam_flash.line.color.rgb = RGBColor(0x00, 0x00, 0x00)
    cam_flash.line.width = Pt(1)

    # Text
    tb1 = slide.shapes.add_textbox(Inches(2.6), Inches(3.4), Inches(9.0), Inches(0.8))
    p1 = tb1.text_frame.paragraphs[0]
    p1.text = "核心观点：风景是背景，有趣才是旅行的灵魂"
    p1.font.size = Pt(28)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p1.font.name = "Microsoft YaHei"
    p1.font.shadow = True

    # 5. Item 2: Novelty Index
    # Star Icon
    star_icon1 = slide.shapes.add_shape(MSO_SHAPE.5_POINT_STAR, Inches(1.7), Inches(4.5), Inches(0.6), Inches(0.6))
    star_icon1.fill.solid()
    star_icon1.fill.fore_color.rgb = RGBColor(0xFF, 0x66, 0x00)
    star_icon1.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    star_icon1.line.width = Pt(1.5)

    # Text
    tb2 = slide.shapes.add_textbox(Inches(2.6), Inches(4.4), Inches(2.5), Inches(0.8))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = "猎奇指数："
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p2.font.name = "Microsoft YaHei"
    p2.font.shadow = True

    # Rating Stars
    for i in range(5):
        star = slide.shapes.add_shape(MSO_SHAPE.5_POINT_STAR, Inches(4.8 + i*0.7), Inches(4.5), Inches(0.6), Inches(0.6))
        star.fill.solid()
        star.fill.fore_color.rgb = RGBColor(0xFF, 0x66, 0x00)
        star.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        star.line.width = Pt(1.5)

    # 6. Item 3: Recommendation Index
    # Star Icon
    star_icon2 = slide.shapes.add_shape(MSO_SHAPE.5_POINT_STAR, Inches(1.7), Inches(5.5), Inches(0.6), Inches(0.6))
    star_icon2.fill.solid()
    star_icon2.fill.fore_color.rgb = RGBColor(0xFF, 0x66, 0x00)
    star_icon2.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    star_icon2.line.width = Pt(1.5)

    # Text
    tb3 = slide.shapes.add_textbox(Inches(2.6), Inches(5.4), Inches(2.5), Inches(0.8))
    p3 = tb3.text_frame.paragraphs[0]
    p3.text = "推荐指数："
    p3.font.size = Pt(28)
    p3.font.bold = True
    p3.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p3.font.name = "Microsoft YaHei"
    p3.font.shadow = True

    # Rating Stars
    for i in range(5):
        star = slide.shapes.add_shape(MSO_SHAPE.5_POINT_STAR, Inches(4.8 + i*0.7), Inches(5.5), Inches(0.6), Inches(0.6))
        if i < 4:
            star.fill.solid()
            star.fill.fore_color.rgb = RGBColor(0xFF, 0x66, 0x00)
            star.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        else:
            star.fill.solid()
            star.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            star.line.color.rgb = RGBColor(0xFF, 0x66, 0x00)
        star.line.width = Pt(1.5)

    # Comment Text
    tb_comment = slide.shapes.add_textbox(Inches(8.2), Inches(5.45), Inches(4.0), Inches(0.8))
    p_comment = tb_comment.text_frame.paragraphs[0]
    p_comment.text = "(扣一星怕它太骄傲)"
    p_comment.font.size = Pt(24)
    p_comment.font.bold = True
    p_comment.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p_comment.font.name = "Microsoft YaHei"
    p_comment.font.shadow = True

    # 7. Decorations
    # OMG!
    omg_box = slide.shapes.add_textbox(Inches(10.2), Inches(3.2), Inches(2.0), Inches(1.0))
    omg_box.rotation = 15
    p_omg = omg_box.text_frame.paragraphs[0]
    p_omg.text = "OMG!"
    p_omg.font.size = Pt(32)
    p_omg.font.bold = True
    p_omg.font.color.rgb = RGBColor(0xFF, 0x66, 0x00)
    p_omg.font.name = "Arial"

    # Boom!
    boom_box = slide.shapes.add_textbox(Inches(10.5), Inches(4.8), Inches(2.0), Inches(1.0))
    boom_box.rotation = -15
    p_boom = boom_box.text_frame.paragraphs[0]
    p_boom.text = "Boom!"
    p_boom.font.size = Pt(32)
    p_boom.font.bold = True
    p_boom.font.color.rgb = RGBColor(0x80, 0x00, 0x80)
    p_boom.font.name = "Arial"

    # Arrow
    arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(7.6), Inches(4.2), Inches(0.3), Inches(0.5))
    arrow.rotation = 225
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RGBColor(0xFF, 0x33, 0x33)
    arrow.line.fill.background()

    # Decorative circles
    circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.5), Inches(3.1), Inches(0.15), Inches(0.15))
    circle1.fill.background()
    circle1.line.color.rgb = RGBColor(0xFF, 0x66, 0x00)
    circle1.line.width = Pt(2)

    circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.8), Inches(4.8), Inches(0.15), Inches(0.15))
    circle2.fill.background()
    circle2.line.color.rgb = RGBColor(0xFF, 0x33, 0x33)
    circle2.line.width = Pt(2)
    
    # Decorative arc
    arc = slide.shapes.add_shape(MSO_SHAPE.ARC, Inches(10.5), Inches(6.0), Inches(0.8), Inches(0.4))
    arc.rotation = 180
    arc.fill.background()
    arc.line.color.rgb = RGBColor(0x80, 0x00, 0x80)
    arc.line.width = Pt(2)

    # 8. Page Number
    page_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(11.8), Inches(6.5), Inches(1.0), Inches(0.4))
    page_box.fill.solid()
    page_box.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    page_box.line.color.rgb = RGBColor(0x2E, 0x9E, 0x7B)
    page_box.line.width = Pt(2)
    
    tf_page = page_box.text_frame
    tf_page.margin_left = 0
    tf_page.margin_right = 0
    tf_page.margin_top = 0
    tf_page.margin_bottom = 0
    p_page = tf_page.paragraphs[0]
    p_page.text = "第10页"
    p_page.alignment = PP_ALIGN.CENTER
    p_page.font.size = Pt(14)
    p_page.font.bold = True
    p_page.font.color.rgb = RGBColor(0x2E, 0x9E, 0x7B)
    p_page.font.name = "Microsoft YaHei"



# ── Slide 11 ──

def build_slide_11(slide):
    # 1. Background
    # Left background (Light beige)
    bg_left = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg_left.fill.solid()
    bg_left.fill.fore_color.rgb = RGBColor(0xF4, 0xF1, 0xEA)
    bg_left.line.fill.background()

    # Right background (Orange, angled)
    bg_right = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(-3), Inches(10), Inches(14))
    bg_right.rotation = 12
    bg_right.fill.solid()
    bg_right.fill.fore_color.rgb = RGBColor(0xFF, 0x7A, 0x00)
    bg_right.line.fill.background()

    # 2. Top Left Tag ("第11页")
    tag_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0.4), Inches(1.5), Inches(0.6))
    tag_box.fill.solid()
    tag_box.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)
    tag_box.line.fill.background()
    
    tf_tag = tag_box.text_frame
    tf_tag.clear()
    p_tag = tf_tag.paragraphs[0]
    p_tag.text = "第11页"
    p_tag.font.name = "Microsoft YaHei"
    p_tag.font.size = Pt(20)
    p_tag.font.bold = True
    p_tag.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p_tag.alignment = PP_ALIGN.CENTER

    # 3. Titles
    # Main Title Background
    title1_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.8), Inches(0.5), Inches(9.5), Inches(1.0))
    title1_bg.rotation = -2
    title1_bg.fill.solid()
    title1_bg.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    title1_bg.line.fill.background()

    # Main Title Text
    title1_tx = slide.shapes.add_textbox(Inches(2.9), Inches(0.5), Inches(9.3), Inches(1.0))
    title1_tx.rotation = -2
    tf1 = title1_tx.text_frame
    tf1.clear()
    p1 = tf1.paragraphs[0]
    p1.text = "感谢观看，记得给个“五星好评”"
    p1.font.name = "Microsoft YaHei"
    p1.font.size = Pt(40)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(0xE6, 0x4A, 0x19) # Deep Orange

    # Subtitle Background
    title2_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(1.5), Inches(7.5), Inches(0.7))
    title2_bg.fill.solid()
    title2_bg.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    title2_bg.line.fill.background()

    # Subtitle Text
    title2_tx = slide.shapes.add_textbox(Inches(5.8), Inches(1.55), Inches(7.0), Inches(0.6))
    tf2 = title2_tx.text_frame
    tf2.clear()
    p2 = tf2.paragraphs[0]
    p2.text = "关注我，带你解锁更多奇葩目的地"
    p2.font.name = "Microsoft YaHei"
    p2.font.size = Pt(26)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    # 4. Left Visuals (Mountains & River Placeholder)
    # Mountains (Overlapping green triangles)
    colors_mtn = [RGBColor(0x81, 0xC7, 0x84), RGBColor(0xA5, 0xD6, 0xA7), RGBColor(0x66, 0xBB, 0x6A)]
    mtn_coords = [(0.2, 3.5, 2.5, 4.0), (1.5, 2.5, 2.5, 5.0), (3.0, 3.5, 2.5, 4.0)]
    
    for i, (x, y, w, h) in enumerate(mtn_coords):
        mtn = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        mtn.fill.solid()
        mtn.fill.fore_color.rgb = colors_mtn[i]
        mtn.line.color.rgb = RGBColor(0x1B, 0x5E, 0x20)
        mtn.line.width = Pt(2)

    # River (Curved shape)
    river = slide.shapes.add_shape(MSO_SHAPE.MOON, Inches(1.0), Inches(5.0), Inches(4.5), Inches(2.5))
    river.rotation = 50
    river.fill.solid()
    river.fill.fore_color.rgb = RGBColor(0xE0, 0xF2, 0xF1)
    river.line.color.rgb = RGBColor(0x00, 0x4D, 0x40)
    river.line.width = Pt(2)

    # Decorations (Stars)
    star_coords = [(0.5, 2.0), (1.0, 2.8), (4.8, 4.2), (1.2, 6.5)]
    for x, y in star_coords:
        star = slide.shapes.add_shape(MSO_SHAPE.5_POINT_STAR, Inches(x), Inches(y), Inches(0.4), Inches(0.4))
        star.rotation = 15
        star.fill.background()
        star.line.color.rgb = RGBColor(0x9C, 0x27, 0xB0)
        star.line.width = Pt(2)

    # Text "BOOM!"
    boom_tx = slide.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(1.2), Inches(0.5))
    boom_tx.rotation = -15
    p_boom = boom_tx.text_frame.paragraphs[0]
    p_boom.text = "BOOM!"
    p_boom.font.name = "Arial"
    p_boom.font.size = Pt(18)
    p_boom.font.bold = True
    p_boom.font.color.rgb = RGBColor(0x9C, 0x27, 0xB0)

    # Text "OMG!"
    omg_coords = [(4.0, 2.2), (4.8, 3.2), (4.5, 5.5)]
    for x, y in omg_coords:
        omg_tx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(1.0), Inches(0.5))
        omg_tx.rotation = -10
        p_omg = omg_tx.text_frame.paragraphs[0]
        p_omg.text = "OMG!"
        p_omg.font.name = "Arial"
        p_omg.font.size = Pt(16)
        p_omg.font.bold = True
        p_omg.font.color.rgb = RGBColor(0xF4, 0x43, 0x36)

    # 5. Center Visual (Silhouette Placeholder)
    # Head
    head = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.5), Inches(2.6), Inches(0.8), Inches(1.0))
    head.fill.solid()
    head.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)
    head.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    head.line.width = Pt(4)

    # Body
    body = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(3.5), Inches(2.0), Inches(3.5))
    body.fill.solid()
    body.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)
    body.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    body.line.width = Pt(4)

    # Arm (waving)
    arm = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(2.8), Inches(0.6), Inches(1.5))
    arm.rotation = -30
    arm.fill.solid()
    arm.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)
    arm.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    arm.line.width = Pt(4)

    # 6. Right Visual (QR Code Area)
    # Sticky Note (Yellow)
    note = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.0), Inches(1.8), Inches(3.8), Inches(3.8))
    note.rotation = 3
    note.fill.solid()
    note.fill.fore_color.rgb = RGBColor(0xFF, 0xEB, 0x3B)
    note.line.fill.background()

    # QR Code Base (Orange)
    qr_base = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.4), Inches(2.2), Inches(3.0), Inches(3.0))
    qr_base.rotation = 3
    qr_base.fill.solid()
    qr_base.fill.fore_color.rgb = RGBColor(0xFF, 0x98, 0x00)
    qr_base.line.fill.background()

    # QR Code Inner details (Yellow squares)
    qr_inners = [(9.6, 2.4), (11.6, 2.5), (9.7, 4.4)]
    for x, y in qr_inners:
        inner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.6), Inches(0.6))
        inner.rotation = 3
        inner.fill.solid()
        inner.fill.fore_color.rgb = RGBColor(0xFF, 0xEB, 0x3B)
        inner.line.fill.background()

    # Hand-drawn circles around QR code
    circ1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.1), Inches(2.1), Inches(3.6), Inches(1.2))
    circ1.rotation = 3
    circ1.fill.background()
    circ1.line.color.rgb = RGBColor(0x00, 0x00, 0x00)
    circ1.line.width = Pt(1.5)

    circ2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.0), Inches(4.3), Inches(1.5), Inches(0.8))
    circ2.rotation = -5
    circ2.fill.background()
    circ2.line.color.rgb = RGBColor(0x00, 0x00, 0x00)
    circ2.line.width = Pt(1.5)

    # 7. Bottom Right Text Block ("内容要点")
    # Title
    tx_points_title = slide.shapes.add_textbox(Inches(8.8), Inches(5.3), Inches(3.0), Inches(0.5))
    p_title = tx_points_title.text_frame.paragraphs[0]
    p_title.text = "内容要点"
    p_title.font.name = "Microsoft YaHei"
    p_title.font.size = Pt(22)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    # Points Data
    points_data = [
        ("?", "互动区：你见过最奇葩的景点在哪里？"),
        ("@", "联系方式：微博/小红书/抖音同名"),
        ("!", "结束语：山水不改，脑洞常在！")
    ]

    start_y = 5.9
    for icon_char, text_content in points_data:
        # Icon circle
        icon_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.8), Inches(start_y), Inches(0.3), Inches(0.3))
        icon_shape.fill.solid()
        icon_shape.fill.fore_color.rgb = RGBColor(0xFF, 0xEB, 0x3B)
        icon_shape.line.color.rgb = RGBColor(0x00, 0x00, 0x00)
        icon_shape.line.width = Pt(1)
        
        icon_p = icon_shape.text_frame.paragraphs[0]
        icon_p.text = icon_char
        icon_p.font.name = "Arial"
        icon_p.font.size = Pt(12)
        icon_p.font.bold = True
        icon_p.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
        icon_p.alignment = PP_ALIGN.CENTER

        # Text
        tx_point = slide.shapes.add_textbox(Inches(9.2), Inches(start_y - 0.05), Inches(4.0), Inches(0.4))
        p_point = tx_point.text_frame.paragraphs[0]
        p_point.text = text_content
        p_point.font.name = "Microsoft YaHei"
        p_point.font.size = Pt(14)
        p_point.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
        
        start_y += 0.45



# ── Main ──
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

s0 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_1(s0)
s1 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_2(s1)
s2 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_3(s2)
s3 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_4(s3)
s4 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_5(s4)
s5 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_6(s5)
s6 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_7(s6)
s7 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_8(s7)
s8 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_9(s8)
s9 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_10(s9)
s10 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_11(s10)
prs.save(OUTPUT_PATH)
