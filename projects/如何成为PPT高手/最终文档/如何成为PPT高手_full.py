from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# ── Color Palette ──
BLUE_HEADER = RGBColor(0x5B, 0x9B, 0xD5)
BLUE_DARK   = RGBColor(0x4A, 0x86, 0xC8)
CYAN        = RGBColor(0x00, 0xBC, 0xD4)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x33, 0x33, 0x33)
GRAY_TEXT   = RGBColor(0x55, 0x55, 0x55)
GRAY_BAR    = RGBColor(0xB0, 0xBE, 0xC5)
RED         = RGBColor(0xE5, 0x39, 0x35)
GREEN       = RGBColor(0x43, 0xA0, 0x47)
ORANGE      = RGBColor(0xFF, 0x98, 0x00)

ICON_BG  = RGBColor(0xE3, 0xE8, 0xED)
ICON_FG  = RGBColor(0x54, 0x6E, 0x7A)
FONT_NAME = "Microsoft YaHei"
SLIDE_WIDTH = Inches(13.333)
HEADER_H    = Inches(0.75)
SUBTITLE_Y  = Inches(0.95)


def add_header_banner(slide, title_text, bg_color=None):
    if bg_color is None:
        bg_color = BLUE_HEADER
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, HEADER_H
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = bg_color
    banner.line.fill.background()
    tf = banner.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(26)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME


def add_subtitle(slide, text, left, top, width=Inches(12), font_size=Pt(18)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.font.name = FONT_NAME
    return txBox


def add_icon_box(slide, left, top, symbol, size=Inches(0.48)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ICON_BG
    shape.line.fill.background()
    shape.adjustments[0] = 0.25
    tf = shape.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = symbol
    p.font.size = Pt(18)
    p.font.color.rgb = ICON_FG
    p.font.bold = False
    return shape


def add_bullet_item(slide, left, top, symbol, label, description,
                    width=Inches(5.5), desc_size=Pt(13)):
    add_icon_box(slide, left, top, symbol)
    text_left = left + Inches(0.65)
    txBox = slide.shapes.add_textbox(text_left, top - Inches(0.02), width, Inches(0.65))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run_label = p.add_run()
    run_label.text = label + "\uff1a"
    run_label.font.size = Pt(14)
    run_label.font.color.rgb = BLACK
    run_label.font.bold = True
    run_label.font.name = FONT_NAME
    run_desc = p.add_run()
    run_desc.text = description
    run_desc.font.size = desc_size
    run_desc.font.color.rgb = GRAY_TEXT
    run_desc.font.bold = False
    run_desc.font.name = FONT_NAME
    return txBox


def add_conclusion_box(slide, left, top, width, text, font_size=Pt(13)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = BLACK
    run.font.bold = True
    run.font.name = FONT_NAME
    return txBox


def add_table(slide, left, top, width, height, rows, cols, data,
              header_color=None, col_widths=None):
    if header_color is None:
        header_color = BLUE_HEADER
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if data[r][c] is not None else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(11)
                    run.font.name = FONT_NAME
                    if r == 0:
                        run.font.color.rgb = WHITE
                        run.font.bold = True
                    else:
                        run.font.color.rgb = BLACK
                        run.font.bold = False
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else RGBColor(0xF5, 0xF5, 0xF5)
            cell.margin_left = Pt(5)
            cell.margin_right = Pt(5)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
    return table_shape


def add_bar_chart(slide, left, top, width, height,
                  categories, values, title="", bar_colors=None):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', values)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    chart.chart_style = 2
    plot = chart.plots[0]
    plot.gap_width = 100
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CYAN
    series.has_data_labels = True
    dl = series.data_labels
    dl.font.size = Pt(13)
    dl.font.bold = True
    dl.font.color.rgb = BLACK
    dl.number_format = '0.#'
    dl.show_value = True
    dl.label_position = XL_LABEL_POSITION.OUTSIDE_END
    if bar_colors:
        for i, color in enumerate(bar_colors):
            pt = series.points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = color
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(12)
    cat_axis.tick_labels.font.name = FONT_NAME
    cat_axis.major_tick_mark = 2
    cat_axis.format.line.fill.background()
    val_axis = chart.value_axis
    val_axis.visible = False
    val_axis.major_tick_mark = 2
    val_axis.format.line.fill.background()
    val_axis.major_gridlines.format.line.fill.background()
    if title:
        chart.has_title = True
        ct = chart.chart_title.text_frame.paragraphs[0]
        ct.text = title
        ct.font.size = Pt(14)
        ct.font.bold = True
        ct.font.name = FONT_NAME
    else:
        chart.has_title = False
    return chart_frame


def add_callout_label(slide, left, top, text, bg_color=None, font_size=Pt(11)):
    if bg_color is None:
        bg_color = CYAN
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.3), Inches(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME
    return shape


def add_data_card(slide, left, top, width, height, value, label,
                  value_color=None, bg_color=None):
    if value_color is None:
        value_color = CYAN
    if bg_color is None:
        bg_color = WHITE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(3)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run1 = p1.add_run()
    run1.text = str(value)
    run1.font.size = Pt(24)
    run1.font.color.rgb = value_color
    run1.font.bold = True
    run1.font.name = FONT_NAME
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(10)
    run2.font.color.rgb = GRAY_TEXT
    run2.font.bold = False
    run2.font.name = FONT_NAME
    return shape


OUTPUT_PATH = r"C:\Users\Administrator\Desktop\notebook\my_notebooklm_with_ppteditable\projects\如何成为PPT高手\最终文档\如何成为PPT高手.pptx"


# ── Slide 1 ──

def build_slide_1(slide):
    # 1. Background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RGBColor(0x05, 0x22, 0x55) # Dark blue
    bg_shape.line.fill.background()

    # Grid lines
    grid_color = RGBColor(0x10, 0x35, 0x70)
    for i in range(1, 14):
        line = slide.shapes.add_connector(1, Inches(i), Inches(0), Inches(i), Inches(7.5))
        line.line.color.rgb = grid_color
        line.line.width = Pt(0.5)
    for i in range(1, 8):
        line = slide.shapes.add_connector(1, Inches(0), Inches(i), Inches(13.333), Inches(i))
        line.line.color.rgb = grid_color
        line.line.width = Pt(0.5)

    # 2. Central Lightbulb Icon
    center_x = 13.333 / 2
    bulb_y = 1.0
    bulb_color = RGBColor(0xAA, 0xD4, 0xFF)
    
    # Glow effect (simulated with a larger faint circle)
    glow = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(center_x - 1.0), Inches(bulb_y - 0.4), Inches(2.0), Inches(2.0))
    glow.fill.solid()
    glow.fill.fore_color.rgb = RGBColor(0x15, 0x45, 0x85)
    glow.line.fill.background()

    # Main bulb (Oval)
    bulb = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(center_x - 0.6), Inches(bulb_y), Inches(1.2), Inches(1.2))
    bulb.fill.background()
    bulb.line.color.rgb = bulb_color
    bulb.line.width = Pt(3)
    
    # Bulb base (Rectangle)
    base = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(center_x - 0.3), Inches(bulb_y + 1.15), Inches(0.6), Inches(0.4))
    base.fill.background()
    base.line.color.rgb = bulb_color
    base.line.width = Pt(3)
    
    # Base screw lines
    for i in range(3):
        line_y = bulb_y + 1.6 + i * 0.12
        line = slide.shapes.add_connector(1, Inches(center_x - 0.25), Inches(line_y), Inches(center_x + 0.25), Inches(line_y))
        line.line.color.rgb = bulb_color
        line.line.width = Pt(2.5)
        
    # Filament (Inner lines)
    fil_left = slide.shapes.add_connector(1, Inches(center_x - 0.2), Inches(bulb_y + 1.15), Inches(center_x - 0.2), Inches(bulb_y + 0.6))
    fil_left.line.color.rgb = bulb_color
    fil_left.line.width = Pt(2)
    
    fil_right = slide.shapes.add_connector(1, Inches(center_x + 0.2), Inches(bulb_y + 1.15), Inches(center_x + 0.2), Inches(bulb_y + 0.6))
    fil_right.line.color.rgb = bulb_color
    fil_right.line.width = Pt(2)
    
    fil_top = slide.shapes.add_shape(MSO_SHAPE.ARC, Inches(center_x - 0.2), Inches(bulb_y + 0.4), Inches(0.4), Inches(0.4))
    fil_top.rotation = 180
    fil_top.fill.background()
    fil_top.line.color.rgb = bulb_color
    fil_top.line.width = Pt(2)

    # 3. Main Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(11.333), Inches(1))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "简而不凡：高效演示文稿的制作之道"
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # 4. Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(2), Inches(4.2), Inches(9.333), Inches(0.8))
    tf_sub = subtitle_box.text_frame
    tf_sub.clear()
    p_sub = tf_sub.paragraphs[0]
    p_sub.alignment = PP_ALIGN.CENTER
    run_sub = p_sub.add_run()
    run_sub.text = "掌握专业PPT的核心逻辑与设计法则"
    run_sub.font.name = "Microsoft YaHei"
    run_sub.font.size = Pt(24)
    run_sub.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # 5. Bullet Points
    bullet_texts = [
        "1. 演示文稿不仅仅是工具，更是思维的视觉化表达",
        "2. 核心目标：降低沟通成本，提升说服力",
        "3. 专家级PPT的两大底层支柱：内容清晰与设计统一"
    ]
    
    start_y = 5.2
    spacing = 0.6
    icon_x = 3.8
    text_x = 4.3
    icon_color = RGBColor(0xAA, 0xD4, 0xFF)
    
    for i, text in enumerate(bullet_texts):
        icon_y = start_y + i * spacing + 0.08
        
        # Draw Icons
        if i == 0: # Idea (Lightbulb)
            icon_bulb = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(icon_x), Inches(icon_y), Inches(0.2), Inches(0.2))
            icon_bulb.fill.background()
            icon_bulb.line.color.rgb = icon_color
            icon_bulb.line.width = Pt(1.5)
            icon_base = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(icon_x + 0.05), Inches(icon_y + 0.18), Inches(0.1), Inches(0.08))
            icon_base.fill.background()
            icon_base.line.color.rgb = icon_color
            icon_base.line.width = Pt(1.5)
            # Rays
            slide.shapes.add_connector(1, Inches(icon_x+0.1), Inches(icon_y-0.05), Inches(icon_x+0.1), Inches(icon_y-0.1)).line.color.rgb = icon_color
            slide.shapes.add_connector(1, Inches(icon_x-0.05), Inches(icon_y+0.1), Inches(icon_x-0.1), Inches(icon_y+0.1)).line.color.rgb = icon_color
            slide.shapes.add_connector(1, Inches(icon_x+0.25), Inches(icon_y+0.1), Inches(icon_x+0.3), Inches(icon_y+0.1)).line.color.rgb = icon_color

        elif i == 1: # Target
            icon1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(icon_x), Inches(icon_y), Inches(0.25), Inches(0.25))
            icon1.fill.background()
            icon1.line.color.rgb = icon_color
            icon1.line.width = Pt(1.5)
            icon2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(icon_x+0.075), Inches(icon_y+0.075), Inches(0.1), Inches(0.1))
            icon2.fill.background()
            icon2.line.color.rgb = icon_color
            icon2.line.width = Pt(1.5)
            # Arrow
            arrow = slide.shapes.add_connector(1, Inches(icon_x+0.15), Inches(icon_y+0.1), Inches(icon_x+0.35), Inches(icon_y-0.1))
            arrow.line.color.rgb = icon_color
            arrow.line.width = Pt(1.5)

        elif i == 2: # Balance
            # Base triangle
            tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(icon_x+0.05), Inches(icon_y+0.1), Inches(0.15), Inches(0.15))
            tri.fill.background()
            tri.line.color.rgb = icon_color
            tri.line.width = Pt(1.5)
            # Top bar
            bar = slide.shapes.add_connector(1, Inches(icon_x-0.05), Inches(icon_y+0.1), Inches(icon_x+0.3), Inches(icon_y+0.1))
            bar.line.color.rgb = icon_color
            bar.line.width = Pt(1.5)
            # Left pan
            slide.shapes.add_connector(1, Inches(icon_x-0.05), Inches(icon_y+0.1), Inches(icon_x-0.05), Inches(icon_y+0.25)).line.color.rgb = icon_color
            slide.shapes.add_connector(1, Inches(icon_x-0.1), Inches(icon_y+0.25), Inches(icon_x), Inches(icon_y+0.25)).line.color.rgb = icon_color
            # Right pan
            slide.shapes.add_connector(1, Inches(icon_x+0.3), Inches(icon_y+0.1), Inches(icon_x+0.3), Inches(icon_y+0.25)).line.color.rgb = icon_color
            slide.shapes.add_connector(1, Inches(icon_x+0.25), Inches(icon_y+0.25), Inches(icon_x+0.35), Inches(icon_y+0.25)).line.color.rgb = icon_color

        # Text
        text_box = slide.shapes.add_textbox(Inches(text_x), Inches(start_y + i * spacing), Inches(8), Inches(0.4))
        tf_bullet = text_box.text_frame
        tf_bullet.clear()
        p_bullet = tf_bullet.paragraphs[0]
        run_bullet = p_bullet.add_run()
        run_bullet.text = text
        run_bullet.font.name = "Microsoft YaHei"
        run_bullet.font.size = Pt(16)
        run_bullet.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)



# ── Slide 2 ──

def build_slide_2(slide):
    from pptx.enum.shapes import MSO_CONNECTOR
    
    # Colors
    BLUE_TITLE = RGBColor(0x1A, 0x56, 0xBA)
    DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
    BLUE_LINE = RGBColor(0x3B, 0x7E, 0xC6)
    ORANGE_TEXT = RGBColor(0xDE, 0x9B, 0x35)
    BLACK_TEXT = RGBColor(0x00, 0x00, 0x00)
    SHADOW_COLOR = RGBColor(0xF0, 0xF4, 0xF8)
    BORDER_COLOR = RGBColor(0xE8, 0xE8, 0xE8)

    # 1. Add Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.0), Inches(0.8))
    tf_title = title_box.text_frame
    p_title = tf_title.paragraphs[0]
    p_title.text = "目录：构建专业PPT的蓝图"
    p_title.font.name = "Microsoft YaHei"
    p_title.font.size = Pt(36)
    p_title.font.bold = True
    p_title.font.color.rgb = BLUE_TITLE

    # 2. Add Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.0), Inches(0.5))
    tf_sub = sub_box.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "本次分享的核心框架"
    p_sub.font.name = "Microsoft YaHei"
    p_sub.font.size = Pt(20)
    p_sub.font.bold = True
    p_sub.font.color.rgb = DARK_GRAY

    # Helper function to create styled text boxes
    def add_node_box(left, top, width, height, text_parts):
        # Shadow
        shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.06), top + Inches(0.06), width, height)
        shadow.fill.solid()
        shadow.fill.fore_color.rgb = SHADOW_COLOR
        shadow.line.fill.background()

        # Main Box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        box.line.color.rgb = BORDER_COLOR
        box.line.width = Pt(1)

        tf = box.text_frame
        tf.margin_left = Inches(0.3)
        tf.margin_top = Inches(0.15)
        p = tf.paragraphs[0]

        for text, color, is_bold in text_parts:
            run = p.add_run()
            run.text = text
            run.font.name = "Microsoft YaHei"
            run.font.size = Pt(18)
            run.font.color.rgb = color
            run.font.bold = is_bold

    # 3. Main Vertical Timeline
    v_line_x = 1.8
    y_nodes = [2.6, 3.8, 5.0, 6.2]
    
    main_v_line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(v_line_x), Inches(y_nodes[0]), Inches(v_line_x), Inches(y_nodes[3]))
    main_v_line.line.color.rgb = BLUE_LINE
    main_v_line.line.width = Pt(1.5)

    # --- Node 1 ---
    y1 = y_nodes[0]
    # Circle
    c1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(v_line_x - 0.08), Inches(y1 - 0.08), Inches(0.16), Inches(0.16))
    c1.fill.solid()
    c1.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    c1.line.color.rgb = BLUE_LINE
    c1.line.width = Pt(2)
    # H-Line
    hl1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(v_line_x + 0.08), Inches(y1), Inches(2.0), Inches(y1))
    hl1.line.color.rgb = BLUE_LINE
    hl1.line.width = Pt(1.5)
    # Text Box
    add_node_box(Inches(2.0), Inches(y1 - 0.35), Inches(5.8), Inches(0.7), [
        ("1. 内容法则：一页一事，", BLACK_TEXT, True),
        ("结论先行", ORANGE_TEXT, True)
    ])
    # Icon 1 (Document)
    doc = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(y1 - 0.4), Inches(0.5), Inches(0.7))
    doc.fill.background()
    doc.line.color.rgb = BLUE_LINE
    doc.line.width = Pt(2)
    for i in range(3):
        l = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.9), Inches(y1 - 0.2 + i*0.15), Inches(1.2), Inches(y1 - 0.2 + i*0.15))
        l.line.color.rgb = BLUE_LINE
        l.line.width = Pt(1.5)
    check_bg = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.15), Inches(y1 - 0.05), Inches(0.3), Inches(0.3))
    check_bg.fill.solid()
    check_bg.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    check_bg.line.color.rgb = BLUE_LINE
    check_bg.line.width = Pt(2)
    ck1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.22), Inches(y1 + 0.1), Inches(1.28), Inches(y1 + 0.16))
    ck1.line.color.rgb = BLUE_LINE
    ck1.line.width = Pt(2)
    ck2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.28), Inches(y1 + 0.16), Inches(1.38), Inches(y1 + 0.02))
    ck2.line.color.rgb = BLUE_LINE
    ck2.line.width = Pt(2)

    # --- Node 2 ---
    y2 = y_nodes[1]
    # Circle
    c2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(v_line_x - 0.08), Inches(y2 - 0.08), Inches(0.16), Inches(0.16))
    c2.fill.solid()
    c2.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    c2.line.color.rgb = BLUE_LINE
    c2.line.width = Pt(2)
    # H-Line
    hl2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(v_line_x + 0.08), Inches(y2), Inches(2.4), Inches(y2))
    hl2.line.color.rgb = BLUE_LINE
    hl2.line.width = Pt(1.5)
    # Text Box
    add_node_box(Inches(3.4), Inches(y2 - 0.35), Inches(6.4), Inches(0.7), [
        ("2. 减法艺术：拒绝文字堆砌，追求", BLACK_TEXT, True),
        ("秒懂", ORANGE_TEXT, True)
    ])
    # Icon 2 (Trash)
    trash = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.45), Inches(y2 - 0.2), Inches(0.35), Inches(0.45))
    trash.fill.background()
    trash.line.color.rgb = BLUE_LINE
    trash.line.width = Pt(2)
    lid = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.35), Inches(y2 - 0.3), Inches(0.55), Inches(0.08))
    lid.fill.background()
    lid.line.color.rgb = BLUE_LINE
    lid.line.width = Pt(2)
    handle = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.55), Inches(y2 - 0.38), Inches(0.15), Inches(0.08))
    handle.fill.background()
    handle.line.color.rgb = BLUE_LINE
    handle.line.width = Pt(1.5)
    for i in range(3):
        vl = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(2.55 + i*0.08), Inches(y2 - 0.15), Inches(2.55 + i*0.08), Inches(y2 + 0.2))
        vl.line.color.rgb = BLUE_LINE
        vl.line.width = Pt(1.5)
    for i in range(3):
        hl = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(2.95), Inches(y2 - 0.1 + i*0.15), Inches(3.25), Inches(y2 - 0.1 + i*0.15))
        hl.line.color.rgb = BLUE_LINE
        hl.line.width = Pt(2)

    # --- Node 3 ---
    y3 = y_nodes[2]
    # Circle
    c3 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(v_line_x - 0.08), Inches(y3 - 0.08), Inches(0.16), Inches(0.16))
    c3.fill.solid()
    c3.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    c3.line.color.rgb = BLUE_LINE
    c3.line.width = Pt(2)
    # H-Line
    hl3 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(v_line_x + 0.08), Inches(y3), Inches(3.4), Inches(y3))
    hl3.line.color.rgb = BLUE_LINE
    hl3.line.width = Pt(1.5)
    # Text Box
    add_node_box(Inches(4.3), Inches(y3 - 0.35), Inches(6.4), Inches(0.7), [
        ("3. 设计规范：高度统一，建立", BLACK_TEXT, True),
        ("专业感", ORANGE_TEXT, True)
    ])
    # Icon 3 (Gears & Ruler)
    g1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.4), Inches(y3 - 0.4), Inches(0.4), Inches(0.4))
    g1.fill.background()
    g1.line.color.rgb = BLUE_LINE
    g1.line.width = Pt(2.5)
    g1_in = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.55), Inches(y3 - 0.25), Inches(0.1), Inches(0.1))
    g1_in.fill.background()
    g1_in.line.color.rgb = BLUE_LINE
    g1_in.line.width = Pt(1.5)
    g2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.75), Inches(y3 - 0.25), Inches(0.3), Inches(0.3))
    g2.fill.background()
    g2.line.color.rgb = BLUE_LINE
    g2.line.width = Pt(2)
    ruler = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.4), Inches(y3 + 0.1), Inches(0.7), Inches(0.15))
    ruler.fill.background()
    ruler.line.color.rgb = BLUE_LINE
    ruler.line.width = Pt(1.5)
    for i in range(6):
        tick = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(3.45 + i*0.1), Inches(y3 + 0.1), Inches(3.45 + i*0.1), Inches(y3 + 0.18))
        tick.line.color.rgb = BLUE_LINE

    # --- Node 4 ---
    y4 = y_nodes[3]
    # Circle
    c4 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(v_line_x - 0.08), Inches(y4 - 0.08), Inches(0.16), Inches(0.16))
    c4.fill.solid()
    c4.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    c4.line.color.rgb = BLUE_LINE
    c4.line.width = Pt(2)
    # H-Line
    hl4 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(v_line_x + 0.08), Inches(y4), Inches(4.2), Inches(y4))
    hl4.line.color.rgb = BLUE_LINE
    hl4.line.width = Pt(1.5)
    # Text Box
    add_node_box(Inches(5.3), Inches(y4 - 0.35), Inches(6.3), Inches(0.7), [
        ("4. 高手境界：简洁有力的视觉哲学", BLACK_TEXT, True)
    ])
    # Icon 4 (Mountain & Scale)
    mt = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(4.2), Inches(y4 - 0.3), Inches(0.65), Inches(0.5))
    mt.fill.background()
    mt.line.color.rgb = BLUE_LINE
    mt.line.width = Pt(2)
    snow = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(4.42), Inches(y4 - 0.3), Inches(0.21), Inches(0.15))
    snow.fill.background()
    snow.line.color.rgb = BLUE_LINE
    snow.line.width = Pt(1.5)
    
    base = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(5.15), Inches(y4 + 0.1), Inches(0.3), Inches(0.1))
    base.fill.background()
    base.line.color.rgb = BLUE_LINE
    base.line.width = Pt(1.5)
    post = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(5.3), Inches(y4 - 0.3), Inches(5.3), Inches(y4 + 0.1))
    post.line.color.rgb = BLUE_LINE
    post.line.width = Pt(2)
    beam = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(5.0), Inches(y4 - 0.2), Inches(5.6), Inches(y4 - 0.2))
    beam.line.color.rgb = BLUE_LINE
    beam.line.width = Pt(2)
    
    p1_v = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(5.05), Inches(y4 - 0.2), Inches(5.05), Inches(y4))
    p1_v.line.color.rgb = BLUE_LINE
    p1_h = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(4.95), Inches(y4), Inches(5.15), Inches(y4))
    p1_h.line.color.rgb = BLUE_LINE
    p1_h.line.width = Pt(1.5)
    
    p2_v = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(5.55), Inches(y4 - 0.2), Inches(5.55), Inches(y4))
    p2_v.line.color.rgb = BLUE_LINE
    p2_h = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(5.45), Inches(y4), Inches(5.65), Inches(y4))
    p2_h.line.color.rgb = BLUE_LINE
    p2_h.line.width = Pt(1.5)

    # 4. Page Number
    page_num = slide.shapes.add_textbox(Inches(12.5), Inches(6.8), Inches(0.5), Inches(0.5))
    tf_num = page_num.text_frame
    p_num = tf_num.paragraphs[0]
    p_num.text = "02"
    p_num.font.name = "Arial"
    p_num.font.size = Pt(14)
    p_num.font.color.rgb = RGBColor(0x66, 0x66, 0x66)



# ── Slide 3 ──

def build_slide_3(slide):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
    from pptx.enum.text import PP_ALIGN

    # Colors
    DARK_BLUE = RGBColor(0x1F, 0x4E, 0x96)
    HIGHLIGHT_BLUE = RGBColor(0x00, 0x70, 0xC0)
    TEXT_BLACK = RGBColor(0x33, 0x33, 0x33)
    TEXT_GRAY = RGBColor(0x66, 0x66, 0x66)
    ORANGE_FILL = RGBColor(0xFF, 0x6B, 0x00)
    ORANGE_LINE = RGBColor(0xE6, 0x51, 0x00)

    # 1. Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(10.0), Inches(0.8))
    tf_title = title_box.text_frame
    p_title = tf_title.paragraphs[0]
    p_title.text = "法则一：内容清晰是PPT的灵魂"
    p_title.font.name = "Microsoft YaHei"
    p_title.font.size = Pt(32)
    p_title.font.bold = True
    p_title.font.color.rgb = DARK_BLUE

    # 2. Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(10.0), Inches(0.5))
    tf_sub = sub_box.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "确保观众的注意力始终聚焦"
    p_sub.font.name = "Microsoft YaHei"
    p_sub.font.size = Pt(18)
    p_sub.font.color.rgb = TEXT_BLACK

    # 3. Central Graphic (Target)
    cx, cy = 5.2, 4.2
    radii = [2.2, 1.75, 1.3, 0.85, 0.4]
    fills = [
        RGBColor(226, 238, 252),
        RGBColor(204, 224, 250),
        RGBColor(182, 210, 248),
        RGBColor(160, 196, 246),
        ORANGE_FILL
    ]
    
    for i, (r, fill) in enumerate(zip(radii, fills)):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, 
            Inches(cx - r), Inches(cy - r), 
            Inches(r * 2), Inches(r * 2)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        if i < 4:
            shape.line.color.rgb = DARK_BLUE
            shape.line.width = Pt(2.5)
        else:
            shape.line.fill.background() # No line for bullseye

    # 4. Central Graphic (Arrow)
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.NOTCHED_RIGHT_ARROW, 
        Inches(cx + 0.4), Inches(cy - 2.4), 
        Inches(2.4), Inches(0.7)
    )
    arrow.rotation = 135
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RGBColor(255, 167, 38)
    arrow.line.color.rgb = ORANGE_LINE
    arrow.line.width = Pt(1.5)

    # 5. Right Side Content - Item 1
    # Icon 1: Magnifying Glass
    mag_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.2), Inches(2.5), Inches(0.28), Inches(0.28))
    mag_circle.fill.background()
    mag_circle.line.color.rgb = DARK_BLUE
    mag_circle.line.width = Pt(2)
    handle = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(8.44), Inches(2.74), Inches(8.6), Inches(2.9))
    handle.line.color.rgb = DARK_BLUE
    handle.line.width = Pt(2.5)

    # Text 1
    tx_box1 = slide.shapes.add_textbox(Inches(8.8), Inches(2.35), Inches(4.0), Inches(1.0))
    tf1 = tx_box1.text_frame
    p1 = tf1.paragraphs[0]
    p1.font.size = Pt(16)
    p1.font.name = "Microsoft YaHei"
    r1_1 = p1.add_run(); r1_1.text = "每一页幻灯片只传达一\n"; r1_1.font.color.rgb = TEXT_BLACK
    r1_2 = p1.add_run(); r1_2.text = "个核心观点"; r1_2.font.color.rgb = HIGHLIGHT_BLUE; r1_2.font.bold = True

    # 6. Right Side Content - Item 2
    # Icon 2: Hierarchy
    box_w, box_h = 0.16, 0.12
    b1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.15), Inches(4.04), Inches(box_w), Inches(box_h))
    b2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.45), Inches(3.84), Inches(box_w), Inches(box_h))
    b3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.45), Inches(4.04), Inches(box_w), Inches(box_h))
    b4 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.45), Inches(4.24), Inches(box_w), Inches(box_h))
    
    for b in [b1, b2, b3, b4]:
        b.fill.background()
        b.line.color.rgb = DARK_BLUE
        b.line.width = Pt(1.5)

    v_line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(8.38), Inches(3.9), Inches(8.38), Inches(4.3))
    h1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(8.31), Inches(4.1), Inches(8.38), Inches(4.1))
    h2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(8.38), Inches(3.9), Inches(8.45), Inches(3.9))
    h3 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(8.38), Inches(4.1), Inches(8.45), Inches(4.1))
    h4 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(8.38), Inches(4.3), Inches(8.45), Inches(4.3))
    
    for l in [v_line, h1, h2, h3, h4]:
        l.line.color.rgb = DARK_BLUE
        l.line.width = Pt(1.5)

    # Text 2
    tx_box2 = slide.shapes.add_textbox(Inches(8.8), Inches(3.75), Inches(4.0), Inches(1.0))
    tf2 = tx_box2.text_frame
    p2 = tf2.paragraphs[0]
    p2.font.size = Pt(16)
    p2.font.name = "Microsoft YaHei"
    r2_1 = p2.add_run(); r2_1.text = "复杂问题"; r2_1.font.color.rgb = TEXT_BLACK
    r2_2 = p2.add_run(); r2_2.text = "拆解化"; r2_2.font.color.rgb = HIGHLIGHT_BLUE; r2_2.font.bold = True
    r2_3 = p2.add_run(); r2_3.text = "，"; r2_3.font.color.rgb = TEXT_BLACK
    r2_4 = p2.add_run(); r2_4.text = "单一\n"; r2_4.font.color.rgb = HIGHLIGHT_BLUE; r2_4.font.bold = True
    r2_5 = p2.add_run(); r2_5.text = "观点"; r2_5.font.color.rgb = HIGHLIGHT_BLUE; r2_5.font.bold = True
    r2_6 = p2.add_run(); r2_6.text = "深度化"; r2_6.font.color.rgb = TEXT_BLACK

    # 7. Right Side Content - Item 3
    # Icon 3: Pyramid/Upload
    up_arrow = slide.shapes.add_shape(MSO_SHAPE.UP_ARROW, Inches(8.3), Inches(5.35), Inches(0.18), Inches(0.22))
    up_arrow.fill.background()
    up_arrow.line.color.rgb = DARK_BLUE
    up_arrow.line.width = Pt(1.5)
    
    base = slide.shapes.add_shape(MSO_SHAPE.TRAPEZOID, Inches(8.15), Inches(5.62), Inches(0.48), Inches(0.15))
    base.fill.background()
    base.line.color.rgb = DARK_BLUE
    base.line.width = Pt(1.5)

    # Text 3
    tx_box3 = slide.shapes.add_textbox(Inches(8.8), Inches(5.25), Inches(4.0), Inches(1.0))
    tf3 = tx_box3.text_frame
    p3 = tf3.paragraphs[0]
    p3.font.size = Pt(16)
    p3.font.name = "Microsoft YaHei"
    r3_1 = p3.add_run(); r3_1.text = "结论先行"; r3_1.font.color.rgb = HIGHLIGHT_BLUE; r3_1.font.bold = True
    r3_2 = p3.add_run(); r3_2.text = "：标题即观点，\n内容即支撑"; r3_2.font.color.rgb = TEXT_BLACK

    # 8. Page Number
    page_num = slide.shapes.add_textbox(Inches(12.2), Inches(6.8), Inches(0.8), Inches(0.4))
    tf_page = page_num.text_frame
    p_page = tf_page.paragraphs[0]
    p_page.text = "3 / 11"
    p_page.font.size = Pt(12)
    p_page.font.color.rgb = TEXT_GRAY
    p_page.font.name = "Microsoft YaHei"
    p_page.alignment = PP_ALIGN.RIGHT



# ── Slide 4 ──

def build_slide_4(slide):
    # 自定义颜色
    LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
    DENSE_TEXT_COLOR = RGBColor(0x99, 0x99, 0x99)
    
    # 1. 标题
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(10), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "减法艺术：拒绝文字堆砌"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK
    p.font.name = "Microsoft YaHei"

    # 2. 副标题
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(10), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "别让你的PPT变成Word搬家"
    p.font.size = Pt(22)
    p.font.color.rgb = GRAY_TEXT
    p.font.name = "Microsoft YaHei"

    # 3. 左侧栏 (错误示例)
    # 标题
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(5.5), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Word搬家（错误示例）"
    p.font.size = Pt(18)
    p.font.color.rgb = RED
    p.font.name = "Microsoft YaHei"
    
    # 红色下划线
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.5), Inches(2.7), Inches(6.0), Inches(2.7))
    line.line.color.rgb = RED
    line.line.width = Pt(1.5)

    # 密集文本块
    dense_text = "这里是一段非常长且密集的文字，用来模拟将Word文档直接复制粘贴到PPT中的错误做法。在实际的演示中，观众根本无法在短时间内阅读并理解这么多文字。这种做法不仅会让幻灯片显得杂乱无章，还会严重分散观众的注意力，导致他们无法专心听讲。优秀的PPT应该只保留核心观点和关键词，通过演讲者的口述来补充细节。如果把所有内容都写在屏幕上，那么演讲者就失去了存在的意义，PPT也就变成了一份阅读材料而不是辅助演示的工具。因此，我们必须学会做减法，拒绝文字堆砌，提炼出最精炼的信息，用视觉化的方式呈现出来，从而提高沟通的效率和效果。" * 4
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.9), Inches(5.5), Inches(3.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = dense_text
    p.font.size = Pt(8)
    p.font.color.rgb = DENSE_TEXT_COLOR
    p.font.name = "Microsoft YaHei"
    p.alignment = PP_ALIGN.JUSTIFY

    # 红色大叉
    cross_center_x = 3.25
    cross_center_y = 4.6
    cross_size = 1.2
    line1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(cross_center_x - cross_size), Inches(cross_center_y - cross_size), Inches(cross_center_x + cross_size), Inches(cross_center_y + cross_size))
    line1.line.color.rgb = RED
    line1.line.width = Pt(25)
    line2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(cross_center_x - cross_size), Inches(cross_center_y + cross_size), Inches(cross_center_x + cross_size), Inches(cross_center_y - cross_size))
    line2.line.color.rgb = RED
    line2.line.width = Pt(25)

    # 左侧底部结论
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(6.0), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "密密麻麻的文字，信息过载，观众无法聚焦。"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLACK
    p.font.name = "Microsoft YaHei"

    # 4. 中间垂直分割线
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(6.66), Inches(2.2), Inches(6.66), Inches(6.8))
    line.line.color.rgb = LIGHT_GRAY
    line.line.width = Pt(1)

    # 5. 右侧栏 (成功示例)
    # 标题
    txBox = slide.shapes.add_textbox(Inches(7.2), Inches(2.2), Inches(5.5), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "极简要点（成功示例）"
    p.font.size = Pt(18)
    p.font.color.rgb = GREEN
    p.font.name = "Microsoft YaHei"
    
    # 绿色下划线
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(7.2), Inches(2.7), Inches(12.8), Inches(2.7))
    line.line.color.rgb = GREEN
    line.line.width = Pt(1.5)

    # 绘制要点条目的内部函数
    def draw_bullet(slide, left, top, icon_text, label, desc):
        # 图标
        txBox = slide.shapes.add_textbox(left, top, Inches(0.6), Inches(0.6))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = icon_text
        p.font.size = Pt(28)
        p.font.name = "Segoe UI Emoji"
        
        # 标签
        txBox = slide.shapes.add_textbox(left + Inches(0.8), top, Inches(4.5), Inches(0.35))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = BLACK
        p.font.name = "Microsoft YaHei"
        
        # 描述
        txBox = slide.shapes.add_textbox(left + Inches(0.8), top + Inches(0.35), Inches(4.5), Inches(0.35))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = BLACK
        p.font.name = "Microsoft YaHei"

    # 添加三个要点
    draw_bullet(slide, Inches(7.2), Inches(3.1), "👂", "专注聆听", "观众阅读文字时，无法同时听取演讲。")
    draw_bullet(slide, Inches(7.2), Inches(4.3), "💎", "提炼金句", "删除冗余的修饰词，只保留核心观点。")
    draw_bullet(slide, Inches(7.2), Inches(5.5), "🖼️", "视觉替代", "用视觉元素（图标/图片）替代长篇大论。")

    # 绿色大勾
    check_x = 9.8
    check_y = 4.8
    line1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(check_x - 0.6), Inches(check_y - 0.1), Inches(check_x), Inches(check_y + 0.5))
    line1.line.color.rgb = GREEN
    line1.line.width = Pt(25)
    line2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(check_x), Inches(check_y + 0.5), Inches(check_x + 1.2), Inches(check_y - 1.0))
    line2.line.color.rgb = GREEN
    line2.line.width = Pt(25)

    # 右侧底部结论
    txBox = slide.shapes.add_textbox(Inches(7.2), Inches(6.5), Inches(6.0), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "精简内容，视觉引导，提升传递效率。"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.font.name = "Microsoft YaHei"

    # 6. 页码
    txBox = slide.shapes.add_textbox(Inches(11.8), Inches(6.8), Inches(1.0), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "4 / 11"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK
    p.font.name = "Microsoft YaHei"
    p.alignment = PP_ALIGN.RIGHT



# ── Slide 5 ──

def build_slide_5(slide):
    # 颜色常量定义
    BLUE_DARK = RGBColor(0x00, 0x55, 0xA4)
    BLUE_VERY_DARK = RGBColor(0x1A, 0x2B, 0x3C)
    ORANGE_TEXT = RGBColor(0xFF, 0x98, 0x00)
    GRAY_DARK = RGBColor(0x33, 0x33, 0x33)
    GRAY_LIGHT = RGBColor(0x55, 0x55, 0x55)
    GRAY_LINE = RGBColor(0xCC, 0xCC, 0xCC)
    GRAY_HANDLE = RGBColor(0x88, 0x88, 0x88)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    # 1. 标题
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(10), Inches(0.8))
    tf_title = title_box.text_frame
    p_title = tf_title.paragraphs[0]
    p_title.text = "秒懂原则：视觉化的力量"
    p_title.font.name = FONT_NAME
    p_title.font.size = Pt(36)
    p_title.font.bold = True
    p_title.font.color.rgb = BLUE_DARK

    # 2. 副标题
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(10), Inches(0.5))
    tf_sub = sub_box.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "让观众在3秒内捕捉核心信息"
    p_sub.font.name = FONT_NAME
    p_sub.font.size = Pt(20)
    p_sub.font.color.rgb = BLUE_DARK

    # 3. 左侧主图：放大镜及内部元素
    # 放大镜手柄连接处 (灰色)
    conn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.8), Inches(5.2), Inches(0.6), Inches(0.8))
    conn.rotation = 45
    conn.fill.solid()
    conn.fill.fore_color.rgb = GRAY_HANDLE
    conn.line.color.rgb = BLUE_VERY_DARK
    conn.line.width = Pt(3)

    # 放大镜手柄主体 (蓝色)
    handle = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(5.8), Inches(0.9), Inches(2.2))
    handle.rotation = 45
    handle.fill.solid()
    handle.fill.fore_color.rgb = BLUE_DARK
    handle.line.color.rgb = BLUE_VERY_DARK
    handle.line.width = Pt(4)

    # 放大镜外圈 (深色粗边框，白色填充以遮挡手柄顶部)
    outer_ring = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.0), Inches(2.2), Inches(4.0), Inches(4.0))
    outer_ring.fill.solid()
    outer_ring.fill.fore_color.rgb = WHITE
    outer_ring.line.color.rgb = BLUE_VERY_DARK
    outer_ring.line.width = Pt(12)

    # 放大镜内圈 (蓝色细边框)
    inner_ring = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.2), Inches(2.4), Inches(3.6), Inches(3.6))
    inner_ring.fill.background()
    inner_ring.line.color.rgb = BLUE_DARK
    inner_ring.line.width = Pt(6)

    # 内部放射状虚线/浅色线
    lines_coords = [
        (4.0, 2.9, 4.0, 2.6), # 上
        (4.0, 5.5, 4.0, 5.8), # 下
        (2.7, 4.2, 2.4, 4.2), # 左
        (5.3, 4.2, 5.6, 4.2), # 右
        (3.1, 3.3, 2.9, 3.1), # 左上
        (4.9, 3.3, 5.1, 3.1), # 右上
        (3.1, 5.1, 2.9, 5.3), # 左下
        (4.9, 5.1, 5.1, 5.3), # 右下
    ]
    for x1, y1, x2, y2 in lines_coords:
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        line.line.color.rgb = GRAY_LINE
        line.line.width = Pt(2)

    # 内部文字 "秒懂"
    text_box_md = slide.shapes.add_textbox(Inches(2.0), Inches(3.4), Inches(4.0), Inches(1.5))
    tf_md = text_box_md.text_frame
    p_md = tf_md.paragraphs[0]
    p_md.text = "秒懂"
    p_md.alignment = PP_ALIGN.CENTER
    p_md.font.name = FONT_NAME
    p_md.font.size = Pt(65)
    p_md.font.bold = True
    p_md.font.color.rgb = ORANGE_TEXT

    # 内部小图标：眼睛 (上方)
    eye_outer = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.3), Inches(3.1), Inches(0.4), Inches(0.25))
    eye_outer.fill.background()
    eye_outer.line.color.rgb = BLUE_DARK
    eye_outer.line.width = Pt(2)
    eye_inner = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.45), Inches(3.17), Inches(0.1), Inches(0.1))
    eye_inner.fill.solid()
    eye_inner.fill.fore_color.rgb = BLUE_DARK
    eye_inner.line.fill.background()

    # 内部小图标：大脑/云朵 (右下方)
    brain = slide.shapes.add_shape(MSO_SHAPE.CLOUD, Inches(4.7), Inches(4.5), Inches(0.4), Inches(0.3))
    brain.fill.background()
    brain.line.color.rgb = BLUE_DARK
    brain.line.width = Pt(1.5)

    # 放大镜外部装饰弧线 (使用简单的曲线近似)
    arc1 = slide.shapes.add_shape(MSO_SHAPE.ARC, Inches(1.6), Inches(1.8), Inches(4.8), Inches(4.8))
    arc1.fill.background()
    arc1.line.color.rgb = BLUE_VERY_DARK
    arc1.line.width = Pt(3)
    arc1.adjustments[0] = 110
    arc1.adjustments[1] = 160

    arc2 = slide.shapes.add_shape(MSO_SHAPE.ARC, Inches(1.8), Inches(2.0), Inches(4.4), Inches(4.4))
    arc2.fill.background()
    arc2.line.color.rgb = BLUE_VERY_DARK
    arc2.line.width = Pt(3)
    arc2.adjustments[0] = 20
    arc2.adjustments[1] = 70

    # 4. 右侧要点列表
    # --- 要点 1：降低认知负荷 ---
    # 图标 1 (人脑齿轮)
    head = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.4), Inches(2.5), Inches(0.5), Inches(0.6))
    head.fill.background()
    head.line.color.rgb = BLUE_DARK
    head.line.width = Pt(2)
    neck = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.55), Inches(3.05), Inches(0.2), Inches(0.15))
    neck.fill.background()
    neck.line.color.rgb = BLUE_DARK
    neck.line.width = Pt(2)
    gear1 = slide.shapes.add_shape(MSO_SHAPE.GEAR_6, Inches(7.45), Inches(2.6), Inches(0.25), Inches(0.25))
    gear1.fill.background()
    gear1.line.color.rgb = BLUE_DARK
    gear1.line.width = Pt(1.5)
    gear2 = slide.shapes.add_shape(MSO_SHAPE.GEAR_6, Inches(7.6), Inches(2.8), Inches(0.2), Inches(0.2))
    gear2.fill.background()
    gear2.line.color.rgb = BLUE_DARK
    gear2.line.width = Pt(1.5)

    # 文本 1
    tb1 = slide.shapes.add_textbox(Inches(8.4), Inches(2.4), Inches(4.5), Inches(1.0))
    tf1 = tb1.text_frame
    p1_1 = tf1.paragraphs[0]
    p1_1.text = "降低认知负荷："
    p1_1.font.name = FONT_NAME
    p1_1.font.size = Pt(18)
    p1_1.font.bold = True
    p1_1.font.color.rgb = GRAY_DARK
    p1_2 = tf1.add_paragraph()
    p1_2.text = "一眼就能看懂逻辑"
    p1_2.font.name = FONT_NAME
    p1_2.font.size = Pt(14)
    p1_2.font.color.rgb = GRAY_LIGHT
    p1_2.space_before = Pt(6)

    # --- 要点 2：视觉层级明确 ---
    # 图标 2 (层级金字塔与箭头)
    base_y2 = 4.8
    rect1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.3), Inches(base_y2), Inches(0.8), Inches(0.15))
    rect1.fill.background()
    rect1.line.color.rgb = BLUE_DARK
    rect1.line.width = Pt(2)
    rect2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.4), Inches(base_y2 - 0.25), Inches(0.6), Inches(0.15))
    rect2.fill.background()
    rect2.line.color.rgb = BLUE_DARK
    rect2.line.width = Pt(2)
    rect3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.5), Inches(base_y2 - 0.5), Inches(0.4), Inches(0.15))
    rect3.fill.background()
    rect3.line.color.rgb = BLUE_DARK
    rect3.line.width = Pt(2)
    arrow = slide.shapes.add_shape(MSO_SHAPE.UP_ARROW, Inches(8.2), Inches(base_y2 - 0.5), Inches(0.15), Inches(0.65))
    arrow.fill.background()
    arrow.line.color.rgb = BLUE_DARK
    arrow.line.width = Pt(1.5)

    # 文本 2
    tb2 = slide.shapes.add_textbox(Inches(8.4), Inches(4.1), Inches(4.5), Inches(1.0))
    tf2 = tb2.text_frame
    p2_1 = tf2.paragraphs[0]
    p2_1.text = "视觉层级明确："
    p2_1.font.name = FONT_NAME
    p2_1.font.size = Pt(18)
    p2_1.font.bold = True
    p2_1.font.color.rgb = GRAY_DARK
    p2_2 = tf2.add_paragraph()
    p2_2.text = "通过大小、颜色区分主次"
    p2_2.font.name = FONT_NAME
    p2_2.font.size = Pt(14)
    p2_2.font.color.rgb = GRAY_LIGHT
    p2_2.space_before = Pt(6)

    # --- 要点 3：善用图表 ---
    # 图标 3 (柱状图与饼图)
    base_y3 = 6.5
    axis_x = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(7.3), Inches(base_y3), Inches(8.0), Inches(base_y3))
    axis_x.line.color.rgb = BLUE_DARK
    axis_x.line.width = Pt(2)
    axis_y = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(7.3), Inches(base_y3), Inches(7.3), Inches(base_y3 - 0.8))
    axis_y.line.color.rgb = BLUE_DARK
    axis_y.line.width = Pt(2)
    
    bar1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.4), Inches(base_y3 - 0.3), Inches(0.12), Inches(0.3))
    bar1.fill.background()
    bar1.line.color.rgb = BLUE_DARK
    bar1.line.width = Pt(1.5)
    bar2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.6), Inches(base_y3 - 0.5), Inches(0.12), Inches(0.5))
    bar2.fill.background()
    bar2.line.color.rgb = BLUE_DARK
    bar2.line.width = Pt(1.5)
    bar3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.8), Inches(base_y3 - 0.2), Inches(0.12), Inches(0.2))
    bar3.fill.background()
    bar3.line.color.rgb = BLUE_DARK
    bar3.line.width = Pt(1.5)
    
    pie = slide.shapes.add_shape(MSO_SHAPE.PIE, Inches(7.7), Inches(base_y3 - 0.9), Inches(0.4), Inches(0.4))
    pie.fill.background()
    pie.line.color.rgb = BLUE_DARK
    pie.line.width = Pt(1.5)

    # 文本 3
    tb3 = slide.shapes.add_textbox(Inches(8.4), Inches(5.8), Inches(4.5), Inches(1.0))
    tf3 = tb3.text_frame
    p3_1 = tf3.paragraphs[0]
    p3_1.text = "善用图表："
    p3_1.font.name = FONT_NAME
    p3_1.font.size = Pt(18)
    p3_1.font.bold = True
    p3_1.font.color.rgb = GRAY_DARK
    p3_2 = tf3.add_paragraph()
    p3_2.text = "数据关系一目了然"
    p3_2.font.name = FONT_NAME
    p3_2.font.size = Pt(14)
    p3_2.font.color.rgb = GRAY_LIGHT
    p3_2.space_before = Pt(6)

    # 5. 页码
    page_num = slide.shapes.add_textbox(Inches(12.0), Inches(6.8), Inches(1.0), Inches(0.5))
    tf_page = page_num.text_frame
    p_page = tf_page.paragraphs[0]
    p_page.text = "5 / 11"
    p_page.font.name = FONT_NAME
    p_page.font.size = Pt(12)
    p_page.font.color.rgb = GRAY_HANDLE
    p_page.alignment = PP_ALIGN.RIGHT



# ── Slide 6 ──

def build_slide_6(slide):
    # Colors
    BLUE_PRIMARY = RGBColor(0x00, 0x52, 0xCC)
    TEXT_BLACK = RGBColor(0x33, 0x33, 0x33)
    TEXT_GRAY = RGBColor(0x7F, 0x7F, 0x7F)
    LIGHT_GRAY = RGBColor(0xD9, 0xD9, 0xD9)
    CUBE_FILL = RGBColor(0xF4, 0xF6, 0xF9)
    CUBE_LINE = RGBColor(0x2F, 0x45, 0x6A)
    HIGHLIGHT_FILL = RGBColor(0xDE, 0xEA, 0xF6)
    HIGHLIGHT_LINE = RGBColor(0x5B, 0x9B, 0xD5)

    # 1. Top Left Page Indicator
    # Small grey dash
    dash = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0.4), Inches(0.15), Inches(0.04))
    dash.fill.solid()
    dash.fill.fore_color.rgb = LIGHT_GRAY
    dash.line.fill.background()

    # "Page 6"
    tx_page = slide.shapes.add_textbox(Inches(0.35), Inches(0.45), Inches(1), Inches(0.3))
    tf_page = tx_page.text_frame
    tf_page.word_wrap = False
    p_page = tf_page.paragraphs[0]
    p_page.text = "Page 6"
    p_page.font.name = "Microsoft YaHei"
    p_page.font.size = Pt(12)
    p_page.font.bold = True
    p_page.font.color.rgb = TEXT_BLACK

    # "6/11"
    tx_num = slide.shapes.add_textbox(Inches(0.35), Inches(0.7), Inches(1), Inches(0.3))
    tf_num = tx_num.text_frame
    p_num = tf_num.paragraphs[0]
    p_num.text = "6/11"
    p_num.font.name = "Microsoft YaHei"
    p_num.font.size = Pt(10)
    p_num.font.color.rgb = TEXT_GRAY

    # 2. Main Title
    tx_title = slide.shapes.add_textbox(Inches(2), Inches(0.8), Inches(9.333), Inches(0.8))
    tf_title = tx_title.text_frame
    p_title = tf_title.paragraphs[0]
    p_title.alignment = PP_ALIGN.CENTER
    
    run1 = p_title.add_run()
    run1.text = "法则二："
    run1.font.name = "Microsoft YaHei"
    run1.font.size = Pt(36)
    run1.font.bold = True
    run1.font.color.rgb = TEXT_BLACK

    run2 = p_title.add_run()
    run2.text = "设计统一建立专业信任"
    run2.font.name = "Microsoft YaHei"
    run2.font.size = Pt(36)
    run2.font.bold = True
    run2.font.color.rgb = BLUE_PRIMARY

    # 3. Subtitle
    tx_sub = slide.shapes.add_textbox(Inches(2), Inches(1.7), Inches(9.333), Inches(0.5))
    tf_sub = tx_sub.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.alignment = PP_ALIGN.CENTER
    p_sub.text = "从视觉一致性中体现专业度"
    p_sub.font.name = "Microsoft YaHei"
    p_sub.font.size = Pt(20)
    p_sub.font.color.rgb = TEXT_BLACK

    # 4. Central Graphics (Cubes)
    cube_y = Inches(2.8)
    cube_w = Inches(2.2)
    cube_h = Inches(2.2)

    # Cube 1 (Left)
    cube1 = slide.shapes.add_shape(MSO_SHAPE.CUBE, Inches(2.8), cube_y, cube_w, cube_h)
    cube1.fill.solid()
    cube1.fill.fore_color.rgb = CUBE_FILL
    cube1.line.color.rgb = CUBE_LINE
    cube1.line.width = Pt(2)

    # Cube 2 (Middle)
    cube2 = slide.shapes.add_shape(MSO_SHAPE.CUBE, Inches(5.5), cube_y, cube_w, cube_h)
    cube2.fill.solid()
    cube2.fill.fore_color.rgb = CUBE_FILL
    cube2.line.color.rgb = CUBE_LINE
    cube2.line.width = Pt(2)

    # Flying piece for Cube 2
    small_cube = slide.shapes.add_shape(MSO_SHAPE.CUBE, Inches(6.8), Inches(2.5), Inches(0.8), Inches(0.8))
    small_cube.fill.solid()
    small_cube.fill.fore_color.rgb = HIGHLIGHT_FILL
    small_cube.line.color.rgb = HIGHLIGHT_LINE
    small_cube.line.width = Pt(1.5)

    # Arrow for flying piece
    arrow_insert = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(7.1), Inches(3.4), Inches(0.2), Inches(0.3))
    arrow_insert.fill.solid()
    arrow_insert.fill.fore_color.rgb = HIGHLIGHT_LINE
    arrow_insert.line.fill.background()

    # Cube 3 (Right)
    cube3 = slide.shapes.add_shape(MSO_SHAPE.CUBE, Inches(8.2), cube_y, cube_w, cube_h)
    cube3.fill.solid()
    cube3.fill.fore_color.rgb = CUBE_FILL
    cube3.line.color.rgb = CUBE_LINE
    cube3.line.width = Pt(2)

    # Highlighted piece on Cube 3 (Simulated by a smaller cube on top right corner)
    hl_cube = slide.shapes.add_shape(MSO_SHAPE.CUBE, Inches(9.2), Inches(3.2), Inches(0.7), Inches(0.7))
    hl_cube.fill.solid()
    hl_cube.fill.fore_color.rgb = HIGHLIGHT_FILL
    hl_cube.line.color.rgb = HIGHLIGHT_LINE
    hl_cube.line.width = Pt(1.5)

    # 5. Bottom Arrow and Text
    # Long arrow line
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(3.8), Inches(5.4), Inches(9.4), Inches(5.4))
    connector.line.color.rgb = HIGHLIGHT_LINE
    connector.line.width = Pt(1.5)
    # Add arrow head (using standard line properties if possible, or draw a small triangle)
    arrow_head = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(9.3), Inches(5.33), Inches(0.15), Inches(0.15))
    arrow_head.rotation = 90
    arrow_head.fill.solid()
    arrow_head.fill.fore_color.rgb = HIGHLIGHT_LINE
    arrow_head.line.fill.background()

    # Text below arrow
    tx_arrow = slide.shapes.add_textbox(Inches(2), Inches(5.6), Inches(9.333), Inches(0.4))
    tf_arrow = tx_arrow.text_frame
    p_arrow = tf_arrow.paragraphs[0]
    p_arrow.alignment = PP_ALIGN.CENTER
    p_arrow.text = "统一感能降低观众的视觉疲劳"
    p_arrow.font.name = "Microsoft YaHei"
    p_arrow.font.size = Pt(14)
    p_arrow.font.bold = True
    p_arrow.font.color.rgb = TEXT_BLACK

    # 6. Bottom 3 Columns
    col_y = Inches(6.3)
    
    # --- Column 1 ---
    # Icon 1: Compass/Ruler
    compass = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(1.0), col_y + Inches(0.1), Inches(0.3), Inches(0.4))
    compass.rotation = -90
    compass.fill.background()
    compass.line.color.rgb = CUBE_LINE
    compass.line.width = Pt(1.5)
    
    ruler = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.4), col_y, Inches(0.15), Inches(0.5))
    ruler.fill.background()
    ruler.line.color.rgb = HIGHLIGHT_LINE
    ruler.line.width = Pt(1.5)

    # Text 1
    tx_col1_title = slide.shapes.add_textbox(Inches(1.7), col_y - Inches(0.1), Inches(3.0), Inches(0.3))
    p_col1_title = tx_col1_title.text_frame.paragraphs[0]
    r1 = p_col1_title.add_run()
    r1.text = "1. "
    r1.font.bold = True
    r1.font.size = Pt(13)
    r2 = p_col1_title.add_run()
    r2.text = "风格漂移"
    r2.font.bold = True
    r2.font.size = Pt(13)
    r2.font.color.rgb = BLUE_PRIMARY
    r3 = p_col1_title.add_run()
    r3.text = "是PPT的大忌"
    r3.font.bold = True
    r3.font.size = Pt(13)

    tx_col1_desc = slide.shapes.add_textbox(Inches(1.7), col_y + Inches(0.2), Inches(3.0), Inches(0.4))
    p_col1_desc = tx_col1_desc.text_frame.paragraphs[0]
    p_col1_desc.text = "避免混乱，保持整体风格的一致性。"
    p_col1_desc.font.size = Pt(11)
    p_col1_desc.font.color.rgb = TEXT_BLACK

    # --- Column 2 ---
    # Icon 2: Eye
    eye_outer = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.0), col_y + Inches(0.1), Inches(0.5), Inches(0.3))
    eye_outer.fill.background()
    eye_outer.line.color.rgb = CUBE_LINE
    eye_outer.line.width = Pt(1.5)
    
    eye_inner = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.15), col_y + Inches(0.15), Inches(0.2), Inches(0.2))
    eye_inner.fill.background()
    eye_inner.line.color.rgb = HIGHLIGHT_LINE
    eye_inner.line.width = Pt(1.5)
    
    pulse = slide.shapes.add_shape(MSO_SHAPE.ZIG_ZAG, Inches(5.0), col_y + Inches(0.45), Inches(0.5), Inches(0.1))
    pulse.fill.background()
    pulse.line.color.rgb = HIGHLIGHT_LINE
    pulse.line.width = Pt(1.5)

    # Text 2
    tx_col2_title = slide.shapes.add_textbox(Inches(5.7), col_y - Inches(0.1), Inches(3.2), Inches(0.3))
    p_col2_title = tx_col2_title.text_frame.paragraphs[0]
    r1 = p_col2_title.add_run()
    r1.text = "2. 统一感能降低观众的"
    r1.font.bold = True
    r1.font.size = Pt(13)
    r2 = p_col2_title.add_run()
    r2.text = "视觉疲劳"
    r2.font.bold = True
    r2.font.size = Pt(13)
    r2.font.color.rgb = BLUE_PRIMARY

    tx_col2_desc = slide.shapes.add_textbox(Inches(5.7), col_y + Inches(0.2), Inches(3.2), Inches(0.4))
    p_col2_desc = tx_col2_desc.text_frame.paragraphs[0]
    p_col2_desc.text = "视觉流畅，让观众更专注于内容。"
    p_col2_desc.font.size = Pt(11)
    p_col2_desc.font.color.rgb = TEXT_BLACK

    # --- Column 3 ---
    # Icon 3: Browser/Window
    browser = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.0), col_y, Inches(0.5), Inches(0.4))
    browser.fill.background()
    browser.line.color.rgb = CUBE_LINE
    browser.line.width = Pt(1.5)
    
    browser_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.0), col_y + Inches(0.1), Inches(0.5), Inches(0.02))
    browser_line.fill.solid()
    browser_line.fill.fore_color.rgb = CUBE_LINE
    browser_line.line.fill.background()
    
    mag_glass = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.3), col_y + Inches(0.2), Inches(0.2), Inches(0.2))
    mag_glass.fill.background()
    mag_glass.line.color.rgb = HIGHLIGHT_LINE
    mag_glass.line.width = Pt(1.5)

    # Text 3
    tx_col3_title = slide.shapes.add_textbox(Inches(9.7), col_y - Inches(0.1), Inches(3.5), Inches(0.3))
    p_col3_title = tx_col3_title.text_frame.paragraphs[0]
    r1 = p_col3_title.add_run()
    r1.text = "3. "
    r1.font.bold = True
    r1.font.size = Pt(13)
    r2 = p_col3_title.add_run()
    r2.text = "专业感"
    r2.font.bold = True
    r2.font.size = Pt(13)
    r2.font.color.rgb = BLUE_PRIMARY
    r3 = p_col3_title.add_run()
    r3.text = "源于对细节的严苛把控"
    r3.font.bold = True
    r3.font.size = Pt(13)

    tx_col3_desc = slide.shapes.add_textbox(Inches(9.7), col_y + Inches(0.2), Inches(3.5), Inches(0.4))
    p_col3_desc = tx_col3_desc.text_frame.paragraphs[0]
    p_col3_desc.text = "对齐、间距、字体、颜色的精准规范。"
    p_col3_desc.font.size = Pt(11)
    p_col3_desc.font.color.rgb = TEXT_BLACK



# ── Slide 7 ──

def build_slide_7(slide):
    # Define Colors
    BLUE_PRIMARY = RGBColor(0x00, 0x52, 0xD9)
    BLUE_DARK = RGBColor(0x00, 0x2B, 0x75)
    GRAY_DARK = RGBColor(0x2B, 0x2F, 0x36)
    ORANGE_ACCENT = RGBColor(0xFF, 0x95, 0x00)
    TEXT_MAIN = RGBColor(0x33, 0x33, 0x33)
    TEXT_SUB = RGBColor(0x66, 0x66, 0x66)
    LINE_COLOR = RGBColor(0xE0, 0xE0, 0xE0)
    BG_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    # 1. Header Area
    # Top thin blue line
    top_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.06))
    top_line.fill.solid()
    top_line.fill.fore_color.rgb = BLUE_PRIMARY
    top_line.line.fill.background()

    # Header Left Text
    tb_header_left = slide.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(2), Inches(0.4))
    p = tb_header_left.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "简而不凡"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = BLUE_PRIMARY

    # Header Right Text
    tb_header_right = slide.shapes.add_textbox(Inches(11.5), Inches(0.15), Inches(1.5), Inches(0.4))
    p = tb_header_right.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "PAGE 7 / 11"
    run.font.size = Pt(12)
    run.font.color.rgb = TEXT_SUB

    # Header Bottom Separator
    header_sep = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.6), Inches(0.6), Inches(12.733), Inches(0.6))
    header_sep.line.color.rgb = LINE_COLOR

    # 2. Main Title & Subtitle
    tb_title = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(8), Inches(0.8))
    p = tb_title.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "视觉规范：配色与字体的秩序"
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = BLUE_PRIMARY

    tb_subtitle = slide.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(8), Inches(0.5))
    p = tb_subtitle.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "建立一套专属的视觉系统"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = TEXT_MAIN

    # 3. Left Bullet Points
    # Item 1
    add_icon_box(slide, Inches(0.6), Inches(3.0), "🎨", size=Inches(0.8))
    tb_item1 = slide.shapes.add_textbox(Inches(1.6), Inches(2.9), Inches(6.0), Inches(1.0))
    tf1 = tb_item1.text_frame
    tf1.word_wrap = True
    p1_1 = tf1.paragraphs[0]
    r1_1 = p1_1.add_run()
    r1_1.text = "1. 全文配色"
    r1_1.font.size = Pt(16)
    r1_1.font.bold = True
    r1_1.font.color.rgb = TEXT_MAIN
    r1_2 = p1_1.add_run()
    r1_2.text = "不超过3种，主次分明"
    r1_2.font.size = Pt(16)
    r1_2.font.bold = True
    r1_2.font.color.rgb = ORANGE_ACCENT
    
    p1_2 = tf1.add_paragraph()
    p1_2.space_before = Pt(6)
    r1_3 = p1_2.add_run()
    r1_3.text = "限制色彩数量，确保页面干净统一，提升专业度。"
    r1_3.font.size = Pt(14)
    r1_3.font.color.rgb = TEXT_SUB

    # Item 2
    add_icon_box(slide, Inches(0.6), Inches(4.4), "Aa", size=Inches(0.8))
    tb_item2 = slide.shapes.add_textbox(Inches(1.6), Inches(4.3), Inches(6.0), Inches(1.0))
    tf2 = tb_item2.text_frame
    tf2.word_wrap = True
    p2_1 = tf2.paragraphs[0]
    r2_1 = p2_1.add_run()
    r2_1.text = "2. 字体选择需统一，建议"
    r2_1.font.size = Pt(16)
    r2_1.font.bold = True
    r2_1.font.color.rgb = TEXT_MAIN
    r2_2 = p2_1.add_run()
    r2_2.text = "不超过2种"
    r2_2.font.size = Pt(16)
    r2_2.font.bold = True
    r2_2.font.color.rgb = ORANGE_ACCENT
    
    p2_2 = tf2.add_paragraph()
    p2_2.space_before = Pt(6)
    r2_3 = p2_2.add_run()
    r2_3.text = "选择易读的无衬线字体（如苹方-简），保持风格一致。"
    r2_3.font.size = Pt(14)
    r2_3.font.color.rgb = TEXT_SUB

    # Item 3
    add_icon_box(slide, Inches(0.6), Inches(5.8), "🖍️", size=Inches(0.8))
    tb_item3 = slide.shapes.add_textbox(Inches(1.6), Inches(5.7), Inches(6.0), Inches(1.0))
    tf3 = tb_item3.text_frame
    tf3.word_wrap = True
    p3_1 = tf3.paragraphs[0]
    r3_1 = p3_1.add_run()
    r3_1.text = "3. 关键信息"
    r3_1.font.size = Pt(16)
    r3_1.font.bold = True
    r3_1.font.color.rgb = TEXT_MAIN
    r3_2 = p3_1.add_run()
    r3_2.text = "加粗或变色"
    r3_2.font.size = Pt(16)
    r3_2.font.bold = True
    r3_2.font.color.rgb = ORANGE_ACCENT
    r3_3 = p3_1.add_run()
    r3_3.text = "，而非随意更改字体"
    r3_3.font.size = Pt(16)
    r3_3.font.bold = True
    r3_3.font.color.rgb = TEXT_MAIN
    
    p3_2 = tf3.add_paragraph()
    p3_2.space_before = Pt(6)
    r3_4 = p3_2.add_run()
    r3_4.text = "通过字重和色彩"
    r3_4.font.size = Pt(14)
    r3_4.font.color.rgb = TEXT_SUB
    r3_5 = p3_2.add_run()
    r3_5.text = "强调重点"
    r3_5.font.size = Pt(14)
    r3_5.font.color.rgb = ORANGE_ACCENT
    r3_6 = p3_2.add_run()
    r3_6.text = "，避免视觉混乱。"
    r3_6.font.size = Pt(14)
    r3_6.font.color.rgb = TEXT_SUB

    # 4. Right Panel (White Card)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.0), Inches(1.2), Inches(4.6), Inches(5.8))
    card.fill.solid()
    card.fill.fore_color.rgb = BG_WHITE
    card.line.color.rgb = LINE_COLOR
    card.line.width = Pt(1)

    # Card Title 1
    tb_card_t1 = slide.shapes.add_textbox(Inches(8.2), Inches(1.4), Inches(4.0), Inches(0.4))
    p = tb_card_t1.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "配色建议"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = TEXT_MAIN

    # Color Swatches
    swatch_data = [
        (Inches(8.25), BLUE_PRIMARY, "主色", "#0052D9"),
        (Inches(9.65), GRAY_DARK, "辅助色", "#2B2F36"),
        (Inches(11.05), ORANGE_ACCENT, "强调色", "#FF9500")
    ]
    
    for left, color, label, hex_code in swatch_data:
        swatch = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.0), Inches(1.2), Inches(0.7))
        swatch.fill.solid()
        swatch.fill.fore_color.rgb = color
        swatch.line.fill.background()
        
        tb_swatch = slide.shapes.add_textbox(left, Inches(2.8), Inches(1.2), Inches(0.5))
        tf = tb_swatch.text_frame
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = label + "\n"
        r1.font.size = Pt(12)
        r1.font.color.rgb = TEXT_MAIN
        r2 = p1.add_run()
        r2.text = hex_code
        r2.font.size = Pt(10)
        r2.font.color.rgb = TEXT_SUB

    # Gradient Bar (Represented as solid primary blue)
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.25), Inches(3.6), Inches(4.0), Inches(0.25))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE_PRIMARY
    bar.line.fill.background()

    tb_bar_left = slide.shapes.add_textbox(Inches(8.15), Inches(3.9), Inches(1.5), Inches(0.3))
    p = tb_bar_left.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "#002B75"
    run.font.size = Pt(10)
    run.font.color.rgb = TEXT_SUB

    tb_bar_right = slide.shapes.add_textbox(Inches(11.15), Inches(3.9), Inches(1.2), Inches(0.3))
    p = tb_bar_right.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "#0052D9"
    run.font.size = Pt(10)
    run.font.color.rgb = TEXT_SUB

    # Card Separator Line
    card_sep = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(8.25), Inches(4.3), Inches(12.25), Inches(4.3))
    card_sep.line.color.rgb = LINE_COLOR

    # Card Title 2
    tb_card_t2 = slide.shapes.add_textbox(Inches(8.2), Inches(4.5), Inches(4.0), Inches(0.4))
    p = tb_card_t2.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "字体样式组合"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = TEXT_MAIN

    # Typography Examples
    tb_ex1 = slide.shapes.add_textbox(Inches(8.2), Inches(5.0), Inches(4.2), Inches(0.5))
    p = tb_ex1.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "一级标题示例 32-40pt"
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = BLUE_PRIMARY

    tb_ex2 = slide.shapes.add_textbox(Inches(8.2), Inches(5.6), Inches(4.2), Inches(0.4))
    p = tb_ex2.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "副标题示例 20-24pt"
    run.font.size = Pt(18)
    run.font.color.rgb = TEXT_MAIN

    tb_ex3 = slide.shapes.add_textbox(Inches(8.2), Inches(6.1), Inches(4.2), Inches(0.4))
    p = tb_ex3.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "正文内容示例，高易读性 16-18pt"
    run.font.size = Pt(14)
    run.font.color.rgb = TEXT_SUB



# ── Slide 8 ──

def build_slide_8(slide):
    # Colors
    BLUE_TITLE = RGBColor(0x1A, 0x66, 0xCC)
    DARK_TEXT = RGBColor(0x22, 0x22, 0x22)
    GRAY_TEXT = RGBColor(0x66, 0x66, 0x66)
    BLUE_LINE = RGBColor(0x5B, 0x9B, 0xD5)
    LIGHT_BLUE_FILL = RGBColor(0xE6, 0xF0, 0xFA)
    GREEN_OK = RGBColor(0x4C, 0xAF, 0x50)
    RED_ERR = RGBColor(0xF4, 0x43, 0x36)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    
    # 1. Title and Subtitle
    tb_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(8.0), Inches(0.8))
    p_title = tb_title.text_frame.paragraphs[0]
    
    run1 = p_title.add_run()
    run1.text = "排版逻辑："
    run1.font.size = Pt(32)
    run1.font.bold = True
    run1.font.color.rgb = BLUE_TITLE
    run1.font.name = "Microsoft YaHei"
    
    run2 = p_title.add_run()
    run2.text = "始终如一的风格表达"
    run2.font.size = Pt(32)
    run2.font.bold = True
    run2.font.color.rgb = DARK_TEXT
    run2.font.name = "Microsoft YaHei"
    
    tb_sub = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.0), Inches(0.5))
    p_sub = tb_sub.text_frame.paragraphs[0]
    p_sub.text = "规范化的布局让阅读更顺畅"
    p_sub.font.size = Pt(18)
    p_sub.font.color.rgb = GRAY_TEXT
    p_sub.font.name = "Microsoft YaHei"

    # 2. Left Section: Alignment & Margins
    # Icon (Grid)
    icon_grid = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.2), Inches(0.6), Inches(0.45))
    icon_grid.fill.solid()
    icon_grid.fill.fore_color.rgb = LIGHT_BLUE_FILL
    icon_grid.line.color.rgb = BLUE_LINE
    icon_grid.line.width = Pt(1.5)
    
    line_v = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.95), Inches(2.1), Inches(0.95), Inches(2.7))
    line_v.line.color.rgb = BLUE_LINE
    line_h = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.7), Inches(2.4), Inches(1.5), Inches(2.4))
    line_h.line.color.rgb = BLUE_LINE

    # Heading
    tb_h1 = slide.shapes.add_textbox(Inches(1.6), Inches(2.15), Inches(4.0), Inches(0.5))
    p_h1 = tb_h1.text_frame.paragraphs[0]
    p_h1.text = "建立统一的页边距与对齐基准"
    p_h1.font.size = Pt(18)
    p_h1.font.bold = True
    p_h1.font.name = "Microsoft YaHei"

    # Wireframe Graphic
    # Outer Box
    box_outer = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.7), Inches(3.1), Inches(3.3), Inches(2.2))
    box_outer.fill.solid()
    box_outer.fill.fore_color.rgb = LIGHT_BLUE_FILL
    box_outer.line.color.rgb = BLUE_LINE
    
    # Inner Dashed Box
    box_inner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.9), Inches(3.3), Inches(2.9), Inches(1.8))
    box_inner.fill.background()
    box_inner.line.color.rgb = BLUE_LINE
    box_inner.line.dash_style = 3 # Dashed
    
    # Content Blocks inside Wireframe
    rect_title = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.95), Inches(3.5), Inches(0.9), Inches(0.2))
    rect_title.fill.solid()
    rect_title.fill.fore_color.rgb = BLUE_TITLE
    rect_title.line.fill.background()
    
    for i in range(3):
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.95), Inches(3.8 + i*0.15), Inches(1.3), Inches(0.05))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(0xBD, 0xD0, 0xE6)
        line.line.fill.background()
        
    rect_img = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.4), Inches(3.5), Inches(1.35), Inches(0.9))
    rect_img.fill.background()
    rect_img.line.color.rgb = BLUE_LINE
    
    # Mountain placeholder inside rect_img
    tri1 = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(3.45), Inches(3.9), Inches(0.6), Inches(0.5))
    tri1.fill.solid()
    tri1.fill.fore_color.rgb = RGBColor(0xBD, 0xD0, 0xE6)
    tri1.line.fill.background()
    
    tri2 = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(3.85), Inches(4.0), Inches(0.5), Inches(0.4))
    tri2.fill.solid()
    tri2.fill.fore_color.rgb = RGBColor(0xBD, 0xD0, 0xE6)
    tri2.line.fill.background()
    
    sun = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4.15), Inches(3.65), Inches(0.15), Inches(0.15))
    sun.fill.solid()
    sun.fill.fore_color.rgb = RGBColor(0xBD, 0xD0, 0xE6)
    sun.line.fill.background()

    # Alignment Guides
    guide_v = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(3.35), Inches(2.9), Inches(3.35), Inches(5.5))
    guide_v.line.color.rgb = BLUE_TITLE
    guide_h = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.5), Inches(4.05), Inches(5.2), Inches(4.05))
    guide_h.line.color.rgb = BLUE_TITLE

    # Labels for Alignment
    lbl_align_l = slide.shapes.add_textbox(Inches(0.7), Inches(3.9), Inches(1.0), Inches(0.3))
    lbl_align_l.text_frame.text = "对齐基准"
    lbl_align_l.text_frame.paragraphs[0].font.size = Pt(12)
    lbl_align_l.text_frame.paragraphs[0].font.color.rgb = BLUE_TITLE

    lbl_align_b = slide.shapes.add_textbox(Inches(3.0), Inches(5.6), Inches(1.0), Inches(0.3))
    lbl_align_b.text_frame.text = "对齐基准"
    lbl_align_b.text_frame.paragraphs[0].font.size = Pt(12)
    lbl_align_b.text_frame.paragraphs[0].font.color.rgb = BLUE_TITLE

    lbl_margin = slide.shapes.add_textbox(Inches(5.2), Inches(3.9), Inches(1.0), Inches(0.3))
    lbl_margin.text_frame.text = "页边距"
    lbl_margin.text_frame.paragraphs[0].font.size = Pt(12)
    lbl_margin.text_frame.paragraphs[0].font.color.rgb = BLUE_TITLE

    # 3. Top Right Section: Icon Consistency
    # Icon (Four squares/circles)
    for r in range(2):
        for c in range(2):
            sq = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.6 + c*0.25), Inches(2.2 + r*0.25), Inches(0.2), Inches(0.2))
            sq.fill.solid()
            sq.fill.fore_color.rgb = LIGHT_BLUE_FILL
            sq.line.color.rgb = BLUE_LINE

    # Heading
    tb_h2 = slide.shapes.add_textbox(Inches(7.2), Inches(2.15), Inches(5.5), Inches(0.5))
    p_h2 = tb_h2.text_frame.paragraphs[0]
    p_h2.text = "保持图标风格一致（全线框或全色块）"
    p_h2.font.size = Pt(18)
    p_h2.font.bold = True
    p_h2.font.name = "Microsoft YaHei"

    # Correct Icons (Outline)
    icon_correct = slide.shapes.add_textbox(Inches(7.2), Inches(2.9), Inches(2.5), Inches(0.8))
    p_ic = icon_correct.text_frame.paragraphs[0]
    p_ic.text = "⚙  💡  📄"
    p_ic.font.size = Pt(36)
    p_ic.font.color.rgb = BLUE_TITLE
    
    # Correct Label
    chk_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.6), Inches(4.0), Inches(0.2), Inches(0.2))
    chk_circle.fill.solid()
    chk_circle.fill.fore_color.rgb = GREEN_OK
    chk_circle.line.fill.background()
    
    lbl_correct = slide.shapes.add_textbox(Inches(7.8), Inches(3.9), Inches(1.5), Inches(0.3))
    lbl_correct.text_frame.text = "正确（一致）"
    lbl_correct.text_frame.paragraphs[0].font.size = Pt(14)

    # Incorrect Icons (Mixed)
    icon_incorrect = slide.shapes.add_textbox(Inches(10.0), Inches(2.9), Inches(2.5), Inches(0.8))
    p_ii = icon_incorrect.text_frame.paragraphs[0]
    p_ii.text = "📢  ✋  ☁"
    p_ii.font.size = Pt(36)
    p_ii.font.color.rgb = BLUE_TITLE

    # Incorrect Label
    err_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.4), Inches(4.0), Inches(0.2), Inches(0.2))
    err_circle.fill.solid()
    err_circle.fill.fore_color.rgb = RED_ERR
    err_circle.line.fill.background()
    
    lbl_incorrect = slide.shapes.add_textbox(Inches(10.6), Inches(3.9), Inches(1.5), Inches(0.3))
    lbl_incorrect.text_frame.text = "错误（混杂）"
    lbl_incorrect.text_frame.paragraphs[0].font.size = Pt(14)

    # 4. Bottom Right Section: Whitespace
    # Icon (Document)
    icon_doc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.6), Inches(4.7), Inches(0.45), Inches(0.5))
    icon_doc.fill.solid()
    icon_doc.fill.fore_color.rgb = LIGHT_BLUE_FILL
    icon_doc.line.color.rgb = BLUE_LINE
    for i in range(3):
        dl = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.7), Inches(4.85 + i*0.1), Inches(0.25), Inches(0.03))
        dl.fill.solid()
        dl.fill.fore_color.rgb = BLUE_LINE
        dl.line.fill.background()

    # Heading
    tb_h3 = slide.shapes.add_textbox(Inches(7.2), Inches(4.75), Inches(5.0), Inches(0.5))
    p_h3 = tb_h3.text_frame.paragraphs[0]
    p_h3.text = "留白艺术：给内容呼吸的空间"
    p_h3.font.size = Pt(18)
    p_h3.font.bold = True
    p_h3.font.name = "Microsoft YaHei"

    # Good Layout Graphic (Whitespace)
    box_ws_outer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.3), Inches(5.5), Inches(2.6), Inches(1.8))
    box_ws_outer.fill.solid()
    box_ws_outer.fill.fore_color.rgb = RGBColor(0xF0, 0xF8, 0xFF)
    box_ws_outer.line.color.rgb = BLUE_LINE
    box_ws_outer.line.dash_style = 3
    
    # Diagonal lines for whitespace indication
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(7.3), Inches(5.5), Inches(7.7), Inches(5.9)).line.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(9.9), Inches(5.5), Inches(9.5), Inches(5.9)).line.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(7.3), Inches(7.3), Inches(7.7), Inches(6.9)).line.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(9.9), Inches(7.3), Inches(9.5), Inches(6.9)).line.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)

    box_ws_inner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.8), Inches(5.9), Inches(1.6), Inches(1.0))
    box_ws_inner.fill.solid()
    box_ws_inner.fill.fore_color.rgb = WHITE
    box_ws_inner.line.fill.background()
    
    # Shadow effect simulation
    box_ws_inner_shadow = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.78), Inches(5.88), Inches(1.64), Inches(1.04))
    box_ws_inner_shadow.fill.background()
    box_ws_inner_shadow.line.color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
    
    tb_core = slide.shapes.add_textbox(Inches(7.8), Inches(6.1), Inches(1.6), Inches(0.4))
    p_core = tb_core.text_frame.paragraphs[0]
    p_core.text = "核心内容"
    p_core.font.size = Pt(16)
    p_core.font.bold = True
    p_core.alignment = PP_ALIGN.CENTER
    
    core_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.2), Inches(6.55), Inches(0.8), Inches(0.05))
    core_line.fill.solid()
    core_line.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    core_line.line.fill.background()

    # Whitespace Labels
    lbl_ws_t = slide.shapes.add_textbox(Inches(8.3), Inches(5.55), Inches(0.8), Inches(0.2))
    lbl_ws_t.text_frame.text = "呼吸空间"
    lbl_ws_t.text_frame.paragraphs[0].font.size = Pt(9)
    lbl_ws_t.text_frame.paragraphs[0].font.color.rgb = BLUE_TITLE
    
    lbl_ws_b = slide.shapes.add_textbox(Inches(8.3), Inches(7.0), Inches(0.8), Inches(0.2))
    lbl_ws_b.text_frame.text = "呼吸空间"
    lbl_ws_b.text_frame.paragraphs[0].font.size = Pt(9)
    lbl_ws_b.text_frame.paragraphs[0].font.color.rgb = BLUE_TITLE

    lbl_ws_l = slide.shapes.add_textbox(Inches(7.35), Inches(6.3), Inches(0.8), Inches(0.2))
    lbl_ws_l.text_frame.text = "呼吸空间"
    lbl_ws_l.text_frame.paragraphs[0].font.size = Pt(9)
    lbl_ws_l.text_frame.paragraphs[0].font.color.rgb = BLUE_TITLE

    lbl_ws_r = slide.shapes.add_textbox(Inches(9.45), Inches(6.3), Inches(0.8), Inches(0.2))
    lbl_ws_r.text_frame.text = "呼吸空间"
    lbl_ws_r.text_frame.paragraphs[0].font.size = Pt(9)
    lbl_ws_r.text_frame.paragraphs[0].font.color.rgb = BLUE_TITLE

    # Bad Layout Graphic (Cluttered)
    box_cl_outer = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.4), Inches(5.5), Inches(1.5), Inches(1.0))
    box_cl_outer.fill.solid()
    box_cl_outer.fill.fore_color.rgb = LIGHT_BLUE_FILL
    box_cl_outer.line.color.rgb = BLUE_LINE

    # Cluttered inner elements
    cl_title = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10.45), Inches(5.55), Inches(0.6), Inches(0.15))
    cl_title.fill.solid()
    cl_title.fill.fore_color.rgb = BLUE_TITLE
    cl_title.line.fill.background()
    
    for i in range(4):
        cl_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10.45), Inches(5.75 + i*0.1), Inches(0.75), Inches(0.05))
        cl_line.fill.solid()
        cl_line.fill.fore_color.rgb = RGBColor(0xBD, 0xD0, 0xE6)
        cl_line.line.fill.background()
        
    cl_img = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(11.3), Inches(5.55), Inches(0.55), Inches(0.4))
    cl_img.fill.solid()
    cl_img.fill.fore_color.rgb = RGBColor(0xBD, 0xD0, 0xE6)
    cl_img.line.color.rgb = BLUE_LINE
    
    cl_box2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10.45), Inches(6.2), Inches(0.75), Inches(0.25))
    cl_box2.fill.solid()
    cl_box2.fill.fore_color.rgb = RGBColor(0xBD, 0xD0, 0xE6)
    cl_box2.line.color.rgb = BLUE_LINE

    cl_box3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(11.3), Inches(6.0), Inches(0.55), Inches(0.45))
    cl_box3.fill.solid()
    cl_box3.fill.fore_color.rgb = RGBColor(0xBD, 0xD0, 0xE6)
    cl_box3.line.color.rgb = BLUE_LINE

    # Cluttered Label
    err_circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.6), Inches(6.75), Inches(0.2), Inches(0.2))
    err_circle2.fill.solid()
    err_circle2.fill.fore_color.rgb = RED_ERR
    err_circle2.line.fill.background()
    
    lbl_cluttered = slide.shapes.add_textbox(Inches(10.8), Inches(6.65), Inches(1.5), Inches(0.3))
    lbl_cluttered.text_frame.text = "拥挤布局"
    lbl_cluttered.text_frame.paragraphs[0].font.size = Pt(14)
    lbl_cluttered.text_frame.paragraphs[0].font.bold = True

    # 5. Page Number
    tb_page = slide.shapes.add_textbox(Inches(11.8), Inches(6.8), Inches(1.2), Inches(0.5))
    p_page = tb_page.text_frame.paragraphs[0]
    p_page.text = "08 / 11"
    p_page.font.size = Pt(20)
    p_page.font.bold = True
    p_page.font.color.rgb = GRAY_TEXT
    p_page.font.name = "Microsoft YaHei"
    p_page.alignment = PP_ALIGN.RIGHT



# ── Slide 9 ──

def build_slide_9(slide):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(10), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "高手境界：简洁有力的视觉哲学"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    # Divider Line
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.8), Inches(1.3), Inches(12.5), Inches(1.3))
    line.line.color.rgb = RGBColor(0x1F, 0x4E, 0x79)
    line.line.width = Pt(1.5)

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(10), Inches(0.6))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "真正的专业不需要花哨的装饰 (20-24pt)"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(0x7F, 0x7F, 0x7F)

    # --- Item 1 ---
    # Icon 1: Eraser
    eraser_line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.9), Inches(3.9), Inches(1.6), Inches(3.9))
    eraser_line.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    eraser_line.line.width = Pt(1.5)

    eraser = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(3.0), Inches(0.8), Inches(0.45))
    eraser.rotation = -45
    eraser.fill.solid()
    eraser.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFA)
    eraser.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    eraser.line.width = Pt(1.5)

    x1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.6), Inches(3.7), Inches(1.8), Inches(3.9))
    x1.line.color.rgb = RGBColor(0xE5, 0x39, 0x35)
    x1.line.width = Pt(2)
    x2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.8), Inches(3.7), Inches(1.6), Inches(3.9))
    x2.line.color.rgb = RGBColor(0xE5, 0x39, 0x35)
    x2.line.width = Pt(2)

    # Bullet 1
    dot1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.4), Inches(3.45), Inches(0.08), Inches(0.08))
    dot1.fill.solid()
    dot1.fill.fore_color.rgb = RGBColor(0x00, 0x70, 0xC0)
    dot1.line.fill.background()

    # Text 1
    tb1 = slide.shapes.add_textbox(Inches(2.6), Inches(3.2), Inches(8), Inches(0.6))
    p1 = tb1.text_frame.paragraphs[0]
    p1.text = "删掉所有不承载信息的装饰性元素"
    p1.font.name = "Microsoft YaHei"
    p1.font.size = Pt(22)
    p1.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    # --- Item 2 ---
    # Icon 2: Circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.9), Inches(4.7), Inches(1.0), Inches(1.0))
    circle.fill.background()
    circle.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    circle.line.width = Pt(1.5)

    dot_inner = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.35), Inches(5.15), Inches(0.1), Inches(0.1))
    dot_inner.fill.solid()
    dot_inner.fill.fore_color.rgb = RGBColor(0xDA, 0xA5, 0x20)
    dot_inner.line.fill.background()

    # Bullet 2
    dot2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.4), Inches(5.15), Inches(0.08), Inches(0.08))
    dot2.fill.solid()
    dot2.fill.fore_color.rgb = RGBColor(0x00, 0x70, 0xC0)
    dot2.line.fill.background()

    # Text 2
    tb2 = slide.shapes.add_textbox(Inches(2.6), Inches(4.9), Inches(8), Inches(0.6))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = "留白不是浪费，而是更高级的强调"
    p2.font.name = "Microsoft YaHei"
    p2.font.size = Pt(22)
    p2.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    # --- Item 3 ---
    # Icon 3: Anchor
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(6.4), Inches(0.3), Inches(1.0))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0xDA, 0xA5, 0x20)
    rect.line.fill.background()

    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.4), Inches(6.4), Inches(0.2), Inches(0.2))
    c.fill.background()
    c.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    c.line.width = Pt(1.5)

    stem = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.5), Inches(6.6), Inches(1.5), Inches(7.3))
    stem.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    stem.line.width = Pt(1.5)

    cross = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.35), Inches(6.75), Inches(1.65), Inches(6.75))
    cross.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    cross.line.width = Pt(1.5)

    l1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.1), Inches(7.0), Inches(1.5), Inches(7.3))
    l1.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    l1.line.width = Pt(1.5)
    
    l2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.9), Inches(7.0), Inches(1.5), Inches(7.3))
    l2.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    l2.line.width = Pt(1.5)

    a1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.1), Inches(7.0), Inches(1.1), Inches(7.15))
    a1.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    a1.line.width = Pt(1.5)
    
    a2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.1), Inches(7.0), Inches(1.25), Inches(7.0))
    a2.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    a2.line.width = Pt(1.5)

    a3 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.9), Inches(7.0), Inches(1.9), Inches(7.15))
    a3.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    a3.line.width = Pt(1.5)
    
    a4 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.9), Inches(7.0), Inches(1.75), Inches(7.0))
    a4.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    a4.line.width = Pt(1.5)

    # Bullet 3
    dot3 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.4), Inches(6.85), Inches(0.08), Inches(0.08))
    dot3.fill.solid()
    dot3.fill.fore_color.rgb = RGBColor(0x00, 0x70, 0xC0)
    dot3.line.fill.background()

    # Text 3
    tb3 = slide.shapes.add_textbox(Inches(2.6), Inches(6.6), Inches(8), Inches(0.6))
    p3 = tb3.text_frame.paragraphs[0]
    p3.text = "简洁即是力量，克制即是专业"
    p3.font.name = "Microsoft YaHei"
    p3.font.size = Pt(22)
    p3.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    # Footer
    footer = slide.shapes.add_textbox(Inches(12.0), Inches(7.0), Inches(1.0), Inches(0.4))
    p_foot = footer.text_frame.paragraphs[0]
    p_foot.text = "第9页"
    p_foot.font.name = "Microsoft YaHei"
    p_foot.font.size = Pt(12)
    p_foot.font.color.rgb = RGBColor(0x7F, 0x7F, 0x7F)



# ── Slide 10 ──

def build_slide_10(slide):
    from pptx.enum.shapes import MSO_CONNECTOR
    from pptx.enum.text import MSO_ANCHOR
    
    # 1. 添加背景 (浅蓝色)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xF4, 0xF8, 0xFC)
    bg.line.fill.background()

    # 2. 添加主标题
    title_box = slide.shapes.add_textbox(Inches(2), Inches(0.6), Inches(9.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "总结回顾：PPT制作的双翼"
    run.font.name = FONT_NAME
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1C, 0x55, 0xBA)

    # 3. 添加副标题
    sub_box = slide.shapes.add_textbox(Inches(4), Inches(1.3), Inches(5.333), Inches(0.5))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "核心要点清单"
    run.font.name = FONT_NAME
    run.font.size = Pt(18)
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # 4. 绘制卡片背景和阴影
    card_width = Inches(5.4)
    card_height = Inches(4.8)
    left_x = Inches(0.8)
    right_x = Inches(7.133)
    card_y = Inches(2.2)

    # 左侧卡片阴影
    shadow_l = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_x + Inches(0.05), card_y + Inches(0.05), card_width, card_height)
    shadow_l.fill.solid()
    shadow_l.fill.fore_color.rgb = RGBColor(0xE0, 0xE8, 0xF0)
    shadow_l.line.fill.background()
    
    # 左侧卡片主体
    card_l = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_x, card_y, card_width, card_height)
    card_l.fill.solid()
    card_l.fill.fore_color.rgb = RGBColor(0xF9, 0xFB, 0xFE)
    card_l.line.color.rgb = RGBColor(0xBD, 0xD7, 0xFA)
    card_l.line.width = Pt(2)

    # 右侧卡片阴影
    shadow_r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, right_x + Inches(0.05), card_y + Inches(0.05), card_width, card_height)
    shadow_r.fill.solid()
    shadow_r.fill.fore_color.rgb = RGBColor(0xE0, 0xE8, 0xF0)
    shadow_r.line.fill.background()

    # 右侧卡片主体
    card_r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, right_x, card_y, card_width, card_height)
    card_r.fill.solid()
    card_r.fill.fore_color.rgb = RGBColor(0xF9, 0xFB, 0xFE)
    card_r.line.color.rgb = RGBColor(0xBD, 0xD7, 0xFA)
    card_r.line.width = Pt(2)

    # 5. 添加卡片标题 (带背景色以遮挡边框)
    def add_card_title(x, y, text):
        title_bg = slide.shapes.add_textbox(x, y, Inches(2.6), Inches(0.4))
        title_bg.fill.solid()
        title_bg.fill.fore_color.rgb = RGBColor(0xF4, 0xF8, 0xFC)
        title_bg.line.fill.background()
        tf = title_bg.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.name = FONT_NAME
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x1C, 0x55, 0xBA)

    add_card_title(left_x + Inches(1.4), card_y - Inches(0.2), "内容要点清单")
    add_card_title(right_x + Inches(1.4), card_y - Inches(0.2), "设计要点清单")

    # 6. 绘制中心天平图标及分割线
    center_x = Inches(6.666)
    
    # 上分割线
    l_top = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, center_x, Inches(2.2), center_x, Inches(3.8))
    l_top.line.color.rgb = RGBColor(0x5B, 0x9B, 0xD5)
    l_top.line.width = Pt(2)

    # 下分割线
    l_bot = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, center_x, Inches(5.0), center_x, Inches(6.6))
    l_bot.line.color.rgb = RGBColor(0x5B, 0x9B, 0xD5)
    l_bot.line.width = Pt(2)

    # 中心圆
    cy = Inches(4.4)
    r = Inches(0.45)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, center_x - r, cy - r, r*2, r*2)
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0xE6, 0xF0, 0xFA)
    circle.line.fill.background()

    # 天平图形
    scale_color = RGBColor(0x1C, 0x55, 0xBA)
    base = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, center_x - Inches(0.15), cy + Inches(0.1), Inches(0.3), Inches(0.15))
    base.fill.solid()
    base.fill.fore_color.rgb = scale_color
    base.line.fill.background()
    
    pillar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, center_x - Inches(0.025), cy - Inches(0.15), Inches(0.05), Inches(0.25))
    pillar.fill.solid()
    pillar.fill.fore_color.rgb = scale_color
    pillar.line.fill.background()
    
    beam = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, center_x - Inches(0.3), cy - Inches(0.15), Inches(0.6), Inches(0.04))
    beam.fill.solid()
    beam.fill.fore_color.rgb = scale_color
    beam.line.fill.background()
    
    pan_l = slide.shapes.add_shape(MSO_SHAPE.OVAL, center_x - Inches(0.4), cy + Inches(0.05), Inches(0.2), Inches(0.05))
    pan_l.fill.solid()
    pan_l.fill.fore_color.rgb = scale_color
    pan_l.line.fill.background()
    
    sl1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, center_x - Inches(0.3), cy - Inches(0.11), center_x - Inches(0.4), cy + Inches(0.05))
    sl1.line.color.rgb = scale_color
    sl2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, center_x - Inches(0.3), cy - Inches(0.11), center_x - Inches(0.2), cy + Inches(0.05))
    sl2.line.color.rgb = scale_color
    
    pan_r = slide.shapes.add_shape(MSO_SHAPE.OVAL, center_x + Inches(0.2), cy + Inches(0.05), Inches(0.2), Inches(0.05))
    pan_r.fill.solid()
    pan_r.fill.fore_color.rgb = scale_color
    pan_r.line.fill.background()
    
    sr1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, center_x + Inches(0.3), cy - Inches(0.11), center_x + Inches(0.2), cy + Inches(0.05))
    sr1.line.color.rgb = scale_color
    sr2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, center_x + Inches(0.3), cy - Inches(0.11), center_x + Inches(0.4), cy + Inches(0.05))
    sr2.line.color.rgb = scale_color

    # 7. 定义绘制列表项的辅助函数
    def draw_item(x, y, icon_type, label_text, desc_text):
        icon_color = RGBColor(0x2B, 0x70, 0xC9)
        
        # 绘制图标
        if icon_type == 'doc':
            shape = slide.shapes.add_shape(MSO_SHAPE.FOLDED_CORNER, x, y, Inches(0.45), Inches(0.55))
            shape.fill.background()
            shape.line.color.rgb = icon_color
            shape.line.width = Pt(2)
            circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x+Inches(0.12), y+Inches(0.25), Inches(0.2), Inches(0.2))
            circle.fill.background()
            circle.line.color.rgb = icon_color
            circle.line.width = Pt(1.5)
        elif icon_type == 'palette':
            shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.55), Inches(0.55))
            shape.fill.background()
            shape.line.color.rgb = icon_color
            shape.line.width = Pt(2)
            for cx, cy_ in [(0.12, 0.12), (0.35, 0.15), (0.12, 0.35)]:
                sc = slide.shapes.add_shape(MSO_SHAPE.OVAL, x+Inches(cx), y+Inches(cy_), Inches(0.08), Inches(0.08))
                sc.fill.solid()
                sc.fill.fore_color.rgb = icon_color
                sc.line.fill.background()
            txt = slide.shapes.add_textbox(x+Inches(0.25), y+Inches(0.25), Inches(0.3), Inches(0.3))
            tf = txt.text_frame
            tf.margin_left, tf.margin_right, tf.margin_top, tf.margin_bottom = 0, 0, 0, 0
            p = txt.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = "A"
            run.font.name = "Arial"
            run.font.size = Pt(12)
            run.font.bold = True
            run.font.color.rgb = icon_color
        elif icon_type == 'funnel':
            shape = slide.shapes.add_shape(MSO_SHAPE.FUNNEL, x, y, Inches(0.55), Inches(0.55))
            shape.fill.background()
            shape.line.color.rgb = icon_color
            shape.line.width = Pt(2)
        elif icon_type == 'layout':
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.55), Inches(0.45))
            shape.fill.background()
            shape.line.color.rgb = icon_color
            shape.line.width = Pt(2)
            l1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, y+Inches(0.15), x+Inches(0.55), y+Inches(0.15))
            l1.line.color.rgb = icon_color
            l1.line.width = Pt(1.5)
            l2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x+Inches(0.2), y+Inches(0.15), x+Inches(0.2), y+Inches(0.45))
            l2.line.color.rgb = icon_color
            l2.line.width = Pt(1.5)
        elif icon_type == 'align':
            for i in range(3):
                r_w = Inches(0.25) if i % 2 == 0 else Inches(0.4)
                r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y+Inches(i*0.18), r_w, Inches(0.1))
                r.fill.background()
                r.line.color.rgb = icon_color
                r.line.width = Pt(1.5)
            vl = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x+Inches(0.55), y, x+Inches(0.55), y+Inches(0.46))
            vl.line.color.rgb = icon_color
            vl.line.width = Pt(1.5)
        elif icon_type == 'hierarchy':
            r1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x+Inches(0.15), y, Inches(0.25), Inches(0.12))
            r1.fill.background()
            r1.line.color.rgb = icon_color
            r1.line.width = Pt(1.5)
            for i in range(3):
                r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x+Inches(i*0.2), y+Inches(0.3), Inches(0.15), Inches(0.1))
                r.fill.background()
                r.line.color.rgb = icon_color
                r.line.width = Pt(1.5)
            c1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x+Inches(0.275), y+Inches(0.12), x+Inches(0.275), y+Inches(0.2))
            c1.line.color.rgb = icon_color
            c2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x+Inches(0.075), y+Inches(0.2), x+Inches(0.475), y+Inches(0.2))
            c2.line.color.rgb = icon_color
            for i in range(3):
                c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x+Inches(i*0.2+0.075), y+Inches(0.2), x+Inches(i*0.2+0.075), y+Inches(0.3))
                c.line.color.rgb = icon_color

        # 绘制对勾
        chk_x = x + Inches(0.7)
        chk_y = y + Inches(0.02)
        chk_bg = slide.shapes.add_shape(MSO_SHAPE.OVAL, chk_x, chk_y, Inches(0.22), Inches(0.22))
        chk_bg.fill.solid()
        chk_bg.fill.fore_color.rgb = RGBColor(0x43, 0xA0, 0x47)
        chk_bg.line.fill.background()
        
        chk_txt = slide.shapes.add_textbox(chk_x - Inches(0.05), chk_y - Inches(0.05), Inches(0.32), Inches(0.32))
        tf = chk_txt.text_frame
        tf.margin_left, tf.margin_right, tf.margin_top, tf.margin_bottom = 0, 0, 0, 0
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = "✔"
        run.font.name = "Segoe UI Symbol"
        run.font.size = Pt(10)
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        # 绘制标签
        lbl_x = chk_x + Inches(0.3)
        lbl_y = y - Inches(0.05)
        lbl_box = slide.shapes.add_textbox(lbl_x, lbl_y, Inches(3.8), Inches(0.4))
        tf = lbl_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = label_text
        run.font.name = FONT_NAME
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x22, 0x22, 0x22)

        # 绘制描述
        desc_y = lbl_y + Inches(0.35)
        desc_box = slide.shapes.add_textbox(lbl_x, desc_y, Inches(3.8), Inches(0.6))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = desc_text
        run.font.name = FONT_NAME
        run.font.size = Pt(11)
        run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    # 8. 填充左侧列表项
    item_l_x = left_x + Inches(0.4)
    draw_item(item_l_x, Inches(2.8), 'doc', "清晰：一页一事，拒绝文字堆砌", "确保每个页面聚焦一个核心观点，避免信息过载。")
    draw_item(item_l_x, Inches(4.0), 'palette', "统一：配色字体，保持风格一致", "建立标准的色彩和字体规范，维护整体视觉的连贯性。")
    draw_item(item_l_x, Inches(5.2), 'funnel', "简洁：去繁就简，追求高效沟通", "删除所有非必要元素，用最短路径传达最关键信息。")

    # 9. 填充右侧列表项
    item_r_x = right_x + Inches(0.4)
    draw_item(item_r_x, Inches(2.8), 'layout', "留白：呼吸空间，引导视觉焦点", "保持页面留白率在40%以上，让内容有足够的呼吸感。")
    draw_item(item_r_x, Inches(4.0), 'align', "对齐：严谨规范，建立秩序美感", "严格遵循左对齐或居中对齐原则，构建清晰的视觉轴线。")
    draw_item(item_r_x, Inches(5.2), 'hierarchy', "层级：主次分明，提升阅读效率", "通过字号、粗细和颜色的对比，明确信息的优先级。")

    # 10. 添加页脚
    footer = slide.shapes.add_textbox(Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "Page 10 of 11"
    run.font.name = "Arial"
    run.font.size = Pt(10)
    run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)



# ── Slide 11 ──

def build_slide_11(slide):
    # 1. Background Grid (Top part - Light Blue)
    for i in range(1, 14):
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(i), 0, Inches(i), Inches(7.5))
        line.line.color.rgb = RGBColor(0xE8, 0xF0, 0xF8)
        line.line.width = Pt(0.5)
    for i in range(1, 8):
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, 0, Inches(i), Inches(13.333), Inches(i))
        line.line.color.rgb = RGBColor(0xE8, 0xF0, 0xF8)
        line.line.width = Pt(0.5)

    # 2. Bottom Dark Blue Background
    bottom_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.2), Inches(13.333), Inches(3.3))
    bottom_rect.fill.solid()
    bottom_rect.fill.fore_color.rgb = RGBColor(0x15, 0x43, 0x85)
    bottom_rect.line.fill.background()

    # Bottom Grid (Overlay on dark blue)
    for i in range(1, 14):
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(i), Inches(4.2), Inches(i), Inches(7.5))
        line.line.color.rgb = RGBColor(0x25, 0x53, 0x95)
        line.line.width = Pt(0.5)
    for i in range(5, 8):
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, 0, Inches(i), Inches(13.333), Inches(i))
        line.line.color.rgb = RGBColor(0x25, 0x53, 0x95)
        line.line.width = Pt(0.5)

    # 3. Title Text
    title_box = slide.shapes.add_textbox(Inches(2.66), Inches(0.8), Inches(8), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "感谢观看"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x00, 0x52, 0xCC)
    p.font.name = "Microsoft YaHei"

    subtitle_box = slide.shapes.add_textbox(Inches(2.66), Inches(2.1), Inches(8), Inches(0.8))
    tf_sub = subtitle_box.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "立即开始你的专业演示之旅"
    p_sub.alignment = PP_ALIGN.CENTER
    p_sub.font.size = Pt(28)
    p_sub.font.color.rgb = RGBColor(0x00, 0x52, 0xCC)
    p_sub.font.name = "Microsoft YaHei"

    # 4. Middle Icons and Text
    icon_color = RGBColor(0x00, 0x52, 0xCC)
    text_color = RGBColor(0x00, 0x00, 0x00)

    # Item 1: Compass
    compass = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.5), Inches(3.2), Inches(0.6), Inches(0.6))
    compass.fill.background()
    compass.line.color.rgb = icon_color
    compass.line.width = Pt(2)
    needle = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.65), Inches(3.65), Inches(1.95), Inches(3.35))
    needle.line.color.rgb = icon_color
    needle.line.width = Pt(2)
    
    tb1 = slide.shapes.add_textbox(Inches(2.2), Inches(3.1), Inches(2.5), Inches(0.8))
    tf1 = tb1.text_frame
    p1 = tf1.paragraphs[0]
    p1.text = "1. 实践是提升PPT能\n力的唯一捷径"
    p1.font.size = Pt(16)
    p1.font.color.rgb = text_color
    p1.font.name = "Microsoft YaHei"

    # Item 2: Lightbulb
    bulb = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.2), Inches(3.2), Inches(0.5), Inches(0.5))
    bulb.fill.background()
    bulb.line.color.rgb = icon_color
    bulb.line.width = Pt(2)
    base = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.35), Inches(3.7), Inches(0.2), Inches(0.15))
    base.fill.background()
    base.line.color.rgb = icon_color
    base.line.width = Pt(2)
    ray1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(5.45), Inches(3.1), Inches(5.45), Inches(2.95))
    ray1.line.color.rgb = icon_color
    ray1.line.width = Pt(1.5)
    ray2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(5.05), Inches(3.45), Inches(4.9), Inches(3.45))
    ray2.line.color.rgb = icon_color
    ray2.line.width = Pt(1.5)
    ray3 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(5.85), Inches(3.45), Inches(6.0), Inches(3.45))
    ray3.line.color.rgb = icon_color
    ray3.line.width = Pt(1.5)
    
    tb2 = slide.shapes.add_textbox(Inches(5.9), Inches(3.1), Inches(2.5), Inches(0.8))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "2. 保持好奇，持续优\n化视觉表达"
    p2.font.size = Pt(16)
    p2.font.color.rgb = text_color
    p2.font.name = "Microsoft YaHei"

    # Item 3: Rocket
    rocket = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.9), Inches(3.2), Inches(0.4), Inches(0.6))
    rocket.rotation = 45
    rocket.fill.background()
    rocket.line.color.rgb = icon_color
    rocket.line.width = Pt(2)
    wing1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(8.9), Inches(3.6), Inches(8.7), Inches(3.8))
    wing1.line.color.rgb = icon_color
    wing1.line.width = Pt(2)
    wing2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(9.2), Inches(3.7), Inches(9.4), Inches(3.9))
    wing2.line.color.rgb = icon_color
    wing2.line.width = Pt(2)
    
    tb3 = slide.shapes.add_textbox(Inches(9.6), Inches(3.2), Inches(2.5), Inches(0.8))
    tf3 = tb3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "3. 期待您的精彩呈现"
    p3.font.size = Pt(16)
    p3.font.color.rgb = text_color
    p3.font.name = "Microsoft YaHei"

    # 5. Bottom Section - Faint Chart
    chart_color = RGBColor(0x4A, 0x76, 0xB5)
    
    # Chart Grid
    for x in [4.5, 5.0, 5.5, 6.0, 6.5, 7.0, 7.5, 8.0, 8.5, 9.0]:
        v_line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(4.8), Inches(x), Inches(6.2))
        v_line.line.color.rgb = chart_color
        v_line.line.width = Pt(0.5)
        
    # X and Y axis
    x_axis = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(4.0), Inches(6.2), Inches(9.3), Inches(6.2))
    x_axis.line.color.rgb = chart_color
    x_axis.line.width = Pt(1)
    y_axis = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(4.0), Inches(4.8), Inches(4.0), Inches(6.2))
    y_axis.line.color.rgb = chart_color
    y_axis.line.width = Pt(1)
    
    # Trend line
    t1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(4.0), Inches(6.0), Inches(5.5), Inches(5.7))
    t1.line.color.rgb = chart_color
    t1.line.width = Pt(1.5)
    t2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(5.5), Inches(5.7), Inches(7.0), Inches(5.9))
    t2.line.color.rgb = chart_color
    t2.line.width = Pt(1.5)
    t3 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(7.0), Inches(5.9), Inches(9.0), Inches(4.8))
    t3.line.color.rgb = chart_color
    t3.line.width = Pt(1.5)

    # 6. "Thank You" Text
    ty_box = slide.shapes.add_textbox(Inches(2.66), Inches(4.6), Inches(8), Inches(1.5))
    tf_ty = ty_box.text_frame
    p_ty = tf_ty.paragraphs[0]
    p_ty.text = "Thank You"
    p_ty.alignment = PP_ALIGN.CENTER
    p_ty.font.size = Pt(72)
    p_ty.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p_ty.font.name = "Arial"

    # 7. Footer
    # QR Code Scanner Icon
    # Brackets
    brackets = [
        (3.5, 6.5, 3.6, 6.5), (3.5, 6.5, 3.5, 6.6), # Top-left
        (4.0, 6.5, 4.1, 6.5), (4.1, 6.5, 4.1, 6.6), # Top-right
        (3.5, 7.0, 3.6, 7.0), (3.5, 6.9, 3.5, 7.0), # Bottom-left
        (4.0, 7.0, 4.1, 7.0), (4.1, 6.9, 4.1, 7.0)  # Bottom-right
    ]
    for x1, y1, x2, y2 in brackets:
        bl = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        bl.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        bl.line.width = Pt(1.5)
        
    # Inner square and line
    sq = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.65), Inches(6.65), Inches(0.3), Inches(0.3))
    sq.fill.background()
    sq.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    sq.line.width = Pt(1)
    scan_line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(3.4), Inches(6.75), Inches(4.2), Inches(6.75))
    scan_line.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    scan_line.line.width = Pt(1.5)

    # Footer Text
    footer_box = slide.shapes.add_textbox(Inches(4.2), Inches(6.55), Inches(7.0), Inches(0.5))
    tf_footer = footer_box.text_frame
    p_footer = tf_footer.paragraphs[0]
    p_footer.text = "扫描二维码联系 | 联系方式: support@example.com | www.example.com"
    p_footer.font.size = Pt(11)
    p_footer.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p_footer.font.name = "Microsoft YaHei"

    # Page Number
    page_box = slide.shapes.add_textbox(Inches(11.8), Inches(6.8), Inches(1.0), Inches(0.5))
    tf_page = page_box.text_frame
    p_page = tf_page.paragraphs[0]
    p_page.text = "11 / 11"
    p_page.alignment = PP_ALIGN.RIGHT
    p_page.font.size = Pt(12)
    p_page.font.color.rgb = RGBColor(0x8A, 0xB4, 0xF8)
    p_page.font.name = "Arial"



# ── Main ──
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

s0 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_1(s0)
s1 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_2(s1)
s2 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_3(s2)
s3 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_4(s3)
s4 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_5(s4)
s5 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_6(s5)
s6 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_7(s6)
s7 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_8(s7)
s8 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_9(s8)
s9 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_10(s9)
s10 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_11(s10)
prs.save(OUTPUT_PATH)
