from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# ── Color Palette ──
BLUE_HEADER = RGBColor(0x5B, 0x9B, 0xD5)
BLUE_DARK   = RGBColor(0x4A, 0x86, 0xC8)
CYAN        = RGBColor(0x00, 0xBC, 0xD4)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x33, 0x33, 0x33)
GRAY_TEXT   = RGBColor(0x55, 0x55, 0x55)
GRAY_BAR    = RGBColor(0xB0, 0xBE, 0xC5)
RED         = RGBColor(0xE5, 0x39, 0x35)
GREEN       = RGBColor(0x43, 0xA0, 0x47)
ORANGE      = RGBColor(0xFF, 0x98, 0x00)

ICON_BG  = RGBColor(0xE3, 0xE8, 0xED)
ICON_FG  = RGBColor(0x54, 0x6E, 0x7A)
FONT_NAME = "Microsoft YaHei"
SLIDE_WIDTH = Inches(13.333)
HEADER_H    = Inches(0.75)
SUBTITLE_Y  = Inches(0.95)


def add_header_banner(slide, title_text, bg_color=None):
    if bg_color is None:
        bg_color = BLUE_HEADER
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, HEADER_H
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = bg_color
    banner.line.fill.background()
    tf = banner.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(26)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME


def add_subtitle(slide, text, left, top, width=Inches(12), font_size=Pt(18)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.font.name = FONT_NAME
    return txBox


def add_icon_box(slide, left, top, symbol, size=Inches(0.48)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ICON_BG
    shape.line.fill.background()
    shape.adjustments[0] = 0.25
    tf = shape.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = symbol
    p.font.size = Pt(18)
    p.font.color.rgb = ICON_FG
    p.font.bold = False
    return shape


def add_bullet_item(slide, left, top, symbol, label, description,
                    width=Inches(5.5), desc_size=Pt(13)):
    add_icon_box(slide, left, top, symbol)
    text_left = left + Inches(0.65)
    txBox = slide.shapes.add_textbox(text_left, top - Inches(0.02), width, Inches(0.65))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run_label = p.add_run()
    run_label.text = label + "\uff1a"
    run_label.font.size = Pt(14)
    run_label.font.color.rgb = BLACK
    run_label.font.bold = True
    run_label.font.name = FONT_NAME
    run_desc = p.add_run()
    run_desc.text = description
    run_desc.font.size = desc_size
    run_desc.font.color.rgb = GRAY_TEXT
    run_desc.font.bold = False
    run_desc.font.name = FONT_NAME
    return txBox


def add_conclusion_box(slide, left, top, width, text, font_size=Pt(13)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = BLACK
    run.font.bold = True
    run.font.name = FONT_NAME
    return txBox


def add_table(slide, left, top, width, height, rows, cols, data,
              header_color=None, col_widths=None):
    if header_color is None:
        header_color = BLUE_HEADER
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if data[r][c] is not None else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(11)
                    run.font.name = FONT_NAME
                    if r == 0:
                        run.font.color.rgb = WHITE
                        run.font.bold = True
                    else:
                        run.font.color.rgb = BLACK
                        run.font.bold = False
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else RGBColor(0xF5, 0xF5, 0xF5)
            cell.margin_left = Pt(5)
            cell.margin_right = Pt(5)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
    return table_shape


def add_bar_chart(slide, left, top, width, height,
                  categories, values, title="", bar_colors=None):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', values)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    chart.chart_style = 2
    plot = chart.plots[0]
    plot.gap_width = 100
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CYAN
    series.has_data_labels = True
    dl = series.data_labels
    dl.font.size = Pt(13)
    dl.font.bold = True
    dl.font.color.rgb = BLACK
    dl.number_format = '0.#'
    dl.show_value = True
    dl.label_position = XL_LABEL_POSITION.OUTSIDE_END
    if bar_colors:
        for i, color in enumerate(bar_colors):
            pt = series.points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = color
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(12)
    cat_axis.tick_labels.font.name = FONT_NAME
    cat_axis.major_tick_mark = 2
    cat_axis.format.line.fill.background()
    val_axis = chart.value_axis
    val_axis.visible = False
    val_axis.major_tick_mark = 2
    val_axis.format.line.fill.background()
    val_axis.major_gridlines.format.line.fill.background()
    if title:
        chart.has_title = True
        ct = chart.chart_title.text_frame.paragraphs[0]
        ct.text = title
        ct.font.size = Pt(14)
        ct.font.bold = True
        ct.font.name = FONT_NAME
    else:
        chart.has_title = False
    return chart_frame


def add_callout_label(slide, left, top, text, bg_color=None, font_size=Pt(11)):
    if bg_color is None:
        bg_color = CYAN
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.3), Inches(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME
    return shape


def add_data_card(slide, left, top, width, height, value, label,
                  value_color=None, bg_color=None):
    if value_color is None:
        value_color = CYAN
    if bg_color is None:
        bg_color = WHITE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(3)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run1 = p1.add_run()
    run1.text = str(value)
    run1.font.size = Pt(24)
    run1.font.color.rgb = value_color
    run1.font.bold = True
    run1.font.name = FONT_NAME
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(10)
    run2.font.color.rgb = GRAY_TEXT
    run2.font.bold = False
    run2.font.name = FONT_NAME
    return shape


OUTPUT_PATH = r"C:\Users\Administrator\Desktop\notebook\my_notebooklm_with_ppteditable\projects\如何成为PPT高手\最终文档\ppt_slides\slide_09.pptx"


# ── Slide 1 ──

def build_slide_1(slide):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(10), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "高手境界：简洁有力的视觉哲学"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    # Divider Line
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.8), Inches(1.3), Inches(12.5), Inches(1.3))
    line.line.color.rgb = RGBColor(0x1F, 0x4E, 0x79)
    line.line.width = Pt(1.5)

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(10), Inches(0.6))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "真正的专业不需要花哨的装饰 (20-24pt)"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(0x7F, 0x7F, 0x7F)

    # --- Item 1 ---
    # Icon 1: Eraser
    eraser_line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.9), Inches(3.9), Inches(1.6), Inches(3.9))
    eraser_line.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    eraser_line.line.width = Pt(1.5)

    eraser = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(3.0), Inches(0.8), Inches(0.45))
    eraser.rotation = -45
    eraser.fill.solid()
    eraser.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFA)
    eraser.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    eraser.line.width = Pt(1.5)

    x1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.6), Inches(3.7), Inches(1.8), Inches(3.9))
    x1.line.color.rgb = RGBColor(0xE5, 0x39, 0x35)
    x1.line.width = Pt(2)
    x2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.8), Inches(3.7), Inches(1.6), Inches(3.9))
    x2.line.color.rgb = RGBColor(0xE5, 0x39, 0x35)
    x2.line.width = Pt(2)

    # Bullet 1
    dot1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.4), Inches(3.45), Inches(0.08), Inches(0.08))
    dot1.fill.solid()
    dot1.fill.fore_color.rgb = RGBColor(0x00, 0x70, 0xC0)
    dot1.line.fill.background()

    # Text 1
    tb1 = slide.shapes.add_textbox(Inches(2.6), Inches(3.2), Inches(8), Inches(0.6))
    p1 = tb1.text_frame.paragraphs[0]
    p1.text = "删掉所有不承载信息的装饰性元素"
    p1.font.name = "Microsoft YaHei"
    p1.font.size = Pt(22)
    p1.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    # --- Item 2 ---
    # Icon 2: Circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.9), Inches(4.7), Inches(1.0), Inches(1.0))
    circle.fill.background()
    circle.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    circle.line.width = Pt(1.5)

    dot_inner = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.35), Inches(5.15), Inches(0.1), Inches(0.1))
    dot_inner.fill.solid()
    dot_inner.fill.fore_color.rgb = RGBColor(0xDA, 0xA5, 0x20)
    dot_inner.line.fill.background()

    # Bullet 2
    dot2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.4), Inches(5.15), Inches(0.08), Inches(0.08))
    dot2.fill.solid()
    dot2.fill.fore_color.rgb = RGBColor(0x00, 0x70, 0xC0)
    dot2.line.fill.background()

    # Text 2
    tb2 = slide.shapes.add_textbox(Inches(2.6), Inches(4.9), Inches(8), Inches(0.6))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = "留白不是浪费，而是更高级的强调"
    p2.font.name = "Microsoft YaHei"
    p2.font.size = Pt(22)
    p2.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    # --- Item 3 ---
    # Icon 3: Anchor
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(6.4), Inches(0.3), Inches(1.0))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0xDA, 0xA5, 0x20)
    rect.line.fill.background()

    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.4), Inches(6.4), Inches(0.2), Inches(0.2))
    c.fill.background()
    c.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    c.line.width = Pt(1.5)

    stem = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.5), Inches(6.6), Inches(1.5), Inches(7.3))
    stem.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    stem.line.width = Pt(1.5)

    cross = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.35), Inches(6.75), Inches(1.65), Inches(6.75))
    cross.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    cross.line.width = Pt(1.5)

    l1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.1), Inches(7.0), Inches(1.5), Inches(7.3))
    l1.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    l1.line.width = Pt(1.5)
    
    l2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.9), Inches(7.0), Inches(1.5), Inches(7.3))
    l2.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    l2.line.width = Pt(1.5)

    a1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.1), Inches(7.0), Inches(1.1), Inches(7.15))
    a1.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    a1.line.width = Pt(1.5)
    
    a2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.1), Inches(7.0), Inches(1.25), Inches(7.0))
    a2.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    a2.line.width = Pt(1.5)

    a3 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.9), Inches(7.0), Inches(1.9), Inches(7.15))
    a3.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    a3.line.width = Pt(1.5)
    
    a4 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.9), Inches(7.0), Inches(1.75), Inches(7.0))
    a4.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    a4.line.width = Pt(1.5)

    # Bullet 3
    dot3 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.4), Inches(6.85), Inches(0.08), Inches(0.08))
    dot3.fill.solid()
    dot3.fill.fore_color.rgb = RGBColor(0x00, 0x70, 0xC0)
    dot3.line.fill.background()

    # Text 3
    tb3 = slide.shapes.add_textbox(Inches(2.6), Inches(6.6), Inches(8), Inches(0.6))
    p3 = tb3.text_frame.paragraphs[0]
    p3.text = "简洁即是力量，克制即是专业"
    p3.font.name = "Microsoft YaHei"
    p3.font.size = Pt(22)
    p3.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    # Footer
    footer = slide.shapes.add_textbox(Inches(12.0), Inches(7.0), Inches(1.0), Inches(0.4))
    p_foot = footer.text_frame.paragraphs[0]
    p_foot.text = "第9页"
    p_foot.font.name = "Microsoft YaHei"
    p_foot.font.size = Pt(12)
    p_foot.font.color.rgb = RGBColor(0x7F, 0x7F, 0x7F)



# ── Main ──
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

s0 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_1(s0)
prs.save(OUTPUT_PATH)
