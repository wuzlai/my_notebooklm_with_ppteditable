from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# ── Color Palette ──
BLUE_HEADER = RGBColor(0x5B, 0x9B, 0xD5)
BLUE_DARK   = RGBColor(0x4A, 0x86, 0xC8)
CYAN        = RGBColor(0x00, 0xBC, 0xD4)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x33, 0x33, 0x33)
GRAY_TEXT   = RGBColor(0x55, 0x55, 0x55)
GRAY_BAR    = RGBColor(0xB0, 0xBE, 0xC5)
RED         = RGBColor(0xE5, 0x39, 0x35)
GREEN       = RGBColor(0x43, 0xA0, 0x47)
ORANGE      = RGBColor(0xFF, 0x98, 0x00)

ICON_BG  = RGBColor(0xE3, 0xE8, 0xED)
ICON_FG  = RGBColor(0x54, 0x6E, 0x7A)
FONT_NAME = "Microsoft YaHei"
SLIDE_WIDTH = Inches(13.333)
HEADER_H    = Inches(0.75)
SUBTITLE_Y  = Inches(0.95)


def add_header_banner(slide, title_text, bg_color=None):
    if bg_color is None:
        bg_color = BLUE_HEADER
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, HEADER_H
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = bg_color
    banner.line.fill.background()
    tf = banner.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(26)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME


def add_subtitle(slide, text, left, top, width=Inches(12), font_size=Pt(18)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.font.name = FONT_NAME
    return txBox


def add_icon_box(slide, left, top, symbol, size=Inches(0.48)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ICON_BG
    shape.line.fill.background()
    shape.adjustments[0] = 0.25
    tf = shape.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = symbol
    p.font.size = Pt(18)
    p.font.color.rgb = ICON_FG
    p.font.bold = False
    return shape


def add_bullet_item(slide, left, top, symbol, label, description,
                    width=Inches(5.5), desc_size=Pt(13)):
    add_icon_box(slide, left, top, symbol)
    text_left = left + Inches(0.65)
    txBox = slide.shapes.add_textbox(text_left, top - Inches(0.02), width, Inches(0.65))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run_label = p.add_run()
    run_label.text = label + "\uff1a"
    run_label.font.size = Pt(14)
    run_label.font.color.rgb = BLACK
    run_label.font.bold = True
    run_label.font.name = FONT_NAME
    run_desc = p.add_run()
    run_desc.text = description
    run_desc.font.size = desc_size
    run_desc.font.color.rgb = GRAY_TEXT
    run_desc.font.bold = False
    run_desc.font.name = FONT_NAME
    return txBox


def add_conclusion_box(slide, left, top, width, text, font_size=Pt(13)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = BLACK
    run.font.bold = True
    run.font.name = FONT_NAME
    return txBox


def add_table(slide, left, top, width, height, rows, cols, data,
              header_color=None, col_widths=None):
    if header_color is None:
        header_color = BLUE_HEADER
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if data[r][c] is not None else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(11)
                    run.font.name = FONT_NAME
                    if r == 0:
                        run.font.color.rgb = WHITE
                        run.font.bold = True
                    else:
                        run.font.color.rgb = BLACK
                        run.font.bold = False
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else RGBColor(0xF5, 0xF5, 0xF5)
            cell.margin_left = Pt(5)
            cell.margin_right = Pt(5)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
    return table_shape


def add_bar_chart(slide, left, top, width, height,
                  categories, values, title="", bar_colors=None):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', values)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    chart.chart_style = 2
    plot = chart.plots[0]
    plot.gap_width = 100
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CYAN
    series.has_data_labels = True
    dl = series.data_labels
    dl.font.size = Pt(13)
    dl.font.bold = True
    dl.font.color.rgb = BLACK
    dl.number_format = '0.#'
    dl.show_value = True
    dl.label_position = XL_LABEL_POSITION.OUTSIDE_END
    if bar_colors:
        for i, color in enumerate(bar_colors):
            pt = series.points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = color
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(12)
    cat_axis.tick_labels.font.name = FONT_NAME
    cat_axis.major_tick_mark = 2
    cat_axis.format.line.fill.background()
    val_axis = chart.value_axis
    val_axis.visible = False
    val_axis.major_tick_mark = 2
    val_axis.format.line.fill.background()
    val_axis.major_gridlines.format.line.fill.background()
    if title:
        chart.has_title = True
        ct = chart.chart_title.text_frame.paragraphs[0]
        ct.text = title
        ct.font.size = Pt(14)
        ct.font.bold = True
        ct.font.name = FONT_NAME
    else:
        chart.has_title = False
    return chart_frame


def add_callout_label(slide, left, top, text, bg_color=None, font_size=Pt(11)):
    if bg_color is None:
        bg_color = CYAN
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.3), Inches(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME
    return shape


def add_data_card(slide, left, top, width, height, value, label,
                  value_color=None, bg_color=None):
    if value_color is None:
        value_color = CYAN
    if bg_color is None:
        bg_color = WHITE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(3)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run1 = p1.add_run()
    run1.text = str(value)
    run1.font.size = Pt(24)
    run1.font.color.rgb = value_color
    run1.font.bold = True
    run1.font.name = FONT_NAME
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(10)
    run2.font.color.rgb = GRAY_TEXT
    run2.font.bold = False
    run2.font.name = FONT_NAME
    return shape


OUTPUT_PATH = r"C:\Users\Administrator\Desktop\notebook\my_notebooklm_with_ppteditable\projects\如何成为PPT高手\最终文档\ppt_slides\slide_02.pptx"


# ── Slide 1 ──

def build_slide_1(slide):
    from pptx.enum.shapes import MSO_CONNECTOR
    
    # Colors
    BLUE_TITLE = RGBColor(0x1A, 0x56, 0xBA)
    DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
    BLUE_LINE = RGBColor(0x3B, 0x7E, 0xC6)
    ORANGE_TEXT = RGBColor(0xDE, 0x9B, 0x35)
    BLACK_TEXT = RGBColor(0x00, 0x00, 0x00)
    SHADOW_COLOR = RGBColor(0xF0, 0xF4, 0xF8)
    BORDER_COLOR = RGBColor(0xE8, 0xE8, 0xE8)

    # 1. Add Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.0), Inches(0.8))
    tf_title = title_box.text_frame
    p_title = tf_title.paragraphs[0]
    p_title.text = "目录：构建专业PPT的蓝图"
    p_title.font.name = "Microsoft YaHei"
    p_title.font.size = Pt(36)
    p_title.font.bold = True
    p_title.font.color.rgb = BLUE_TITLE

    # 2. Add Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.0), Inches(0.5))
    tf_sub = sub_box.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "本次分享的核心框架"
    p_sub.font.name = "Microsoft YaHei"
    p_sub.font.size = Pt(20)
    p_sub.font.bold = True
    p_sub.font.color.rgb = DARK_GRAY

    # Helper function to create styled text boxes
    def add_node_box(left, top, width, height, text_parts):
        # Shadow
        shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.06), top + Inches(0.06), width, height)
        shadow.fill.solid()
        shadow.fill.fore_color.rgb = SHADOW_COLOR
        shadow.line.fill.background()

        # Main Box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        box.line.color.rgb = BORDER_COLOR
        box.line.width = Pt(1)

        tf = box.text_frame
        tf.margin_left = Inches(0.3)
        tf.margin_top = Inches(0.15)
        p = tf.paragraphs[0]

        for text, color, is_bold in text_parts:
            run = p.add_run()
            run.text = text
            run.font.name = "Microsoft YaHei"
            run.font.size = Pt(18)
            run.font.color.rgb = color
            run.font.bold = is_bold

    # 3. Main Vertical Timeline
    v_line_x = 1.8
    y_nodes = [2.6, 3.8, 5.0, 6.2]
    
    main_v_line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(v_line_x), Inches(y_nodes[0]), Inches(v_line_x), Inches(y_nodes[3]))
    main_v_line.line.color.rgb = BLUE_LINE
    main_v_line.line.width = Pt(1.5)

    # --- Node 1 ---
    y1 = y_nodes[0]
    # Circle
    c1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(v_line_x - 0.08), Inches(y1 - 0.08), Inches(0.16), Inches(0.16))
    c1.fill.solid()
    c1.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    c1.line.color.rgb = BLUE_LINE
    c1.line.width = Pt(2)
    # H-Line
    hl1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(v_line_x + 0.08), Inches(y1), Inches(2.0), Inches(y1))
    hl1.line.color.rgb = BLUE_LINE
    hl1.line.width = Pt(1.5)
    # Text Box
    add_node_box(Inches(2.0), Inches(y1 - 0.35), Inches(5.8), Inches(0.7), [
        ("1. 内容法则：一页一事，", BLACK_TEXT, True),
        ("结论先行", ORANGE_TEXT, True)
    ])
    # Icon 1 (Document)
    doc = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(y1 - 0.4), Inches(0.5), Inches(0.7))
    doc.fill.background()
    doc.line.color.rgb = BLUE_LINE
    doc.line.width = Pt(2)
    for i in range(3):
        l = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.9), Inches(y1 - 0.2 + i*0.15), Inches(1.2), Inches(y1 - 0.2 + i*0.15))
        l.line.color.rgb = BLUE_LINE
        l.line.width = Pt(1.5)
    check_bg = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.15), Inches(y1 - 0.05), Inches(0.3), Inches(0.3))
    check_bg.fill.solid()
    check_bg.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    check_bg.line.color.rgb = BLUE_LINE
    check_bg.line.width = Pt(2)
    ck1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.22), Inches(y1 + 0.1), Inches(1.28), Inches(y1 + 0.16))
    ck1.line.color.rgb = BLUE_LINE
    ck1.line.width = Pt(2)
    ck2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.28), Inches(y1 + 0.16), Inches(1.38), Inches(y1 + 0.02))
    ck2.line.color.rgb = BLUE_LINE
    ck2.line.width = Pt(2)

    # --- Node 2 ---
    y2 = y_nodes[1]
    # Circle
    c2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(v_line_x - 0.08), Inches(y2 - 0.08), Inches(0.16), Inches(0.16))
    c2.fill.solid()
    c2.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    c2.line.color.rgb = BLUE_LINE
    c2.line.width = Pt(2)
    # H-Line
    hl2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(v_line_x + 0.08), Inches(y2), Inches(2.4), Inches(y2))
    hl2.line.color.rgb = BLUE_LINE
    hl2.line.width = Pt(1.5)
    # Text Box
    add_node_box(Inches(3.4), Inches(y2 - 0.35), Inches(6.4), Inches(0.7), [
        ("2. 减法艺术：拒绝文字堆砌，追求", BLACK_TEXT, True),
        ("秒懂", ORANGE_TEXT, True)
    ])
    # Icon 2 (Trash)
    trash = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.45), Inches(y2 - 0.2), Inches(0.35), Inches(0.45))
    trash.fill.background()
    trash.line.color.rgb = BLUE_LINE
    trash.line.width = Pt(2)
    lid = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.35), Inches(y2 - 0.3), Inches(0.55), Inches(0.08))
    lid.fill.background()
    lid.line.color.rgb = BLUE_LINE
    lid.line.width = Pt(2)
    handle = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.55), Inches(y2 - 0.38), Inches(0.15), Inches(0.08))
    handle.fill.background()
    handle.line.color.rgb = BLUE_LINE
    handle.line.width = Pt(1.5)
    for i in range(3):
        vl = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(2.55 + i*0.08), Inches(y2 - 0.15), Inches(2.55 + i*0.08), Inches(y2 + 0.2))
        vl.line.color.rgb = BLUE_LINE
        vl.line.width = Pt(1.5)
    for i in range(3):
        hl = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(2.95), Inches(y2 - 0.1 + i*0.15), Inches(3.25), Inches(y2 - 0.1 + i*0.15))
        hl.line.color.rgb = BLUE_LINE
        hl.line.width = Pt(2)

    # --- Node 3 ---
    y3 = y_nodes[2]
    # Circle
    c3 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(v_line_x - 0.08), Inches(y3 - 0.08), Inches(0.16), Inches(0.16))
    c3.fill.solid()
    c3.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    c3.line.color.rgb = BLUE_LINE
    c3.line.width = Pt(2)
    # H-Line
    hl3 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(v_line_x + 0.08), Inches(y3), Inches(3.4), Inches(y3))
    hl3.line.color.rgb = BLUE_LINE
    hl3.line.width = Pt(1.5)
    # Text Box
    add_node_box(Inches(4.3), Inches(y3 - 0.35), Inches(6.4), Inches(0.7), [
        ("3. 设计规范：高度统一，建立", BLACK_TEXT, True),
        ("专业感", ORANGE_TEXT, True)
    ])
    # Icon 3 (Gears & Ruler)
    g1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.4), Inches(y3 - 0.4), Inches(0.4), Inches(0.4))
    g1.fill.background()
    g1.line.color.rgb = BLUE_LINE
    g1.line.width = Pt(2.5)
    g1_in = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.55), Inches(y3 - 0.25), Inches(0.1), Inches(0.1))
    g1_in.fill.background()
    g1_in.line.color.rgb = BLUE_LINE
    g1_in.line.width = Pt(1.5)
    g2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.75), Inches(y3 - 0.25), Inches(0.3), Inches(0.3))
    g2.fill.background()
    g2.line.color.rgb = BLUE_LINE
    g2.line.width = Pt(2)
    ruler = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.4), Inches(y3 + 0.1), Inches(0.7), Inches(0.15))
    ruler.fill.background()
    ruler.line.color.rgb = BLUE_LINE
    ruler.line.width = Pt(1.5)
    for i in range(6):
        tick = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(3.45 + i*0.1), Inches(y3 + 0.1), Inches(3.45 + i*0.1), Inches(y3 + 0.18))
        tick.line.color.rgb = BLUE_LINE

    # --- Node 4 ---
    y4 = y_nodes[3]
    # Circle
    c4 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(v_line_x - 0.08), Inches(y4 - 0.08), Inches(0.16), Inches(0.16))
    c4.fill.solid()
    c4.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    c4.line.color.rgb = BLUE_LINE
    c4.line.width = Pt(2)
    # H-Line
    hl4 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(v_line_x + 0.08), Inches(y4), Inches(4.2), Inches(y4))
    hl4.line.color.rgb = BLUE_LINE
    hl4.line.width = Pt(1.5)
    # Text Box
    add_node_box(Inches(5.3), Inches(y4 - 0.35), Inches(6.3), Inches(0.7), [
        ("4. 高手境界：简洁有力的视觉哲学", BLACK_TEXT, True)
    ])
    # Icon 4 (Mountain & Scale)
    mt = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(4.2), Inches(y4 - 0.3), Inches(0.65), Inches(0.5))
    mt.fill.background()
    mt.line.color.rgb = BLUE_LINE
    mt.line.width = Pt(2)
    snow = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(4.42), Inches(y4 - 0.3), Inches(0.21), Inches(0.15))
    snow.fill.background()
    snow.line.color.rgb = BLUE_LINE
    snow.line.width = Pt(1.5)
    
    base = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(5.15), Inches(y4 + 0.1), Inches(0.3), Inches(0.1))
    base.fill.background()
    base.line.color.rgb = BLUE_LINE
    base.line.width = Pt(1.5)
    post = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(5.3), Inches(y4 - 0.3), Inches(5.3), Inches(y4 + 0.1))
    post.line.color.rgb = BLUE_LINE
    post.line.width = Pt(2)
    beam = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(5.0), Inches(y4 - 0.2), Inches(5.6), Inches(y4 - 0.2))
    beam.line.color.rgb = BLUE_LINE
    beam.line.width = Pt(2)
    
    p1_v = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(5.05), Inches(y4 - 0.2), Inches(5.05), Inches(y4))
    p1_v.line.color.rgb = BLUE_LINE
    p1_h = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(4.95), Inches(y4), Inches(5.15), Inches(y4))
    p1_h.line.color.rgb = BLUE_LINE
    p1_h.line.width = Pt(1.5)
    
    p2_v = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(5.55), Inches(y4 - 0.2), Inches(5.55), Inches(y4))
    p2_v.line.color.rgb = BLUE_LINE
    p2_h = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(5.45), Inches(y4), Inches(5.65), Inches(y4))
    p2_h.line.color.rgb = BLUE_LINE
    p2_h.line.width = Pt(1.5)

    # 4. Page Number
    page_num = slide.shapes.add_textbox(Inches(12.5), Inches(6.8), Inches(0.5), Inches(0.5))
    tf_num = page_num.text_frame
    p_num = tf_num.paragraphs[0]
    p_num.text = "02"
    p_num.font.name = "Arial"
    p_num.font.size = Pt(14)
    p_num.font.color.rgb = RGBColor(0x66, 0x66, 0x66)



# ── Main ──
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

s0 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_1(s0)
prs.save(OUTPUT_PATH)
