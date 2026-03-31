from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# ── Color Palette ──
BLUE_HEADER = RGBColor(0x5B, 0x9B, 0xD5)
BLUE_DARK   = RGBColor(0x4A, 0x86, 0xC8)
CYAN        = RGBColor(0x00, 0xBC, 0xD4)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x33, 0x33, 0x33)
GRAY_TEXT   = RGBColor(0x55, 0x55, 0x55)
GRAY_BAR    = RGBColor(0xB0, 0xBE, 0xC5)
RED         = RGBColor(0xE5, 0x39, 0x35)
GREEN       = RGBColor(0x43, 0xA0, 0x47)
ORANGE      = RGBColor(0xFF, 0x98, 0x00)

ICON_BG  = RGBColor(0xE3, 0xE8, 0xED)
ICON_FG  = RGBColor(0x54, 0x6E, 0x7A)
FONT_NAME = "Microsoft YaHei"
SLIDE_WIDTH = Inches(13.333)
HEADER_H    = Inches(0.75)
SUBTITLE_Y  = Inches(0.95)


def add_header_banner(slide, title_text, bg_color=None):
    if bg_color is None:
        bg_color = BLUE_HEADER
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, HEADER_H
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = bg_color
    banner.line.fill.background()
    tf = banner.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(26)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME


def add_subtitle(slide, text, left, top, width=Inches(12), font_size=Pt(18)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.font.name = FONT_NAME
    return txBox


def add_icon_box(slide, left, top, symbol, size=Inches(0.48)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ICON_BG
    shape.line.fill.background()
    shape.adjustments[0] = 0.25
    tf = shape.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = symbol
    p.font.size = Pt(18)
    p.font.color.rgb = ICON_FG
    p.font.bold = False
    return shape


def add_bullet_item(slide, left, top, symbol, label, description,
                    width=Inches(5.5), desc_size=Pt(13)):
    add_icon_box(slide, left, top, symbol)
    text_left = left + Inches(0.65)
    txBox = slide.shapes.add_textbox(text_left, top - Inches(0.02), width, Inches(0.65))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run_label = p.add_run()
    run_label.text = label + "\uff1a"
    run_label.font.size = Pt(14)
    run_label.font.color.rgb = BLACK
    run_label.font.bold = True
    run_label.font.name = FONT_NAME
    run_desc = p.add_run()
    run_desc.text = description
    run_desc.font.size = desc_size
    run_desc.font.color.rgb = GRAY_TEXT
    run_desc.font.bold = False
    run_desc.font.name = FONT_NAME
    return txBox


def add_conclusion_box(slide, left, top, width, text, font_size=Pt(13)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = BLACK
    run.font.bold = True
    run.font.name = FONT_NAME
    return txBox


def add_table(slide, left, top, width, height, rows, cols, data,
              header_color=None, col_widths=None):
    if header_color is None:
        header_color = BLUE_HEADER
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if data[r][c] is not None else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(11)
                    run.font.name = FONT_NAME
                    if r == 0:
                        run.font.color.rgb = WHITE
                        run.font.bold = True
                    else:
                        run.font.color.rgb = BLACK
                        run.font.bold = False
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else RGBColor(0xF5, 0xF5, 0xF5)
            cell.margin_left = Pt(5)
            cell.margin_right = Pt(5)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
    return table_shape


def add_bar_chart(slide, left, top, width, height,
                  categories, values, title="", bar_colors=None):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', values)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    chart.chart_style = 2
    plot = chart.plots[0]
    plot.gap_width = 100
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CYAN
    series.has_data_labels = True
    dl = series.data_labels
    dl.font.size = Pt(13)
    dl.font.bold = True
    dl.font.color.rgb = BLACK
    dl.number_format = '0.#'
    dl.show_value = True
    dl.label_position = XL_LABEL_POSITION.OUTSIDE_END
    if bar_colors:
        for i, color in enumerate(bar_colors):
            pt = series.points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = color
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(12)
    cat_axis.tick_labels.font.name = FONT_NAME
    cat_axis.major_tick_mark = 2
    cat_axis.format.line.fill.background()
    val_axis = chart.value_axis
    val_axis.visible = False
    val_axis.major_tick_mark = 2
    val_axis.format.line.fill.background()
    val_axis.major_gridlines.format.line.fill.background()
    if title:
        chart.has_title = True
        ct = chart.chart_title.text_frame.paragraphs[0]
        ct.text = title
        ct.font.size = Pt(14)
        ct.font.bold = True
        ct.font.name = FONT_NAME
    else:
        chart.has_title = False
    return chart_frame


def add_callout_label(slide, left, top, text, bg_color=None, font_size=Pt(11)):
    if bg_color is None:
        bg_color = CYAN
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.3), Inches(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME
    return shape


def add_data_card(slide, left, top, width, height, value, label,
                  value_color=None, bg_color=None):
    if value_color is None:
        value_color = CYAN
    if bg_color is None:
        bg_color = WHITE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(3)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run1 = p1.add_run()
    run1.text = str(value)
    run1.font.size = Pt(24)
    run1.font.color.rgb = value_color
    run1.font.bold = True
    run1.font.name = FONT_NAME
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(10)
    run2.font.color.rgb = GRAY_TEXT
    run2.font.bold = False
    run2.font.name = FONT_NAME
    return shape


OUTPUT_PATH = r"C:\Users\Administrator\Desktop\notebook\my_notebooklm_with_ppteditable\projects\如何成为PPT高手\最终文档\ppt_slides\slide_06.pptx"


# ── Slide 1 ──

def build_slide_1(slide):
    # Colors
    BLUE_PRIMARY = RGBColor(0x00, 0x52, 0xCC)
    TEXT_BLACK = RGBColor(0x33, 0x33, 0x33)
    TEXT_GRAY = RGBColor(0x7F, 0x7F, 0x7F)
    LIGHT_GRAY = RGBColor(0xD9, 0xD9, 0xD9)
    CUBE_FILL = RGBColor(0xF4, 0xF6, 0xF9)
    CUBE_LINE = RGBColor(0x2F, 0x45, 0x6A)
    HIGHLIGHT_FILL = RGBColor(0xDE, 0xEA, 0xF6)
    HIGHLIGHT_LINE = RGBColor(0x5B, 0x9B, 0xD5)

    # 1. Top Left Page Indicator
    # Small grey dash
    dash = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0.4), Inches(0.15), Inches(0.04))
    dash.fill.solid()
    dash.fill.fore_color.rgb = LIGHT_GRAY
    dash.line.fill.background()

    # "Page 6"
    tx_page = slide.shapes.add_textbox(Inches(0.35), Inches(0.45), Inches(1), Inches(0.3))
    tf_page = tx_page.text_frame
    tf_page.word_wrap = False
    p_page = tf_page.paragraphs[0]
    p_page.text = "Page 6"
    p_page.font.name = "Microsoft YaHei"
    p_page.font.size = Pt(12)
    p_page.font.bold = True
    p_page.font.color.rgb = TEXT_BLACK

    # "6/11"
    tx_num = slide.shapes.add_textbox(Inches(0.35), Inches(0.7), Inches(1), Inches(0.3))
    tf_num = tx_num.text_frame
    p_num = tf_num.paragraphs[0]
    p_num.text = "6/11"
    p_num.font.name = "Microsoft YaHei"
    p_num.font.size = Pt(10)
    p_num.font.color.rgb = TEXT_GRAY

    # 2. Main Title
    tx_title = slide.shapes.add_textbox(Inches(2), Inches(0.8), Inches(9.333), Inches(0.8))
    tf_title = tx_title.text_frame
    p_title = tf_title.paragraphs[0]
    p_title.alignment = PP_ALIGN.CENTER
    
    run1 = p_title.add_run()
    run1.text = "法则二："
    run1.font.name = "Microsoft YaHei"
    run1.font.size = Pt(36)
    run1.font.bold = True
    run1.font.color.rgb = TEXT_BLACK

    run2 = p_title.add_run()
    run2.text = "设计统一建立专业信任"
    run2.font.name = "Microsoft YaHei"
    run2.font.size = Pt(36)
    run2.font.bold = True
    run2.font.color.rgb = BLUE_PRIMARY

    # 3. Subtitle
    tx_sub = slide.shapes.add_textbox(Inches(2), Inches(1.7), Inches(9.333), Inches(0.5))
    tf_sub = tx_sub.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.alignment = PP_ALIGN.CENTER
    p_sub.text = "从视觉一致性中体现专业度"
    p_sub.font.name = "Microsoft YaHei"
    p_sub.font.size = Pt(20)
    p_sub.font.color.rgb = TEXT_BLACK

    # 4. Central Graphics (Cubes)
    cube_y = Inches(2.8)
    cube_w = Inches(2.2)
    cube_h = Inches(2.2)

    # Cube 1 (Left)
    cube1 = slide.shapes.add_shape(MSO_SHAPE.CUBE, Inches(2.8), cube_y, cube_w, cube_h)
    cube1.fill.solid()
    cube1.fill.fore_color.rgb = CUBE_FILL
    cube1.line.color.rgb = CUBE_LINE
    cube1.line.width = Pt(2)

    # Cube 2 (Middle)
    cube2 = slide.shapes.add_shape(MSO_SHAPE.CUBE, Inches(5.5), cube_y, cube_w, cube_h)
    cube2.fill.solid()
    cube2.fill.fore_color.rgb = CUBE_FILL
    cube2.line.color.rgb = CUBE_LINE
    cube2.line.width = Pt(2)

    # Flying piece for Cube 2
    small_cube = slide.shapes.add_shape(MSO_SHAPE.CUBE, Inches(6.8), Inches(2.5), Inches(0.8), Inches(0.8))
    small_cube.fill.solid()
    small_cube.fill.fore_color.rgb = HIGHLIGHT_FILL
    small_cube.line.color.rgb = HIGHLIGHT_LINE
    small_cube.line.width = Pt(1.5)

    # Arrow for flying piece
    arrow_insert = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(7.1), Inches(3.4), Inches(0.2), Inches(0.3))
    arrow_insert.fill.solid()
    arrow_insert.fill.fore_color.rgb = HIGHLIGHT_LINE
    arrow_insert.line.fill.background()

    # Cube 3 (Right)
    cube3 = slide.shapes.add_shape(MSO_SHAPE.CUBE, Inches(8.2), cube_y, cube_w, cube_h)
    cube3.fill.solid()
    cube3.fill.fore_color.rgb = CUBE_FILL
    cube3.line.color.rgb = CUBE_LINE
    cube3.line.width = Pt(2)

    # Highlighted piece on Cube 3 (Simulated by a smaller cube on top right corner)
    hl_cube = slide.shapes.add_shape(MSO_SHAPE.CUBE, Inches(9.2), Inches(3.2), Inches(0.7), Inches(0.7))
    hl_cube.fill.solid()
    hl_cube.fill.fore_color.rgb = HIGHLIGHT_FILL
    hl_cube.line.color.rgb = HIGHLIGHT_LINE
    hl_cube.line.width = Pt(1.5)

    # 5. Bottom Arrow and Text
    # Long arrow line
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(3.8), Inches(5.4), Inches(9.4), Inches(5.4))
    connector.line.color.rgb = HIGHLIGHT_LINE
    connector.line.width = Pt(1.5)
    # Add arrow head (using standard line properties if possible, or draw a small triangle)
    arrow_head = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(9.3), Inches(5.33), Inches(0.15), Inches(0.15))
    arrow_head.rotation = 90
    arrow_head.fill.solid()
    arrow_head.fill.fore_color.rgb = HIGHLIGHT_LINE
    arrow_head.line.fill.background()

    # Text below arrow
    tx_arrow = slide.shapes.add_textbox(Inches(2), Inches(5.6), Inches(9.333), Inches(0.4))
    tf_arrow = tx_arrow.text_frame
    p_arrow = tf_arrow.paragraphs[0]
    p_arrow.alignment = PP_ALIGN.CENTER
    p_arrow.text = "统一感能降低观众的视觉疲劳"
    p_arrow.font.name = "Microsoft YaHei"
    p_arrow.font.size = Pt(14)
    p_arrow.font.bold = True
    p_arrow.font.color.rgb = TEXT_BLACK

    # 6. Bottom 3 Columns
    col_y = Inches(6.3)
    
    # --- Column 1 ---
    # Icon 1: Compass/Ruler
    compass = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(1.0), col_y + Inches(0.1), Inches(0.3), Inches(0.4))
    compass.rotation = -90
    compass.fill.background()
    compass.line.color.rgb = CUBE_LINE
    compass.line.width = Pt(1.5)
    
    ruler = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.4), col_y, Inches(0.15), Inches(0.5))
    ruler.fill.background()
    ruler.line.color.rgb = HIGHLIGHT_LINE
    ruler.line.width = Pt(1.5)

    # Text 1
    tx_col1_title = slide.shapes.add_textbox(Inches(1.7), col_y - Inches(0.1), Inches(3.0), Inches(0.3))
    p_col1_title = tx_col1_title.text_frame.paragraphs[0]
    r1 = p_col1_title.add_run()
    r1.text = "1. "
    r1.font.bold = True
    r1.font.size = Pt(13)
    r2 = p_col1_title.add_run()
    r2.text = "风格漂移"
    r2.font.bold = True
    r2.font.size = Pt(13)
    r2.font.color.rgb = BLUE_PRIMARY
    r3 = p_col1_title.add_run()
    r3.text = "是PPT的大忌"
    r3.font.bold = True
    r3.font.size = Pt(13)

    tx_col1_desc = slide.shapes.add_textbox(Inches(1.7), col_y + Inches(0.2), Inches(3.0), Inches(0.4))
    p_col1_desc = tx_col1_desc.text_frame.paragraphs[0]
    p_col1_desc.text = "避免混乱，保持整体风格的一致性。"
    p_col1_desc.font.size = Pt(11)
    p_col1_desc.font.color.rgb = TEXT_BLACK

    # --- Column 2 ---
    # Icon 2: Eye
    eye_outer = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.0), col_y + Inches(0.1), Inches(0.5), Inches(0.3))
    eye_outer.fill.background()
    eye_outer.line.color.rgb = CUBE_LINE
    eye_outer.line.width = Pt(1.5)
    
    eye_inner = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.15), col_y + Inches(0.15), Inches(0.2), Inches(0.2))
    eye_inner.fill.background()
    eye_inner.line.color.rgb = HIGHLIGHT_LINE
    eye_inner.line.width = Pt(1.5)
    
    pulse = slide.shapes.add_shape(MSO_SHAPE.ZIG_ZAG, Inches(5.0), col_y + Inches(0.45), Inches(0.5), Inches(0.1))
    pulse.fill.background()
    pulse.line.color.rgb = HIGHLIGHT_LINE
    pulse.line.width = Pt(1.5)

    # Text 2
    tx_col2_title = slide.shapes.add_textbox(Inches(5.7), col_y - Inches(0.1), Inches(3.2), Inches(0.3))
    p_col2_title = tx_col2_title.text_frame.paragraphs[0]
    r1 = p_col2_title.add_run()
    r1.text = "2. 统一感能降低观众的"
    r1.font.bold = True
    r1.font.size = Pt(13)
    r2 = p_col2_title.add_run()
    r2.text = "视觉疲劳"
    r2.font.bold = True
    r2.font.size = Pt(13)
    r2.font.color.rgb = BLUE_PRIMARY

    tx_col2_desc = slide.shapes.add_textbox(Inches(5.7), col_y + Inches(0.2), Inches(3.2), Inches(0.4))
    p_col2_desc = tx_col2_desc.text_frame.paragraphs[0]
    p_col2_desc.text = "视觉流畅，让观众更专注于内容。"
    p_col2_desc.font.size = Pt(11)
    p_col2_desc.font.color.rgb = TEXT_BLACK

    # --- Column 3 ---
    # Icon 3: Browser/Window
    browser = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.0), col_y, Inches(0.5), Inches(0.4))
    browser.fill.background()
    browser.line.color.rgb = CUBE_LINE
    browser.line.width = Pt(1.5)
    
    browser_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.0), col_y + Inches(0.1), Inches(0.5), Inches(0.02))
    browser_line.fill.solid()
    browser_line.fill.fore_color.rgb = CUBE_LINE
    browser_line.line.fill.background()
    
    mag_glass = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.3), col_y + Inches(0.2), Inches(0.2), Inches(0.2))
    mag_glass.fill.background()
    mag_glass.line.color.rgb = HIGHLIGHT_LINE
    mag_glass.line.width = Pt(1.5)

    # Text 3
    tx_col3_title = slide.shapes.add_textbox(Inches(9.7), col_y - Inches(0.1), Inches(3.5), Inches(0.3))
    p_col3_title = tx_col3_title.text_frame.paragraphs[0]
    r1 = p_col3_title.add_run()
    r1.text = "3. "
    r1.font.bold = True
    r1.font.size = Pt(13)
    r2 = p_col3_title.add_run()
    r2.text = "专业感"
    r2.font.bold = True
    r2.font.size = Pt(13)
    r2.font.color.rgb = BLUE_PRIMARY
    r3 = p_col3_title.add_run()
    r3.text = "源于对细节的严苛把控"
    r3.font.bold = True
    r3.font.size = Pt(13)

    tx_col3_desc = slide.shapes.add_textbox(Inches(9.7), col_y + Inches(0.2), Inches(3.5), Inches(0.4))
    p_col3_desc = tx_col3_desc.text_frame.paragraphs[0]
    p_col3_desc.text = "对齐、间距、字体、颜色的精准规范。"
    p_col3_desc.font.size = Pt(11)
    p_col3_desc.font.color.rgb = TEXT_BLACK



# ── Main ──
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

s0 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_1(s0)
prs.save(OUTPUT_PATH)
