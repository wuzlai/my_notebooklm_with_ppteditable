from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# ── Color Palette ──
BLUE_HEADER = RGBColor(0x5B, 0x9B, 0xD5)
BLUE_DARK   = RGBColor(0x4A, 0x86, 0xC8)
CYAN        = RGBColor(0x00, 0xBC, 0xD4)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x33, 0x33, 0x33)
GRAY_TEXT   = RGBColor(0x55, 0x55, 0x55)
GRAY_BAR    = RGBColor(0xB0, 0xBE, 0xC5)
RED         = RGBColor(0xE5, 0x39, 0x35)
GREEN       = RGBColor(0x43, 0xA0, 0x47)
ORANGE      = RGBColor(0xFF, 0x98, 0x00)

ICON_BG  = RGBColor(0xE3, 0xE8, 0xED)
ICON_FG  = RGBColor(0x54, 0x6E, 0x7A)
FONT_NAME = "Microsoft YaHei"
SLIDE_WIDTH = Inches(13.333)
HEADER_H    = Inches(0.75)
SUBTITLE_Y  = Inches(0.95)


def add_header_banner(slide, title_text, bg_color=None):
    if bg_color is None:
        bg_color = BLUE_HEADER
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, HEADER_H
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = bg_color
    banner.line.fill.background()
    tf = banner.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(26)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME


def add_subtitle(slide, text, left, top, width=Inches(12), font_size=Pt(18)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.font.name = FONT_NAME
    return txBox


def add_icon_box(slide, left, top, symbol, size=Inches(0.48)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ICON_BG
    shape.line.fill.background()
    shape.adjustments[0] = 0.25
    tf = shape.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = symbol
    p.font.size = Pt(18)
    p.font.color.rgb = ICON_FG
    p.font.bold = False
    return shape


def add_bullet_item(slide, left, top, symbol, label, description,
                    width=Inches(5.5), desc_size=Pt(13)):
    add_icon_box(slide, left, top, symbol)
    text_left = left + Inches(0.65)
    txBox = slide.shapes.add_textbox(text_left, top - Inches(0.02), width, Inches(0.65))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run_label = p.add_run()
    run_label.text = label + "\uff1a"
    run_label.font.size = Pt(14)
    run_label.font.color.rgb = BLACK
    run_label.font.bold = True
    run_label.font.name = FONT_NAME
    run_desc = p.add_run()
    run_desc.text = description
    run_desc.font.size = desc_size
    run_desc.font.color.rgb = GRAY_TEXT
    run_desc.font.bold = False
    run_desc.font.name = FONT_NAME
    return txBox


def add_conclusion_box(slide, left, top, width, text, font_size=Pt(13)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = BLACK
    run.font.bold = True
    run.font.name = FONT_NAME
    return txBox


def add_table(slide, left, top, width, height, rows, cols, data,
              header_color=None, col_widths=None):
    if header_color is None:
        header_color = BLUE_HEADER
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if data[r][c] is not None else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(11)
                    run.font.name = FONT_NAME
                    if r == 0:
                        run.font.color.rgb = WHITE
                        run.font.bold = True
                    else:
                        run.font.color.rgb = BLACK
                        run.font.bold = False
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else RGBColor(0xF5, 0xF5, 0xF5)
            cell.margin_left = Pt(5)
            cell.margin_right = Pt(5)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
    return table_shape


def add_bar_chart(slide, left, top, width, height,
                  categories, values, title="", bar_colors=None):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', values)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    chart.chart_style = 2
    plot = chart.plots[0]
    plot.gap_width = 100
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CYAN
    series.has_data_labels = True
    dl = series.data_labels
    dl.font.size = Pt(13)
    dl.font.bold = True
    dl.font.color.rgb = BLACK
    dl.number_format = '0.#'
    dl.show_value = True
    dl.label_position = XL_LABEL_POSITION.OUTSIDE_END
    if bar_colors:
        for i, color in enumerate(bar_colors):
            pt = series.points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = color
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(12)
    cat_axis.tick_labels.font.name = FONT_NAME
    cat_axis.major_tick_mark = 2
    cat_axis.format.line.fill.background()
    val_axis = chart.value_axis
    val_axis.visible = False
    val_axis.major_tick_mark = 2
    val_axis.format.line.fill.background()
    val_axis.major_gridlines.format.line.fill.background()
    if title:
        chart.has_title = True
        ct = chart.chart_title.text_frame.paragraphs[0]
        ct.text = title
        ct.font.size = Pt(14)
        ct.font.bold = True
        ct.font.name = FONT_NAME
    else:
        chart.has_title = False
    return chart_frame


def add_callout_label(slide, left, top, text, bg_color=None, font_size=Pt(11)):
    if bg_color is None:
        bg_color = CYAN
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.3), Inches(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME
    return shape


def add_data_card(slide, left, top, width, height, value, label,
                  value_color=None, bg_color=None):
    if value_color is None:
        value_color = CYAN
    if bg_color is None:
        bg_color = WHITE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(3)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run1 = p1.add_run()
    run1.text = str(value)
    run1.font.size = Pt(24)
    run1.font.color.rgb = value_color
    run1.font.bold = True
    run1.font.name = FONT_NAME
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(10)
    run2.font.color.rgb = GRAY_TEXT
    run2.font.bold = False
    run2.font.name = FONT_NAME
    return shape


OUTPUT_PATH = r"C:\Users\Administrator\Desktop\notebook\my_notebooklm_with_ppteditable\projects\如何成为PPT高手\最终文档\ppt_slides\slide_01.pptx"


# ── Slide 1 ──

def build_slide_1(slide):
    # 1. Background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RGBColor(0x05, 0x22, 0x55) # Dark blue
    bg_shape.line.fill.background()

    # Grid lines
    grid_color = RGBColor(0x10, 0x35, 0x70)
    for i in range(1, 14):
        line = slide.shapes.add_connector(1, Inches(i), Inches(0), Inches(i), Inches(7.5))
        line.line.color.rgb = grid_color
        line.line.width = Pt(0.5)
    for i in range(1, 8):
        line = slide.shapes.add_connector(1, Inches(0), Inches(i), Inches(13.333), Inches(i))
        line.line.color.rgb = grid_color
        line.line.width = Pt(0.5)

    # 2. Central Lightbulb Icon
    center_x = 13.333 / 2
    bulb_y = 1.0
    bulb_color = RGBColor(0xAA, 0xD4, 0xFF)
    
    # Glow effect (simulated with a larger faint circle)
    glow = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(center_x - 1.0), Inches(bulb_y - 0.4), Inches(2.0), Inches(2.0))
    glow.fill.solid()
    glow.fill.fore_color.rgb = RGBColor(0x15, 0x45, 0x85)
    glow.line.fill.background()

    # Main bulb (Oval)
    bulb = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(center_x - 0.6), Inches(bulb_y), Inches(1.2), Inches(1.2))
    bulb.fill.background()
    bulb.line.color.rgb = bulb_color
    bulb.line.width = Pt(3)
    
    # Bulb base (Rectangle)
    base = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(center_x - 0.3), Inches(bulb_y + 1.15), Inches(0.6), Inches(0.4))
    base.fill.background()
    base.line.color.rgb = bulb_color
    base.line.width = Pt(3)
    
    # Base screw lines
    for i in range(3):
        line_y = bulb_y + 1.6 + i * 0.12
        line = slide.shapes.add_connector(1, Inches(center_x - 0.25), Inches(line_y), Inches(center_x + 0.25), Inches(line_y))
        line.line.color.rgb = bulb_color
        line.line.width = Pt(2.5)
        
    # Filament (Inner lines)
    fil_left = slide.shapes.add_connector(1, Inches(center_x - 0.2), Inches(bulb_y + 1.15), Inches(center_x - 0.2), Inches(bulb_y + 0.6))
    fil_left.line.color.rgb = bulb_color
    fil_left.line.width = Pt(2)
    
    fil_right = slide.shapes.add_connector(1, Inches(center_x + 0.2), Inches(bulb_y + 1.15), Inches(center_x + 0.2), Inches(bulb_y + 0.6))
    fil_right.line.color.rgb = bulb_color
    fil_right.line.width = Pt(2)
    
    fil_top = slide.shapes.add_shape(MSO_SHAPE.ARC, Inches(center_x - 0.2), Inches(bulb_y + 0.4), Inches(0.4), Inches(0.4))
    fil_top.rotation = 180
    fil_top.fill.background()
    fil_top.line.color.rgb = bulb_color
    fil_top.line.width = Pt(2)

    # 3. Main Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(11.333), Inches(1))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "简而不凡：高效演示文稿的制作之道"
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # 4. Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(2), Inches(4.2), Inches(9.333), Inches(0.8))
    tf_sub = subtitle_box.text_frame
    tf_sub.clear()
    p_sub = tf_sub.paragraphs[0]
    p_sub.alignment = PP_ALIGN.CENTER
    run_sub = p_sub.add_run()
    run_sub.text = "掌握专业PPT的核心逻辑与设计法则"
    run_sub.font.name = "Microsoft YaHei"
    run_sub.font.size = Pt(24)
    run_sub.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # 5. Bullet Points
    bullet_texts = [
        "1. 演示文稿不仅仅是工具，更是思维的视觉化表达",
        "2. 核心目标：降低沟通成本，提升说服力",
        "3. 专家级PPT的两大底层支柱：内容清晰与设计统一"
    ]
    
    start_y = 5.2
    spacing = 0.6
    icon_x = 3.8
    text_x = 4.3
    icon_color = RGBColor(0xAA, 0xD4, 0xFF)
    
    for i, text in enumerate(bullet_texts):
        icon_y = start_y + i * spacing + 0.08
        
        # Draw Icons
        if i == 0: # Idea (Lightbulb)
            icon_bulb = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(icon_x), Inches(icon_y), Inches(0.2), Inches(0.2))
            icon_bulb.fill.background()
            icon_bulb.line.color.rgb = icon_color
            icon_bulb.line.width = Pt(1.5)
            icon_base = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(icon_x + 0.05), Inches(icon_y + 0.18), Inches(0.1), Inches(0.08))
            icon_base.fill.background()
            icon_base.line.color.rgb = icon_color
            icon_base.line.width = Pt(1.5)
            # Rays
            slide.shapes.add_connector(1, Inches(icon_x+0.1), Inches(icon_y-0.05), Inches(icon_x+0.1), Inches(icon_y-0.1)).line.color.rgb = icon_color
            slide.shapes.add_connector(1, Inches(icon_x-0.05), Inches(icon_y+0.1), Inches(icon_x-0.1), Inches(icon_y+0.1)).line.color.rgb = icon_color
            slide.shapes.add_connector(1, Inches(icon_x+0.25), Inches(icon_y+0.1), Inches(icon_x+0.3), Inches(icon_y+0.1)).line.color.rgb = icon_color

        elif i == 1: # Target
            icon1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(icon_x), Inches(icon_y), Inches(0.25), Inches(0.25))
            icon1.fill.background()
            icon1.line.color.rgb = icon_color
            icon1.line.width = Pt(1.5)
            icon2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(icon_x+0.075), Inches(icon_y+0.075), Inches(0.1), Inches(0.1))
            icon2.fill.background()
            icon2.line.color.rgb = icon_color
            icon2.line.width = Pt(1.5)
            # Arrow
            arrow = slide.shapes.add_connector(1, Inches(icon_x+0.15), Inches(icon_y+0.1), Inches(icon_x+0.35), Inches(icon_y-0.1))
            arrow.line.color.rgb = icon_color
            arrow.line.width = Pt(1.5)

        elif i == 2: # Balance
            # Base triangle
            tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(icon_x+0.05), Inches(icon_y+0.1), Inches(0.15), Inches(0.15))
            tri.fill.background()
            tri.line.color.rgb = icon_color
            tri.line.width = Pt(1.5)
            # Top bar
            bar = slide.shapes.add_connector(1, Inches(icon_x-0.05), Inches(icon_y+0.1), Inches(icon_x+0.3), Inches(icon_y+0.1))
            bar.line.color.rgb = icon_color
            bar.line.width = Pt(1.5)
            # Left pan
            slide.shapes.add_connector(1, Inches(icon_x-0.05), Inches(icon_y+0.1), Inches(icon_x-0.05), Inches(icon_y+0.25)).line.color.rgb = icon_color
            slide.shapes.add_connector(1, Inches(icon_x-0.1), Inches(icon_y+0.25), Inches(icon_x), Inches(icon_y+0.25)).line.color.rgb = icon_color
            # Right pan
            slide.shapes.add_connector(1, Inches(icon_x+0.3), Inches(icon_y+0.1), Inches(icon_x+0.3), Inches(icon_y+0.25)).line.color.rgb = icon_color
            slide.shapes.add_connector(1, Inches(icon_x+0.25), Inches(icon_y+0.25), Inches(icon_x+0.35), Inches(icon_y+0.25)).line.color.rgb = icon_color

        # Text
        text_box = slide.shapes.add_textbox(Inches(text_x), Inches(start_y + i * spacing), Inches(8), Inches(0.4))
        tf_bullet = text_box.text_frame
        tf_bullet.clear()
        p_bullet = tf_bullet.paragraphs[0]
        run_bullet = p_bullet.add_run()
        run_bullet.text = text
        run_bullet.font.name = "Microsoft YaHei"
        run_bullet.font.size = Pt(16)
        run_bullet.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)



# ── Main ──
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

s0 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_1(s0)
prs.save(OUTPUT_PATH)
