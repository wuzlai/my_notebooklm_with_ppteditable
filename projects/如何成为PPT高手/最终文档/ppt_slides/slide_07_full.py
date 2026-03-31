from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# ── Color Palette ──
BLUE_HEADER = RGBColor(0x5B, 0x9B, 0xD5)
BLUE_DARK   = RGBColor(0x4A, 0x86, 0xC8)
CYAN        = RGBColor(0x00, 0xBC, 0xD4)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x33, 0x33, 0x33)
GRAY_TEXT   = RGBColor(0x55, 0x55, 0x55)
GRAY_BAR    = RGBColor(0xB0, 0xBE, 0xC5)
RED         = RGBColor(0xE5, 0x39, 0x35)
GREEN       = RGBColor(0x43, 0xA0, 0x47)
ORANGE      = RGBColor(0xFF, 0x98, 0x00)

ICON_BG  = RGBColor(0xE3, 0xE8, 0xED)
ICON_FG  = RGBColor(0x54, 0x6E, 0x7A)
FONT_NAME = "Microsoft YaHei"
SLIDE_WIDTH = Inches(13.333)
HEADER_H    = Inches(0.75)
SUBTITLE_Y  = Inches(0.95)


def add_header_banner(slide, title_text, bg_color=None):
    if bg_color is None:
        bg_color = BLUE_HEADER
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, HEADER_H
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = bg_color
    banner.line.fill.background()
    tf = banner.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(26)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME


def add_subtitle(slide, text, left, top, width=Inches(12), font_size=Pt(18)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.font.name = FONT_NAME
    return txBox


def add_icon_box(slide, left, top, symbol, size=Inches(0.48)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ICON_BG
    shape.line.fill.background()
    shape.adjustments[0] = 0.25
    tf = shape.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = symbol
    p.font.size = Pt(18)
    p.font.color.rgb = ICON_FG
    p.font.bold = False
    return shape


def add_bullet_item(slide, left, top, symbol, label, description,
                    width=Inches(5.5), desc_size=Pt(13)):
    add_icon_box(slide, left, top, symbol)
    text_left = left + Inches(0.65)
    txBox = slide.shapes.add_textbox(text_left, top - Inches(0.02), width, Inches(0.65))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run_label = p.add_run()
    run_label.text = label + "\uff1a"
    run_label.font.size = Pt(14)
    run_label.font.color.rgb = BLACK
    run_label.font.bold = True
    run_label.font.name = FONT_NAME
    run_desc = p.add_run()
    run_desc.text = description
    run_desc.font.size = desc_size
    run_desc.font.color.rgb = GRAY_TEXT
    run_desc.font.bold = False
    run_desc.font.name = FONT_NAME
    return txBox


def add_conclusion_box(slide, left, top, width, text, font_size=Pt(13)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = BLACK
    run.font.bold = True
    run.font.name = FONT_NAME
    return txBox


def add_table(slide, left, top, width, height, rows, cols, data,
              header_color=None, col_widths=None):
    if header_color is None:
        header_color = BLUE_HEADER
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if data[r][c] is not None else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(11)
                    run.font.name = FONT_NAME
                    if r == 0:
                        run.font.color.rgb = WHITE
                        run.font.bold = True
                    else:
                        run.font.color.rgb = BLACK
                        run.font.bold = False
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else RGBColor(0xF5, 0xF5, 0xF5)
            cell.margin_left = Pt(5)
            cell.margin_right = Pt(5)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
    return table_shape


def add_bar_chart(slide, left, top, width, height,
                  categories, values, title="", bar_colors=None):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', values)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    chart.chart_style = 2
    plot = chart.plots[0]
    plot.gap_width = 100
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CYAN
    series.has_data_labels = True
    dl = series.data_labels
    dl.font.size = Pt(13)
    dl.font.bold = True
    dl.font.color.rgb = BLACK
    dl.number_format = '0.#'
    dl.show_value = True
    dl.label_position = XL_LABEL_POSITION.OUTSIDE_END
    if bar_colors:
        for i, color in enumerate(bar_colors):
            pt = series.points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = color
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(12)
    cat_axis.tick_labels.font.name = FONT_NAME
    cat_axis.major_tick_mark = 2
    cat_axis.format.line.fill.background()
    val_axis = chart.value_axis
    val_axis.visible = False
    val_axis.major_tick_mark = 2
    val_axis.format.line.fill.background()
    val_axis.major_gridlines.format.line.fill.background()
    if title:
        chart.has_title = True
        ct = chart.chart_title.text_frame.paragraphs[0]
        ct.text = title
        ct.font.size = Pt(14)
        ct.font.bold = True
        ct.font.name = FONT_NAME
    else:
        chart.has_title = False
    return chart_frame


def add_callout_label(slide, left, top, text, bg_color=None, font_size=Pt(11)):
    if bg_color is None:
        bg_color = CYAN
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.3), Inches(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME
    return shape


def add_data_card(slide, left, top, width, height, value, label,
                  value_color=None, bg_color=None):
    if value_color is None:
        value_color = CYAN
    if bg_color is None:
        bg_color = WHITE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(3)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run1 = p1.add_run()
    run1.text = str(value)
    run1.font.size = Pt(24)
    run1.font.color.rgb = value_color
    run1.font.bold = True
    run1.font.name = FONT_NAME
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(10)
    run2.font.color.rgb = GRAY_TEXT
    run2.font.bold = False
    run2.font.name = FONT_NAME
    return shape


OUTPUT_PATH = r"C:\Users\Administrator\Desktop\notebook\my_notebooklm_with_ppteditable\projects\如何成为PPT高手\最终文档\ppt_slides\slide_07.pptx"


# ── Slide 1 ──

def build_slide_1(slide):
    # Define Colors
    BLUE_PRIMARY = RGBColor(0x00, 0x52, 0xD9)
    BLUE_DARK = RGBColor(0x00, 0x2B, 0x75)
    GRAY_DARK = RGBColor(0x2B, 0x2F, 0x36)
    ORANGE_ACCENT = RGBColor(0xFF, 0x95, 0x00)
    TEXT_MAIN = RGBColor(0x33, 0x33, 0x33)
    TEXT_SUB = RGBColor(0x66, 0x66, 0x66)
    LINE_COLOR = RGBColor(0xE0, 0xE0, 0xE0)
    BG_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    # 1. Header Area
    # Top thin blue line
    top_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.06))
    top_line.fill.solid()
    top_line.fill.fore_color.rgb = BLUE_PRIMARY
    top_line.line.fill.background()

    # Header Left Text
    tb_header_left = slide.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(2), Inches(0.4))
    p = tb_header_left.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "简而不凡"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = BLUE_PRIMARY

    # Header Right Text
    tb_header_right = slide.shapes.add_textbox(Inches(11.5), Inches(0.15), Inches(1.5), Inches(0.4))
    p = tb_header_right.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "PAGE 7 / 11"
    run.font.size = Pt(12)
    run.font.color.rgb = TEXT_SUB

    # Header Bottom Separator
    header_sep = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.6), Inches(0.6), Inches(12.733), Inches(0.6))
    header_sep.line.color.rgb = LINE_COLOR

    # 2. Main Title & Subtitle
    tb_title = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(8), Inches(0.8))
    p = tb_title.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "视觉规范：配色与字体的秩序"
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = BLUE_PRIMARY

    tb_subtitle = slide.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(8), Inches(0.5))
    p = tb_subtitle.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "建立一套专属的视觉系统"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = TEXT_MAIN

    # 3. Left Bullet Points
    # Item 1
    add_icon_box(slide, Inches(0.6), Inches(3.0), "🎨", size=Inches(0.8))
    tb_item1 = slide.shapes.add_textbox(Inches(1.6), Inches(2.9), Inches(6.0), Inches(1.0))
    tf1 = tb_item1.text_frame
    tf1.word_wrap = True
    p1_1 = tf1.paragraphs[0]
    r1_1 = p1_1.add_run()
    r1_1.text = "1. 全文配色"
    r1_1.font.size = Pt(16)
    r1_1.font.bold = True
    r1_1.font.color.rgb = TEXT_MAIN
    r1_2 = p1_1.add_run()
    r1_2.text = "不超过3种，主次分明"
    r1_2.font.size = Pt(16)
    r1_2.font.bold = True
    r1_2.font.color.rgb = ORANGE_ACCENT
    
    p1_2 = tf1.add_paragraph()
    p1_2.space_before = Pt(6)
    r1_3 = p1_2.add_run()
    r1_3.text = "限制色彩数量，确保页面干净统一，提升专业度。"
    r1_3.font.size = Pt(14)
    r1_3.font.color.rgb = TEXT_SUB

    # Item 2
    add_icon_box(slide, Inches(0.6), Inches(4.4), "Aa", size=Inches(0.8))
    tb_item2 = slide.shapes.add_textbox(Inches(1.6), Inches(4.3), Inches(6.0), Inches(1.0))
    tf2 = tb_item2.text_frame
    tf2.word_wrap = True
    p2_1 = tf2.paragraphs[0]
    r2_1 = p2_1.add_run()
    r2_1.text = "2. 字体选择需统一，建议"
    r2_1.font.size = Pt(16)
    r2_1.font.bold = True
    r2_1.font.color.rgb = TEXT_MAIN
    r2_2 = p2_1.add_run()
    r2_2.text = "不超过2种"
    r2_2.font.size = Pt(16)
    r2_2.font.bold = True
    r2_2.font.color.rgb = ORANGE_ACCENT
    
    p2_2 = tf2.add_paragraph()
    p2_2.space_before = Pt(6)
    r2_3 = p2_2.add_run()
    r2_3.text = "选择易读的无衬线字体（如苹方-简），保持风格一致。"
    r2_3.font.size = Pt(14)
    r2_3.font.color.rgb = TEXT_SUB

    # Item 3
    add_icon_box(slide, Inches(0.6), Inches(5.8), "🖍️", size=Inches(0.8))
    tb_item3 = slide.shapes.add_textbox(Inches(1.6), Inches(5.7), Inches(6.0), Inches(1.0))
    tf3 = tb_item3.text_frame
    tf3.word_wrap = True
    p3_1 = tf3.paragraphs[0]
    r3_1 = p3_1.add_run()
    r3_1.text = "3. 关键信息"
    r3_1.font.size = Pt(16)
    r3_1.font.bold = True
    r3_1.font.color.rgb = TEXT_MAIN
    r3_2 = p3_1.add_run()
    r3_2.text = "加粗或变色"
    r3_2.font.size = Pt(16)
    r3_2.font.bold = True
    r3_2.font.color.rgb = ORANGE_ACCENT
    r3_3 = p3_1.add_run()
    r3_3.text = "，而非随意更改字体"
    r3_3.font.size = Pt(16)
    r3_3.font.bold = True
    r3_3.font.color.rgb = TEXT_MAIN
    
    p3_2 = tf3.add_paragraph()
    p3_2.space_before = Pt(6)
    r3_4 = p3_2.add_run()
    r3_4.text = "通过字重和色彩"
    r3_4.font.size = Pt(14)
    r3_4.font.color.rgb = TEXT_SUB
    r3_5 = p3_2.add_run()
    r3_5.text = "强调重点"
    r3_5.font.size = Pt(14)
    r3_5.font.color.rgb = ORANGE_ACCENT
    r3_6 = p3_2.add_run()
    r3_6.text = "，避免视觉混乱。"
    r3_6.font.size = Pt(14)
    r3_6.font.color.rgb = TEXT_SUB

    # 4. Right Panel (White Card)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.0), Inches(1.2), Inches(4.6), Inches(5.8))
    card.fill.solid()
    card.fill.fore_color.rgb = BG_WHITE
    card.line.color.rgb = LINE_COLOR
    card.line.width = Pt(1)

    # Card Title 1
    tb_card_t1 = slide.shapes.add_textbox(Inches(8.2), Inches(1.4), Inches(4.0), Inches(0.4))
    p = tb_card_t1.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "配色建议"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = TEXT_MAIN

    # Color Swatches
    swatch_data = [
        (Inches(8.25), BLUE_PRIMARY, "主色", "#0052D9"),
        (Inches(9.65), GRAY_DARK, "辅助色", "#2B2F36"),
        (Inches(11.05), ORANGE_ACCENT, "强调色", "#FF9500")
    ]
    
    for left, color, label, hex_code in swatch_data:
        swatch = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.0), Inches(1.2), Inches(0.7))
        swatch.fill.solid()
        swatch.fill.fore_color.rgb = color
        swatch.line.fill.background()
        
        tb_swatch = slide.shapes.add_textbox(left, Inches(2.8), Inches(1.2), Inches(0.5))
        tf = tb_swatch.text_frame
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = label + "\n"
        r1.font.size = Pt(12)
        r1.font.color.rgb = TEXT_MAIN
        r2 = p1.add_run()
        r2.text = hex_code
        r2.font.size = Pt(10)
        r2.font.color.rgb = TEXT_SUB

    # Gradient Bar (Represented as solid primary blue)
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.25), Inches(3.6), Inches(4.0), Inches(0.25))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE_PRIMARY
    bar.line.fill.background()

    tb_bar_left = slide.shapes.add_textbox(Inches(8.15), Inches(3.9), Inches(1.5), Inches(0.3))
    p = tb_bar_left.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "#002B75"
    run.font.size = Pt(10)
    run.font.color.rgb = TEXT_SUB

    tb_bar_right = slide.shapes.add_textbox(Inches(11.15), Inches(3.9), Inches(1.2), Inches(0.3))
    p = tb_bar_right.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "#0052D9"
    run.font.size = Pt(10)
    run.font.color.rgb = TEXT_SUB

    # Card Separator Line
    card_sep = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(8.25), Inches(4.3), Inches(12.25), Inches(4.3))
    card_sep.line.color.rgb = LINE_COLOR

    # Card Title 2
    tb_card_t2 = slide.shapes.add_textbox(Inches(8.2), Inches(4.5), Inches(4.0), Inches(0.4))
    p = tb_card_t2.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "字体样式组合"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = TEXT_MAIN

    # Typography Examples
    tb_ex1 = slide.shapes.add_textbox(Inches(8.2), Inches(5.0), Inches(4.2), Inches(0.5))
    p = tb_ex1.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "一级标题示例 32-40pt"
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = BLUE_PRIMARY

    tb_ex2 = slide.shapes.add_textbox(Inches(8.2), Inches(5.6), Inches(4.2), Inches(0.4))
    p = tb_ex2.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "副标题示例 20-24pt"
    run.font.size = Pt(18)
    run.font.color.rgb = TEXT_MAIN

    tb_ex3 = slide.shapes.add_textbox(Inches(8.2), Inches(6.1), Inches(4.2), Inches(0.4))
    p = tb_ex3.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "正文内容示例，高易读性 16-18pt"
    run.font.size = Pt(14)
    run.font.color.rgb = TEXT_SUB



# ── Main ──
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

s0 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_1(s0)
prs.save(OUTPUT_PATH)
