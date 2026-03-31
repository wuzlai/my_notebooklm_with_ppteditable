def build_slide(slide):
    from pptx.enum.shapes import MSO_CONNECTOR
    
    # Colors
    BLUE_TITLE = RGBColor(0x1A, 0x56, 0xBA)
    DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
    BLUE_LINE = RGBColor(0x3B, 0x7E, 0xC6)
    ORANGE_TEXT = RGBColor(0xDE, 0x9B, 0x35)
    BLACK_TEXT = RGBColor(0x00, 0x00, 0x00)
    SHADOW_COLOR = RGBColor(0xF0, 0xF4, 0xF8)
    BORDER_COLOR = RGBColor(0xE8, 0xE8, 0xE8)

    # 1. Add Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.0), Inches(0.8))
    tf_title = title_box.text_frame
    p_title = tf_title.paragraphs[0]
    p_title.text = "目录：构建专业PPT的蓝图"
    p_title.font.name = "Microsoft YaHei"
    p_title.font.size = Pt(36)
    p_title.font.bold = True
    p_title.font.color.rgb = BLUE_TITLE

    # 2. Add Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.0), Inches(0.5))
    tf_sub = sub_box.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "本次分享的核心框架"
    p_sub.font.name = "Microsoft YaHei"
    p_sub.font.size = Pt(20)
    p_sub.font.bold = True
    p_sub.font.color.rgb = DARK_GRAY

    # Helper function to create styled text boxes
    def add_node_box(left, top, width, height, text_parts):
        # Shadow
        shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.06), top + Inches(0.06), width, height)
        shadow.fill.solid()
        shadow.fill.fore_color.rgb = SHADOW_COLOR
        shadow.line.fill.background()

        # Main Box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        box.line.color.rgb = BORDER_COLOR
        box.line.width = Pt(1)

        tf = box.text_frame
        tf.margin_left = Inches(0.3)
        tf.margin_top = Inches(0.15)
        p = tf.paragraphs[0]

        for text, color, is_bold in text_parts:
            run = p.add_run()
            run.text = text
            run.font.name = "Microsoft YaHei"
            run.font.size = Pt(18)
            run.font.color.rgb = color
            run.font.bold = is_bold

    # 3. Main Vertical Timeline
    v_line_x = 1.8
    y_nodes = [2.6, 3.8, 5.0, 6.2]
    
    main_v_line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(v_line_x), Inches(y_nodes[0]), Inches(v_line_x), Inches(y_nodes[3]))
    main_v_line.line.color.rgb = BLUE_LINE
    main_v_line.line.width = Pt(1.5)

    # --- Node 1 ---
    y1 = y_nodes[0]
    # Circle
    c1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(v_line_x - 0.08), Inches(y1 - 0.08), Inches(0.16), Inches(0.16))
    c1.fill.solid()
    c1.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    c1.line.color.rgb = BLUE_LINE
    c1.line.width = Pt(2)
    # H-Line
    hl1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(v_line_x + 0.08), Inches(y1), Inches(2.0), Inches(y1))
    hl1.line.color.rgb = BLUE_LINE
    hl1.line.width = Pt(1.5)
    # Text Box
    add_node_box(Inches(2.0), Inches(y1 - 0.35), Inches(5.8), Inches(0.7), [
        ("1. 内容法则：一页一事，", BLACK_TEXT, True),
        ("结论先行", ORANGE_TEXT, True)
    ])
    # Icon 1 (Document)
    doc = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(y1 - 0.4), Inches(0.5), Inches(0.7))
    doc.fill.background()
    doc.line.color.rgb = BLUE_LINE
    doc.line.width = Pt(2)
    for i in range(3):
        l = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.9), Inches(y1 - 0.2 + i*0.15), Inches(1.2), Inches(y1 - 0.2 + i*0.15))
        l.line.color.rgb = BLUE_LINE
        l.line.width = Pt(1.5)
    check_bg = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.15), Inches(y1 - 0.05), Inches(0.3), Inches(0.3))
    check_bg.fill.solid()
    check_bg.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    check_bg.line.color.rgb = BLUE_LINE
    check_bg.line.width = Pt(2)
    ck1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.22), Inches(y1 + 0.1), Inches(1.28), Inches(y1 + 0.16))
    ck1.line.color.rgb = BLUE_LINE
    ck1.line.width = Pt(2)
    ck2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.28), Inches(y1 + 0.16), Inches(1.38), Inches(y1 + 0.02))
    ck2.line.color.rgb = BLUE_LINE
    ck2.line.width = Pt(2)

    # --- Node 2 ---
    y2 = y_nodes[1]
    # Circle
    c2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(v_line_x - 0.08), Inches(y2 - 0.08), Inches(0.16), Inches(0.16))
    c2.fill.solid()
    c2.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    c2.line.color.rgb = BLUE_LINE
    c2.line.width = Pt(2)
    # H-Line
    hl2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(v_line_x + 0.08), Inches(y2), Inches(2.4), Inches(y2))
    hl2.line.color.rgb = BLUE_LINE
    hl2.line.width = Pt(1.5)
    # Text Box
    add_node_box(Inches(3.4), Inches(y2 - 0.35), Inches(6.4), Inches(0.7), [
        ("2. 减法艺术：拒绝文字堆砌，追求", BLACK_TEXT, True),
        ("秒懂", ORANGE_TEXT, True)
    ])
    # Icon 2 (Trash)
    trash = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.45), Inches(y2 - 0.2), Inches(0.35), Inches(0.45))
    trash.fill.background()
    trash.line.color.rgb = BLUE_LINE
    trash.line.width = Pt(2)
    lid = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.35), Inches(y2 - 0.3), Inches(0.55), Inches(0.08))
    lid.fill.background()
    lid.line.color.rgb = BLUE_LINE
    lid.line.width = Pt(2)
    handle = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.55), Inches(y2 - 0.38), Inches(0.15), Inches(0.08))
    handle.fill.background()
    handle.line.color.rgb = BLUE_LINE
    handle.line.width = Pt(1.5)
    for i in range(3):
        vl = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(2.55 + i*0.08), Inches(y2 - 0.15), Inches(2.55 + i*0.08), Inches(y2 + 0.2))
        vl.line.color.rgb = BLUE_LINE
        vl.line.width = Pt(1.5)
    for i in range(3):
        hl = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(2.95), Inches(y2 - 0.1 + i*0.15), Inches(3.25), Inches(y2 - 0.1 + i*0.15))
        hl.line.color.rgb = BLUE_LINE
        hl.line.width = Pt(2)

    # --- Node 3 ---
    y3 = y_nodes[2]
    # Circle
    c3 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(v_line_x - 0.08), Inches(y3 - 0.08), Inches(0.16), Inches(0.16))
    c3.fill.solid()
    c3.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    c3.line.color.rgb = BLUE_LINE
    c3.line.width = Pt(2)
    # H-Line
    hl3 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(v_line_x + 0.08), Inches(y3), Inches(3.4), Inches(y3))
    hl3.line.color.rgb = BLUE_LINE
    hl3.line.width = Pt(1.5)
    # Text Box
    add_node_box(Inches(4.3), Inches(y3 - 0.35), Inches(6.4), Inches(0.7), [
        ("3. 设计规范：高度统一，建立", BLACK_TEXT, True),
        ("专业感", ORANGE_TEXT, True)
    ])
    # Icon 3 (Gears & Ruler)
    g1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.4), Inches(y3 - 0.4), Inches(0.4), Inches(0.4))
    g1.fill.background()
    g1.line.color.rgb = BLUE_LINE
    g1.line.width = Pt(2.5)
    g1_in = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.55), Inches(y3 - 0.25), Inches(0.1), Inches(0.1))
    g1_in.fill.background()
    g1_in.line.color.rgb = BLUE_LINE
    g1_in.line.width = Pt(1.5)
    g2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.75), Inches(y3 - 0.25), Inches(0.3), Inches(0.3))
    g2.fill.background()
    g2.line.color.rgb = BLUE_LINE
    g2.line.width = Pt(2)
    ruler = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.4), Inches(y3 + 0.1), Inches(0.7), Inches(0.15))
    ruler.fill.background()
    ruler.line.color.rgb = BLUE_LINE
    ruler.line.width = Pt(1.5)
    for i in range(6):
        tick = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(3.45 + i*0.1), Inches(y3 + 0.1), Inches(3.45 + i*0.1), Inches(y3 + 0.18))
        tick.line.color.rgb = BLUE_LINE

    # --- Node 4 ---
    y4 = y_nodes[3]
    # Circle
    c4 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(v_line_x - 0.08), Inches(y4 - 0.08), Inches(0.16), Inches(0.16))
    c4.fill.solid()
    c4.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    c4.line.color.rgb = BLUE_LINE
    c4.line.width = Pt(2)
    # H-Line
    hl4 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(v_line_x + 0.08), Inches(y4), Inches(4.2), Inches(y4))
    hl4.line.color.rgb = BLUE_LINE
    hl4.line.width = Pt(1.5)
    # Text Box
    add_node_box(Inches(5.3), Inches(y4 - 0.35), Inches(6.3), Inches(0.7), [
        ("4. 高手境界：简洁有力的视觉哲学", BLACK_TEXT, True)
    ])
    # Icon 4 (Mountain & Scale)
    mt = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(4.2), Inches(y4 - 0.3), Inches(0.65), Inches(0.5))
    mt.fill.background()
    mt.line.color.rgb = BLUE_LINE
    mt.line.width = Pt(2)
    snow = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(4.42), Inches(y4 - 0.3), Inches(0.21), Inches(0.15))
    snow.fill.background()
    snow.line.color.rgb = BLUE_LINE
    snow.line.width = Pt(1.5)
    
    base = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(5.15), Inches(y4 + 0.1), Inches(0.3), Inches(0.1))
    base.fill.background()
    base.line.color.rgb = BLUE_LINE
    base.line.width = Pt(1.5)
    post = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(5.3), Inches(y4 - 0.3), Inches(5.3), Inches(y4 + 0.1))
    post.line.color.rgb = BLUE_LINE
    post.line.width = Pt(2)
    beam = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(5.0), Inches(y4 - 0.2), Inches(5.6), Inches(y4 - 0.2))
    beam.line.color.rgb = BLUE_LINE
    beam.line.width = Pt(2)
    
    p1_v = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(5.05), Inches(y4 - 0.2), Inches(5.05), Inches(y4))
    p1_v.line.color.rgb = BLUE_LINE
    p1_h = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(4.95), Inches(y4), Inches(5.15), Inches(y4))
    p1_h.line.color.rgb = BLUE_LINE
    p1_h.line.width = Pt(1.5)
    
    p2_v = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(5.55), Inches(y4 - 0.2), Inches(5.55), Inches(y4))
    p2_v.line.color.rgb = BLUE_LINE
    p2_h = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(5.45), Inches(y4), Inches(5.65), Inches(y4))
    p2_h.line.color.rgb = BLUE_LINE
    p2_h.line.width = Pt(1.5)

    # 4. Page Number
    page_num = slide.shapes.add_textbox(Inches(12.5), Inches(6.8), Inches(0.5), Inches(0.5))
    tf_num = page_num.text_frame
    p_num = tf_num.paragraphs[0]
    p_num.text = "02"
    p_num.font.name = "Arial"
    p_num.font.size = Pt(14)
    p_num.font.color.rgb = RGBColor(0x66, 0x66, 0x66)