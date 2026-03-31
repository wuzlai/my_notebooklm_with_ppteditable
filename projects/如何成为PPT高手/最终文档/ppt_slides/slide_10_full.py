from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# ── Color Palette ──
BLUE_HEADER = RGBColor(0x5B, 0x9B, 0xD5)
BLUE_DARK   = RGBColor(0x4A, 0x86, 0xC8)
CYAN        = RGBColor(0x00, 0xBC, 0xD4)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x33, 0x33, 0x33)
GRAY_TEXT   = RGBColor(0x55, 0x55, 0x55)
GRAY_BAR    = RGBColor(0xB0, 0xBE, 0xC5)
RED         = RGBColor(0xE5, 0x39, 0x35)
GREEN       = RGBColor(0x43, 0xA0, 0x47)
ORANGE      = RGBColor(0xFF, 0x98, 0x00)

ICON_BG  = RGBColor(0xE3, 0xE8, 0xED)
ICON_FG  = RGBColor(0x54, 0x6E, 0x7A)
FONT_NAME = "Microsoft YaHei"
SLIDE_WIDTH = Inches(13.333)
HEADER_H    = Inches(0.75)
SUBTITLE_Y  = Inches(0.95)


def add_header_banner(slide, title_text, bg_color=None):
    if bg_color is None:
        bg_color = BLUE_HEADER
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, HEADER_H
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = bg_color
    banner.line.fill.background()
    tf = banner.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(26)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME


def add_subtitle(slide, text, left, top, width=Inches(12), font_size=Pt(18)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.font.name = FONT_NAME
    return txBox


def add_icon_box(slide, left, top, symbol, size=Inches(0.48)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ICON_BG
    shape.line.fill.background()
    shape.adjustments[0] = 0.25
    tf = shape.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = symbol
    p.font.size = Pt(18)
    p.font.color.rgb = ICON_FG
    p.font.bold = False
    return shape


def add_bullet_item(slide, left, top, symbol, label, description,
                    width=Inches(5.5), desc_size=Pt(13)):
    add_icon_box(slide, left, top, symbol)
    text_left = left + Inches(0.65)
    txBox = slide.shapes.add_textbox(text_left, top - Inches(0.02), width, Inches(0.65))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run_label = p.add_run()
    run_label.text = label + "\uff1a"
    run_label.font.size = Pt(14)
    run_label.font.color.rgb = BLACK
    run_label.font.bold = True
    run_label.font.name = FONT_NAME
    run_desc = p.add_run()
    run_desc.text = description
    run_desc.font.size = desc_size
    run_desc.font.color.rgb = GRAY_TEXT
    run_desc.font.bold = False
    run_desc.font.name = FONT_NAME
    return txBox


def add_conclusion_box(slide, left, top, width, text, font_size=Pt(13)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = BLACK
    run.font.bold = True
    run.font.name = FONT_NAME
    return txBox


def add_table(slide, left, top, width, height, rows, cols, data,
              header_color=None, col_widths=None):
    if header_color is None:
        header_color = BLUE_HEADER
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if data[r][c] is not None else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(11)
                    run.font.name = FONT_NAME
                    if r == 0:
                        run.font.color.rgb = WHITE
                        run.font.bold = True
                    else:
                        run.font.color.rgb = BLACK
                        run.font.bold = False
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else RGBColor(0xF5, 0xF5, 0xF5)
            cell.margin_left = Pt(5)
            cell.margin_right = Pt(5)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
    return table_shape


def add_bar_chart(slide, left, top, width, height,
                  categories, values, title="", bar_colors=None):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', values)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    chart.chart_style = 2
    plot = chart.plots[0]
    plot.gap_width = 100
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CYAN
    series.has_data_labels = True
    dl = series.data_labels
    dl.font.size = Pt(13)
    dl.font.bold = True
    dl.font.color.rgb = BLACK
    dl.number_format = '0.#'
    dl.show_value = True
    dl.label_position = XL_LABEL_POSITION.OUTSIDE_END
    if bar_colors:
        for i, color in enumerate(bar_colors):
            pt = series.points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = color
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(12)
    cat_axis.tick_labels.font.name = FONT_NAME
    cat_axis.major_tick_mark = 2
    cat_axis.format.line.fill.background()
    val_axis = chart.value_axis
    val_axis.visible = False
    val_axis.major_tick_mark = 2
    val_axis.format.line.fill.background()
    val_axis.major_gridlines.format.line.fill.background()
    if title:
        chart.has_title = True
        ct = chart.chart_title.text_frame.paragraphs[0]
        ct.text = title
        ct.font.size = Pt(14)
        ct.font.bold = True
        ct.font.name = FONT_NAME
    else:
        chart.has_title = False
    return chart_frame


def add_callout_label(slide, left, top, text, bg_color=None, font_size=Pt(11)):
    if bg_color is None:
        bg_color = CYAN
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.3), Inches(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME
    return shape


def add_data_card(slide, left, top, width, height, value, label,
                  value_color=None, bg_color=None):
    if value_color is None:
        value_color = CYAN
    if bg_color is None:
        bg_color = WHITE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(3)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run1 = p1.add_run()
    run1.text = str(value)
    run1.font.size = Pt(24)
    run1.font.color.rgb = value_color
    run1.font.bold = True
    run1.font.name = FONT_NAME
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(10)
    run2.font.color.rgb = GRAY_TEXT
    run2.font.bold = False
    run2.font.name = FONT_NAME
    return shape


OUTPUT_PATH = r"C:\Users\Administrator\Desktop\notebook\my_notebooklm_with_ppteditable\projects\如何成为PPT高手\最终文档\ppt_slides\slide_10.pptx"


# ── Slide 1 ──

def build_slide_1(slide):
    from pptx.enum.shapes import MSO_CONNECTOR
    from pptx.enum.text import MSO_ANCHOR
    
    # 1. 添加背景 (浅蓝色)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xF4, 0xF8, 0xFC)
    bg.line.fill.background()

    # 2. 添加主标题
    title_box = slide.shapes.add_textbox(Inches(2), Inches(0.6), Inches(9.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "总结回顾：PPT制作的双翼"
    run.font.name = FONT_NAME
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1C, 0x55, 0xBA)

    # 3. 添加副标题
    sub_box = slide.shapes.add_textbox(Inches(4), Inches(1.3), Inches(5.333), Inches(0.5))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "核心要点清单"
    run.font.name = FONT_NAME
    run.font.size = Pt(18)
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # 4. 绘制卡片背景和阴影
    card_width = Inches(5.4)
    card_height = Inches(4.8)
    left_x = Inches(0.8)
    right_x = Inches(7.133)
    card_y = Inches(2.2)

    # 左侧卡片阴影
    shadow_l = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_x + Inches(0.05), card_y + Inches(0.05), card_width, card_height)
    shadow_l.fill.solid()
    shadow_l.fill.fore_color.rgb = RGBColor(0xE0, 0xE8, 0xF0)
    shadow_l.line.fill.background()
    
    # 左侧卡片主体
    card_l = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_x, card_y, card_width, card_height)
    card_l.fill.solid()
    card_l.fill.fore_color.rgb = RGBColor(0xF9, 0xFB, 0xFE)
    card_l.line.color.rgb = RGBColor(0xBD, 0xD7, 0xFA)
    card_l.line.width = Pt(2)

    # 右侧卡片阴影
    shadow_r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, right_x + Inches(0.05), card_y + Inches(0.05), card_width, card_height)
    shadow_r.fill.solid()
    shadow_r.fill.fore_color.rgb = RGBColor(0xE0, 0xE8, 0xF0)
    shadow_r.line.fill.background()

    # 右侧卡片主体
    card_r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, right_x, card_y, card_width, card_height)
    card_r.fill.solid()
    card_r.fill.fore_color.rgb = RGBColor(0xF9, 0xFB, 0xFE)
    card_r.line.color.rgb = RGBColor(0xBD, 0xD7, 0xFA)
    card_r.line.width = Pt(2)

    # 5. 添加卡片标题 (带背景色以遮挡边框)
    def add_card_title(x, y, text):
        title_bg = slide.shapes.add_textbox(x, y, Inches(2.6), Inches(0.4))
        title_bg.fill.solid()
        title_bg.fill.fore_color.rgb = RGBColor(0xF4, 0xF8, 0xFC)
        title_bg.line.fill.background()
        tf = title_bg.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.name = FONT_NAME
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x1C, 0x55, 0xBA)

    add_card_title(left_x + Inches(1.4), card_y - Inches(0.2), "内容要点清单")
    add_card_title(right_x + Inches(1.4), card_y - Inches(0.2), "设计要点清单")

    # 6. 绘制中心天平图标及分割线
    center_x = Inches(6.666)
    
    # 上分割线
    l_top = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, center_x, Inches(2.2), center_x, Inches(3.8))
    l_top.line.color.rgb = RGBColor(0x5B, 0x9B, 0xD5)
    l_top.line.width = Pt(2)

    # 下分割线
    l_bot = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, center_x, Inches(5.0), center_x, Inches(6.6))
    l_bot.line.color.rgb = RGBColor(0x5B, 0x9B, 0xD5)
    l_bot.line.width = Pt(2)

    # 中心圆
    cy = Inches(4.4)
    r = Inches(0.45)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, center_x - r, cy - r, r*2, r*2)
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0xE6, 0xF0, 0xFA)
    circle.line.fill.background()

    # 天平图形
    scale_color = RGBColor(0x1C, 0x55, 0xBA)
    base = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, center_x - Inches(0.15), cy + Inches(0.1), Inches(0.3), Inches(0.15))
    base.fill.solid()
    base.fill.fore_color.rgb = scale_color
    base.line.fill.background()
    
    pillar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, center_x - Inches(0.025), cy - Inches(0.15), Inches(0.05), Inches(0.25))
    pillar.fill.solid()
    pillar.fill.fore_color.rgb = scale_color
    pillar.line.fill.background()
    
    beam = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, center_x - Inches(0.3), cy - Inches(0.15), Inches(0.6), Inches(0.04))
    beam.fill.solid()
    beam.fill.fore_color.rgb = scale_color
    beam.line.fill.background()
    
    pan_l = slide.shapes.add_shape(MSO_SHAPE.OVAL, center_x - Inches(0.4), cy + Inches(0.05), Inches(0.2), Inches(0.05))
    pan_l.fill.solid()
    pan_l.fill.fore_color.rgb = scale_color
    pan_l.line.fill.background()
    
    sl1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, center_x - Inches(0.3), cy - Inches(0.11), center_x - Inches(0.4), cy + Inches(0.05))
    sl1.line.color.rgb = scale_color
    sl2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, center_x - Inches(0.3), cy - Inches(0.11), center_x - Inches(0.2), cy + Inches(0.05))
    sl2.line.color.rgb = scale_color
    
    pan_r = slide.shapes.add_shape(MSO_SHAPE.OVAL, center_x + Inches(0.2), cy + Inches(0.05), Inches(0.2), Inches(0.05))
    pan_r.fill.solid()
    pan_r.fill.fore_color.rgb = scale_color
    pan_r.line.fill.background()
    
    sr1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, center_x + Inches(0.3), cy - Inches(0.11), center_x + Inches(0.2), cy + Inches(0.05))
    sr1.line.color.rgb = scale_color
    sr2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, center_x + Inches(0.3), cy - Inches(0.11), center_x + Inches(0.4), cy + Inches(0.05))
    sr2.line.color.rgb = scale_color

    # 7. 定义绘制列表项的辅助函数
    def draw_item(x, y, icon_type, label_text, desc_text):
        icon_color = RGBColor(0x2B, 0x70, 0xC9)
        
        # 绘制图标
        if icon_type == 'doc':
            shape = slide.shapes.add_shape(MSO_SHAPE.FOLDED_CORNER, x, y, Inches(0.45), Inches(0.55))
            shape.fill.background()
            shape.line.color.rgb = icon_color
            shape.line.width = Pt(2)
            circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x+Inches(0.12), y+Inches(0.25), Inches(0.2), Inches(0.2))
            circle.fill.background()
            circle.line.color.rgb = icon_color
            circle.line.width = Pt(1.5)
        elif icon_type == 'palette':
            shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.55), Inches(0.55))
            shape.fill.background()
            shape.line.color.rgb = icon_color
            shape.line.width = Pt(2)
            for cx, cy_ in [(0.12, 0.12), (0.35, 0.15), (0.12, 0.35)]:
                sc = slide.shapes.add_shape(MSO_SHAPE.OVAL, x+Inches(cx), y+Inches(cy_), Inches(0.08), Inches(0.08))
                sc.fill.solid()
                sc.fill.fore_color.rgb = icon_color
                sc.line.fill.background()
            txt = slide.shapes.add_textbox(x+Inches(0.25), y+Inches(0.25), Inches(0.3), Inches(0.3))
            tf = txt.text_frame
            tf.margin_left, tf.margin_right, tf.margin_top, tf.margin_bottom = 0, 0, 0, 0
            p = txt.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = "A"
            run.font.name = "Arial"
            run.font.size = Pt(12)
            run.font.bold = True
            run.font.color.rgb = icon_color
        elif icon_type == 'funnel':
            shape = slide.shapes.add_shape(MSO_SHAPE.FUNNEL, x, y, Inches(0.55), Inches(0.55))
            shape.fill.background()
            shape.line.color.rgb = icon_color
            shape.line.width = Pt(2)
        elif icon_type == 'layout':
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.55), Inches(0.45))
            shape.fill.background()
            shape.line.color.rgb = icon_color
            shape.line.width = Pt(2)
            l1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, y+Inches(0.15), x+Inches(0.55), y+Inches(0.15))
            l1.line.color.rgb = icon_color
            l1.line.width = Pt(1.5)
            l2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x+Inches(0.2), y+Inches(0.15), x+Inches(0.2), y+Inches(0.45))
            l2.line.color.rgb = icon_color
            l2.line.width = Pt(1.5)
        elif icon_type == 'align':
            for i in range(3):
                r_w = Inches(0.25) if i % 2 == 0 else Inches(0.4)
                r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y+Inches(i*0.18), r_w, Inches(0.1))
                r.fill.background()
                r.line.color.rgb = icon_color
                r.line.width = Pt(1.5)
            vl = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x+Inches(0.55), y, x+Inches(0.55), y+Inches(0.46))
            vl.line.color.rgb = icon_color
            vl.line.width = Pt(1.5)
        elif icon_type == 'hierarchy':
            r1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x+Inches(0.15), y, Inches(0.25), Inches(0.12))
            r1.fill.background()
            r1.line.color.rgb = icon_color
            r1.line.width = Pt(1.5)
            for i in range(3):
                r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x+Inches(i*0.2), y+Inches(0.3), Inches(0.15), Inches(0.1))
                r.fill.background()
                r.line.color.rgb = icon_color
                r.line.width = Pt(1.5)
            c1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x+Inches(0.275), y+Inches(0.12), x+Inches(0.275), y+Inches(0.2))
            c1.line.color.rgb = icon_color
            c2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x+Inches(0.075), y+Inches(0.2), x+Inches(0.475), y+Inches(0.2))
            c2.line.color.rgb = icon_color
            for i in range(3):
                c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x+Inches(i*0.2+0.075), y+Inches(0.2), x+Inches(i*0.2+0.075), y+Inches(0.3))
                c.line.color.rgb = icon_color

        # 绘制对勾
        chk_x = x + Inches(0.7)
        chk_y = y + Inches(0.02)
        chk_bg = slide.shapes.add_shape(MSO_SHAPE.OVAL, chk_x, chk_y, Inches(0.22), Inches(0.22))
        chk_bg.fill.solid()
        chk_bg.fill.fore_color.rgb = RGBColor(0x43, 0xA0, 0x47)
        chk_bg.line.fill.background()
        
        chk_txt = slide.shapes.add_textbox(chk_x - Inches(0.05), chk_y - Inches(0.05), Inches(0.32), Inches(0.32))
        tf = chk_txt.text_frame
        tf.margin_left, tf.margin_right, tf.margin_top, tf.margin_bottom = 0, 0, 0, 0
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = "✔"
        run.font.name = "Segoe UI Symbol"
        run.font.size = Pt(10)
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        # 绘制标签
        lbl_x = chk_x + Inches(0.3)
        lbl_y = y - Inches(0.05)
        lbl_box = slide.shapes.add_textbox(lbl_x, lbl_y, Inches(3.8), Inches(0.4))
        tf = lbl_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = label_text
        run.font.name = FONT_NAME
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x22, 0x22, 0x22)

        # 绘制描述
        desc_y = lbl_y + Inches(0.35)
        desc_box = slide.shapes.add_textbox(lbl_x, desc_y, Inches(3.8), Inches(0.6))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = desc_text
        run.font.name = FONT_NAME
        run.font.size = Pt(11)
        run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    # 8. 填充左侧列表项
    item_l_x = left_x + Inches(0.4)
    draw_item(item_l_x, Inches(2.8), 'doc', "清晰：一页一事，拒绝文字堆砌", "确保每个页面聚焦一个核心观点，避免信息过载。")
    draw_item(item_l_x, Inches(4.0), 'palette', "统一：配色字体，保持风格一致", "建立标准的色彩和字体规范，维护整体视觉的连贯性。")
    draw_item(item_l_x, Inches(5.2), 'funnel', "简洁：去繁就简，追求高效沟通", "删除所有非必要元素，用最短路径传达最关键信息。")

    # 9. 填充右侧列表项
    item_r_x = right_x + Inches(0.4)
    draw_item(item_r_x, Inches(2.8), 'layout', "留白：呼吸空间，引导视觉焦点", "保持页面留白率在40%以上，让内容有足够的呼吸感。")
    draw_item(item_r_x, Inches(4.0), 'align', "对齐：严谨规范，建立秩序美感", "严格遵循左对齐或居中对齐原则，构建清晰的视觉轴线。")
    draw_item(item_r_x, Inches(5.2), 'hierarchy', "层级：主次分明，提升阅读效率", "通过字号、粗细和颜色的对比，明确信息的优先级。")

    # 10. 添加页脚
    footer = slide.shapes.add_textbox(Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "Page 10 of 11"
    run.font.name = "Arial"
    run.font.size = Pt(10)
    run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)



# ── Main ──
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

s0 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_1(s0)
prs.save(OUTPUT_PATH)
