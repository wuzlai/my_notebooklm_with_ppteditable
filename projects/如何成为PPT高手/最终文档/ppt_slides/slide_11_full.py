from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# ── Color Palette ──
BLUE_HEADER = RGBColor(0x5B, 0x9B, 0xD5)
BLUE_DARK   = RGBColor(0x4A, 0x86, 0xC8)
CYAN        = RGBColor(0x00, 0xBC, 0xD4)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x33, 0x33, 0x33)
GRAY_TEXT   = RGBColor(0x55, 0x55, 0x55)
GRAY_BAR    = RGBColor(0xB0, 0xBE, 0xC5)
RED         = RGBColor(0xE5, 0x39, 0x35)
GREEN       = RGBColor(0x43, 0xA0, 0x47)
ORANGE      = RGBColor(0xFF, 0x98, 0x00)

ICON_BG  = RGBColor(0xE3, 0xE8, 0xED)
ICON_FG  = RGBColor(0x54, 0x6E, 0x7A)
FONT_NAME = "Microsoft YaHei"
SLIDE_WIDTH = Inches(13.333)
HEADER_H    = Inches(0.75)
SUBTITLE_Y  = Inches(0.95)


def add_header_banner(slide, title_text, bg_color=None):
    if bg_color is None:
        bg_color = BLUE_HEADER
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, HEADER_H
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = bg_color
    banner.line.fill.background()
    tf = banner.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(26)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME


def add_subtitle(slide, text, left, top, width=Inches(12), font_size=Pt(18)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.font.name = FONT_NAME
    return txBox


def add_icon_box(slide, left, top, symbol, size=Inches(0.48)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ICON_BG
    shape.line.fill.background()
    shape.adjustments[0] = 0.25
    tf = shape.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = symbol
    p.font.size = Pt(18)
    p.font.color.rgb = ICON_FG
    p.font.bold = False
    return shape


def add_bullet_item(slide, left, top, symbol, label, description,
                    width=Inches(5.5), desc_size=Pt(13)):
    add_icon_box(slide, left, top, symbol)
    text_left = left + Inches(0.65)
    txBox = slide.shapes.add_textbox(text_left, top - Inches(0.02), width, Inches(0.65))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run_label = p.add_run()
    run_label.text = label + "\uff1a"
    run_label.font.size = Pt(14)
    run_label.font.color.rgb = BLACK
    run_label.font.bold = True
    run_label.font.name = FONT_NAME
    run_desc = p.add_run()
    run_desc.text = description
    run_desc.font.size = desc_size
    run_desc.font.color.rgb = GRAY_TEXT
    run_desc.font.bold = False
    run_desc.font.name = FONT_NAME
    return txBox


def add_conclusion_box(slide, left, top, width, text, font_size=Pt(13)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = BLACK
    run.font.bold = True
    run.font.name = FONT_NAME
    return txBox


def add_table(slide, left, top, width, height, rows, cols, data,
              header_color=None, col_widths=None):
    if header_color is None:
        header_color = BLUE_HEADER
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if data[r][c] is not None else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(11)
                    run.font.name = FONT_NAME
                    if r == 0:
                        run.font.color.rgb = WHITE
                        run.font.bold = True
                    else:
                        run.font.color.rgb = BLACK
                        run.font.bold = False
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else RGBColor(0xF5, 0xF5, 0xF5)
            cell.margin_left = Pt(5)
            cell.margin_right = Pt(5)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
    return table_shape


def add_bar_chart(slide, left, top, width, height,
                  categories, values, title="", bar_colors=None):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', values)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    chart.chart_style = 2
    plot = chart.plots[0]
    plot.gap_width = 100
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CYAN
    series.has_data_labels = True
    dl = series.data_labels
    dl.font.size = Pt(13)
    dl.font.bold = True
    dl.font.color.rgb = BLACK
    dl.number_format = '0.#'
    dl.show_value = True
    dl.label_position = XL_LABEL_POSITION.OUTSIDE_END
    if bar_colors:
        for i, color in enumerate(bar_colors):
            pt = series.points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = color
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(12)
    cat_axis.tick_labels.font.name = FONT_NAME
    cat_axis.major_tick_mark = 2
    cat_axis.format.line.fill.background()
    val_axis = chart.value_axis
    val_axis.visible = False
    val_axis.major_tick_mark = 2
    val_axis.format.line.fill.background()
    val_axis.major_gridlines.format.line.fill.background()
    if title:
        chart.has_title = True
        ct = chart.chart_title.text_frame.paragraphs[0]
        ct.text = title
        ct.font.size = Pt(14)
        ct.font.bold = True
        ct.font.name = FONT_NAME
    else:
        chart.has_title = False
    return chart_frame


def add_callout_label(slide, left, top, text, bg_color=None, font_size=Pt(11)):
    if bg_color is None:
        bg_color = CYAN
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.3), Inches(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME
    return shape


def add_data_card(slide, left, top, width, height, value, label,
                  value_color=None, bg_color=None):
    if value_color is None:
        value_color = CYAN
    if bg_color is None:
        bg_color = WHITE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(3)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run1 = p1.add_run()
    run1.text = str(value)
    run1.font.size = Pt(24)
    run1.font.color.rgb = value_color
    run1.font.bold = True
    run1.font.name = FONT_NAME
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(10)
    run2.font.color.rgb = GRAY_TEXT
    run2.font.bold = False
    run2.font.name = FONT_NAME
    return shape


OUTPUT_PATH = r"C:\Users\Administrator\Desktop\notebook\my_notebooklm_with_ppteditable\projects\如何成为PPT高手\最终文档\ppt_slides\slide_11.pptx"


# ── Slide 1 ──

def build_slide_1(slide):
    # 1. Background Grid (Top part - Light Blue)
    for i in range(1, 14):
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(i), 0, Inches(i), Inches(7.5))
        line.line.color.rgb = RGBColor(0xE8, 0xF0, 0xF8)
        line.line.width = Pt(0.5)
    for i in range(1, 8):
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, 0, Inches(i), Inches(13.333), Inches(i))
        line.line.color.rgb = RGBColor(0xE8, 0xF0, 0xF8)
        line.line.width = Pt(0.5)

    # 2. Bottom Dark Blue Background
    bottom_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.2), Inches(13.333), Inches(3.3))
    bottom_rect.fill.solid()
    bottom_rect.fill.fore_color.rgb = RGBColor(0x15, 0x43, 0x85)
    bottom_rect.line.fill.background()

    # Bottom Grid (Overlay on dark blue)
    for i in range(1, 14):
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(i), Inches(4.2), Inches(i), Inches(7.5))
        line.line.color.rgb = RGBColor(0x25, 0x53, 0x95)
        line.line.width = Pt(0.5)
    for i in range(5, 8):
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, 0, Inches(i), Inches(13.333), Inches(i))
        line.line.color.rgb = RGBColor(0x25, 0x53, 0x95)
        line.line.width = Pt(0.5)

    # 3. Title Text
    title_box = slide.shapes.add_textbox(Inches(2.66), Inches(0.8), Inches(8), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "感谢观看"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x00, 0x52, 0xCC)
    p.font.name = "Microsoft YaHei"

    subtitle_box = slide.shapes.add_textbox(Inches(2.66), Inches(2.1), Inches(8), Inches(0.8))
    tf_sub = subtitle_box.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "立即开始你的专业演示之旅"
    p_sub.alignment = PP_ALIGN.CENTER
    p_sub.font.size = Pt(28)
    p_sub.font.color.rgb = RGBColor(0x00, 0x52, 0xCC)
    p_sub.font.name = "Microsoft YaHei"

    # 4. Middle Icons and Text
    icon_color = RGBColor(0x00, 0x52, 0xCC)
    text_color = RGBColor(0x00, 0x00, 0x00)

    # Item 1: Compass
    compass = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.5), Inches(3.2), Inches(0.6), Inches(0.6))
    compass.fill.background()
    compass.line.color.rgb = icon_color
    compass.line.width = Pt(2)
    needle = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.65), Inches(3.65), Inches(1.95), Inches(3.35))
    needle.line.color.rgb = icon_color
    needle.line.width = Pt(2)
    
    tb1 = slide.shapes.add_textbox(Inches(2.2), Inches(3.1), Inches(2.5), Inches(0.8))
    tf1 = tb1.text_frame
    p1 = tf1.paragraphs[0]
    p1.text = "1. 实践是提升PPT能\n力的唯一捷径"
    p1.font.size = Pt(16)
    p1.font.color.rgb = text_color
    p1.font.name = "Microsoft YaHei"

    # Item 2: Lightbulb
    bulb = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.2), Inches(3.2), Inches(0.5), Inches(0.5))
    bulb.fill.background()
    bulb.line.color.rgb = icon_color
    bulb.line.width = Pt(2)
    base = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.35), Inches(3.7), Inches(0.2), Inches(0.15))
    base.fill.background()
    base.line.color.rgb = icon_color
    base.line.width = Pt(2)
    ray1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(5.45), Inches(3.1), Inches(5.45), Inches(2.95))
    ray1.line.color.rgb = icon_color
    ray1.line.width = Pt(1.5)
    ray2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(5.05), Inches(3.45), Inches(4.9), Inches(3.45))
    ray2.line.color.rgb = icon_color
    ray2.line.width = Pt(1.5)
    ray3 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(5.85), Inches(3.45), Inches(6.0), Inches(3.45))
    ray3.line.color.rgb = icon_color
    ray3.line.width = Pt(1.5)
    
    tb2 = slide.shapes.add_textbox(Inches(5.9), Inches(3.1), Inches(2.5), Inches(0.8))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "2. 保持好奇，持续优\n化视觉表达"
    p2.font.size = Pt(16)
    p2.font.color.rgb = text_color
    p2.font.name = "Microsoft YaHei"

    # Item 3: Rocket
    rocket = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.9), Inches(3.2), Inches(0.4), Inches(0.6))
    rocket.rotation = 45
    rocket.fill.background()
    rocket.line.color.rgb = icon_color
    rocket.line.width = Pt(2)
    wing1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(8.9), Inches(3.6), Inches(8.7), Inches(3.8))
    wing1.line.color.rgb = icon_color
    wing1.line.width = Pt(2)
    wing2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(9.2), Inches(3.7), Inches(9.4), Inches(3.9))
    wing2.line.color.rgb = icon_color
    wing2.line.width = Pt(2)
    
    tb3 = slide.shapes.add_textbox(Inches(9.6), Inches(3.2), Inches(2.5), Inches(0.8))
    tf3 = tb3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "3. 期待您的精彩呈现"
    p3.font.size = Pt(16)
    p3.font.color.rgb = text_color
    p3.font.name = "Microsoft YaHei"

    # 5. Bottom Section - Faint Chart
    chart_color = RGBColor(0x4A, 0x76, 0xB5)
    
    # Chart Grid
    for x in [4.5, 5.0, 5.5, 6.0, 6.5, 7.0, 7.5, 8.0, 8.5, 9.0]:
        v_line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(4.8), Inches(x), Inches(6.2))
        v_line.line.color.rgb = chart_color
        v_line.line.width = Pt(0.5)
        
    # X and Y axis
    x_axis = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(4.0), Inches(6.2), Inches(9.3), Inches(6.2))
    x_axis.line.color.rgb = chart_color
    x_axis.line.width = Pt(1)
    y_axis = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(4.0), Inches(4.8), Inches(4.0), Inches(6.2))
    y_axis.line.color.rgb = chart_color
    y_axis.line.width = Pt(1)
    
    # Trend line
    t1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(4.0), Inches(6.0), Inches(5.5), Inches(5.7))
    t1.line.color.rgb = chart_color
    t1.line.width = Pt(1.5)
    t2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(5.5), Inches(5.7), Inches(7.0), Inches(5.9))
    t2.line.color.rgb = chart_color
    t2.line.width = Pt(1.5)
    t3 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(7.0), Inches(5.9), Inches(9.0), Inches(4.8))
    t3.line.color.rgb = chart_color
    t3.line.width = Pt(1.5)

    # 6. "Thank You" Text
    ty_box = slide.shapes.add_textbox(Inches(2.66), Inches(4.6), Inches(8), Inches(1.5))
    tf_ty = ty_box.text_frame
    p_ty = tf_ty.paragraphs[0]
    p_ty.text = "Thank You"
    p_ty.alignment = PP_ALIGN.CENTER
    p_ty.font.size = Pt(72)
    p_ty.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p_ty.font.name = "Arial"

    # 7. Footer
    # QR Code Scanner Icon
    # Brackets
    brackets = [
        (3.5, 6.5, 3.6, 6.5), (3.5, 6.5, 3.5, 6.6), # Top-left
        (4.0, 6.5, 4.1, 6.5), (4.1, 6.5, 4.1, 6.6), # Top-right
        (3.5, 7.0, 3.6, 7.0), (3.5, 6.9, 3.5, 7.0), # Bottom-left
        (4.0, 7.0, 4.1, 7.0), (4.1, 6.9, 4.1, 7.0)  # Bottom-right
    ]
    for x1, y1, x2, y2 in brackets:
        bl = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        bl.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        bl.line.width = Pt(1.5)
        
    # Inner square and line
    sq = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.65), Inches(6.65), Inches(0.3), Inches(0.3))
    sq.fill.background()
    sq.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    sq.line.width = Pt(1)
    scan_line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(3.4), Inches(6.75), Inches(4.2), Inches(6.75))
    scan_line.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    scan_line.line.width = Pt(1.5)

    # Footer Text
    footer_box = slide.shapes.add_textbox(Inches(4.2), Inches(6.55), Inches(7.0), Inches(0.5))
    tf_footer = footer_box.text_frame
    p_footer = tf_footer.paragraphs[0]
    p_footer.text = "扫描二维码联系 | 联系方式: support@example.com | www.example.com"
    p_footer.font.size = Pt(11)
    p_footer.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p_footer.font.name = "Microsoft YaHei"

    # Page Number
    page_box = slide.shapes.add_textbox(Inches(11.8), Inches(6.8), Inches(1.0), Inches(0.5))
    tf_page = page_box.text_frame
    p_page = tf_page.paragraphs[0]
    p_page.text = "11 / 11"
    p_page.alignment = PP_ALIGN.RIGHT
    p_page.font.size = Pt(12)
    p_page.font.color.rgb = RGBColor(0x8A, 0xB4, 0xF8)
    p_page.font.name = "Arial"



# ── Main ──
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

s0 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_1(s0)
prs.save(OUTPUT_PATH)
