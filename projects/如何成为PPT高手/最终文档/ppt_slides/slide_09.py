def build_slide(slide):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(10), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "高手境界：简洁有力的视觉哲学"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    # Divider Line
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.8), Inches(1.3), Inches(12.5), Inches(1.3))
    line.line.color.rgb = RGBColor(0x1F, 0x4E, 0x79)
    line.line.width = Pt(1.5)

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(10), Inches(0.6))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "真正的专业不需要花哨的装饰 (20-24pt)"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(0x7F, 0x7F, 0x7F)

    # --- Item 1 ---
    # Icon 1: Eraser
    eraser_line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.9), Inches(3.9), Inches(1.6), Inches(3.9))
    eraser_line.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    eraser_line.line.width = Pt(1.5)

    eraser = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(3.0), Inches(0.8), Inches(0.45))
    eraser.rotation = -45
    eraser.fill.solid()
    eraser.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFA)
    eraser.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    eraser.line.width = Pt(1.5)

    x1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.6), Inches(3.7), Inches(1.8), Inches(3.9))
    x1.line.color.rgb = RGBColor(0xE5, 0x39, 0x35)
    x1.line.width = Pt(2)
    x2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.8), Inches(3.7), Inches(1.6), Inches(3.9))
    x2.line.color.rgb = RGBColor(0xE5, 0x39, 0x35)
    x2.line.width = Pt(2)

    # Bullet 1
    dot1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.4), Inches(3.45), Inches(0.08), Inches(0.08))
    dot1.fill.solid()
    dot1.fill.fore_color.rgb = RGBColor(0x00, 0x70, 0xC0)
    dot1.line.fill.background()

    # Text 1
    tb1 = slide.shapes.add_textbox(Inches(2.6), Inches(3.2), Inches(8), Inches(0.6))
    p1 = tb1.text_frame.paragraphs[0]
    p1.text = "删掉所有不承载信息的装饰性元素"
    p1.font.name = "Microsoft YaHei"
    p1.font.size = Pt(22)
    p1.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    # --- Item 2 ---
    # Icon 2: Circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.9), Inches(4.7), Inches(1.0), Inches(1.0))
    circle.fill.background()
    circle.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    circle.line.width = Pt(1.5)

    dot_inner = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.35), Inches(5.15), Inches(0.1), Inches(0.1))
    dot_inner.fill.solid()
    dot_inner.fill.fore_color.rgb = RGBColor(0xDA, 0xA5, 0x20)
    dot_inner.line.fill.background()

    # Bullet 2
    dot2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.4), Inches(5.15), Inches(0.08), Inches(0.08))
    dot2.fill.solid()
    dot2.fill.fore_color.rgb = RGBColor(0x00, 0x70, 0xC0)
    dot2.line.fill.background()

    # Text 2
    tb2 = slide.shapes.add_textbox(Inches(2.6), Inches(4.9), Inches(8), Inches(0.6))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = "留白不是浪费，而是更高级的强调"
    p2.font.name = "Microsoft YaHei"
    p2.font.size = Pt(22)
    p2.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    # --- Item 3 ---
    # Icon 3: Anchor
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(6.4), Inches(0.3), Inches(1.0))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0xDA, 0xA5, 0x20)
    rect.line.fill.background()

    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.4), Inches(6.4), Inches(0.2), Inches(0.2))
    c.fill.background()
    c.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    c.line.width = Pt(1.5)

    stem = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.5), Inches(6.6), Inches(1.5), Inches(7.3))
    stem.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    stem.line.width = Pt(1.5)

    cross = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.35), Inches(6.75), Inches(1.65), Inches(6.75))
    cross.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    cross.line.width = Pt(1.5)

    l1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.1), Inches(7.0), Inches(1.5), Inches(7.3))
    l1.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    l1.line.width = Pt(1.5)
    
    l2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.9), Inches(7.0), Inches(1.5), Inches(7.3))
    l2.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    l2.line.width = Pt(1.5)

    a1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.1), Inches(7.0), Inches(1.1), Inches(7.15))
    a1.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    a1.line.width = Pt(1.5)
    
    a2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.1), Inches(7.0), Inches(1.25), Inches(7.0))
    a2.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    a2.line.width = Pt(1.5)

    a3 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.9), Inches(7.0), Inches(1.9), Inches(7.15))
    a3.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    a3.line.width = Pt(1.5)
    
    a4 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.9), Inches(7.0), Inches(1.75), Inches(7.0))
    a4.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    a4.line.width = Pt(1.5)

    # Bullet 3
    dot3 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.4), Inches(6.85), Inches(0.08), Inches(0.08))
    dot3.fill.solid()
    dot3.fill.fore_color.rgb = RGBColor(0x00, 0x70, 0xC0)
    dot3.line.fill.background()

    # Text 3
    tb3 = slide.shapes.add_textbox(Inches(2.6), Inches(6.6), Inches(8), Inches(0.6))
    p3 = tb3.text_frame.paragraphs[0]
    p3.text = "简洁即是力量，克制即是专业"
    p3.font.name = "Microsoft YaHei"
    p3.font.size = Pt(22)
    p3.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    # Footer
    footer = slide.shapes.add_textbox(Inches(12.0), Inches(7.0), Inches(1.0), Inches(0.4))
    p_foot = footer.text_frame.paragraphs[0]
    p_foot.text = "第9页"
    p_foot.font.name = "Microsoft YaHei"
    p_foot.font.size = Pt(12)
    p_foot.font.color.rgb = RGBColor(0x7F, 0x7F, 0x7F)