from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# ── Color Palette ──
BLUE_HEADER = RGBColor(0x5B, 0x9B, 0xD5)
BLUE_DARK   = RGBColor(0x4A, 0x86, 0xC8)
CYAN        = RGBColor(0x00, 0xBC, 0xD4)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x33, 0x33, 0x33)
GRAY_TEXT   = RGBColor(0x55, 0x55, 0x55)
GRAY_BAR    = RGBColor(0xB0, 0xBE, 0xC5)
RED         = RGBColor(0xE5, 0x39, 0x35)
GREEN       = RGBColor(0x43, 0xA0, 0x47)
ORANGE      = RGBColor(0xFF, 0x98, 0x00)

ICON_BG  = RGBColor(0xE3, 0xE8, 0xED)
ICON_FG  = RGBColor(0x54, 0x6E, 0x7A)
FONT_NAME = "Microsoft YaHei"
SLIDE_WIDTH = Inches(13.333)
HEADER_H    = Inches(0.75)
SUBTITLE_Y  = Inches(0.95)


def add_header_banner(slide, title_text, bg_color=None):
    if bg_color is None:
        bg_color = BLUE_HEADER
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, HEADER_H
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = bg_color
    banner.line.fill.background()
    tf = banner.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(26)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME


def add_subtitle(slide, text, left, top, width=Inches(12), font_size=Pt(18)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.font.name = FONT_NAME
    return txBox


def add_icon_box(slide, left, top, symbol, size=Inches(0.48)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ICON_BG
    shape.line.fill.background()
    shape.adjustments[0] = 0.25
    tf = shape.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = symbol
    p.font.size = Pt(18)
    p.font.color.rgb = ICON_FG
    p.font.bold = False
    return shape


def add_bullet_item(slide, left, top, symbol, label, description,
                    width=Inches(5.5), desc_size=Pt(13)):
    add_icon_box(slide, left, top, symbol)
    text_left = left + Inches(0.65)
    txBox = slide.shapes.add_textbox(text_left, top - Inches(0.02), width, Inches(0.65))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run_label = p.add_run()
    run_label.text = label + "\uff1a"
    run_label.font.size = Pt(14)
    run_label.font.color.rgb = BLACK
    run_label.font.bold = True
    run_label.font.name = FONT_NAME
    run_desc = p.add_run()
    run_desc.text = description
    run_desc.font.size = desc_size
    run_desc.font.color.rgb = GRAY_TEXT
    run_desc.font.bold = False
    run_desc.font.name = FONT_NAME
    return txBox


def add_conclusion_box(slide, left, top, width, text, font_size=Pt(13)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = BLACK
    run.font.bold = True
    run.font.name = FONT_NAME
    return txBox


def add_table(slide, left, top, width, height, rows, cols, data,
              header_color=None, col_widths=None):
    if header_color is None:
        header_color = BLUE_HEADER
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if data[r][c] is not None else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(11)
                    run.font.name = FONT_NAME
                    if r == 0:
                        run.font.color.rgb = WHITE
                        run.font.bold = True
                    else:
                        run.font.color.rgb = BLACK
                        run.font.bold = False
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else RGBColor(0xF5, 0xF5, 0xF5)
            cell.margin_left = Pt(5)
            cell.margin_right = Pt(5)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
    return table_shape


def add_bar_chart(slide, left, top, width, height,
                  categories, values, title="", bar_colors=None):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', values)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    chart.chart_style = 2
    plot = chart.plots[0]
    plot.gap_width = 100
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CYAN
    series.has_data_labels = True
    dl = series.data_labels
    dl.font.size = Pt(13)
    dl.font.bold = True
    dl.font.color.rgb = BLACK
    dl.number_format = '0.#'
    dl.show_value = True
    dl.label_position = XL_LABEL_POSITION.OUTSIDE_END
    if bar_colors:
        for i, color in enumerate(bar_colors):
            pt = series.points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = color
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(12)
    cat_axis.tick_labels.font.name = FONT_NAME
    cat_axis.major_tick_mark = 2
    cat_axis.format.line.fill.background()
    val_axis = chart.value_axis
    val_axis.visible = False
    val_axis.major_tick_mark = 2
    val_axis.format.line.fill.background()
    val_axis.major_gridlines.format.line.fill.background()
    if title:
        chart.has_title = True
        ct = chart.chart_title.text_frame.paragraphs[0]
        ct.text = title
        ct.font.size = Pt(14)
        ct.font.bold = True
        ct.font.name = FONT_NAME
    else:
        chart.has_title = False
    return chart_frame


def add_callout_label(slide, left, top, text, bg_color=None, font_size=Pt(11)):
    if bg_color is None:
        bg_color = CYAN
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.3), Inches(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME
    return shape


def add_data_card(slide, left, top, width, height, value, label,
                  value_color=None, bg_color=None):
    if value_color is None:
        value_color = CYAN
    if bg_color is None:
        bg_color = WHITE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(3)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run1 = p1.add_run()
    run1.text = str(value)
    run1.font.size = Pt(24)
    run1.font.color.rgb = value_color
    run1.font.bold = True
    run1.font.name = FONT_NAME
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(10)
    run2.font.color.rgb = GRAY_TEXT
    run2.font.bold = False
    run2.font.name = FONT_NAME
    return shape


OUTPUT_PATH = r"C:\Users\Administrator\Desktop\notebook\my_notebooklm_with_ppteditable\projects\如何成为PPT高手\最终文档\ppt_slides\slide_05.pptx"


# ── Slide 1 ──

def build_slide_1(slide):
    # 颜色常量定义
    BLUE_DARK = RGBColor(0x00, 0x55, 0xA4)
    BLUE_VERY_DARK = RGBColor(0x1A, 0x2B, 0x3C)
    ORANGE_TEXT = RGBColor(0xFF, 0x98, 0x00)
    GRAY_DARK = RGBColor(0x33, 0x33, 0x33)
    GRAY_LIGHT = RGBColor(0x55, 0x55, 0x55)
    GRAY_LINE = RGBColor(0xCC, 0xCC, 0xCC)
    GRAY_HANDLE = RGBColor(0x88, 0x88, 0x88)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    # 1. 标题
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(10), Inches(0.8))
    tf_title = title_box.text_frame
    p_title = tf_title.paragraphs[0]
    p_title.text = "秒懂原则：视觉化的力量"
    p_title.font.name = FONT_NAME
    p_title.font.size = Pt(36)
    p_title.font.bold = True
    p_title.font.color.rgb = BLUE_DARK

    # 2. 副标题
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(10), Inches(0.5))
    tf_sub = sub_box.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "让观众在3秒内捕捉核心信息"
    p_sub.font.name = FONT_NAME
    p_sub.font.size = Pt(20)
    p_sub.font.color.rgb = BLUE_DARK

    # 3. 左侧主图：放大镜及内部元素
    # 放大镜手柄连接处 (灰色)
    conn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.8), Inches(5.2), Inches(0.6), Inches(0.8))
    conn.rotation = 45
    conn.fill.solid()
    conn.fill.fore_color.rgb = GRAY_HANDLE
    conn.line.color.rgb = BLUE_VERY_DARK
    conn.line.width = Pt(3)

    # 放大镜手柄主体 (蓝色)
    handle = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(5.8), Inches(0.9), Inches(2.2))
    handle.rotation = 45
    handle.fill.solid()
    handle.fill.fore_color.rgb = BLUE_DARK
    handle.line.color.rgb = BLUE_VERY_DARK
    handle.line.width = Pt(4)

    # 放大镜外圈 (深色粗边框，白色填充以遮挡手柄顶部)
    outer_ring = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.0), Inches(2.2), Inches(4.0), Inches(4.0))
    outer_ring.fill.solid()
    outer_ring.fill.fore_color.rgb = WHITE
    outer_ring.line.color.rgb = BLUE_VERY_DARK
    outer_ring.line.width = Pt(12)

    # 放大镜内圈 (蓝色细边框)
    inner_ring = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.2), Inches(2.4), Inches(3.6), Inches(3.6))
    inner_ring.fill.background()
    inner_ring.line.color.rgb = BLUE_DARK
    inner_ring.line.width = Pt(6)

    # 内部放射状虚线/浅色线
    lines_coords = [
        (4.0, 2.9, 4.0, 2.6), # 上
        (4.0, 5.5, 4.0, 5.8), # 下
        (2.7, 4.2, 2.4, 4.2), # 左
        (5.3, 4.2, 5.6, 4.2), # 右
        (3.1, 3.3, 2.9, 3.1), # 左上
        (4.9, 3.3, 5.1, 3.1), # 右上
        (3.1, 5.1, 2.9, 5.3), # 左下
        (4.9, 5.1, 5.1, 5.3), # 右下
    ]
    for x1, y1, x2, y2 in lines_coords:
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        line.line.color.rgb = GRAY_LINE
        line.line.width = Pt(2)

    # 内部文字 "秒懂"
    text_box_md = slide.shapes.add_textbox(Inches(2.0), Inches(3.4), Inches(4.0), Inches(1.5))
    tf_md = text_box_md.text_frame
    p_md = tf_md.paragraphs[0]
    p_md.text = "秒懂"
    p_md.alignment = PP_ALIGN.CENTER
    p_md.font.name = FONT_NAME
    p_md.font.size = Pt(65)
    p_md.font.bold = True
    p_md.font.color.rgb = ORANGE_TEXT

    # 内部小图标：眼睛 (上方)
    eye_outer = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.3), Inches(3.1), Inches(0.4), Inches(0.25))
    eye_outer.fill.background()
    eye_outer.line.color.rgb = BLUE_DARK
    eye_outer.line.width = Pt(2)
    eye_inner = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.45), Inches(3.17), Inches(0.1), Inches(0.1))
    eye_inner.fill.solid()
    eye_inner.fill.fore_color.rgb = BLUE_DARK
    eye_inner.line.fill.background()

    # 内部小图标：大脑/云朵 (右下方)
    brain = slide.shapes.add_shape(MSO_SHAPE.CLOUD, Inches(4.7), Inches(4.5), Inches(0.4), Inches(0.3))
    brain.fill.background()
    brain.line.color.rgb = BLUE_DARK
    brain.line.width = Pt(1.5)

    # 放大镜外部装饰弧线 (使用简单的曲线近似)
    arc1 = slide.shapes.add_shape(MSO_SHAPE.ARC, Inches(1.6), Inches(1.8), Inches(4.8), Inches(4.8))
    arc1.fill.background()
    arc1.line.color.rgb = BLUE_VERY_DARK
    arc1.line.width = Pt(3)
    arc1.adjustments[0] = 110
    arc1.adjustments[1] = 160

    arc2 = slide.shapes.add_shape(MSO_SHAPE.ARC, Inches(1.8), Inches(2.0), Inches(4.4), Inches(4.4))
    arc2.fill.background()
    arc2.line.color.rgb = BLUE_VERY_DARK
    arc2.line.width = Pt(3)
    arc2.adjustments[0] = 20
    arc2.adjustments[1] = 70

    # 4. 右侧要点列表
    # --- 要点 1：降低认知负荷 ---
    # 图标 1 (人脑齿轮)
    head = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.4), Inches(2.5), Inches(0.5), Inches(0.6))
    head.fill.background()
    head.line.color.rgb = BLUE_DARK
    head.line.width = Pt(2)
    neck = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.55), Inches(3.05), Inches(0.2), Inches(0.15))
    neck.fill.background()
    neck.line.color.rgb = BLUE_DARK
    neck.line.width = Pt(2)
    gear1 = slide.shapes.add_shape(MSO_SHAPE.GEAR_6, Inches(7.45), Inches(2.6), Inches(0.25), Inches(0.25))
    gear1.fill.background()
    gear1.line.color.rgb = BLUE_DARK
    gear1.line.width = Pt(1.5)
    gear2 = slide.shapes.add_shape(MSO_SHAPE.GEAR_6, Inches(7.6), Inches(2.8), Inches(0.2), Inches(0.2))
    gear2.fill.background()
    gear2.line.color.rgb = BLUE_DARK
    gear2.line.width = Pt(1.5)

    # 文本 1
    tb1 = slide.shapes.add_textbox(Inches(8.4), Inches(2.4), Inches(4.5), Inches(1.0))
    tf1 = tb1.text_frame
    p1_1 = tf1.paragraphs[0]
    p1_1.text = "降低认知负荷："
    p1_1.font.name = FONT_NAME
    p1_1.font.size = Pt(18)
    p1_1.font.bold = True
    p1_1.font.color.rgb = GRAY_DARK
    p1_2 = tf1.add_paragraph()
    p1_2.text = "一眼就能看懂逻辑"
    p1_2.font.name = FONT_NAME
    p1_2.font.size = Pt(14)
    p1_2.font.color.rgb = GRAY_LIGHT
    p1_2.space_before = Pt(6)

    # --- 要点 2：视觉层级明确 ---
    # 图标 2 (层级金字塔与箭头)
    base_y2 = 4.8
    rect1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.3), Inches(base_y2), Inches(0.8), Inches(0.15))
    rect1.fill.background()
    rect1.line.color.rgb = BLUE_DARK
    rect1.line.width = Pt(2)
    rect2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.4), Inches(base_y2 - 0.25), Inches(0.6), Inches(0.15))
    rect2.fill.background()
    rect2.line.color.rgb = BLUE_DARK
    rect2.line.width = Pt(2)
    rect3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.5), Inches(base_y2 - 0.5), Inches(0.4), Inches(0.15))
    rect3.fill.background()
    rect3.line.color.rgb = BLUE_DARK
    rect3.line.width = Pt(2)
    arrow = slide.shapes.add_shape(MSO_SHAPE.UP_ARROW, Inches(8.2), Inches(base_y2 - 0.5), Inches(0.15), Inches(0.65))
    arrow.fill.background()
    arrow.line.color.rgb = BLUE_DARK
    arrow.line.width = Pt(1.5)

    # 文本 2
    tb2 = slide.shapes.add_textbox(Inches(8.4), Inches(4.1), Inches(4.5), Inches(1.0))
    tf2 = tb2.text_frame
    p2_1 = tf2.paragraphs[0]
    p2_1.text = "视觉层级明确："
    p2_1.font.name = FONT_NAME
    p2_1.font.size = Pt(18)
    p2_1.font.bold = True
    p2_1.font.color.rgb = GRAY_DARK
    p2_2 = tf2.add_paragraph()
    p2_2.text = "通过大小、颜色区分主次"
    p2_2.font.name = FONT_NAME
    p2_2.font.size = Pt(14)
    p2_2.font.color.rgb = GRAY_LIGHT
    p2_2.space_before = Pt(6)

    # --- 要点 3：善用图表 ---
    # 图标 3 (柱状图与饼图)
    base_y3 = 6.5
    axis_x = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(7.3), Inches(base_y3), Inches(8.0), Inches(base_y3))
    axis_x.line.color.rgb = BLUE_DARK
    axis_x.line.width = Pt(2)
    axis_y = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(7.3), Inches(base_y3), Inches(7.3), Inches(base_y3 - 0.8))
    axis_y.line.color.rgb = BLUE_DARK
    axis_y.line.width = Pt(2)
    
    bar1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.4), Inches(base_y3 - 0.3), Inches(0.12), Inches(0.3))
    bar1.fill.background()
    bar1.line.color.rgb = BLUE_DARK
    bar1.line.width = Pt(1.5)
    bar2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.6), Inches(base_y3 - 0.5), Inches(0.12), Inches(0.5))
    bar2.fill.background()
    bar2.line.color.rgb = BLUE_DARK
    bar2.line.width = Pt(1.5)
    bar3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.8), Inches(base_y3 - 0.2), Inches(0.12), Inches(0.2))
    bar3.fill.background()
    bar3.line.color.rgb = BLUE_DARK
    bar3.line.width = Pt(1.5)
    
    pie = slide.shapes.add_shape(MSO_SHAPE.PIE, Inches(7.7), Inches(base_y3 - 0.9), Inches(0.4), Inches(0.4))
    pie.fill.background()
    pie.line.color.rgb = BLUE_DARK
    pie.line.width = Pt(1.5)

    # 文本 3
    tb3 = slide.shapes.add_textbox(Inches(8.4), Inches(5.8), Inches(4.5), Inches(1.0))
    tf3 = tb3.text_frame
    p3_1 = tf3.paragraphs[0]
    p3_1.text = "善用图表："
    p3_1.font.name = FONT_NAME
    p3_1.font.size = Pt(18)
    p3_1.font.bold = True
    p3_1.font.color.rgb = GRAY_DARK
    p3_2 = tf3.add_paragraph()
    p3_2.text = "数据关系一目了然"
    p3_2.font.name = FONT_NAME
    p3_2.font.size = Pt(14)
    p3_2.font.color.rgb = GRAY_LIGHT
    p3_2.space_before = Pt(6)

    # 5. 页码
    page_num = slide.shapes.add_textbox(Inches(12.0), Inches(6.8), Inches(1.0), Inches(0.5))
    tf_page = page_num.text_frame
    p_page = tf_page.paragraphs[0]
    p_page.text = "5 / 11"
    p_page.font.name = FONT_NAME
    p_page.font.size = Pt(12)
    p_page.font.color.rgb = GRAY_HANDLE
    p_page.alignment = PP_ALIGN.RIGHT



# ── Main ──
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

s0 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_1(s0)
prs.save(OUTPUT_PATH)
