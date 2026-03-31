from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# ── Color Palette ──
BLUE_HEADER = RGBColor(0x5B, 0x9B, 0xD5)
BLUE_DARK   = RGBColor(0x4A, 0x86, 0xC8)
CYAN        = RGBColor(0x00, 0xBC, 0xD4)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x33, 0x33, 0x33)
GRAY_TEXT   = RGBColor(0x55, 0x55, 0x55)
GRAY_BAR    = RGBColor(0xB0, 0xBE, 0xC5)
RED         = RGBColor(0xE5, 0x39, 0x35)
GREEN       = RGBColor(0x43, 0xA0, 0x47)
ORANGE      = RGBColor(0xFF, 0x98, 0x00)

ICON_BG  = RGBColor(0xE3, 0xE8, 0xED)
ICON_FG  = RGBColor(0x54, 0x6E, 0x7A)
FONT_NAME = "Microsoft YaHei"
SLIDE_WIDTH = Inches(13.333)
HEADER_H    = Inches(0.75)
SUBTITLE_Y  = Inches(0.95)


def add_header_banner(slide, title_text, bg_color=None):
    if bg_color is None:
        bg_color = BLUE_HEADER
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, HEADER_H
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = bg_color
    banner.line.fill.background()
    tf = banner.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(26)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME


def add_subtitle(slide, text, left, top, width=Inches(12), font_size=Pt(18)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.font.name = FONT_NAME
    return txBox


def add_icon_box(slide, left, top, symbol, size=Inches(0.48)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ICON_BG
    shape.line.fill.background()
    shape.adjustments[0] = 0.25
    tf = shape.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = symbol
    p.font.size = Pt(18)
    p.font.color.rgb = ICON_FG
    p.font.bold = False
    return shape


def add_bullet_item(slide, left, top, symbol, label, description,
                    width=Inches(5.5), desc_size=Pt(13)):
    add_icon_box(slide, left, top, symbol)
    text_left = left + Inches(0.65)
    txBox = slide.shapes.add_textbox(text_left, top - Inches(0.02), width, Inches(0.65))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run_label = p.add_run()
    run_label.text = label + "\uff1a"
    run_label.font.size = Pt(14)
    run_label.font.color.rgb = BLACK
    run_label.font.bold = True
    run_label.font.name = FONT_NAME
    run_desc = p.add_run()
    run_desc.text = description
    run_desc.font.size = desc_size
    run_desc.font.color.rgb = GRAY_TEXT
    run_desc.font.bold = False
    run_desc.font.name = FONT_NAME
    return txBox


def add_conclusion_box(slide, left, top, width, text, font_size=Pt(13)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = BLACK
    run.font.bold = True
    run.font.name = FONT_NAME
    return txBox


def add_table(slide, left, top, width, height, rows, cols, data,
              header_color=None, col_widths=None):
    if header_color is None:
        header_color = BLUE_HEADER
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if data[r][c] is not None else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(11)
                    run.font.name = FONT_NAME
                    if r == 0:
                        run.font.color.rgb = WHITE
                        run.font.bold = True
                    else:
                        run.font.color.rgb = BLACK
                        run.font.bold = False
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else RGBColor(0xF5, 0xF5, 0xF5)
            cell.margin_left = Pt(5)
            cell.margin_right = Pt(5)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
    return table_shape


def add_bar_chart(slide, left, top, width, height,
                  categories, values, title="", bar_colors=None):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', values)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    chart.chart_style = 2
    plot = chart.plots[0]
    plot.gap_width = 100
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CYAN
    series.has_data_labels = True
    dl = series.data_labels
    dl.font.size = Pt(13)
    dl.font.bold = True
    dl.font.color.rgb = BLACK
    dl.number_format = '0.#'
    dl.show_value = True
    dl.label_position = XL_LABEL_POSITION.OUTSIDE_END
    if bar_colors:
        for i, color in enumerate(bar_colors):
            pt = series.points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = color
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(12)
    cat_axis.tick_labels.font.name = FONT_NAME
    cat_axis.major_tick_mark = 2
    cat_axis.format.line.fill.background()
    val_axis = chart.value_axis
    val_axis.visible = False
    val_axis.major_tick_mark = 2
    val_axis.format.line.fill.background()
    val_axis.major_gridlines.format.line.fill.background()
    if title:
        chart.has_title = True
        ct = chart.chart_title.text_frame.paragraphs[0]
        ct.text = title
        ct.font.size = Pt(14)
        ct.font.bold = True
        ct.font.name = FONT_NAME
    else:
        chart.has_title = False
    return chart_frame


def add_callout_label(slide, left, top, text, bg_color=None, font_size=Pt(11)):
    if bg_color is None:
        bg_color = CYAN
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.3), Inches(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME
    return shape


def add_data_card(slide, left, top, width, height, value, label,
                  value_color=None, bg_color=None):
    if value_color is None:
        value_color = CYAN
    if bg_color is None:
        bg_color = WHITE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(3)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run1 = p1.add_run()
    run1.text = str(value)
    run1.font.size = Pt(24)
    run1.font.color.rgb = value_color
    run1.font.bold = True
    run1.font.name = FONT_NAME
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(10)
    run2.font.color.rgb = GRAY_TEXT
    run2.font.bold = False
    run2.font.name = FONT_NAME
    return shape


OUTPUT_PATH = r"C:\Users\Administrator\Desktop\notebook\my_notebooklm_with_ppteditable\projects\如何成为PPT高手\最终文档\ppt_slides\slide_03.pptx"


# ── Slide 1 ──

def build_slide_1(slide):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
    from pptx.enum.text import PP_ALIGN

    # Colors
    DARK_BLUE = RGBColor(0x1F, 0x4E, 0x96)
    HIGHLIGHT_BLUE = RGBColor(0x00, 0x70, 0xC0)
    TEXT_BLACK = RGBColor(0x33, 0x33, 0x33)
    TEXT_GRAY = RGBColor(0x66, 0x66, 0x66)
    ORANGE_FILL = RGBColor(0xFF, 0x6B, 0x00)
    ORANGE_LINE = RGBColor(0xE6, 0x51, 0x00)

    # 1. Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(10.0), Inches(0.8))
    tf_title = title_box.text_frame
    p_title = tf_title.paragraphs[0]
    p_title.text = "法则一：内容清晰是PPT的灵魂"
    p_title.font.name = "Microsoft YaHei"
    p_title.font.size = Pt(32)
    p_title.font.bold = True
    p_title.font.color.rgb = DARK_BLUE

    # 2. Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(10.0), Inches(0.5))
    tf_sub = sub_box.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "确保观众的注意力始终聚焦"
    p_sub.font.name = "Microsoft YaHei"
    p_sub.font.size = Pt(18)
    p_sub.font.color.rgb = TEXT_BLACK

    # 3. Central Graphic (Target)
    cx, cy = 5.2, 4.2
    radii = [2.2, 1.75, 1.3, 0.85, 0.4]
    fills = [
        RGBColor(226, 238, 252),
        RGBColor(204, 224, 250),
        RGBColor(182, 210, 248),
        RGBColor(160, 196, 246),
        ORANGE_FILL
    ]
    
    for i, (r, fill) in enumerate(zip(radii, fills)):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, 
            Inches(cx - r), Inches(cy - r), 
            Inches(r * 2), Inches(r * 2)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        if i < 4:
            shape.line.color.rgb = DARK_BLUE
            shape.line.width = Pt(2.5)
        else:
            shape.line.fill.background() # No line for bullseye

    # 4. Central Graphic (Arrow)
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.NOTCHED_RIGHT_ARROW, 
        Inches(cx + 0.4), Inches(cy - 2.4), 
        Inches(2.4), Inches(0.7)
    )
    arrow.rotation = 135
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RGBColor(255, 167, 38)
    arrow.line.color.rgb = ORANGE_LINE
    arrow.line.width = Pt(1.5)

    # 5. Right Side Content - Item 1
    # Icon 1: Magnifying Glass
    mag_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.2), Inches(2.5), Inches(0.28), Inches(0.28))
    mag_circle.fill.background()
    mag_circle.line.color.rgb = DARK_BLUE
    mag_circle.line.width = Pt(2)
    handle = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(8.44), Inches(2.74), Inches(8.6), Inches(2.9))
    handle.line.color.rgb = DARK_BLUE
    handle.line.width = Pt(2.5)

    # Text 1
    tx_box1 = slide.shapes.add_textbox(Inches(8.8), Inches(2.35), Inches(4.0), Inches(1.0))
    tf1 = tx_box1.text_frame
    p1 = tf1.paragraphs[0]
    p1.font.size = Pt(16)
    p1.font.name = "Microsoft YaHei"
    r1_1 = p1.add_run(); r1_1.text = "每一页幻灯片只传达一\n"; r1_1.font.color.rgb = TEXT_BLACK
    r1_2 = p1.add_run(); r1_2.text = "个核心观点"; r1_2.font.color.rgb = HIGHLIGHT_BLUE; r1_2.font.bold = True

    # 6. Right Side Content - Item 2
    # Icon 2: Hierarchy
    box_w, box_h = 0.16, 0.12
    b1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.15), Inches(4.04), Inches(box_w), Inches(box_h))
    b2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.45), Inches(3.84), Inches(box_w), Inches(box_h))
    b3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.45), Inches(4.04), Inches(box_w), Inches(box_h))
    b4 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.45), Inches(4.24), Inches(box_w), Inches(box_h))
    
    for b in [b1, b2, b3, b4]:
        b.fill.background()
        b.line.color.rgb = DARK_BLUE
        b.line.width = Pt(1.5)

    v_line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(8.38), Inches(3.9), Inches(8.38), Inches(4.3))
    h1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(8.31), Inches(4.1), Inches(8.38), Inches(4.1))
    h2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(8.38), Inches(3.9), Inches(8.45), Inches(3.9))
    h3 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(8.38), Inches(4.1), Inches(8.45), Inches(4.1))
    h4 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(8.38), Inches(4.3), Inches(8.45), Inches(4.3))
    
    for l in [v_line, h1, h2, h3, h4]:
        l.line.color.rgb = DARK_BLUE
        l.line.width = Pt(1.5)

    # Text 2
    tx_box2 = slide.shapes.add_textbox(Inches(8.8), Inches(3.75), Inches(4.0), Inches(1.0))
    tf2 = tx_box2.text_frame
    p2 = tf2.paragraphs[0]
    p2.font.size = Pt(16)
    p2.font.name = "Microsoft YaHei"
    r2_1 = p2.add_run(); r2_1.text = "复杂问题"; r2_1.font.color.rgb = TEXT_BLACK
    r2_2 = p2.add_run(); r2_2.text = "拆解化"; r2_2.font.color.rgb = HIGHLIGHT_BLUE; r2_2.font.bold = True
    r2_3 = p2.add_run(); r2_3.text = "，"; r2_3.font.color.rgb = TEXT_BLACK
    r2_4 = p2.add_run(); r2_4.text = "单一\n"; r2_4.font.color.rgb = HIGHLIGHT_BLUE; r2_4.font.bold = True
    r2_5 = p2.add_run(); r2_5.text = "观点"; r2_5.font.color.rgb = HIGHLIGHT_BLUE; r2_5.font.bold = True
    r2_6 = p2.add_run(); r2_6.text = "深度化"; r2_6.font.color.rgb = TEXT_BLACK

    # 7. Right Side Content - Item 3
    # Icon 3: Pyramid/Upload
    up_arrow = slide.shapes.add_shape(MSO_SHAPE.UP_ARROW, Inches(8.3), Inches(5.35), Inches(0.18), Inches(0.22))
    up_arrow.fill.background()
    up_arrow.line.color.rgb = DARK_BLUE
    up_arrow.line.width = Pt(1.5)
    
    base = slide.shapes.add_shape(MSO_SHAPE.TRAPEZOID, Inches(8.15), Inches(5.62), Inches(0.48), Inches(0.15))
    base.fill.background()
    base.line.color.rgb = DARK_BLUE
    base.line.width = Pt(1.5)

    # Text 3
    tx_box3 = slide.shapes.add_textbox(Inches(8.8), Inches(5.25), Inches(4.0), Inches(1.0))
    tf3 = tx_box3.text_frame
    p3 = tf3.paragraphs[0]
    p3.font.size = Pt(16)
    p3.font.name = "Microsoft YaHei"
    r3_1 = p3.add_run(); r3_1.text = "结论先行"; r3_1.font.color.rgb = HIGHLIGHT_BLUE; r3_1.font.bold = True
    r3_2 = p3.add_run(); r3_2.text = "：标题即观点，\n内容即支撑"; r3_2.font.color.rgb = TEXT_BLACK

    # 8. Page Number
    page_num = slide.shapes.add_textbox(Inches(12.2), Inches(6.8), Inches(0.8), Inches(0.4))
    tf_page = page_num.text_frame
    p_page = tf_page.paragraphs[0]
    p_page.text = "3 / 11"
    p_page.font.size = Pt(12)
    p_page.font.color.rgb = TEXT_GRAY
    p_page.font.name = "Microsoft YaHei"
    p_page.alignment = PP_ALIGN.RIGHT



# ── Main ──
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

s0 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_1(s0)
prs.save(OUTPUT_PATH)
