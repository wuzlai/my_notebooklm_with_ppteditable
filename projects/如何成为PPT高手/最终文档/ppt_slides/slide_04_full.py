from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# ── Color Palette ──
BLUE_HEADER = RGBColor(0x5B, 0x9B, 0xD5)
BLUE_DARK   = RGBColor(0x4A, 0x86, 0xC8)
CYAN        = RGBColor(0x00, 0xBC, 0xD4)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x33, 0x33, 0x33)
GRAY_TEXT   = RGBColor(0x55, 0x55, 0x55)
GRAY_BAR    = RGBColor(0xB0, 0xBE, 0xC5)
RED         = RGBColor(0xE5, 0x39, 0x35)
GREEN       = RGBColor(0x43, 0xA0, 0x47)
ORANGE      = RGBColor(0xFF, 0x98, 0x00)

ICON_BG  = RGBColor(0xE3, 0xE8, 0xED)
ICON_FG  = RGBColor(0x54, 0x6E, 0x7A)
FONT_NAME = "Microsoft YaHei"
SLIDE_WIDTH = Inches(13.333)
HEADER_H    = Inches(0.75)
SUBTITLE_Y  = Inches(0.95)


def add_header_banner(slide, title_text, bg_color=None):
    if bg_color is None:
        bg_color = BLUE_HEADER
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, HEADER_H
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = bg_color
    banner.line.fill.background()
    tf = banner.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(26)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME


def add_subtitle(slide, text, left, top, width=Inches(12), font_size=Pt(18)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.font.name = FONT_NAME
    return txBox


def add_icon_box(slide, left, top, symbol, size=Inches(0.48)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ICON_BG
    shape.line.fill.background()
    shape.adjustments[0] = 0.25
    tf = shape.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = symbol
    p.font.size = Pt(18)
    p.font.color.rgb = ICON_FG
    p.font.bold = False
    return shape


def add_bullet_item(slide, left, top, symbol, label, description,
                    width=Inches(5.5), desc_size=Pt(13)):
    add_icon_box(slide, left, top, symbol)
    text_left = left + Inches(0.65)
    txBox = slide.shapes.add_textbox(text_left, top - Inches(0.02), width, Inches(0.65))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run_label = p.add_run()
    run_label.text = label + "\uff1a"
    run_label.font.size = Pt(14)
    run_label.font.color.rgb = BLACK
    run_label.font.bold = True
    run_label.font.name = FONT_NAME
    run_desc = p.add_run()
    run_desc.text = description
    run_desc.font.size = desc_size
    run_desc.font.color.rgb = GRAY_TEXT
    run_desc.font.bold = False
    run_desc.font.name = FONT_NAME
    return txBox


def add_conclusion_box(slide, left, top, width, text, font_size=Pt(13)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = BLACK
    run.font.bold = True
    run.font.name = FONT_NAME
    return txBox


def add_table(slide, left, top, width, height, rows, cols, data,
              header_color=None, col_widths=None):
    if header_color is None:
        header_color = BLUE_HEADER
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if data[r][c] is not None else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(11)
                    run.font.name = FONT_NAME
                    if r == 0:
                        run.font.color.rgb = WHITE
                        run.font.bold = True
                    else:
                        run.font.color.rgb = BLACK
                        run.font.bold = False
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else RGBColor(0xF5, 0xF5, 0xF5)
            cell.margin_left = Pt(5)
            cell.margin_right = Pt(5)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
    return table_shape


def add_bar_chart(slide, left, top, width, height,
                  categories, values, title="", bar_colors=None):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', values)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    chart.chart_style = 2
    plot = chart.plots[0]
    plot.gap_width = 100
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CYAN
    series.has_data_labels = True
    dl = series.data_labels
    dl.font.size = Pt(13)
    dl.font.bold = True
    dl.font.color.rgb = BLACK
    dl.number_format = '0.#'
    dl.show_value = True
    dl.label_position = XL_LABEL_POSITION.OUTSIDE_END
    if bar_colors:
        for i, color in enumerate(bar_colors):
            pt = series.points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = color
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(12)
    cat_axis.tick_labels.font.name = FONT_NAME
    cat_axis.major_tick_mark = 2
    cat_axis.format.line.fill.background()
    val_axis = chart.value_axis
    val_axis.visible = False
    val_axis.major_tick_mark = 2
    val_axis.format.line.fill.background()
    val_axis.major_gridlines.format.line.fill.background()
    if title:
        chart.has_title = True
        ct = chart.chart_title.text_frame.paragraphs[0]
        ct.text = title
        ct.font.size = Pt(14)
        ct.font.bold = True
        ct.font.name = FONT_NAME
    else:
        chart.has_title = False
    return chart_frame


def add_callout_label(slide, left, top, text, bg_color=None, font_size=Pt(11)):
    if bg_color is None:
        bg_color = CYAN
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.3), Inches(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME
    return shape


def add_data_card(slide, left, top, width, height, value, label,
                  value_color=None, bg_color=None):
    if value_color is None:
        value_color = CYAN
    if bg_color is None:
        bg_color = WHITE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(3)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run1 = p1.add_run()
    run1.text = str(value)
    run1.font.size = Pt(24)
    run1.font.color.rgb = value_color
    run1.font.bold = True
    run1.font.name = FONT_NAME
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(10)
    run2.font.color.rgb = GRAY_TEXT
    run2.font.bold = False
    run2.font.name = FONT_NAME
    return shape


OUTPUT_PATH = r"C:\Users\Administrator\Desktop\notebook\my_notebooklm_with_ppteditable\projects\如何成为PPT高手\最终文档\ppt_slides\slide_04.pptx"


# ── Slide 1 ──

def build_slide_1(slide):
    # 自定义颜色
    LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
    DENSE_TEXT_COLOR = RGBColor(0x99, 0x99, 0x99)
    
    # 1. 标题
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(10), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "减法艺术：拒绝文字堆砌"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK
    p.font.name = "Microsoft YaHei"

    # 2. 副标题
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(10), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "别让你的PPT变成Word搬家"
    p.font.size = Pt(22)
    p.font.color.rgb = GRAY_TEXT
    p.font.name = "Microsoft YaHei"

    # 3. 左侧栏 (错误示例)
    # 标题
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(5.5), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Word搬家（错误示例）"
    p.font.size = Pt(18)
    p.font.color.rgb = RED
    p.font.name = "Microsoft YaHei"
    
    # 红色下划线
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.5), Inches(2.7), Inches(6.0), Inches(2.7))
    line.line.color.rgb = RED
    line.line.width = Pt(1.5)

    # 密集文本块
    dense_text = "这里是一段非常长且密集的文字，用来模拟将Word文档直接复制粘贴到PPT中的错误做法。在实际的演示中，观众根本无法在短时间内阅读并理解这么多文字。这种做法不仅会让幻灯片显得杂乱无章，还会严重分散观众的注意力，导致他们无法专心听讲。优秀的PPT应该只保留核心观点和关键词，通过演讲者的口述来补充细节。如果把所有内容都写在屏幕上，那么演讲者就失去了存在的意义，PPT也就变成了一份阅读材料而不是辅助演示的工具。因此，我们必须学会做减法，拒绝文字堆砌，提炼出最精炼的信息，用视觉化的方式呈现出来，从而提高沟通的效率和效果。" * 4
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.9), Inches(5.5), Inches(3.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = dense_text
    p.font.size = Pt(8)
    p.font.color.rgb = DENSE_TEXT_COLOR
    p.font.name = "Microsoft YaHei"
    p.alignment = PP_ALIGN.JUSTIFY

    # 红色大叉
    cross_center_x = 3.25
    cross_center_y = 4.6
    cross_size = 1.2
    line1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(cross_center_x - cross_size), Inches(cross_center_y - cross_size), Inches(cross_center_x + cross_size), Inches(cross_center_y + cross_size))
    line1.line.color.rgb = RED
    line1.line.width = Pt(25)
    line2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(cross_center_x - cross_size), Inches(cross_center_y + cross_size), Inches(cross_center_x + cross_size), Inches(cross_center_y - cross_size))
    line2.line.color.rgb = RED
    line2.line.width = Pt(25)

    # 左侧底部结论
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(6.0), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "密密麻麻的文字，信息过载，观众无法聚焦。"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLACK
    p.font.name = "Microsoft YaHei"

    # 4. 中间垂直分割线
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(6.66), Inches(2.2), Inches(6.66), Inches(6.8))
    line.line.color.rgb = LIGHT_GRAY
    line.line.width = Pt(1)

    # 5. 右侧栏 (成功示例)
    # 标题
    txBox = slide.shapes.add_textbox(Inches(7.2), Inches(2.2), Inches(5.5), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "极简要点（成功示例）"
    p.font.size = Pt(18)
    p.font.color.rgb = GREEN
    p.font.name = "Microsoft YaHei"
    
    # 绿色下划线
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(7.2), Inches(2.7), Inches(12.8), Inches(2.7))
    line.line.color.rgb = GREEN
    line.line.width = Pt(1.5)

    # 绘制要点条目的内部函数
    def draw_bullet(slide, left, top, icon_text, label, desc):
        # 图标
        txBox = slide.shapes.add_textbox(left, top, Inches(0.6), Inches(0.6))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = icon_text
        p.font.size = Pt(28)
        p.font.name = "Segoe UI Emoji"
        
        # 标签
        txBox = slide.shapes.add_textbox(left + Inches(0.8), top, Inches(4.5), Inches(0.35))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = BLACK
        p.font.name = "Microsoft YaHei"
        
        # 描述
        txBox = slide.shapes.add_textbox(left + Inches(0.8), top + Inches(0.35), Inches(4.5), Inches(0.35))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = BLACK
        p.font.name = "Microsoft YaHei"

    # 添加三个要点
    draw_bullet(slide, Inches(7.2), Inches(3.1), "👂", "专注聆听", "观众阅读文字时，无法同时听取演讲。")
    draw_bullet(slide, Inches(7.2), Inches(4.3), "💎", "提炼金句", "删除冗余的修饰词，只保留核心观点。")
    draw_bullet(slide, Inches(7.2), Inches(5.5), "🖼️", "视觉替代", "用视觉元素（图标/图片）替代长篇大论。")

    # 绿色大勾
    check_x = 9.8
    check_y = 4.8
    line1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(check_x - 0.6), Inches(check_y - 0.1), Inches(check_x), Inches(check_y + 0.5))
    line1.line.color.rgb = GREEN
    line1.line.width = Pt(25)
    line2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(check_x), Inches(check_y + 0.5), Inches(check_x + 1.2), Inches(check_y - 1.0))
    line2.line.color.rgb = GREEN
    line2.line.width = Pt(25)

    # 右侧底部结论
    txBox = slide.shapes.add_textbox(Inches(7.2), Inches(6.5), Inches(6.0), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "精简内容，视觉引导，提升传递效率。"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.font.name = "Microsoft YaHei"

    # 6. 页码
    txBox = slide.shapes.add_textbox(Inches(11.8), Inches(6.8), Inches(1.0), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "4 / 11"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK
    p.font.name = "Microsoft YaHei"
    p.alignment = PP_ALIGN.RIGHT



# ── Main ──
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

s0 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_1(s0)
prs.save(OUTPUT_PATH)
