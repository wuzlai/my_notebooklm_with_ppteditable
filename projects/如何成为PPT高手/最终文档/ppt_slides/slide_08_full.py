from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# ── Color Palette ──
BLUE_HEADER = RGBColor(0x5B, 0x9B, 0xD5)
BLUE_DARK   = RGBColor(0x4A, 0x86, 0xC8)
CYAN        = RGBColor(0x00, 0xBC, 0xD4)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x33, 0x33, 0x33)
GRAY_TEXT   = RGBColor(0x55, 0x55, 0x55)
GRAY_BAR    = RGBColor(0xB0, 0xBE, 0xC5)
RED         = RGBColor(0xE5, 0x39, 0x35)
GREEN       = RGBColor(0x43, 0xA0, 0x47)
ORANGE      = RGBColor(0xFF, 0x98, 0x00)

ICON_BG  = RGBColor(0xE3, 0xE8, 0xED)
ICON_FG  = RGBColor(0x54, 0x6E, 0x7A)
FONT_NAME = "Microsoft YaHei"
SLIDE_WIDTH = Inches(13.333)
HEADER_H    = Inches(0.75)
SUBTITLE_Y  = Inches(0.95)


def add_header_banner(slide, title_text, bg_color=None):
    if bg_color is None:
        bg_color = BLUE_HEADER
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, HEADER_H
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = bg_color
    banner.line.fill.background()
    tf = banner.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(26)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME


def add_subtitle(slide, text, left, top, width=Inches(12), font_size=Pt(18)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.font.name = FONT_NAME
    return txBox


def add_icon_box(slide, left, top, symbol, size=Inches(0.48)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ICON_BG
    shape.line.fill.background()
    shape.adjustments[0] = 0.25
    tf = shape.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = symbol
    p.font.size = Pt(18)
    p.font.color.rgb = ICON_FG
    p.font.bold = False
    return shape


def add_bullet_item(slide, left, top, symbol, label, description,
                    width=Inches(5.5), desc_size=Pt(13)):
    add_icon_box(slide, left, top, symbol)
    text_left = left + Inches(0.65)
    txBox = slide.shapes.add_textbox(text_left, top - Inches(0.02), width, Inches(0.65))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run_label = p.add_run()
    run_label.text = label + "\uff1a"
    run_label.font.size = Pt(14)
    run_label.font.color.rgb = BLACK
    run_label.font.bold = True
    run_label.font.name = FONT_NAME
    run_desc = p.add_run()
    run_desc.text = description
    run_desc.font.size = desc_size
    run_desc.font.color.rgb = GRAY_TEXT
    run_desc.font.bold = False
    run_desc.font.name = FONT_NAME
    return txBox


def add_conclusion_box(slide, left, top, width, text, font_size=Pt(13)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = BLACK
    run.font.bold = True
    run.font.name = FONT_NAME
    return txBox


def add_table(slide, left, top, width, height, rows, cols, data,
              header_color=None, col_widths=None):
    if header_color is None:
        header_color = BLUE_HEADER
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if data[r][c] is not None else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(11)
                    run.font.name = FONT_NAME
                    if r == 0:
                        run.font.color.rgb = WHITE
                        run.font.bold = True
                    else:
                        run.font.color.rgb = BLACK
                        run.font.bold = False
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else RGBColor(0xF5, 0xF5, 0xF5)
            cell.margin_left = Pt(5)
            cell.margin_right = Pt(5)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
    return table_shape


def add_bar_chart(slide, left, top, width, height,
                  categories, values, title="", bar_colors=None):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', values)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    chart.chart_style = 2
    plot = chart.plots[0]
    plot.gap_width = 100
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CYAN
    series.has_data_labels = True
    dl = series.data_labels
    dl.font.size = Pt(13)
    dl.font.bold = True
    dl.font.color.rgb = BLACK
    dl.number_format = '0.#'
    dl.show_value = True
    dl.label_position = XL_LABEL_POSITION.OUTSIDE_END
    if bar_colors:
        for i, color in enumerate(bar_colors):
            pt = series.points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = color
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(12)
    cat_axis.tick_labels.font.name = FONT_NAME
    cat_axis.major_tick_mark = 2
    cat_axis.format.line.fill.background()
    val_axis = chart.value_axis
    val_axis.visible = False
    val_axis.major_tick_mark = 2
    val_axis.format.line.fill.background()
    val_axis.major_gridlines.format.line.fill.background()
    if title:
        chart.has_title = True
        ct = chart.chart_title.text_frame.paragraphs[0]
        ct.text = title
        ct.font.size = Pt(14)
        ct.font.bold = True
        ct.font.name = FONT_NAME
    else:
        chart.has_title = False
    return chart_frame


def add_callout_label(slide, left, top, text, bg_color=None, font_size=Pt(11)):
    if bg_color is None:
        bg_color = CYAN
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.3), Inches(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME
    return shape


def add_data_card(slide, left, top, width, height, value, label,
                  value_color=None, bg_color=None):
    if value_color is None:
        value_color = CYAN
    if bg_color is None:
        bg_color = WHITE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(3)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run1 = p1.add_run()
    run1.text = str(value)
    run1.font.size = Pt(24)
    run1.font.color.rgb = value_color
    run1.font.bold = True
    run1.font.name = FONT_NAME
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(10)
    run2.font.color.rgb = GRAY_TEXT
    run2.font.bold = False
    run2.font.name = FONT_NAME
    return shape


OUTPUT_PATH = r"C:\Users\Administrator\Desktop\notebook\my_notebooklm_with_ppteditable\projects\如何成为PPT高手\最终文档\ppt_slides\slide_08.pptx"


# ── Slide 1 ──

def build_slide_1(slide):
    # Colors
    BLUE_TITLE = RGBColor(0x1A, 0x66, 0xCC)
    DARK_TEXT = RGBColor(0x22, 0x22, 0x22)
    GRAY_TEXT = RGBColor(0x66, 0x66, 0x66)
    BLUE_LINE = RGBColor(0x5B, 0x9B, 0xD5)
    LIGHT_BLUE_FILL = RGBColor(0xE6, 0xF0, 0xFA)
    GREEN_OK = RGBColor(0x4C, 0xAF, 0x50)
    RED_ERR = RGBColor(0xF4, 0x43, 0x36)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    
    # 1. Title and Subtitle
    tb_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(8.0), Inches(0.8))
    p_title = tb_title.text_frame.paragraphs[0]
    
    run1 = p_title.add_run()
    run1.text = "排版逻辑："
    run1.font.size = Pt(32)
    run1.font.bold = True
    run1.font.color.rgb = BLUE_TITLE
    run1.font.name = "Microsoft YaHei"
    
    run2 = p_title.add_run()
    run2.text = "始终如一的风格表达"
    run2.font.size = Pt(32)
    run2.font.bold = True
    run2.font.color.rgb = DARK_TEXT
    run2.font.name = "Microsoft YaHei"
    
    tb_sub = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.0), Inches(0.5))
    p_sub = tb_sub.text_frame.paragraphs[0]
    p_sub.text = "规范化的布局让阅读更顺畅"
    p_sub.font.size = Pt(18)
    p_sub.font.color.rgb = GRAY_TEXT
    p_sub.font.name = "Microsoft YaHei"

    # 2. Left Section: Alignment & Margins
    # Icon (Grid)
    icon_grid = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.2), Inches(0.6), Inches(0.45))
    icon_grid.fill.solid()
    icon_grid.fill.fore_color.rgb = LIGHT_BLUE_FILL
    icon_grid.line.color.rgb = BLUE_LINE
    icon_grid.line.width = Pt(1.5)
    
    line_v = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.95), Inches(2.1), Inches(0.95), Inches(2.7))
    line_v.line.color.rgb = BLUE_LINE
    line_h = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.7), Inches(2.4), Inches(1.5), Inches(2.4))
    line_h.line.color.rgb = BLUE_LINE

    # Heading
    tb_h1 = slide.shapes.add_textbox(Inches(1.6), Inches(2.15), Inches(4.0), Inches(0.5))
    p_h1 = tb_h1.text_frame.paragraphs[0]
    p_h1.text = "建立统一的页边距与对齐基准"
    p_h1.font.size = Pt(18)
    p_h1.font.bold = True
    p_h1.font.name = "Microsoft YaHei"

    # Wireframe Graphic
    # Outer Box
    box_outer = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.7), Inches(3.1), Inches(3.3), Inches(2.2))
    box_outer.fill.solid()
    box_outer.fill.fore_color.rgb = LIGHT_BLUE_FILL
    box_outer.line.color.rgb = BLUE_LINE
    
    # Inner Dashed Box
    box_inner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.9), Inches(3.3), Inches(2.9), Inches(1.8))
    box_inner.fill.background()
    box_inner.line.color.rgb = BLUE_LINE
    box_inner.line.dash_style = 3 # Dashed
    
    # Content Blocks inside Wireframe
    rect_title = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.95), Inches(3.5), Inches(0.9), Inches(0.2))
    rect_title.fill.solid()
    rect_title.fill.fore_color.rgb = BLUE_TITLE
    rect_title.line.fill.background()
    
    for i in range(3):
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.95), Inches(3.8 + i*0.15), Inches(1.3), Inches(0.05))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(0xBD, 0xD0, 0xE6)
        line.line.fill.background()
        
    rect_img = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.4), Inches(3.5), Inches(1.35), Inches(0.9))
    rect_img.fill.background()
    rect_img.line.color.rgb = BLUE_LINE
    
    # Mountain placeholder inside rect_img
    tri1 = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(3.45), Inches(3.9), Inches(0.6), Inches(0.5))
    tri1.fill.solid()
    tri1.fill.fore_color.rgb = RGBColor(0xBD, 0xD0, 0xE6)
    tri1.line.fill.background()
    
    tri2 = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(3.85), Inches(4.0), Inches(0.5), Inches(0.4))
    tri2.fill.solid()
    tri2.fill.fore_color.rgb = RGBColor(0xBD, 0xD0, 0xE6)
    tri2.line.fill.background()
    
    sun = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4.15), Inches(3.65), Inches(0.15), Inches(0.15))
    sun.fill.solid()
    sun.fill.fore_color.rgb = RGBColor(0xBD, 0xD0, 0xE6)
    sun.line.fill.background()

    # Alignment Guides
    guide_v = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(3.35), Inches(2.9), Inches(3.35), Inches(5.5))
    guide_v.line.color.rgb = BLUE_TITLE
    guide_h = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.5), Inches(4.05), Inches(5.2), Inches(4.05))
    guide_h.line.color.rgb = BLUE_TITLE

    # Labels for Alignment
    lbl_align_l = slide.shapes.add_textbox(Inches(0.7), Inches(3.9), Inches(1.0), Inches(0.3))
    lbl_align_l.text_frame.text = "对齐基准"
    lbl_align_l.text_frame.paragraphs[0].font.size = Pt(12)
    lbl_align_l.text_frame.paragraphs[0].font.color.rgb = BLUE_TITLE

    lbl_align_b = slide.shapes.add_textbox(Inches(3.0), Inches(5.6), Inches(1.0), Inches(0.3))
    lbl_align_b.text_frame.text = "对齐基准"
    lbl_align_b.text_frame.paragraphs[0].font.size = Pt(12)
    lbl_align_b.text_frame.paragraphs[0].font.color.rgb = BLUE_TITLE

    lbl_margin = slide.shapes.add_textbox(Inches(5.2), Inches(3.9), Inches(1.0), Inches(0.3))
    lbl_margin.text_frame.text = "页边距"
    lbl_margin.text_frame.paragraphs[0].font.size = Pt(12)
    lbl_margin.text_frame.paragraphs[0].font.color.rgb = BLUE_TITLE

    # 3. Top Right Section: Icon Consistency
    # Icon (Four squares/circles)
    for r in range(2):
        for c in range(2):
            sq = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.6 + c*0.25), Inches(2.2 + r*0.25), Inches(0.2), Inches(0.2))
            sq.fill.solid()
            sq.fill.fore_color.rgb = LIGHT_BLUE_FILL
            sq.line.color.rgb = BLUE_LINE

    # Heading
    tb_h2 = slide.shapes.add_textbox(Inches(7.2), Inches(2.15), Inches(5.5), Inches(0.5))
    p_h2 = tb_h2.text_frame.paragraphs[0]
    p_h2.text = "保持图标风格一致（全线框或全色块）"
    p_h2.font.size = Pt(18)
    p_h2.font.bold = True
    p_h2.font.name = "Microsoft YaHei"

    # Correct Icons (Outline)
    icon_correct = slide.shapes.add_textbox(Inches(7.2), Inches(2.9), Inches(2.5), Inches(0.8))
    p_ic = icon_correct.text_frame.paragraphs[0]
    p_ic.text = "⚙  💡  📄"
    p_ic.font.size = Pt(36)
    p_ic.font.color.rgb = BLUE_TITLE
    
    # Correct Label
    chk_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.6), Inches(4.0), Inches(0.2), Inches(0.2))
    chk_circle.fill.solid()
    chk_circle.fill.fore_color.rgb = GREEN_OK
    chk_circle.line.fill.background()
    
    lbl_correct = slide.shapes.add_textbox(Inches(7.8), Inches(3.9), Inches(1.5), Inches(0.3))
    lbl_correct.text_frame.text = "正确（一致）"
    lbl_correct.text_frame.paragraphs[0].font.size = Pt(14)

    # Incorrect Icons (Mixed)
    icon_incorrect = slide.shapes.add_textbox(Inches(10.0), Inches(2.9), Inches(2.5), Inches(0.8))
    p_ii = icon_incorrect.text_frame.paragraphs[0]
    p_ii.text = "📢  ✋  ☁"
    p_ii.font.size = Pt(36)
    p_ii.font.color.rgb = BLUE_TITLE

    # Incorrect Label
    err_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.4), Inches(4.0), Inches(0.2), Inches(0.2))
    err_circle.fill.solid()
    err_circle.fill.fore_color.rgb = RED_ERR
    err_circle.line.fill.background()
    
    lbl_incorrect = slide.shapes.add_textbox(Inches(10.6), Inches(3.9), Inches(1.5), Inches(0.3))
    lbl_incorrect.text_frame.text = "错误（混杂）"
    lbl_incorrect.text_frame.paragraphs[0].font.size = Pt(14)

    # 4. Bottom Right Section: Whitespace
    # Icon (Document)
    icon_doc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.6), Inches(4.7), Inches(0.45), Inches(0.5))
    icon_doc.fill.solid()
    icon_doc.fill.fore_color.rgb = LIGHT_BLUE_FILL
    icon_doc.line.color.rgb = BLUE_LINE
    for i in range(3):
        dl = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.7), Inches(4.85 + i*0.1), Inches(0.25), Inches(0.03))
        dl.fill.solid()
        dl.fill.fore_color.rgb = BLUE_LINE
        dl.line.fill.background()

    # Heading
    tb_h3 = slide.shapes.add_textbox(Inches(7.2), Inches(4.75), Inches(5.0), Inches(0.5))
    p_h3 = tb_h3.text_frame.paragraphs[0]
    p_h3.text = "留白艺术：给内容呼吸的空间"
    p_h3.font.size = Pt(18)
    p_h3.font.bold = True
    p_h3.font.name = "Microsoft YaHei"

    # Good Layout Graphic (Whitespace)
    box_ws_outer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.3), Inches(5.5), Inches(2.6), Inches(1.8))
    box_ws_outer.fill.solid()
    box_ws_outer.fill.fore_color.rgb = RGBColor(0xF0, 0xF8, 0xFF)
    box_ws_outer.line.color.rgb = BLUE_LINE
    box_ws_outer.line.dash_style = 3
    
    # Diagonal lines for whitespace indication
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(7.3), Inches(5.5), Inches(7.7), Inches(5.9)).line.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(9.9), Inches(5.5), Inches(9.5), Inches(5.9)).line.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(7.3), Inches(7.3), Inches(7.7), Inches(6.9)).line.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(9.9), Inches(7.3), Inches(9.5), Inches(6.9)).line.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)

    box_ws_inner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.8), Inches(5.9), Inches(1.6), Inches(1.0))
    box_ws_inner.fill.solid()
    box_ws_inner.fill.fore_color.rgb = WHITE
    box_ws_inner.line.fill.background()
    
    # Shadow effect simulation
    box_ws_inner_shadow = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.78), Inches(5.88), Inches(1.64), Inches(1.04))
    box_ws_inner_shadow.fill.background()
    box_ws_inner_shadow.line.color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
    
    tb_core = slide.shapes.add_textbox(Inches(7.8), Inches(6.1), Inches(1.6), Inches(0.4))
    p_core = tb_core.text_frame.paragraphs[0]
    p_core.text = "核心内容"
    p_core.font.size = Pt(16)
    p_core.font.bold = True
    p_core.alignment = PP_ALIGN.CENTER
    
    core_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.2), Inches(6.55), Inches(0.8), Inches(0.05))
    core_line.fill.solid()
    core_line.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    core_line.line.fill.background()

    # Whitespace Labels
    lbl_ws_t = slide.shapes.add_textbox(Inches(8.3), Inches(5.55), Inches(0.8), Inches(0.2))
    lbl_ws_t.text_frame.text = "呼吸空间"
    lbl_ws_t.text_frame.paragraphs[0].font.size = Pt(9)
    lbl_ws_t.text_frame.paragraphs[0].font.color.rgb = BLUE_TITLE
    
    lbl_ws_b = slide.shapes.add_textbox(Inches(8.3), Inches(7.0), Inches(0.8), Inches(0.2))
    lbl_ws_b.text_frame.text = "呼吸空间"
    lbl_ws_b.text_frame.paragraphs[0].font.size = Pt(9)
    lbl_ws_b.text_frame.paragraphs[0].font.color.rgb = BLUE_TITLE

    lbl_ws_l = slide.shapes.add_textbox(Inches(7.35), Inches(6.3), Inches(0.8), Inches(0.2))
    lbl_ws_l.text_frame.text = "呼吸空间"
    lbl_ws_l.text_frame.paragraphs[0].font.size = Pt(9)
    lbl_ws_l.text_frame.paragraphs[0].font.color.rgb = BLUE_TITLE

    lbl_ws_r = slide.shapes.add_textbox(Inches(9.45), Inches(6.3), Inches(0.8), Inches(0.2))
    lbl_ws_r.text_frame.text = "呼吸空间"
    lbl_ws_r.text_frame.paragraphs[0].font.size = Pt(9)
    lbl_ws_r.text_frame.paragraphs[0].font.color.rgb = BLUE_TITLE

    # Bad Layout Graphic (Cluttered)
    box_cl_outer = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.4), Inches(5.5), Inches(1.5), Inches(1.0))
    box_cl_outer.fill.solid()
    box_cl_outer.fill.fore_color.rgb = LIGHT_BLUE_FILL
    box_cl_outer.line.color.rgb = BLUE_LINE

    # Cluttered inner elements
    cl_title = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10.45), Inches(5.55), Inches(0.6), Inches(0.15))
    cl_title.fill.solid()
    cl_title.fill.fore_color.rgb = BLUE_TITLE
    cl_title.line.fill.background()
    
    for i in range(4):
        cl_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10.45), Inches(5.75 + i*0.1), Inches(0.75), Inches(0.05))
        cl_line.fill.solid()
        cl_line.fill.fore_color.rgb = RGBColor(0xBD, 0xD0, 0xE6)
        cl_line.line.fill.background()
        
    cl_img = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(11.3), Inches(5.55), Inches(0.55), Inches(0.4))
    cl_img.fill.solid()
    cl_img.fill.fore_color.rgb = RGBColor(0xBD, 0xD0, 0xE6)
    cl_img.line.color.rgb = BLUE_LINE
    
    cl_box2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10.45), Inches(6.2), Inches(0.75), Inches(0.25))
    cl_box2.fill.solid()
    cl_box2.fill.fore_color.rgb = RGBColor(0xBD, 0xD0, 0xE6)
    cl_box2.line.color.rgb = BLUE_LINE

    cl_box3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(11.3), Inches(6.0), Inches(0.55), Inches(0.45))
    cl_box3.fill.solid()
    cl_box3.fill.fore_color.rgb = RGBColor(0xBD, 0xD0, 0xE6)
    cl_box3.line.color.rgb = BLUE_LINE

    # Cluttered Label
    err_circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.6), Inches(6.75), Inches(0.2), Inches(0.2))
    err_circle2.fill.solid()
    err_circle2.fill.fore_color.rgb = RED_ERR
    err_circle2.line.fill.background()
    
    lbl_cluttered = slide.shapes.add_textbox(Inches(10.8), Inches(6.65), Inches(1.5), Inches(0.3))
    lbl_cluttered.text_frame.text = "拥挤布局"
    lbl_cluttered.text_frame.paragraphs[0].font.size = Pt(14)
    lbl_cluttered.text_frame.paragraphs[0].font.bold = True

    # 5. Page Number
    tb_page = slide.shapes.add_textbox(Inches(11.8), Inches(6.8), Inches(1.2), Inches(0.5))
    p_page = tb_page.text_frame.paragraphs[0]
    p_page.text = "08 / 11"
    p_page.font.size = Pt(20)
    p_page.font.bold = True
    p_page.font.color.rgb = GRAY_TEXT
    p_page.font.name = "Microsoft YaHei"
    p_page.alignment = PP_ALIGN.RIGHT



# ── Main ──
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

s0 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_1(s0)
prs.save(OUTPUT_PATH)
