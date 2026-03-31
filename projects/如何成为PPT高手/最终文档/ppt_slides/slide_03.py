def build_slide(slide):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
    from pptx.enum.text import PP_ALIGN

    # Colors
    DARK_BLUE = RGBColor(0x1F, 0x4E, 0x96)
    HIGHLIGHT_BLUE = RGBColor(0x00, 0x70, 0xC0)
    TEXT_BLACK = RGBColor(0x33, 0x33, 0x33)
    TEXT_GRAY = RGBColor(0x66, 0x66, 0x66)
    ORANGE_FILL = RGBColor(0xFF, 0x6B, 0x00)
    ORANGE_LINE = RGBColor(0xE6, 0x51, 0x00)

    # 1. Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(10.0), Inches(0.8))
    tf_title = title_box.text_frame
    p_title = tf_title.paragraphs[0]
    p_title.text = "法则一：内容清晰是PPT的灵魂"
    p_title.font.name = "Microsoft YaHei"
    p_title.font.size = Pt(32)
    p_title.font.bold = True
    p_title.font.color.rgb = DARK_BLUE

    # 2. Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(10.0), Inches(0.5))
    tf_sub = sub_box.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "确保观众的注意力始终聚焦"
    p_sub.font.name = "Microsoft YaHei"
    p_sub.font.size = Pt(18)
    p_sub.font.color.rgb = TEXT_BLACK

    # 3. Central Graphic (Target)
    cx, cy = 5.2, 4.2
    radii = [2.2, 1.75, 1.3, 0.85, 0.4]
    fills = [
        RGBColor(226, 238, 252),
        RGBColor(204, 224, 250),
        RGBColor(182, 210, 248),
        RGBColor(160, 196, 246),
        ORANGE_FILL
    ]
    
    for i, (r, fill) in enumerate(zip(radii, fills)):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, 
            Inches(cx - r), Inches(cy - r), 
            Inches(r * 2), Inches(r * 2)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        if i < 4:
            shape.line.color.rgb = DARK_BLUE
            shape.line.width = Pt(2.5)
        else:
            shape.line.fill.background() # No line for bullseye

    # 4. Central Graphic (Arrow)
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.NOTCHED_RIGHT_ARROW, 
        Inches(cx + 0.4), Inches(cy - 2.4), 
        Inches(2.4), Inches(0.7)
    )
    arrow.rotation = 135
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RGBColor(255, 167, 38)
    arrow.line.color.rgb = ORANGE_LINE
    arrow.line.width = Pt(1.5)

    # 5. Right Side Content - Item 1
    # Icon 1: Magnifying Glass
    mag_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.2), Inches(2.5), Inches(0.28), Inches(0.28))
    mag_circle.fill.background()
    mag_circle.line.color.rgb = DARK_BLUE
    mag_circle.line.width = Pt(2)
    handle = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(8.44), Inches(2.74), Inches(8.6), Inches(2.9))
    handle.line.color.rgb = DARK_BLUE
    handle.line.width = Pt(2.5)

    # Text 1
    tx_box1 = slide.shapes.add_textbox(Inches(8.8), Inches(2.35), Inches(4.0), Inches(1.0))
    tf1 = tx_box1.text_frame
    p1 = tf1.paragraphs[0]
    p1.font.size = Pt(16)
    p1.font.name = "Microsoft YaHei"
    r1_1 = p1.add_run(); r1_1.text = "每一页幻灯片只传达一\n"; r1_1.font.color.rgb = TEXT_BLACK
    r1_2 = p1.add_run(); r1_2.text = "个核心观点"; r1_2.font.color.rgb = HIGHLIGHT_BLUE; r1_2.font.bold = True

    # 6. Right Side Content - Item 2
    # Icon 2: Hierarchy
    box_w, box_h = 0.16, 0.12
    b1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.15), Inches(4.04), Inches(box_w), Inches(box_h))
    b2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.45), Inches(3.84), Inches(box_w), Inches(box_h))
    b3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.45), Inches(4.04), Inches(box_w), Inches(box_h))
    b4 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.45), Inches(4.24), Inches(box_w), Inches(box_h))
    
    for b in [b1, b2, b3, b4]:
        b.fill.background()
        b.line.color.rgb = DARK_BLUE
        b.line.width = Pt(1.5)

    v_line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(8.38), Inches(3.9), Inches(8.38), Inches(4.3))
    h1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(8.31), Inches(4.1), Inches(8.38), Inches(4.1))
    h2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(8.38), Inches(3.9), Inches(8.45), Inches(3.9))
    h3 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(8.38), Inches(4.1), Inches(8.45), Inches(4.1))
    h4 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(8.38), Inches(4.3), Inches(8.45), Inches(4.3))
    
    for l in [v_line, h1, h2, h3, h4]:
        l.line.color.rgb = DARK_BLUE
        l.line.width = Pt(1.5)

    # Text 2
    tx_box2 = slide.shapes.add_textbox(Inches(8.8), Inches(3.75), Inches(4.0), Inches(1.0))
    tf2 = tx_box2.text_frame
    p2 = tf2.paragraphs[0]
    p2.font.size = Pt(16)
    p2.font.name = "Microsoft YaHei"
    r2_1 = p2.add_run(); r2_1.text = "复杂问题"; r2_1.font.color.rgb = TEXT_BLACK
    r2_2 = p2.add_run(); r2_2.text = "拆解化"; r2_2.font.color.rgb = HIGHLIGHT_BLUE; r2_2.font.bold = True
    r2_3 = p2.add_run(); r2_3.text = "，"; r2_3.font.color.rgb = TEXT_BLACK
    r2_4 = p2.add_run(); r2_4.text = "单一\n"; r2_4.font.color.rgb = HIGHLIGHT_BLUE; r2_4.font.bold = True
    r2_5 = p2.add_run(); r2_5.text = "观点"; r2_5.font.color.rgb = HIGHLIGHT_BLUE; r2_5.font.bold = True
    r2_6 = p2.add_run(); r2_6.text = "深度化"; r2_6.font.color.rgb = TEXT_BLACK

    # 7. Right Side Content - Item 3
    # Icon 3: Pyramid/Upload
    up_arrow = slide.shapes.add_shape(MSO_SHAPE.UP_ARROW, Inches(8.3), Inches(5.35), Inches(0.18), Inches(0.22))
    up_arrow.fill.background()
    up_arrow.line.color.rgb = DARK_BLUE
    up_arrow.line.width = Pt(1.5)
    
    base = slide.shapes.add_shape(MSO_SHAPE.TRAPEZOID, Inches(8.15), Inches(5.62), Inches(0.48), Inches(0.15))
    base.fill.background()
    base.line.color.rgb = DARK_BLUE
    base.line.width = Pt(1.5)

    # Text 3
    tx_box3 = slide.shapes.add_textbox(Inches(8.8), Inches(5.25), Inches(4.0), Inches(1.0))
    tf3 = tx_box3.text_frame
    p3 = tf3.paragraphs[0]
    p3.font.size = Pt(16)
    p3.font.name = "Microsoft YaHei"
    r3_1 = p3.add_run(); r3_1.text = "结论先行"; r3_1.font.color.rgb = HIGHLIGHT_BLUE; r3_1.font.bold = True
    r3_2 = p3.add_run(); r3_2.text = "：标题即观点，\n内容即支撑"; r3_2.font.color.rgb = TEXT_BLACK

    # 8. Page Number
    page_num = slide.shapes.add_textbox(Inches(12.2), Inches(6.8), Inches(0.8), Inches(0.4))
    tf_page = page_num.text_frame
    p_page = tf_page.paragraphs[0]
    p_page.text = "3 / 11"
    p_page.font.size = Pt(12)
    p_page.font.color.rgb = TEXT_GRAY
    p_page.font.name = "Microsoft YaHei"
    p_page.alignment = PP_ALIGN.RIGHT