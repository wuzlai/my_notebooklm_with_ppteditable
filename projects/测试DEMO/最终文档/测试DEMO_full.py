from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# ── Color Palette ──
BLUE_HEADER = RGBColor(0x5B, 0x9B, 0xD5)
BLUE_DARK   = RGBColor(0x4A, 0x86, 0xC8)
CYAN        = RGBColor(0x00, 0xBC, 0xD4)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x33, 0x33, 0x33)
GRAY_TEXT   = RGBColor(0x55, 0x55, 0x55)
GRAY_BAR    = RGBColor(0xB0, 0xBE, 0xC5)
RED         = RGBColor(0xE5, 0x39, 0x35)
GREEN       = RGBColor(0x43, 0xA0, 0x47)
ORANGE      = RGBColor(0xFF, 0x98, 0x00)

ICON_BG  = RGBColor(0xE3, 0xE8, 0xED)
ICON_FG  = RGBColor(0x54, 0x6E, 0x7A)
FONT_NAME = "Microsoft YaHei"
SLIDE_WIDTH = Inches(13.333)
HEADER_H    = Inches(0.75)
SUBTITLE_Y  = Inches(0.95)


def add_header_banner(slide, title_text, bg_color=None):
    if bg_color is None:
        bg_color = BLUE_HEADER
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, HEADER_H
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = bg_color
    banner.line.fill.background()
    tf = banner.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(26)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME


def add_subtitle(slide, text, left, top, width=Inches(12), font_size=Pt(18)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.font.name = FONT_NAME
    return txBox


def add_icon_box(slide, left, top, symbol, size=Inches(0.48)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ICON_BG
    shape.line.fill.background()
    shape.adjustments[0] = 0.25
    tf = shape.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = symbol
    p.font.size = Pt(18)
    p.font.color.rgb = ICON_FG
    p.font.bold = False
    return shape


def add_bullet_item(slide, left, top, symbol, label, description,
                    width=Inches(5.5), desc_size=Pt(13)):
    add_icon_box(slide, left, top, symbol)
    text_left = left + Inches(0.65)
    txBox = slide.shapes.add_textbox(text_left, top - Inches(0.02), width, Inches(0.65))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run_label = p.add_run()
    run_label.text = label + "\uff1a"
    run_label.font.size = Pt(14)
    run_label.font.color.rgb = BLACK
    run_label.font.bold = True
    run_label.font.name = FONT_NAME
    run_desc = p.add_run()
    run_desc.text = description
    run_desc.font.size = desc_size
    run_desc.font.color.rgb = GRAY_TEXT
    run_desc.font.bold = False
    run_desc.font.name = FONT_NAME
    return txBox


def add_conclusion_box(slide, left, top, width, text, font_size=Pt(13)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = BLACK
    run.font.bold = True
    run.font.name = FONT_NAME
    return txBox


def add_table(slide, left, top, width, height, rows, cols, data,
              header_color=None, col_widths=None):
    if header_color is None:
        header_color = BLUE_HEADER
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if data[r][c] is not None else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(11)
                    run.font.name = FONT_NAME
                    if r == 0:
                        run.font.color.rgb = WHITE
                        run.font.bold = True
                    else:
                        run.font.color.rgb = BLACK
                        run.font.bold = False
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else RGBColor(0xF5, 0xF5, 0xF5)
            cell.margin_left = Pt(5)
            cell.margin_right = Pt(5)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
    return table_shape


def add_bar_chart(slide, left, top, width, height,
                  categories, values, title="", bar_colors=None):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', values)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    chart.chart_style = 2
    plot = chart.plots[0]
    plot.gap_width = 100
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CYAN
    series.has_data_labels = True
    dl = series.data_labels
    dl.font.size = Pt(13)
    dl.font.bold = True
    dl.font.color.rgb = BLACK
    dl.number_format = '0.#'
    dl.show_value = True
    dl.label_position = XL_LABEL_POSITION.OUTSIDE_END
    if bar_colors:
        for i, color in enumerate(bar_colors):
            pt = series.points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = color
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(12)
    cat_axis.tick_labels.font.name = FONT_NAME
    cat_axis.major_tick_mark = 2
    cat_axis.format.line.fill.background()
    val_axis = chart.value_axis
    val_axis.visible = False
    val_axis.major_tick_mark = 2
    val_axis.format.line.fill.background()
    val_axis.major_gridlines.format.line.fill.background()
    if title:
        chart.has_title = True
        ct = chart.chart_title.text_frame.paragraphs[0]
        ct.text = title
        ct.font.size = Pt(14)
        ct.font.bold = True
        ct.font.name = FONT_NAME
    else:
        chart.has_title = False
    return chart_frame


def add_callout_label(slide, left, top, text, bg_color=None, font_size=Pt(11)):
    if bg_color is None:
        bg_color = CYAN
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.3), Inches(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME
    return shape


def add_data_card(slide, left, top, width, height, value, label,
                  value_color=None, bg_color=None):
    if value_color is None:
        value_color = CYAN
    if bg_color is None:
        bg_color = WHITE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(3)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run1 = p1.add_run()
    run1.text = str(value)
    run1.font.size = Pt(24)
    run1.font.color.rgb = value_color
    run1.font.bold = True
    run1.font.name = FONT_NAME
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(10)
    run2.font.color.rgb = GRAY_TEXT
    run2.font.bold = False
    run2.font.name = FONT_NAME
    return shape


OUTPUT_PATH = r"C:\Users\Administrator\Desktop\notebook\my_notebooklm_with_ppteditable\projects\测试DEMO\最终文档\测试DEMO.pptx"


# ── Slide 1 ──

def build_slide_1(slide):
    # Colors
    DARK_BLUE = RGBColor(0x1B, 0x2A, 0x49)
    GRAY_TEXT = RGBColor(0x66, 0x66, 0x66)
    LIGHT_GRAY = RGBColor(0x99, 0x99, 0x99)
    GREEN = RGBColor(0x2E, 0xA1, 0x54)
    RED = RGBColor(0xD9, 0x3A, 0x36)
    BLUE_ICON = RGBColor(0x4A, 0x86, 0xC8)
    BORDER_COLOR = RGBColor(0xE5, 0xE5, 0xE5)
    BLACK = RGBColor(0x00, 0x00, 0x00)

    # 1. Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(8.0), Inches(0.6))
    tf_title = title_box.text_frame
    p_title = tf_title.paragraphs[0]
    run_title = p_title.add_run()
    run_title.text = "简单场景验证：Claude Code 胜出"
    run_title.font.size = Pt(28)
    run_title.font.bold = True
    run_title.font.color.rgb = DARK_BLUE
    run_title.font.name = "Microsoft YaHei"

    # 2. Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(8.0), Inches(0.4))
    tf_sub = sub_box.text_frame
    p_sub = tf_sub.paragraphs[0]
    run_sub = p_sub.add_run()
    run_sub.text = "销售订单报表查询（低复杂度）效率对比"
    run_sub.font.size = Pt(16)
    run_sub.font.color.rgb = GRAY_TEXT
    run_sub.font.name = "Microsoft YaHei"

    # 3. Left Panel (White Box)
    left_panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.7), Inches(5.4), Inches(4.8))
    left_panel.fill.solid()
    left_panel.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    left_panel.line.color.rgb = BORDER_COLOR
    left_panel.line.width = Pt(1)

    # Helper for left panel text
    def add_left_text(y, runs):
        tb = slide.shapes.add_textbox(Inches(1.2), y, Inches(4.5), Inches(0.8))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.line_spacing = 1.3
        for text, is_bold, color in runs:
            r = p.add_run()
            r.text = text
            r.font.name = "Microsoft YaHei"
            r.font.size = Pt(13)
            if is_bold:
                r.font.bold = True
            if color:
                r.font.color.rgb = color

    # Item 1: Up Arrow Icon & Text
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.75), Inches(2.05), Inches(0.28), Inches(0.28))
    circle.fill.background()
    circle.line.color.rgb = GREEN
    circle.line.width = Pt(1.5)
    arrow = slide.shapes.add_shape(MSO_SHAPE.UP_ARROW, Inches(0.84), Inches(2.1), Inches(0.1), Inches(0.18))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = GREEN
    arrow.line.fill.background()
    
    add_left_text(Inches(1.95), [
        ("效率拐点：", True, BLACK),
        ("Claude Code 耗时 ", False, BLACK),
        ("30 分钟", True, BLACK),
        ("，较手写开发效率", False, BLACK),
        ("提升 50%", True, BLACK),
        ("。", False, BLACK)
    ])

    # Item 2: Warning Icon & Text
    triangle = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(0.75), Inches(3.05), Inches(0.28), Inches(0.28))
    triangle.fill.background()
    triangle.line.color.rgb = RED
    triangle.line.width = Pt(1.5)
    warn_tb = slide.shapes.add_textbox(Inches(0.75), Inches(3.1), Inches(0.28), Inches(0.28))
    warn_p = warn_tb.text_frame.paragraphs[0]
    warn_p.alignment = PP_ALIGN.CENTER
    warn_r = warn_p.add_run()
    warn_r.text = "!"
    warn_r.font.size = Pt(12)
    warn_r.font.bold = True
    warn_r.font.color.rgb = RED

    add_left_text(Inches(2.95), [
        ("稳定性差异：", True, BLACK),
        ("GitHub Copilot 运行出现 ", False, BLACK),
        ("Short Dump", True, RED),
        ("，而 Claude Code ", False, BLACK),
        ("运行正常", True, GREEN),
        ("。", False, BLACK)
    ])

    # Item 3: Chat Icon & Text
    chat1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGULAR_CALLOUT, Inches(0.75), Inches(4.1), Inches(0.22), Inches(0.18))
    chat1.fill.background()
    chat1.line.color.rgb = BLUE_ICON
    chat1.line.width = Pt(1.5)
    chat2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGULAR_CALLOUT, Inches(0.82), Inches(4.18), Inches(0.22), Inches(0.18))
    chat2.fill.background()
    chat2.line.color.rgb = BLUE_ICON
    chat2.line.width = Pt(1.5)

    add_left_text(Inches(3.95), [
        ("交互成本：", True, BLACK),
        ("Claude Code 仅需 ", False, BLACK),
        ("2 轮", True, BLACK),
        ("人工干预，远优于 Copilot 的 ", False, BLACK),
        ("5 轮", True, BLACK),
        ("以上。", False, BLACK)
    ])

    # Item 4: Link Icon & Text
    link1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.78), Inches(5.15), Inches(0.18), Inches(0.1))
    link1.rotation = 45
    link1.fill.background()
    link1.line.color.rgb = LIGHT_GRAY
    link1.line.width = Pt(1.5)
    link2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.86), Inches(5.23), Inches(0.18), Inches(0.1))
    link2.rotation = 45
    link2.fill.background()
    link2.line.color.rgb = LIGHT_GRAY
    link2.line.width = Pt(1.5)

    add_left_text(Inches(4.95), [
        ("核心瓶颈：", True, BLACK),
        ("两者初次生成均不可直接运行，仍需人工修正 ", False, BLACK),
        ("SQL", True, BLACK),
        (" 逻辑。", False, BLACK)
    ])

    # 4. Right Top Panel (Chart)
    chart_title = slide.shapes.add_textbox(Inches(6.4), Inches(1.7), Inches(4.0), Inches(0.4))
    p_ct = chart_title.text_frame.paragraphs[0]
    r_ct = p_ct.add_run()
    r_ct.text = "开发耗时对比（分钟）"
    r_ct.font.size = Pt(14)
    r_ct.font.bold = True
    r_ct.font.name = "Microsoft YaHei"

    chart_data = CategoryChartData()
    chart_data.categories = ['Claude Code', 'GitHub Copilot', '手写开发']
    chart_data.add_series('耗时', (30, 65, 60))

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, Inches(6.4), Inches(2.2), Inches(6.4), Inches(2.0), chart_data
    )
    chart = chart_shape.chart
    chart.has_legend = False

    val_axis = chart.value_axis
    val_axis.maximum_scale = 70
    val_axis.major_unit = 15
    val_axis.has_major_gridlines = True
    val_axis.major_gridlines.format.line.color.rgb = BORDER_COLOR
    val_axis.tick_labels.font.size = Pt(10)
    val_axis.tick_labels.font.color.rgb = GRAY_TEXT

    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(11)
    cat_axis.tick_labels.font.color.rgb = GRAY_TEXT
    cat_axis.has_major_gridlines = False

    series = chart.series[0]
    series.has_data_labels = True

    # Claude Code (Green)
    p0 = series.points[0]
    p0.format.fill.solid()
    p0.format.fill.fore_color.rgb = GREEN
    p0.data_label.has_text_frame = True
    p0.data_label.text_frame.text = "30 min ✅"
    p0.data_label.font.size = Pt(10)
    p0.data_label.font.bold = True

    # GitHub Copilot (Red)
    p1 = series.points[1]
    p1.format.fill.solid()
    p1.format.fill.fore_color.rgb = RED
    p1.data_label.has_text_frame = True
    p1.data_label.text_frame.text = "> 60 min ⚠️"
    p1.data_label.font.size = Pt(10)
    p1.data_label.font.bold = True
    p1.data_label.font.color.rgb = RED

    # 手写开发 (Red)
    p2 = series.points[2]
    p2.format.fill.solid()
    p2.format.fill.fore_color.rgb = RED
    p2.data_label.has_text_frame = True
    p2.data_label.text_frame.text = "60 min"
    p2.data_label.font.size = Pt(10)
    p2.data_label.font.bold = True

    # 5. Right Bottom Left Panel (Data Card)
    card1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.4), Inches(4.5), Inches(2.6), Inches(2.0))
    card1.fill.solid()
    card1.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    card1.line.color.rgb = BORDER_COLOR

    tb_50 = slide.shapes.add_textbox(Inches(6.5), Inches(4.8), Inches(2.0), Inches(0.6))
    p_50 = tb_50.text_frame.paragraphs[0]
    r_50 = p_50.add_run()
    r_50.text = "50%"
    r_50.font.size = Pt(40)
    r_50.font.bold = True
    r_50.font.color.rgb = GREEN

    tb_eff = slide.shapes.add_textbox(Inches(6.5), Inches(5.5), Inches(2.0), Inches(0.4))
    p_eff = tb_eff.text_frame.paragraphs[0]
    r_eff1 = p_eff.add_run()
    r_eff1.text = "效率提升 "
    r_eff1.font.size = Pt(16)
    r_eff1.font.bold = True
    r_eff1.font.name = "Microsoft YaHei"
    r_eff2 = p_eff.add_run()
    r_eff2.text = "↗"
    r_eff2.font.size = Pt(16)
    r_eff2.font.bold = True
    r_eff2.font.color.rgb = GREEN

    tb_vs = slide.shapes.add_textbox(Inches(6.5), Inches(6.0), Inches(2.4), Inches(0.3))
    p_vs = tb_vs.text_frame.paragraphs[0]
    r_vs = p_vs.add_run()
    r_vs.text = "Claude Code vs 手写开发"
    r_vs.font.size = Pt(10)
    r_vs.font.color.rgb = GRAY_TEXT
    r_vs.font.name = "Microsoft YaHei"

    # 6. Right Bottom Right Panel (Grid)
    card2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.3), Inches(4.5), Inches(3.5), Inches(2.0))
    card2.fill.solid()
    card2.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    card2.line.color.rgb = BORDER_COLOR

    line_v = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(11.05), Inches(4.5), Inches(11.05), Inches(6.5))
    line_v.line.color.rgb = BORDER_COLOR
    line_h = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(9.3), Inches(5.5), Inches(12.8), Inches(5.5))
    line_h.line.color.rgb = BORDER_COLOR

    def add_grid_cell(x, y, icon_text, icon_color, text, icon_size=24):
        tb = slide.shapes.add_textbox(x, y, Inches(1.75), Inches(1.0))
        tf = tb.text_frame
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = icon_text + "\n"
        r1.font.size = Pt(icon_size)
        r1.font.color.rgb = icon_color
        r1.font.name = "Segoe UI Emoji"

        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = text
        r2.font.size = Pt(11)
        r2.font.color.rgb = BLACK
        r2.font.name = "Microsoft YaHei"

    add_grid_cell(Inches(9.3), Inches(4.6), "❌", RED, "GitHub Copilot")
    add_grid_cell(Inches(11.05), Inches(4.6), "❌", RED, "Short Dump")
    add_grid_cell(Inches(9.3), Inches(5.6), "✅", GREEN, "Copilot: 5+ 轮")
    add_grid_cell(Inches(11.05), Inches(5.6), "👥", BLUE_ICON, "Claude Code: 2 轮")

    # 7. Page Number
    page_num = slide.shapes.add_textbox(Inches(11.8), Inches(6.9), Inches(1.0), Inches(0.3))
    p_page = page_num.text_frame.paragraphs[0]
    p_page.alignment = PP_ALIGN.RIGHT
    r_page = p_page.add_run()
    r_page.text = "Page 1 / 4"
    r_page.font.size = Pt(10)
    r_page.font.color.rgb = LIGHT_GRAY
    r_page.font.name = "Microsoft YaHei"



# ── Slide 2 ──

def build_slide_2(slide):
    # Colors
    DARK_BLUE = RGBColor(0x1A, 0x36, 0x5D)
    RED = RGBColor(0xE5, 0x39, 0x35)
    LIGHT_RED = RGBColor(0xFD, 0xED, 0xEC)
    LIGHT_GREEN = RGBColor(0xE8, 0xF5, 0xE9)
    GRAY_TEXT = RGBColor(0x55, 0x55, 0x55)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xE0)
    
    # 1. Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(10), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "中等场景瓶颈：AI “幻觉”引发崩溃"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = "Microsoft YaHei"

    # 2. Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(10), Inches(0.5))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "采购配额维护（中等复杂度）开发验证"
    p.font.size = Pt(18)
    p.font.color.rgb = GRAY_TEXT
    p.font.name = "Microsoft YaHei"

    # 3. Left Bullets
    bullets_data = [
        ("🔗", "逻辑误区", "Copilot 完全混淆业务概念（配额误作货源），代码完全不可用。"),
        ("⚠️", "幻觉严重", "Claude Code 虚构字段比例高达 50%，导致 21 个连锁语法错误。"),
        ("🚫", "效率归零", "AI 修复成本大于重写成本，整体效率对比手写无任何提升。"),
        ("⚠️", "规则缺失", "AI 无法准确遵循 SAP 函数接口规范，直接忽略“优先 API”的指令。")
    ]
    
    start_y = 1.8
    for icon, label, desc in bullets_data:
        # Icon
        icon_box = slide.shapes.add_textbox(Inches(0.5), Inches(start_y), Inches(0.6), Inches(0.6))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(28)
        p.font.color.rgb = RED
        p.alignment = PP_ALIGN.CENTER
        
        # Label
        label_box = slide.shapes.add_textbox(Inches(1.3), Inches(start_y), Inches(4.5), Inches(0.4))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.name = "Microsoft YaHei"
        
        # Desc
        desc_box = slide.shapes.add_textbox(Inches(1.3), Inches(start_y + 0.35), Inches(4.8), Inches(0.8))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY_TEXT
        p.font.name = "Microsoft YaHei"
        
        start_y += 1.3

    # 4. Right Data Card
    card_left = Inches(6.8)
    card_top = Inches(1.4)
    card_width = Inches(6.0)
    card_height = Inches(1.4)
    
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_left, card_top, card_width, card_height)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = LIGHT_GRAY
    card.line.width = Pt(1)

    # Card Text - 50%
    val_box = slide.shapes.add_textbox(card_left + Inches(0.2), card_top + Inches(0.1), Inches(3), Inches(0.8))
    tf = val_box.text_frame
    p = tf.paragraphs[0]
    p.text = "50% ↗"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RED
    p.font.name = "Arial"

    # Card Text - Label
    lbl_box = slide.shapes.add_textbox(card_left + Inches(0.2), card_top + Inches(0.9), Inches(4), Inches(0.4))
    tf = lbl_box.text_frame
    p = tf.paragraphs[0]
    p.text = "虚构字段比例 (Fictional Field Ratio)"
    p.font.size = Pt(12)
    p.font.color.rgb = GRAY_TEXT
    p.font.name = "Microsoft YaHei"

    # Card Icon - Warning
    warn_box = slide.shapes.add_textbox(card_left + Inches(4.8), card_top + Inches(0.2), Inches(1), Inches(1))
    tf = warn_box.text_frame
    p = tf.paragraphs[0]
    p.text = "⚠️"
    p.font.size = Pt(60)
    p.font.color.rgb = RED
    p.alignment = PP_ALIGN.CENTER

    # 5. Right Table
    table_top = Inches(3.0)
    
    # Table Title
    tt_box = slide.shapes.add_textbox(card_left, table_top, card_width, Inches(0.4))
    tf = tt_box.text_frame
    p = tf.paragraphs[0]
    p.text = "字段对比示例"
    p.font.size = Pt(14)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Microsoft YaHei"

    # Table Shape
    rows = 6
    cols = 2
    table_shape = slide.shapes.add_table(rows, cols, card_left, table_top + Inches(0.4), card_width, Inches(2.0)).table
    table_shape.columns[0].width = Inches(3.0)
    table_shape.columns[1].width = Inches(3.0)

    table_data = [
        ("❌ 虚构字段", "✅ 正确字段"),
        ("MSEG-QUOTA_ID (不存在)", "EKKO-EBELN (采购凭证)"),
        ("EKPO-ALLOC_QTY (错误逻辑)", "EKPO-MENGE (数量)"),
        ("EKPO-ALLOC_QTY (错误逻辑)", "EKPO-MENGE (数量)"),
        ("EKPO-BLG_ID (不存在)", "EKKO-QUOTA (数量)"),
        ("...", "...")
    ]

    for r in range(rows):
        for c in range(cols):
            cell = table_shape.cell(r, c)
            cell.text = table_data[r][c]
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.font.name = "Microsoft YaHei"
            
            if r == 0:
                p.font.bold = True
            
            # Set background colors
            if r == 0:
                if c == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = LIGHT_RED
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = LIGHT_GREEN
            else:
                if c == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xFE, 0xF5, 0xF5)
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xF1, 0xF8, 0xF1)

    # 6. Right Flowchart
    flow_top = Inches(5.7)
    
    # Box 1
    b1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_left, flow_top, Inches(1.2), Inches(1.2))
    b1.fill.solid()
    b1.fill.fore_color.rgb = RED
    b1.line.fill.background()
    tf = b1.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = "50%"
    p1.font.size = Pt(20)
    p1.font.color.rgb = WHITE
    p1.font.bold = True
    p1.alignment = PP_ALIGN.CENTER
    p1.font.name = "Arial"
    p2 = tf.add_paragraph()
    p2.text = "虚构字段"
    p2.font.size = Pt(12)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER
    p2.font.name = "Microsoft YaHei"

    # Arrow 1
    a1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, card_left + Inches(1.3), flow_top + Inches(0.5), Inches(0.3), Inches(0.2))
    a1.fill.solid()
    a1.fill.fore_color.rgb = GRAY_TEXT
    a1.line.fill.background()

    # Box 2
    b2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_left + Inches(1.7), flow_top, Inches(1.6), Inches(1.2))
    b2.fill.solid()
    b2.fill.fore_color.rgb = RED
    b2.line.fill.background()
    tf = b2.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = "🔗"
    p1.font.size = Pt(20)
    p1.font.color.rgb = WHITE
    p1.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "21 个连锁语法错误"
    p2.font.size = Pt(11)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER
    p2.font.name = "Microsoft YaHei"

    # Arrow 2
    a2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, card_left + Inches(3.4), flow_top + Inches(0.5), Inches(0.3), Inches(0.2))
    a2.fill.solid()
    a2.fill.fore_color.rgb = GRAY_TEXT
    a2.line.fill.background()

    # Box 3
    b3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_left + Inches(3.8), flow_top, Inches(2.2), Inches(1.2))
    b3.fill.solid()
    b3.fill.fore_color.rgb = RED
    b3.line.fill.background()
    tf = b3.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = "🚫"
    p1.font.size = Pt(20)
    p1.font.color.rgb = WHITE
    p1.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "AI 修复成本 >> 重写成本，\n整体效率对比手写无提升"
    p2.font.size = Pt(10)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER
    p2.font.name = "Microsoft YaHei"

    # 7. Footer
    footer = slide.shapes.add_textbox(Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = "PAGE 2 OF 4"
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY_TEXT
    p.alignment = PP_ALIGN.RIGHT



# ── Slide 3 ──

def build_slide_3(slide):
    from pptx.enum.shapes import MSO_CONNECTOR
    
    # 1. 添加背景卡片（可选，为了更好的视觉效果，这里添加一个全屏的浅色背景和白色卡片）
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)
    bg.line.fill.background()
    
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(0.4), Inches(12.533), Inches(6.7))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.fill.background()

    # 2. 标题和副标题
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(8), Inches(0.6))
    p = title_box.text_frame.paragraphs[0]
    p.text = "复杂场景失效：修复-爆炸模式"
    p.font.name = FONT_NAME
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1C, 0x2A, 0x43)

    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8), Inches(0.4))
    p = sub_box.text_frame.paragraphs[0]
    p.text = "跨工厂 STO 报表（高复杂度）深度测试"
    p.font.name = FONT_NAME
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY_TEXT

    # ==================== 左侧栏 ====================

    # 效率倒挂
    h1 = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(4), Inches(0.4))
    p = h1.text_frame.paragraphs[0]
    p.text = "效率倒挂"
    p.font.bold = True
    p.font.size = Pt(18)

    # 红色数据框
    red_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.3), Inches(5.2), Inches(0.9))
    red_bg.fill.solid()
    red_bg.fill.fore_color.rgb = RGBColor(0xFC, 0xE9, 0xE9)
    red_bg.line.fill.background()

    red_border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.3), Inches(0.08), Inches(0.9))
    red_border.fill.solid()
    red_border.fill.fore_color.rgb = RED
    red_border.line.fill.background()

    rb_text = slide.shapes.add_textbox(Inches(1.0), Inches(2.35), Inches(5.0), Inches(0.8))
    tf = rb_text.text_frame
    p1 = tf.paragraphs[0]
    p1.text = "Claude Code 耗时 8 人天"
    p1.font.size = Pt(22)
    p1.font.bold = True
    p1.font.color.rgb = RED
    p2 = tf.add_paragraph()
    p2.text = "比手写开发（5天）慢 60%"
    p2.font.size = Pt(12)
    p2.font.color.rgb = GRAY_TEXT

    # 修复陷阱
    h2 = slide.shapes.add_textbox(Inches(0.8), Inches(3.4), Inches(4), Inches(0.4))
    p = h2.text_frame.paragraphs[0]
    p.text = "修复陷阱"
    p.font.bold = True
    p.font.size = Pt(18)

    icon1 = slide.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(0.5), Inches(0.5))
    p = icon1.text_frame.paragraphs[0]
    p.text = "🔗"
    p.font.size = Pt(24)
    p.font.color.rgb = RED

    desc1 = slide.shapes.add_textbox(Inches(1.4), Inches(3.85), Inches(4.6), Inches(0.6))
    desc1.text_frame.word_wrap = True
    p = desc1.text_frame.paragraphs[0]
    p.text = "语法错误呈现“3→2→9”反向增长，陷入越修越错的逻辑矛盾。"
    p.font.size = Pt(13)

    # 字典缺失
    h3 = slide.shapes.add_textbox(Inches(0.8), Inches(4.6), Inches(4), Inches(0.4))
    p = h3.text_frame.paragraphs[0]
    p.text = "字典缺失"
    p.font.bold = True
    p.font.size = Pt(18)

    icon2 = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(0.9), Inches(5.1), Inches(0.3), Inches(0.3))
    icon2.fill.solid()
    icon2.fill.fore_color.rgb = RED
    icon2.line.fill.background()
    exc2 = slide.shapes.add_textbox(Inches(0.9), Inches(5.05), Inches(0.3), Inches(0.3))
    p = exc2.text_frame.paragraphs[0]
    p.text = "!"
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    desc2 = slide.shapes.add_textbox(Inches(1.4), Inches(5.05), Inches(4.6), Inches(0.6))
    desc2.text_frame.word_wrap = True
    p = desc2.text_frame.paragraphs[0]
    p.text = "累计虚构 10 余项数据字典组件，AI 无法理解复杂的 SAP 关联逻辑。"
    p.font.size = Pt(13)

    # 开发结论
    h4 = slide.shapes.add_textbox(Inches(0.8), Inches(5.8), Inches(4), Inches(0.4))
    p = h4.text_frame.paragraphs[0]
    p.text = "开发结论"
    p.font.bold = True
    p.font.size = Pt(18)

    desc3 = slide.shapes.add_textbox(Inches(0.8), Inches(6.2), Inches(5.2), Inches(0.6))
    desc3.text_frame.word_wrap = True
    p = desc3.text_frame.paragraphs[0]
    p.text = "•  在涉及 18+ 张表及多级状态追踪时，AI 框架生成能力失效。"
    p.font.size = Pt(13)

    # ==================== 右侧栏 ====================

    # 折线图
    chart_data = CategoryChartData()
    chart_data.categories = ['初始生成', '第一次修复', '第二次修复']
    chart_data.add_series('错误数量', (3, 2, 9))

    x, y, cx, cy = Inches(6.5), Inches(1.5), Inches(6.0), Inches(2.8)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS, x, y, cx, cy, chart_data
    ).chart

    chart.has_title = True
    chart.chart_title.text_frame.text = "修复-爆炸趋势折线图"
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)

    series = chart.series[0]
    series.format.line.color.rgb = RED
    series.format.line.width = Pt(3)
    series.marker.style = 8
    series.marker.size = 8
    series.marker.format.fill.solid()
    series.marker.format.fill.fore_color.rgb = RED
    series.marker.format.line.fill.background()

    series.has_data_labels = True
    for dl in series.data_labels:
        dl.font.size = Pt(12)
        dl.font.color.rgb = RED
        dl.position = XL_LABEL_POSITION.ABOVE

    value_axis = chart.value_axis
    value_axis.has_title = True
    value_axis.axis_title.text_frame.text = "错误数量"
    value_axis.maximum_scale = 10
    value_axis.minimum_scale = 0
    value_axis.major_unit = 2
    value_axis.has_major_gridlines = True
    value_axis.major_gridlines.format.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    chart.category_axis.has_major_gridlines = False

    # 图表标注箭头和文字
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(9.15), Inches(2.6), Inches(2.4), Inches(0.15))
    arrow.rotation = -45
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RED
    arrow.line.fill.background()

    anno_text = slide.shapes.add_textbox(Inches(8.5), Inches(2.0), Inches(2.5), Inches(0.4))
    p = anno_text.text_frame.paragraphs[0]
    p.text = "反向增长：越修越错"
    p.font.size = Pt(11)
    p.font.color.rgb = RED

    # ==================== 底部对比图 ====================

    # 左侧 AI生成 框
    box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(4.8), Inches(2.6), Inches(2.0))
    box1.fill.solid()
    box1.fill.fore_color.rgb = RGBColor(0xFC, 0xE9, 0xE9)
    box1.line.fill.background()

    header1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(4.8), Inches(2.6), Inches(0.4))
    header1.fill.solid()
    header1.fill.fore_color.rgb = RED
    header1.line.fill.background()
    h1_mask = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(5.0), Inches(2.6), Inches(0.2))
    h1_mask.fill.solid()
    h1_mask.fill.fore_color.rgb = RED
    h1_mask.line.fill.background()

    h1_text = slide.shapes.add_textbox(Inches(6.5), Inches(4.8), Inches(2.6), Inches(0.4))
    p = h1_text.text_frame.paragraphs[0]
    p.text = "AI生成"
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = WHITE
    p.font.bold = True

    # 断链图标 (绘制)
    l_link = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.3), Inches(5.3), Inches(0.6), Inches(0.3))
    l_link.rotation = 45
    l_link.fill.background()
    l_link.line.color.rgb = BLACK
    l_link.line.width = Pt(3)
    r_link = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.7), Inches(5.5), Inches(0.6), Inches(0.3))
    r_link.rotation = 45
    r_link.fill.background()
    r_link.line.color.rgb = BLACK
    r_link.line.width = Pt(3)
    b1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.5), Inches(5.4), Inches(0.2), Inches(0.5))
    b1.rotation = 45
    b1.fill.solid()
    b1.fill.fore_color.rgb = RGBColor(0xFC, 0xE9, 0xE9)
    b1.line.fill.background()

    # 警告文字
    warn_icon = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(6.7), Inches(6.3), Inches(0.2), Inches(0.2))
    warn_icon.fill.solid()
    warn_icon.fill.fore_color.rgb = RED
    warn_icon.line.fill.background()
    exc = slide.shapes.add_textbox(Inches(6.7), Inches(6.25), Inches(0.2), Inches(0.2))
    p = exc.text_frame.paragraphs[0]
    p.text = "!"
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    warn_text = slide.shapes.add_textbox(Inches(7.0), Inches(6.2), Inches(2.0), Inches(0.6))
    p = warn_text.text_frame.paragraphs[0]
    p.text = "18+ 张表关联失败\n10+ 幻觉组件"
    p.font.size = Pt(11)
    p.font.color.rgb = BLACK

    # 中间箭头
    mid_arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(9.2), Inches(5.6), Inches(0.8), Inches(0.4))
    mid_arrow.fill.solid()
    mid_arrow.fill.fore_color.rgb = RGBColor(0xCF, 0xD8, 0xDC)
    mid_arrow.line.fill.background()

    cross = slide.shapes.add_shape(MSO_SHAPE.MATH_MULTIPLY, Inches(9.4), Inches(5.6), Inches(0.4), Inches(0.4))
    cross.fill.solid()
    cross.fill.fore_color.rgb = RED
    cross.line.fill.background()

    # 右侧 预期目标 框
    box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.1), Inches(4.8), Inches(2.6), Inches(2.0))
    box2.fill.solid()
    box2.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
    box2.line.fill.background()

    header2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.1), Inches(4.8), Inches(2.6), Inches(0.4))
    header2.fill.solid()
    header2.fill.fore_color.rgb = GREEN
    header2.line.fill.background()
    h2_mask = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10.1), Inches(5.0), Inches(2.6), Inches(0.2))
    h2_mask.fill.solid()
    h2_mask.fill.fore_color.rgb = GREEN
    h2_mask.line.fill.background()

    h2_text = slide.shapes.add_textbox(Inches(10.1), Inches(4.8), Inches(2.6), Inches(0.4))
    p = h2_text.text_frame.paragraphs[0]
    p.text = "预期目标（手动）"
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = WHITE
    p.font.bold = True

    # 架构图标 (绘制)
    for i, (px, py) in enumerate([(10.7, 5.3), (11.3, 5.3), (11.3, 5.7)]):
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(px), Inches(py), Inches(0.4), Inches(0.25))
        rect.fill.background()
        rect.line.color.rgb = GREEN
        rect.line.width = Pt(2)
        l = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(px+0.05), Inches(py+0.1), Inches(0.3), Inches(0.02))
        l.fill.solid()
        l.fill.fore_color.rgb = GREEN
        l.line.fill.background()
    
    c1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(11.1), Inches(5.425), Inches(11.3), Inches(5.425))
    c1.line.color.rgb = GREEN
    c1.line.width = Pt(2)
    c2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10.9), Inches(5.55), Inches(10.9), Inches(5.825))
    c2.line.color.rgb = GREEN
    c2.line.width = Pt(2)
    c3 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10.9), Inches(5.825), Inches(11.3), Inches(5.825))
    c3.line.color.rgb = GREEN
    c3.line.width = Pt(2)

    # 成功文字
    check_icon = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.3), Inches(6.3), Inches(0.2), Inches(0.2))
    check_icon.fill.solid()
    check_icon.fill.fore_color.rgb = GREEN
    check_icon.line.fill.background()
    chk = slide.shapes.add_textbox(Inches(10.3), Inches(6.25), Inches(0.2), Inches(0.2))
    p = chk.text_frame.paragraphs[0]
    p.text = "✓"
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    check_text = slide.shapes.add_textbox(Inches(10.6), Inches(6.2), Inches(2.0), Inches(0.6))
    p = check_text.text_frame.paragraphs[0]
    p.text = "清晰逻辑架构\n准确数据字典"
    p.font.size = Pt(11)
    p.font.color.rgb = BLACK

    # ==================== 页脚 ====================
    footer = slide.shapes.add_textbox(Inches(0), Inches(7.1), Inches(13.333), Inches(0.4))
    p = footer.text_frame.paragraphs[0]
    p.text = "Page 3 of 4"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)



# ── Slide 4 ──

def build_slide_4(slide):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

    # Colors
    DARK_BLUE = RGBColor(0x0F, 0x2B, 0x5B)
    GRAY_TEXT = RGBColor(0x59, 0x59, 0x59)
    GREEN = RGBColor(0x2E, 0x9E, 0x66)
    GRAY_BLUE = RGBColor(0x7C, 0x8A, 0x9C)
    RED = RGBColor(0xD9, 0x53, 0x4F)
    BLACK = RGBColor(0x00, 0x00, 0x00)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    BG_COLOR = RGBColor(0xF8, 0xF9, 0xFA)

    # Background Canvas
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(0.3), Inches(12.733), Inches(6.9))
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(8), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "结论：复杂度决定 AI 价值"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(8), Inches(0.4))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "SAP ABAP AI Coding 效率与质量总览"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY_TEXT

    # Axes
    center_x, center_y = 6.66, 3.3

    # Horizontal Axis
    h_line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.5), Inches(center_y), Inches(11.8), Inches(center_y))
    h_line.line.color.rgb = DARK_BLUE
    h_line.line.width = Pt(1.5)
    
    l_arrow = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(1.45), Inches(center_y - 0.05), Inches(0.1), Inches(0.1))
    l_arrow.rotation = -90
    l_arrow.fill.solid()
    l_arrow.fill.fore_color.rgb = DARK_BLUE
    l_arrow.line.fill.background()
    
    r_arrow = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(11.75), Inches(center_y - 0.05), Inches(0.1), Inches(0.1))
    r_arrow.rotation = 90
    r_arrow.fill.solid()
    r_arrow.fill.fore_color.rgb = DARK_BLUE
    r_arrow.line.fill.background()

    # Vertical Axis
    v_line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(center_x), Inches(1.8), Inches(center_x), Inches(4.8))
    v_line.line.color.rgb = DARK_BLUE
    v_line.line.width = Pt(1.5)
    
    t_arrow = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(center_x - 0.05), Inches(1.75), Inches(0.1), Inches(0.1))
    t_arrow.fill.solid()
    t_arrow.fill.fore_color.rgb = DARK_BLUE
    t_arrow.line.fill.background()
    
    b_arrow = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(center_x - 0.05), Inches(4.75), Inches(0.1), Inches(0.1))
    b_arrow.rotation = 180
    b_arrow.fill.solid()
    b_arrow.fill.fore_color.rgb = DARK_BLUE
    b_arrow.line.fill.background()

    # Center Circle
    c_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(center_x - 0.1), Inches(center_y - 0.1), Inches(0.2), Inches(0.2))
    c_circle.fill.solid()
    c_circle.fill.fore_color.rgb = BG_COLOR
    c_circle.line.color.rgb = DARK_BLUE
    c_circle.line.width = Pt(1.5)

    # Axis Labels
    def add_axis_label(left, top, width, text, align):
        box = slide.shapes.add_textbox(left, top, width, Inches(0.6))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(12)
        p.alignment = align

    add_axis_label(Inches(0.5), Inches(3.0), Inches(0.9), "任务度\nSimple 简单", PP_ALIGN.RIGHT)
    add_axis_label(Inches(11.9), Inches(3.0), Inches(1.2), "复杂度\nComplex 复杂", PP_ALIGN.LEFT)
    add_axis_label(Inches(5.66), Inches(1.3), Inches(2.0), "效率影响\nPositive 正面", PP_ALIGN.CENTER)
    add_axis_label(Inches(5.66), Inches(4.9), Inches(2.0), "Negative 负面\n效率影响", PP_ALIGN.CENTER)

    # Quadrant Data
    def add_quadrant_data(left, top, icon_char, icon_color, title, title_color, subtitle, desc):
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top + Inches(0.05), Inches(0.35), Inches(0.35))
        circle.fill.solid()
        circle.fill.fore_color.rgb = icon_color
        circle.line.fill.background()
        tf = circle.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = icon_char
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        box = slide.shapes.add_textbox(left + Inches(0.4), top, Inches(2.5), Inches(1.0))
        tf = box.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.name = "Microsoft YaHei"
        p1.font.size = Pt(20)
        p1.font.bold = True
        p1.font.color.rgb = title_color

        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.name = "Microsoft YaHei"
        p2.font.size = Pt(14)
        p2.font.bold = True
        p2.font.color.rgb = BLACK
        p2.space_before = Pt(5)

        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.name = "Microsoft YaHei"
        p3.font.size = Pt(12)
        p3.font.color.rgb = GRAY_TEXT

    add_quadrant_data(Inches(2.9), Inches(2.0), "↑", GREEN, "+50% 提效", GREEN, "🕒 简单场景", "快速生成工具代码")
    add_quadrant_data(Inches(6.8), Inches(2.3), "→", GRAY_BLUE, "~0% 持平", GRAY_BLUE, "⚖️ 中等场景", "需人工修正逻辑")
    add_quadrant_data(Inches(8.5), Inches(3.6), "↓", RED, "-60% 效率反降", RED, "⚙️ 复杂场景", "频繁幻觉与重构")

    # Bottom Section - Title
    bot_title = slide.shapes.add_textbox(Inches(0.6), Inches(4.8), Inches(3), Inches(0.4))
    p = bot_title.text_frame.paragraphs[0]
    p.text = "核心洞察与评价"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLACK

    # Bottom Left Block - Chart Icon
    chart_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(5.4), Inches(0.8), Inches(0.6))
    chart_box.fill.background()
    chart_box.line.color.rgb = BLACK
    chart_box.line.width = Pt(1.5)

    points = [(1.0, 5.5), (1.2, 5.6), (1.4, 5.8), (1.6, 5.9)]
    for i in range(len(points)-1):
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(points[i][0]), Inches(points[i][1]), Inches(points[i+1][0]), Inches(points[i+1][1]))
        line.line.color.rgb = BLACK
        line.line.width = Pt(1.5)
    for pt in points:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(pt[0]-0.05), Inches(pt[1]-0.05), Inches(0.1), Inches(0.1))
        dot.fill.solid()
        dot.fill.fore_color.rgb = WHITE
        dot.line.color.rgb = BLACK
        dot.line.width = Pt(1.5)

    # Bottom Left Block - Text
    left_text = slide.shapes.add_textbox(Inches(1.9), Inches(5.2), Inches(4.8), Inches(1.2))
    tf = left_text.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = "效率曲线与工具表现"
    p1.font.name = "Microsoft YaHei"
    p1.font.size = Pt(13)
    p1.font.bold = True

    p2 = tf.add_paragraph()
    p2.text = "• 效率趋势：简单场景提效50%，中等场景持平，复杂场景效率反降60%。"
    p2.font.name = "Microsoft YaHei"
    p2.font.size = Pt(11)
    p2.space_before = Pt(3)

    p3 = tf.add_paragraph()
    p3.text = "• 工具评价：Claude Code 在文档解析与逻辑框架上优于 Copilot，但均不精通 SAP 领域。 ✅ ， ❌ Copilot"
    p3.font.name = "Microsoft YaHei"
    p3.font.size = Pt(11)
    p3.space_before = Pt(3)

    # Bottom Right Block - Warning Icon
    warn_tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(7.0), Inches(5.3), Inches(0.5), Inches(0.45))
    warn_tri.fill.background()
    warn_tri.line.color.rgb = BLACK
    warn_tri.line.width = Pt(1.5)
    warn_ex = slide.shapes.add_textbox(Inches(7.0), Inches(5.35), Inches(0.5), Inches(0.4))
    p = warn_ex.text_frame.paragraphs[0]
    p.text = "!"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(16)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Bottom Right Block - Broken Link Icon
    link_oval1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.05), Inches(5.9), Inches(0.25), Inches(0.15))
    link_oval1.rotation = 45
    link_oval1.fill.background()
    link_oval1.line.color.rgb = BLACK
    link_oval1.line.width = Pt(1.5)

    link_oval2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(6.05), Inches(0.25), Inches(0.15))
    link_oval2.rotation = 45
    link_oval2.fill.background()
    link_oval2.line.color.rgb = BLACK
    link_oval2.line.width = Pt(1.5)

    break_line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(7.1), Inches(6.1), Inches(7.35), Inches(5.9))
    break_line.line.color.rgb = BLACK
    break_line.line.width = Pt(1.5)

    # Bottom Right Block - Text
    right_text = slide.shapes.add_textbox(Inches(7.7), Inches(5.2), Inches(4.8), Inches(1.2))
    tf = right_text.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = "核心痛点与建议"
    p1.font.name = "Microsoft YaHei"
    p1.font.size = Pt(13)
    p1.font.bold = True

    p2 = tf.add_paragraph()
    p2.text = "• 核心痛点：虚构数据字典与忽略提示词是限制 AI 在 SAP 领域应用的两大死穴。"
    p2.font.name = "Microsoft YaHei"
    p2.font.size = Pt(11)
    p2.space_before = Pt(3)

    p3 = tf.add_paragraph()
    p3.text = "• 专家建议：现阶段 AI 仅适用于辅助编写简单工具脚本，无法应对中高复杂度业务。 ❌ 中高复杂度业务。"
    p3.font.name = "Microsoft YaHei"
    p3.font.size = Pt(11)
    p3.space_before = Pt(3)

    # Conclusion Box
    conc_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(6.6), Inches(11.833), Inches(0.5))
    conc_box.fill.solid()
    conc_box.fill.fore_color.rgb = BG_COLOR
    conc_box.line.color.rgb = BLACK
    conc_box.line.width = Pt(1.5)

    tf = conc_box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = "现阶段 AI 仅适用于辅助编写简单工具脚本，无法应对中高复杂度业务。"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLACK
    p.alignment = PP_ALIGN.CENTER



# ── Main ──
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

s0 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_1(s0)
s1 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_2(s1)
s2 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_3(s2)
s3 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_4(s3)
prs.save(OUTPUT_PATH)
