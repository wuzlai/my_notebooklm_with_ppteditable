from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# ── Color Palette ──
BLUE_HEADER = RGBColor(0x5B, 0x9B, 0xD5)
BLUE_DARK   = RGBColor(0x4A, 0x86, 0xC8)
CYAN        = RGBColor(0x00, 0xBC, 0xD4)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x33, 0x33, 0x33)
GRAY_TEXT   = RGBColor(0x55, 0x55, 0x55)
GRAY_BAR    = RGBColor(0xB0, 0xBE, 0xC5)
RED         = RGBColor(0xE5, 0x39, 0x35)
GREEN       = RGBColor(0x43, 0xA0, 0x47)
ORANGE      = RGBColor(0xFF, 0x98, 0x00)

ICON_BG  = RGBColor(0xE3, 0xE8, 0xED)
ICON_FG  = RGBColor(0x54, 0x6E, 0x7A)
FONT_NAME = "Microsoft YaHei"
SLIDE_WIDTH = Inches(13.333)
HEADER_H    = Inches(0.75)
SUBTITLE_Y  = Inches(0.95)


def add_header_banner(slide, title_text, bg_color=None):
    if bg_color is None:
        bg_color = BLUE_HEADER
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, HEADER_H
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = bg_color
    banner.line.fill.background()
    tf = banner.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(26)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME


def add_subtitle(slide, text, left, top, width=Inches(12), font_size=Pt(18)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.font.name = FONT_NAME
    return txBox


def add_icon_box(slide, left, top, symbol, size=Inches(0.48)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ICON_BG
    shape.line.fill.background()
    shape.adjustments[0] = 0.25
    tf = shape.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = symbol
    p.font.size = Pt(18)
    p.font.color.rgb = ICON_FG
    p.font.bold = False
    return shape


def add_bullet_item(slide, left, top, symbol, label, description,
                    width=Inches(5.5), desc_size=Pt(13)):
    add_icon_box(slide, left, top, symbol)
    text_left = left + Inches(0.65)
    txBox = slide.shapes.add_textbox(text_left, top - Inches(0.02), width, Inches(0.65))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run_label = p.add_run()
    run_label.text = label + "\uff1a"
    run_label.font.size = Pt(14)
    run_label.font.color.rgb = BLACK
    run_label.font.bold = True
    run_label.font.name = FONT_NAME
    run_desc = p.add_run()
    run_desc.text = description
    run_desc.font.size = desc_size
    run_desc.font.color.rgb = GRAY_TEXT
    run_desc.font.bold = False
    run_desc.font.name = FONT_NAME
    return txBox


def add_conclusion_box(slide, left, top, width, text, font_size=Pt(13)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = BLACK
    run.font.bold = True
    run.font.name = FONT_NAME
    return txBox


def add_table(slide, left, top, width, height, rows, cols, data,
              header_color=None, col_widths=None):
    if header_color is None:
        header_color = BLUE_HEADER
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if data[r][c] is not None else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(11)
                    run.font.name = FONT_NAME
                    if r == 0:
                        run.font.color.rgb = WHITE
                        run.font.bold = True
                    else:
                        run.font.color.rgb = BLACK
                        run.font.bold = False
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else RGBColor(0xF5, 0xF5, 0xF5)
            cell.margin_left = Pt(5)
            cell.margin_right = Pt(5)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
    return table_shape


def add_bar_chart(slide, left, top, width, height,
                  categories, values, title="", bar_colors=None):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', values)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    chart.chart_style = 2
    plot = chart.plots[0]
    plot.gap_width = 100
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CYAN
    series.has_data_labels = True
    dl = series.data_labels
    dl.font.size = Pt(13)
    dl.font.bold = True
    dl.font.color.rgb = BLACK
    dl.number_format = '0.#'
    dl.show_value = True
    dl.label_position = XL_LABEL_POSITION.OUTSIDE_END
    if bar_colors:
        for i, color in enumerate(bar_colors):
            pt = series.points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = color
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(12)
    cat_axis.tick_labels.font.name = FONT_NAME
    cat_axis.major_tick_mark = 2
    cat_axis.format.line.fill.background()
    val_axis = chart.value_axis
    val_axis.visible = False
    val_axis.major_tick_mark = 2
    val_axis.format.line.fill.background()
    val_axis.major_gridlines.format.line.fill.background()
    if title:
        chart.has_title = True
        ct = chart.chart_title.text_frame.paragraphs[0]
        ct.text = title
        ct.font.size = Pt(14)
        ct.font.bold = True
        ct.font.name = FONT_NAME
    else:
        chart.has_title = False
    return chart_frame


def add_callout_label(slide, left, top, text, bg_color=None, font_size=Pt(11)):
    if bg_color is None:
        bg_color = CYAN
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.3), Inches(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME
    return shape


def add_data_card(slide, left, top, width, height, value, label,
                  value_color=None, bg_color=None):
    if value_color is None:
        value_color = CYAN
    if bg_color is None:
        bg_color = WHITE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(3)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run1 = p1.add_run()
    run1.text = str(value)
    run1.font.size = Pt(24)
    run1.font.color.rgb = value_color
    run1.font.bold = True
    run1.font.name = FONT_NAME
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(10)
    run2.font.color.rgb = GRAY_TEXT
    run2.font.bold = False
    run2.font.name = FONT_NAME
    return shape


OUTPUT_PATH = r"C:\Users\Administrator\Desktop\notebook\my_notebooklm_with_ppteditable\projects\测试DEMO\最终文档\ppt_slides\slide_04.pptx"


# ── Slide 1 ──

def build_slide_1(slide):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

    # Colors
    DARK_BLUE = RGBColor(0x0F, 0x2B, 0x5B)
    GRAY_TEXT = RGBColor(0x59, 0x59, 0x59)
    GREEN = RGBColor(0x2E, 0x9E, 0x66)
    GRAY_BLUE = RGBColor(0x7C, 0x8A, 0x9C)
    RED = RGBColor(0xD9, 0x53, 0x4F)
    BLACK = RGBColor(0x00, 0x00, 0x00)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    BG_COLOR = RGBColor(0xF8, 0xF9, 0xFA)

    # Background Canvas
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(0.3), Inches(12.733), Inches(6.9))
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(8), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "结论：复杂度决定 AI 价值"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(8), Inches(0.4))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "SAP ABAP AI Coding 效率与质量总览"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY_TEXT

    # Axes
    center_x, center_y = 6.66, 3.3

    # Horizontal Axis
    h_line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.5), Inches(center_y), Inches(11.8), Inches(center_y))
    h_line.line.color.rgb = DARK_BLUE
    h_line.line.width = Pt(1.5)
    
    l_arrow = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(1.45), Inches(center_y - 0.05), Inches(0.1), Inches(0.1))
    l_arrow.rotation = -90
    l_arrow.fill.solid()
    l_arrow.fill.fore_color.rgb = DARK_BLUE
    l_arrow.line.fill.background()
    
    r_arrow = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(11.75), Inches(center_y - 0.05), Inches(0.1), Inches(0.1))
    r_arrow.rotation = 90
    r_arrow.fill.solid()
    r_arrow.fill.fore_color.rgb = DARK_BLUE
    r_arrow.line.fill.background()

    # Vertical Axis
    v_line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(center_x), Inches(1.8), Inches(center_x), Inches(4.8))
    v_line.line.color.rgb = DARK_BLUE
    v_line.line.width = Pt(1.5)
    
    t_arrow = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(center_x - 0.05), Inches(1.75), Inches(0.1), Inches(0.1))
    t_arrow.fill.solid()
    t_arrow.fill.fore_color.rgb = DARK_BLUE
    t_arrow.line.fill.background()
    
    b_arrow = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(center_x - 0.05), Inches(4.75), Inches(0.1), Inches(0.1))
    b_arrow.rotation = 180
    b_arrow.fill.solid()
    b_arrow.fill.fore_color.rgb = DARK_BLUE
    b_arrow.line.fill.background()

    # Center Circle
    c_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(center_x - 0.1), Inches(center_y - 0.1), Inches(0.2), Inches(0.2))
    c_circle.fill.solid()
    c_circle.fill.fore_color.rgb = BG_COLOR
    c_circle.line.color.rgb = DARK_BLUE
    c_circle.line.width = Pt(1.5)

    # Axis Labels
    def add_axis_label(left, top, width, text, align):
        box = slide.shapes.add_textbox(left, top, width, Inches(0.6))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(12)
        p.alignment = align

    add_axis_label(Inches(0.5), Inches(3.0), Inches(0.9), "任务度\nSimple 简单", PP_ALIGN.RIGHT)
    add_axis_label(Inches(11.9), Inches(3.0), Inches(1.2), "复杂度\nComplex 复杂", PP_ALIGN.LEFT)
    add_axis_label(Inches(5.66), Inches(1.3), Inches(2.0), "效率影响\nPositive 正面", PP_ALIGN.CENTER)
    add_axis_label(Inches(5.66), Inches(4.9), Inches(2.0), "Negative 负面\n效率影响", PP_ALIGN.CENTER)

    # Quadrant Data
    def add_quadrant_data(left, top, icon_char, icon_color, title, title_color, subtitle, desc):
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top + Inches(0.05), Inches(0.35), Inches(0.35))
        circle.fill.solid()
        circle.fill.fore_color.rgb = icon_color
        circle.line.fill.background()
        tf = circle.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = icon_char
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        box = slide.shapes.add_textbox(left + Inches(0.4), top, Inches(2.5), Inches(1.0))
        tf = box.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.name = "Microsoft YaHei"
        p1.font.size = Pt(20)
        p1.font.bold = True
        p1.font.color.rgb = title_color

        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.name = "Microsoft YaHei"
        p2.font.size = Pt(14)
        p2.font.bold = True
        p2.font.color.rgb = BLACK
        p2.space_before = Pt(5)

        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.name = "Microsoft YaHei"
        p3.font.size = Pt(12)
        p3.font.color.rgb = GRAY_TEXT

    add_quadrant_data(Inches(2.9), Inches(2.0), "↑", GREEN, "+50% 提效", GREEN, "🕒 简单场景", "快速生成工具代码")
    add_quadrant_data(Inches(6.8), Inches(2.3), "→", GRAY_BLUE, "~0% 持平", GRAY_BLUE, "⚖️ 中等场景", "需人工修正逻辑")
    add_quadrant_data(Inches(8.5), Inches(3.6), "↓", RED, "-60% 效率反降", RED, "⚙️ 复杂场景", "频繁幻觉与重构")

    # Bottom Section - Title
    bot_title = slide.shapes.add_textbox(Inches(0.6), Inches(4.8), Inches(3), Inches(0.4))
    p = bot_title.text_frame.paragraphs[0]
    p.text = "核心洞察与评价"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLACK

    # Bottom Left Block - Chart Icon
    chart_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(5.4), Inches(0.8), Inches(0.6))
    chart_box.fill.background()
    chart_box.line.color.rgb = BLACK
    chart_box.line.width = Pt(1.5)

    points = [(1.0, 5.5), (1.2, 5.6), (1.4, 5.8), (1.6, 5.9)]
    for i in range(len(points)-1):
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(points[i][0]), Inches(points[i][1]), Inches(points[i+1][0]), Inches(points[i+1][1]))
        line.line.color.rgb = BLACK
        line.line.width = Pt(1.5)
    for pt in points:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(pt[0]-0.05), Inches(pt[1]-0.05), Inches(0.1), Inches(0.1))
        dot.fill.solid()
        dot.fill.fore_color.rgb = WHITE
        dot.line.color.rgb = BLACK
        dot.line.width = Pt(1.5)

    # Bottom Left Block - Text
    left_text = slide.shapes.add_textbox(Inches(1.9), Inches(5.2), Inches(4.8), Inches(1.2))
    tf = left_text.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = "效率曲线与工具表现"
    p1.font.name = "Microsoft YaHei"
    p1.font.size = Pt(13)
    p1.font.bold = True

    p2 = tf.add_paragraph()
    p2.text = "• 效率趋势：简单场景提效50%，中等场景持平，复杂场景效率反降60%。"
    p2.font.name = "Microsoft YaHei"
    p2.font.size = Pt(11)
    p2.space_before = Pt(3)

    p3 = tf.add_paragraph()
    p3.text = "• 工具评价：Claude Code 在文档解析与逻辑框架上优于 Copilot，但均不精通 SAP 领域。 ✅ ， ❌ Copilot"
    p3.font.name = "Microsoft YaHei"
    p3.font.size = Pt(11)
    p3.space_before = Pt(3)

    # Bottom Right Block - Warning Icon
    warn_tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(7.0), Inches(5.3), Inches(0.5), Inches(0.45))
    warn_tri.fill.background()
    warn_tri.line.color.rgb = BLACK
    warn_tri.line.width = Pt(1.5)
    warn_ex = slide.shapes.add_textbox(Inches(7.0), Inches(5.35), Inches(0.5), Inches(0.4))
    p = warn_ex.text_frame.paragraphs[0]
    p.text = "!"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(16)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Bottom Right Block - Broken Link Icon
    link_oval1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.05), Inches(5.9), Inches(0.25), Inches(0.15))
    link_oval1.rotation = 45
    link_oval1.fill.background()
    link_oval1.line.color.rgb = BLACK
    link_oval1.line.width = Pt(1.5)

    link_oval2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(6.05), Inches(0.25), Inches(0.15))
    link_oval2.rotation = 45
    link_oval2.fill.background()
    link_oval2.line.color.rgb = BLACK
    link_oval2.line.width = Pt(1.5)

    break_line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(7.1), Inches(6.1), Inches(7.35), Inches(5.9))
    break_line.line.color.rgb = BLACK
    break_line.line.width = Pt(1.5)

    # Bottom Right Block - Text
    right_text = slide.shapes.add_textbox(Inches(7.7), Inches(5.2), Inches(4.8), Inches(1.2))
    tf = right_text.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = "核心痛点与建议"
    p1.font.name = "Microsoft YaHei"
    p1.font.size = Pt(13)
    p1.font.bold = True

    p2 = tf.add_paragraph()
    p2.text = "• 核心痛点：虚构数据字典与忽略提示词是限制 AI 在 SAP 领域应用的两大死穴。"
    p2.font.name = "Microsoft YaHei"
    p2.font.size = Pt(11)
    p2.space_before = Pt(3)

    p3 = tf.add_paragraph()
    p3.text = "• 专家建议：现阶段 AI 仅适用于辅助编写简单工具脚本，无法应对中高复杂度业务。 ❌ 中高复杂度业务。"
    p3.font.name = "Microsoft YaHei"
    p3.font.size = Pt(11)
    p3.space_before = Pt(3)

    # Conclusion Box
    conc_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(6.6), Inches(11.833), Inches(0.5))
    conc_box.fill.solid()
    conc_box.fill.fore_color.rgb = BG_COLOR
    conc_box.line.color.rgb = BLACK
    conc_box.line.width = Pt(1.5)

    tf = conc_box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = "现阶段 AI 仅适用于辅助编写简单工具脚本，无法应对中高复杂度业务。"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLACK
    p.alignment = PP_ALIGN.CENTER



# ── Main ──
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

s0 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_1(s0)
prs.save(OUTPUT_PATH)
