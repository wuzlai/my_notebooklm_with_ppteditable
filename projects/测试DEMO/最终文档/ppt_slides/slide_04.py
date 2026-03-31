def build_slide(slide):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

    # Colors
    DARK_BLUE = RGBColor(0x0F, 0x2B, 0x5B)
    GRAY_TEXT = RGBColor(0x59, 0x59, 0x59)
    GREEN = RGBColor(0x2E, 0x9E, 0x66)
    GRAY_BLUE = RGBColor(0x7C, 0x8A, 0x9C)
    RED = RGBColor(0xD9, 0x53, 0x4F)
    BLACK = RGBColor(0x00, 0x00, 0x00)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    BG_COLOR = RGBColor(0xF8, 0xF9, 0xFA)

    # Background Canvas
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(0.3), Inches(12.733), Inches(6.9))
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(8), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "结论：复杂度决定 AI 价值"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(8), Inches(0.4))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "SAP ABAP AI Coding 效率与质量总览"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY_TEXT

    # Axes
    center_x, center_y = 6.66, 3.3

    # Horizontal Axis
    h_line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.5), Inches(center_y), Inches(11.8), Inches(center_y))
    h_line.line.color.rgb = DARK_BLUE
    h_line.line.width = Pt(1.5)
    
    l_arrow = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(1.45), Inches(center_y - 0.05), Inches(0.1), Inches(0.1))
    l_arrow.rotation = -90
    l_arrow.fill.solid()
    l_arrow.fill.fore_color.rgb = DARK_BLUE
    l_arrow.line.fill.background()
    
    r_arrow = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(11.75), Inches(center_y - 0.05), Inches(0.1), Inches(0.1))
    r_arrow.rotation = 90
    r_arrow.fill.solid()
    r_arrow.fill.fore_color.rgb = DARK_BLUE
    r_arrow.line.fill.background()

    # Vertical Axis
    v_line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(center_x), Inches(1.8), Inches(center_x), Inches(4.8))
    v_line.line.color.rgb = DARK_BLUE
    v_line.line.width = Pt(1.5)
    
    t_arrow = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(center_x - 0.05), Inches(1.75), Inches(0.1), Inches(0.1))
    t_arrow.fill.solid()
    t_arrow.fill.fore_color.rgb = DARK_BLUE
    t_arrow.line.fill.background()
    
    b_arrow = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(center_x - 0.05), Inches(4.75), Inches(0.1), Inches(0.1))
    b_arrow.rotation = 180
    b_arrow.fill.solid()
    b_arrow.fill.fore_color.rgb = DARK_BLUE
    b_arrow.line.fill.background()

    # Center Circle
    c_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(center_x - 0.1), Inches(center_y - 0.1), Inches(0.2), Inches(0.2))
    c_circle.fill.solid()
    c_circle.fill.fore_color.rgb = BG_COLOR
    c_circle.line.color.rgb = DARK_BLUE
    c_circle.line.width = Pt(1.5)

    # Axis Labels
    def add_axis_label(left, top, width, text, align):
        box = slide.shapes.add_textbox(left, top, width, Inches(0.6))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(12)
        p.alignment = align

    add_axis_label(Inches(0.5), Inches(3.0), Inches(0.9), "任务度\nSimple 简单", PP_ALIGN.RIGHT)
    add_axis_label(Inches(11.9), Inches(3.0), Inches(1.2), "复杂度\nComplex 复杂", PP_ALIGN.LEFT)
    add_axis_label(Inches(5.66), Inches(1.3), Inches(2.0), "效率影响\nPositive 正面", PP_ALIGN.CENTER)
    add_axis_label(Inches(5.66), Inches(4.9), Inches(2.0), "Negative 负面\n效率影响", PP_ALIGN.CENTER)

    # Quadrant Data
    def add_quadrant_data(left, top, icon_char, icon_color, title, title_color, subtitle, desc):
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top + Inches(0.05), Inches(0.35), Inches(0.35))
        circle.fill.solid()
        circle.fill.fore_color.rgb = icon_color
        circle.line.fill.background()
        tf = circle.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = icon_char
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        box = slide.shapes.add_textbox(left + Inches(0.4), top, Inches(2.5), Inches(1.0))
        tf = box.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.name = "Microsoft YaHei"
        p1.font.size = Pt(20)
        p1.font.bold = True
        p1.font.color.rgb = title_color

        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.name = "Microsoft YaHei"
        p2.font.size = Pt(14)
        p2.font.bold = True
        p2.font.color.rgb = BLACK
        p2.space_before = Pt(5)

        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.name = "Microsoft YaHei"
        p3.font.size = Pt(12)
        p3.font.color.rgb = GRAY_TEXT

    add_quadrant_data(Inches(2.9), Inches(2.0), "↑", GREEN, "+50% 提效", GREEN, "🕒 简单场景", "快速生成工具代码")
    add_quadrant_data(Inches(6.8), Inches(2.3), "→", GRAY_BLUE, "~0% 持平", GRAY_BLUE, "⚖️ 中等场景", "需人工修正逻辑")
    add_quadrant_data(Inches(8.5), Inches(3.6), "↓", RED, "-60% 效率反降", RED, "⚙️ 复杂场景", "频繁幻觉与重构")

    # Bottom Section - Title
    bot_title = slide.shapes.add_textbox(Inches(0.6), Inches(4.8), Inches(3), Inches(0.4))
    p = bot_title.text_frame.paragraphs[0]
    p.text = "核心洞察与评价"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLACK

    # Bottom Left Block - Chart Icon
    chart_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(5.4), Inches(0.8), Inches(0.6))
    chart_box.fill.background()
    chart_box.line.color.rgb = BLACK
    chart_box.line.width = Pt(1.5)

    points = [(1.0, 5.5), (1.2, 5.6), (1.4, 5.8), (1.6, 5.9)]
    for i in range(len(points)-1):
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(points[i][0]), Inches(points[i][1]), Inches(points[i+1][0]), Inches(points[i+1][1]))
        line.line.color.rgb = BLACK
        line.line.width = Pt(1.5)
    for pt in points:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(pt[0]-0.05), Inches(pt[1]-0.05), Inches(0.1), Inches(0.1))
        dot.fill.solid()
        dot.fill.fore_color.rgb = WHITE
        dot.line.color.rgb = BLACK
        dot.line.width = Pt(1.5)

    # Bottom Left Block - Text
    left_text = slide.shapes.add_textbox(Inches(1.9), Inches(5.2), Inches(4.8), Inches(1.2))
    tf = left_text.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = "效率曲线与工具表现"
    p1.font.name = "Microsoft YaHei"
    p1.font.size = Pt(13)
    p1.font.bold = True

    p2 = tf.add_paragraph()
    p2.text = "• 效率趋势：简单场景提效50%，中等场景持平，复杂场景效率反降60%。"
    p2.font.name = "Microsoft YaHei"
    p2.font.size = Pt(11)
    p2.space_before = Pt(3)

    p3 = tf.add_paragraph()
    p3.text = "• 工具评价：Claude Code 在文档解析与逻辑框架上优于 Copilot，但均不精通 SAP 领域。 ✅ ， ❌ Copilot"
    p3.font.name = "Microsoft YaHei"
    p3.font.size = Pt(11)
    p3.space_before = Pt(3)

    # Bottom Right Block - Warning Icon
    warn_tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(7.0), Inches(5.3), Inches(0.5), Inches(0.45))
    warn_tri.fill.background()
    warn_tri.line.color.rgb = BLACK
    warn_tri.line.width = Pt(1.5)
    warn_ex = slide.shapes.add_textbox(Inches(7.0), Inches(5.35), Inches(0.5), Inches(0.4))
    p = warn_ex.text_frame.paragraphs[0]
    p.text = "!"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(16)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Bottom Right Block - Broken Link Icon
    link_oval1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.05), Inches(5.9), Inches(0.25), Inches(0.15))
    link_oval1.rotation = 45
    link_oval1.fill.background()
    link_oval1.line.color.rgb = BLACK
    link_oval1.line.width = Pt(1.5)

    link_oval2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(6.05), Inches(0.25), Inches(0.15))
    link_oval2.rotation = 45
    link_oval2.fill.background()
    link_oval2.line.color.rgb = BLACK
    link_oval2.line.width = Pt(1.5)

    break_line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(7.1), Inches(6.1), Inches(7.35), Inches(5.9))
    break_line.line.color.rgb = BLACK
    break_line.line.width = Pt(1.5)

    # Bottom Right Block - Text
    right_text = slide.shapes.add_textbox(Inches(7.7), Inches(5.2), Inches(4.8), Inches(1.2))
    tf = right_text.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = "核心痛点与建议"
    p1.font.name = "Microsoft YaHei"
    p1.font.size = Pt(13)
    p1.font.bold = True

    p2 = tf.add_paragraph()
    p2.text = "• 核心痛点：虚构数据字典与忽略提示词是限制 AI 在 SAP 领域应用的两大死穴。"
    p2.font.name = "Microsoft YaHei"
    p2.font.size = Pt(11)
    p2.space_before = Pt(3)

    p3 = tf.add_paragraph()
    p3.text = "• 专家建议：现阶段 AI 仅适用于辅助编写简单工具脚本，无法应对中高复杂度业务。 ❌ 中高复杂度业务。"
    p3.font.name = "Microsoft YaHei"
    p3.font.size = Pt(11)
    p3.space_before = Pt(3)

    # Conclusion Box
    conc_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(6.6), Inches(11.833), Inches(0.5))
    conc_box.fill.solid()
    conc_box.fill.fore_color.rgb = BG_COLOR
    conc_box.line.color.rgb = BLACK
    conc_box.line.width = Pt(1.5)

    tf = conc_box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = "现阶段 AI 仅适用于辅助编写简单工具脚本，无法应对中高复杂度业务。"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLACK
    p.alignment = PP_ALIGN.CENTER