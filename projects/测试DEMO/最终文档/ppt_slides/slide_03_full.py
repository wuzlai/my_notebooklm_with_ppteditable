from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# ── Color Palette ──
BLUE_HEADER = RGBColor(0x5B, 0x9B, 0xD5)
BLUE_DARK   = RGBColor(0x4A, 0x86, 0xC8)
CYAN        = RGBColor(0x00, 0xBC, 0xD4)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x33, 0x33, 0x33)
GRAY_TEXT   = RGBColor(0x55, 0x55, 0x55)
GRAY_BAR    = RGBColor(0xB0, 0xBE, 0xC5)
RED         = RGBColor(0xE5, 0x39, 0x35)
GREEN       = RGBColor(0x43, 0xA0, 0x47)
ORANGE      = RGBColor(0xFF, 0x98, 0x00)

ICON_BG  = RGBColor(0xE3, 0xE8, 0xED)
ICON_FG  = RGBColor(0x54, 0x6E, 0x7A)
FONT_NAME = "Microsoft YaHei"
SLIDE_WIDTH = Inches(13.333)
HEADER_H    = Inches(0.75)
SUBTITLE_Y  = Inches(0.95)


def add_header_banner(slide, title_text, bg_color=None):
    if bg_color is None:
        bg_color = BLUE_HEADER
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, HEADER_H
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = bg_color
    banner.line.fill.background()
    tf = banner.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(26)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME


def add_subtitle(slide, text, left, top, width=Inches(12), font_size=Pt(18)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.font.name = FONT_NAME
    return txBox


def add_icon_box(slide, left, top, symbol, size=Inches(0.48)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ICON_BG
    shape.line.fill.background()
    shape.adjustments[0] = 0.25
    tf = shape.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = symbol
    p.font.size = Pt(18)
    p.font.color.rgb = ICON_FG
    p.font.bold = False
    return shape


def add_bullet_item(slide, left, top, symbol, label, description,
                    width=Inches(5.5), desc_size=Pt(13)):
    add_icon_box(slide, left, top, symbol)
    text_left = left + Inches(0.65)
    txBox = slide.shapes.add_textbox(text_left, top - Inches(0.02), width, Inches(0.65))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run_label = p.add_run()
    run_label.text = label + "\uff1a"
    run_label.font.size = Pt(14)
    run_label.font.color.rgb = BLACK
    run_label.font.bold = True
    run_label.font.name = FONT_NAME
    run_desc = p.add_run()
    run_desc.text = description
    run_desc.font.size = desc_size
    run_desc.font.color.rgb = GRAY_TEXT
    run_desc.font.bold = False
    run_desc.font.name = FONT_NAME
    return txBox


def add_conclusion_box(slide, left, top, width, text, font_size=Pt(13)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = BLACK
    run.font.bold = True
    run.font.name = FONT_NAME
    return txBox


def add_table(slide, left, top, width, height, rows, cols, data,
              header_color=None, col_widths=None):
    if header_color is None:
        header_color = BLUE_HEADER
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if data[r][c] is not None else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(11)
                    run.font.name = FONT_NAME
                    if r == 0:
                        run.font.color.rgb = WHITE
                        run.font.bold = True
                    else:
                        run.font.color.rgb = BLACK
                        run.font.bold = False
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else RGBColor(0xF5, 0xF5, 0xF5)
            cell.margin_left = Pt(5)
            cell.margin_right = Pt(5)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
    return table_shape


def add_bar_chart(slide, left, top, width, height,
                  categories, values, title="", bar_colors=None):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', values)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    chart.chart_style = 2
    plot = chart.plots[0]
    plot.gap_width = 100
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CYAN
    series.has_data_labels = True
    dl = series.data_labels
    dl.font.size = Pt(13)
    dl.font.bold = True
    dl.font.color.rgb = BLACK
    dl.number_format = '0.#'
    dl.show_value = True
    dl.label_position = XL_LABEL_POSITION.OUTSIDE_END
    if bar_colors:
        for i, color in enumerate(bar_colors):
            pt = series.points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = color
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(12)
    cat_axis.tick_labels.font.name = FONT_NAME
    cat_axis.major_tick_mark = 2
    cat_axis.format.line.fill.background()
    val_axis = chart.value_axis
    val_axis.visible = False
    val_axis.major_tick_mark = 2
    val_axis.format.line.fill.background()
    val_axis.major_gridlines.format.line.fill.background()
    if title:
        chart.has_title = True
        ct = chart.chart_title.text_frame.paragraphs[0]
        ct.text = title
        ct.font.size = Pt(14)
        ct.font.bold = True
        ct.font.name = FONT_NAME
    else:
        chart.has_title = False
    return chart_frame


def add_callout_label(slide, left, top, text, bg_color=None, font_size=Pt(11)):
    if bg_color is None:
        bg_color = CYAN
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.3), Inches(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME
    return shape


def add_data_card(slide, left, top, width, height, value, label,
                  value_color=None, bg_color=None):
    if value_color is None:
        value_color = CYAN
    if bg_color is None:
        bg_color = WHITE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(3)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run1 = p1.add_run()
    run1.text = str(value)
    run1.font.size = Pt(24)
    run1.font.color.rgb = value_color
    run1.font.bold = True
    run1.font.name = FONT_NAME
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(10)
    run2.font.color.rgb = GRAY_TEXT
    run2.font.bold = False
    run2.font.name = FONT_NAME
    return shape


OUTPUT_PATH = r"C:\Users\Administrator\Desktop\notebook\my_notebooklm_with_ppteditable\projects\测试DEMO\最终文档\ppt_slides\slide_03.pptx"


# ── Slide 1 ──

def build_slide_1(slide):
    from pptx.enum.shapes import MSO_CONNECTOR
    
    # 1. 添加背景卡片（可选，为了更好的视觉效果，这里添加一个全屏的浅色背景和白色卡片）
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)
    bg.line.fill.background()
    
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(0.4), Inches(12.533), Inches(6.7))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.fill.background()

    # 2. 标题和副标题
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(8), Inches(0.6))
    p = title_box.text_frame.paragraphs[0]
    p.text = "复杂场景失效：修复-爆炸模式"
    p.font.name = FONT_NAME
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1C, 0x2A, 0x43)

    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8), Inches(0.4))
    p = sub_box.text_frame.paragraphs[0]
    p.text = "跨工厂 STO 报表（高复杂度）深度测试"
    p.font.name = FONT_NAME
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY_TEXT

    # ==================== 左侧栏 ====================

    # 效率倒挂
    h1 = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(4), Inches(0.4))
    p = h1.text_frame.paragraphs[0]
    p.text = "效率倒挂"
    p.font.bold = True
    p.font.size = Pt(18)

    # 红色数据框
    red_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.3), Inches(5.2), Inches(0.9))
    red_bg.fill.solid()
    red_bg.fill.fore_color.rgb = RGBColor(0xFC, 0xE9, 0xE9)
    red_bg.line.fill.background()

    red_border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.3), Inches(0.08), Inches(0.9))
    red_border.fill.solid()
    red_border.fill.fore_color.rgb = RED
    red_border.line.fill.background()

    rb_text = slide.shapes.add_textbox(Inches(1.0), Inches(2.35), Inches(5.0), Inches(0.8))
    tf = rb_text.text_frame
    p1 = tf.paragraphs[0]
    p1.text = "Claude Code 耗时 8 人天"
    p1.font.size = Pt(22)
    p1.font.bold = True
    p1.font.color.rgb = RED
    p2 = tf.add_paragraph()
    p2.text = "比手写开发（5天）慢 60%"
    p2.font.size = Pt(12)
    p2.font.color.rgb = GRAY_TEXT

    # 修复陷阱
    h2 = slide.shapes.add_textbox(Inches(0.8), Inches(3.4), Inches(4), Inches(0.4))
    p = h2.text_frame.paragraphs[0]
    p.text = "修复陷阱"
    p.font.bold = True
    p.font.size = Pt(18)

    icon1 = slide.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(0.5), Inches(0.5))
    p = icon1.text_frame.paragraphs[0]
    p.text = "🔗"
    p.font.size = Pt(24)
    p.font.color.rgb = RED

    desc1 = slide.shapes.add_textbox(Inches(1.4), Inches(3.85), Inches(4.6), Inches(0.6))
    desc1.text_frame.word_wrap = True
    p = desc1.text_frame.paragraphs[0]
    p.text = "语法错误呈现“3→2→9”反向增长，陷入越修越错的逻辑矛盾。"
    p.font.size = Pt(13)

    # 字典缺失
    h3 = slide.shapes.add_textbox(Inches(0.8), Inches(4.6), Inches(4), Inches(0.4))
    p = h3.text_frame.paragraphs[0]
    p.text = "字典缺失"
    p.font.bold = True
    p.font.size = Pt(18)

    icon2 = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(0.9), Inches(5.1), Inches(0.3), Inches(0.3))
    icon2.fill.solid()
    icon2.fill.fore_color.rgb = RED
    icon2.line.fill.background()
    exc2 = slide.shapes.add_textbox(Inches(0.9), Inches(5.05), Inches(0.3), Inches(0.3))
    p = exc2.text_frame.paragraphs[0]
    p.text = "!"
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    desc2 = slide.shapes.add_textbox(Inches(1.4), Inches(5.05), Inches(4.6), Inches(0.6))
    desc2.text_frame.word_wrap = True
    p = desc2.text_frame.paragraphs[0]
    p.text = "累计虚构 10 余项数据字典组件，AI 无法理解复杂的 SAP 关联逻辑。"
    p.font.size = Pt(13)

    # 开发结论
    h4 = slide.shapes.add_textbox(Inches(0.8), Inches(5.8), Inches(4), Inches(0.4))
    p = h4.text_frame.paragraphs[0]
    p.text = "开发结论"
    p.font.bold = True
    p.font.size = Pt(18)

    desc3 = slide.shapes.add_textbox(Inches(0.8), Inches(6.2), Inches(5.2), Inches(0.6))
    desc3.text_frame.word_wrap = True
    p = desc3.text_frame.paragraphs[0]
    p.text = "•  在涉及 18+ 张表及多级状态追踪时，AI 框架生成能力失效。"
    p.font.size = Pt(13)

    # ==================== 右侧栏 ====================

    # 折线图
    chart_data = CategoryChartData()
    chart_data.categories = ['初始生成', '第一次修复', '第二次修复']
    chart_data.add_series('错误数量', (3, 2, 9))

    x, y, cx, cy = Inches(6.5), Inches(1.5), Inches(6.0), Inches(2.8)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS, x, y, cx, cy, chart_data
    ).chart

    chart.has_title = True
    chart.chart_title.text_frame.text = "修复-爆炸趋势折线图"
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)

    series = chart.series[0]
    series.format.line.color.rgb = RED
    series.format.line.width = Pt(3)
    series.marker.style = 8
    series.marker.size = 8
    series.marker.format.fill.solid()
    series.marker.format.fill.fore_color.rgb = RED
    series.marker.format.line.fill.background()

    series.has_data_labels = True
    for dl in series.data_labels:
        dl.font.size = Pt(12)
        dl.font.color.rgb = RED
        dl.position = XL_LABEL_POSITION.ABOVE

    value_axis = chart.value_axis
    value_axis.has_title = True
    value_axis.axis_title.text_frame.text = "错误数量"
    value_axis.maximum_scale = 10
    value_axis.minimum_scale = 0
    value_axis.major_unit = 2
    value_axis.has_major_gridlines = True
    value_axis.major_gridlines.format.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    chart.category_axis.has_major_gridlines = False

    # 图表标注箭头和文字
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(9.15), Inches(2.6), Inches(2.4), Inches(0.15))
    arrow.rotation = -45
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RED
    arrow.line.fill.background()

    anno_text = slide.shapes.add_textbox(Inches(8.5), Inches(2.0), Inches(2.5), Inches(0.4))
    p = anno_text.text_frame.paragraphs[0]
    p.text = "反向增长：越修越错"
    p.font.size = Pt(11)
    p.font.color.rgb = RED

    # ==================== 底部对比图 ====================

    # 左侧 AI生成 框
    box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(4.8), Inches(2.6), Inches(2.0))
    box1.fill.solid()
    box1.fill.fore_color.rgb = RGBColor(0xFC, 0xE9, 0xE9)
    box1.line.fill.background()

    header1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(4.8), Inches(2.6), Inches(0.4))
    header1.fill.solid()
    header1.fill.fore_color.rgb = RED
    header1.line.fill.background()
    h1_mask = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(5.0), Inches(2.6), Inches(0.2))
    h1_mask.fill.solid()
    h1_mask.fill.fore_color.rgb = RED
    h1_mask.line.fill.background()

    h1_text = slide.shapes.add_textbox(Inches(6.5), Inches(4.8), Inches(2.6), Inches(0.4))
    p = h1_text.text_frame.paragraphs[0]
    p.text = "AI生成"
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = WHITE
    p.font.bold = True

    # 断链图标 (绘制)
    l_link = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.3), Inches(5.3), Inches(0.6), Inches(0.3))
    l_link.rotation = 45
    l_link.fill.background()
    l_link.line.color.rgb = BLACK
    l_link.line.width = Pt(3)
    r_link = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.7), Inches(5.5), Inches(0.6), Inches(0.3))
    r_link.rotation = 45
    r_link.fill.background()
    r_link.line.color.rgb = BLACK
    r_link.line.width = Pt(3)
    b1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.5), Inches(5.4), Inches(0.2), Inches(0.5))
    b1.rotation = 45
    b1.fill.solid()
    b1.fill.fore_color.rgb = RGBColor(0xFC, 0xE9, 0xE9)
    b1.line.fill.background()

    # 警告文字
    warn_icon = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(6.7), Inches(6.3), Inches(0.2), Inches(0.2))
    warn_icon.fill.solid()
    warn_icon.fill.fore_color.rgb = RED
    warn_icon.line.fill.background()
    exc = slide.shapes.add_textbox(Inches(6.7), Inches(6.25), Inches(0.2), Inches(0.2))
    p = exc.text_frame.paragraphs[0]
    p.text = "!"
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    warn_text = slide.shapes.add_textbox(Inches(7.0), Inches(6.2), Inches(2.0), Inches(0.6))
    p = warn_text.text_frame.paragraphs[0]
    p.text = "18+ 张表关联失败\n10+ 幻觉组件"
    p.font.size = Pt(11)
    p.font.color.rgb = BLACK

    # 中间箭头
    mid_arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(9.2), Inches(5.6), Inches(0.8), Inches(0.4))
    mid_arrow.fill.solid()
    mid_arrow.fill.fore_color.rgb = RGBColor(0xCF, 0xD8, 0xDC)
    mid_arrow.line.fill.background()

    cross = slide.shapes.add_shape(MSO_SHAPE.MATH_MULTIPLY, Inches(9.4), Inches(5.6), Inches(0.4), Inches(0.4))
    cross.fill.solid()
    cross.fill.fore_color.rgb = RED
    cross.line.fill.background()

    # 右侧 预期目标 框
    box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.1), Inches(4.8), Inches(2.6), Inches(2.0))
    box2.fill.solid()
    box2.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
    box2.line.fill.background()

    header2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.1), Inches(4.8), Inches(2.6), Inches(0.4))
    header2.fill.solid()
    header2.fill.fore_color.rgb = GREEN
    header2.line.fill.background()
    h2_mask = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10.1), Inches(5.0), Inches(2.6), Inches(0.2))
    h2_mask.fill.solid()
    h2_mask.fill.fore_color.rgb = GREEN
    h2_mask.line.fill.background()

    h2_text = slide.shapes.add_textbox(Inches(10.1), Inches(4.8), Inches(2.6), Inches(0.4))
    p = h2_text.text_frame.paragraphs[0]
    p.text = "预期目标（手动）"
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = WHITE
    p.font.bold = True

    # 架构图标 (绘制)
    for i, (px, py) in enumerate([(10.7, 5.3), (11.3, 5.3), (11.3, 5.7)]):
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(px), Inches(py), Inches(0.4), Inches(0.25))
        rect.fill.background()
        rect.line.color.rgb = GREEN
        rect.line.width = Pt(2)
        l = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(px+0.05), Inches(py+0.1), Inches(0.3), Inches(0.02))
        l.fill.solid()
        l.fill.fore_color.rgb = GREEN
        l.line.fill.background()
    
    c1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(11.1), Inches(5.425), Inches(11.3), Inches(5.425))
    c1.line.color.rgb = GREEN
    c1.line.width = Pt(2)
    c2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10.9), Inches(5.55), Inches(10.9), Inches(5.825))
    c2.line.color.rgb = GREEN
    c2.line.width = Pt(2)
    c3 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10.9), Inches(5.825), Inches(11.3), Inches(5.825))
    c3.line.color.rgb = GREEN
    c3.line.width = Pt(2)

    # 成功文字
    check_icon = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.3), Inches(6.3), Inches(0.2), Inches(0.2))
    check_icon.fill.solid()
    check_icon.fill.fore_color.rgb = GREEN
    check_icon.line.fill.background()
    chk = slide.shapes.add_textbox(Inches(10.3), Inches(6.25), Inches(0.2), Inches(0.2))
    p = chk.text_frame.paragraphs[0]
    p.text = "✓"
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    check_text = slide.shapes.add_textbox(Inches(10.6), Inches(6.2), Inches(2.0), Inches(0.6))
    p = check_text.text_frame.paragraphs[0]
    p.text = "清晰逻辑架构\n准确数据字典"
    p.font.size = Pt(11)
    p.font.color.rgb = BLACK

    # ==================== 页脚 ====================
    footer = slide.shapes.add_textbox(Inches(0), Inches(7.1), Inches(13.333), Inches(0.4))
    p = footer.text_frame.paragraphs[0]
    p.text = "Page 3 of 4"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)



# ── Main ──
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

s0 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_1(s0)
prs.save(OUTPUT_PATH)
