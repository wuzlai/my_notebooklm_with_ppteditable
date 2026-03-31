from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# ── Color Palette ──
BLUE_HEADER = RGBColor(0x5B, 0x9B, 0xD5)
BLUE_DARK   = RGBColor(0x4A, 0x86, 0xC8)
CYAN        = RGBColor(0x00, 0xBC, 0xD4)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x33, 0x33, 0x33)
GRAY_TEXT   = RGBColor(0x55, 0x55, 0x55)
GRAY_BAR    = RGBColor(0xB0, 0xBE, 0xC5)
RED         = RGBColor(0xE5, 0x39, 0x35)
GREEN       = RGBColor(0x43, 0xA0, 0x47)
ORANGE      = RGBColor(0xFF, 0x98, 0x00)

ICON_BG  = RGBColor(0xE3, 0xE8, 0xED)
ICON_FG  = RGBColor(0x54, 0x6E, 0x7A)
FONT_NAME = "Microsoft YaHei"
SLIDE_WIDTH = Inches(13.333)
HEADER_H    = Inches(0.75)
SUBTITLE_Y  = Inches(0.95)


def add_header_banner(slide, title_text, bg_color=None):
    if bg_color is None:
        bg_color = BLUE_HEADER
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, HEADER_H
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = bg_color
    banner.line.fill.background()
    tf = banner.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(26)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME


def add_subtitle(slide, text, left, top, width=Inches(12), font_size=Pt(18)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.font.name = FONT_NAME
    return txBox


def add_icon_box(slide, left, top, symbol, size=Inches(0.48)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ICON_BG
    shape.line.fill.background()
    shape.adjustments[0] = 0.25
    tf = shape.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = symbol
    p.font.size = Pt(18)
    p.font.color.rgb = ICON_FG
    p.font.bold = False
    return shape


def add_bullet_item(slide, left, top, symbol, label, description,
                    width=Inches(5.5), desc_size=Pt(13)):
    add_icon_box(slide, left, top, symbol)
    text_left = left + Inches(0.65)
    txBox = slide.shapes.add_textbox(text_left, top - Inches(0.02), width, Inches(0.65))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run_label = p.add_run()
    run_label.text = label + "\uff1a"
    run_label.font.size = Pt(14)
    run_label.font.color.rgb = BLACK
    run_label.font.bold = True
    run_label.font.name = FONT_NAME
    run_desc = p.add_run()
    run_desc.text = description
    run_desc.font.size = desc_size
    run_desc.font.color.rgb = GRAY_TEXT
    run_desc.font.bold = False
    run_desc.font.name = FONT_NAME
    return txBox


def add_conclusion_box(slide, left, top, width, text, font_size=Pt(13)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = BLACK
    run.font.bold = True
    run.font.name = FONT_NAME
    return txBox


def add_table(slide, left, top, width, height, rows, cols, data,
              header_color=None, col_widths=None):
    if header_color is None:
        header_color = BLUE_HEADER
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if data[r][c] is not None else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(11)
                    run.font.name = FONT_NAME
                    if r == 0:
                        run.font.color.rgb = WHITE
                        run.font.bold = True
                    else:
                        run.font.color.rgb = BLACK
                        run.font.bold = False
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else RGBColor(0xF5, 0xF5, 0xF5)
            cell.margin_left = Pt(5)
            cell.margin_right = Pt(5)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
    return table_shape


def add_bar_chart(slide, left, top, width, height,
                  categories, values, title="", bar_colors=None):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', values)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    chart.chart_style = 2
    plot = chart.plots[0]
    plot.gap_width = 100
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CYAN
    series.has_data_labels = True
    dl = series.data_labels
    dl.font.size = Pt(13)
    dl.font.bold = True
    dl.font.color.rgb = BLACK
    dl.number_format = '0.#'
    dl.show_value = True
    dl.label_position = XL_LABEL_POSITION.OUTSIDE_END
    if bar_colors:
        for i, color in enumerate(bar_colors):
            pt = series.points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = color
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(12)
    cat_axis.tick_labels.font.name = FONT_NAME
    cat_axis.major_tick_mark = 2
    cat_axis.format.line.fill.background()
    val_axis = chart.value_axis
    val_axis.visible = False
    val_axis.major_tick_mark = 2
    val_axis.format.line.fill.background()
    val_axis.major_gridlines.format.line.fill.background()
    if title:
        chart.has_title = True
        ct = chart.chart_title.text_frame.paragraphs[0]
        ct.text = title
        ct.font.size = Pt(14)
        ct.font.bold = True
        ct.font.name = FONT_NAME
    else:
        chart.has_title = False
    return chart_frame


def add_callout_label(slide, left, top, text, bg_color=None, font_size=Pt(11)):
    if bg_color is None:
        bg_color = CYAN
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.3), Inches(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME
    return shape


def add_data_card(slide, left, top, width, height, value, label,
                  value_color=None, bg_color=None):
    if value_color is None:
        value_color = CYAN
    if bg_color is None:
        bg_color = WHITE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(3)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run1 = p1.add_run()
    run1.text = str(value)
    run1.font.size = Pt(24)
    run1.font.color.rgb = value_color
    run1.font.bold = True
    run1.font.name = FONT_NAME
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(10)
    run2.font.color.rgb = GRAY_TEXT
    run2.font.bold = False
    run2.font.name = FONT_NAME
    return shape


OUTPUT_PATH = r"C:\Users\Administrator\Desktop\notebook\my_notebooklm_with_ppteditable\projects\测试DEMO\最终文档\ppt_slides\slide_01.pptx"


# ── Slide 1 ──

def build_slide_1(slide):
    # Colors
    DARK_BLUE = RGBColor(0x1B, 0x2A, 0x49)
    GRAY_TEXT = RGBColor(0x66, 0x66, 0x66)
    LIGHT_GRAY = RGBColor(0x99, 0x99, 0x99)
    GREEN = RGBColor(0x2E, 0xA1, 0x54)
    RED = RGBColor(0xD9, 0x3A, 0x36)
    BLUE_ICON = RGBColor(0x4A, 0x86, 0xC8)
    BORDER_COLOR = RGBColor(0xE5, 0xE5, 0xE5)
    BLACK = RGBColor(0x00, 0x00, 0x00)

    # 1. Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(8.0), Inches(0.6))
    tf_title = title_box.text_frame
    p_title = tf_title.paragraphs[0]
    run_title = p_title.add_run()
    run_title.text = "简单场景验证：Claude Code 胜出"
    run_title.font.size = Pt(28)
    run_title.font.bold = True
    run_title.font.color.rgb = DARK_BLUE
    run_title.font.name = "Microsoft YaHei"

    # 2. Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(8.0), Inches(0.4))
    tf_sub = sub_box.text_frame
    p_sub = tf_sub.paragraphs[0]
    run_sub = p_sub.add_run()
    run_sub.text = "销售订单报表查询（低复杂度）效率对比"
    run_sub.font.size = Pt(16)
    run_sub.font.color.rgb = GRAY_TEXT
    run_sub.font.name = "Microsoft YaHei"

    # 3. Left Panel (White Box)
    left_panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.7), Inches(5.4), Inches(4.8))
    left_panel.fill.solid()
    left_panel.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    left_panel.line.color.rgb = BORDER_COLOR
    left_panel.line.width = Pt(1)

    # Helper for left panel text
    def add_left_text(y, runs):
        tb = slide.shapes.add_textbox(Inches(1.2), y, Inches(4.5), Inches(0.8))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.line_spacing = 1.3
        for text, is_bold, color in runs:
            r = p.add_run()
            r.text = text
            r.font.name = "Microsoft YaHei"
            r.font.size = Pt(13)
            if is_bold:
                r.font.bold = True
            if color:
                r.font.color.rgb = color

    # Item 1: Up Arrow Icon & Text
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.75), Inches(2.05), Inches(0.28), Inches(0.28))
    circle.fill.background()
    circle.line.color.rgb = GREEN
    circle.line.width = Pt(1.5)
    arrow = slide.shapes.add_shape(MSO_SHAPE.UP_ARROW, Inches(0.84), Inches(2.1), Inches(0.1), Inches(0.18))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = GREEN
    arrow.line.fill.background()
    
    add_left_text(Inches(1.95), [
        ("效率拐点：", True, BLACK),
        ("Claude Code 耗时 ", False, BLACK),
        ("30 分钟", True, BLACK),
        ("，较手写开发效率", False, BLACK),
        ("提升 50%", True, BLACK),
        ("。", False, BLACK)
    ])

    # Item 2: Warning Icon & Text
    triangle = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(0.75), Inches(3.05), Inches(0.28), Inches(0.28))
    triangle.fill.background()
    triangle.line.color.rgb = RED
    triangle.line.width = Pt(1.5)
    warn_tb = slide.shapes.add_textbox(Inches(0.75), Inches(3.1), Inches(0.28), Inches(0.28))
    warn_p = warn_tb.text_frame.paragraphs[0]
    warn_p.alignment = PP_ALIGN.CENTER
    warn_r = warn_p.add_run()
    warn_r.text = "!"
    warn_r.font.size = Pt(12)
    warn_r.font.bold = True
    warn_r.font.color.rgb = RED

    add_left_text(Inches(2.95), [
        ("稳定性差异：", True, BLACK),
        ("GitHub Copilot 运行出现 ", False, BLACK),
        ("Short Dump", True, RED),
        ("，而 Claude Code ", False, BLACK),
        ("运行正常", True, GREEN),
        ("。", False, BLACK)
    ])

    # Item 3: Chat Icon & Text
    chat1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGULAR_CALLOUT, Inches(0.75), Inches(4.1), Inches(0.22), Inches(0.18))
    chat1.fill.background()
    chat1.line.color.rgb = BLUE_ICON
    chat1.line.width = Pt(1.5)
    chat2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGULAR_CALLOUT, Inches(0.82), Inches(4.18), Inches(0.22), Inches(0.18))
    chat2.fill.background()
    chat2.line.color.rgb = BLUE_ICON
    chat2.line.width = Pt(1.5)

    add_left_text(Inches(3.95), [
        ("交互成本：", True, BLACK),
        ("Claude Code 仅需 ", False, BLACK),
        ("2 轮", True, BLACK),
        ("人工干预，远优于 Copilot 的 ", False, BLACK),
        ("5 轮", True, BLACK),
        ("以上。", False, BLACK)
    ])

    # Item 4: Link Icon & Text
    link1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.78), Inches(5.15), Inches(0.18), Inches(0.1))
    link1.rotation = 45
    link1.fill.background()
    link1.line.color.rgb = LIGHT_GRAY
    link1.line.width = Pt(1.5)
    link2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.86), Inches(5.23), Inches(0.18), Inches(0.1))
    link2.rotation = 45
    link2.fill.background()
    link2.line.color.rgb = LIGHT_GRAY
    link2.line.width = Pt(1.5)

    add_left_text(Inches(4.95), [
        ("核心瓶颈：", True, BLACK),
        ("两者初次生成均不可直接运行，仍需人工修正 ", False, BLACK),
        ("SQL", True, BLACK),
        (" 逻辑。", False, BLACK)
    ])

    # 4. Right Top Panel (Chart)
    chart_title = slide.shapes.add_textbox(Inches(6.4), Inches(1.7), Inches(4.0), Inches(0.4))
    p_ct = chart_title.text_frame.paragraphs[0]
    r_ct = p_ct.add_run()
    r_ct.text = "开发耗时对比（分钟）"
    r_ct.font.size = Pt(14)
    r_ct.font.bold = True
    r_ct.font.name = "Microsoft YaHei"

    chart_data = CategoryChartData()
    chart_data.categories = ['Claude Code', 'GitHub Copilot', '手写开发']
    chart_data.add_series('耗时', (30, 65, 60))

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, Inches(6.4), Inches(2.2), Inches(6.4), Inches(2.0), chart_data
    )
    chart = chart_shape.chart
    chart.has_legend = False

    val_axis = chart.value_axis
    val_axis.maximum_scale = 70
    val_axis.major_unit = 15
    val_axis.has_major_gridlines = True
    val_axis.major_gridlines.format.line.color.rgb = BORDER_COLOR
    val_axis.tick_labels.font.size = Pt(10)
    val_axis.tick_labels.font.color.rgb = GRAY_TEXT

    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(11)
    cat_axis.tick_labels.font.color.rgb = GRAY_TEXT
    cat_axis.has_major_gridlines = False

    series = chart.series[0]
    series.has_data_labels = True

    # Claude Code (Green)
    p0 = series.points[0]
    p0.format.fill.solid()
    p0.format.fill.fore_color.rgb = GREEN
    p0.data_label.has_text_frame = True
    p0.data_label.text_frame.text = "30 min ✅"
    p0.data_label.font.size = Pt(10)
    p0.data_label.font.bold = True

    # GitHub Copilot (Red)
    p1 = series.points[1]
    p1.format.fill.solid()
    p1.format.fill.fore_color.rgb = RED
    p1.data_label.has_text_frame = True
    p1.data_label.text_frame.text = "> 60 min ⚠️"
    p1.data_label.font.size = Pt(10)
    p1.data_label.font.bold = True
    p1.data_label.font.color.rgb = RED

    # 手写开发 (Red)
    p2 = series.points[2]
    p2.format.fill.solid()
    p2.format.fill.fore_color.rgb = RED
    p2.data_label.has_text_frame = True
    p2.data_label.text_frame.text = "60 min"
    p2.data_label.font.size = Pt(10)
    p2.data_label.font.bold = True

    # 5. Right Bottom Left Panel (Data Card)
    card1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.4), Inches(4.5), Inches(2.6), Inches(2.0))
    card1.fill.solid()
    card1.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    card1.line.color.rgb = BORDER_COLOR

    tb_50 = slide.shapes.add_textbox(Inches(6.5), Inches(4.8), Inches(2.0), Inches(0.6))
    p_50 = tb_50.text_frame.paragraphs[0]
    r_50 = p_50.add_run()
    r_50.text = "50%"
    r_50.font.size = Pt(40)
    r_50.font.bold = True
    r_50.font.color.rgb = GREEN

    tb_eff = slide.shapes.add_textbox(Inches(6.5), Inches(5.5), Inches(2.0), Inches(0.4))
    p_eff = tb_eff.text_frame.paragraphs[0]
    r_eff1 = p_eff.add_run()
    r_eff1.text = "效率提升 "
    r_eff1.font.size = Pt(16)
    r_eff1.font.bold = True
    r_eff1.font.name = "Microsoft YaHei"
    r_eff2 = p_eff.add_run()
    r_eff2.text = "↗"
    r_eff2.font.size = Pt(16)
    r_eff2.font.bold = True
    r_eff2.font.color.rgb = GREEN

    tb_vs = slide.shapes.add_textbox(Inches(6.5), Inches(6.0), Inches(2.4), Inches(0.3))
    p_vs = tb_vs.text_frame.paragraphs[0]
    r_vs = p_vs.add_run()
    r_vs.text = "Claude Code vs 手写开发"
    r_vs.font.size = Pt(10)
    r_vs.font.color.rgb = GRAY_TEXT
    r_vs.font.name = "Microsoft YaHei"

    # 6. Right Bottom Right Panel (Grid)
    card2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.3), Inches(4.5), Inches(3.5), Inches(2.0))
    card2.fill.solid()
    card2.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    card2.line.color.rgb = BORDER_COLOR

    line_v = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(11.05), Inches(4.5), Inches(11.05), Inches(6.5))
    line_v.line.color.rgb = BORDER_COLOR
    line_h = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(9.3), Inches(5.5), Inches(12.8), Inches(5.5))
    line_h.line.color.rgb = BORDER_COLOR

    def add_grid_cell(x, y, icon_text, icon_color, text, icon_size=24):
        tb = slide.shapes.add_textbox(x, y, Inches(1.75), Inches(1.0))
        tf = tb.text_frame
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = icon_text + "\n"
        r1.font.size = Pt(icon_size)
        r1.font.color.rgb = icon_color
        r1.font.name = "Segoe UI Emoji"

        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = text
        r2.font.size = Pt(11)
        r2.font.color.rgb = BLACK
        r2.font.name = "Microsoft YaHei"

    add_grid_cell(Inches(9.3), Inches(4.6), "❌", RED, "GitHub Copilot")
    add_grid_cell(Inches(11.05), Inches(4.6), "❌", RED, "Short Dump")
    add_grid_cell(Inches(9.3), Inches(5.6), "✅", GREEN, "Copilot: 5+ 轮")
    add_grid_cell(Inches(11.05), Inches(5.6), "👥", BLUE_ICON, "Claude Code: 2 轮")

    # 7. Page Number
    page_num = slide.shapes.add_textbox(Inches(11.8), Inches(6.9), Inches(1.0), Inches(0.3))
    p_page = page_num.text_frame.paragraphs[0]
    p_page.alignment = PP_ALIGN.RIGHT
    r_page = p_page.add_run()
    r_page.text = "Page 1 / 4"
    r_page.font.size = Pt(10)
    r_page.font.color.rgb = LIGHT_GRAY
    r_page.font.name = "Microsoft YaHei"



# ── Main ──
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

s0 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_1(s0)
prs.save(OUTPUT_PATH)
