from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# ── Color Palette ──
BLUE_HEADER = RGBColor(0x5B, 0x9B, 0xD5)
BLUE_DARK   = RGBColor(0x4A, 0x86, 0xC8)
CYAN        = RGBColor(0x00, 0xBC, 0xD4)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x33, 0x33, 0x33)
GRAY_TEXT   = RGBColor(0x55, 0x55, 0x55)
GRAY_BAR    = RGBColor(0xB0, 0xBE, 0xC5)
RED         = RGBColor(0xE5, 0x39, 0x35)
GREEN       = RGBColor(0x43, 0xA0, 0x47)
ORANGE      = RGBColor(0xFF, 0x98, 0x00)

ICON_BG  = RGBColor(0xE3, 0xE8, 0xED)
ICON_FG  = RGBColor(0x54, 0x6E, 0x7A)
FONT_NAME = "Microsoft YaHei"
SLIDE_WIDTH = Inches(13.333)
HEADER_H    = Inches(0.75)
SUBTITLE_Y  = Inches(0.95)


def add_header_banner(slide, title_text, bg_color=None):
    if bg_color is None:
        bg_color = BLUE_HEADER
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, HEADER_H
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = bg_color
    banner.line.fill.background()
    tf = banner.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(26)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME


def add_subtitle(slide, text, left, top, width=Inches(12), font_size=Pt(18)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.font.name = FONT_NAME
    return txBox


def add_icon_box(slide, left, top, symbol, size=Inches(0.48)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ICON_BG
    shape.line.fill.background()
    shape.adjustments[0] = 0.25
    tf = shape.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = symbol
    p.font.size = Pt(18)
    p.font.color.rgb = ICON_FG
    p.font.bold = False
    return shape


def add_bullet_item(slide, left, top, symbol, label, description,
                    width=Inches(5.5), desc_size=Pt(13)):
    add_icon_box(slide, left, top, symbol)
    text_left = left + Inches(0.65)
    txBox = slide.shapes.add_textbox(text_left, top - Inches(0.02), width, Inches(0.65))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run_label = p.add_run()
    run_label.text = label + "\uff1a"
    run_label.font.size = Pt(14)
    run_label.font.color.rgb = BLACK
    run_label.font.bold = True
    run_label.font.name = FONT_NAME
    run_desc = p.add_run()
    run_desc.text = description
    run_desc.font.size = desc_size
    run_desc.font.color.rgb = GRAY_TEXT
    run_desc.font.bold = False
    run_desc.font.name = FONT_NAME
    return txBox


def add_conclusion_box(slide, left, top, width, text, font_size=Pt(13)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = BLACK
    run.font.bold = True
    run.font.name = FONT_NAME
    return txBox


def add_table(slide, left, top, width, height, rows, cols, data,
              header_color=None, col_widths=None):
    if header_color is None:
        header_color = BLUE_HEADER
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if data[r][c] is not None else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(11)
                    run.font.name = FONT_NAME
                    if r == 0:
                        run.font.color.rgb = WHITE
                        run.font.bold = True
                    else:
                        run.font.color.rgb = BLACK
                        run.font.bold = False
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else RGBColor(0xF5, 0xF5, 0xF5)
            cell.margin_left = Pt(5)
            cell.margin_right = Pt(5)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
    return table_shape


def add_bar_chart(slide, left, top, width, height,
                  categories, values, title="", bar_colors=None):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', values)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    chart.chart_style = 2
    plot = chart.plots[0]
    plot.gap_width = 100
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CYAN
    series.has_data_labels = True
    dl = series.data_labels
    dl.font.size = Pt(13)
    dl.font.bold = True
    dl.font.color.rgb = BLACK
    dl.number_format = '0.#'
    dl.show_value = True
    dl.label_position = XL_LABEL_POSITION.OUTSIDE_END
    if bar_colors:
        for i, color in enumerate(bar_colors):
            pt = series.points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = color
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(12)
    cat_axis.tick_labels.font.name = FONT_NAME
    cat_axis.major_tick_mark = 2
    cat_axis.format.line.fill.background()
    val_axis = chart.value_axis
    val_axis.visible = False
    val_axis.major_tick_mark = 2
    val_axis.format.line.fill.background()
    val_axis.major_gridlines.format.line.fill.background()
    if title:
        chart.has_title = True
        ct = chart.chart_title.text_frame.paragraphs[0]
        ct.text = title
        ct.font.size = Pt(14)
        ct.font.bold = True
        ct.font.name = FONT_NAME
    else:
        chart.has_title = False
    return chart_frame


def add_callout_label(slide, left, top, text, bg_color=None, font_size=Pt(11)):
    if bg_color is None:
        bg_color = CYAN
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.3), Inches(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME
    return shape


def add_data_card(slide, left, top, width, height, value, label,
                  value_color=None, bg_color=None):
    if value_color is None:
        value_color = CYAN
    if bg_color is None:
        bg_color = WHITE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(3)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run1 = p1.add_run()
    run1.text = str(value)
    run1.font.size = Pt(24)
    run1.font.color.rgb = value_color
    run1.font.bold = True
    run1.font.name = FONT_NAME
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(10)
    run2.font.color.rgb = GRAY_TEXT
    run2.font.bold = False
    run2.font.name = FONT_NAME
    return shape


OUTPUT_PATH = r"C:\Users\Administrator\Desktop\notebook\my_notebooklm_with_ppteditable\projects\测试DEMO\最终文档\ppt_slides\slide_02.pptx"


# ── Slide 1 ──

def build_slide_1(slide):
    # Colors
    DARK_BLUE = RGBColor(0x1A, 0x36, 0x5D)
    RED = RGBColor(0xE5, 0x39, 0x35)
    LIGHT_RED = RGBColor(0xFD, 0xED, 0xEC)
    LIGHT_GREEN = RGBColor(0xE8, 0xF5, 0xE9)
    GRAY_TEXT = RGBColor(0x55, 0x55, 0x55)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xE0)
    
    # 1. Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(10), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "中等场景瓶颈：AI “幻觉”引发崩溃"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = "Microsoft YaHei"

    # 2. Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(10), Inches(0.5))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "采购配额维护（中等复杂度）开发验证"
    p.font.size = Pt(18)
    p.font.color.rgb = GRAY_TEXT
    p.font.name = "Microsoft YaHei"

    # 3. Left Bullets
    bullets_data = [
        ("🔗", "逻辑误区", "Copilot 完全混淆业务概念（配额误作货源），代码完全不可用。"),
        ("⚠️", "幻觉严重", "Claude Code 虚构字段比例高达 50%，导致 21 个连锁语法错误。"),
        ("🚫", "效率归零", "AI 修复成本大于重写成本，整体效率对比手写无任何提升。"),
        ("⚠️", "规则缺失", "AI 无法准确遵循 SAP 函数接口规范，直接忽略“优先 API”的指令。")
    ]
    
    start_y = 1.8
    for icon, label, desc in bullets_data:
        # Icon
        icon_box = slide.shapes.add_textbox(Inches(0.5), Inches(start_y), Inches(0.6), Inches(0.6))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(28)
        p.font.color.rgb = RED
        p.alignment = PP_ALIGN.CENTER
        
        # Label
        label_box = slide.shapes.add_textbox(Inches(1.3), Inches(start_y), Inches(4.5), Inches(0.4))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.name = "Microsoft YaHei"
        
        # Desc
        desc_box = slide.shapes.add_textbox(Inches(1.3), Inches(start_y + 0.35), Inches(4.8), Inches(0.8))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY_TEXT
        p.font.name = "Microsoft YaHei"
        
        start_y += 1.3

    # 4. Right Data Card
    card_left = Inches(6.8)
    card_top = Inches(1.4)
    card_width = Inches(6.0)
    card_height = Inches(1.4)
    
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_left, card_top, card_width, card_height)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = LIGHT_GRAY
    card.line.width = Pt(1)

    # Card Text - 50%
    val_box = slide.shapes.add_textbox(card_left + Inches(0.2), card_top + Inches(0.1), Inches(3), Inches(0.8))
    tf = val_box.text_frame
    p = tf.paragraphs[0]
    p.text = "50% ↗"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RED
    p.font.name = "Arial"

    # Card Text - Label
    lbl_box = slide.shapes.add_textbox(card_left + Inches(0.2), card_top + Inches(0.9), Inches(4), Inches(0.4))
    tf = lbl_box.text_frame
    p = tf.paragraphs[0]
    p.text = "虚构字段比例 (Fictional Field Ratio)"
    p.font.size = Pt(12)
    p.font.color.rgb = GRAY_TEXT
    p.font.name = "Microsoft YaHei"

    # Card Icon - Warning
    warn_box = slide.shapes.add_textbox(card_left + Inches(4.8), card_top + Inches(0.2), Inches(1), Inches(1))
    tf = warn_box.text_frame
    p = tf.paragraphs[0]
    p.text = "⚠️"
    p.font.size = Pt(60)
    p.font.color.rgb = RED
    p.alignment = PP_ALIGN.CENTER

    # 5. Right Table
    table_top = Inches(3.0)
    
    # Table Title
    tt_box = slide.shapes.add_textbox(card_left, table_top, card_width, Inches(0.4))
    tf = tt_box.text_frame
    p = tf.paragraphs[0]
    p.text = "字段对比示例"
    p.font.size = Pt(14)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Microsoft YaHei"

    # Table Shape
    rows = 6
    cols = 2
    table_shape = slide.shapes.add_table(rows, cols, card_left, table_top + Inches(0.4), card_width, Inches(2.0)).table
    table_shape.columns[0].width = Inches(3.0)
    table_shape.columns[1].width = Inches(3.0)

    table_data = [
        ("❌ 虚构字段", "✅ 正确字段"),
        ("MSEG-QUOTA_ID (不存在)", "EKKO-EBELN (采购凭证)"),
        ("EKPO-ALLOC_QTY (错误逻辑)", "EKPO-MENGE (数量)"),
        ("EKPO-ALLOC_QTY (错误逻辑)", "EKPO-MENGE (数量)"),
        ("EKPO-BLG_ID (不存在)", "EKKO-QUOTA (数量)"),
        ("...", "...")
    ]

    for r in range(rows):
        for c in range(cols):
            cell = table_shape.cell(r, c)
            cell.text = table_data[r][c]
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.font.name = "Microsoft YaHei"
            
            if r == 0:
                p.font.bold = True
            
            # Set background colors
            if r == 0:
                if c == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = LIGHT_RED
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = LIGHT_GREEN
            else:
                if c == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xFE, 0xF5, 0xF5)
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xF1, 0xF8, 0xF1)

    # 6. Right Flowchart
    flow_top = Inches(5.7)
    
    # Box 1
    b1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_left, flow_top, Inches(1.2), Inches(1.2))
    b1.fill.solid()
    b1.fill.fore_color.rgb = RED
    b1.line.fill.background()
    tf = b1.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = "50%"
    p1.font.size = Pt(20)
    p1.font.color.rgb = WHITE
    p1.font.bold = True
    p1.alignment = PP_ALIGN.CENTER
    p1.font.name = "Arial"
    p2 = tf.add_paragraph()
    p2.text = "虚构字段"
    p2.font.size = Pt(12)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER
    p2.font.name = "Microsoft YaHei"

    # Arrow 1
    a1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, card_left + Inches(1.3), flow_top + Inches(0.5), Inches(0.3), Inches(0.2))
    a1.fill.solid()
    a1.fill.fore_color.rgb = GRAY_TEXT
    a1.line.fill.background()

    # Box 2
    b2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_left + Inches(1.7), flow_top, Inches(1.6), Inches(1.2))
    b2.fill.solid()
    b2.fill.fore_color.rgb = RED
    b2.line.fill.background()
    tf = b2.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = "🔗"
    p1.font.size = Pt(20)
    p1.font.color.rgb = WHITE
    p1.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "21 个连锁语法错误"
    p2.font.size = Pt(11)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER
    p2.font.name = "Microsoft YaHei"

    # Arrow 2
    a2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, card_left + Inches(3.4), flow_top + Inches(0.5), Inches(0.3), Inches(0.2))
    a2.fill.solid()
    a2.fill.fore_color.rgb = GRAY_TEXT
    a2.line.fill.background()

    # Box 3
    b3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_left + Inches(3.8), flow_top, Inches(2.2), Inches(1.2))
    b3.fill.solid()
    b3.fill.fore_color.rgb = RED
    b3.line.fill.background()
    tf = b3.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = "🚫"
    p1.font.size = Pt(20)
    p1.font.color.rgb = WHITE
    p1.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "AI 修复成本 >> 重写成本，\n整体效率对比手写无提升"
    p2.font.size = Pt(10)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER
    p2.font.name = "Microsoft YaHei"

    # 7. Footer
    footer = slide.shapes.add_textbox(Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = "PAGE 2 OF 4"
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY_TEXT
    p.alignment = PP_ALIGN.RIGHT



# ── Main ──
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

s0 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_1(s0)
prs.save(OUTPUT_PATH)
