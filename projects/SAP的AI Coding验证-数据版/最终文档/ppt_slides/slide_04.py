def build_slide(slide):
    from pptx.enum.shapes import MSO_CONNECTOR
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE

    # 颜色定义
    DARK_BLUE = RGBColor(0x1B, 0x3B, 0x5A)
    GRAY_TEXT = RGBColor(0x55, 0x55, 0x55)
    LIGHT_GRAY_BG = RGBColor(0xF9, 0xFA, 0xFB)
    
    GREEN_FILL = RGBColor(0x6E, 0xE7, 0xB7)
    GREEN_BORDER = RGBColor(0x34, 0xD3, 0x99)
    
    YELLOW_BORDER = RGBColor(0xFB, 0xBF, 0x24)
    
    RED_FILL = RGBColor(0xFC, 0xCA, 0xCA)
    RED_BORDER = RGBColor(0xEF, 0x44, 0x44)
    
    ORANGE_TEXT = RGBColor(0xD9, 0x77, 0x06)
    ORANGE_TREND = RGBColor(0xED, 0x7D, 0x31)

    # 1. 标题和副标题
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8.0), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "结论：复杂度决定 AI 的价值边界"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = "Microsoft YaHei"

    subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(8.0), Inches(0.4))
    tf_sub = subtitle_box.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "SAP ABAP AI Coding 效率与质量总览"
    p_sub.font.size = Pt(16)
    p_sub.font.color.rgb = GRAY_TEXT
    p_sub.font.name = "Microsoft YaHei"

    # 页码
    page_box = slide.shapes.add_textbox(Inches(12.0), Inches(0.4), Inches(1.0), Inches(0.4))
    p_page = page_box.text_frame.paragraphs[0]
    p_page.text = "第 4 页"
    p_page.font.size = Pt(12)
    p_page.font.color.rgb = GRAY_TEXT
    p_page.font.name = "Microsoft YaHei"

    # 2. 坐标轴
    # Y轴
    y_axis = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.3), Inches(5.4), Inches(1.3), Inches(1.7))
    y_axis.line.color.rgb = DARK_BLUE
    y_axis.line.width = Pt(2)
    y_axis.line.end_arrowhead = True

    # X轴
    x_axis = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.3), Inches(5.4), Inches(8.5), Inches(5.4))
    x_axis.line.color.rgb = DARK_BLUE
    x_axis.line.width = Pt(2)
    x_axis.line.end_arrowhead = True

    # 坐标轴标签
    labels = [
        ("正", 0.9, 1.9),
        ("负", 0.9, 5.1),
        ("低", 1.4, 5.5),
        ("高", 8.1, 5.5),
        ("复杂度 (Complexity)", 3.8, 5.5)
    ]
    for text, left, top in labels:
        tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(2.0), Inches(0.4))
        p = tb.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.name = "Microsoft YaHei"

    # Y轴标题 (垂直排列)
    y_title_box = slide.shapes.add_textbox(Inches(0.7), Inches(3.0), Inches(0.5), Inches(2.0))
    p_y = y_title_box.text_frame.paragraphs[0]
    p_y.text = "提\n效\n程\n度"
    p_y.font.size = Pt(12)
    p_y.font.bold = True
    p_y.alignment = PP_ALIGN.CENTER
    p_y.font.name = "Microsoft YaHei"

    # 3. 矩阵象限 (圆角矩形)
    def add_quadrant(left, top, text, fill_color, border_color):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(3.2), Inches(1.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.color.rgb = border_color
        shape.line.width = Pt(2)
        
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Microsoft YaHei"
        return shape

    # 左上：简单场景
    add_quadrant(1.5, 1.9, "简单场景\n(提效 50%)", GREEN_FILL, GREEN_BORDER)
    
    # 左下：中等场景
    add_quadrant(1.5, 3.6, "中等场景\n(零提升)", LIGHT_GRAY_BG, YELLOW_BORDER)
    
    # 右上：空白背景框
    shape_tr = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.9), Inches(1.9), Inches(3.2), Inches(1.5))
    shape_tr.fill.solid()
    shape_tr.fill.fore_color.rgb = LIGHT_GRAY_BG
    shape_tr.line.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)
    shape_tr.line.width = Pt(1)
    
    # 右下：复杂场景
    add_quadrant(4.9, 3.6, "复杂场景\n(反降 60%)", RED_FILL, RED_BORDER)

    # 4. 趋势线与拐点
    trend_line = slide.shapes.add_connector(MSO_CONNECTOR.CURVE, Inches(3.2), Inches(2.6), Inches(7.5), Inches(4.8))
    trend_line.line.color.rgb = ORANGE_TREND
    trend_line.line.width = Pt(6)
    trend_line.line.end_arrowhead = True

    # 拐点黄点
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.7), Inches(3.2), Inches(0.25), Inches(0.25))
    dot.fill.solid()
    dot.fill.fore_color.rgb = YELLOW_BORDER
    dot.line.fill.background()

    # 拐点文字
    trend_text = slide.shapes.add_textbox(Inches(5.3), Inches(2.8), Inches(1.5), Inches(0.4))
    p_trend = trend_text.text_frame.paragraphs[0]
    p_trend.text = "价值拐点趋势"
    p_trend.font.size = Pt(12)
    p_trend.font.bold = True
    p_trend.font.name = "Microsoft YaHei"

    # 5. 右侧金句卡片
    # 主卡片框
    card_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.1), Inches(1.8), Inches(3.6), Inches(3.8))
    card_bg.fill.solid()
    card_bg.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    card_bg.line.color.rgb = DARK_BLUE
    card_bg.line.width = Pt(5)

    # 顶部标签 "金句卡片"
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.1), Inches(1.6), Inches(1.6), Inches(0.45))
    badge.fill.solid()
    badge.fill.fore_color.rgb = DARK_BLUE
    badge.line.fill.background()
    p_badge = badge.text_frame.paragraphs[0]
    p_badge.text = "金句卡片"
    p_badge.font.size = Pt(14)
    p_badge.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p_badge.alignment = PP_ALIGN.CENTER
    p_badge.font.name = "Microsoft YaHei"

    # 左上角引号图标 (用文字模拟)
    quote_mark = slide.shapes.add_textbox(Inches(9.3), Inches(2.0), Inches(0.8), Inches(0.8))
    p_quote = quote_mark.text_frame.paragraphs[0]
    p_quote.text = "“"
    p_quote.font.size = Pt(60)
    p_quote.font.color.rgb = DARK_BLUE
    p_quote.font.bold = True
    p_quote.font.name = "Microsoft YaHei"

    # 右上角警告图标 (三角形 + 感叹号)
    warning_tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(12.0), Inches(2.2), Inches(0.4), Inches(0.35))
    warning_tri.fill.solid()
    warning_tri.fill.fore_color.rgb = ORANGE_TREND
    warning_tri.line.fill.background()
    
    warning_text = slide.shapes.add_textbox(Inches(12.0), Inches(2.2), Inches(0.4), Inches(0.35))
    p_warn = warning_text.text_frame.paragraphs[0]
    p_warn.text = "!"
    p_warn.font.size = Pt(14)
    p_warn.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p_warn.font.bold = True
    p_warn.alignment = PP_ALIGN.CENTER
    p_warn.font.name = "Microsoft YaHei"

    # 卡片正文
    card_text = slide.shapes.add_textbox(Inches(9.4), Inches(2.8), Inches(3.0), Inches(2.5))
    tf_card = card_text.text_frame
    tf_card.word_wrap = True
    
    # 第一段
    p1 = tf_card.paragraphs[0]
    run1 = p1.add_run()
    run1.text = "专家建议："
    run1.font.size = Pt(14)
    run1.font.bold = True
    run1.font.color.rgb = RGBColor(0, 0, 0)
    run1.font.name = "Microsoft YaHei"
    
    run2 = p1.add_run()
    run2.text = "现阶段 AI 仅适用于简单逻辑片段，中高复杂度开发仍需资深顾问人工把控。"
    run2.font.size = Pt(14)
    run2.font.bold = True
    run2.font.color.rgb = ORANGE_TEXT
    run2.font.name = "Microsoft YaHei"
    
    # 换行间距
    p_space = tf_card.add_paragraph()
    p_space.text = ""
    p_space.font.size = Pt(8)

    # 第二段
    p2 = tf_card.add_paragraph()
    run3 = p2.add_run()
    run3.text = "人类智慧在核心业务逻辑中不可替代。"
    run3.font.size = Pt(14)
    run3.font.bold = True
    run3.font.color.rgb = RGBColor(0, 0, 0)
    run3.font.name = "Microsoft YaHei"

    # 6. 底部要点说明
    bullet_box = slide.shapes.add_textbox(Inches(1.5), Inches(6.0), Inches(11.0), Inches(1.2))
    tf_bullet = bullet_box.text_frame
    tf_bullet.word_wrap = True

    # 要点 1
    p_b1 = tf_bullet.paragraphs[0]
    p_b1.space_after = Pt(10)
    run_b1_1 = p_b1.add_run()
    run_b1_1.text = "• 核心瓶颈："
    run_b1_1.font.size = Pt(14)
    run_b1_1.font.bold = True
    run_b1_1.font.name = "Microsoft YaHei"
    
    run_b1_2 = p_b1.add_run()
    run_b1_2.text = "AI 对 SAP 专用数据字典（DDIC）的“幻觉”是阻碍生产力的首要因素。"
    run_b1_2.font.size = Pt(14)
    run_b1_2.font.name = "Microsoft YaHei"

    # 要点 2
    p_b2 = tf_bullet.add_paragraph()
    run_b2_1 = p_b2.add_run()
    run_b2_1.text = "• 工具选型："
    run_b2_1.font.size = Pt(14)
    run_b2_1.font.bold = True
    run_b2_1.font.name = "Microsoft YaHei"
    
    run_b2_2 = p_b2.add_run()
    run_b2_2.text = "Claude Code 在理解力、多文件处理及文档解析上全面优于 GitHub Copilot。"
    run_b2_2.font.size = Pt(14)
    run_b2_2.font.name = "Microsoft YaHei"