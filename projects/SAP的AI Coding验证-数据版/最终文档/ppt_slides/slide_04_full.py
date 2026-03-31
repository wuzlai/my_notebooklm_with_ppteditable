from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# ── Color Palette ──
BLUE_HEADER = RGBColor(0x5B, 0x9B, 0xD5)
BLUE_DARK   = RGBColor(0x4A, 0x86, 0xC8)
CYAN        = RGBColor(0x00, 0xBC, 0xD4)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x33, 0x33, 0x33)
GRAY_TEXT   = RGBColor(0x55, 0x55, 0x55)
GRAY_BAR    = RGBColor(0xB0, 0xBE, 0xC5)
RED         = RGBColor(0xE5, 0x39, 0x35)
GREEN       = RGBColor(0x43, 0xA0, 0x47)
ORANGE      = RGBColor(0xFF, 0x98, 0x00)

ICON_BG  = RGBColor(0xE3, 0xE8, 0xED)
ICON_FG  = RGBColor(0x54, 0x6E, 0x7A)
FONT_NAME = "Microsoft YaHei"
SLIDE_WIDTH = Inches(13.333)
HEADER_H    = Inches(0.75)
SUBTITLE_Y  = Inches(0.95)


def add_header_banner(slide, title_text, bg_color=None):
    if bg_color is None:
        bg_color = BLUE_HEADER
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, HEADER_H
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = bg_color
    banner.line.fill.background()
    tf = banner.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(26)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME


def add_subtitle(slide, text, left, top, width=Inches(12), font_size=Pt(18)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.font.name = FONT_NAME
    return txBox


def add_icon_box(slide, left, top, symbol, size=Inches(0.48)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ICON_BG
    shape.line.fill.background()
    shape.adjustments[0] = 0.25
    tf = shape.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = symbol
    p.font.size = Pt(18)
    p.font.color.rgb = ICON_FG
    p.font.bold = False
    return shape


def add_bullet_item(slide, left, top, symbol, label, description,
                    width=Inches(5.5), desc_size=Pt(13)):
    add_icon_box(slide, left, top, symbol)
    text_left = left + Inches(0.65)
    txBox = slide.shapes.add_textbox(text_left, top - Inches(0.02), width, Inches(0.65))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run_label = p.add_run()
    run_label.text = label + "\uff1a"
    run_label.font.size = Pt(14)
    run_label.font.color.rgb = BLACK
    run_label.font.bold = True
    run_label.font.name = FONT_NAME
    run_desc = p.add_run()
    run_desc.text = description
    run_desc.font.size = desc_size
    run_desc.font.color.rgb = GRAY_TEXT
    run_desc.font.bold = False
    run_desc.font.name = FONT_NAME
    return txBox


def add_conclusion_box(slide, left, top, width, text, font_size=Pt(13)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = BLACK
    run.font.bold = True
    run.font.name = FONT_NAME
    return txBox


def add_table(slide, left, top, width, height, rows, cols, data,
              header_color=None, col_widths=None):
    if header_color is None:
        header_color = BLUE_HEADER
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if data[r][c] is not None else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(11)
                    run.font.name = FONT_NAME
                    if r == 0:
                        run.font.color.rgb = WHITE
                        run.font.bold = True
                    else:
                        run.font.color.rgb = BLACK
                        run.font.bold = False
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else RGBColor(0xF5, 0xF5, 0xF5)
            cell.margin_left = Pt(5)
            cell.margin_right = Pt(5)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
    return table_shape


def add_bar_chart(slide, left, top, width, height,
                  categories, values, title="", bar_colors=None):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', values)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    chart.chart_style = 2
    plot = chart.plots[0]
    plot.gap_width = 100
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CYAN
    series.has_data_labels = True
    dl = series.data_labels
    dl.font.size = Pt(13)
    dl.font.bold = True
    dl.font.color.rgb = BLACK
    dl.number_format = '0.#'
    dl.show_value = True
    dl.label_position = XL_LABEL_POSITION.OUTSIDE_END
    if bar_colors:
        for i, color in enumerate(bar_colors):
            pt = series.points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = color
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(12)
    cat_axis.tick_labels.font.name = FONT_NAME
    cat_axis.major_tick_mark = 2
    cat_axis.format.line.fill.background()
    val_axis = chart.value_axis
    val_axis.visible = False
    val_axis.major_tick_mark = 2
    val_axis.format.line.fill.background()
    val_axis.major_gridlines.format.line.fill.background()
    if title:
        chart.has_title = True
        ct = chart.chart_title.text_frame.paragraphs[0]
        ct.text = title
        ct.font.size = Pt(14)
        ct.font.bold = True
        ct.font.name = FONT_NAME
    else:
        chart.has_title = False
    return chart_frame


def add_callout_label(slide, left, top, text, bg_color=None, font_size=Pt(11)):
    if bg_color is None:
        bg_color = CYAN
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.3), Inches(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME
    return shape


def add_data_card(slide, left, top, width, height, value, label,
                  value_color=None, bg_color=None):
    if value_color is None:
        value_color = CYAN
    if bg_color is None:
        bg_color = WHITE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(3)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run1 = p1.add_run()
    run1.text = str(value)
    run1.font.size = Pt(24)
    run1.font.color.rgb = value_color
    run1.font.bold = True
    run1.font.name = FONT_NAME
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(10)
    run2.font.color.rgb = GRAY_TEXT
    run2.font.bold = False
    run2.font.name = FONT_NAME
    return shape


OUTPUT_PATH = r"C:\Users\Administrator\Desktop\notebook\my_notebooklm_with_ppteditable\projects\SAP的AI Coding验证-数据版\最终文档\ppt_slides\slide_04.pptx"


# ── Slide 1 ──

def build_slide_1(slide):
    from pptx.enum.shapes import MSO_CONNECTOR
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE

    # 颜色定义
    DARK_BLUE = RGBColor(0x1B, 0x3B, 0x5A)
    GRAY_TEXT = RGBColor(0x55, 0x55, 0x55)
    LIGHT_GRAY_BG = RGBColor(0xF9, 0xFA, 0xFB)
    
    GREEN_FILL = RGBColor(0x6E, 0xE7, 0xB7)
    GREEN_BORDER = RGBColor(0x34, 0xD3, 0x99)
    
    YELLOW_BORDER = RGBColor(0xFB, 0xBF, 0x24)
    
    RED_FILL = RGBColor(0xFC, 0xCA, 0xCA)
    RED_BORDER = RGBColor(0xEF, 0x44, 0x44)
    
    ORANGE_TEXT = RGBColor(0xD9, 0x77, 0x06)
    ORANGE_TREND = RGBColor(0xED, 0x7D, 0x31)

    # 1. 标题和副标题
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8.0), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "结论：复杂度决定 AI 的价值边界"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = "Microsoft YaHei"

    subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(8.0), Inches(0.4))
    tf_sub = subtitle_box.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "SAP ABAP AI Coding 效率与质量总览"
    p_sub.font.size = Pt(16)
    p_sub.font.color.rgb = GRAY_TEXT
    p_sub.font.name = "Microsoft YaHei"

    # 页码
    page_box = slide.shapes.add_textbox(Inches(12.0), Inches(0.4), Inches(1.0), Inches(0.4))
    p_page = page_box.text_frame.paragraphs[0]
    p_page.text = "第 4 页"
    p_page.font.size = Pt(12)
    p_page.font.color.rgb = GRAY_TEXT
    p_page.font.name = "Microsoft YaHei"

    # 2. 坐标轴
    # Y轴
    y_axis = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.3), Inches(5.4), Inches(1.3), Inches(1.7))
    y_axis.line.color.rgb = DARK_BLUE
    y_axis.line.width = Pt(2)
    y_axis.line.end_arrowhead = True

    # X轴
    x_axis = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.3), Inches(5.4), Inches(8.5), Inches(5.4))
    x_axis.line.color.rgb = DARK_BLUE
    x_axis.line.width = Pt(2)
    x_axis.line.end_arrowhead = True

    # 坐标轴标签
    labels = [
        ("正", 0.9, 1.9),
        ("负", 0.9, 5.1),
        ("低", 1.4, 5.5),
        ("高", 8.1, 5.5),
        ("复杂度 (Complexity)", 3.8, 5.5)
    ]
    for text, left, top in labels:
        tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(2.0), Inches(0.4))
        p = tb.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.name = "Microsoft YaHei"

    # Y轴标题 (垂直排列)
    y_title_box = slide.shapes.add_textbox(Inches(0.7), Inches(3.0), Inches(0.5), Inches(2.0))
    p_y = y_title_box.text_frame.paragraphs[0]
    p_y.text = "提\n效\n程\n度"
    p_y.font.size = Pt(12)
    p_y.font.bold = True
    p_y.alignment = PP_ALIGN.CENTER
    p_y.font.name = "Microsoft YaHei"

    # 3. 矩阵象限 (圆角矩形)
    def add_quadrant(left, top, text, fill_color, border_color):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(3.2), Inches(1.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.color.rgb = border_color
        shape.line.width = Pt(2)
        
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Microsoft YaHei"
        return shape

    # 左上：简单场景
    add_quadrant(1.5, 1.9, "简单场景\n(提效 50%)", GREEN_FILL, GREEN_BORDER)
    
    # 左下：中等场景
    add_quadrant(1.5, 3.6, "中等场景\n(零提升)", LIGHT_GRAY_BG, YELLOW_BORDER)
    
    # 右上：空白背景框
    shape_tr = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.9), Inches(1.9), Inches(3.2), Inches(1.5))
    shape_tr.fill.solid()
    shape_tr.fill.fore_color.rgb = LIGHT_GRAY_BG
    shape_tr.line.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)
    shape_tr.line.width = Pt(1)
    
    # 右下：复杂场景
    add_quadrant(4.9, 3.6, "复杂场景\n(反降 60%)", RED_FILL, RED_BORDER)

    # 4. 趋势线与拐点
    trend_line = slide.shapes.add_connector(MSO_CONNECTOR.CURVE, Inches(3.2), Inches(2.6), Inches(7.5), Inches(4.8))
    trend_line.line.color.rgb = ORANGE_TREND
    trend_line.line.width = Pt(6)
    trend_line.line.end_arrowhead = True

    # 拐点黄点
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.7), Inches(3.2), Inches(0.25), Inches(0.25))
    dot.fill.solid()
    dot.fill.fore_color.rgb = YELLOW_BORDER
    dot.line.fill.background()

    # 拐点文字
    trend_text = slide.shapes.add_textbox(Inches(5.3), Inches(2.8), Inches(1.5), Inches(0.4))
    p_trend = trend_text.text_frame.paragraphs[0]
    p_trend.text = "价值拐点趋势"
    p_trend.font.size = Pt(12)
    p_trend.font.bold = True
    p_trend.font.name = "Microsoft YaHei"

    # 5. 右侧金句卡片
    # 主卡片框
    card_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.1), Inches(1.8), Inches(3.6), Inches(3.8))
    card_bg.fill.solid()
    card_bg.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    card_bg.line.color.rgb = DARK_BLUE
    card_bg.line.width = Pt(5)

    # 顶部标签 "金句卡片"
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.1), Inches(1.6), Inches(1.6), Inches(0.45))
    badge.fill.solid()
    badge.fill.fore_color.rgb = DARK_BLUE
    badge.line.fill.background()
    p_badge = badge.text_frame.paragraphs[0]
    p_badge.text = "金句卡片"
    p_badge.font.size = Pt(14)
    p_badge.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p_badge.alignment = PP_ALIGN.CENTER
    p_badge.font.name = "Microsoft YaHei"

    # 左上角引号图标 (用文字模拟)
    quote_mark = slide.shapes.add_textbox(Inches(9.3), Inches(2.0), Inches(0.8), Inches(0.8))
    p_quote = quote_mark.text_frame.paragraphs[0]
    p_quote.text = "“"
    p_quote.font.size = Pt(60)
    p_quote.font.color.rgb = DARK_BLUE
    p_quote.font.bold = True
    p_quote.font.name = "Microsoft YaHei"

    # 右上角警告图标 (三角形 + 感叹号)
    warning_tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(12.0), Inches(2.2), Inches(0.4), Inches(0.35))
    warning_tri.fill.solid()
    warning_tri.fill.fore_color.rgb = ORANGE_TREND
    warning_tri.line.fill.background()
    
    warning_text = slide.shapes.add_textbox(Inches(12.0), Inches(2.2), Inches(0.4), Inches(0.35))
    p_warn = warning_text.text_frame.paragraphs[0]
    p_warn.text = "!"
    p_warn.font.size = Pt(14)
    p_warn.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p_warn.font.bold = True
    p_warn.alignment = PP_ALIGN.CENTER
    p_warn.font.name = "Microsoft YaHei"

    # 卡片正文
    card_text = slide.shapes.add_textbox(Inches(9.4), Inches(2.8), Inches(3.0), Inches(2.5))
    tf_card = card_text.text_frame
    tf_card.word_wrap = True
    
    # 第一段
    p1 = tf_card.paragraphs[0]
    run1 = p1.add_run()
    run1.text = "专家建议："
    run1.font.size = Pt(14)
    run1.font.bold = True
    run1.font.color.rgb = RGBColor(0, 0, 0)
    run1.font.name = "Microsoft YaHei"
    
    run2 = p1.add_run()
    run2.text = "现阶段 AI 仅适用于简单逻辑片段，中高复杂度开发仍需资深顾问人工把控。"
    run2.font.size = Pt(14)
    run2.font.bold = True
    run2.font.color.rgb = ORANGE_TEXT
    run2.font.name = "Microsoft YaHei"
    
    # 换行间距
    p_space = tf_card.add_paragraph()
    p_space.text = ""
    p_space.font.size = Pt(8)

    # 第二段
    p2 = tf_card.add_paragraph()
    run3 = p2.add_run()
    run3.text = "人类智慧在核心业务逻辑中不可替代。"
    run3.font.size = Pt(14)
    run3.font.bold = True
    run3.font.color.rgb = RGBColor(0, 0, 0)
    run3.font.name = "Microsoft YaHei"

    # 6. 底部要点说明
    bullet_box = slide.shapes.add_textbox(Inches(1.5), Inches(6.0), Inches(11.0), Inches(1.2))
    tf_bullet = bullet_box.text_frame
    tf_bullet.word_wrap = True

    # 要点 1
    p_b1 = tf_bullet.paragraphs[0]
    p_b1.space_after = Pt(10)
    run_b1_1 = p_b1.add_run()
    run_b1_1.text = "• 核心瓶颈："
    run_b1_1.font.size = Pt(14)
    run_b1_1.font.bold = True
    run_b1_1.font.name = "Microsoft YaHei"
    
    run_b1_2 = p_b1.add_run()
    run_b1_2.text = "AI 对 SAP 专用数据字典（DDIC）的“幻觉”是阻碍生产力的首要因素。"
    run_b1_2.font.size = Pt(14)
    run_b1_2.font.name = "Microsoft YaHei"

    # 要点 2
    p_b2 = tf_bullet.add_paragraph()
    run_b2_1 = p_b2.add_run()
    run_b2_1.text = "• 工具选型："
    run_b2_1.font.size = Pt(14)
    run_b2_1.font.bold = True
    run_b2_1.font.name = "Microsoft YaHei"
    
    run_b2_2 = p_b2.add_run()
    run_b2_2.text = "Claude Code 在理解力、多文件处理及文档解析上全面优于 GitHub Copilot。"
    run_b2_2.font.size = Pt(14)
    run_b2_2.font.name = "Microsoft YaHei"



# ── Main ──
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

s0 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_1(s0)
prs.save(OUTPUT_PATH)
