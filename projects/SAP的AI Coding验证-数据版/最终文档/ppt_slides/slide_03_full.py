from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# ── Color Palette ──
BLUE_HEADER = RGBColor(0x5B, 0x9B, 0xD5)
BLUE_DARK   = RGBColor(0x4A, 0x86, 0xC8)
CYAN        = RGBColor(0x00, 0xBC, 0xD4)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x33, 0x33, 0x33)
GRAY_TEXT   = RGBColor(0x55, 0x55, 0x55)
GRAY_BAR    = RGBColor(0xB0, 0xBE, 0xC5)
RED         = RGBColor(0xE5, 0x39, 0x35)
GREEN       = RGBColor(0x43, 0xA0, 0x47)
ORANGE      = RGBColor(0xFF, 0x98, 0x00)

ICON_BG  = RGBColor(0xE3, 0xE8, 0xED)
ICON_FG  = RGBColor(0x54, 0x6E, 0x7A)
FONT_NAME = "Microsoft YaHei"
SLIDE_WIDTH = Inches(13.333)
HEADER_H    = Inches(0.75)
SUBTITLE_Y  = Inches(0.95)


def add_header_banner(slide, title_text, bg_color=None):
    if bg_color is None:
        bg_color = BLUE_HEADER
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, HEADER_H
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = bg_color
    banner.line.fill.background()
    tf = banner.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(26)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME


def add_subtitle(slide, text, left, top, width=Inches(12), font_size=Pt(18)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.font.name = FONT_NAME
    return txBox


def add_icon_box(slide, left, top, symbol, size=Inches(0.48)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ICON_BG
    shape.line.fill.background()
    shape.adjustments[0] = 0.25
    tf = shape.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = symbol
    p.font.size = Pt(18)
    p.font.color.rgb = ICON_FG
    p.font.bold = False
    return shape


def add_bullet_item(slide, left, top, symbol, label, description,
                    width=Inches(5.5), desc_size=Pt(13)):
    add_icon_box(slide, left, top, symbol)
    text_left = left + Inches(0.65)
    txBox = slide.shapes.add_textbox(text_left, top - Inches(0.02), width, Inches(0.65))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run_label = p.add_run()
    run_label.text = label + "\uff1a"
    run_label.font.size = Pt(14)
    run_label.font.color.rgb = BLACK
    run_label.font.bold = True
    run_label.font.name = FONT_NAME
    run_desc = p.add_run()
    run_desc.text = description
    run_desc.font.size = desc_size
    run_desc.font.color.rgb = GRAY_TEXT
    run_desc.font.bold = False
    run_desc.font.name = FONT_NAME
    return txBox


def add_conclusion_box(slide, left, top, width, text, font_size=Pt(13)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = BLACK
    run.font.bold = True
    run.font.name = FONT_NAME
    return txBox


def add_table(slide, left, top, width, height, rows, cols, data,
              header_color=None, col_widths=None):
    if header_color is None:
        header_color = BLUE_HEADER
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if data[r][c] is not None else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(11)
                    run.font.name = FONT_NAME
                    if r == 0:
                        run.font.color.rgb = WHITE
                        run.font.bold = True
                    else:
                        run.font.color.rgb = BLACK
                        run.font.bold = False
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else RGBColor(0xF5, 0xF5, 0xF5)
            cell.margin_left = Pt(5)
            cell.margin_right = Pt(5)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
    return table_shape


def add_bar_chart(slide, left, top, width, height,
                  categories, values, title="", bar_colors=None):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', values)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    chart.chart_style = 2
    plot = chart.plots[0]
    plot.gap_width = 100
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CYAN
    series.has_data_labels = True
    dl = series.data_labels
    dl.font.size = Pt(13)
    dl.font.bold = True
    dl.font.color.rgb = BLACK
    dl.number_format = '0.#'
    dl.show_value = True
    dl.label_position = XL_LABEL_POSITION.OUTSIDE_END
    if bar_colors:
        for i, color in enumerate(bar_colors):
            pt = series.points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = color
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(12)
    cat_axis.tick_labels.font.name = FONT_NAME
    cat_axis.major_tick_mark = 2
    cat_axis.format.line.fill.background()
    val_axis = chart.value_axis
    val_axis.visible = False
    val_axis.major_tick_mark = 2
    val_axis.format.line.fill.background()
    val_axis.major_gridlines.format.line.fill.background()
    if title:
        chart.has_title = True
        ct = chart.chart_title.text_frame.paragraphs[0]
        ct.text = title
        ct.font.size = Pt(14)
        ct.font.bold = True
        ct.font.name = FONT_NAME
    else:
        chart.has_title = False
    return chart_frame


def add_callout_label(slide, left, top, text, bg_color=None, font_size=Pt(11)):
    if bg_color is None:
        bg_color = CYAN
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.3), Inches(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME
    return shape


def add_data_card(slide, left, top, width, height, value, label,
                  value_color=None, bg_color=None):
    if value_color is None:
        value_color = CYAN
    if bg_color is None:
        bg_color = WHITE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(3)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run1 = p1.add_run()
    run1.text = str(value)
    run1.font.size = Pt(24)
    run1.font.color.rgb = value_color
    run1.font.bold = True
    run1.font.name = FONT_NAME
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(10)
    run2.font.color.rgb = GRAY_TEXT
    run2.font.bold = False
    run2.font.name = FONT_NAME
    return shape


OUTPUT_PATH = r"C:\Users\Administrator\Desktop\notebook\my_notebooklm_with_ppteditable\projects\SAP的AI Coding验证-数据版\最终文档\ppt_slides\slide_03.pptx"


# ── Slide 1 ──

def build_slide_1(slide):
    # Colors
    DARK_BLUE = RGBColor(0x1A, 0x3B, 0x66)
    GRAY_TEXT = RGBColor(0x55, 0x55, 0x55)
    LIGHT_GRAY_BG = RGBColor(0xF8, 0xF9, 0xFA)
    BORDER_GRAY = RGBColor(0xE0, 0xE0, 0xE0)
    RED = RGBColor(0xD3, 0x2F, 0x2F)
    BLUE = RGBColor(0x19, 0x76, 0xD2)
    GREEN_BG = RGBColor(0xEA, 0xF8, 0xE6)
    GREEN_BORDER = RGBColor(0x8B, 0xF3, 0x69)
    GREEN_ICON = RGBColor(0x4C, 0xAF, 0x50)
    BLACK = RGBColor(0x00, 0x00, 0x00)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    # 1. Title Area
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.0), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "第3页 | 高复杂度：效率反降 60% 的“修复爆炸”"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.0), Inches(0.5))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "案例C - 跨工厂 STO 报表（高复杂度）验证"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(20)
    p.font.color.rgb = GRAY_TEXT

    # 2. Left Panel (Comparison)
    left_panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.8), Inches(4.6), Inches(5.2))
    left_panel.fill.solid()
    left_panel.fill.fore_color.rgb = LIGHT_GRAY_BG
    left_panel.line.color.rgb = BORDER_GRAY
    left_panel.line.width = Pt(1)

    # Vertical Divider in Left Panel
    divider = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(2.9), Inches(2.2), Inches(2.9), Inches(6.6))
    divider.line.color.rgb = BORDER_GRAY
    divider.line.width = Pt(1)

    # Left Column (Claude Code)
    tx_claude = slide.shapes.add_textbox(Inches(0.6), Inches(2.2), Inches(2.3), Inches(0.8))
    tf = tx_claude.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Claude Code\n"
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(16)
    run.font.bold = True
    run = p.add_run()
    run.text = "修复耗时"
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(16)

    # Red Stopwatch Icon
    stopwatch_red = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.55), Inches(3.2), Inches(0.4), Inches(0.4))
    stopwatch_red.fill.background()
    stopwatch_red.line.color.rgb = RED
    stopwatch_red.line.width = Pt(3)
    sw_top_red = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.7), Inches(3.1), Inches(0.1), Inches(0.1))
    sw_top_red.fill.solid()
    sw_top_red.fill.fore_color.rgb = RED
    sw_top_red.line.fill.background()
    sw_hand_red = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.75), Inches(3.4), Inches(1.75), Inches(3.25))
    sw_hand_red.line.color.rgb = RED
    sw_hand_red.line.width = Pt(2)

    # 8 人天
    tx_8 = slide.shapes.add_textbox(Inches(0.6), Inches(3.8), Inches(2.3), Inches(0.8))
    tf = tx_8.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "8 "
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(44)
    run.font.bold = True
    run = p.add_run()
    run.text = "人天"
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(20)

    # -60% Efficiency
    arrow_down = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(1.0), Inches(5.1), Inches(0.2), Inches(0.3))
    arrow_down.fill.solid()
    arrow_down.fill.fore_color.rgb = RED
    arrow_down.line.fill.background()

    tx_eff = slide.shapes.add_textbox(Inches(1.2), Inches(4.9), Inches(1.7), Inches(0.8))
    tf = tx_eff.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "-60%\n"
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RED
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "效率"
    run2.font.name = "Microsoft YaHei"
    run2.font.size = Pt(16)

    # Right Column (Manual)
    tx_manual = slide.shapes.add_textbox(Inches(2.9), Inches(2.2), Inches(2.3), Inches(0.8))
    tf = tx_manual.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "\n手写开发耗时"
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(16)
    run.font.bold = True

    # Blue Stopwatch Icon
    stopwatch_blue = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.85), Inches(3.2), Inches(0.4), Inches(0.4))
    stopwatch_blue.fill.background()
    stopwatch_blue.line.color.rgb = BLUE
    stopwatch_blue.line.width = Pt(3)
    sw_top_blue = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.0), Inches(3.1), Inches(0.1), Inches(0.1))
    sw_top_blue.fill.solid()
    sw_top_blue.fill.fore_color.rgb = BLUE
    sw_top_blue.line.fill.background()
    sw_hand_blue = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(4.05), Inches(3.4), Inches(4.05), Inches(3.25))
    sw_hand_blue.line.color.rgb = BLUE
    sw_hand_blue.line.width = Pt(2)

    # 5 人天
    tx_5 = slide.shapes.add_textbox(Inches(2.9), Inches(3.8), Inches(2.3), Inches(0.8))
    tf = tx_5.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "5 "
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(44)
    run.font.bold = True
    run = p.add_run()
    run.text = "人天"
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(20)

    # 3. Right Top Panel (Trend Chart)
    right_top_panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.5), Inches(1.8), Inches(7.3), Inches(3.6))
    right_top_panel.fill.solid()
    right_top_panel.fill.fore_color.rgb = LIGHT_GRAY_BG
    right_top_panel.line.color.rgb = BORDER_GRAY
    right_top_panel.line.width = Pt(1)

    tx_trend_title = slide.shapes.add_textbox(Inches(5.7), Inches(2.0), Inches(6.0), Inches(0.4))
    tf = tx_trend_title.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "修复-爆炸模式：错误数量非线性增长趋势"
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(16)
    run.font.bold = True

    # Trend Line Segments
    points = [
        (6.7, 4.2), (7.9, 4.2), (7.9, 4.5), (9.1, 4.5), 
        (10.2, 3.2), (10.6, 3.8), (11.8, 2.6)
    ]
    for i in range(len(points) - 1):
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(points[i][0]), Inches(points[i][1]), Inches(points[i+1][0]), Inches(points[i+1][1]))
        line.line.color.rgb = RED
        line.line.width = Pt(4)
        if i == len(points) - 2:
            line.line.end_arrowhead = 2 # Triangle arrowhead

    # Data Points (Circles)
    def add_data_point(x, y, text1, text2):
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x-0.08), Inches(y-0.08), Inches(0.16), Inches(0.16))
        circle.fill.solid()
        circle.fill.fore_color.rgb = RED
        circle.line.color.rgb = WHITE
        circle.line.width = Pt(2)
        
        tx = slide.shapes.add_textbox(Inches(x-0.6), Inches(y-0.7) if y > 3.5 else Inches(y-0.1), Inches(1.2), Inches(0.6))
        tf = tx.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text1 + "\n"
        run.font.name = "Microsoft YaHei"
        run.font.size = Pt(14)
        run.font.bold = True
        if "9" in text1:
            run.font.color.rgb = RED
        run2 = p.add_run()
        run2.text = text2
        run2.font.name = "Microsoft YaHei"
        run2.font.size = Pt(12)

    add_data_point(6.7, 4.2, "3 错误", "(起始)")
    add_data_point(9.1, 4.5, "2 错误", "(初次修复)")
    add_data_point(11.2, 3.3, "9 错误", "(引发新矛盾)")

    # Bomb Icon
    bomb_body = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.7), Inches(2.0), Inches(0.45), Inches(0.45))
    bomb_body.fill.solid()
    bomb_body.fill.fore_color.rgb = BLACK
    bomb_body.line.fill.background()
    
    bomb_cap = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(11.9), Inches(1.9), Inches(0.1), Inches(0.15))
    bomb_cap.fill.solid()
    bomb_cap.fill.fore_color.rgb = GRAY_TEXT
    bomb_cap.line.fill.background()
    
    explosion = slide.shapes.add_shape(MSO_SHAPE.EXPLOSION1, Inches(12.1), Inches(1.8), Inches(0.35), Inches(0.35))
    explosion.fill.solid()
    explosion.fill.fore_color.rgb = RGBColor(0xFF, 0x57, 0x22)
    explosion.line.color.rgb = RGBColor(0xFF, 0xEB, 0x3B)

    # JetBrains Mono Text
    tx_jb = slide.shapes.add_textbox(Inches(8.5), Inches(4.9), Inches(2.0), Inches(0.3))
    tf = tx_jb.text_frame
    p = tf.paragraphs[0]
    p.text = "JetBrains Mono"
    p.font.name = "Consolas"
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY_TEXT

    # 4. Right Bottom Left (Data Dictionary)
    # Database Icon
    for i in range(3):
        cyl = slide.shapes.add_shape(MSO_SHAPE.CAN, Inches(5.5), Inches(5.7 + i*0.12), Inches(0.3), Inches(0.18))
        cyl.fill.solid()
        cyl.fill.fore_color.rgb = RGBColor(0x78, 0x90, 0x9C)
        cyl.line.color.rgb = WHITE
    
    cross = slide.shapes.add_shape(MSO_SHAPE.MATH_MULTIPLY, Inches(5.7), Inches(5.9), Inches(0.15), Inches(0.15))
    cross.fill.solid()
    cross.fill.fore_color.rgb = RED
    cross.line.fill.background()

    tx_dict_title = slide.shapes.add_textbox(Inches(5.9), Inches(5.65), Inches(3.0), Inches(0.4))
    tf = tx_dict_title.text_frame
    p = tf.paragraphs[0]
    p.text = "数据字典重灾区"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(14)
    p.font.bold = True

    tx_dict_desc = slide.shapes.add_textbox(Inches(5.4), Inches(6.1), Inches(3.5), Inches(0.8))
    tf = tx_dict_desc.text_frame
    p1 = tf.paragraphs[0]
    run1 = p1.add_run()
    run1.text = "累计虚构 10+ 项表字段和数据类型\n"
    run1.font.name = "Microsoft YaHei"
    run1.font.size = Pt(12)
    run1.font.bold = True
    run1.font.color.rgb = RED
    
    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = "开发者需大量时间查表对数"
    run2.font.name = "Microsoft YaHei"
    run2.font.size = Pt(12)
    run2.font.color.rgb = BLACK

    # 5. Right Bottom Right (Advantage Box)
    adv_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.2), Inches(5.5), Inches(3.6), Inches(1.4))
    adv_box.fill.solid()
    adv_box.fill.fore_color.rgb = GREEN_BG
    adv_box.line.color.rgb = GREEN_BORDER
    adv_box.line.width = Pt(2)

    # Document Icon
    doc = slide.shapes.add_shape(MSO_SHAPE.FOLDED_CORNER, Inches(9.4), Inches(5.7), Inches(0.25), Inches(0.35))
    doc.fill.solid()
    doc.fill.fore_color.rgb = WHITE
    doc.line.color.rgb = GRAY_TEXT
    
    chk = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.55), Inches(5.9), Inches(0.15), Inches(0.15))
    chk.fill.solid()
    chk.fill.fore_color.rgb = GREEN_ICON
    chk.line.fill.background()

    tx_adv_title = slide.shapes.add_textbox(Inches(9.8), Inches(5.65), Inches(2.8), Inches(0.4))
    tf = tx_adv_title.text_frame
    p = tf.paragraphs[0]
    p.text = "文档解析优势 (Claude)"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(14)
    p.font.bold = True

    tx_adv_desc = slide.shapes.add_textbox(Inches(9.3), Inches(6.1), Inches(3.4), Inches(0.7))
    tf = tx_adv_desc.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "● Claude 虽支持文档解析，但生成的代码架构因底层逻辑矛盾而无法运行。"
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(11)
    run.font.color.rgb = BLACK



# ── Main ──
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

s0 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_1(s0)
prs.save(OUTPUT_PATH)
