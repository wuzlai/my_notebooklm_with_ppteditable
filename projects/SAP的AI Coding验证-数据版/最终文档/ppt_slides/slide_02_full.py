from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# ── Color Palette ──
BLUE_HEADER = RGBColor(0x5B, 0x9B, 0xD5)
BLUE_DARK   = RGBColor(0x4A, 0x86, 0xC8)
CYAN        = RGBColor(0x00, 0xBC, 0xD4)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x33, 0x33, 0x33)
GRAY_TEXT   = RGBColor(0x55, 0x55, 0x55)
GRAY_BAR    = RGBColor(0xB0, 0xBE, 0xC5)
RED         = RGBColor(0xE5, 0x39, 0x35)
GREEN       = RGBColor(0x43, 0xA0, 0x47)
ORANGE      = RGBColor(0xFF, 0x98, 0x00)

ICON_BG  = RGBColor(0xE3, 0xE8, 0xED)
ICON_FG  = RGBColor(0x54, 0x6E, 0x7A)
FONT_NAME = "Microsoft YaHei"
SLIDE_WIDTH = Inches(13.333)
HEADER_H    = Inches(0.75)
SUBTITLE_Y  = Inches(0.95)


def add_header_banner(slide, title_text, bg_color=None):
    if bg_color is None:
        bg_color = BLUE_HEADER
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, HEADER_H
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = bg_color
    banner.line.fill.background()
    tf = banner.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(26)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME


def add_subtitle(slide, text, left, top, width=Inches(12), font_size=Pt(18)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.font.name = FONT_NAME
    return txBox


def add_icon_box(slide, left, top, symbol, size=Inches(0.48)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ICON_BG
    shape.line.fill.background()
    shape.adjustments[0] = 0.25
    tf = shape.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = symbol
    p.font.size = Pt(18)
    p.font.color.rgb = ICON_FG
    p.font.bold = False
    return shape


def add_bullet_item(slide, left, top, symbol, label, description,
                    width=Inches(5.5), desc_size=Pt(13)):
    add_icon_box(slide, left, top, symbol)
    text_left = left + Inches(0.65)
    txBox = slide.shapes.add_textbox(text_left, top - Inches(0.02), width, Inches(0.65))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run_label = p.add_run()
    run_label.text = label + "\uff1a"
    run_label.font.size = Pt(14)
    run_label.font.color.rgb = BLACK
    run_label.font.bold = True
    run_label.font.name = FONT_NAME
    run_desc = p.add_run()
    run_desc.text = description
    run_desc.font.size = desc_size
    run_desc.font.color.rgb = GRAY_TEXT
    run_desc.font.bold = False
    run_desc.font.name = FONT_NAME
    return txBox


def add_conclusion_box(slide, left, top, width, text, font_size=Pt(13)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = BLACK
    run.font.bold = True
    run.font.name = FONT_NAME
    return txBox


def add_table(slide, left, top, width, height, rows, cols, data,
              header_color=None, col_widths=None):
    if header_color is None:
        header_color = BLUE_HEADER
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if data[r][c] is not None else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(11)
                    run.font.name = FONT_NAME
                    if r == 0:
                        run.font.color.rgb = WHITE
                        run.font.bold = True
                    else:
                        run.font.color.rgb = BLACK
                        run.font.bold = False
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else RGBColor(0xF5, 0xF5, 0xF5)
            cell.margin_left = Pt(5)
            cell.margin_right = Pt(5)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
    return table_shape


def add_bar_chart(slide, left, top, width, height,
                  categories, values, title="", bar_colors=None):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', values)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    chart.chart_style = 2
    plot = chart.plots[0]
    plot.gap_width = 100
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CYAN
    series.has_data_labels = True
    dl = series.data_labels
    dl.font.size = Pt(13)
    dl.font.bold = True
    dl.font.color.rgb = BLACK
    dl.number_format = '0.#'
    dl.show_value = True
    dl.label_position = XL_LABEL_POSITION.OUTSIDE_END
    if bar_colors:
        for i, color in enumerate(bar_colors):
            pt = series.points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = color
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(12)
    cat_axis.tick_labels.font.name = FONT_NAME
    cat_axis.major_tick_mark = 2
    cat_axis.format.line.fill.background()
    val_axis = chart.value_axis
    val_axis.visible = False
    val_axis.major_tick_mark = 2
    val_axis.format.line.fill.background()
    val_axis.major_gridlines.format.line.fill.background()
    if title:
        chart.has_title = True
        ct = chart.chart_title.text_frame.paragraphs[0]
        ct.text = title
        ct.font.size = Pt(14)
        ct.font.bold = True
        ct.font.name = FONT_NAME
    else:
        chart.has_title = False
    return chart_frame


def add_callout_label(slide, left, top, text, bg_color=None, font_size=Pt(11)):
    if bg_color is None:
        bg_color = CYAN
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.3), Inches(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_NAME
    return shape


def add_data_card(slide, left, top, width, height, value, label,
                  value_color=None, bg_color=None):
    if value_color is None:
        value_color = CYAN
    if bg_color is None:
        bg_color = WHITE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(3)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run1 = p1.add_run()
    run1.text = str(value)
    run1.font.size = Pt(24)
    run1.font.color.rgb = value_color
    run1.font.bold = True
    run1.font.name = FONT_NAME
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(10)
    run2.font.color.rgb = GRAY_TEXT
    run2.font.bold = False
    run2.font.name = FONT_NAME
    return shape


OUTPUT_PATH = r"C:\Users\Administrator\Desktop\notebook\my_notebooklm_with_ppteditable\projects\SAP的AI Coding验证-数据版\最终文档\ppt_slides\slide_02.pptx"


# ── Slide 1 ──

def build_slide_1(slide):
    # Colors
    DARK_BLUE = RGBColor(0x1B, 0x36, 0x5D)
    DARK_GRAY = RGBColor(0x55, 0x55, 0x55)
    DARK_RED = RGBColor(0x9E, 0x2A, 0x2B)
    LIGHT_RED = RGBColor(0xC1, 0x3C, 0x3D)
    GRAY = RGBColor(0xB0, 0xB0, 0xB0)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    BLACK = RGBColor(0x00, 0x00, 0x00)
    BG_GRAY = RGBColor(0xEB, 0xEF, 0xF2)

    # 1. Slide Background
    slide_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    slide_bg.fill.solid()
    slide_bg.fill.fore_color.rgb = BG_GRAY
    slide_bg.line.fill.background()

    # 2. Header Banner
    add_header_banner(slide, "SAP ABAP AI 效率测评报告", bg_color=DARK_BLUE)

    # 3. Main White Card
    bg_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(0.8), Inches(12.533), Inches(6.4))
    bg_card.fill.solid()
    bg_card.fill.fore_color.rgb = WHITE
    bg_card.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)

    # 4. Main Title & Subtitle
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(10), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "第2页 中等复杂度：深陷“虚构字段”泥潭"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(10), Inches(0.5))
    tf_sub = sub_box.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "案例B - 采购配额维护（中等复杂度）验证"
    p_sub.font.size = Pt(20)
    p_sub.font.bold = True
    p_sub.font.color.rgb = DARK_GRAY

    # 5. Section 1: 业务理解偏差
    icon1 = slide.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(0.8), Inches(0.8))
    icon1.text_frame.text = "🧠❌"
    icon1.text_frame.paragraphs[0].font.size = Pt(32)

    t1 = slide.shapes.add_textbox(Inches(1.8), Inches(2.3), Inches(4), Inches(0.4))
    t1.text_frame.text = "业务理解偏差"
    t1.text_frame.paragraphs[0].font.size = Pt(20)
    t1.text_frame.paragraphs[0].font.bold = True

    b1 = slide.shapes.add_textbox(Inches(1.8), Inches(2.7), Inches(4.5), Inches(1.0))
    tf1 = b1.text_frame
    tf1.word_wrap = True
    p1_1 = tf1.paragraphs[0]
    p1_1.text = "• Copilot 完全混淆“配额”与“货源”概念。"
    p1_1.font.size = Pt(14)
    p1_2 = tf1.add_paragraph()
    p1_2.text = "• 数据模型从底层开始错误，导致逻辑无法构建。"
    p1_2.font.size = Pt(14)

    # 6. Section 2: 严重的“幻觉”现象
    icon2 = slide.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(0.8), Inches(0.8))
    icon2.text_frame.text = "🔗💥"
    icon2.text_frame.paragraphs[0].font.size = Pt(32)

    t2 = slide.shapes.add_textbox(Inches(1.8), Inches(3.7), Inches(4), Inches(0.4))
    t2.text_frame.text = "严重的“幻觉”现象"
    t2.text_frame.paragraphs[0].font.size = Pt(20)
    t2.text_frame.paragraphs[0].font.bold = True

    lbl_real = slide.shapes.add_textbox(Inches(1.8), Inches(4.2), Inches(1.2), Inches(0.3))
    lbl_real.text_frame.text = "实际字段"
    lbl_real.text_frame.paragraphs[0].font.size = Pt(12)

    lbl_fake = slide.shapes.add_textbox(Inches(3.0), Inches(4.2), Inches(1.2), Inches(0.3))
    lbl_fake.text_frame.text = "虚构字段"
    lbl_fake.text_frame.paragraphs[0].font.size = Pt(12)
    lbl_fake.text_frame.paragraphs[0].font.color.rgb = DARK_RED

    box_real = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.8), Inches(4.6), Inches(1.2), Inches(1.0))
    box_real.fill.solid()
    box_real.fill.fore_color.rgb = GRAY
    box_real.line.fill.background()

    box_fake = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.0), Inches(4.6), Inches(2.0), Inches(1.0))
    box_fake.fill.solid()
    box_fake.fill.fore_color.rgb = DARK_RED
    box_fake.line.fill.background()
    tf_fake = box_fake.text_frame
    tf_fake.text = "9处关键字段\n(占比 50%)"
    tf_fake.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf_fake.paragraphs[0].font.size = Pt(14)
    tf_fake.paragraphs[0].font.color.rgb = WHITE
    if len(tf_fake.paragraphs) > 1:
        tf_fake.paragraphs[1].alignment = PP_ALIGN.CENTER
        tf_fake.paragraphs[1].font.size = Pt(14)
        tf_fake.paragraphs[1].font.color.rgb = WHITE

    arrow1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(5.1), Inches(4.7), Inches(1.2), Inches(0.8))
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = LIGHT_RED
    arrow1.line.fill.background()
    tf_arr1 = arrow1.text_frame
    tf_arr1.text = "连锁关键字段\n(21个)"
    tf_arr1.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf_arr1.paragraphs[0].font.size = Pt(11)
    tf_arr1.paragraphs[0].font.color.rgb = WHITE
    if len(tf_arr1.paragraphs) > 1:
        tf_arr1.paragraphs[1].alignment = PP_ALIGN.CENTER
        tf_arr1.paragraphs[1].font.size = Pt(11)
        tf_arr1.paragraphs[1].font.color.rgb = WHITE

    # 7. Section 3: 接口规范缺失
    icon3 = slide.shapes.add_textbox(Inches(0.8), Inches(5.8), Inches(0.8), Inches(0.8))
    icon3.text_frame.text = "🔌❌"
    icon3.text_frame.paragraphs[0].font.size = Pt(32)

    t3 = slide.shapes.add_textbox(Inches(1.8), Inches(5.7), Inches(4), Inches(0.4))
    t3.text_frame.text = "接口规范缺失"
    t3.text_frame.paragraphs[0].font.size = Pt(20)
    t3.text_frame.paragraphs[0].font.bold = True

    b3 = slide.shapes.add_textbox(Inches(1.8), Inches(6.1), Inches(4.5), Inches(1.0))
    tf3 = b3.text_frame
    tf3.word_wrap = True
    p3_1 = tf3.paragraphs[0]
    p3_1.text = "• AI 无法识别 SAP 函数模块 (FM) 仅接受 DDIC 类型的硬性规则。"
    p3_1.font.size = Pt(14)
    p3_2 = tf3.add_paragraph()
    p3_2.text = "• 数据类型不匹配导致接口调用必然失败。"
    p3_2.font.size = Pt(14)

    # 8. Section 4: 开发成本倒挂
    icon4 = slide.shapes.add_textbox(Inches(6.2), Inches(5.8), Inches(0.8), Inches(0.8))
    icon4.text_frame.text = "⚖️❌"
    icon4.text_frame.paragraphs[0].font.size = Pt(32)

    t4 = slide.shapes.add_textbox(Inches(7.2), Inches(5.7), Inches(4), Inches(0.4))
    t4.text_frame.text = "开发成本倒挂"
    t4.text_frame.paragraphs[0].font.size = Pt(20)
    t4.text_frame.paragraphs[0].font.bold = True

    b4 = slide.shapes.add_textbox(Inches(7.2), Inches(6.1), Inches(5.0), Inches(1.0))
    tf4 = b4.text_frame
    tf4.word_wrap = True
    p4_1 = tf4.paragraphs[0]
    p4_1.text = "• 修复 AI 错误的代码成本已超过直接重写。"
    p4_1.font.size = Pt(14)
    p4_2 = tf4.add_paragraph()
    p4_2.text = "• AI 辅助在该场景下失去价值，带来负面效率。"
    p4_2.font.size = Pt(14)

    # 9. Large Diagram
    diag_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.4), Inches(2.2), Inches(6.3), Inches(3.3))
    diag_bg.fill.solid()
    diag_bg.fill.fore_color.rgb = RGBColor(0xF9, 0xF9, 0xF9)
    diag_bg.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)

    l_real = slide.shapes.add_textbox(Inches(6.5), Inches(2.3), Inches(1.3), Inches(0.3))
    l_real.text_frame.text = "实际字段"
    l_real.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    l_real.text_frame.paragraphs[0].font.size = Pt(12)

    l_fake = slide.shapes.add_textbox(Inches(7.85), Inches(2.3), Inches(1.8), Inches(0.3))
    l_fake.text_frame.text = "虚构字段"
    l_fake.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    l_fake.text_frame.paragraphs[0].font.size = Pt(12)
    l_fake.text_frame.paragraphs[0].font.color.rgb = DARK_RED

    l_err = slide.shapes.add_textbox(Inches(10.0), Inches(2.3), Inches(2.5), Inches(0.3))
    l_err.text_frame.text = "连锁语法错误"
    l_err.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    l_err.text_frame.paragraphs[0].font.size = Pt(12)
    l_err.text_frame.paragraphs[0].font.color.rgb = DARK_RED

    # Gray Stack
    gray_top = 2.7
    for i in range(4):
        gb = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(gray_top + i*0.55), Inches(1.3), Inches(0.5))
        gb.fill.solid()
        gb.fill.fore_color.rgb = GRAY
        gb.line.color.rgb = WHITE
        gb.line.width = Pt(1)

    # Red Stack
    red_texts = ["虚构: QUOTA_MATNR", "虚构: SOURCE_VENDOR", "虚构: VALID_DATE_FROM"]
    red_top = 2.7
    for i, text in enumerate(red_texts):
        rb = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.85), Inches(red_top + i*0.733), Inches(1.8), Inches(0.68))
        rb.fill.solid()
        rb.fill.fore_color.rgb = DARK_RED
        rb.line.color.rgb = WHITE
        rb.line.width = Pt(1)
        tf = rb.text_frame
        tf.text = text
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.size = Pt(10)
        tf.paragraphs[0].font.color.rgb = WHITE

    # Arrow
    arrow2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(9.7), Inches(3.6), Inches(0.25), Inches(0.4))
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = LIGHT_RED
    arrow2.line.fill.background()

    # Red Grid
    err_left = 10.0
    err_top = 2.7
    err_w = 2.5
    err_h = 2.2
    
    base_red = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(err_left), Inches(err_top), Inches(err_w), Inches(err_h))
    base_red.fill.solid()
    base_red.fill.fore_color.rgb = LIGHT_RED
    base_red.line.fill.background()

    # Grid Lines
    for i in range(1, 4):
        y = err_top + i * (err_h / 4)
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(err_left), Inches(y), Inches(err_left + err_w), Inches(y))
        line.line.color.rgb = WHITE
        line.line.width = Pt(1.5)
        
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(err_left + err_w*0.25), Inches(err_top), Inches(err_left + err_w*0.25), Inches(err_top + err_h/4)).line.color.rgb = WHITE
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(err_left + err_w*0.5), Inches(err_top), Inches(err_left + err_w*0.5), Inches(err_top + err_h/4)).line.color.rgb = WHITE
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(err_left + err_w*0.75), Inches(err_top), Inches(err_left + err_w*0.75), Inches(err_top + err_h/4)).line.color.rgb = WHITE
    
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(err_left + err_w*0.33), Inches(err_top + err_h/4), Inches(err_left + err_w*0.33), Inches(err_top + err_h/2)).line.color.rgb = WHITE
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(err_left + err_w*0.66), Inches(err_top + err_h/4), Inches(err_left + err_w*0.66), Inches(err_top + err_h/2)).line.color.rgb = WHITE
    
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(err_left + err_w*0.25), Inches(err_top + err_h/2), Inches(err_left + err_w*0.25), Inches(err_top + err_h*0.75)).line.color.rgb = WHITE
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(err_left + err_w*0.5), Inches(err_top + err_h/2), Inches(err_left + err_w*0.5), Inches(err_top + err_h*0.75)).line.color.rgb = WHITE
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(err_left + err_w*0.75), Inches(err_top + err_h/2), Inches(err_left + err_w*0.75), Inches(err_top + err_h*0.75)).line.color.rgb = WHITE
    
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(err_left + err_w*0.33), Inches(err_top + err_h*0.75), Inches(err_left + err_w*0.33), Inches(err_top + err_h)).line.color.rgb = WHITE
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(err_left + err_w*0.66), Inches(err_top + err_h*0.75), Inches(err_left + err_w*0.66), Inches(err_top + err_h)).line.color.rgb = WHITE

    # Center Block
    err_center = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10.2), Inches(3.1), Inches(2.1), Inches(1.4))
    err_center.fill.solid()
    err_center.fill.fore_color.rgb = DARK_RED
    err_center.line.color.rgb = WHITE
    err_center.line.width = Pt(1.5)
    tf_err = err_center.text_frame
    tf_err.text = "连锁语法错误\n(21个)"
    tf_err.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf_err.paragraphs[0].font.size = Pt(16)
    tf_err.paragraphs[0].font.bold = True
    tf_err.paragraphs[0].font.color.rgb = WHITE
    if len(tf_err.paragraphs) > 1:
        tf_err.paragraphs[1].alignment = PP_ALIGN.CENTER
        tf_err.paragraphs[1].font.size = Pt(16)
        tf_err.paragraphs[1].font.bold = True
        tf_err.paragraphs[1].font.color.rgb = WHITE

    # Bottom Text
    bot_text = slide.shapes.add_textbox(Inches(6.4), Inches(5.0), Inches(6.3), Inches(0.4))
    tf_bot = bot_text.text_frame
    tf_bot.text = "Claude Code 虚构了关键字段，引发 21 个连锁语法错误。"
    tf_bot.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf_bot.paragraphs[0].font.size = Pt(14)
    tf_bot.paragraphs[0].font.color.rgb = BLACK



# ── Main ──
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

s0 = prs.slides.add_slide(prs.slide_layouts[6])
build_slide_1(s0)
prs.save(OUTPUT_PATH)
